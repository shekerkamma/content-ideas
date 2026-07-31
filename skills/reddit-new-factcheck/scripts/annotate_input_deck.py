#!/usr/bin/env python3
"""Create an annotated review copy of an input PPTX from fact-check artifacts."""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - exercised by CLI users
    raise SystemExit("python-pptx is required to annotate PPTX files") from exc


COLORS = {
    "navy": RGBColor(10, 14, 20),
    "panel": RGBColor(18, 24, 38),
    "white": RGBColor(245, 248, 252),
    "muted": RGBColor(142, 153, 178),
    "green": RGBColor(52, 211, 153),
    "amber": RGBColor(251, 191, 36),
    "red": RGBColor(251, 113, 133),
    "blue": RGBColor(56, 189, 248),
}

BRANDED_TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE") or os.path.expanduser("~/.claude/templates/branded-template.pptx"))


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_branded_template() -> None:
    if not BRANDED_TEMPLATE.exists():
        raise SystemExit(
            "Branded Canva template is required for client-ready PPTX output. "
            "Set BRANDED_PPTX_TEMPLATE or install ~/.claude/templates/branded-template.pptx."
        )


def short(text: str, limit: int) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "..."


def display_claim_id(value: str) -> str:
    match = re.search(r"(\d+)$", value)
    return f"C-{int(match.group(1)):03d}" if match else value.upper()


def slide_index(locator: str, slide_count: int) -> int | None:
    match = re.search(r"\b(?:Slide|Section)\s+(\d+)\b", locator, re.I)
    if not match:
        return None
    index = int(match.group(1)) - 1
    if 0 <= index < slide_count:
        return index
    return None


def evidence_lookup(evidence_pack: dict[str, Any] | None) -> dict[str, dict[str, Any]]:
    if not evidence_pack:
        return {}
    return {item["claim_id"]: item for item in evidence_pack.get("claims", [])}


def verdict_for(claim: dict[str, Any], evidence: dict[str, Any] | None) -> tuple[str, str, RGBColor]:
    raw_verdict = (evidence or {}).get("suggested_reddit_verdict", "")
    if "contradict" in raw_verdict:
        return "Contradicted on Reddit", "review", COLORS["red"]
    if claim.get("primary_source_required"):
        if "reddit_signal" in raw_verdict or "weak_reddit_support" in raw_verdict or "weak_qualified_reddit_support" in raw_verdict:
            return "Reddit signal + primary source", "mixed", COLORS["blue"]
        return "Primary source required", "source", COLORS["blue"]
    if raw_verdict in {"supported_by_reddit_evidence", "qualified_reddit_support", "evidence_found"}:
        count = int((evidence or {}).get("evidence_count") or 0)
        confidence = "high" if count >= 3 else "medium"
        return "Reddit-supported", confidence, COLORS["green"]
    if "weak_reddit_support" in raw_verdict or "weak_qualified_reddit_support" in raw_verdict:
        return "Weak Reddit support", "low", COLORS["amber"]
    if raw_verdict in {"no_reddit_evidence_found", "no_qualified_reddit_evidence"}:
        return "No Reddit evidence", "low", COLORS["amber"]
    return "Pending Reddit evidence", "pending", COLORS["amber"]


def grouped_claims(claims: list[dict[str, Any]], slide_count: int) -> tuple[dict[int, list[dict[str, Any]]], list[dict[str, Any]]]:
    by_slide: dict[int, list[dict[str, Any]]] = {}
    unmapped: list[dict[str, Any]] = []
    for claim in claims:
        index = slide_index(claim.get("locator", ""), slide_count)
        if index is None:
            unmapped.append(claim)
            continue
        by_slide.setdefault(index, []).append(claim)
    return by_slide, unmapped


def set_text_style(run, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_rail(slide, claims: list[dict[str, Any]], evidence_by_id: dict[str, dict[str, Any]], slide_width: int) -> None:
    visible = claims[:3]
    more = len(claims) - len(visible)
    x = slide_width - Inches(3.05)
    y = Inches(0.38)
    w = Inches(2.75)
    h = Inches(0.78 + (0.78 * len(visible)) + (0.24 if more else 0))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["panel"]
    box.line.color.rgb = COLORS["blue"]
    box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "FACT CHECK"
    set_text_style(r, 9, COLORS["blue"], True)

    for claim in visible:
        evidence = evidence_by_id.get(claim["id"])
        verdict, confidence, color = verdict_for(claim, evidence)
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"{display_claim_id(claim['id'])}  {verdict} ({confidence})"
        set_text_style(r, 7, color, True)

        p = tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = short(claim["claim"], 92)
        set_text_style(r, 6, COLORS["white"])

    if more:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"+ {more} more in Fact Check Notes"
        set_text_style(r, 6, COLORS["muted"], True)


def add_text_box(slide, text: str, x, y, w, h, size: int = 14, color: RGBColor | None = None, bold: bool = False):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_text_style(r, size, color or COLORS["white"], bold)
    return shape


def rect(slide, x, y, w, h, color, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def blank_layout(prs):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def add_notes_slide(prs, source_slide_number: int | str, claims: list[dict[str, Any]], evidence_by_id: dict[str, dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]

    add_text_box(
        slide,
        f"Fact Check Notes - Source Slide {source_slide_number}",
        Inches(0.55),
        Inches(0.35),
        Inches(8.8),
        Inches(0.4),
        size=18,
        color=COLORS["blue"],
        bold=True,
    )

    y = 0.95
    for claim in claims[:4]:
        evidence = evidence_by_id.get(claim["id"])
        verdict, confidence, color = verdict_for(claim, evidence)
        evidence_items = (evidence or {}).get("evidence") or []
        top = evidence_items[0] if evidence_items else None
        source = top.get("source_json", "") if top else ""
        excerpt = top.get("excerpt", "") if top else ""
        body = [
            f"{display_claim_id(claim['id'])}: {verdict} ({confidence})",
            short(claim["claim"], 185),
            f"Reddit query: {short(claim.get('reddit_query', ''), 130)}",
        ]
        if top:
            body.append(f"Top Reddit evidence: {short(excerpt, 165)}")
            body.append(f"Source: {short(source, 120)}")
        elif claim.get("primary_source_required"):
            body.append("Reddit can provide community signal; this claim still needs a primary source.")
        else:
            body.append("No scored Reddit evidence is attached yet.")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(12.2), Inches(1.18))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["panel"]
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        for index, line in enumerate(body):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            set_text_style(r, 8 if index else 9, color if index == 0 else COLORS["white"], index == 0)
        y += 1.34


def verdict_label(raw_verdict: str, claim: dict[str, Any]) -> str:
    if "contradict" in raw_verdict:
        return "Contradicted"
    if claim.get("primary_source_required"):
        if "reddit_signal" in raw_verdict or "weak_reddit_support" in raw_verdict or "weak_qualified_reddit_support" in raw_verdict:
            return "Reddit + Primary"
        return "Primary Required"
    if raw_verdict in {"supported_by_reddit_evidence", "qualified_reddit_support", "evidence_found"}:
        return "Reddit-Supported"
    if "weak_reddit_support" in raw_verdict or "weak_qualified_reddit_support" in raw_verdict:
        return "Weak Reddit Support"
    if raw_verdict in {"no_reddit_evidence_found", "no_qualified_reddit_evidence"}:
        return "No Reddit Evidence"
    return "Pending"


def add_summary_slide_one(prs, claims: list[dict[str, Any]], evidence_by_id: dict[str, dict[str, Any]], evidence_pack: dict[str, Any] | None) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]

    add_text_box(slide, "Fact Check Summary", Inches(0.55), Inches(0.38), Inches(7.2), Inches(0.45), size=20, color=COLORS["blue"], bold=True)
    add_text_box(slide, "Reddit evidence has been attached. Hard facts still need official-source review.", Inches(0.55), Inches(0.82), Inches(7.7), Inches(0.28), size=9.5, color=COLORS["white"])

    verdict_counts = {
        "Reddit-Supported": 0,
        "Weak Reddit Support": 0,
        "No Reddit Evidence": 0,
        "Primary Required": 0,
        "Contradicted": 0,
        "Reddit + Primary": 0,
        "Pending": 0,
    }
    for claim in claims:
        evidence = evidence_by_id.get(claim["id"])
        raw_verdict = (evidence or {}).get("suggested_reddit_verdict", "")
        verdict_counts[verdict_label(raw_verdict, claim)] += 1

    cards = [
        (str(verdict_counts["Reddit-Supported"]), "Reddit-supported", COLORS["green"]),
        (str(verdict_counts["Weak Reddit Support"]), "Weak Reddit support", COLORS["amber"]),
        (str(verdict_counts["No Reddit Evidence"]), "No Reddit evidence", COLORS["red"]),
        (str(verdict_counts["Primary Required"]), "Primary required", COLORS["blue"]),
    ]
    x = 0.55
    for value, label, color in cards:
        rect(slide, Inches(x), Inches(1.28), Inches(1.55), Inches(0.82), COLORS["panel"])
        rect(slide, Inches(x), Inches(1.28), Inches(0.08), Inches(0.82), color)
        add_text_box(slide, value, Inches(x + 0.1), Inches(1.42), Inches(0.48), Inches(0.18), size=18, color=color, bold=True)
        add_text_box(slide, label, Inches(x + 0.58), Inches(1.46), Inches(0.85), Inches(0.18), size=6.7, color=COLORS["white"], bold=True)
        x += 1.75

    source_count = len((evidence_pack or {}).get("source_files", []))
    comment_count = 0
    for source_file in (evidence_pack or {}).get("source_files", []):
        try:
            data = load_json(Path(source_file))
            comment_count += len(data.get("comments", []))
        except Exception:
            pass
    add_text_box(slide, f"{source_count} Reddit threads parsed", Inches(0.58), Inches(2.32), Inches(2.4), Inches(0.2), size=11, color=COLORS["white"], bold=True)
    add_text_box(slide, f"{comment_count} comments scored", Inches(0.58), Inches(2.56), Inches(2.4), Inches(0.2), size=11, color=COLORS["white"], bold=True)
    add_text_box(slide, "Evidence is strongest for operator pain, objections, and governance language.", Inches(0.58), Inches(2.86), Inches(5.5), Inches(0.22), size=10, color=COLORS["muted"])

    top = sorted(
        (
            (claim, evidence_by_id.get(claim["id"]))
            for claim in claims
            if evidence_by_id.get(claim["id"])
        ),
        key=lambda pair: int(pair[1].get("evidence_count", 0)),
        reverse=True,
    )[:5]
    y = 3.28
    for claim, evidence in top:
        raw_verdict = (evidence or {}).get("suggested_reddit_verdict", "")
        add_text_box(
            slide,
            f"{display_claim_id(claim['id'])} {verdict_label(raw_verdict, claim)}",
            Inches(0.58),
            Inches(y),
            Inches(2.1),
            Inches(0.18),
            size=7.4,
            color=COLORS["green"] if "Supported" in verdict_label(raw_verdict, claim) else COLORS["amber"],
            bold=True,
        )
        add_text_box(slide, short(claim["claim"], 108), Inches(2.1), Inches(y), Inches(5.55), Inches(0.18), size=7.2, color=COLORS["white"])
        y += 0.28


def add_summary_slide_two(prs, claims: list[dict[str, Any]], evidence_by_id: dict[str, dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]

    add_text_box(slide, "What Changed In The Deck", Inches(0.55), Inches(0.38), Inches(7.2), Inches(0.45), size=20, color=COLORS["blue"], bold=True)
    add_text_box(slide, "The source deck is annotated slide-by-slide, and the first slides summarize the verdict mix.", Inches(0.55), Inches(0.82), Inches(7.7), Inches(0.28), size=9.5, color=COLORS["white"])

    groups = [
        ("Strong Reddit support", [claim for claim in claims if verdict_label((evidence_by_id.get(claim["id"]) or {}).get("suggested_reddit_verdict", ""), claim) == "Reddit-Supported"][:4], COLORS["green"]),
        ("Weak / no Reddit evidence", [claim for claim in claims if verdict_label((evidence_by_id.get(claim["id"]) or {}).get("suggested_reddit_verdict", ""), claim) in {"Weak Reddit Support", "No Reddit Evidence"}][:4], COLORS["amber"]),
        ("Primary-source required", [claim for claim in claims if claim.get("primary_source_required")][:4], COLORS["blue"]),
    ]

    x = 0.55
    for title, items, color in groups:
        rect(slide, Inches(x), Inches(1.28), Inches(2.38), Inches(4.9), COLORS["panel"])
        rect(slide, Inches(x), Inches(1.28), Inches(2.38), Inches(0.1), color)
        add_text_box(slide, title, Inches(x + 0.12), Inches(1.47), Inches(2.05), Inches(0.2), size=9.5, color=COLORS["white"], bold=True)
        y = 1.82
        for claim in items:
            evidence = evidence_by_id.get(claim["id"])
            raw_verdict = (evidence or {}).get("suggested_reddit_verdict", "")
            add_text_box(slide, display_claim_id(claim["id"]), Inches(x + 0.12), Inches(y), Inches(0.42), Inches(0.14), size=6.8, color=color, bold=True)
            add_text_box(slide, short(claim["claim"], 42), Inches(x + 0.52), Inches(y), Inches(1.65), Inches(0.18), size=6.6, color=COLORS["white"])
            add_text_box(slide, verdict_label(raw_verdict, claim), Inches(x + 0.12), Inches(y + 0.16), Inches(1.4), Inches(0.12), size=5.8, color=COLORS["muted"])
            y += 0.78
        x += 2.58

    add_text_box(slide, "Slides 3+ preserve the original deck structure with rails and notes appended where claims need detail.", Inches(0.55), Inches(6.55), Inches(7.3), Inches(0.2), size=8.6, color=COLORS["muted"])


def move_last_n_slides_to_front(prs, count: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    moved = [sldIdLst[-count + offset] for offset in range(count)]
    for _ in range(count):
        del sldIdLst[-1]
    for item in moved:
        sldIdLst.insert(0, item)


def annotate_deck(input_pptx: Path, claim_pack: Path, out: Path, evidence_pack: Path | None = None) -> dict[str, Any]:
    prs = Presentation(str(input_pptx))
    claims = load_json(claim_pack).get("claims", [])
    evidence_data = load_json(evidence_pack) if evidence_pack and evidence_pack.exists() else None
    evidence_by_id = evidence_lookup(evidence_data)
    by_slide, unmapped = grouped_claims(claims, len(prs.slides))

    for index, slide_claims in by_slide.items():
        add_rail(prs.slides[index], slide_claims, evidence_by_id, prs.slide_width)

    for index, slide_claims in sorted(by_slide.items()):
        for start in range(0, len(slide_claims), 4):
            add_notes_slide(prs, index + 1, slide_claims[start : start + 4], evidence_by_id)

    if unmapped:
        for start in range(0, len(unmapped), 4):
            add_notes_slide(prs, "Unmapped", unmapped[start : start + 4], evidence_by_id)

    add_summary_slide_two(prs, claims, evidence_by_id)
    add_summary_slide_one(prs, claims, evidence_by_id, evidence_data)
    move_last_n_slides_to_front(prs, 2)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {
        "input_pptx": str(input_pptx),
        "output_pptx": str(out),
        "claims": len(claims),
        "annotated_source_slides": len(by_slide),
        "unmapped_claims": len(unmapped),
    }


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input-pptx", required=True, type=Path)
    parser.add_argument("--claim-pack", required=True, type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--evidence-pack", type=Path)
    parser.add_argument("--summary-json", type=Path)
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv or sys.argv[1:])
    ensure_branded_template()
    summary = annotate_deck(args.input_pptx, args.claim_pack, args.out, args.evidence_pack)
    if args.summary_json:
        args.summary_json.parent.mkdir(parents=True, exist_ok=True)
        args.summary_json.write_text(json.dumps(summary, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
