#!/usr/bin/env python3
"""Native editable v14 deck driven by the validated per-slide design contracts."""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

RUN=Path(__file__).resolve().parent
TOTAL_SLIDES=9

NAVY=RGBColor(25,42,59); INK=RGBColor(35,49,64); WHITE=RGBColor(255,255,255)
GREEN=RGBColor(27,132,77); PALEGREEN=RGBColor(228,241,234); DGREEN=RGBColor(15,91,53)
CORAL=RGBColor(205,70,53); PALECORAL=RGBColor(249,232,229); GOLD=RGBColor(205,157,35)
PALEGOLD=RGBColor(249,243,222); SLATE=RGBColor(91,111,132); PALE=RGBColor(246,248,248); LINE=RGBColor(211,220,217)
STATUS={'verified':GREEN,'attributed company claim':SLATE,'qualified interpretation':GOLD,'insufficient evidence':CORAL}

def clean(s): return ' '.join(str(s).split())
def words(s,n):
 a=clean(s).split(); return ' '.join(a if len(a)<=n else a[:n]+['…'])
def rect(slide,x,y,w,h,fill=WHITE,line=LINE,round=False,width=.7):
 sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh
def tx(slide,s,x,y,w,h,pt=12,color=INK,bold=False,align=PP_ALIGN.LEFT,font='Aptos',margin=.02):
 sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(margin); tf.margin_top=tf.margin_bottom=Inches(.01); tf.vertical_anchor=MSO_ANCHOR.TOP
 p=tf.paragraphs[0]; p.text=clean(s); p.font.name=font; p.font.size=Pt(pt); p.font.bold=bold; p.font.color.rgb=color; p.alignment=align; return sh
def line(slide,x1,y1,x2,y2,color=LINE,width=1):
 sh=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2)); sh.line.color.rgb=color; sh.line.width=Pt(width); return sh
def label(slide,s,x,y,w,color=GREEN): tx(slide,s.upper(),x,y,w,.20,9,color,True)
def header(slide,s,n):
 rect(slide,0,0,13.333,.055,NAVY,NAVY); tx(slide,'SYNTHETIC NATIVE ARCHETYPE · REPLACE THROUGH CONTENT ENVELOPE',.62,.15,8.2,.18,8,CORAL,True)
 title=words(s['assertion_title'],9); pt=24
 tx(slide,title,.62,.38,12.0,.85,pt,INK,True,font='Aptos Display'); tx(slide,words(s['analytical_question'],10),.62,1.28,11.8,.22,10,SLATE)
 rect(slide,.62,1.57,1.05,.035,GREEN,GREEN); tx(slide,'Compound competitor analysis · v14-derived native archetype',.62,7.17,7,.14,7,SLATE); tx(slide,f'{n} / {TOTAL_SLIDES}',12.02,7.17,.65,.14,7,INK,align=PP_ALIGN.RIGHT)
def decision_rail(slide,s):
 y=5.67; line(slide,.62,y,12.70,y,LINE,.9)
 cols=[(.62,3.55,'IMPLICATION',s['implication'],GREEN,10),(4.42,3.62,'DECISION',s['decision'],DGREEN,10),(8.28,4.40,'TRIGGER · OWNER',s['trigger']+' · '+s['owner'],CORAL,12)]
 for x,w,lab,body,c,nw in cols:
  label(slide,lab,x,y+.14,w,c); tx(slide,words(body,nw),x,y+.39,w,.54,10.0,INK,lab=='DECISION')
def status_dot(slide,status,x,y):
 c=STATUS[status]; sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(.12),Inches(.12)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.color.rgb=c
def evidence_row(slide,e,x,y,w,h=.60):
 c=STATUS[e['evidence_status']]; rect(slide,x,y,w,h,WHITE,LINE,True); rect(slide,x,y,.08,h,c,c)
 tx(slide,words(e['label'],4),x+.22,y+.10,1.50,.30,8,c,True); tx(slide,words(e['claim'],14),x+1.76,y+.08,w-1.96,.40,9.5,INK)

def cover(slide,s,n):
 rect(slide,0,0,13.333,7.5,WHITE,WHITE); rect(slide,0,0,13.333,.055,NAVY,NAVY)
 tx(slide,'SYNTHETIC NATIVE ARCHETYPE · REPLACE THROUGH CONTENT ENVELOPE',.62,.20,8.2,.18,8,CORAL,True)
 tx(slide,words(s['assertion_title'],8),.70,.88,7.35,1.02,28,INK,True,font='Aptos Display'); tx(slide,words(s['executive_answer'],14),.70,2.15,7.0,.90,16,GREEN,True)
 tx(slide,words(s['comparison_logic'],27),.70,3.32,6.95,.90,11.5,SLATE)
 rect(slide,8.55,.82,4.10,4.25,PALE,LINE,True); label(slide,'DECISION LOGIC',8.90,1.12,2.2)
 for i,e in enumerate(s['exhibit']['elements'][:3]):
  y=1.66+i*.95; rect(slide,8.90,y,.40,.40,[GREEN,GOLD,CORAL][i],[GREEN,GOLD,CORAL][i],True); tx(slide,str(i+1),8.90,y+.09,.40,.16,8,WHITE,True,PP_ALIGN.CENTER); tx(slide,words(e,11),9.52,y-.01,2.70,.54,10.5,INK,i==2)
 decision_rail(slide,s); tx(slide,'Compound competitor analysis · v14-derived native archetype',.62,7.17,7,.14,7,SLATE); tx(slide,f'1 / {TOTAL_SLIDES}',12.02,7.17,.65,.14,7,INK,align=PP_ALIGN.RIGHT)

def ladder(slide,s,c):
 label(slide,'EVIDENCE LADDER',.72,1.77,1.7)
 for i,e in enumerate(s['evidence_blocks'][:4]): evidence_row(slide,e,.72,2.08+i*.72,7.20,.58)
 rect(slide,8.25,1.80,4.38,3.55,PALE,LINE,True); label(slide,'BOUNDED VERDICT',8.55,2.05,1.8); tx(slide,words(s['executive_answer'],20),8.55,2.42,3.75,1.30,14,INK,True,font='Aptos Display')
 line(slide,8.55,3.90,12.15,3.90,GOLD,1.3); label(slide,'FALSIFIER',8.55,4.13,1.0,CORAL); tx(slide,words(s['counterargument'],18),8.55,4.43,3.68,.64,9.8,SLATE)
 decision_rail(slide,s)

def mechanism(slide,s,c):
 label(slide,'CAUSAL MECHANISM',.72,1.78,2.35)
 blocks=list(s['evidence_blocks'][:3]); bodies=[(e['label'],e['claim'],STATUS[e['evidence_status']]) for e in blocks]+[('DECISION GATE',s['trigger'],GREEN)]
 for i,(lab,body,col) in enumerate(bodies):
  x=.72+i*3.07; rect(slide,x,2.18,2.63,2.40,WHITE,LINE,True); tx(slide,f'{i+1:02d}',x+.18,2.38,.42,.20,9,col,True); label(slide,words(lab,4),x+.18,2.80,2.12,col); tx(slide,words(body,10),x+.18,3.18,2.18,.94,10.5,INK,i==3)
  if i<3: tx(slide,'→',x+2.68,3.17,.30,.28,16,GREEN,True,PP_ALIGN.CENTER)
 label(slide,'COUNTERARGUMENT',.72,4.91,1.65,CORAL); tx(slide,words(s['counterargument'],14),2.38,4.86,9.98,.34,9.7,SLATE)
 decision_rail(slide,s)

def heatmap(slide,s,c):
 label(slide,'THREAT × BUYING ARENA',.72,1.78,2.65)
 cols=['GOVT / PSU','OEM / TIER-1','FLEET / MINING']; x0=4.35
 for j,h in enumerate(cols): tx(slide,h,x0+j*2.63,1.79,2.28,.18,8.5,SLATE,True,PP_ALIGN.CENTER)
 for i,e in enumerate(s['evidence_blocks'][:4]):
  y=2.20+i*.66; col=STATUS[e['evidence_status']]; status_dot(slide,e['evidence_status'],.74,y+.16); tx(slide,words(e['label'],5),.94,y+.08,3.12,.34,9.2,INK,True)
  for j in range(3):
   k=(i+j)%3; fills=[PALE,PALEGOLD,PALEGREEN]; rect(slide,x0+j*2.63,y,2.28,.48,fills[k],WHITE,True); tx(slide,['LOW','WATCH','HIGH'][k],x0+j*2.63,y+.13,2.28,.18,8.5,[SLATE,GOLD,DGREEN][k],True,PP_ALIGN.CENTER)
 rect(slide,.72,4.94,11.90,.43,PALE,LINE,True); tx(slide,words(s['comparison_logic'],18),.90,5.04,11.50,.23,9.2,INK)
 decision_rail(slide,s)

def matrix(slide,s,c):
 label(slide,'COMPARISON FIELD',.72,1.78,1.8)
 pos=[(.72,2.12),(6.75,2.12),(.72,3.58),(6.75,3.58)]
 for i,(e,(x,y)) in enumerate(zip(s['evidence_blocks'][:4],pos)):
  col=STATUS[e['evidence_status']]; rect(slide,x,y,5.70,1.18,WHITE,LINE,True); rect(slide,x,y,.09,1.18,col,col); tx(slide,words(e['label'],5),x+.25,y+.15,1.55,.34,8.3,col,True); tx(slide,words(e['claim'],16),x+1.88,y+.13,3.53,.70,10.1,INK)
 label(slide,'COMPARISON RULE',.72,5.02,1.7); tx(slide,words(s['comparison_logic'],16),2.32,4.97,10.04,.34,9.4,SLATE)
 decision_rail(slide,s)

def funnel(slide,s,c):
 label(slide,'QUALIFICATION FUNNEL',.72,1.78,1.9)
 widths=[10.8,8.8,6.8,4.8]
 for i,e in enumerate(s['evidence_blocks'][:4]):
  w=widths[i]; x=(13.333-w)/2; y=2.12+i*.68; col=STATUS[e['evidence_status']]; sh=slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID,Inches(x),Inches(y),Inches(w),Inches(.52)); sh.fill.solid(); sh.fill.fore_color.rgb=PALE; sh.line.color.rgb=col
  tx(slide,words(e['label'],3),x+.26,y+.12,1.62,.18,8,col,True); tx(slide,words(e['claim'],max(4,8-i)),x+1.92,y+.07,w-2.18,.36,9.3,INK,i==3)
 tx(slide,'SCOPE NARROWS AS ACCOUNTABILITY RISES',4.05,5.03,5.20,.18,8,GREEN,True,PP_ALIGN.CENTER); decision_rail(slide,s)

def timeline(slide,s,c):
 label(slide,'EVENT SEQUENCE',.72,1.78,1.6); line(slide,1.02,3.42,12.25,3.42,SLATE,1.6)
 for i,e in enumerate(s['evidence_blocks'][:4]):
  x=1.02+i*3.02; col=STATUS[e['evidence_status']]; rect(slide,x,3.17,.50,.50,col,col,True); tx(slide,str(i+1),x,3.30,.50,.17,8,WHITE,True,PP_ALIGN.CENTER); tx(slide,words(e['label'],4),x-.10,2.18,2.48,.36,8.2,col,True); tx(slide,words(e['claim'],11),x-.10,2.62,2.50,.58,9.2,INK)
  tx(slide,words(s['exhibit']['elements'][min(i,len(s['exhibit']['elements'])-1)],8),x-.10,3.92,2.48,.42,8.5,SLATE,i==3)
 label(slide,'COUNTERARGUMENT',.72,4.92,1.65,CORAL); tx(slide,words(s['counterargument'],14),2.38,4.87,9.98,.36,9.5,SLATE); decision_rail(slide,s)

def table(slide,s,c):
 headers=['ANALYTICAL LENS','EVIDENCE','BOUNDARY / STATUS']; xs=[.72,3.17,9.67]; ws=[2.32,6.38,2.95]
 for x,w,h in zip(xs,ws,headers): rect(slide,x,1.86,w,.38,NAVY,NAVY); tx(slide,h,x+.12,1.95,w-.24,.17,8.5,WHITE,True)
 for i,e in enumerate(s['evidence_blocks'][:4]):
  y=2.31+i*.65; fill=WHITE if i%2==0 else PALE; col=STATUS[e['evidence_status']]
  for x,w in zip(xs,ws): rect(slide,x,y,w,.58,fill,LINE)
  tx(slide,words(e['label'],5),xs[0]+.12,y+.11,ws[0]-.24,.31,8.8,INK,True); tx(slide,words(e['claim'],16),xs[1]+.12,y+.08,ws[1]-.24,.38,9.1,INK); status_dot(slide,e['evidence_status'],xs[2]+.14,y+.18); tx(slide,e['evidence_status'],xs[2]+.35,y+.12,ws[2]-.48,.27,8.2,col,True)
 label(slide,'ANALYTICAL CONCLUSION',.72,5.06,1.90); tx(slide,words(s['comparison_logic'],15),2.62,5.01,9.72,.32,9.3,SLATE); decision_rail(slide,s)

def waterfall(slide,s,c):
 label(slide,'DIRECTIONAL BRIDGE',.72,1.78,1.8); line(slide,.95,4.87,12.30,4.87,SLATE,1)
 for i,e in enumerate(s['evidence_blocks'][:4]):
  x=1.05+i*2.98; h=.75+i*.35; col=STATUS[e['evidence_status']]; rect(slide,x,4.87-h,2.26,h,col,col); tx(slide,words(e['label'],3),x,4.42-h,2.26,.30,8.2,col,True,PP_ALIGN.CENTER); tx(slide,words(e['claim'],7),x+.14,4.97-h,1.98,h-.17,8.8,WHITE,True,PP_ALIGN.CENTER)
 tx(slide,'DIRECTIONAL ONLY · BAR HEIGHT DOES NOT REPRESENT VERIFIED MAGNITUDE',3.52,1.82,6.30,.20,8.5,GOLD,True,PP_ALIGN.CENTER); decision_rail(slide,s)

RENDER={'cover':cover,'ladder':ladder,'mechanism':mechanism,'heatmap':heatmap,'matrix':matrix,'funnel':funnel,'timeline':timeline,'table':table,'waterfall':waterfall}
