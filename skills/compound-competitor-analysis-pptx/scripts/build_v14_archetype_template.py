#!/usr/bin/env python3
"""Build sanitized native archetypes using the reviewed DeepGrid v14 renderer and theme."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation

import v14_design_renderer as v14


FAMILIES = ["cover", "ladder", "table", "funnel", "matrix", "heatmap", "mechanism", "timeline", "waterfall"]
TITLES = {
    "cover": "Asterion changes near-term threat—but NovaGrid still controls scalable economics",
    "ladder": "Field access supports a high retrofit threat; OEM qualification remains unproven",
    "table": "Leadership should fund three evidence gates—and stop scaling claims ahead of proof",
    "funnel": "Commercial claims must narrow as product accountability rises",
    "matrix": "Asterion wins today's fleet purchase; NovaGrid can win scaled OEM economics",
    "heatmap": "Threat varies by arena because access, qualification and economics control different gates",
    "mechanism": "Owned compute matters only when it converts into qualified cost and supply control",
    "timeline": "A 90-day sequence converts competitor claims into reversible commercial tests",
    "waterfall": "Commitment should rise only as field, qualification and economics evidence accumulates",
}


def evidence(label: str, claim: str, status: str) -> dict:
    return {"label": label, "claim": claim, "evidence_status": status}


def content(family: str) -> dict:
    return {
        "assertion_title": TITLES[family],
        "analytical_question": "What decision does this evidence support, and what would change it?",
        "executive_answer": "Use the bounded conclusion now; release further commitment only when the next evidence gate passes.",
        "comparison_logic": "Buyer control points—not generic feature breadth—determine where each competitor can win.",
        "counterargument": "A verified programme nomination or homologation filing would materially raise the current threat rating.",
        "implication": "Separate the retrofit and OEM plays.",
        "decision": "Partner-test in retrofit; defend through qualification in OEM.",
        "trigger": "Resolve the named evidence gap within 30 days",
        "owner": "Commercial + Engineering",
        "source_note": "Synthetic example; no claim is reusable market evidence.",
        "evidence_blocks": [
            evidence("Field proof", "Independent evidence supports operating deployments in the retrofit segment.", "verified"),
            evidence("Readiness", "Production readiness remains an attributed company claim.", "attributed company claim"),
            evidence("Economics", "Owned compute may improve cost and supply control at qualified volume.", "qualified interpretation"),
            evidence("OEM gate", "No verified nomination or full homologation filing is available.", "insufficient evidence"),
        ],
        "exhibit": {"elements": ["Position by buyer gate", "Test the partner path", "Prove before scaling"]},
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path, help="Reviewed DeepGrid v14 native PPTX")
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    prs = Presentation(args.source)
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)
    blank = prs.slide_layouts[6]
    for number, family in enumerate(FAMILIES, 1):
        slide = prs.slides.add_slide(blank)
        slide_content = content(family)
        contract = {"slide_number": number, "family": family}
        if family != "cover":
            v14.header(slide, slide_content, number)
        v14.RENDER[family](slide, slide_content, contract)
        slide.notes_slide.notes_text_frame.text = (
            f"Archetype: {family}. Theme, geometry and native composition derive from the reviewed "
            "DeepGrid v14 design-contract deck. Visible content is synthetic and must be replaced "
            "through the governed content envelope."
        )
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"wrote {args.output} with {len(prs.slides)} v14-derived native archetypes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
