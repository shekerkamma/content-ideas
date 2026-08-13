#!/usr/bin/env python3
from __future__ import annotations

import hashlib
import json
from pathlib import Path
import sys

from pptx import Presentation


def fail(message: str) -> None:
    raise SystemExit(message)


def load(path: Path) -> dict:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        fail(f"invalid JSON {path}: {exc}")


def main() -> int:
    root = Path(sys.argv[1] if len(sys.argv) > 1 else ".").resolve()
    required = [
        root / "examples/prompts/sample-prompts.md",
        root / "examples/completed-contracts/slide-content-contract.json",
        root / "examples/completed-contracts/slide-design-contract.json",
        root / "examples/review-control-traces/grill-to-meta-loop.json",
        root / "examples/golden-path/README.md",
        root / "assets/slide-archetypes/archetype-catalog.json",
        root / "assets/slide-archetypes/deepgrid-v14-native-archetypes-draft.pptx",
        root / "assets/slide-archetypes/qa/officecli-final/issues.json",
        root / "assets/reference-decks/accenture-style-claude-guide-draft.pptx",
        root / "references/accenture-guide-content-envelope.md",
        root / "examples/reference-source/presentation-evidence.json",
        root / "examples/reference-source/slide_ascii.md",
    ]
    missing = [str(path) for path in required if not path.is_file()]
    if missing:
        fail("missing example artifacts: " + ", ".join(missing))

    content = load(required[1])
    design = load(required[2])
    trace = load(required[3])
    catalog = load(required[5])
    archetypes = {item["id"] for item in catalog.get("archetypes", [])}
    if len(archetypes) != 9:
        fail("v14 archetype catalog must contain exactly nine unique design families")
    if content.get("archetype") not in archetypes or design.get("archetype") not in archetypes:
        fail("completed contracts must reference a catalog archetype")
    if content.get("slide_id") != design.get("slide_id") or content.get("slide_id") != trace.get("slide_id"):
        fail("content, design, and review trace slide IDs must match")
    controls = set(content.get("review_control_ids", []))
    expected = {trace["grill_finding"]["control_id"], trace["meta_loop_disposition"]["control_id"]}
    if not expected.issubset(controls):
        fail("completed content contract must cite Grill-Me and Meta LOOP controls")
    prs = Presentation(required[6])
    if len(prs.slides) != len(archetypes):
        fail("native v14 design reference slide count must match archetype count")
    if any(len(slide.shapes) < 10 for slide in prs.slides):
        fail("native v14 archetypes must remain editable multi-object compositions")
    issues = load(required[7])
    if ((issues.get("data") or {}).get("count")) != 0:
        fail("native v14 archetype deck has unresolved OfficeCLI issues")
    expected_hash = "961d4685f42c29709bb9582e10de2486c78c92cd82fe4596a3942d8317ccfd02"
    actual_hash = hashlib.sha256(required[8].read_bytes()).hexdigest()
    if actual_hash != expected_hash:
        fail("supplied Accenture-style guide does not match its registered checksum")
    envelope = required[9].read_text(encoding="utf-8")
    for heading in ("WHEN TO USE", "WORKFLOW", "KEY PROMPT", "OUTPUT INCLUDES"):
        if heading not in envelope:
            fail(f"content-envelope reference missing {heading}")
    evidence = load(required[10])
    if len(evidence.get("slides", [])) != 24:
        fail("normalized Accenture-style reference must contain 24 source slides")
    source_hashes = {item.get("sha256") for item in evidence.get("sources", [])}
    if expected_hash not in source_hashes:
        fail("presentation evidence is not hash-linked to the supplied reference deck")
    print(json.dumps({"status": "healthy", "archetypes": len(archetypes), "example_slide": content["slide_id"]}))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
