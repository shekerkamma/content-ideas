#!/usr/bin/env python3
"""preview_pptx.py's overflow rule, with the effective font size resolved correctly.

preview_pptx.py reads `para.runs[0].font.size` and falls back to 14pt when it is
None. artifact-tool writes the size to the paragraph's a:defRPr and emits runs with
no a:rPr, so that fallback fires for EVERY paragraph -- over-measuring small text
and under-measuring display text. This resolves defRPr first, then the run, then 14.
Everything else (wrap width, 1.18 line stacking, 0.03in tolerance) is unchanged.
"""
import sys, textwrap
from pptx import Presentation
from pptx.util import Emu
A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def eff_size(p):
    if p.runs and p.runs[0].font.size is not None:
        return p.runs[0].font.size.pt
    if p._pPr is not None:
        d = p._pPr.find(A + 'defRPr')
        if d is not None and d.get('sz'):
            return int(d.get('sz')) / 100
    return 14.0

def main(path):
    prs = Presentation(path); bad = 0
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            total = 0.0
            for para in sh.text_frame.paragraphs:
                if not para.runs:
                    continue
                text = "".join(r.text for r in para.runs)
                size = eff_size(para)
                max_chars = max(4, int(w / (size * 0.0085)))
                lines = textwrap.wrap(text, max_chars) or [text]
                sb = para.space_before.pt / 72.0 if para.space_before else 0
                total += sb + len(lines) * (size / 72.0) * 1.18
            if total > h + 0.03:
                bad += 1
                print(f"slide {idx:2d}  needs {total:.2f}in in {h:.2f}in  "
                      f":: {sh.text_frame.text.strip()[:56]!r}")
    print(f"\n{bad} overflowing shape(s)")
    return 0 if bad == 0 else 2

if __name__ == '__main__':
    sys.exit(main(sys.argv[1]))
