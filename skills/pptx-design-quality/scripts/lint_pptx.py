#!/usr/bin/env python3
"""Deterministic design-quality linting for native PowerPoint artifacts."""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter
from dataclasses import asdict, dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError


EMU_PER_INCH = 914400
DEFAULTS = {
    "deck": {"output_mode": "native"},
    "typography": {
        "min_body_font_pt": 14,
        "min_caption_font_pt": 9,
        "max_font_families": 3,
    },
    "color": {"max_distinct_colors": 10, "min_contrast_ratio": 4.5},
    "layout": {
        "aspect_ratio": "16:9",
        "max_words_per_slide": 80,
        "max_shapes_per_slide": 36,
        "max_consecutive_same_archetype": 2,
        "min_image_dpi": 120,
    },
    "qa": {"ignore_rules": [], "waivers": {}},
}


@dataclass(frozen=True)
class Finding:
    rule_id: str
    severity: str
    message: str
    slide: int | None = None
    shape: str | None = None


@dataclass(frozen=True)
class ShapeView:
    shape: Any
    left: int
    top: int
    width: int
    height: int
    rotation: float = 0
    flip_h: bool = False
    flip_v: bool = False

    def __getattr__(self, name: str) -> Any:
        return getattr(self.shape, name)


@dataclass(frozen=True)
class TableCellProxy:
    cell: Any
    name: str
    shape_type: MSO_SHAPE_TYPE = MSO_SHAPE_TYPE.TABLE
    has_text_frame: bool = True
    is_placeholder: bool = False
    is_table_cell: bool = True

    @property
    def text_frame(self) -> Any:
        return self.cell.text_frame

    @property
    def text(self) -> str:
        return self.cell.text

    @property
    def fill(self) -> Any:
        return self.cell.fill

    @property
    def _element(self) -> Any:
        return self.cell._tc


def _merge_defaults(config: dict[str, Any] | None) -> dict[str, Any]:
    merged = json.loads(json.dumps(DEFAULTS))
    if not config:
        return merged
    for section, values in config.items():
        if isinstance(values, dict) and isinstance(merged.get(section), dict):
            merged[section].update(values)
        else:
            merged[section] = values
    return merged


def load_config(path: Path | None) -> dict[str, Any]:
    if path is None:
        return _merge_defaults(None)
    return _merge_defaults(json.loads(path.read_text(encoding="utf-8")))


def _inches(value: int) -> float:
    return value / EMU_PER_INCH


def _shape_name(shape: Any, index: int) -> str:
    return getattr(shape, "name", None) or f"shape-{index}"


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join(shape.text.split())


def _runs(shape: Any) -> Iterable[Any]:
    if not getattr(shape, "has_text_frame", False):
        return ()
    return (
        run
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    )


def _xml_default_run_properties(shape: Any, paragraph: Any, slide: Any) -> Iterable[Any]:
    """Yield default run properties from nearest to broadest inheritance scope."""
    level = min(max(getattr(paragraph, "level", 0), 0), 8) + 1
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    level_tag = f"{{{drawing_ns}}}lvl{level}pPr"
    def_run_tag = f"{{{drawing_ns}}}defRPr"
    list_style_tag = f"{{{drawing_ns}}}lstStyle"

    paragraph_properties = paragraph._p.pPr
    if paragraph_properties is not None:
        default_run = paragraph_properties.find(def_run_tag)
        if default_run is not None:
            yield default_run

    def defaults_from_text_body(candidate: Any) -> Iterable[Any]:
        text_body = getattr(candidate._element, "txBody", None)
        if text_body is None:
            return ()
        list_style = text_body.find(list_style_tag)
        level_properties = list_style.find(level_tag) if list_style is not None else None
        default_run = (
            level_properties.find(def_run_tag) if level_properties is not None else None
        )
        return (default_run,) if default_run is not None else ()

    yield from defaults_from_text_body(shape)

    if not getattr(shape, "is_placeholder", False):
        return

    placeholder_index = shape.placeholder_format.idx
    for layout_shape in slide.slide_layout.placeholders:
        if layout_shape.placeholder_format.idx == placeholder_index:
            yield from defaults_from_text_body(layout_shape)
            break

    placeholder_type = shape.placeholder_format.type
    style_name = (
        "titleStyle"
        if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
        else "bodyStyle"
        if placeholder_type
        in {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
            PP_PLACEHOLDER.VERTICAL_BODY,
            PP_PLACEHOLDER.VERTICAL_OBJECT,
        }
        else "otherStyle"
    )
    master = slide.slide_layout.slide_master._element
    text_styles = master.find(f"{{{presentation_ns}}}txStyles")
    style = (
        text_styles.find(f"{{{presentation_ns}}}{style_name}")
        if text_styles is not None
        else None
    )
    level_properties = style.find(level_tag) if style is not None else None
    default_run = (
        level_properties.find(def_run_tag) if level_properties is not None else None
    )
    if default_run is not None:
        yield default_run


def _inherited_run_attribute(
    shape: Any, paragraph: Any, slide: Any, attribute: str
) -> str | None:
    for properties in _xml_default_run_properties(shape, paragraph, slide):
        value = properties.get(attribute)
        if value:
            return value
    return None


@lru_cache(maxsize=8)
def _xml_root(blob: bytes) -> ET.Element:
    return ET.fromstring(blob)


def _theme_root(slide: Any) -> ET.Element | None:
    master_part = slide.slide_layout.slide_master.part
    for relationship in master_part.rels.values():
        if relationship.reltype.endswith("/theme"):
            return _xml_root(relationship.target_part.blob)
    return None


def _theme_font(slide: Any, major: bool) -> str | None:
    root = _theme_root(slide)
    if root is None:
        return None
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    family = "majorFont" if major else "minorFont"
    latin = root.find(
        f".//{{{drawing_ns}}}fontScheme/{{{drawing_ns}}}{family}/{{{drawing_ns}}}latin"
    )
    return latin.get("typeface") if latin is not None else None


def _font_name(run: Any, paragraph: Any, shape: Any, slide: Any) -> str | None:
    if run.font.name:
        return run.font.name
    if paragraph.font.name:
        return paragraph.font.name
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for properties in _xml_default_run_properties(shape, paragraph, slide):
        latin = properties.find(f"{{{drawing_ns}}}latin")
        if latin is not None and latin.get("typeface"):
            typeface = latin.get("typeface")
            if typeface == "+mj-lt":
                return _theme_font(slide, major=True)
            if typeface == "+mn-lt":
                return _theme_font(slide, major=False)
            return typeface
    is_title_placeholder = getattr(shape, "is_placeholder", False) and (
        shape.placeholder_format.type
        in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    )
    return _theme_font(slide, major=is_title_placeholder)


def _theme_color(slide: Any, scheme_name: str) -> str | None:
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    master = slide.slide_layout.slide_master._element
    color_map = master.find(f"{{{presentation_ns}}}clrMap")
    master_mapping = dict(color_map.items()) if color_map is not None else {}
    active_mapping = master_mapping
    for owner in (slide.slide_layout, slide):
        override = owner._element.find(f"{{{presentation_ns}}}clrMapOvr")
        if override is None:
            continue
        explicit = override.find(f"{{{drawing_ns}}}overrideClrMapping")
        if explicit is not None:
            active_mapping = dict(explicit.items())
        elif override.find(f"{{{drawing_ns}}}masterClrMapping") is not None:
            active_mapping = master_mapping
    mapped_name = active_mapping.get(scheme_name, scheme_name)
    root = _theme_root(slide)
    if root is None:
        return None
    scheme = root.find(
        f".//{{{drawing_ns}}}clrScheme/{{{drawing_ns}}}{mapped_name}"
    )
    if scheme is None or len(scheme) == 0:
        return None
    color = scheme[0]
    value = (color.get("val") or "").upper()
    if re.fullmatch(r"[0-9A-F]{6}", value):
        return value
    fallback = (color.get("lastClr") or "").upper()
    return fallback if re.fullmatch(r"[0-9A-F]{6}", fallback) else None


def _font_rgb(run: Any, paragraph: Any, shape: Any, slide: Any) -> str | None:
    explicit = _rgb_from_color(run.font.color, slide) or _rgb_from_color(
        paragraph.font.color, slide
    )
    if explicit:
        return explicit
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for properties in _xml_default_run_properties(shape, paragraph, slide):
        solid_fill = properties.find(f"{{{drawing_ns}}}solidFill")
        if solid_fill is None or len(solid_fill) == 0:
            continue
        color = solid_fill[0]
        if color.tag == f"{{{drawing_ns}}}srgbClr":
            return (color.get("val") or "").upper() or None
        if color.tag == f"{{{drawing_ns}}}schemeClr" and color.get("val"):
            return _theme_color(slide, color.get("val"))
    return _theme_color(slide, "tx1")


def _font_size(
    run: Any,
    paragraph: Any | None = None,
    shape: Any | None = None,
    slide: Any | None = None,
) -> float | None:
    if run.font.size is not None:
        return float(run.font.size.pt)
    if paragraph is not None and paragraph.font.size is not None:
        return float(paragraph.font.size.pt)
    if paragraph is not None and shape is not None and slide is not None:
        inherited = _inherited_run_attribute(shape, paragraph, slide, "sz")
        if inherited is not None:
            return int(inherited) / 100
    return None


def _shape_font_sizes(shape: Any, slide: Any | None = None) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            size = _font_size(run, paragraph, shape, slide)
            if size is not None:
                sizes.append(size)
    return sizes


def _apply_brightness(rgb: str, brightness: float) -> str:
    channels = [int(rgb[index : index + 2], 16) for index in (0, 2, 4)]
    adjusted = [
        round(channel + (255 - channel) * brightness)
        if brightness >= 0
        else round(channel * (1 + brightness))
        for channel in channels
    ]
    return "".join(f"{max(0, min(255, channel)):02X}" for channel in adjusted)


def _rgb_from_color(color: Any, slide: Any | None = None) -> str | None:
    try:
        if color.type == MSO_COLOR_TYPE.RGB and color.rgb is not None:
            return str(color.rgb).upper()
        if color.type == MSO_COLOR_TYPE.SCHEME and slide is not None:
            scheme_names = {
                MSO_THEME_COLOR.DARK_1: "dk1",
                MSO_THEME_COLOR.LIGHT_1: "lt1",
                MSO_THEME_COLOR.DARK_2: "dk2",
                MSO_THEME_COLOR.LIGHT_2: "lt2",
                MSO_THEME_COLOR.ACCENT_1: "accent1",
                MSO_THEME_COLOR.ACCENT_2: "accent2",
                MSO_THEME_COLOR.ACCENT_3: "accent3",
                MSO_THEME_COLOR.ACCENT_4: "accent4",
                MSO_THEME_COLOR.ACCENT_5: "accent5",
                MSO_THEME_COLOR.ACCENT_6: "accent6",
                MSO_THEME_COLOR.HYPERLINK: "hlink",
                MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "folHlink",
                MSO_THEME_COLOR.TEXT_1: "tx1",
                MSO_THEME_COLOR.TEXT_2: "tx2",
                MSO_THEME_COLOR.BACKGROUND_1: "bg1",
                MSO_THEME_COLOR.BACKGROUND_2: "bg2",
            }
            scheme_name = scheme_names.get(color.theme_color)
            resolved = _theme_color(slide, scheme_name) if scheme_name else None
            if resolved:
                return _apply_brightness(resolved, color.brightness)
    except (AttributeError, ValueError):
        pass
    return None


def _solid_fill_rgb(fill: Any, slide: Any | None = None) -> str | None:
    try:
        if fill.type == MSO_FILL_TYPE.SOLID:
            fill_element = getattr(fill, "_fill", None)
            fill_element = getattr(fill_element, "_solidFill", fill_element)
            if fill_element is not None:
                for descendant in fill_element.iter():
                    local_name = descendant.tag.rsplit("}", 1)[-1]
                    if not local_name.startswith("alpha"):
                        continue
                    value = descendant.get("val")
                    if local_name == "alphaOff":
                        if value is None or value != "0":
                            return None
                    elif value is None or value != "100000":
                        return None
            return _rgb_from_color(fill.fore_color, slide)
    except (AttributeError, ValueError):
        pass
    return None


def _xml_color_rgb(color: Any, slide: Any) -> str | None:
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    local_name = color.tag.rsplit("}", 1)[-1]
    value = (color.get("val") or color.get("lastClr") or "").upper()
    if local_name in {"srgbClr", "sysClr"}:
        return value if re.fullmatch(r"[0-9A-F]{6}", value) else None
    if local_name == "schemeClr" and color.get("val"):
        return _theme_color(slide, color.get("val"))
    if local_name == "scrgbClr":
        try:
            channels = [round(int(color.get(name)) * 255 / 100000) for name in ("r", "g", "b")]
        except (TypeError, ValueError):
            return None
        return "".join(f"{max(0, min(255, channel)):02X}" for channel in channels)
    # Preset and HSL colors are intentionally unresolved; rendered review owns them.
    return None


def _background_element_rgb(background: Any, slide: Any) -> str | None:
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    solid_fill = background.find(
        f"./{{{presentation_ns}}}bgPr/{{{drawing_ns}}}solidFill"
    )
    if solid_fill is not None and len(solid_fill):
        return _xml_color_rgb(solid_fill[0], slide)
    background_ref = background.find(f"./{{{presentation_ns}}}bgRef")
    if background_ref is not None and len(background_ref):
        return _xml_color_rgb(background_ref[0], slide)
    return None


def _slide_background_rgb(slide: Any) -> str | None:
    owners = (slide, slide.slide_layout, slide.slide_layout.slide_master)
    for owner in owners:
        background = owner._element.cSld.bg
        if background is not None:
            return _background_element_rgb(background, slide)
    return _theme_color(slide, "bg1")


def _relative_luminance(rgb: str) -> float:
    channels = [int(rgb[index : index + 2], 16) / 255 for index in (0, 2, 4)]
    converted = [
        channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4
        for channel in channels
    ]
    return 0.2126 * converted[0] + 0.7152 * converted[1] + 0.0722 * converted[2]


def _contrast_ratio(first: str, second: str) -> float:
    light, dark = sorted(
        (_relative_luminance(first), _relative_luminance(second)), reverse=True
    )
    return (light + 0.05) / (dark + 0.05)


def _is_title(shape: Any, slide_height: int, slide: Any | None = None) -> bool:
    if getattr(shape, "is_table_cell", False):
        return False
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type in {
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            } and bool(_shape_text(shape))
    except (AttributeError, ValueError):
        pass
    sizes = _shape_font_sizes(shape, slide)
    return (
        bool(sizes)
        and _rendered_bounds(shape)[1] < slide_height * 0.28
        and max(sizes) >= 24
    )


def _is_footer(shape: Any, slide_height: int) -> bool:
    return _rendered_bounds(shape)[1] >= slide_height * 0.88


def _is_caption(
    shape: Any, slide_height: int, earlier_shapes: Iterable[Any]
) -> bool:
    if getattr(shape, "is_table_cell", False):
        return False
    if _is_footer(shape, slide_height):
        return True
    label = f"{getattr(shape, 'name', '')} {_shape_text(shape)}"
    if re.search(r"\b(caption|source|footnote|note)\b", label, re.IGNORECASE):
        return True
    if shape.height > 0.65 * EMU_PER_INCH:
        return False
    shape_left, shape_top, shape_right, _ = _rendered_bounds(shape)
    for candidate in reversed(list(earlier_shapes)):
        if not (
            candidate.shape_type == MSO_SHAPE_TYPE.PICTURE
            or getattr(candidate, "has_chart", False)
            or getattr(candidate, "has_table", False)
        ):
            continue
        candidate_left, _, candidate_right, candidate_bottom = _rendered_bounds(
            candidate
        )
        gap = shape_top - candidate_bottom
        horizontal_overlap = max(
            0, min(shape_right, candidate_right) - max(shape_left, candidate_left)
        )
        if (
            0 <= gap <= 0.35 * EMU_PER_INCH
            and horizontal_overlap >= shape.width * 0.5
        ):
            return True
    return False


def _rendered_bounds(shape: Any) -> tuple[float, float, float, float]:
    rotation = math.radians(float(getattr(shape, "rotation", 0) or 0))
    rendered_width = abs(shape.width * math.cos(rotation)) + abs(
        shape.height * math.sin(rotation)
    )
    rendered_height = abs(shape.width * math.sin(rotation)) + abs(
        shape.height * math.cos(rotation)
    )
    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2
    return (
        center_x - rendered_width / 2,
        center_y - rendered_height / 2,
        center_x + rendered_width / 2,
        center_y + rendered_height / 2,
    )


def _overlap_ratio(first: Any, second: Any) -> float:
    first_left, first_top, first_right, first_bottom = _rendered_bounds(first)
    second_left, second_top, second_right, second_bottom = _rendered_bounds(second)
    left = max(first_left, second_left)
    top = max(first_top, second_top)
    right = min(first_right, second_right)
    bottom = min(first_bottom, second_bottom)
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    smaller = min(
        (first_right - first_left) * (first_bottom - first_top),
        (second_right - second_left) * (second_bottom - second_top),
    )
    return intersection / smaller if smaller else 0.0


def _coverage_ratio(target: Any, background: Any) -> float:
    target_left, target_top, target_right, target_bottom = _rendered_bounds(target)
    bg_left, bg_top, bg_right, bg_bottom = _rendered_bounds(background)
    left = max(target_left, bg_left)
    top = max(target_top, bg_top)
    right = min(target_right, bg_right)
    bottom = min(target_bottom, bg_bottom)
    if right <= left or bottom <= top:
        return 0.0
    target_area = (target_right - target_left) * (target_bottom - target_top)
    return ((right - left) * (bottom - top)) / target_area if target_area else 0.0


def _effective_background_rgb(
    shape: Any, earlier_shapes: Iterable[Any], slide: Any, slide_background: str | None
) -> str | None:
    direct = _solid_fill_rgb(getattr(shape, "fill", None), slide)
    if direct:
        return direct
    for underlying in reversed(list(earlier_shapes)):
        if _coverage_ratio(shape, underlying) < 0.5:
            continue
        # A photograph or screenshot is the visible background, but its local
        # color cannot be inferred reliably from OOXML geometry alone. Skip the
        # automated contrast assertion and leave it to rendered-slide review.
        if underlying.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return None
        underlying_fill = _solid_fill_rgb(
            getattr(underlying, "fill", None), slide
        )
        if underlying_fill:
            return underlying_fill
    return slide_background


def _overflow_risk(shape: Any, slide: Any | None = None) -> bool:
    sizes = _shape_font_sizes(shape, slide)
    if not sizes or shape.width <= 0 or shape.height <= 0:
        return False
    font_pt = max(8.0, sum(sizes) / len(sizes))
    char_width = (font_pt / 72) * 0.52
    line_height = (font_pt / 72) * 1.2
    chars_per_line = max(1, int(_inches(shape.width) / char_width))
    text_lines = [
        line
        for paragraph in shape.text_frame.paragraphs
        for line in "".join(run.text for run in paragraph.runs)
        .replace("\v", "\n")
        .split("\n")
    ]
    if not any(line.strip() for line in text_lines):
        return False
    required_lines = sum(
        max(1, math.ceil(len(line) / chars_per_line)) for line in text_lines
    )
    available_lines = max(1, int(_inches(shape.height) / line_height))
    return required_lines > available_lines * 1.2


def _group_children(group: ShapeView) -> Iterable[ShapeView]:
    transform = group.shape._element.xfrm
    child_width = transform.chExt.cx
    child_height = transform.chExt.cy
    scale_x = group.width / child_width if child_width else 1
    scale_y = group.height / child_height if child_height else 1
    group_rotation = float(group.rotation or 0)
    group_center_x = group.left + group.width / 2
    group_center_y = group.top + group.height / 2
    flip_h = group.flip_h ^ bool(getattr(transform, "flipH", False))
    flip_v = group.flip_v ^ bool(getattr(transform, "flipV", False))
    angle = math.radians(group_rotation)
    for child in group.shape.shapes:
        width = round(child.width * scale_x)
        height = round(child.height * scale_y)
        center_x = group.left + (
            child.left + child.width / 2 - transform.chOff.x
        ) * scale_x
        center_y = group.top + (
            child.top + child.height / 2 - transform.chOff.y
        ) * scale_y
        offset_x = center_x - group_center_x
        offset_y = center_y - group_center_y
        if flip_h:
            offset_x = -offset_x
        if flip_v:
            offset_y = -offset_y
        rotated_center_x = (
            group_center_x + offset_x * math.cos(angle) - offset_y * math.sin(angle)
        )
        rotated_center_y = (
            group_center_y + offset_x * math.sin(angle) + offset_y * math.cos(angle)
        )
        view = ShapeView(
            child,
            round(rotated_center_x - width / 2),
            round(rotated_center_y - height / 2),
            width,
            height,
            (group_rotation + float(getattr(child, "rotation", 0) or 0)) % 360,
            flip_h,
            flip_v,
        )
        if child.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _group_children(view)
        else:
            yield view


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[ShapeView]:
    for shape in shapes:
        view = ShapeView(
            shape,
            shape.left,
            shape.top,
            shape.width,
            shape.height,
            float(getattr(shape, "rotation", 0) or 0),
        )
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _group_children(view)
        else:
            yield view


def _table_cells(table_view: ShapeView) -> Iterable[ShapeView]:
    table = table_view.shape.table
    total_width = sum(column.width for column in table.columns)
    total_height = sum(row.height for row in table.rows)
    scale_x = table_view.width / total_width if total_width else 1
    scale_y = table_view.height / total_height if total_height else 1
    top = table_view.top
    for row_index, row in enumerate(table.rows, 1):
        left = table_view.left
        for column_index, column in enumerate(table.columns, 1):
            cell = row.cells[column_index - 1]
            if not getattr(cell, "is_spanned", False):
                span_width = max(1, int(getattr(cell, "span_width", 1)))
                span_height = max(1, int(getattr(cell, "span_height", 1)))
                merged_width = sum(
                    table.columns[offset].width
                    for offset in range(
                        column_index - 1,
                        min(column_index - 1 + span_width, len(table.columns)),
                    )
                )
                merged_height = sum(
                    table.rows[offset].height
                    for offset in range(
                        row_index - 1,
                        min(row_index - 1 + span_height, len(table.rows)),
                    )
                )
                yield ShapeView(
                    TableCellProxy(cell, f"{table_view.name} cell {row_index},{column_index}"),
                    round(left),
                    round(top),
                    round(merged_width * scale_x),
                    round(merged_height * scale_y),
                    table_view.rotation,
                )
            left += column.width * scale_x
        top += row.height * scale_y


def _layout_signature(slide: Any, slide_width: int, slide_height: int) -> tuple[Any, ...]:
    signature: list[tuple[Any, ...]] = []
    for shape in _iter_shapes(slide.shapes):
        if _is_footer(shape, slide_height):
            continue
        left, top, right, bottom = _rendered_bounds(shape)
        signature.append(
            (
                int(shape.shape_type),
                round(left / slide_width, 1),
                round(top / slide_height, 1),
                round((right - left) / slide_width, 1),
                round((bottom - top) / slide_height, 1),
                bool(_shape_text(shape)),
            )
        )
    return tuple(sorted(signature))


def _expected_ratio(label: str) -> float | None:
    return {"16:9": 16 / 9, "4:3": 4 / 3}.get(label)


def lint_presentation(
    pptx_path: Path, config: dict[str, Any] | None = None, fast: bool = False
) -> dict[str, Any]:
    cfg = _merge_defaults(config)
    presentation = Presentation(str(pptx_path))
    output_mode = cfg["deck"]["output_mode"]
    findings: list[Finding] = []
    fonts: set[str] = set()
    colors: set[str] = set()
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    expected_ratio = _expected_ratio(cfg["layout"]["aspect_ratio"])
    actual_ratio = slide_width / slide_height
    if expected_ratio and abs(actual_ratio - expected_ratio) > 0.02:
        findings.append(
            Finding(
                "DECK_ASPECT_RATIO",
                "error",
                f"Expected {cfg['layout']['aspect_ratio']}; found {actual_ratio:.3f}:1.",
            )
        )
    if not presentation.slides:
        findings.append(
            Finding(
                "DECK_NO_SLIDES",
                "error",
                "Presentation contains no slides.",
            )
        )

    signatures: list[tuple[Any, ...]] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        words = 0
        text_shapes: list[tuple[int, Any]] = []
        has_title = False
        background = _slide_background_rgb(slide)

        shapes = list(_iter_shapes(slide.shapes))
        if not shapes:
            findings.append(
                Finding(
                    "SLIDE_EMPTY",
                    "error",
                    "Slide contains no visible or editable objects.",
                    slide_number,
                )
            )
        inspection_shapes = [
            inspection
            for shape in shapes
            for inspection in (
                list(_table_cells(shape)) if getattr(shape, "has_table", False) else [shape]
            )
        ]
        template_background_shapes = [
            *list(_iter_shapes(slide.slide_layout.slide_master.shapes)),
            *list(_iter_shapes(slide.slide_layout.shapes)),
        ]
        if len(shapes) > cfg["layout"]["max_shapes_per_slide"]:
            findings.append(
                Finding(
                    "SLIDE_SHAPE_COUNT",
                    "warning",
                    f"{len(shapes)} shapes exceed the configured maximum "
                    f"{cfg['layout']['max_shapes_per_slide']}.",
                    slide_number,
                )
            )

        for shape_index, shape in enumerate(inspection_shapes, 1):
            name = _shape_name(shape, shape_index)
            text = _shape_text(shape)
            words += len(re.findall(r"\b[\w’'-]+\b", text))
            has_title = has_title or _is_title(shape, slide_height, slide)
            if text:
                text_shapes.append((shape_index, shape))

            tolerance = 0.02 * EMU_PER_INCH
            rendered_left, rendered_top, rendered_right, rendered_bottom = (
                _rendered_bounds(shape)
            )
            if (
                rendered_left < -tolerance
                or rendered_top < -tolerance
                or rendered_right > slide_width + tolerance
                or rendered_bottom > slide_height + tolerance
            ):
                findings.append(
                    Finding(
                        "SHAPE_OUT_OF_BOUNDS",
                        "error",
                        "Shape extends beyond the slide canvas.",
                        slide_number,
                        name,
                    )
                )

            if text and _overflow_risk(shape, slide):
                findings.append(
                    Finding(
                        "TEXT_OVERFLOW_RISK",
                        "warning",
                        "Text density likely exceeds the available box height.",
                        slide_number,
                        name,
                    )
                )

            if text:
                title = _is_title(shape, slide_height, slide)
                caption = _is_caption(
                    shape, slide_height, inspection_shapes[: shape_index - 1]
                )
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip():
                            continue
                        font_name = _font_name(run, paragraph, shape, slide)
                        if font_name:
                            fonts.add(font_name)
                        size = _font_size(run, paragraph, shape, slide)
                        if size is not None:
                            minimum = (
                                cfg["typography"]["min_caption_font_pt"]
                                if caption
                                else cfg["typography"]["min_body_font_pt"]
                            )
                            if not title and size < minimum:
                                findings.append(
                                    Finding(
                                        "TEXT_TOO_SMALL",
                                        "warning",
                                        f"{size:g}pt text is below the {minimum:g}pt minimum.",
                                        slide_number,
                                        name,
                                    )
                                )
                        font_rgb = _font_rgb(run, paragraph, shape, slide)
                        if font_rgb:
                            colors.add(font_rgb)
                            if not fast:
                                fill_rgb = _effective_background_rgb(
                                    shape,
                                    [
                                        *template_background_shapes,
                                        *inspection_shapes[: shape_index - 1],
                                    ],
                                    slide,
                                    background,
                                )
                                if fill_rgb:
                                    ratio = _contrast_ratio(font_rgb, fill_rgb)
                                    if ratio < cfg["color"]["min_contrast_ratio"]:
                                        findings.append(
                                            Finding(
                                                "TEXT_LOW_CONTRAST",
                                                "warning",
                                                f"Contrast {ratio:.2f}:1 is below "
                                                f"{cfg['color']['min_contrast_ratio']}:1.",
                                                slide_number,
                                                name,
                                            )
                                        )

            fill_rgb = _solid_fill_rgb(getattr(shape, "fill", None), slide)
            if fill_rgb:
                colors.add(fill_rgb)

            if not fast and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    width_px, height_px = shape.image.size
                    visible_width_px = width_px * max(
                        0.0, min(1.0, 1 - shape.crop_left - shape.crop_right)
                    )
                    visible_height_px = height_px * max(
                        0.0, min(1.0, 1 - shape.crop_top - shape.crop_bottom)
                    )
                    effective_dpi = min(
                        visible_width_px / _inches(shape.width),
                        visible_height_px / _inches(shape.height),
                    )
                    if effective_dpi < cfg["layout"]["min_image_dpi"]:
                        findings.append(
                            Finding(
                                "IMAGE_LOW_DPI",
                                "warning",
                                f"Effective resolution {effective_dpi:.0f} DPI is below "
                                f"{cfg['layout']['min_image_dpi']:g} DPI.",
                                slide_number,
                                name,
                            )
                        )
                except (AttributeError, ZeroDivisionError):
                    pass

        if words > cfg["layout"]["max_words_per_slide"]:
            findings.append(
                Finding(
                    "SLIDE_WORD_COUNT",
                    "warning",
                    f"{words} words exceed the configured maximum "
                    f"{cfg['layout']['max_words_per_slide']}.",
                    slide_number,
                )
            )
        if not has_title and output_mode != "image-per-slide":
            findings.append(
                Finding(
                    "SLIDE_MISSING_TITLE",
                    "warning",
                    "No title-like text was found in the upper slide region.",
                    slide_number,
                )
            )

        for index, (first_index, first) in enumerate(text_shapes):
            for second_index, second in text_shapes[index + 1 :]:
                if _overlap_ratio(first, second) >= 0.2:
                    findings.append(
                        Finding(
                            "TEXT_BOX_OVERLAP",
                            "error",
                            f"Text boxes overlap materially with "
                            f"{_shape_name(second, second_index)}.",
                            slide_number,
                            _shape_name(first, first_index),
                        )
                    )

        signatures.append(_layout_signature(slide, slide_width, slide_height))

    if len(fonts) > cfg["typography"]["max_font_families"]:
        findings.append(
            Finding(
                "DECK_FONT_COUNT",
                "warning",
                f"{len(fonts)} font families exceed the configured maximum "
                f"{cfg['typography']['max_font_families']}: {', '.join(sorted(fonts))}.",
            )
        )
    if len(colors) > cfg["color"]["max_distinct_colors"]:
        findings.append(
            Finding(
                "DECK_COLOR_COUNT",
                "warning",
                f"{len(colors)} explicit colors exceed the configured maximum "
                f"{cfg['color']['max_distinct_colors']}.",
            )
        )

    if output_mode != "image-per-slide":
        allowed_repeat = cfg["layout"]["max_consecutive_same_archetype"]
        run_start = 0
        for index in range(1, len(signatures) + 1):
            if index < len(signatures) and signatures[index] == signatures[run_start]:
                continue
            run_length = index - run_start
            if run_length > allowed_repeat:
                findings.append(
                    Finding(
                        "LAYOUT_REPETITION",
                        "warning",
                        f"Slides {run_start + 1}-{index} repeat the same layout "
                        f"{run_length} times; configured maximum is {allowed_repeat}.",
                        run_start + 1,
                    )
                )
            run_start = index

    ignored_rules = set(cfg["qa"].get("ignore_rules", []))
    waiver_reasons = cfg["qa"].get("waivers", {})

    def is_valid_waiver(finding: Finding) -> bool:
        reason = waiver_reasons.get(finding.rule_id)
        return (
            finding.severity == "warning"
            and finding.rule_id in ignored_rules
            and isinstance(reason, str)
            and bool(reason.strip())
        )

    active = [finding for finding in findings if not is_valid_waiver(finding)]
    ignored = [finding for finding in findings if is_valid_waiver(finding)]
    counts = Counter(finding.severity for finding in active)
    return {
        "contract_version": "1.0",
        "file": str(pptx_path),
        "fast": fast,
        "status": "clean" if not active else "findings",
        "counts": {
            "error": counts.get("error", 0),
            "warning": counts.get("warning", 0),
            "ignored": len(ignored),
        },
        "findings": [asdict(finding) for finding in active],
        "ignored_findings": [asdict(finding) for finding in ignored],
    }


def _print_text(report: dict[str, Any]) -> None:
    print(
        f"{report['status']}: {report['file']} "
        f"({report['counts']['error']} errors, "
        f"{report['counts']['warning']} warnings)"
    )
    for finding in report["findings"]:
        location = ""
        if finding["slide"] is not None:
            location += f" slide {finding['slide']}"
        if finding["shape"]:
            location += f" / {finding['shape']}"
        print(
            f"- {finding['severity'].upper()} {finding['rule_id']}{location}: "
            f"{finding['message']}"
        )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", help="PowerPoint file to inspect")
    parser.add_argument("--config", type=Path, help="deck-design.json path")
    parser.add_argument("--fast", action="store_true", help="Skip contrast and image DPI")
    parser.add_argument("--json", action="store_true", help="Print JSON instead of text")
    parser.add_argument("--out", type=Path, help="Write the JSON report to this path")
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    try:
        if not pptx_path.is_file():
            raise FileNotFoundError(pptx_path)
        report = lint_presentation(pptx_path, load_config(args.config), args.fast)
    except (
        OSError,
        KeyError,
        TypeError,
        ValueError,
        json.JSONDecodeError,
        PackageNotFoundError,
        zipfile.BadZipFile,
    ) as exc:
        print(f"lint failed: {exc}", file=sys.stderr)
        return 1

    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(
            json.dumps(report, indent=2, ensure_ascii=False) + "\n",
            encoding="utf-8",
        )
    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        _print_text(report)
    return 0 if report["status"] == "clean" else 2


if __name__ == "__main__":
    raise SystemExit(main())
