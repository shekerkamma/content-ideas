#!/usr/bin/env python3
"""Derive a draft template profile from a reference deck.

Never writes the canonical template-profile.json; always writes a distinctly
named draft that must be reviewed, copied over, and validated with
validate_template_profile.py. See references/template-derivation.md for the
heuristics used.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


EMU_PER_INCH = 914400

SKILL_DIR = Path(__file__).resolve().parents[1]
ASSETS_DIR = SKILL_DIR / "assets"
REFERENCES_DIR = SKILL_DIR / "references"

STRUCTURAL_DEFAULT_ARCHETYPES = {"cover", "executive-summary", "section-divider"}
DEFAULTED_GEOMETRY_FIELDS = ("grid_columns", "gutter_inches")
DEFAULTED_COMPOSITION_FIELDS = (
    "density",
    "corner_radius",
    "shadow_policy",
    "accent_policy",
    "whitespace_policy",
)


def _inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 3)


def _load_default_template() -> dict[str, Any]:
    return json.loads(
        (ASSETS_DIR / "template-profile.template.json").read_text(encoding="utf-8")
    )


def _load_archetype_catalog() -> list[dict[str, Any]]:
    catalog = json.loads(
        (REFERENCES_DIR / "slide-archetypes.json").read_text(encoding="utf-8")
    )
    return catalog["archetypes"]


def _luminance(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _saturation(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    hi, lo = max(r, g, b), min(r, g, b)
    return (hi - lo) / 255.0


def derive_colors(image_paths: list[Path], notes: list[str]) -> dict[str, str] | None:
    """Sample dominant page/ink/accent colors from rendered slide images."""
    try:
        from PIL import Image
    except ImportError:
        notes.append("defaulted: brand.colors (Pillow not installed)")
        return None

    counts: dict[tuple[int, int, int], int] = {}
    sampled = 0
    for image_path in image_paths:
        try:
            with Image.open(image_path) as img:
                small = img.convert("RGB").resize((48, 27))
                quantized = small.quantize(colors=8).convert("RGB")
                for count, rgb in quantized.getcolors(48 * 27) or []:
                    counts[rgb] = counts.get(rgb, 0) + count
            sampled += 1
        except (OSError, ValueError):
            continue

    if not counts:
        notes.append("defaulted: brand.colors (no readable slide images)")
        return None

    palette = [rgb for rgb, _ in sorted(counts.items(), key=lambda item: item[1], reverse=True)]
    page = max(palette, key=_luminance)
    ink = min(palette, key=_luminance)
    neutral_candidates = [c for c in palette if c not in (page, ink)]
    accent = max(neutral_candidates, key=_saturation) if neutral_candidates else ink

    def hexcode(rgb: tuple[int, int, int]) -> str:
        return "#{:02X}{:02X}{:02X}".format(*rgb)

    notes.append(f"derived: brand.colors (sampled {sampled} slide image(s))")
    return {"page": hexcode(page), "ink": hexcode(ink), "accent": hexcode(accent)}


CANVAS_FILE = "canvas.json"
ARTBOARD_SUFFIX = ".dc.html"
PT_PER_INCH = 72.0
ROOT_FONT_PX = 16.0

# PowerPoint's own stage widths. A canvas artboard represents one slide, so the
# artboard's pixel width and the slide's real width fix the px scale exactly --
# no 96dpi assumption, no hardcoded 0.75pt/px.
SLIDE_WIDTH_INCHES = {"16:9": 13.333, "4:3": 10.0}

_RULE_RE = re.compile(r"([^{}]*)\{([^{}]*)\}")
_INLINE_STYLE_RE = re.compile(r"""style\s*=\s*(?:"([^"]*)"|'([^']*)')""", re.IGNORECASE)
_HEX_RE = re.compile(r"#([0-9A-Fa-f]{6})\b|#([0-9A-Fa-f]{3})\b")
_RGB_RE = re.compile(r"rgba?\(\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})", re.IGNORECASE)
_REPEAT_RE = re.compile(r"repeat\(\s*(\d{1,2})\s*,", re.IGNORECASE)
_LENGTH_RE = re.compile(r"(-?[\d.]+)\s*(px|pt|rem|em)\b", re.IGNORECASE)
_CLASS_ATTR_RE = re.compile(r"""class\s*=\s*(?:"([^"]*)"|'([^']*)')""", re.IGNORECASE)
_SELECTOR_CLASS_RE = re.compile(r"\.(-?[A-Za-z_][A-Za-z0-9_-]*)")


def _css_px(value: str) -> float | None:
    """First length in a CSS value, normalized to px. None when unitless."""
    match = _LENGTH_RE.search(value)
    if not match:
        return None
    try:
        number = float(match.group(1))
    except ValueError:
        return None
    unit = match.group(2).lower()
    if unit == "px":
        return number
    if unit == "pt":
        return number / 0.75
    return number * ROOT_FONT_PX  # rem and em, against the 16px root default


def _used_classes(text: str) -> set[str]:
    used: set[str] = set()
    for quoted, single in _CLASS_ATTR_RE.findall(text):
        used.update((quoted or single).split())
    return used


def _selector_applies(selector: str, used: set[str]) -> bool:
    """True when some comma-branch of the selector could match this document.

    A rule whose classes never appear in the markup is dead weight: a shared
    stylesheet copied across artboards would otherwise put a cover headline's
    64px into the type ramp of every body slide that never uses it.
    """
    for branch in selector.split(","):
        classes = _SELECTOR_CLASS_RE.findall(branch)
        if not classes or all(name in used for name in classes):
            return True
    return False


def _declaration_blocks(text: str) -> list[dict[str, str]]:
    """Every applicable CSS rule body and inline style, as property maps."""
    used = _used_classes(text)
    bodies = [body for selector, body in _RULE_RE.findall(text) if _selector_applies(selector, used)]
    for quoted, single in _INLINE_STYLE_RE.findall(text):
        bodies.append(quoted or single)

    blocks: list[dict[str, str]] = []
    for body in bodies:
        props: dict[str, str] = {}
        for declaration in body.split(";"):
            name, sep, value = declaration.partition(":")
            if not sep:
                continue
            props[name.strip().lower()] = value.strip()
        if props:
            blocks.append(props)
    return blocks


def _first_font_family(value: str) -> str | None:
    first = value.split(",")[0].strip().strip("\"'")
    return first or None


def _collect_colors(blocks: list[dict[str, str]]) -> dict[tuple[int, int, int], int]:
    counts: dict[tuple[int, int, int], int] = {}
    for props in blocks:
        for value in props.values():
            for six, three in _HEX_RE.findall(value):
                code = six or "".join(ch * 2 for ch in three)
                rgb = (int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))
                counts[rgb] = counts.get(rgb, 0) + 1
            for r, g, b in _RGB_RE.findall(value):
                rgb = (min(int(r), 255), min(int(g), 255), min(int(b), 255))
                counts[rgb] = counts.get(rgb, 0) + 1
    return counts


def _resolve_canvas_dir(canvas_arg: str) -> tuple[Path, dict[str, Any] | None]:
    """Accept either the working directory or the canvas.json inside it."""
    path = Path(canvas_arg)
    if path.is_dir():
        canvas_dir, manifest_path = path, path / CANVAS_FILE
    else:
        canvas_dir, manifest_path = path.parent, path
    manifest = None
    if manifest_path.is_file():
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    return canvas_dir, manifest


def _reference_artboard(
    manifest: dict[str, Any] | None, artboards: list[Path]
) -> tuple[Path | None, dict[str, Any] | None]:
    """The artboard whose frame sets the px scale: Main, then launch, then first listed."""
    by_name = {path.name: path for path in artboards}
    entries = (manifest or {}).get("artboards") or []
    listed = [e for e in entries if isinstance(e, dict) and e.get("file") in by_name]

    preferred = ["Main" + ARTBOARD_SUFFIX]
    launch = (manifest or {}).get("launch") or {}
    if isinstance(launch, dict) and isinstance(launch.get("file"), str):
        preferred.append(launch["file"])

    for name in preferred:
        for entry in listed:
            if entry["file"] == name:
                return by_name[name], entry
    if listed:
        return by_name[listed[0]["file"]], listed[0]
    return (artboards[0], None) if artboards else (None, None)


def derive_from_canvas(canvas_arg: str, notes: list[str]) -> dict[str, Any]:
    """Measure fonts, colors, grid and radius from a /design canvas working tree.

    The canvas is authored in CSS px and a slide renders on a fixed-inch stage,
    so the reference artboard's frame width fixes px -> inches for the whole
    profile. That is the point of this input: the same numbers the designer saw.
    """
    canvas_dir, manifest = _resolve_canvas_dir(canvas_arg)
    if not canvas_dir.is_dir():
        raise OSError(f"canvas directory not found: {canvas_dir}")

    artboards = sorted(p for p in canvas_dir.glob("*" + ARTBOARD_SUFFIX) if p.is_file())
    if not artboards:
        raise OSError(f"no *{ARTBOARD_SUFFIX} artboards in {canvas_dir}")
    if manifest is None:
        notes.append(f"defaulted: no {CANVAS_FILE} beside the artboards; frame size unknown")

    result: dict[str, Any] = {}
    reference, entry = _reference_artboard(manifest, artboards)
    notes.append(f"derived: reference artboard ({reference.name}, {len(artboards)} on canvas)")

    # --- scale: artboard px -> slide inches -------------------------------
    frame_w = float(entry["w"]) if entry and isinstance(entry.get("w"), (int, float)) else None
    frame_h = float(entry["h"]) if entry and isinstance(entry.get("h"), (int, float)) else None

    aspect = None
    if frame_w and frame_h:
        ratio = frame_w / frame_h
        if abs(ratio - (16 / 9)) < 0.02:
            aspect = "16:9"
        elif abs(ratio - (4 / 3)) < 0.02:
            aspect = "4:3"
        else:
            aspect = "custom"
        result["aspect_ratio"] = aspect
        notes.append(
            f"derived: template.aspect_ratio ({aspect}, frame {frame_w:g}x{frame_h:g}px)"
        )

    inches_per_px: float | None = None
    if frame_w and aspect in SLIDE_WIDTH_INCHES:
        inches_per_px = SLIDE_WIDTH_INCHES[aspect] / frame_w
        notes.append(
            f"derived: px scale (1px = {inches_per_px:.5f}in = "
            f"{inches_per_px * PT_PER_INCH:.3f}pt at {frame_w:g}px per "
            f"{SLIDE_WIDTH_INCHES[aspect]:g}in slide)"
        )
    else:
        notes.append(
            "defaulted: px scale (reference artboard has no 16:9/4:3 frame in "
            f"{CANVAS_FILE}); typography and geometry left at template defaults"
        )

    def to_pt(px: float) -> float:
        return round(px * inches_per_px * PT_PER_INCH, 1)  # type: ignore[operator]

    def to_inches(px: float) -> float:
        return round(px * inches_per_px, 3)  # type: ignore[operator]

    # --- read every artboard ----------------------------------------------
    blocks: list[dict[str, str]] = []
    reference_blocks: list[dict[str, str]] = []
    for artboard in artboards:
        parsed = _declaration_blocks(artboard.read_text(encoding="utf-8", errors="replace"))
        blocks.extend(parsed)
        if artboard == reference:
            reference_blocks = parsed

    sizes_px: list[float] = []
    families: list[str] = []
    size_family: list[tuple[float, str]] = []
    radii_px: list[float] = []
    grid_columns: list[int] = []
    gutters_px: list[float] = []
    root_padding_px: list[float] = []

    for props in blocks:
        size = _css_px(props["font-size"]) if "font-size" in props else None
        family = _first_font_family(props["font-family"]) if "font-family" in props else None
        if family:
            families.append(family)
        if size and size > 0 and family:
            size_family.append((size, family))

        if "border-radius" in props:
            radius = _css_px(props["border-radius"])
            if radius is not None and radius > 0:
                radii_px.append(radius)

        if "grid-template-columns" in props:
            value = props["grid-template-columns"]
            repeat = _REPEAT_RE.search(value)
            if repeat:
                grid_columns.append(int(repeat.group(1)))
            elif value.strip():
                grid_columns.append(len(value.split()))

        for gap_prop in ("gap", "column-gap", "grid-gap", "grid-column-gap"):
            if gap_prop in props:
                gap = _css_px(props[gap_prop])
                if gap is not None and gap > 0:
                    gutters_px.append(gap)

        # A root that fills the frame carries the safe margin as its padding.
        if "padding" in props and frame_w:
            width_px = _css_px(props.get("width", "")) if "width" in props else None
            if width_px is not None and abs(width_px - frame_w) < 2:
                pad = _css_px(props["padding"])
                if pad is not None and pad > 0:
                    root_padding_px.append(pad)

    for props in reference_blocks:
        size = _css_px(props["font-size"]) if "font-size" in props else None
        if size and size > 0:
            sizes_px.append(size)

    # --- typography --------------------------------------------------------
    if families:
        distinct = list(dict.fromkeys(families))
        result["max_font_families"] = len(distinct)
        heading = max(size_family, key=lambda pair: pair[0])[1] if size_family else distinct[0]
        body_pairs = [pair for pair in size_family if pair[1] != heading]
        body = min(body_pairs, key=lambda pair: pair[0])[1] if body_pairs else heading
        result["heading_font"] = heading
        result["body_font"] = body
        notes.append(
            f"derived: brand.heading_font/body_font ({heading} / {body}; "
            f"{len(distinct)} family(ies) declared)"
        )

    if sizes_px and inches_per_px:
        ordered = sorted(sizes_px)
        title_pt = to_pt(ordered[-1])
        caption_pt = to_pt(ordered[0])
        below_title = [s for s in ordered if s < ordered[-1]] or ordered
        body_pt = to_pt(below_title[len(below_title) // 2])
        if title_pt >= 18:
            result["title_pt"] = title_pt
            notes.append(
                f"derived: typography.title_pt ({title_pt}, largest text on "
                f"{reference.name})"
            )
        else:
            notes.append(f"defaulted: typography.title_pt (largest text is {title_pt}pt, below 18)")
        if body_pt >= 8:
            result["body_pt"] = body_pt
            notes.append(f"derived: typography.body_pt ({body_pt}, median below the title size)")
        if caption_pt >= 6:
            result["caption_pt"] = caption_pt
            notes.append(
                f"derived: typography.caption_pt ({caption_pt}, smallest text on "
                f"{reference.name})"
            )
    notes.append("defaulted: typography.metric_pt (a KPI number is not distinguishable from a heading in markup)")

    # --- colors ------------------------------------------------------------
    counts = _collect_colors(blocks)
    if counts:
        palette = [rgb for rgb, _ in sorted(counts.items(), key=lambda item: item[1], reverse=True)][:8]
        page = max(palette, key=_luminance)
        ink = min(palette, key=_luminance)
        remaining = [c for c in palette if c not in (page, ink)]
        accent = max(remaining, key=_saturation) if remaining else ink
        result["colors"] = {
            "page": "#{:02X}{:02X}{:02X}".format(*page),
            "ink": "#{:02X}{:02X}{:02X}".format(*ink),
            "accent": "#{:02X}{:02X}{:02X}".format(*accent),
        }
        notes.append(f"derived: brand.colors ({len(counts)} distinct color(s) declared, top 8 classified)")

    # --- geometry and composition -----------------------------------------
    # These three are exactly what the .pptx path documents as underivable.
    if grid_columns:
        columns = max(set(grid_columns), key=grid_columns.count)
        result["grid_columns"] = columns
        notes.append(f"derived: geometry.grid_columns ({columns}, most common grid track count)")
    if gutters_px and inches_per_px:
        gutter = to_inches(max(set(gutters_px), key=gutters_px.count))
        result["gutter_inches"] = gutter
        notes.append(f"derived: geometry.gutter_inches ({gutter}, most common grid gap)")
    if radii_px and inches_per_px:
        radius = to_inches(max(set(radii_px), key=radii_px.count))
        result["corner_radius"] = radius
        notes.append(f"derived: composition.corner_radius ({radius}, most common border-radius)")
    if root_padding_px and inches_per_px:
        margin = to_inches(min(root_padding_px))
        result["safe_margin_inches"] = margin
        notes.append(f"derived: geometry.safe_margin_inches ({margin}, padding on a frame-width root)")

    # --- archetypes from artboard names ------------------------------------
    known = {entry["id"] for entry in _load_archetype_catalog()}
    matched = set()
    for artboard in artboards:
        stem = artboard.name[: -len(ARTBOARD_SUFFIX)]
        slug = re.sub(r"[^a-z0-9]+", "-", stem.lower()).strip("-")
        if slug in known:
            matched.add(slug)
    if matched:
        result["archetypes"] = sorted(matched | STRUCTURAL_DEFAULT_ARCHETYPES)
        notes.append(f"derived: archetypes ({len(matched)} artboard name(s) matched the catalog)")

    return result


def derive_from_pptx(pptx_path: Path, notes: list[str]) -> dict[str, Any]:
    """Measure aspect ratio, fonts, and placeholder geometry from a reference .pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation(str(pptx_path))
    result: dict[str, Any] = {}

    slide_w, slide_h = prs.slide_width, prs.slide_height
    ratio = (slide_w / slide_h) if slide_h else 0
    if abs(ratio - (16 / 9)) < 0.02:
        aspect = "16:9"
    elif abs(ratio - (4 / 3)) < 0.02:
        aspect = "4:3"
    else:
        aspect = "custom"
    result["aspect_ratio"] = aspect
    notes.append(f"derived: template.aspect_ratio ({aspect})")

    title_bottoms_emu: list[int] = []
    footer_heights_emu: list[int] = []
    margin_candidates_in: list[float] = []
    title_sizes_pt: list[float] = []
    body_sizes_pt: list[float] = []
    font_names: set[str] = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                continue
            ph_type = shape.placeholder_format.type
            is_title = ph_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

            if slide_w and slide_h:
                margin_candidates_in.extend(
                    _inches(v)
                    for v in (
                        shape.left,
                        shape.top,
                        slide_w - (shape.left + shape.width),
                        slide_h - (shape.top + shape.height),
                    )
                    if v >= 0
                )

            if is_title:
                title_bottoms_emu.append(shape.top + shape.height)
            elif slide_h and (slide_h - (shape.top + shape.height)) < 0.6 * EMU_PER_INCH:
                footer_heights_emu.append(shape.height)

            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        font_names.add(run.font.name)
                    if run.font.size is not None:
                        (title_sizes_pt if is_title else body_sizes_pt).append(run.font.size.pt)

    if title_sizes_pt:
        title_pt = round(max(title_sizes_pt))
        if title_pt >= 18:
            result["title_pt"] = title_pt
            notes.append(f"derived: typography.title_pt ({title_pt})")
    if body_sizes_pt:
        body_pt = round(min(body_sizes_pt))
        if body_pt >= 8:
            result["body_pt"] = body_pt
            notes.append(f"derived: typography.body_pt ({body_pt})")
    if font_names:
        sorted_fonts = sorted(font_names)
        result["max_font_families"] = len(sorted_fonts)
        result["heading_font"] = sorted_fonts[0]
        result["body_font"] = sorted_fonts[-1] if len(sorted_fonts) > 1 else sorted_fonts[0]
        notes.append(
            f"derived: typography.max_font_families and brand.heading_font/body_font "
            f"({len(sorted_fonts)} font(s) found)"
        )

    if title_bottoms_emu:
        title_zone = _inches(max(title_bottoms_emu))
        if title_zone > 0:
            result["title_zone_inches"] = round(title_zone, 2)
            notes.append(f"derived: geometry.title_zone_inches ({result['title_zone_inches']})")
    if margin_candidates_in:
        safe_margin = round(min(margin_candidates_in), 2)
        if safe_margin >= 0:
            result["safe_margin_inches"] = safe_margin
            notes.append(f"derived: geometry.safe_margin_inches ({safe_margin})")
    if footer_heights_emu:
        footer_zone = round(_inches(max(footer_heights_emu)), 2)
        result["footer_zone_inches"] = footer_zone
        notes.append(f"derived: geometry.footer_zone_inches ({footer_zone})")

    return result


def build_profile(args: argparse.Namespace) -> tuple[dict[str, Any], list[str]]:
    profile = _load_default_template()
    notes: list[str] = []

    source_label = args.canvas or args.pptx or args.evidence
    profile["template"]["source"] = source_label
    notes.append(f"derived: template.source ({source_label})")

    evidence: dict[str, Any] | None = None
    evidence_dir: Path | None = None
    if args.evidence:
        evidence_path = Path(args.evidence)
        evidence = json.loads(evidence_path.read_text(encoding="utf-8"))
        evidence_dir = evidence_path.parent

    if args.pptx:
        pptx_fields = derive_from_pptx(Path(args.pptx), notes)
        if "aspect_ratio" in pptx_fields:
            profile["template"]["aspect_ratio"] = pptx_fields["aspect_ratio"]
        for key in ("title_pt", "body_pt", "max_font_families"):
            if key in pptx_fields:
                profile["typography"][key] = pptx_fields[key]
        for key in ("heading_font", "body_font"):
            if key in pptx_fields:
                profile["brand"][key] = pptx_fields[key]
        for key in ("title_zone_inches", "safe_margin_inches", "footer_zone_inches"):
            if key in pptx_fields:
                profile["geometry"][key] = pptx_fields[key]

    canvas_fields: dict[str, Any] = {}
    if args.canvas:
        # Applied after --pptx on purpose: a canvas is the design you intend,
        # a reference deck is the design you have.
        canvas_fields = derive_from_canvas(args.canvas, notes)
        if "aspect_ratio" in canvas_fields:
            profile["template"]["aspect_ratio"] = canvas_fields["aspect_ratio"]
        for key in ("title_pt", "body_pt", "caption_pt", "max_font_families"):
            if key in canvas_fields:
                profile["typography"][key] = canvas_fields[key]
        for key in ("heading_font", "body_font", "colors"):
            if key in canvas_fields:
                profile["brand"][key] = canvas_fields[key]
        for key in ("safe_margin_inches", "grid_columns", "gutter_inches"):
            if key in canvas_fields:
                profile["geometry"][key] = canvas_fields[key]
        if "corner_radius" in canvas_fields:
            profile["composition"]["corner_radius"] = canvas_fields["corner_radius"]

    for field in DEFAULTED_GEOMETRY_FIELDS:
        if field in canvas_fields:
            continue
        notes.append(f"defaulted: geometry.{field} (no reliable signal; kept template default)")
    for field in DEFAULTED_COMPOSITION_FIELDS:
        if field in canvas_fields:
            continue
        notes.append(f"defaulted: composition.{field} (no reliable signal; kept template default)")

    image_paths: list[Path] = []
    slide_word_counts: list[int] = []
    if evidence:
        for slide in evidence.get("slides", []):
            image = slide.get("image")
            if image:
                image_paths.append((evidence_dir / image) if evidence_dir else Path(image))
            text = slide.get("extracted_text") or ""
            if text:
                slide_word_counts.append(len(text.split()))

    if "colors" in canvas_fields:
        pass  # the canvas palette already won; slide images would only blur it
    else:
        colors = derive_colors(image_paths, notes) if image_paths else None
        if colors:
            profile["brand"]["colors"] = colors
        elif not image_paths:
            notes.append(
                "defaulted: brand.colors (no slide images in evidence; kept template default)"
            )

    catalog = _load_archetype_catalog()
    if "archetypes" in canvas_fields:
        profile["archetypes"] = canvas_fields["archetypes"]
    elif slide_word_counts:
        candidate_ids = set(STRUCTURAL_DEFAULT_ARCHETYPES)
        for archetype in catalog:
            max_words = archetype["max_words"]
            if any(count <= max_words for count in slide_word_counts):
                candidate_ids.add(archetype["id"])
        profile["archetypes"] = sorted(candidate_ids)
        notes.append(f"derived: archetypes ({len(candidate_ids)} candidate(s) from slide word counts)")
    else:
        notes.append("defaulted: archetypes (no evidence slide text; kept template default)")

    return profile, notes


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--run", required=True, help="Deck run directory")
    parser.add_argument("--evidence", help="Path to presentation-evidence.json")
    parser.add_argument("--pptx", help="Path to a reference .pptx file")
    parser.add_argument(
        "--canvas",
        help="Path to a /design canvas working tree (the directory holding "
        "canvas.json and *.dc.html), or to the canvas.json itself",
    )
    parser.add_argument(
        "--out", help="Draft output path (default: <run>/draft-template-profile.json)"
    )
    args = parser.parse_args()

    if not args.evidence and not args.pptx and not args.canvas:
        print(
            "usage: derive_template_profile.py --run <run-dir> "
            "[--evidence <presentation-evidence.json>] [--pptx <reference.pptx>] "
            "[--canvas <canvas-dir>]",
            file=sys.stderr,
        )
        print(
            "error: at least one of --evidence, --pptx or --canvas is required",
            file=sys.stderr,
        )
        return 2

    run_dir = Path(args.run)
    run_dir.mkdir(parents=True, exist_ok=True)
    out_path = Path(args.out) if args.out else run_dir / "draft-template-profile.json"

    try:
        profile, notes = build_profile(args)
    except (OSError, json.JSONDecodeError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    out_path.write_text(json.dumps(profile, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    for note in notes:
        print(note)
    print(f"draft written: {out_path}")
    print("Review, then: cp", out_path, "<run-dir>/template-profile.json")
    print("Validate with: python3 scripts/validate_template_profile.py <run-dir>/template-profile.json")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
