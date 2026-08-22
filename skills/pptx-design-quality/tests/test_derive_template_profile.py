from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
PIL = pytest.importorskip("PIL")

from pptx.util import Emu, Pt  # noqa: E402
from PIL import Image  # noqa: E402


SKILL_DIR = Path(__file__).resolve().parents[1]
DERIVER = SKILL_DIR / "scripts/derive_template_profile.py"
VALIDATOR = SKILL_DIR / "scripts/validate_template_profile.py"


def run_deriver(args: list[str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(DERIVER), *args],
        check=False,
        capture_output=True,
        text=True,
    )


def run_validator(path: Path) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(VALIDATOR), str(path)],
        check=False,
        capture_output=True,
        text=True,
    )


def _write_evidence(tmp_path: Path, slides: list[dict]) -> Path:
    evidence = {
        "contract_version": "1.0",
        "bundle": {"id": "test-bundle", "title": "Test Deck", "root": str(tmp_path)},
        "sources": [],
        "slides": slides,
        "transcript_segments": [],
        "frames": [],
    }
    path = tmp_path / "presentation-evidence.json"
    path.write_text(json.dumps(evidence), encoding="utf-8")
    return path


def _write_slide_image(tmp_path: Path, name: str, color: tuple[int, int, int]) -> str:
    image_dir = tmp_path / "slide_images"
    image_dir.mkdir(exist_ok=True)
    Image.new("RGB", (100, 56), color).save(image_dir / name)
    return f"slide_images/{name}"


def test_missing_inputs_errors_cleanly(tmp_path: Path) -> None:
    result = run_deriver(["--run", str(tmp_path)])

    assert result.returncode == 2
    assert "at least one of --evidence, --pptx or --canvas is required" in result.stderr


def test_derive_from_evidence_only_produces_valid_draft(tmp_path: Path) -> None:
    image = _write_slide_image(tmp_path, "slide_1.png", (255, 255, 255))
    evidence_path = _write_evidence(
        tmp_path,
        [
            {
                "id": "slide-1",
                "page": 1,
                "source_id": "src-1",
                "image": image,
                "image_sha256": None,
                "extracted_text": "Short title slide",
                "summary": None,
                "links": [],
                "transcript_segment_ids": [],
                "frame_ids": [],
            }
        ],
    )

    result = run_deriver(["--run", str(tmp_path), "--evidence", str(evidence_path)])

    assert result.returncode == 0, result.stderr
    draft_path = tmp_path / "draft-template-profile.json"
    assert draft_path.is_file()
    assert "derived: brand.colors" in result.stdout
    assert "derived: archetypes" in result.stdout

    validated = run_validator(draft_path)
    assert validated.returncode == 0, validated.stderr


def _build_reference_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 16:9 at 13.33in wide
    prs.slide_height = Emu(6858000)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_shape = slide.shapes.title
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.3)
    title_shape.width = Inches(11.0)
    title_shape.height = Inches(1.0)
    title_shape.text_frame.text = "Reference Deck Title"
    run = title_shape.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(32)
    run.font.name = "Aptos Display"

    body_box = slide.shapes.placeholders[1] if len(slide.placeholders) > 1 else None
    if body_box is None:
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(10.0), Inches(4.0))
    body_box.text_frame.text = "Body content example"
    body_run = body_box.text_frame.paragraphs[0].runs[0]
    body_run.font.size = Pt(16)
    body_run.font.name = "Aptos"

    out_path = tmp_path / "reference.pptx"
    prs.save(str(out_path))
    return out_path


def test_derive_from_pptx_extracts_geometry_and_fonts(tmp_path: Path) -> None:
    pptx_path = _build_reference_pptx(tmp_path)

    result = run_deriver(["--run", str(tmp_path), "--pptx", str(pptx_path)])

    assert result.returncode == 0, result.stderr
    draft = json.loads((tmp_path / "draft-template-profile.json").read_text(encoding="utf-8"))

    assert draft["template"]["aspect_ratio"] == "16:9"
    assert draft["typography"]["title_pt"] == 32
    assert "Aptos Display" in {draft["brand"]["heading_font"], draft["brand"]["body_font"]}

    validated = run_validator(tmp_path / "draft-template-profile.json")
    assert validated.returncode == 0, validated.stderr


def test_output_never_overwrites_template_profile(tmp_path: Path) -> None:
    canonical = tmp_path / "template-profile.json"
    canonical.write_text("{}", encoding="utf-8")

    image = _write_slide_image(tmp_path, "slide_1.png", (10, 10, 10))
    evidence_path = _write_evidence(
        tmp_path,
        [
            {
                "id": "slide-1",
                "page": 1,
                "source_id": "src-1",
                "image": image,
                "image_sha256": None,
                "extracted_text": "Slide text",
                "summary": None,
                "links": [],
                "transcript_segment_ids": [],
                "frame_ids": [],
            }
        ],
    )

    result = run_deriver(["--run", str(tmp_path), "--evidence", str(evidence_path)])

    assert result.returncode == 0, result.stderr
    assert canonical.read_text(encoding="utf-8") == "{}"
    assert (tmp_path / "draft-template-profile.json").is_file()
