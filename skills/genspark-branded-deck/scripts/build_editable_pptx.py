#!/usr/bin/env python3
"""build_editable_pptx.py — hybrid EDITABLE deck.

Background = text-free design render (bg/slide-XX.png). On top, native PowerPoint
text boxes rebuilt from pos/slide-XX.json (bbox + runs + colour/bold/font/align).
Result: the exact design, with every heading and paragraph editable in PowerPoint.

  python3 build_editable_pptx.py --src <dir-with-bg-and-pos> --out <deck.pptx>
"""
import argparse, glob, json, os, re, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# Stage width is a CLI flag: this skill's own template is a 1280 stage, but a
# recovered Genspark deck is 1920. Hardcoding 1280 silently halves every
# coordinate and font size on a 1920 deck.
DEFAULT_STAGE = 1280
# Fallbacks only — the real family is captured per box by render_hybrid.mjs.
FONT_FALLBACK = {"mono": "Consolas", "display": "Segoe UI", "body": "Segoe UI"}
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def rgb(css):
    m = re.findall(r"[\d.]+", css or "")
    if len(m) >= 3:
        return RGBColor(int(float(m[0])), int(float(m[1])), int(float(m[2])))
    return RGBColor(0xE8, 0xEE, 0xF4)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", default="build")
    ap.add_argument("--out", default="build/deck-editable.pptx")
    ap.add_argument("--stage", type=int, default=DEFAULT_STAGE,
                    help="CSS width of the deck stage in px (1280 default, 1920 for recovered Genspark decks)")
    a = ap.parse_args()
    in_per_px = 13.333 / a.stage      # stage px -> slide inches (16:9 exact)
    pt_per_px = 960.0 / a.stage       # stage px -> points

    bgs = sorted(glob.glob(os.path.join(a.src, "bg", "slide-*.png")))
    if not bgs:
        sys.exit(f"No bg/slide-*.png in {a.src} — run render_hybrid.mjs first.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    nboxes = 0

    for bg in bgs:
        n = os.path.basename(bg)
        pos = os.path.join(a.src, "pos", n.replace(".png", ".json"))
        boxes = json.load(open(pos)) if os.path.exists(pos) else []
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        for b in boxes:
            tb = s.shapes.add_textbox(Inches(b["x"] * in_per_px), Inches(b["y"] * in_per_px),
                                      Inches(b["w"] * in_per_px), Inches(max(b["h"], b["fontPx"]) * in_per_px))
            tf = tb.text_frame
            # Preformatted boxes (code, file trees, ASCII diagrams) must NOT wrap:
            # their line breaks and spacing are the content. Wrapping reflows a
            # directory tree into prose while the PNG still looks correct.
            tf.word_wrap = not b.get("pre", False)
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink so text can't overflow its slot
            tf.vertical_anchor = MSO_ANCHOR.TOP
            for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(tf, m, Emu(0))
            for li, line in enumerate(b["lines"]):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = ALIGN.get(b["align"], PP_ALIGN.LEFT)
                # exact point spacing beats a ratio: see lhPx note in render_hybrid.mjs
                if b.get("lhPx"):
                    p.line_spacing = Pt(round(b["lhPx"] * pt_per_px, 2))
                else:
                    p.line_spacing = round(b.get("lh", 1.2), 3)
                for rn in line:
                    run = p.add_run()
                    run.text = rn["text"]
                    f = run.font
                    f.size = Pt(round(b["fontPx"] * pt_per_px, 1))
                    f.bold = bool(rn["bold"])
                    f.name = b.get("family") or FONT_FALLBACK.get(b.get("role"), "Segoe UI")
                    f.color.rgb = rgb(rn["color"])
            nboxes += 1

    os.makedirs(os.path.dirname(os.path.abspath(a.out)), exist_ok=True)
    prs.save(a.out)
    # round-trip sanity check
    Presentation(a.out)
    print(f"DONE: {len(bgs)} slides, {nboxes} editable text boxes, {a.stage}px stage -> {a.out}")


if __name__ == "__main__":
    main()
