#!/usr/bin/env python3
"""Compare visible slide content before promoting a PPTX redesign."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


def slide_text(slide: Any) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def editable_text_shapes(prs: Presentation) -> int:
    return sum(
        1
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def compare(previous: Path, candidate: Path) -> dict[str, Any]:
    old = Presentation(previous)
    new = Presentation(candidate)
    old_text = [slide_text(slide) for slide in old.slides]
    new_text = [slide_text(slide) for slide in new.slides]
    overlap = min(len(old_text), len(new_text))
    same_positions = [index + 1 for index in range(overlap) if old_text[index] == new_text[index]]
    changed_positions = [index + 1 for index in range(overlap) if old_text[index] != new_text[index]]
    old_text_set = set(old_text)
    identical_anywhere = [index + 1 for index, text in enumerate(new_text) if text in old_text_set]
    changed_or_replaced = len(changed_positions) + abs(len(old_text) - len(new_text))
    denominator = max(len(old_text), len(new_text), 1)
    material_change_ratio = changed_or_replaced / denominator
    material = (
        len(old_text) != len(new_text)
        or material_change_ratio >= 0.30
        or len(identical_anywhere) == 0
    )
    return {
        "previous": str(previous),
        "candidate": str(candidate),
        "previous_slide_count": len(old.slides),
        "candidate_slide_count": len(new.slides),
        "same_position_full_text_count": len(same_positions),
        "same_position_full_text_slides": same_positions,
        "changed_position_count": len(changed_positions),
        "changed_position_slides": changed_positions,
        "identical_slide_text_anywhere_count": len(identical_anywhere),
        "candidate_slides_identical_anywhere": identical_anywhere,
        "previous_editable_text_shape_count": editable_text_shapes(old),
        "candidate_editable_text_shape_count": editable_text_shapes(new),
        "material_change_ratio": round(material_change_ratio, 4),
        "material_change_passed": material,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("previous", type=Path)
    parser.add_argument("candidate", type=Path)
    parser.add_argument("--require-material", action="store_true")
    parser.add_argument("--json-out", type=Path)
    args = parser.parse_args()

    for path in (args.previous, args.candidate):
        if not path.is_file():
            parser.error(f"PPTX not found: {path}")

    result = compare(args.previous, args.candidate)
    rendered = json.dumps(result, indent=2)
    print(rendered)
    if args.json_out:
        args.json_out.parent.mkdir(parents=True, exist_ok=True)
        args.json_out.write_text(rendered + "\n", encoding="utf-8")
    return 1 if args.require_material and not result["material_change_passed"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
