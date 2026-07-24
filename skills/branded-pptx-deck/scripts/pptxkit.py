#!/usr/bin/env python3
"""pptxkit — a small, reusable toolkit for building branded 16:9 PPTX decks.

Why this exists: hand-positioned python-pptx decks repeatedly hit two problems —
(1) text boxes overflow into each other, and (2) a shape-shadow bug that makes
PowerPoint show the "repair" dialog. This module solves both:

  * `text(...)` supports shrink-to-fit and never relies on a fixed-height box it
    can overflow silently.
  * `rect(..., shadow=True)` reuses the single <a:effectLst> instead of writing a
    second one (the repair-prompt bug).
  * `save(...)` runs a structural validation before returning.

Import this from a deck builder script:

    from pptxkit import Brand, Deck, PP_ALIGN, Inches, Pt
    d = Deck()                       # 16:9, blank
    s = d.slide(fill=d.b.NAVY)
    d.rect(s, 0, 0, d.W, d.H, d.b.NAVY)
    d.text(s, "Hello", d.M, Inches(1), d.CW, Inches(1), size=40, color=d.b.WHITE, bold=True)
    d.save("out.pptx")               # validates, raises on malformed XML

Keep slide *content* in your builder; keep *mechanics* here.
"""
from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

__all__ = ["Brand", "Deck", "RGBColor", "PP_ALIGN", "MSO_ANCHOR", "Inches", "Pt", "Emu"]


def hx(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


@dataclass
class Brand:
    """Swap this to re-skin a whole deck. Defaults = the user's Canva-Pro palette."""
    NAVY: RGBColor = hx("0A1628")
    NAVY_2: RGBColor = hx("12243A")
    TEAL: RGBColor = hx("00C9A7")
    ACCENT: RGBColor = hx("009B82")
    DARK_TEAL: RGBColor = hx("008F75")
    LIGHT_TEAL: RGBColor = hx("E0F7F1")
    GOLD: RGBColor = hx("FFB800")
    AMBER: RGBColor = hx("F2A83B")
    CORAL: RGBColor = hx("E05A6B")
    WHITE: RGBColor = hx("FFFFFF")
    SOFT: RGBColor = hx("F4F7F8")
    INK: RGBColor = hx("1B2B3C")
    MUTED: RGBColor = hx("5B6B7C")
    GRID: RGBColor = hx("D9DFE5")
    FONT: str = "Calibri"
    FONT_H: str = "Calibri"
    # matplotlib hex mirrors
    HX_TEAL: str = "#00C9A7"
    HX_TEALD: str = "#009B82"
    HX_NAVY: str = "#0A1628"
    HX_INK: str = "#1B2B3C"
    HX_MUTED: str = "#5B6B7C"
    HX_GRID: str = "#D9DFE5"
    HX_GOLD: str = "#FFB800"


class Deck:
    def __init__(self, brand: Brand | None = None, footer: str = ""):
        self.b = brand or Brand()
        self.prs = Presentation()
        self.W = Inches(13.333)
        self.H = Inches(7.5)
        self.prs.slide_width = self.W
        self.prs.slide_height = self.H
        self.M = Inches(0.6)                      # margin
        self.CW = self.W - self.M * 2             # content width
        self.footer_text = footer
        self._wipe()

    # ── lifecycle ────────────────────────────────────────────────────────────
    def _wipe(self):
        for i in range(len(self.prs.slides._sldIdLst) - 1, -1, -1):
            rid = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rid)
            del self.prs.slides._sldIdLst[i]

    def slide(self, fill: RGBColor | None = None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        if fill is not None:
            self.rect(s, 0, 0, self.W, self.H, fill)
        return s

    @property
    def n(self):
        return len(self.prs.slides._sldIdLst)

    def save(self, path: str | Path, *, validate: bool = True):
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        if validate:
            problems = validate_pptx(path)
            if problems:
                raise RuntimeError(
                    "Deck saved but FAILED validation (would trigger PowerPoint repair):\n  - "
                    + "\n  - ".join(problems)
                )
        print(f"Saved {path} with {self.n} slides" + ("  [validated]" if validate else ""))
        return path

    # ── primitives ───────────────────────────────────────────────────────────
    def rect(self, s, left, top, width, height, fill, *, line=None, radius=0.0, shadow=False):
        st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        sh = s.shapes.add_shape(st, left, top, width, height)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Pt(1)
        if radius:
            try:
                sh.adjustments[0] = radius
            except Exception:
                pass
        sh.shadow.inherit = False  # writes one empty <a:effectLst/>
        if shadow:
            spPr = sh._element.spPr
            eff = spPr.find(qn("a:effectLst"))  # REUSE it — never append a 2nd one
            if eff is None:
                eff = spPr.makeelement(qn("a:effectLst"), {})
                spPr.append(eff)
            sdw = eff.makeelement(
                qn("a:outerShdw"),
                {"blurRad": "80000", "dist": "30000", "dir": "5400000", "rotWithShape": "0"},
            )
            c = sdw.makeelement(qn("a:srgbClr"), {"val": "0A1628"})
            a = c.makeelement(qn("a:alpha"), {"val": "20000"})
            c.append(a)
            sdw.append(c)
            eff.append(sdw)
        return sh

    def text(self, s, runs, left, top, width, height, *, size=14, color=None, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP, ls=1.05,
             shrink=False, wrap=True):
        """runs: a str, or a list of str / dicts. dict keys: text,size,color,bold,
        italic,align,font,bullet,space_before. Use shrink=True for any box whose
        text length is data-driven (prevents silent overflow)."""
        color = self.b.INK if color is None else color
        font = self.b.FONT if font is None else font
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if shrink else MSO_AUTO_SIZE.NONE
        paras = runs if isinstance(runs, list) else [runs]
        for i, spec in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(spec, dict):
                content = spec.get("text", "")
                ps, pc = spec.get("size", size), spec.get("color", color)
                pb, pi = spec.get("bold", bold), spec.get("italic", italic)
                pa, pf = spec.get("align", align), spec.get("font", font)
                bullet = spec.get("bullet", False)
                p.space_before = Pt(spec.get("space_before", 0))
            else:
                content, ps, pc, pb, pi, pa, pf, bullet = spec, size, color, bold, italic, align, font, False
                p.space_before = Pt(0)
            p.alignment = pa
            try:
                p.line_spacing = ls
            except Exception:
                pass
            r = p.add_run()
            r.text = ("•  " + content) if bullet else content
            r.font.name = pf
            r.font.size = Pt(ps)
            r.font.bold = pb
            r.font.italic = pi
            r.font.color.rgb = pc
        return box

    # ── common furniture ──────────────────────────────────────────────────────
    def header(self, s, title, subtitle=None, *, band=True):
        b = self.b
        if band:
            self.rect(s, 0, 0, self.W, Inches(0.16), b.TEAL)
        self.text(s, title, self.M, Inches(0.42), self.CW, Inches(0.66), size=30,
                  color=b.NAVY, bold=True, font=b.FONT_H, shrink=True)
        self.rect(s, self.M, Inches(1.12), Inches(1.45), Inches(0.05), b.TEAL)
        if subtitle:
            self.text(s, subtitle, self.M, Inches(1.26), self.CW, Inches(0.4), size=14.5, color=b.MUTED)

    def footer(self, s, page, total, *, dark=False):
        col = RGBColor(0x88, 0x90, 0x98) if dark else self.b.MUTED
        if self.footer_text:
            self.text(s, self.footer_text, self.M, Inches(7.08), Inches(9), Inches(0.3), size=9, color=col)
        self.text(s, f"{page} / {total}", self.W - Inches(1.3), Inches(7.06), Inches(0.7), Inches(0.3),
                  size=10, color=col, align=PP_ALIGN.RIGHT, bold=True)

    # ── charts (matplotlib -> transparent PNG) ────────────────────────────────
    def chart_barh(self, labels, values, out_path, *, highlight_at=6):
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        b = self.b
        order = sorted(zip(labels, values), key=lambda t: t[1], reverse=True)
        labels = [t[0] for t in order]
        values = [t[1] for t in order]
        fig, ax = plt.subplots(figsize=(9.6, 5.1), dpi=200)
        fig.patch.set_alpha(0); ax.set_facecolor("none")
        y = range(len(labels))
        colors = [b.HX_TEAL if v >= highlight_at else b.HX_TEALD for v in values]
        ax.barh(y, values, color=colors, height=0.6, zorder=3)
        ax.set_yticks(list(y)); ax.set_yticklabels(labels, fontsize=12.5, color=b.HX_INK, fontweight="bold")
        ax.invert_yaxis(); ax.set_xlim(0, max(values) + 1.2); ax.set_xticks([])
        for i, v in enumerate(values):
            ax.text(v + 0.15, i, str(v), va="center", fontsize=13.5, color=b.HX_NAVY, fontweight="bold")
        for sp in ax.spines.values():
            sp.set_visible(False)
        ax.tick_params(length=0)
        plt.tight_layout(pad=0.4)
        fig.savefig(out_path, transparent=True, bbox_inches="tight"); plt.close(fig)
        return out_path

    def chart_matrix(self, points, out_path, *, xlabel="", ylabel="", note=""):
        """points: list of (label, x[0..1], y[0..1], hex_color)."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.patches import FancyBboxPatch
        b = self.b
        fig, ax = plt.subplots(figsize=(9.4, 5.3), dpi=200)
        fig.patch.set_alpha(0); ax.set_facecolor("none")
        ax.add_patch(FancyBboxPatch((0.5, 0.5), 0.5, 0.5, boxstyle="square,pad=0",
                                    facecolor=b.HX_TEAL, alpha=0.10, edgecolor="none"))
        ax.axhline(0.5, color=b.HX_GRID, lw=1.2); ax.axvline(0.5, color=b.HX_GRID, lw=1.2)
        for lab, x, y, c in points:
            ax.scatter(x, y, s=240, color=c, alpha=0.92, zorder=3, edgecolors="white", linewidths=1.5)
            ax.annotate(lab, (x, y), xytext=(0, 12), textcoords="offset points", ha="center",
                        fontsize=10.5, color=b.HX_INK, fontweight="bold")
        ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.set_xticks([]); ax.set_yticks([])
        ax.set_xlabel(xlabel, fontsize=11.5, color=b.HX_MUTED, labelpad=10)
        ax.set_ylabel(ylabel, fontsize=11.5, color=b.HX_MUTED, labelpad=10)
        if note:
            ax.text(0.985, 0.96, note, ha="right", va="top", fontsize=11, color=b.HX_TEALD,
                    fontweight="bold", style="italic")
        for sp in ax.spines.values():
            sp.set_color(b.HX_GRID)
        plt.tight_layout(pad=0.6)
        fig.savefig(out_path, transparent=True, bbox_inches="tight"); plt.close(fig)
        return out_path

    def picture_centered(self, s, img, *, top=Inches(1.6), width=Inches(9.4), max_bottom=Inches(6.1)):
        pic = s.shapes.add_picture(str(img), 0, 0, width=width)
        pic.left = int((self.W - pic.width) / 2)
        pic.top = top
        if pic.top + pic.height > max_bottom:
            sc = (max_bottom - top) / pic.height
            pic.width = int(pic.width * sc); pic.height = int(pic.height * sc)
            pic.left = int((self.W - pic.width) / 2)
        return pic


# ── standalone validators (also importable) ──────────────────────────────────


def validate_pptx(path: str | Path) -> list[str]:
    """Return a list of structural problems that would trigger PowerPoint repair.
    Empty list == clean."""
    path = Path(path)
    problems: list[str] = []
    try:
        z = zipfile.ZipFile(path)
    except Exception as e:  # not a valid zip
        return [f"not a readable .pptx (zip): {e}"]
    for n in z.namelist():
        if re.match(r"ppt/slides/slide\d+\.xml$", n):
            xml = z.read(n).decode("utf-8", "ignore")
            for sp in re.findall(r"<p:spPr>.*?</p:spPr>", xml, re.S):
                if sp.count("<a:effectLst") > 1:
                    problems.append(f"{n}: a shape has >1 <a:effectLst> (shadow bug)")
                if "<a:effectLst" in sp and "<a:ln" in sp and sp.index("<a:ln") > sp.index("<a:effectLst"):
                    problems.append(f"{n}: <a:effectLst> precedes <a:ln> (bad child order)")
    # round-trip
    try:
        Presentation(str(path))
    except Exception as e:
        problems.append(f"python-pptx cannot reopen: {e}")
    return problems


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        probs = validate_pptx(sys.argv[1])
        print("CLEAN" if not probs else "PROBLEMS:\n  - " + "\n  - ".join(probs))
    else:
        # smoke test
        d = Deck(footer="pptxkit demo")
        s = d.slide(fill=d.b.NAVY)
        d.text(s, "pptxkit", d.M, Inches(2.6), d.CW, Inches(1.2), size=60, color=d.b.WHITE, bold=True)
        d.text(s, "branded deck toolkit", d.M, Inches(3.8), d.CW, Inches(0.6), size=22, color=d.b.TEAL)
        d.footer(s, 1, 1, dark=True)
        d.save("/tmp/pptxkit_demo.pptx")
