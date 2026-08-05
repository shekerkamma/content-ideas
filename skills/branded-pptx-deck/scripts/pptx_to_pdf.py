#!/usr/bin/env python3
"""Render a .pptx to PDF using only python-pptx and the stdlib.

LAST-RESORT FALLBACK. Prefer `render_pptx.sh`, which drives LibreOffice and
produces a true render. Reach for this only on a host with no LibreOffice at
all (no WSL interop, no native install), where the alternative is no visual
QA whatsoever and decks stuck at `draft` forever.

Why it can exist: PDF is a text-based format, so one can be emitted directly.

This is a QA proxy, not a faithful renderer. It draws shape rectangles, fills,
and text with base-14 fonts at approximate metrics. It is good enough to see
collisions, overflow, off-canvas content, and layout balance — which is what
visual QA is for. Native charts are drawn as labelled placeholders (their
geometry is real; their contents are not).

Usage:
    python3 pptx_to_pdf.py DECK.pptx [-o OUT.pdf]
"""

import argparse
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
PT_IN = 72.0
AVG_CW = 0.50          # base-14 average char width as fraction of size
LINE_H = 1.20


def pt(v):
    return (v / EMU_IN) * PT_IN


def esc(s):
    return s.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def ascii_only(s):
    """Base-14 WinAnsi can't carry the em dashes / rupee signs used in the deck."""
    repl = {"—": "-", "–": "-", "‘": "'", "’": "'",
            "“": '"', "”": '"', "•": "*", "₹": "Rs ",
            "×": "x", "→": "->", "≥": ">=", "≤": "<=",
            "±": "+/-", "…": "...", " ": " "}
    for k, v in repl.items():
        s = s.replace(k, v)
    return "".join(c if 32 <= ord(c) < 127 else "?" for c in s)


def shape_fill_rgb(shape):
    try:
        f = shape.fill
        if f.type is not None and f.type == 1:  # solid
            c = f.fore_color
            if c.type == 1:  # RGB
                r, g, b = c.rgb[0], c.rgb[1], c.rgb[2]
                return r / 255, g / 255, b / 255
    except Exception:
        pass
    return None


def run_color(run):
    try:
        c = run.font.color
        if c and c.type == 1:
            return c.rgb[0] / 255, c.rgb[1] / 255, c.rgb[2] / 255
    except Exception:
        pass
    return 0.1, 0.17, 0.24


def wrap(text, width_pt, size):
    cw = max(size * AVG_CW, 0.1)
    n = max(int(width_pt / cw), 1)
    words, lines, cur = text.split(), [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if len(trial) <= n:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


def render_slide(slide, W, H):
    """Return a PDF content stream for one slide."""
    out = ["q", "1 1 1 rg", f"0 0 {W:.1f} {H:.1f} re f"]

    for shape in slide.shapes:
        try:
            x, y = pt(shape.left), pt(shape.top)
            w, h = pt(shape.width), pt(shape.height)
        except Exception:
            continue
        yp = H - y - h  # flip to PDF origin

        rgb = shape_fill_rgb(shape)
        if rgb and not shape.has_text_frame:
            out.append(f"{rgb[0]:.3f} {rgb[1]:.3f} {rgb[2]:.3f} rg")
            out.append(f"{x:.1f} {yp:.1f} {w:.1f} {h:.1f} re f")
        elif rgb:
            out.append(f"{rgb[0]:.3f} {rgb[1]:.3f} {rgb[2]:.3f} rg")
            out.append(f"{x:.1f} {yp:.1f} {w:.1f} {h:.1f} re f")

        if getattr(shape, "has_chart", False) and shape.has_chart:
            out.append("0.85 0.89 0.92 rg")
            out.append(f"{x:.1f} {yp:.1f} {w:.1f} {h:.1f} re f")
            out.append("0.35 0.42 0.49 RG 1 w")
            out.append(f"{x:.1f} {yp:.1f} {w:.1f} {h:.1f} re S")
            label = "[ native chart ]"
            out += ["BT", "/F2 11 Tf", "0.2 0.28 0.35 rg",
                    f"{x + w / 2 - len(label) * 2.6:.1f} {yp + h / 2:.1f} Td",
                    f"({esc(label)}) Tj", "ET"]
            continue

        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.text.strip():
            continue

        cursor = H - y - 2
        for para in tf.paragraphs:
            txt = "".join(r.text or "" for r in para.runs)
            if not txt.strip():
                cursor -= 6
                continue
            sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
            size = max(sizes) if sizes else 12.0
            bold = any(r.font.bold for r in para.runs)
            col = run_color(para.runs[0]) if para.runs else (0.1, 0.17, 0.24)
            try:
                if para.space_before is not None:
                    cursor -= para.space_before.pt
            except Exception:
                pass
            font = "/F2" if bold else "/F1"
            align = str(getattr(para, "alignment", "") or "")
            for line in wrap(ascii_only(txt), w, size):
                cursor -= size * LINE_H
                if cursor < -20:
                    break
                lx = x
                if "CENTER" in align:
                    lx = x + (w - len(line) * size * AVG_CW) / 2
                elif "RIGHT" in align:
                    lx = x + w - len(line) * size * AVG_CW
                out += ["BT", f"{font} {size:.1f} Tf",
                        f"{col[0]:.3f} {col[1]:.3f} {col[2]:.3f} rg",
                        f"{lx:.1f} {cursor:.1f} Td", f"({esc(line)}) Tj", "ET"]
    out.append("Q")
    return "\n".join(out).encode("latin-1", "replace")


def build_pdf(prs, out_path):
    W, H = pt(prs.slide_width), pt(prs.slide_height)
    objs = {}          # number -> bytes
    streams = []
    for slide in prs.slides:
        streams.append(zlib.compress(render_slide(slide, W, H)))

    n_pages = len(streams)
    # object plan: 1 catalog, 2 pages, 3 F1, 4 F2, then per page: page obj + content
    first_page_obj = 5
    page_ids = [first_page_obj + 2 * i for i in range(n_pages)]
    content_ids = [first_page_obj + 2 * i + 1 for i in range(n_pages)]

    objs[1] = b"<< /Type /Catalog /Pages 2 0 R >>"
    kids = " ".join(f"{i} 0 R" for i in page_ids)
    objs[2] = f"<< /Type /Pages /Count {n_pages} /Kids [{kids}] >>".encode()
    objs[3] = b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica /Encoding /WinAnsiEncoding >>"
    objs[4] = b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold /Encoding /WinAnsiEncoding >>"

    for i in range(n_pages):
        objs[page_ids[i]] = (
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 {W:.1f} {H:.1f}] "
            f"/Resources << /Font << /F1 3 0 R /F2 4 0 R >> >> "
            f"/Contents {content_ids[i]} 0 R >>"
        ).encode()
        objs[content_ids[i]] = (
            f"<< /Length {len(streams[i])} /Filter /FlateDecode >>\nstream\n".encode()
            + streams[i] + b"\nendstream"
        )

    buf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = {}
    for num in sorted(objs):
        offsets[num] = len(buf)
        buf += f"{num} 0 obj\n".encode() + objs[num] + b"\nendobj\n"
    xref_at = len(buf)
    top = max(objs) + 1
    buf += f"xref\n0 {top}\n".encode()
    buf += b"0000000000 65535 f \n"
    for num in range(1, top):
        if num in offsets:
            buf += f"{offsets[num]:010d} 00000 n \n".encode()
        else:
            buf += b"0000000000 65535 f \n"
    buf += (f"trailer\n<< /Size {top} /Root 1 0 R >>\nstartxref\n{xref_at}\n"
            "%%EOF\n").encode()

    Path(out_path).write_bytes(bytes(buf))
    return out_path, n_pages


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck")
    ap.add_argument("-o", "--out")
    args = ap.parse_args()
    src = Path(args.deck)
    out = Path(args.out) if args.out else src.with_suffix(".preview.pdf")
    prs = Presentation(str(src))
    path, n = build_pdf(prs, out)
    print(f"wrote {path}  ({n} pages)")
    print("QA proxy only — shapes/text approximate, native charts are placeholders.")


if __name__ == "__main__":
    main()
