"""
Universal Branded PPTX Compiler — Client-Ready Template Builder
================================================================
Creates visually branded slides from scratch. Dynamically supports
any of the 21 strategy-consulting frameworks through a unified JSON schema.
"""
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG_DARK      = RGBColor(0x0A, 0x16, 0x28)
BG_TABLE_ROW = RGBColor(0x14, 0x26, 0x3E)
BG_TABLE_ALT = RGBColor(0x0E, 0x1C, 0x30)
ACCENT       = RGBColor(0x00, 0x9B, 0x82)
ACCENT_LIGHT = RGBColor(0x00, 0xC9, 0xA7)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT   = RGBColor(0x8A, 0x9B, 0xAE)
GREY_MID     = RGBColor(0x5B, 0x6B, 0x7C)
FONT         = 'Calibri'

def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def _add_text_box(slide, left, top, width, height, text, font_size=Pt(14), font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def _add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    return shape

def _add_finding_bullet(slide, left, top, width, number, text):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Pt(2), Pt(22), Pt(22))
    circle.line.fill.background()
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_text_box(slide, left + Pt(30), top, width - Pt(30), Pt(50), str(text), font_size=Pt(12), font_color=GREY_LIGHT)

def _add_table_row(slide, top, row_data, bg_color, is_header=False, col_widths=None):
    left_margin = Inches(0.6)
    row_h = Pt(36) if is_header else Pt(40)
    
    # Calculate dynamic column widths if not provided
    if not col_widths:
        total_width = Inches(12.1)
        col_widths = [total_width / len(row_data)] * len(row_data)
        
    row_w = sum(col_widths)
    _add_shape(slide, left_margin, top, row_w, row_h, bg_color)
    
    font_sz = Pt(10)
    font_clr = ACCENT_LIGHT if is_header else WHITE
    font_b = True if is_header else False
    
    x = left_margin + Pt(8)
    for i, (text, cw) in enumerate(zip(row_data, col_widths)):
        _add_text_box(slide, x, top + Pt(4), cw - Pt(16), row_h - Pt(8), str(text), font_size=font_sz, font_color=font_clr, bold=font_b)
        x += cw

def _add_footer(slide, data, slide_num, total_slides):
    _add_accent_line(slide, Inches(0.6), Inches(6.5), Inches(12.1))
    _add_text_box(slide, Inches(0.6), Inches(6.7), Inches(8), Pt(18), f"{data['company']} Analysis  |  Strategy Skill Pipe  |  June 2026", font_size=Pt(8), font_color=GREY_MID)
    _add_text_box(slide, Inches(11.5), Inches(6.7), Inches(1.2), Pt(18), f"{slide_num} / {total_slides}", font_size=Pt(8), font_color=GREY_MID, alignment=PP_ALIGN.RIGHT)

def build_title_slide(prs, data, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(2))
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(5), Pt(20), f"{data['company'].upper()}  \u00b7  STRATEGY ANALYSIS", font_size=Pt(9), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(1.5), f"{data['company']}\n{data.get('domain', 'Assessment').title()}", font_size=Pt(44), font_color=WHITE, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(3.8), Inches(8), Inches(0.8), f"Powered by Strategy Skill Pipe", font_size=Pt(16), font_color=GREY_LIGHT)
    _add_footer(slide, data, 1, total_slides)

def build_findings_slide(prs, data, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.35), Inches(8), Pt(16), f"DOMAIN: {data.get('domain', '').upper()}", font_size=Pt(8.5), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.9), data.get('headline', ''), font_size=Pt(22), font_color=WHITE, bold=True)
    _add_accent_line(slide, Inches(0.6), Inches(1.65), Inches(1.5))
    _add_text_box(slide, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.2), data.get('executive_read', ''), font_size=Pt(11), font_color=GREY_LIGHT)
    _add_text_box(slide, Inches(0.6), Inches(3.2), Inches(3), Pt(18), "KEY FINDINGS", font_size=Pt(8.5), font_color=ACCENT_LIGHT, bold=True)
    y_start = Inches(3.55)
    for i, finding in enumerate(data.get('key_findings', [])[:4]):
        _add_finding_bullet(slide, Inches(0.6), y_start + Inches(i * 0.55), Inches(11.5), i + 1, finding)
    _add_footer(slide, data, 2, total_slides)

def build_generic_table_slide(prs, data, slide_data, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.35), Inches(8), Pt(16), f"{data['company'].upper()} OUTPUT", font_size=Pt(8.5), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.6), slide_data.get('title', 'Data Matrix'), font_size=Pt(22), font_color=WHITE, bold=True)
    _add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(1.5))
    
    rows = slide_data.get('rows', [])
    headers = slide_data.get('headers', [])
    row_count = len(rows) + (1 if headers else 0)
    col_count = len(headers) if headers else (len(rows[0]) if rows else 1)
    
    if row_count > 0:
        table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.5))
        table = table_shape.table
        
        current_row = 0
        if headers:
            for c_idx, text in enumerate(headers):
                cell = table.cell(0, c_idx)
                cell.text = str(text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = BG_DARK
                    p.font.bold = True
            current_row += 1
            
        for r_idx, row in enumerate(rows[:6]):
            bg_color = BG_TABLE_ROW if r_idx % 2 == 0 else BG_TABLE_ALT
            for c_idx, text in enumerate(row):
                if c_idx < col_count:
                    cell = table.cell(current_row, c_idx)
                    cell.text = str(text)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg_color
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(11)
                        p.font.color.rgb = WHITE
            current_row += 1
    
    _add_footer(slide, data, slide_num, total_slides)

def build_generic_split_slide(prs, data, slide_data, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.35), Inches(8), Pt(16), f"{data['company'].upper()} OUTPUT", font_size=Pt(8.5), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.6), slide_data.get('title', 'Strategic Analysis'), font_size=Pt(22), font_color=WHITE, bold=True)
    _add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(1.5))
    
    # Left Column
    _add_text_box(slide, Inches(0.6), Inches(1.7), Inches(5.5), Pt(18), slide_data.get('left_title', '').upper(), font_size=Pt(12), font_color=ACCENT_LIGHT, bold=True)
    y_start = Inches(2.1)
    for i, item in enumerate(slide_data.get('left_bullets', [])[:4]):
        _add_finding_bullet(slide, Inches(0.6), y_start + Inches(i * 0.7), Inches(5.5), i + 1, item)
        
    # Right Column
    _add_text_box(slide, Inches(6.5), Inches(1.7), Inches(5.5), Pt(18), slide_data.get('right_title', '').upper(), font_size=Pt(12), font_color=ACCENT_LIGHT, bold=True)
    for i, item in enumerate(slide_data.get('right_bullets', [])[:4]):
        _add_finding_bullet(slide, Inches(6.5), y_start + Inches(i * 0.7), Inches(5.5), i + 1, item)
        
    _add_footer(slide, data, slide_num, total_slides)

def build_generic_quotes_grid(prs, data, slide_data, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.35), Inches(8), Pt(16), f"{data['company'].upper()} OSINT", font_size=Pt(8.5), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.6), slide_data.get('title', 'Voice of Customer'), font_size=Pt(22), font_color=WHITE, bold=True)
    _add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(1.5))
    
    quotes = slide_data.get('quotes', [])
    
    for i, quote in enumerate(quotes[:6]):
        col = i % 2
        row = i // 2
        left = Inches(0.6) if col == 0 else Inches(6.8)
        top = Inches(1.6 + (row * 1.6))
        
        _add_shape(slide, left, top, Inches(5.9), Inches(1.4), BG_TABLE_ROW)
        
        text = f"\"{quote}\""
        _add_text_box(slide, left + Pt(10), top + Pt(10), Inches(5.6), Inches(1.2), text, font_size=Pt(11), font_color=WHITE, alignment=PP_ALIGN.LEFT)
        
    _add_footer(slide, data, slide_num, total_slides)

def build_generic_bullets_slide(prs, data, slide_data, slide_num, total_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.35), Inches(8), Pt(16), f"{data['company'].upper()} STRATEGY", font_size=Pt(8.5), font_color=ACCENT, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.6), slide_data.get('title', 'Analysis'), font_size=Pt(22), font_color=WHITE, bold=True)
    _add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(1.5))
    
    y_start = Inches(1.7)
    for i, item in enumerate(slide_data.get('bullets', [])[:6]):
        _add_finding_bullet(slide, Inches(0.6), y_start + Inches(i * 0.7), Inches(11), i + 1, item)
        
    _add_footer(slide, data, slide_num, total_slides)

def compile_deck(json_path, template_path, output_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    prs = Presentation(template_path)
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.rId
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)
    
    dynamic_slides = data.get('slides', [])
    total_slides = 2 + len(dynamic_slides)
    
    build_title_slide(prs, data, total_slides)
    build_findings_slide(prs, data, total_slides)
    
    current_slide = 3
    for s_data in dynamic_slides:
        if s_data.get('type') == 'table':
            build_generic_table_slide(prs, data, s_data, current_slide, total_slides)
        elif s_data.get('type') == 'split':
            build_generic_split_slide(prs, data, s_data, current_slide, total_slides)
        elif s_data.get('type') == 'quotes_grid':
            build_generic_quotes_grid(prs, data, s_data, current_slide, total_slides)
        elif s_data.get('type') == 'bullets':
            build_generic_bullets_slide(prs, data, s_data, current_slide, total_slides)
        current_slide += 1
    
    prs.save(output_path)
    print(f"Generated {output_path} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python compile.py <json_path> <template_pptx> <output_pptx>")
        sys.exit(1)
    compile_deck(sys.argv[1], sys.argv[2], sys.argv[3])
