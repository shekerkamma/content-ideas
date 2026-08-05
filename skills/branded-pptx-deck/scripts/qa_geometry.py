#!/usr/bin/env python3
"""Mechanical PPTX layout QA without a renderer.

`preview_pptx.py` needs matplotlib, and rendering via LibreOffice needs WSL
interop — neither is guaranteed. This checks what can be checked from the file's
own geometry: text that cannot fit its box, overlapping text frames, off-slide
shapes, and missing footers.

It is NOT a substitute for looking at contact sheets. It catches the mechanical
failures (overflow, collision) but cannot judge visual balance, so a deck that
passes here is still `draft` until someone sees it rendered.

Usage:
    python3 qa_geometry.py DECK.pptx [--verbose]
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_PER_IN = 914400
PT_PER_IN = 72.0

# Calibri-ish metrics. Deliberately slightly generous so we flag real problems
# rather than drowning the report in borderline cases.
AVG_CHAR_W = 0.47      # of font size, in points
LINE_H = 1.22          # of font size, times line spacing
BOLD_FACTOR = 1.04


def emu_to_pt(v):
    return (v / EMU_PER_IN) * PT_PER_IN


def estimate_text_height(tf, box_w_pt):
    """Estimate rendered height in points for a text frame at a given width."""
    total = 0.0
    for para in tf.paragraphs:
        text = "".join(r.text or "" for r in para.runs)
        sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
        size = max(sizes) if sizes else 14.0
        bold = any(r.font.bold for r in para.runs)
        cw = size * AVG_CHAR_W * (BOLD_FACTOR if bold else 1.0)
        usable = max(box_w_pt, 1.0)
        chars_per_line = max(int(usable / cw), 1)
        lines = max(1, -(-len(text) // chars_per_line)) if text else 1
        ls = 1.0
        try:
            if para.line_spacing and isinstance(para.line_spacing, float):
                ls = para.line_spacing
        except Exception:
            pass
        space_before = 0.0
        try:
            if para.space_before is not None:
                space_before = para.space_before.pt
        except Exception:
            pass
        total += lines * size * LINE_H * ls + space_before
    return total


def rects_overlap(a, b, tol_pt=2.0):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ox = min(ax2, bx2) - max(ax1, bx1)
    oy = min(ay2, by2) - max(ay1, by1)
    return ox > tol_pt and oy > tol_pt, max(ox, 0), max(oy, 0)


def check(path, verbose=False):
    prs = Presentation(str(path))
    slide_w = emu_to_pt(prs.slide_width)
    slide_h = emu_to_pt(prs.slide_height)
    problems = []
    stats = {"slides": 0, "textboxes": 0, "charts": 0}

    for idx, slide in enumerate(prs.slides, 1):
        stats["slides"] += 1
        text_boxes = []
        charts = []
        opaque = []
        has_footer = False

        for shape in slide.shapes:
            geom = (emu_to_pt(shape.left), emu_to_pt(shape.top),
                    emu_to_pt(shape.left) + emu_to_pt(shape.width),
                    emu_to_pt(shape.top) + emu_to_pt(shape.height))
            if shape.has_chart:
                stats["charts"] += 1
                charts.append((geom, "chart"))
                continue
            # An opaque filled shape drawn over a chart hides data. This is the
            # bug class that pure text-vs-text checking cannot see: KPI cards
            # and panels are rectangles, not text frames.
            if not shape.has_text_frame:
                try:
                    if shape.fill.type == 1:
                        opaque.append((geom, "filled shape"))
                except Exception:
                    pass
                continue
            try:
                if shape.fill.type == 1:
                    opaque.append((geom, "filled shape"))
            except Exception:
                pass
            tf = shape.text_frame
            raw = tf.text.strip()
            if not raw:
                continue
            stats["textboxes"] += 1

            left, top = emu_to_pt(shape.left), emu_to_pt(shape.top)
            w, h = emu_to_pt(shape.width), emu_to_pt(shape.height)

            if "/" in raw and len(raw) <= 8:
                has_footer = True

            # off-slide
            if left < -2 or top < -2 or left + w > slide_w + 2 or top + h > slide_h + 2:
                problems.append(
                    f"slide {idx}: shape extends off-canvas — "
                    f"box({left:.0f},{top:.0f},{w:.0f}x{h:.0f}) vs slide {slide_w:.0f}x{slide_h:.0f}"
                    f"  [{raw[:40]!r}]"
                )

            # overflow — skip boxes set to shrink-to-fit, which self-correct
            autofit = None
            try:
                autofit = tf.auto_size
            except Exception:
                pass
            needed = estimate_text_height(tf, w)
            if autofit is None or int(getattr(autofit, "value", -1)) != 2:
                if needed > h * 1.18:
                    problems.append(
                        f"slide {idx}: text likely OVERFLOWS its box — needs ~{needed:.0f}pt "
                        f"in {h:.0f}pt  [{raw[:46]!r}]"
                    )

            text_boxes.append(((left, top, left + w, top + h), raw, needed, h))

        # collisions between text frames that actually hold text
        for i in range(len(text_boxes)):
            for j in range(i + 1, len(text_boxes)):
                (ra, ta, na, ha) = text_boxes[i]
                (rb, tb, nb, hb) = text_boxes[j]
                hit, ox, oy = rects_overlap(ra, rb)
                # 3pt, not 6: a footnote sitting 5.8pt inside a table's last
                # row passed at the looser threshold and was only caught by eye.
                if hit and ox > 3 and oy > 3:
                    problems.append(
                        f"slide {idx}: text frames overlap by {ox:.0f}x{oy:.0f}pt — "
                        f"{ta[:26]!r} vs {tb[:26]!r}"
                    )

        # opaque shapes covering text that began above them — e.g. a summary
        # panel laid over the last rows of a table. Text drawn *inside* a card
        # is normal and must not be flagged, so only text starting above the
        # shape and extending into it counts.
        for ogeom, kind in opaque:
            for tgeom, ttxt, _, _ in text_boxes:
                if tgeom[1] >= ogeom[1] - 1:
                    continue
                hit, ox, oy = rects_overlap(tgeom, ogeom)
                if hit and ox > 20 and oy > 4:
                    problems.append(
                        f"slide {idx}: {kind} covers text that starts above it — "
                        f"{oy:.0f}pt into {ttxt[:30]!r}"
                    )

        # opaque shapes covering a chart — the failure text-only checks miss
        for cgeom, _ in charts:
            for ogeom, kind in opaque:
                hit, ox, oy = rects_overlap(cgeom, ogeom)
                if hit and ox > 10 and oy > 10:
                    cw_ = cgeom[2] - cgeom[0]
                    chh = cgeom[3] - cgeom[1]
                    frac = (ox * oy) / max(cw_ * chh, 1) * 100
                    problems.append(
                        f"slide {idx}: {kind} COVERS chart — {ox:.0f}x{oy:.0f}pt "
                        f"({frac:.0f}% of the chart area is obscured)"
                    )

        if not has_footer:
            problems.append(f"slide {idx}: no page-number footer found")

        if verbose:
            print(f"  slide {idx}: {len(text_boxes)} text frames")

    return problems, stats


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck")
    ap.add_argument("--verbose", action="store_true")
    args = ap.parse_args()

    problems, stats = check(Path(args.deck), args.verbose)

    print(f"\nGEOMETRY QA — {Path(args.deck).name}")
    print(f"  slides {stats['slides']}  text frames {stats['textboxes']}  "
          f"native charts {stats['charts']}")
    if not problems:
        print("\n  no overflow, collision, off-canvas or missing-footer issues detected")
        print("\n  NOTE: mechanical checks only. Visual balance is unverified — the deck "
              "stays `draft` until contact sheets are reviewed.")
        return 0
    print(f"\n  {len(problems)} issue(s):")
    for p in problems:
        print(f"   - {p}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
