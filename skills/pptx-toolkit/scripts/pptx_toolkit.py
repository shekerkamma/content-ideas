#!/usr/bin/env python3
"""Inspect and make controlled edits to existing PPTX files."""

from __future__ import annotations

import argparse
import json
import shutil
import sys
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any, Iterable
from xml.etree import ElementTree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


REQUIRED_MEMBERS = {"[Content_Types].xml", "ppt/presentation.xml"}


class ToolkitError(ValueError):
    """Raised for invalid user input or an unsafe PPTX operation."""


def _enum_name(value: Any) -> str | None:
    return getattr(value, "name", None) if value is not None else None


def _inches(value: Any) -> float:
    return round(float(value) / 914400, 4)


def _color(run: Any) -> str | None:
    try:
        rgb = run.font.color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb) if rgb is not None else None


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _run_record(run: Any) -> dict[str, Any]:
    size = run.font.size
    return {
        "text": run.text,
        "font_name": run.font.name,
        "font_size_pt": round(size.pt, 2) if size is not None else None,
        "bold": run.font.bold,
        "italic": run.font.italic,
        "underline": run.font.underline,
        "color_rgb": _color(run),
    }


def _shape_record(shape: Any) -> dict[str, Any]:
    record: dict[str, Any] = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "type": _enum_name(shape.shape_type),
        "left_in": _inches(shape.left),
        "top_in": _inches(shape.top),
        "width_in": _inches(shape.width),
        "height_in": _inches(shape.height),
    }
    if shape.is_placeholder:
        record["placeholder_type"] = _enum_name(shape.placeholder_format.type)
    if shape.has_text_frame:
        record["text"] = shape.text
        record["paragraphs"] = [
            {
                "text": paragraph.text,
                "level": paragraph.level,
                "runs": [_run_record(run) for run in paragraph.runs],
            }
            for paragraph in shape.text_frame.paragraphs
        ]
    return record


def inspect_pptx(path: str | Path) -> dict[str, Any]:
    source = Path(path)
    problems = validate_pptx(source)
    if problems:
        raise ToolkitError("Cannot inspect invalid PPTX: " + "; ".join(problems))

    prs = Presentation(str(source))
    with zipfile.ZipFile(source) as archive:
        comment_parts = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/comments/") and name.endswith(".xml")
        )

    slides = []
    for number, slide in enumerate(prs.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        title = slide.shapes.title.text if slide.shapes.title is not None else ""
        slides.append(
            {
                "slide": number,
                "title": title,
                "layout": slide.slide_layout.name,
                "notes": notes,
                "shapes": [_shape_record(shape) for shape in _iter_shapes(slide.shapes)],
            }
        )

    return {
        "file": str(source.resolve()),
        "slide_count": len(prs.slides),
        "width_in": _inches(prs.slide_width),
        "height_in": _inches(prs.slide_height),
        "comment_parts": comment_parts,
        "slides": slides,
    }


def _find_shape(slide: Any, operation: dict[str, Any]) -> Any:
    selectors = [key for key in ("shape_id", "shape_name") if key in operation]
    if len(selectors) != 1:
        raise ToolkitError("Each replacement requires exactly one of shape_id or shape_name")

    key = selectors[0]
    wanted = operation[key]
    matches = [
        shape
        for shape in _iter_shapes(slide.shapes)
        if (shape.shape_id if key == "shape_id" else shape.name) == wanted
    ]
    if not matches:
        raise ToolkitError(f"No shape matched {key}={wanted!r}")
    if len(matches) > 1:
        raise ToolkitError(f"Multiple shapes matched {key}={wanted!r}; use shape_id")
    return matches[0]


def _replace_shape_text(shape: Any, text: str) -> None:
    if not shape.has_text_frame:
        raise ToolkitError(f"Shape {shape.shape_id} ({shape.name}) has no text frame")

    paragraphs = shape.text_frame.paragraphs
    first_run = next((run for paragraph in paragraphs for run in paragraph.runs), None)
    if first_run is None:
        first_run = paragraphs[0].add_run()
    first_run.text = text
    first_run_element = first_run._r

    found_first = False
    for paragraph in paragraphs:
        for run in paragraph.runs:
            if run._r is first_run_element and not found_first:
                found_first = True
                continue
            run.text = ""


def edit_pptx(
    source: str | Path,
    output: str | Path,
    spec: dict[str, Any],
    *,
    overwrite: bool = False,
) -> dict[str, Any]:
    source_path = Path(source)
    output_path = _prepare_output(source_path, Path(output), overwrite)
    prs = Presentation(str(source_path))
    slide_specs = spec.get("slides")
    if not isinstance(slide_specs, list) or not slide_specs:
        raise ToolkitError("Edit spec requires a non-empty slides list")

    changed_shapes = 0
    changed_notes = 0
    seen_slides: set[int] = set()
    for slide_spec in slide_specs:
        if not isinstance(slide_spec, dict):
            raise ToolkitError("Each slides entry must be an object")
        number = slide_spec.get("slide")
        if not isinstance(number, int) or not 1 <= number <= len(prs.slides):
            raise ToolkitError(f"Invalid one-based slide number: {number!r}")
        if number in seen_slides:
            raise ToolkitError(f"Slide {number} appears more than once in the edit spec")
        seen_slides.add(number)
        slide = prs.slides[number - 1]

        replacements = slide_spec.get("replacements", [])
        if not isinstance(replacements, list):
            raise ToolkitError(f"Slide {number} replacements must be a list")
        for replacement in replacements:
            if not isinstance(replacement, dict) or not isinstance(replacement.get("text"), str):
                raise ToolkitError(f"Slide {number} replacement requires string text")
            shape = _find_shape(slide, replacement)
            _replace_shape_text(shape, replacement["text"])
            changed_shapes += 1

        if "notes" in slide_spec:
            notes = slide_spec["notes"]
            if not isinstance(notes, str):
                raise ToolkitError(f"Slide {number} notes must be a string")
            slide.notes_slide.notes_text_frame.text = notes
            changed_notes += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    _raise_if_invalid(output_path)
    return {
        "output": str(output_path),
        "slides_touched": len(seen_slides),
        "shapes_replaced": changed_shapes,
        "notes_updated": changed_notes,
    }


def reorder_pptx(
    source: str | Path,
    output: str | Path,
    order: list[int],
    *,
    overwrite: bool = False,
) -> dict[str, Any]:
    source_path = Path(source)
    output_path = _prepare_output(source_path, Path(output), overwrite)
    prs = Presentation(str(source_path))
    slide_count = len(prs.slides)
    if not order:
        raise ToolkitError("Order must contain at least one slide")
    if len(order) != len(set(order)):
        raise ToolkitError("Slide duplication is unsupported; order values must be unique")
    invalid = [number for number in order if number < 1 or number > slide_count]
    if invalid:
        raise ToolkitError(f"Slide numbers out of range 1..{slide_count}: {invalid}")

    slide_ids = list(prs.slides._sldIdLst)
    selected = [slide_ids[number - 1] for number in order]
    omitted = [slide_id for index, slide_id in enumerate(slide_ids, start=1) if index not in order]
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_id in selected:
        slide_id_list.append(slide_id)
    for slide_id in omitted:
        prs.part.drop_rel(slide_id.rId)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    _raise_if_invalid(output_path)
    return {
        "output": str(output_path),
        "slide_count": len(order),
        "removed": slide_count - len(order),
        "order": order,
    }


def validate_pptx(path: str | Path) -> list[str]:
    source = Path(path)
    problems: list[str] = []
    if not source.is_file():
        return [f"file not found: {source}"]
    if source.suffix.lower() != ".pptx":
        problems.append("file extension must be .pptx")
    if not zipfile.is_zipfile(source):
        return problems + ["file is not a ZIP-based Office package"]

    try:
        with zipfile.ZipFile(source) as archive:
            names = set(archive.namelist())
            missing = sorted(REQUIRED_MEMBERS - names)
            if missing:
                problems.append("missing required members: " + ", ".join(missing))
            corrupt = archive.testzip()
            if corrupt:
                problems.append(f"corrupt ZIP member: {corrupt}")
            for name in names:
                if name.endswith((".xml", ".rels")):
                    try:
                        ElementTree.fromstring(archive.read(name))
                    except ElementTree.ParseError as exc:
                        problems.append(f"invalid XML in {name}: {exc}")
    except (OSError, zipfile.BadZipFile) as exc:
        problems.append(f"cannot read package: {exc}")

    if not problems:
        try:
            Presentation(str(source))
        except Exception as exc:  # python-pptx exposes several package/XML exceptions
            problems.append(f"python-pptx cannot open package: {exc}")
    return problems


def unpack_pptx(source: str | Path, output_dir: str | Path) -> dict[str, Any]:
    source_path = Path(source)
    _raise_if_invalid(source_path)
    target = Path(output_dir)
    if target.exists() and any(target.iterdir()):
        raise ToolkitError(f"Unpack destination must be absent or empty: {target}")
    target.mkdir(parents=True, exist_ok=True)

    extracted = 0
    with zipfile.ZipFile(source_path) as archive:
        for member in archive.infolist():
            member_path = PurePosixPath(member.filename)
            if member_path.is_absolute() or ".." in member_path.parts:
                raise ToolkitError(f"Unsafe archive path: {member.filename}")
            destination = target.joinpath(*member_path.parts)
            if member.is_dir():
                destination.mkdir(parents=True, exist_ok=True)
                continue
            destination.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(member) as source_file, destination.open("wb") as output_file:
                shutil.copyfileobj(source_file, output_file)
            extracted += 1
    return {"output_dir": str(target), "files_extracted": extracted}


def pack_pptx(source_dir: str | Path, output: str | Path, *, overwrite: bool = False) -> dict[str, Any]:
    source = Path(source_dir)
    if not source.is_dir():
        raise ToolkitError(f"Pack source directory not found: {source}")
    missing = sorted(name for name in REQUIRED_MEMBERS if not (source / name).is_file())
    if missing:
        raise ToolkitError("Pack source is missing: " + ", ".join(missing))
    output_path = Path(output)
    if output_path.exists() and not overwrite:
        raise ToolkitError(f"Output already exists: {output_path}; pass --force to replace it")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    files = []
    for path in source.rglob("*"):
        if path.is_symlink():
            raise ToolkitError(f"Pack source contains unsupported symlink: {path}")
        if path.is_file():
            files.append(path)
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path in sorted(files):
            archive.write(path, path.relative_to(source).as_posix())
    _raise_if_invalid(output_path)
    return {"output": str(output_path), "files_packed": len(files)}


def _prepare_output(source: Path, output: Path, overwrite: bool) -> Path:
    if source.resolve() == output.resolve():
        raise ToolkitError("Input and output must be different files")
    _raise_if_invalid(source)
    if output.exists() and not overwrite:
        raise ToolkitError(f"Output already exists: {output}; pass --force to replace it")
    return output


def _raise_if_invalid(path: Path) -> None:
    problems = validate_pptx(path)
    if problems:
        raise ToolkitError("Invalid PPTX: " + "; ".join(problems))


def _write_json(payload: dict[str, Any], output: str | None) -> None:
    text = json.dumps(payload, indent=2, ensure_ascii=False) + "\n"
    if output:
        path = Path(output)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text, encoding="utf-8")
    else:
        sys.stdout.write(text)


def _parse_order(value: str) -> list[int]:
    try:
        return [int(item.strip()) for item in value.split(",") if item.strip()]
    except ValueError as exc:
        raise argparse.ArgumentTypeError("order must be comma-separated slide numbers") from exc


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)

    inspect_parser = subparsers.add_parser("inspect", help="Write a PPTX inventory")
    inspect_parser.add_argument("source")
    inspect_parser.add_argument("--out")

    edit_parser = subparsers.add_parser("edit", help="Replace shape text or native notes")
    edit_parser.add_argument("source")
    edit_parser.add_argument("output")
    edit_parser.add_argument("--spec", required=True)
    edit_parser.add_argument("--force", action="store_true")

    reorder_parser = subparsers.add_parser("reorder", help="Reorder or remove slides")
    reorder_parser.add_argument("source")
    reorder_parser.add_argument("output")
    reorder_parser.add_argument("--order", required=True, type=_parse_order)
    reorder_parser.add_argument("--force", action="store_true")

    validate_parser = subparsers.add_parser("validate", help="Validate PPTX package structure")
    validate_parser.add_argument("source")

    unpack_parser = subparsers.add_parser("unpack", help="Safely unpack PPTX OOXML")
    unpack_parser.add_argument("source")
    unpack_parser.add_argument("output_dir")

    pack_parser = subparsers.add_parser("pack", help="Pack an OOXML directory as PPTX")
    pack_parser.add_argument("source_dir")
    pack_parser.add_argument("output")
    pack_parser.add_argument("--force", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        if args.command == "inspect":
            _write_json(inspect_pptx(args.source), args.out)
        elif args.command == "edit":
            spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
            _write_json(edit_pptx(args.source, args.output, spec, overwrite=args.force), None)
        elif args.command == "reorder":
            _write_json(reorder_pptx(args.source, args.output, args.order, overwrite=args.force), None)
        elif args.command == "validate":
            problems = validate_pptx(args.source)
            _write_json({"valid": not problems, "problems": problems}, None)
            return 1 if problems else 0
        elif args.command == "unpack":
            _write_json(unpack_pptx(args.source, args.output_dir), None)
        elif args.command == "pack":
            _write_json(pack_pptx(args.source_dir, args.output, overwrite=args.force), None)
    except (ToolkitError, json.JSONDecodeError, OSError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
