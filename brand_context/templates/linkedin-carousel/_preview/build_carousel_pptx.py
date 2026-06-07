"""
Build a 7-slide LinkedIn carousel PPTX with real stock photo backgrounds.
Slide size: 1080x1350 (portrait). Colors/fonts from tokens.json.
Background images from the stock library (Pexels) or ./bg/ fallback.
"""

import json
import random
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Paths ────────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).resolve().parent
TOKENS_PATH = SCRIPT_DIR.parent.parent.parent / "visual-identity" / "tokens.json"
BG_DIR      = SCRIPT_DIR / "bg"
STOCK_DIR   = SCRIPT_DIR.parent.parent.parent / "visual-identity" / "stock-library" / "themes"

def pick_stock_bg(theme, index=0):
    """Pick a stock photo from the library, falling back to ./bg/ generated images."""
    theme_dir = STOCK_DIR / theme
    if theme_dir.exists():
        jpgs = sorted(theme_dir.glob("*.jpg"))
        if jpgs:
            return jpgs[index % len(jpgs)]
    # Fallback to generated backgrounds
    fallback_map = {
        "dark-tech": "bg-cover.png",
        "gradients-abstract": "bg-body-01.png",
        "textures": "bg-body-02.png",
        "code-terminal": "bg-body-03.png",
    }
    return BG_DIR / fallback_map.get(theme, "bg-cover.png")

with open(TOKENS_PATH) as f:
    tokens = json.load(f)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# Colors
PRIMARY     = hex_to_rgb(tokens["palette"]["primary"]["hex"])
SECONDARY   = hex_to_rgb(tokens["palette"]["secondary"]["hex"])
ACCENT      = hex_to_rgb(tokens["palette"]["accent"]["hex"])
TEXT_ON_DARK = hex_to_rgb(tokens["palette"]["text_on_dark"]["hex"])
MUTED       = hex_to_rgb(tokens["palette"]["muted"]["hex"])
SUCCESS     = hex_to_rgb(tokens["palette"]["success"]["hex"])

# Fonts
FONT_DISPLAY = tokens["typography"]["display"]["family"]
FONT_BODY    = tokens["typography"]["body"]["family"]
FONT_MONO    = tokens["typography"]["mono"]["family"]

# ── Presentation setup ───────────────────────────────────────
SLIDE_W = Emu(int(1080 / 96 * 914400))
SLIDE_H = Emu(int(1350 / 96 * 914400))

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK_LAYOUT = prs.slide_layouts[6]

MARGIN   = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ── Helpers ──────────────────────────────────────────────────
def add_bg_image(slide, image_path):
    """Add a full-bleed background image to the slide."""
    slide.shapes.add_picture(str(image_path), 0, 0, SLIDE_W, SLIDE_H)

def add_semi_overlay(slide, opacity=0.55):
    """Add a semi-transparent dark overlay for text readability."""
    from lxml import etree
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    # Set transparency via XML — navigate the spPr element
    spPr = shape._element.find(qn("p:spPr"))
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(opacity * 100000)))
    shape.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                color=TEXT_ON_DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txbox

def add_multiline(slide, left, top, width, height, lines, font_name, font_size,
                  color=TEXT_ON_DARK, bold=False, line_spacing=1.3):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txbox

def add_accent_bar(slide, left, top, width=Inches(0.6), height=Pt(5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_masthead(slide):
    masthead = tokens.get("chrome", {}).get("masthead", {})
    labels = masthead.get("labels", ["JUNE 2026", "{{handle}}", "{{tagline}}"])
    labels = [l.replace("{{month_year}}", "JUNE 2026") for l in labels]
    y = Pt(28)
    positions = [MARGIN, SLIDE_W // 2 - Inches(1), SLIDE_W - MARGIN - Inches(2.5)]
    for label, x in zip(labels, positions):
        add_textbox(slide, x, y, Inches(2.5), Pt(20), label,
                    FONT_BODY, 10, MUTED, bold=True)

def add_pagination(slide, active_idx):
    dot_size = Pt(10)
    gap = Pt(14)
    total_w = 7 * dot_size + 6 * gap
    start_x = (SLIDE_W - total_w) // 2
    y = SLIDE_H - Pt(48)
    for i in range(7):
        x = start_x + i * (dot_size + gap)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dot_size, dot_size)
        dot.fill.solid()
        if i == active_idx:
            dot.fill.fore_color.rgb = ACCENT
        else:
            dot.fill.fore_color.rgb = MUTED
        dot.line.fill.background()

def add_slide_number(slide, num_text, y_pos=Inches(1.1)):
    """Large watermark-style numeral."""
    add_textbox(slide, MARGIN, y_pos, CONTENT_W, Inches(1.5), num_text,
                FONT_DISPLAY, 110, SECONDARY, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s1, pick_stock_bg("dark-tech", 5))  # circuit board
add_semi_overlay(s1, 0.65)
add_masthead(s1)

y = Inches(3.2)
add_accent_bar(s1, MARGIN, y, width=Inches(0.8), height=Pt(5))
y += Pt(36)
add_textbox(s1, MARGIN, y, CONTENT_W, Inches(3),
            "Stop\nOverengineering\nYour Automation",
            FONT_DISPLAY, 56, TEXT_ON_DARK, bold=True)
y += Inches(3.0)
add_textbox(s1, MARGIN, y, CONTENT_W, Inches(0.8),
            "The best solution has the fewest\nmoving parts. Here's why.",
            FONT_BODY, 20, MUTED)

# CTA
y += Inches(1.4)
cta_circle = s1.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y, Pt(42), Pt(42))
cta_circle.fill.solid()
cta_circle.fill.fore_color.rgb = ACCENT
cta_circle.line.fill.background()
add_textbox(s1, MARGIN + Pt(6), y + Pt(6), Pt(42), Pt(30), "\u2192",
            FONT_DISPLAY, 22, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s1, MARGIN + Pt(52), y + Pt(8), Inches(2), Pt(30),
            "SWIPE TO LEARN", FONT_BODY, 13, TEXT_ON_DARK, bold=True)
add_pagination(s1, 0)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s2, pick_stock_bg("gradients-abstract", 0))
add_semi_overlay(s2, 0.70)
add_masthead(s2)
add_slide_number(s2, "01")

y = Inches(2.8)
add_accent_bar(s2, MARGIN, y)
y += Pt(32)
add_textbox(s2, MARGIN, y, CONTENT_W, Inches(0.7), "The Problem",
            FONT_DISPLAY, 38, TEXT_ON_DARK, bold=True)
y += Inches(1.0)

# Content card
card_h = Inches(4.5)
add_rounded_rect(s2, MARGIN, y, CONTENT_W, card_h, SECONDARY)
add_multiline(s2, MARGIN + Pt(28), y + Pt(28), CONTENT_W - Pt(56), card_h - Pt(56), [
    'Everyone wants the "scalable,',
    'enterprise-ready, future-proof" version.',
    '',
    'So they spend 3 months building',
    'infrastructure instead of solving',
    'the actual problem.',
    '',
    'Meanwhile, a 20-line script',
    "would've been done in 20 minutes.",
], FONT_BODY, 21, MUTED, line_spacing=1.5)

# Highlight "20 minutes" with accent
add_textbox(s2, MARGIN + Pt(28), y + card_h - Pt(100), CONTENT_W - Pt(56), Pt(40),
            "20 lines. 20 minutes. Done.",
            FONT_DISPLAY, 22, ACCENT, bold=True)

add_pagination(s2, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: THE REAL QUESTION
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s3, pick_stock_bg("textures", 1))
add_semi_overlay(s3, 0.72)
add_masthead(s3)
add_slide_number(s3, "02")

y = Inches(2.8)
add_accent_bar(s3, MARGIN, y)
y += Pt(32)
add_textbox(s3, MARGIN, y, CONTENT_W, Inches(0.7), "The Real Question",
            FONT_DISPLAY, 38, TEXT_ON_DARK, bold=True)
y += Inches(1.0)

add_multiline(s3, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Here's the thing: before you pick",
    "a framework, ask one question.",
], FONT_BODY, 21, MUTED, line_spacing=1.5)

# Big question in accent card
y += Inches(2.0)
q_card_h = Inches(1.8)
add_rounded_rect(s3, MARGIN, y, CONTENT_W, q_card_h, SECONDARY)
# Accent left border
add_accent_bar(s3, MARGIN, y, width=Pt(4), height=q_card_h)
add_textbox(s3, MARGIN + Pt(32), y + Pt(24), CONTENT_W - Pt(64), q_card_h - Pt(48),
            '"Does a simple script\nsolve this?"',
            FONT_DISPLAY, 32, ACCENT, bold=True)

y += q_card_h + Pt(40)
add_textbox(s3, MARGIN, y, CONTENT_W, Inches(0.8),
            "If yes, you're done. Ship it.",
            FONT_DISPLAY, 24, TEXT_ON_DARK, bold=True)

add_pagination(s3, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: WHAT SIMPLE LOOKS LIKE
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s4, pick_stock_bg("code-terminal", 6))  # code editor screenshot
add_semi_overlay(s4, 0.68)
add_masthead(s4)
add_slide_number(s4, "03")

y = Inches(2.8)
add_accent_bar(s4, MARGIN, y)
y += Pt(32)
add_textbox(s4, MARGIN, y, CONTENT_W, Inches(0.7), "What Simple Looks Like",
            FONT_DISPLAY, 38, TEXT_ON_DARK, bold=True)
y += Inches(1.0)

# Terminal window
term_h = Inches(4.2)
# Terminal chrome bar
chrome_h = Pt(36)
add_rounded_rect(s4, MARGIN, y, CONTENT_W, chrome_h, hex_to_rgb("#162032"))
# Traffic light dots
dot_y = y + Pt(12)
for i, c in enumerate(["#EF4444", "#F59E0B", "#10B981"]):
    dx = MARGIN + Pt(16) + i * Pt(18)
    d = s4.shapes.add_shape(MSO_SHAPE.OVAL, dx, dot_y, Pt(10), Pt(10))
    d.fill.solid()
    d.fill.fore_color.rgb = hex_to_rgb(c)
    d.line.fill.background()

# Terminal body
tbody_y = y + chrome_h
tbody_h = term_h - chrome_h
add_rounded_rect(s4, MARGIN, tbody_y, CONTENT_W, tbody_h, SECONDARY)
# Accent left border
add_accent_bar(s4, MARGIN, tbody_y, width=Pt(3), height=tbody_h)

code_lines = [
    "import requests",
    "from bs4 import BeautifulSoup",
    "",
    'page = requests.get(',
    '    "https://example.com/prices")',
    "prices = [t.text for t in",
    '    BeautifulSoup(page.text)',
    '    .select(".price")]',
]
add_multiline(s4, MARGIN + Pt(24), tbody_y + Pt(20),
              CONTENT_W - Pt(48), tbody_h - Pt(40),
              code_lines, FONT_MONO, 17, TEXT_ON_DARK, line_spacing=1.6)

y += term_h + Pt(24)
add_multiline(s4, MARGIN, y, CONTENT_W, Inches(1), [
    "Four lines. Runs anywhere.",
    "No framework required.",
], FONT_BODY, 20, MUTED, line_spacing=1.4)

add_pagination(s4, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: THE FRAMEWORK TRAP
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s5, pick_stock_bg("dark-tech", 8))
add_semi_overlay(s5, 0.72)
add_masthead(s5)
add_slide_number(s5, "04")

y = Inches(2.8)
add_accent_bar(s5, MARGIN, y)
y += Pt(32)
add_textbox(s5, MARGIN, y, CONTENT_W, Inches(0.7), "The Framework Trap",
            FONT_DISPLAY, 38, TEXT_ON_DARK, bold=True)
y += Inches(1.0)

# Two cards: DO vs DON'T
card_w = (CONTENT_W - Pt(20)) // 2
card_h = Inches(4.5)

# DON'T card
add_rounded_rect(s5, MARGIN, y, card_w, card_h, SECONDARY)
add_textbox(s5, MARGIN + Pt(20), y + Pt(16), card_w - Pt(40), Pt(30),
            "THE TRAP", FONT_DISPLAY, 13, hex_to_rgb("#EF4444"), bold=True)
add_multiline(s5, MARGIN + Pt(20), y + Pt(56), card_w - Pt(40), card_h - Pt(72), [
    "Pick a framework first",
    "",
    "Spend weeks on",
    "infrastructure",
    "",
    "Over-architect for",
    "\"future scale\"",
    "",
    "Never actually ship",
], FONT_BODY, 17, MUTED, line_spacing=1.4)

# DO card
do_x = MARGIN + card_w + Pt(20)
add_rounded_rect(s5, do_x, y, card_w, card_h, SECONDARY)
# Accent top border
add_accent_bar(s5, do_x, y, width=card_w, height=Pt(3))
add_textbox(s5, do_x + Pt(20), y + Pt(16), card_w - Pt(40), Pt(30),
            "THE FIX", FONT_DISPLAY, 13, ACCENT, bold=True)
add_multiline(s5, do_x + Pt(20), y + Pt(56), card_w - Pt(40), card_h - Pt(72), [
    "Start with a script",
    "",
    "Solve the actual",
    "problem first",
    "",
    "Add structure only",
    "when it breaks",
    "",
    "Ship today",
], FONT_BODY, 17, TEXT_ON_DARK, bold=False, line_spacing=1.4)

add_pagination(s5, 4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: DEAD SIMPLE TEST
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s6, pick_stock_bg("gradients-abstract", 3))
add_semi_overlay(s6, 0.70)
add_masthead(s6)
add_slide_number(s6, "05")

y = Inches(2.8)
add_accent_bar(s6, MARGIN, y)
y += Pt(32)
add_textbox(s6, MARGIN, y, CONTENT_W, Inches(0.7), "Dead Simple Test",
            FONT_DISPLAY, 38, TEXT_ON_DARK, bold=True)
y += Inches(1.2)

# Three checklist cards
checks = [
    ("01", "Does it run?"),
    ("02", "Does it save someone time?"),
    ("03", "Can you explain it\nin one sentence?"),
]
for i, (num, text) in enumerate(checks):
    cy = y + i * Pt(130)
    # Card
    ch = Pt(110)
    add_rounded_rect(s6, MARGIN, cy, CONTENT_W, ch, SECONDARY)
    # Number circle
    circle = s6.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN + Pt(20), cy + Pt(24), Pt(52), Pt(52))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_textbox(s6, MARGIN + Pt(20), cy + Pt(30), Pt(52), Pt(40), num,
                FONT_DISPLAY, 18, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    # Text
    add_textbox(s6, MARGIN + Pt(88), cy + Pt(28), CONTENT_W - Pt(120), ch - Pt(56), text,
                FONT_BODY, 22, TEXT_ON_DARK)

y += Pt(430)
add_rounded_rect(s6, MARGIN, y, CONTENT_W, Pt(60), ACCENT)
add_textbox(s6, MARGIN, y + Pt(12), CONTENT_W, Pt(40),
            "If yes to all three \u2014 ship it. Stop adding things.",
            FONT_DISPLAY, 18, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_pagination(s6, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: CLOSER
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_image(s7, pick_stock_bg("textures", 5))
add_semi_overlay(s7, 0.65)
add_masthead(s7)

y = Inches(3.0)
add_accent_bar(s7, MARGIN, y, width=Inches(1.0), height=Pt(5))
y += Pt(40)
add_textbox(s7, MARGIN, y, CONTENT_W, Inches(2.5),
            "The best scraper\nis the one that\nactually runs.",
            FONT_DISPLAY, 46, ACCENT, bold=True)
y += Inches(3.0)
add_multiline(s7, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Not the one with the cleanest architecture,",
    "the most tests, or the fanciest dashboard.",
    "",
    "The one that runs. That's it.",
], FONT_BODY, 20, MUTED, line_spacing=1.5)
y += Inches(2.2)

# Branded footer
add_accent_bar(s7, MARGIN, y, width=Inches(6.5), height=Pt(1))
y += Pt(20)
brand_name = tokens.get("brand", "{{brand_name}}")
add_textbox(s7, MARGIN, y, CONTENT_W, Inches(0.5), brand_name,
            FONT_DISPLAY, 24, TEXT_ON_DARK, bold=True)
masthead_labels = tokens.get("chrome", {}).get("masthead", {}).get("labels", [])
tagline = masthead_labels[2] if len(masthead_labels) > 2 else ""
add_textbox(s7, MARGIN, y + Pt(32), CONTENT_W, Inches(0.3),
            tagline, FONT_BODY, 14, MUTED)

add_pagination(s7, 6)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
OUT = SCRIPT_DIR / "carousel-simplicity-beats-sophistication.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT.stat().st_size / 1024:.0f} KB")
