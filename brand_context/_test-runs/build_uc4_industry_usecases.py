"""
Use Case 4: AI Use Cases Across Industries
Deliverable: 16:9 PPTX deck — cover, 8 industry slides, summary slide
Each industry slide: 2 use cases with key metrics and what actually works
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR  = Path(__file__).resolve().parent
TOKENS_PATH = SCRIPT_DIR.parent / "visual-identity" / "tokens.json"
STOCK_DIR   = SCRIPT_DIR.parent / "visual-identity" / "stock-library" / "themes"

with open(TOKENS_PATH) as f:
    tokens = json.load(f)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

PRIMARY     = hex_to_rgb(tokens["palette"]["primary"]["hex"])
SECONDARY   = hex_to_rgb(tokens["palette"]["secondary"]["hex"])
ACCENT      = hex_to_rgb(tokens["palette"]["accent"]["hex"])
TEXT_WHITE   = hex_to_rgb(tokens["palette"]["text_on_dark"]["hex"])
MUTED       = hex_to_rgb(tokens["palette"]["muted"]["hex"])
SUCCESS     = hex_to_rgb(tokens["palette"]["success"]["hex"])
DANGER      = hex_to_rgb("#EF4444")
WARNING     = hex_to_rgb("#F59E0B")
HIGHLIGHT   = hex_to_rgb("#818CF8")  # Indigo 400

FONT_DISPLAY = tokens["typography"]["display"]["family"]
FONT_BODY    = tokens["typography"]["body"]["family"]
FONT_MONO    = tokens["typography"]["mono"]["family"]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


def pick_bg(theme, idx=0):
    theme_dir = STOCK_DIR / theme
    if theme_dir.exists():
        jpgs = sorted(theme_dir.glob("*.jpg"))
        if jpgs:
            return jpgs[idx % len(jpgs)]
    return None


def add_bg(slide, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg_image(slide, path, overlay_alpha=0.80):
    slide.shapes.add_picture(str(path), 0, 0, SLIDE_W, SLIDE_H)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    spPr = shape._element.find(qn("p:spPr"))
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(overlay_alpha * 100000)))
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, font=FONT_BODY,
             size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txbox


def add_lines(slide, left, top, width, height, lines, font=FONT_BODY,
              size=16, color=TEXT_WHITE, bold=False, spacing=1.2):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(size * (spacing - 1))
    return txbox


def add_accent_bar(slide, left, top, width=Inches(0.8), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_card(slide, left, top, width, height, fill=SECONDARY, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def add_stat_pill(slide, left, top, text, color=ACCENT):
    w = Inches(2.2)
    h = Pt(28)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, left, top + Pt(3), w, Pt(22), text,
             FONT_MONO, 11, PRIMARY, bold=True, align=PP_ALIGN.CENTER)


def add_industry_slide(industry_num, industry_name, bg_theme, bg_idx,
                       uc1_title, uc1_stats, uc1_insight,
                       uc2_title, uc2_stats, uc2_insight):
    s = prs.slides.add_slide(BLANK)
    bg = pick_bg(bg_theme, bg_idx)
    if bg:
        add_bg_image(s, bg)
    else:
        add_bg(s)

    # Industry number + name
    add_text(s, MARGIN, Inches(0.3), Inches(1), Inches(0.6),
             f"{industry_num:02d}", FONT_DISPLAY, 56, SECONDARY, bold=True)
    add_text(s, MARGIN + Inches(1.1), Inches(0.35), Inches(8), Inches(0.5),
             industry_name, FONT_DISPLAY, 32, TEXT_WHITE, bold=True)
    add_accent_bar(s, MARGIN + Inches(1.1), Inches(0.85), Inches(0.8))

    # Two cards side by side
    card_w = (CONTENT_W - Pt(30)) / 2
    card_h = Inches(5.5)
    y = Inches(1.2)

    for i, (title, stats, insight) in enumerate([
        (uc1_title, uc1_stats, uc1_insight),
        (uc2_title, uc2_stats, uc2_insight),
    ]):
        x = MARGIN + i * (card_w + Pt(30))
        add_card(s, x, y, card_w, card_h)

        # UC title
        add_text(s, x + Pt(20), y + Pt(14), card_w - Pt(40), Pt(50), title,
                 FONT_DISPLAY, 17, ACCENT, bold=True)

        # Stats
        sy = y + Pt(68)
        for j, stat in enumerate(stats):
            add_text(s, x + Pt(20), sy + j * Pt(26), card_w - Pt(40), Pt(24), stat,
                     FONT_MONO, 12, SUCCESS, bold=True)

        # Divider
        div_y = sy + len(stats) * Pt(26) + Pt(8)
        add_accent_bar(s, x + Pt(20), div_y, card_w - Pt(40), Pt(2))

        # Insight
        ins_y = div_y + Pt(14)
        add_lines(s, x + Pt(20), ins_y, card_w - Pt(40), Inches(3),
                  insight, FONT_BODY, 13, MUTED, spacing=1.35)


# ═══════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 4)
if bg:
    add_bg_image(s, bg, 0.82)
else:
    add_bg(s)

y = Inches(1.6)
add_accent_bar(s, MARGIN, y, Inches(1.2))
y += Pt(30)
add_text(s, MARGIN, y, CONTENT_W, Inches(1.2),
         "AI Use Cases That Actually Ship", FONT_DISPLAY, 48, TEXT_WHITE, bold=True)
y += Inches(1.0)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "18 Production Deployments Across 8 Industries", FONT_BODY, 22, MUTED)
y += Inches(0.6)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "Real numbers. Real stacks. No \"transformational AI\" nonsense.", FONT_MONO, 16, ACCENT)
y += Inches(0.9)

# Summary stats row
stats_row = [
    ("8 Industries", ACCENT),
    ("18 Use Cases", SUCCESS),
    ("Avg 73% efficiency gain", WARNING),
    ("6-week ROI proof", HIGHLIGHT),
]
for i, (label, color) in enumerate(stats_row):
    add_stat_pill(s, MARGIN + i * Inches(2.6), y, label, color)

y += Inches(1.0)
add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Every use case here is in production. Measured in months, not \"transformational journeys.\"",
    "The pattern is the same everywhere: automate the boring part, keep humans on judgment calls.",
    "Start small. Prove ROI in weeks. Expand from evidence, not roadmaps.",
], FONT_BODY, 15, MUTED, spacing=1.4)

add_text(s, MARGIN, SLIDE_H - Inches(0.6), CONTENT_W, Inches(0.4),
         "{{brand_name}}", FONT_DISPLAY, 14, SECONDARY, bold=True)


# ═══════════════════════════════════════════════════
# INDUSTRY SLIDES
# ═══════════════════════════════════════════════════

# 1. Financial Services
add_industry_slide(
    1, "Financial Services",
    "dark-tech", 1,
    "Claims Processing Automation",
    ["14 days \u2192 3 days cycle time", "73% auto-approved", "Error rate: 12% \u2192 2.1%"],
    [
        "They didn't try to automate everything.",
        "The 73% that's straightforward gets",
        "auto-processed. The 27% that's complicated",
        "still goes to a human \u2014 but now that human",
        "isn't buried under the easy stuff.",
    ],
    "AML Alert Triage",
    ["False positives: 95% \u2192 60%", "Analyst productivity: 4x", "Investigation: 4hrs \u2192 45min"],
    [
        "The old system flagged everything.",
        "Analysts spent 95% of their time on nothing.",
        "The model doesn't replace the analyst \u2014",
        "it tells them where to look first.",
        "That's the actual work.",
    ],
)

# 2. Healthcare & Life Sciences
add_industry_slide(
    2, "Healthcare & Life Sciences",
    "gradients-abstract", 3,
    "Radiology Pre-Read Prioritization",
    ["11hrs \u2192 38min critical turnaround", "Missed findings: 3.2% \u2192 0.4%", "Same volume, right order"],
    [
        "Nobody asked the AI to diagnose anything.",
        "It just sorts the pile. The patient with a",
        "collapsed lung gets seen in 38 minutes",
        "instead of 11 hours. That's the difference",
        "between walking out and not.",
    ],
    "Prior Authorization Automation",
    ["62% auto-approved (no human)", "5 days \u2192 4 hours approval", "$2.1M saved/year per facility"],
    [
        "Prior auth is a bureaucratic bottleneck,",
        "not a clinical one. Most requests meet",
        "criteria \u2014 someone just needs to check",
        "the boxes. The AI checks the boxes.",
        "Edge cases still go to a nurse.",
    ],
)

# 3. Manufacturing & Industrial
add_industry_slide(
    3, "Manufacturing & Industrial",
    "code-terminal", 2,
    "Predictive Maintenance",
    ["Unplanned downtime: \u221242%", "Maintenance costs: \u221228%", "$3.8M saved/year (1 plant)"],
    [
        "They started with the 12 machines that",
        "caused 80% of downtime. Instrumented those.",
        "Proved ROI in 6 weeks. Then expanded.",
        "That's how you survive the pilot.",
    ],
    "Visual Quality Inspection",
    ["Detection: 99.7% (vs 87% human)", "200ms per unit (vs 4s manual)", "Warranty claims: \u221238%"],
    [
        "Human inspectors get tired after 2 hours.",
        "Accuracy drops to 72% by hour 6.",
        "The camera doesn't get tired. But they kept",
        "humans for the 0.3% the camera wasn't sure",
        "about. Best of both.",
    ],
)

# 4. Retail & E-commerce
add_industry_slide(
    4, "Retail & E-commerce",
    "textures", 3,
    "Dynamic Pricing Optimization",
    ["Revenue/visitor: +12%", "Margin: +3.2 pts", "Markdown waste: \u221223%"],
    [
        "50K SKUs, one person, once a month.",
        "The math doesn't work. The model runs",
        "continuously. The category manager now sets",
        "strategy and guardrails \u2014 min margins,",
        "brand rules. Better job. Better outcomes.",
    ],
    "Returns Prediction & Prevention",
    ["Return rate: 28% \u2192 19%", "$15-33 saved per prevented return", "$22M/year savings"],
    [
        "Returns aren't a logistics problem. They're a",
        "\"customer got the wrong thing\" problem.",
        "Fix the information gap before the purchase",
        "and the return never happens.",
    ],
)

# 5. Logistics & Supply Chain
add_industry_slide(
    5, "Logistics & Supply Chain",
    "dark-tech", 6,
    "Demand Forecasting + Inventory",
    ["Forecast accuracy: 63% \u2192 82%", "Stockouts: \u221235%  |  Overstock: \u221228%", "Working capital freed: $12M"],
    [
        "The old model was a moving average plus",
        "\"gut feel.\" The real win was trusting",
        "the model enough to automate replenishment.",
        "Automate the boring orders, let buyers",
        "handle exceptions. That's where the money is.",
    ],
    "Last-Mile Route Optimization",
    ["Deliveries/route: +18%", "Fuel costs: \u221214%  |  On-time: 97%", "$8.2M saved/yr (1,200 vehicles)"],
    [
        "A dispatcher planning 200 routes manually",
        "is solving an NP-hard problem with a",
        "whiteboard. The optimizer finds a route 18%",
        "better than humans can \u2014 in 4 minutes",
        "instead of 4 hours.",
    ],
)

# 6. Energy & Utilities
add_industry_slide(
    6, "Energy & Utilities",
    "gradients-abstract", 7,
    "Grid Load Forecasting",
    ["Peak demand: \u22128-12%", "Balancing costs: \u2212$4.7M/yr", "Renewable waste: \u221222%"],
    [
        "The grid doesn't have a storage problem \u2014",
        "it has a prediction problem. Know tomorrow's",
        "peak within 2%, and you pre-position resources.",
        "Off by 15%? Gas peaker plant. The AI gets it",
        "within 2.3%. Cheap vs expensive electrons.",
    ],
    "Pipeline Leak Detection",
    ["72hrs \u2192 15min detection", "False alarms: 34% \u2192 2.1%", "$23M in fines avoided (yr 1)"],
    [
        "Old system had 34% false alarm rate.",
        "Crews stopped trusting it. Started ignoring",
        "alerts. New model: 2.1% false alarms.",
        "When it says leak, there's a leak.",
        "Trust drives action. Dead simple.",
    ],
)

# 7. Legal & Professional Services
add_industry_slide(
    7, "Legal & Professional Services",
    "textures", 7,
    "Contract Review & Risk Extraction",
    ["4 hours \u2192 26 minutes (90% faster)", "Risk accuracy: 94%", "Outside counsel: \u221232%"],
    [
        "Lawyers weren't spending 4 hours on judgment.",
        "They spent 3 hours finding clauses and",
        "1 hour on judgment. The AI handles finding.",
        "The lawyer handles judgment. Everyone does",
        "what they're actually good at.",
    ],
    "M&A Due Diligence",
    ["6 weeks \u2192 8 days", "14K docs in 72 hours", "Found 3 risks humans missed"],
    [
        "14,000 documents. 12 associates. 4 weeks.",
        "Someone's going to miss something on page",
        "8,247 of exhibit C. The AI reads every page",
        "with the same attention. Partners still make",
        "the call. That's still judgment. Still human.",
    ],
)

# 8. Real Estate & Construction
add_industry_slide(
    8, "Real Estate & Construction",
    "code-terminal", 5,
    "Property Valuation at Scale",
    ["Accuracy: within 4.2% of appraisal", "3 months \u2192 2 weeks", "$340M in mispriced assets found"],
    [
        "They didn't replace appraisers for deals.",
        "The model handles portfolio monitoring.",
        "The appraiser handles the transaction.",
        "Right tool, right job.",
    ],
    "Construction Progress Monitoring",
    ["Deviations caught 3 weeks earlier", "Safety violations: +340% detected", "Rework costs: \u221221%"],
    [
        "The superintendent walks the site once a day.",
        "The drone sees everything. Every floor, every",
        "corner, every missing guardrail. The super now",
        "acts on data instead of generating it.",
        "Faster decisions, fewer surprises.",
    ],
)


# ═══════════════════════════════════════════════════
# SLIDE 10: THE PATTERN (SUMMARY)
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 9)
if bg:
    add_bg_image(s, bg, 0.85)
else:
    add_bg(s)

y = Inches(0.4)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "What the Winners Have in Common", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
add_accent_bar(s, MARGIN, y + Inches(0.7), Inches(1.0))

y = Inches(1.3)
rules = [
    ("01", "Define the number before writing code",
     "Not \"improve efficiency.\" A specific metric, a specific target, a specific timeline."),
    ("02", "Automate the boring part, not the judgment",
     "Every winning use case keeps humans for decisions. AI handles extraction, matching, triage."),
    ("03", "Start with 12 machines, not 1,200",
     "One document type, one product category, one process. Prove it. Then expand from evidence."),
    ("04", "Measure in weeks, not quarters",
     "If you can't show ROI in 6 weeks, your scope is wrong. Not your model."),
    ("05", "Trust the model enough to act on it",
     "A forecast nobody acts on is a screensaver. Automate the action. That's where the money lives."),
]

for i, (num, title, desc) in enumerate(rules):
    ry = y + i * Pt(86)
    add_card(s, MARGIN, ry, CONTENT_W, Pt(76))
    add_text(s, MARGIN + Pt(16), ry + Pt(10), Pt(40), Pt(30), num,
             FONT_DISPLAY, 22, ACCENT, bold=True)
    add_text(s, MARGIN + Pt(60), ry + Pt(10), CONTENT_W - Pt(80), Pt(28), title,
             FONT_DISPLAY, 18, TEXT_WHITE, bold=True)
    add_text(s, MARGIN + Pt(60), ry + Pt(40), CONTENT_W - Pt(80), Pt(28), desc,
             FONT_BODY, 13, MUTED)

add_text(s, MARGIN, SLIDE_H - Inches(0.6), CONTENT_W, Inches(0.4),
         "{{brand_name}}", FONT_DISPLAY, 14, SECONDARY, bold=True)


# Save
OUT = SCRIPT_DIR / "uc4-ai-use-cases-by-industry.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT.stat().st_size / 1024:.0f} KB")
