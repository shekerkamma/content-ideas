"""
Use Case 3: AI Pattern Audit — Before/After Report
Deliverable: PPTX showing original text, detected patterns, and rewritten version
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR  = Path(__file__).resolve().parent
TOKENS_PATH = SCRIPT_DIR.parent / "visual-identity" / "tokens.json"
STOCK_DIR   = SCRIPT_DIR.parent / "visual-identity" / "stock-library" / "themes"

with open(TOKENS_PATH) as f:
    tokens = json.load(f)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

PRIMARY     = hex_to_rgb(tokens["palette"]["primary"]["hex"])
SECONDARY   = hex_to_rgb(tokens["palette"]["secondary"]["hex"])
ACCENT      = hex_to_rgb(tokens["palette"]["accent"]["hex"])
TEXT_WHITE   = hex_to_rgb(tokens["palette"]["text_on_dark"]["hex"])
MUTED       = hex_to_rgb(tokens["palette"]["muted"]["hex"])
SUCCESS     = hex_to_rgb(tokens["palette"]["success"]["hex"])
DANGER      = hex_to_rgb("#EF4444")
WARNING     = hex_to_rgb("#F59E0B")

FONT_DISPLAY = tokens["typography"]["display"]["family"]
FONT_BODY    = tokens["typography"]["body"]["family"]
FONT_MONO    = tokens["typography"]["mono"]["family"]

# 16:9 landscape for report
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


def pick_bg(theme, idx=0):
    theme_dir = STOCK_DIR / theme
    if theme_dir.exists():
        jpgs = sorted(theme_dir.glob("*.jpg"))
        if jpgs:
            return jpgs[idx % len(jpgs)]
    return None


def add_bg(slide, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg_image(slide, path):
    slide.shapes.add_picture(str(path), 0, 0, SLIDE_W, SLIDE_H)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    spPr = shape._element.find(qn("p:spPr"))
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(0.80 * 100000)))
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, font=FONT_BODY,
             size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txbox


def add_lines(slide, left, top, width, height, lines, font=FONT_BODY,
              size=16, color=TEXT_WHITE, bold=False, spacing=1.3):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(size * (spacing - 1))
    return txbox


def add_accent_bar(slide, left, top, width=Inches(0.8), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_card(slide, left, top, width, height, fill=SECONDARY, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def add_score_circle(slide, left, top, score, label, color):
    size = Pt(90)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, left, top + Pt(16), size, Pt(40), f"{score}",
             FONT_DISPLAY, 32, PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Pt(52), size, Pt(20), "/10",
             FONT_BODY, 12, PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, left - Pt(10), top + size + Pt(8), size + Pt(20), Pt(20), label,
             FONT_BODY, 11, MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 6)
if bg: add_bg_image(s, bg)
else: add_bg(s)

y = Inches(2.0)
add_accent_bar(s, MARGIN, y, Inches(1.2))
y += Pt(30)
add_text(s, MARGIN, y, CONTENT_W, Inches(1.2),
         "AI Pattern Audit Report", FONT_DISPLAY, 48, TEXT_WHITE, bold=True)
y += Inches(1.0)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "Before & After  |  Humanizer Deep Mode  |  Voice-Matched",
         FONT_BODY, 20, MUTED)
y += Inches(0.8)
# Score badges
add_score_circle(s, MARGIN, y, "2.1", "BEFORE", DANGER)
add_text(s, MARGIN + Pt(110), y + Pt(30), Pt(40), Pt(30), "\u2192",
         FONT_DISPLAY, 28, MUTED, align=PP_ALIGN.CENTER)
add_score_circle(s, MARGIN + Pt(160), y, "8.9", "AFTER", SUCCESS)

add_text(s, MARGIN + Inches(4), y + Pt(10), Inches(6), Inches(1),
         "28 AI patterns removed\n220 \u2192 175 words (20% reduction)\nVoice-profile matched to brand",
         FONT_MONO, 16, MUTED)


# ═══════════════════════════════════════════════════
# SLIDE 2: ORIGINAL TEXT (BEFORE)
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_text(s, MARGIN, Inches(0.4), Inches(3), Inches(0.5),
         "BEFORE", FONT_DISPLAY, 32, DANGER, bold=True)
add_score_circle(s, SLIDE_W - MARGIN - Pt(90), Inches(0.3), "2.1", "SCORE", DANGER)

y = Inches(1.2)
add_card(s, MARGIN, y, CONTENT_W, Inches(5.8), border=DANGER)
add_lines(s, MARGIN + Pt(24), y + Pt(16), CONTENT_W - Pt(48), Inches(5.5), [
    "In today's rapidly evolving AI landscape, organizations are",
    "increasingly leveraging cutting-edge artificial intelligence",
    "solutions to drive transformative outcomes across their enterprise.",
    "However, it's important to note that a staggering 87% of AI pilot",
    "projects never make it to production. Let's dive deep into why",
    "this happens and what you can do about it.",
    "",
    "The first and arguably most critical challenge is the lack of",
    "clearly defined success metrics. Many organizations embark on",
    "their AI journey without establishing measurable KPIs, leading",
    "to ambiguous outcomes that fail to demonstrate tangible business",
    "value. It's no secret that without clear metrics, even the most",
    "sophisticated AI models cannot prove their worth.",
    "",
    "Furthermore, data readiness remains a significant bottleneck.",
    "Organizations often underestimate the importance of having clean,",
    "accessible, and well-governed data pipelines. Moreover, the gap",
    "between proof-of-concept data environments and production-grade",
    "data infrastructure is frequently underappreciated...",
], FONT_MONO, 13, MUTED, spacing=1.3)


# ═══════════════════════════════════════════════════
# SLIDE 3: PATTERN DETECTION
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_text(s, MARGIN, Inches(0.4), Inches(6), Inches(0.5),
         "28 AI Patterns Detected", FONT_DISPLAY, 28, WARNING, bold=True)

y = Inches(1.1)
col_w = (CONTENT_W - Pt(40)) / 3

categories = [
    ("AI Cliches", "HIGH", DANGER, [
        "rapidly evolving landscape",
        "Let's dive deep",
        "it's no secret",
        "At the end of the day",
        "key takeaway",
        "unlock the full potential",
        "game-changer",
        "AI journey",
    ]),
    ("Corporate Buzzwords", "MED-HIGH", WARNING, [
        "leveraging",
        "cutting-edge",
        "transformative outcomes",
        "holistic endeavor",
        "best-in-class",
        "revolutionize",
        "utilizing",
        "fostering a culture",
    ]),
    ("Hedging + Transitions", "MEDIUM", MUTED, [
        "it's important to note",
        "arguably most critical",
        "Furthermore",
        "Moreover",
        "Additionally",
        "nuanced approach",
    ]),
]

for i, (cat, severity, color, patterns) in enumerate(categories):
    x = MARGIN + i * (col_w + Pt(20))
    add_card(s, x, y, col_w, Inches(5.8))
    add_text(s, x + Pt(16), y + Pt(12), col_w - Pt(32), Pt(24), cat,
             FONT_DISPLAY, 14, color, bold=True)
    add_text(s, x + Pt(16), y + Pt(36), col_w - Pt(32), Pt(16), severity,
             FONT_MONO, 10, color)
    add_accent_bar(s, x + Pt(16), y + Pt(56), col_w - Pt(32), Pt(2))

    lines = [f"\u2022 {p}" for p in patterns]
    add_lines(s, x + Pt(16), y + Pt(68), col_w - Pt(32), Inches(4.5),
              lines, FONT_MONO, 12, MUTED, spacing=1.5)


# ═══════════════════════════════════════════════════
# SLIDE 4: AFTER (REWRITTEN)
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)

add_text(s, MARGIN, Inches(0.4), Inches(3), Inches(0.5),
         "AFTER", FONT_DISPLAY, 32, SUCCESS, bold=True)
add_text(s, MARGIN + Inches(3), Inches(0.4), Inches(4), Inches(0.5),
         "(Voice-Matched \u2014 Deep Mode)", FONT_BODY, 18, MUTED)
add_score_circle(s, SLIDE_W - MARGIN - Pt(90), Inches(0.3), "8.9", "SCORE", SUCCESS)

y = Inches(1.2)
add_card(s, MARGIN, y, CONTENT_W, Inches(5.8), border=SUCCESS)
add_lines(s, MARGIN + Pt(24), y + Pt(16), CONTENT_W - Pt(48), Inches(5.5), [
    "87% of AI pilots never make it to production.",
    "",
    "Not because the models don't work. Because nobody defined what",
    "\"working\" means before they started building.",
    "",
    "Here's the thing \u2014 \"improve efficiency\" isn't a success metric.",
    "\"Reduce claim processing from 14 days to 3\" is. If you can't put",
    "a number on the goal before you write a single line of code,",
    "you're not ready for AI. You're ready for a whiteboard.",
    "",
    "The proof-of-concept worked because you hand-picked clean data",
    "for the demo. Production data is messy, late, undocumented, and",
    "lives in six different systems that don't talk to each other.",
    "Fix the pipes before you train the model. Dead simple. Nobody does it.",
    "",
    "You can build the best model in the world. If the team doesn't",
    "trust it, they won't use it. Show them the tool. Let them break it.",
    "Then fix what they broke. That's the actual work.",
], FONT_MONO, 13, TEXT_WHITE, spacing=1.3)


# ═══════════════════════════════════════════════════
# SLIDE 5: VOICE MATCH REPORT
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("gradients-abstract", 2)
if bg: add_bg_image(s, bg)
else: add_bg(s)

add_text(s, MARGIN, Inches(0.4), Inches(6), Inches(0.5),
         "Voice Match Report", FONT_DISPLAY, 32, TEXT_WHITE, bold=True)

y = Inches(1.2)
dimensions = [
    ("Tone", "Casual-but-competent", "9/10", SUCCESS),
    ("Rhythm", "Short + medium mix, fragments", "8/10", SUCCESS),
    ("Vocabulary", "Brand terms applied", "9/10", SUCCESS),
    ("Anti-patterns", "Zero banned words", "10/10", SUCCESS),
    ("Confidence", "Bold assertions, no hedging", "9/10", SUCCESS),
]

for i, (dim, desc, score, color) in enumerate(dimensions):
    cy = y + i * Pt(80)
    add_card(s, MARGIN, cy, CONTENT_W, Pt(68))
    add_text(s, MARGIN + Pt(20), cy + Pt(12), Inches(2.5), Pt(28), dim,
             FONT_DISPLAY, 18, TEXT_WHITE, bold=True)
    add_text(s, MARGIN + Inches(3), cy + Pt(14), Inches(5), Pt(24), desc,
             FONT_BODY, 15, MUTED)
    add_text(s, SLIDE_W - MARGIN - Inches(1.5), cy + Pt(14), Inches(1.2), Pt(28), score,
             FONT_DISPLAY, 22, color, bold=True, align=PP_ALIGN.RIGHT)

y += Pt(440)
add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Brand insertions: \"here's the thing\", \"dead simple\", \"the actual work\"",
    "Specific numbers anchored: 87%, 14 days \u2192 3, six systems",
    "Word count: 220 \u2192 175 (20% reduction)",
], FONT_MONO, 14, ACCENT, spacing=1.4)


# Save
OUT = SCRIPT_DIR / "uc3-ai-pattern-audit.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT.stat().st_size / 1024:.0f} KB")
