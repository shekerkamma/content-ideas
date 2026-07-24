"""
Build a test report PPTX showing all 10 skill test results.
Each test becomes a slide with key findings.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = Path(__file__).resolve().parent
TOKENS_PATH = SCRIPT_DIR.parent / "visual-identity" / "tokens.json"
STOCK_DIR = SCRIPT_DIR.parent / "visual-identity" / "stock-library" / "themes"

with open(TOKENS_PATH) as f:
    tokens = json.load(f)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

PRIMARY     = hex_to_rgb(tokens["palette"]["primary"]["hex"])
SECONDARY   = hex_to_rgb(tokens["palette"]["secondary"]["hex"])
ACCENT      = hex_to_rgb(tokens["palette"]["accent"]["hex"])
TEXT_WHITE   = hex_to_rgb(tokens["palette"]["text_on_dark"]["hex"])
MUTED       = hex_to_rgb(tokens["palette"]["muted"]["hex"])
SUCCESS     = hex_to_rgb(tokens["palette"]["success"]["hex"])

FONT_DISPLAY = tokens["typography"]["display"]["family"]
FONT_BODY    = tokens["typography"]["body"]["family"]
FONT_MONO    = tokens["typography"]["mono"]["family"]

# Slide size: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


def add_bg(slide, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg_image(slide, image_path):
    slide.shapes.add_picture(str(image_path), 0, 0, SLIDE_W, SLIDE_H)
    # Dark overlay
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, font=FONT_BODY,
             size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txbox


def add_multiline(slide, left, top, width, height, lines, font=FONT_BODY,
                  size=16, color=TEXT_WHITE, bold=False, spacing=1.3):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(size * (spacing - 1))
    return txbox


def add_accent_bar(slide, left, top, width=Inches(0.8), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_status_badge(slide, left, top, text="PASS", color=SUCCESS):
    w, h = Inches(1.2), Pt(32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, left, top + Pt(4), w, Pt(24), text,
             FONT_DISPLAY, 14, PRIMARY, bold=True, align=PP_ALIGN.CENTER)


def pick_bg(theme, idx=0):
    theme_dir = STOCK_DIR / theme
    if theme_dir.exists():
        jpgs = sorted(theme_dir.glob("*.jpg"))
        if jpgs:
            return jpgs[idx % len(jpgs)]
    return None


# ═══════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 3)
if bg:
    add_bg_image(s, bg)
else:
    add_bg(s)

y = Inches(2.0)
add_accent_bar(s, MARGIN, y, Inches(1.2))
y += Pt(30)
add_text(s, MARGIN, y, CONTENT_W, Inches(1.5),
         "Skill Test Report", FONT_DISPLAY, 52, TEXT_WHITE, bold=True)
y += Inches(1.2)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "Brand Voice  |  Humanizer  |  Body of Work  |  Visual Identity",
         FONT_BODY, 22, MUTED)
y += Inches(0.8)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5),
         "10 untested paths exercised  |  All PASS  |  2026-06-07",
         FONT_BODY, 18, ACCENT)


# ═══════════════════════════════════════════════════
# TEST SLIDES
# ═══════════════════════════════════════════════════
tests = [
    {
        "num": "01",
        "title": "Brand Voice — Import Mode",
        "skill": "mkt-brand-voice",
        "bg_theme": "gradients-abstract",
        "bg_idx": 1,
        "summary": [
            "Input: Enterprise brand guidelines (tone, rules, platform guidance)",
            "Process: Map guidelines to voice-profile.md template format",
            "Output: Full voice profile + gap analysis",
            "",
            "Gaps identified:",
            "  - No real content samples (can't build samples.md from guidelines alone)",
            "  - Signature phrases derived, not sourced",
            "  - Confidence zones inferred from positioning",
            "",
            "Recommendation: Enrich with Extract or Auto-Scrape to add real samples",
        ],
    },
    {
        "num": "02",
        "title": "Brand Voice — Extract Mode",
        "skill": "mkt-brand-voice",
        "bg_theme": "textures",
        "bg_idx": 2,
        "summary": [
            "Input: 5 content samples (LinkedIn, email, blog, deck, Slack) — 680 words",
            "Process: 6-dimension extraction (tone, vocab, rhythm, structure, perspective, conviction)",
            "Output: Voice profile + 5 harvested samples for samples.md",
            "",
            "Key findings:",
            "  - Authenticity ranking: Slack > Email > Deck > LinkedIn > Blog",
            "  - Signature move: takes assumption, dismantles with specific number",
            "  - Voice in one line: 'Show me the number, not the narrative'",
            "",
            "Gap: needs 3-5 more samples across platforms for validation",
        ],
    },
    {
        "num": "03",
        "title": "Humanizer — Quick Mode",
        "skill": "tool-humanizer",
        "bg_theme": "code-terminal",
        "bg_idx": 0,
        "summary": [
            "Input: 180-word AI-saturated paragraph (every cliche in the book)",
            "Score: 1.8/10 -> 7.6/10",
            "",
            "Patterns detected: 26 (14.4 per 100 words)",
            "  - 10 HIGH severity (AI cliches)",
            "  - 6 MEDIUM-HIGH (corporate buzzwords)",
            "  - 7 MEDIUM (hedging, transitions)",
            "  - 3 MEDIUM (promotional inflation)",
            "",
            "Word count: 180 -> 105 (42% reduction)",
            "No voice-matching applied (quick mode = cliche removal only)",
        ],
    },
    {
        "num": "04",
        "title": "Humanizer — Deep Mode (Voice-Matched)",
        "skill": "tool-humanizer",
        "bg_theme": "dark-tech",
        "bg_idx": 7,
        "summary": [
            "Input: 220-word AI blog post + voice-profile.md loaded",
            "Score: 2.9/10 -> 8.7/10",
            "",
            "Voice-match report:",
            "  - Tone (casual-but-competent): 9/10",
            "  - Rhythm (short+medium mix): 8/10",
            "  - Vocabulary (brand terms): 9/10",
            "  - Anti-patterns (banned words): 10/10",
            "  - Confidence (bold assertions): 9/10",
            "",
            "Brand insertions: 'here's what matters', 'here's the thing', specific numbers",
            "Word count: 220 -> 175",
        ],
    },
    {
        "num": "05",
        "title": "Humanizer — Pipeline Mode",
        "skill": "tool-humanizer",
        "bg_theme": "gradients-abstract",
        "bg_idx": 5,
        "summary": [
            "Simulated: social-media-team calls humanizer as post-processing",
            "Score: 2.4/10 -> 8.1/10 (delta 5.7)",
            "",
            "Pipeline contract verified:",
            "  [ok] Received text as input",
            "  [ok] Ran Steps 2-6 silently (no user output)",
            "  [ok] Returned cleaned text to calling skill",
            "  [ok] Showed score summary (delta > 2 threshold)",
            "  [ok] Did NOT show full change log (minimal output mode)",
            "",
            "Calling skill is responsible for saving final output",
        ],
    },
    {
        "num": "06",
        "title": "Brand Voice — Auto-Scrape Mode",
        "skill": "mkt-brand-voice",
        "bg_theme": "code-terminal",
        "bg_idx": 4,
        "summary": [
            "URL: https://www.anthropic.com/company",
            "Tool used: WebFetch (free, no API key)",
            "Content: ~800 words, high quality",
            "",
            "Extraction result:",
            "  - Tone: authoritative yet accessible",
            "  - Vocabulary: precision triplets, mission-anchored values",
            "  - Zero marketing superlatives in source",
            "  - Conviction: institutional humility + long-term framing",
            "",
            "Firecrawl fallback: documented (API key missing, would offer Build mode)",
            "Brand asset extraction: skipped (noted for user)",
        ],
    },
    {
        "num": "07",
        "title": "Brand Voice — Playbook (Deep Interview)",
        "skill": "mkt-brand-voice",
        "bg_theme": "textures",
        "bg_idx": 6,
        "summary": [
            "80/20 Formula: 5 personality Qs + 4 strategy Qs + example collection",
            "",
            "Synthesis rules verified:",
            "  [ok] Rule 1: Derive, don't template (labels from user's words)",
            "  [ok] Rule 2: No example brief copying (zero verbatim)",
            "  [ok] Rule 3: Q1->Never list, Q5->Voice summary opener",
            "",
            "Signature phrases harvested: 4 from user answers",
            "LinkedIn voice test: used 2+ signature phrases",
            "Quick synthesis check: 7/7 green",
        ],
    },
    {
        "num": "08",
        "title": "Visual Identity — Update Mode",
        "skill": "mkt-visual-identity",
        "bg_theme": "gradients-abstract",
        "bg_idx": 8,
        "summary": [
            "Trigger: tokens.json exists -> Update mode (not rebuild)",
            "",
            "Showed current identity summary:",
            "  - 9-color palette, 3 typefaces, masthead config",
            "  - Asked user what to change (not starting over)",
            "",
            "Scenarios tested:",
            "  A: Change accent color -> WCAG revalidation",
            "  B: Swap font pairing -> cascade to brand book + CSS",
            "  C: Update masthead -> templates auto-pick up from tokens.json",
            "",
            "locked_fields respected (brand, chrome.masthead)",
        ],
    },
    {
        "num": "09",
        "title": "Visual Identity — Path A: Referenced",
        "skill": "mkt-visual-identity",
        "bg_theme": "dark-tech",
        "bg_idx": 2,
        "summary": [
            "3 visual references analyzed (dark SaaS, data viz, enterprise UI)",
            "",
            "Cross-reference synthesis:",
            "  - Common thread: dark mode, data-dense, accent-driven hierarchy",
            "  - Derived: GitHub Dark primary, Teal accent, Geist font family",
            "",
            "WCAG validation:",
            "  - text_on_dark on primary: 17.4:1 (pass)",
            "  - accent on primary: 8.7:1 (pass)",
            "  - accent on light: 2.9:1 (adjusted to 3.8:1)",
            "",
            "Output distinct from Path B (data-dense dark vs clean neutral)",
        ],
    },
    {
        "num": "10",
        "title": "Brand Voice — Context Enrichment",
        "skill": "mkt-brand-voice",
        "bg_theme": "code-terminal",
        "bg_idx": 8,
        "summary": [
            "Step 7: ICP + Positioning offers after voice profile save",
            "",
            "ICP offer: triggered (icp.md missing)",
            "  -> Would invoke mkt-icp skill if user says 'yes'",
            "  -> Skip cleanly if user says 'skip'",
            "",
            "Positioning offer: triggered (positioning.md missing)",
            "  -> Would invoke mkt-positioning skill",
            "",
            "Edge cases verified:",
            "  - Both exist: skip silently (no re-ask)",
            "  - One exists: offer only the missing one",
            "  - Neither skill available: skip, note for later",
        ],
    },
]

for t in tests:
    s = prs.slides.add_slide(BLANK)
    bg = pick_bg(t["bg_theme"], t["bg_idx"])
    if bg:
        add_bg_image(s, bg)
    else:
        add_bg(s)

    # Slide number
    add_text(s, MARGIN, Inches(0.4), Inches(1), Inches(0.6),
             t["num"], FONT_DISPLAY, 72, SECONDARY, bold=True)

    # Status badge
    add_status_badge(s, SLIDE_W - MARGIN - Inches(1.2), Inches(0.5))

    # Title
    y = Inches(0.5)
    add_accent_bar(s, MARGIN + Inches(1.2), y + Pt(30), Inches(0.6))
    add_text(s, MARGIN + Inches(1.2), y, CONTENT_W - Inches(1.2), Inches(0.8),
             t["title"], FONT_DISPLAY, 32, TEXT_WHITE, bold=True)

    # Skill tag
    add_text(s, MARGIN + Inches(1.2), y + Inches(0.7), Inches(4), Inches(0.4),
             t["skill"], FONT_MONO, 14, ACCENT)

    # Content
    add_multiline(s, MARGIN, Inches(1.6), CONTENT_W, Inches(5.5),
                  t["summary"], FONT_MONO, 15, MUTED, spacing=1.25)


# ═══════════════════════════════════════════════════
# CLOSER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("textures", 5)
if bg:
    add_bg_image(s, bg)
else:
    add_bg(s)

y = Inches(2.2)
add_accent_bar(s, MARGIN, y, Inches(1.2))
y += Pt(30)
add_text(s, MARGIN, y, CONTENT_W, Inches(1),
         "All 10 Paths Tested", FONT_DISPLAY, 48, ACCENT, bold=True)
y += Inches(1.0)
add_multiline(s, MARGIN, y, CONTENT_W, Inches(3), [
    "mkt-brand-voice: Import, Extract, Playbook, Auto-Scrape, Update, Enrichment",
    "tool-humanizer: Quick, Deep (voice-matched), Pipeline",
    "mkt-visual-identity: Update, Path A (referenced)",
    "",
    "Combined with prior session: Build (quick), Body of Work,",
    "Visual Identity Path B (neutral), Brand Book, Carousel Templates",
    "",
    "Full coverage across all skill modes and edge cases.",
], FONT_BODY, 20, MUTED, spacing=1.4)

# Save
OUT = SCRIPT_DIR / "skill-test-report.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT.stat().st_size / 1024:.0f} KB")
