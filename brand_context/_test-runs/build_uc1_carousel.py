"""
Use Case 1: Rough Draft → Branded LinkedIn Carousel PPTX
Content: "Why 87% of AI Pilots Die Before Production"
Voice: Humanized with brand voice profile (deep mode)
Backgrounds: Real stock photos from the library
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR  = Path(__file__).resolve().parent
TOKENS_PATH = SCRIPT_DIR.parent / "visual-identity" / "tokens.json"
STOCK_DIR   = SCRIPT_DIR.parent / "visual-identity" / "stock-library" / "themes"

with open(TOKENS_PATH) as f:
    tokens = json.load(f)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

PRIMARY     = hex_to_rgb(tokens["palette"]["primary"]["hex"])
SECONDARY   = hex_to_rgb(tokens["palette"]["secondary"]["hex"])
ACCENT      = hex_to_rgb(tokens["palette"]["accent"]["hex"])
TEXT_WHITE   = hex_to_rgb(tokens["palette"]["text_on_dark"]["hex"])
MUTED       = hex_to_rgb(tokens["palette"]["muted"]["hex"])
SUCCESS     = hex_to_rgb(tokens["palette"]["success"]["hex"])
DANGER      = hex_to_rgb("#EF4444")

FONT_DISPLAY = tokens["typography"]["display"]["family"]
FONT_BODY    = tokens["typography"]["body"]["family"]
FONT_MONO    = tokens["typography"]["mono"]["family"]

# Portrait carousel
SLIDE_W = Emu(int(1080 / 96 * 914400))
SLIDE_H = Emu(int(1350 / 96 * 914400))

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


def pick_bg(theme, idx=0):
    theme_dir = STOCK_DIR / theme
    if theme_dir.exists():
        jpgs = sorted(theme_dir.glob("*.jpg"))
        if jpgs:
            return jpgs[idx % len(jpgs)]
    return None


def add_bg_image(slide, path):
    slide.shapes.add_picture(str(path), 0, 0, SLIDE_W, SLIDE_H)


def add_overlay(slide, opacity=0.70):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    spPr = shape._element.find(qn("p:spPr"))
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(opacity * 100000)))
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, font=FONT_BODY,
             size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txbox


def add_lines(slide, left, top, width, height, lines, font=FONT_BODY,
              size=18, color=TEXT_WHITE, bold=False, spacing=1.4):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(size * (spacing - 1))
    return txbox


def add_accent_bar(slide, left, top, width=Inches(0.8), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_card(slide, left, top, width, height, fill=SECONDARY, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_masthead(slide):
    masthead = tokens.get("chrome", {}).get("masthead", {})
    labels = masthead.get("labels", ["JUNE 2026", "{{handle}}", "{{tagline}}"])
    labels = [l.replace("{{month_year}}", "JUNE 2026") for l in labels]
    y = Pt(28)
    positions = [MARGIN, SLIDE_W // 2 - Inches(1), SLIDE_W - MARGIN - Inches(2.5)]
    for label, x in zip(labels, positions):
        add_text(slide, x, y, Inches(2.5), Pt(20), label,
                 FONT_BODY, 10, MUTED, bold=True)


def add_pagination(slide, active):
    dot_size = Pt(10)
    gap = Pt(14)
    total_w = 7 * dot_size + 6 * gap
    start_x = (SLIDE_W - total_w) // 2
    y = SLIDE_H - Pt(48)
    for i in range(7):
        x = start_x + i * (dot_size + gap)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dot_size, dot_size)
        d.fill.solid()
        d.fill.fore_color.rgb = ACCENT if i == active else MUTED
        d.line.fill.background()


def add_slide_num(slide, num, y=Inches(1.1)):
    add_text(slide, MARGIN, y, CONTENT_W, Inches(1.5), num,
             FONT_DISPLAY, 110, SECONDARY, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 4)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.68)
add_masthead(s)

y = Inches(2.8)
add_accent_bar(s, MARGIN, y, Inches(1.0))
y += Pt(36)
add_text(s, MARGIN, y, CONTENT_W, Inches(3),
         "Why 87% of\nAI Pilots Die\nBefore Production",
         FONT_DISPLAY, 52, TEXT_WHITE, bold=True)
y += Inches(2.8)
add_lines(s, MARGIN, y, CONTENT_W, Inches(0.8), [
    "Three problems. All fixable.",
    "None require more AI.",
], FONT_BODY, 20, MUTED, spacing=1.3)
y += Inches(1.4)
# CTA
cta = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y, Pt(42), Pt(42))
cta.fill.solid()
cta.fill.fore_color.rgb = ACCENT
cta.line.fill.background()
add_text(s, MARGIN + Pt(8), y + Pt(6), Pt(42), Pt(30), "\u2192",
         FONT_DISPLAY, 20, PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, MARGIN + Pt(52), y + Pt(8), Inches(2), Pt(30),
         "SWIPE TO LEARN", FONT_BODY, 13, TEXT_WHITE, bold=True)
add_pagination(s, 0)


# ═══════════════════════════════════════════════════
# SLIDE 2: THE STAT
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("gradients-abstract", 3)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.72)
add_masthead(s)
add_slide_num(s, "01")

y = Inches(2.8)
add_accent_bar(s, MARGIN, y)
y += Pt(32)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.7),
         "The Number Nobody Talks About", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
y += Inches(1.2)

# Big stat card
card_h = Inches(4.0)
add_card(s, MARGIN, y, CONTENT_W, card_h)
# Big number
add_text(s, MARGIN + Pt(28), y + Pt(20), Inches(3), Inches(1.5),
         "87%", FONT_DISPLAY, 80, DANGER, bold=True)
add_lines(s, MARGIN + Pt(28), y + Pt(120), CONTENT_W - Pt(56), Inches(2.5), [
    "of AI pilots never make it to production.",
    "",
    "Not because the models don't work.",
    "Because nobody defined what",
    "\"working\" means before they",
    "started building.",
], FONT_BODY, 20, MUTED, spacing=1.4)
y += card_h + Pt(24)
add_text(s, MARGIN, y, CONTENT_W, Pt(30),
         "That's not a technology problem.", FONT_DISPLAY, 20, ACCENT, bold=True)
add_pagination(s, 1)


# ═══════════════════════════════════════════════════
# SLIDE 3: PROBLEM 1 — NO METRIC
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("textures", 0)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.74)
add_masthead(s)
add_slide_num(s, "02")

y = Inches(2.8)
add_accent_bar(s, MARGIN, y)
y += Pt(32)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.7),
         "No Metric, No Project", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
y += Inches(1.0)

add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Here's the thing \u2014",
    "\"improve efficiency\" isn't a success metric.",
], FONT_BODY, 21, MUTED, spacing=1.5)
y += Inches(1.8)

# Quote card with accent border
qh = Inches(1.6)
add_card(s, MARGIN, y, CONTENT_W, qh)
add_accent_bar(s, MARGIN, y, width=Pt(4), height=qh)
add_text(s, MARGIN + Pt(32), y + Pt(20), CONTENT_W - Pt(64), qh - Pt(40),
         "\"Reduce claim processing\nfrom 14 days to 3.\"",
         FONT_DISPLAY, 28, ACCENT, bold=True)

y += qh + Pt(40)
add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "If you can't put a number on the goal",
    "before you write a single line of code,",
    "you're not ready for AI.",
    "You're ready for a whiteboard.",
], FONT_BODY, 20, MUTED, spacing=1.4)
add_pagination(s, 2)


# ═══════════════════════════════════════════════════
# SLIDE 4: PROBLEM 2 — DATA
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("code-terminal", 3)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.70)
add_masthead(s)
add_slide_num(s, "03")

y = Inches(2.8)
add_accent_bar(s, MARGIN, y)
y += Pt(32)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.7),
         "Your Data Isn't Ready", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
y += Inches(1.0)

add_lines(s, MARGIN, y, CONTENT_W, Inches(2), [
    "The proof-of-concept worked because you",
    "hand-picked clean data for the demo.",
], FONT_BODY, 21, MUTED, spacing=1.5)
y += Inches(1.6)

card_h = Inches(2.5)
add_card(s, MARGIN, y, CONTENT_W, card_h)
add_lines(s, MARGIN + Pt(28), y + Pt(20), CONTENT_W - Pt(56), card_h - Pt(40), [
    "Production data is messy, late,",
    "undocumented, and lives in six",
    "different systems that don't",
    "talk to each other.",
], FONT_BODY, 20, TEXT_WHITE, spacing=1.5)

y += card_h + Pt(24)
add_lines(s, MARGIN, y, CONTENT_W, Inches(0.6), [
    "Fix the pipes before you train the model.",
    "Dead simple. Nobody does it.",
], FONT_DISPLAY, 19, ACCENT, bold=True, spacing=1.3)
add_pagination(s, 3)


# ═══════════════════════════════════════════════════
# SLIDE 5: PROBLEM 3 — HUMANS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("dark-tech", 8)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.72)
add_masthead(s)
add_slide_num(s, "04")

y = Inches(2.8)
add_accent_bar(s, MARGIN, y)
y += Pt(32)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.7),
         "You Forgot the Humans", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
y += Inches(1.0)

add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "You can build the best model in the world.",
    "If the team doesn't trust it, they won't use it.",
], FONT_BODY, 21, MUTED, spacing=1.5)
y += Inches(1.8)

card_h = Inches(2.0)
add_card(s, MARGIN, y, CONTENT_W, card_h)
add_accent_bar(s, MARGIN, y, width=Pt(4), height=card_h)
add_lines(s, MARGIN + Pt(32), y + Pt(20), CONTENT_W - Pt(64), card_h - Pt(40), [
    "This isn't a \"change management\" problem \u2014",
    "it's a \"did you involve the actual users",
    "before building?\" problem.",
], FONT_BODY, 20, TEXT_WHITE, spacing=1.5)

y += card_h + Pt(32)
add_lines(s, MARGIN, y, CONTENT_W, Inches(1.2), [
    "Show them the tool. Let them break it.",
    "Then fix what they broke.",
    "That's the actual work.",
], FONT_DISPLAY, 19, ACCENT, bold=True, spacing=1.3)
add_pagination(s, 4)


# ═══════════════════════════════════════════════════
# SLIDE 6: THE FIX — CHECKLIST
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("gradients-abstract", 7)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.72)
add_masthead(s)
add_slide_num(s, "05")

y = Inches(2.8)
add_accent_bar(s, MARGIN, y)
y += Pt(32)
add_text(s, MARGIN, y, CONTENT_W, Inches(0.7),
         "The Dead Simple Checklist", FONT_DISPLAY, 36, TEXT_WHITE, bold=True)
y += Inches(1.2)

checks = [
    ("01", "Define a measurable outcome\nbefore day 1"),
    ("02", "Audit your data pipelines\nbefore picking a model"),
    ("03", "Put the tool in front of real\nusers in week 2, not month 6"),
]
for i, (num, text) in enumerate(checks):
    cy = y + i * Pt(130)
    ch = Pt(110)
    add_card(s, MARGIN, cy, CONTENT_W, ch)
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN + Pt(20), cy + Pt(22), Pt(56), Pt(56))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_text(s, MARGIN + Pt(20), cy + Pt(28), Pt(56), Pt(44), num,
             FONT_DISPLAY, 20, PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN + Pt(92), cy + Pt(22), CONTENT_W - Pt(130), ch - Pt(44), text,
             FONT_BODY, 20, TEXT_WHITE)

y += Pt(420)
add_card(s, MARGIN, y, CONTENT_W, Pt(56), fill=ACCENT)
add_text(s, MARGIN, y + Pt(12), CONTENT_W, Pt(36),
         "Skip any of these and your pilot fails. That's it.",
         FONT_DISPLAY, 17, PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_pagination(s, 5)


# ═══════════════════════════════════════════════════
# SLIDE 7: CLOSER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = pick_bg("textures", 5)
if bg: add_bg_image(s, bg)
add_overlay(s, 0.65)
add_masthead(s)

y = Inches(2.8)
add_accent_bar(s, MARGIN, y, Inches(1.2))
y += Pt(40)
add_text(s, MARGIN, y, CONTENT_W, Inches(2.5),
         "The best AI project\nis the one that\nactually ships.",
         FONT_DISPLAY, 46, ACCENT, bold=True)
y += Inches(2.8)
add_lines(s, MARGIN, y, CONTENT_W, Inches(1.5), [
    "Not the one with the cleanest architecture,",
    "the most sophisticated model,",
    "or the biggest budget.",
    "",
    "The one that ships. With a number attached.",
], FONT_BODY, 20, MUTED, spacing=1.4)

y += Inches(2.4)
add_accent_bar(s, MARGIN, y, Inches(6.5), Pt(1))
y += Pt(20)
brand_name = tokens.get("brand", "{{brand_name}}")
add_text(s, MARGIN, y, CONTENT_W, Inches(0.5), brand_name,
         FONT_DISPLAY, 24, TEXT_WHITE, bold=True)
masthead_labels = tokens.get("chrome", {}).get("masthead", {}).get("labels", [])
tagline = masthead_labels[2] if len(masthead_labels) > 2 else ""
add_text(s, MARGIN, y + Pt(32), CONTENT_W, Inches(0.3), tagline,
         FONT_BODY, 14, MUTED)
add_pagination(s, 6)


# ═══════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════
OUT = SCRIPT_DIR / "uc1-why-ai-pilots-fail.pptx"
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT.stat().st_size / 1024:.0f} KB")
