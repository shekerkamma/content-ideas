#!/usr/bin/env python3
"""Build a reviewed, local-only agent-native apps PPTX deck.

This version uses the branded pptxkit workflow and PPT-native vector-like
diagrams. It intentionally avoids decorative AI-generated architecture images.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

RUN_DIR = Path(__file__).resolve().parent
TOOLKIT_DIR = Path(
    "/home/shekerk/snap/codex/64/plugins/cache/shekerk-toolkit/"
    "shekerk-toolkit/1.0.0/skills/branded-pptx-deck/scripts"
)
sys.path.insert(0, str(TOOLKIT_DIR))

from pptxkit import Brand, Deck, RGBColor, validate_pptx  # noqa: E402

TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE", "/home/shekerk/.claude/templates/branded-template.pptx"))
OUT = RUN_DIR / "agent-native-apps-local-rebuild-v5-reviewed.pptx"


def hx(value: str) -> RGBColor:
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


class TemplateDeck(Deck):
    """Use the branded template as the PowerPoint package source, then wipe slides."""

    def __init__(self, template: Path, brand: Brand, footer: str):
        self.b = brand
        self.prs = Presentation(str(template))
        self.W = Inches(13.333)
        self.H = Inches(7.5)
        self.prs.slide_width = self.W
        self.prs.slide_height = self.H
        self.M = Inches(0.62)
        self.CW = self.W - self.M * 2
        self.footer_text = footer
        self._wipe()


brand = Brand(
    NAVY=hx("0A1628"),
    NAVY_2=hx("12243A"),
    TEAL=hx("00C9A7"),
    ACCENT=hx("009B82"),
    DARK_TEAL=hx("008F75"),
    LIGHT_TEAL=hx("E0F7F1"),
    GOLD=hx("FFB800"),
    AMBER=hx("F2A83B"),
    CORAL=hx("E05A6B"),
    WHITE=hx("FFFFFF"),
    SOFT=hx("F4F7F8"),
    INK=hx("1B2B3C"),
    MUTED=hx("5B6B7C"),
    GRID=hx("D9DFE5"),
    FONT="Aptos",
    FONT_H="Aptos Display",
)

d = TemplateDeck(TEMPLATE, brand, "Agent-Native Apps Opportunity | Local reviewed rebuild | June 2026")
b = d.b


def set_text_fit(shape, size: float) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = b.FONT
            run.font.size = Pt(size)
            run.font.color.rgb = b.INK


def rounded(slide, x, y, w, h, fill, line=None, radius=0.12, shadow=False):
    return d.rect(slide, x, y, w, h, fill, line=line or b.GRID, radius=radius, shadow=shadow)


def label(slide, text, x, y, w, h, *, size=13, color=None, bold=False, align=PP_ALIGN.CENTER, font=None):
    hin = h / 914400
    if hin < 0.22:
        size = min(size, 8.8)
    return d.text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color or b.INK,
        bold=bold,
        align=align,
        font=font or b.FONT,
        anchor=MSO_ANCHOR.MIDDLE,
        shrink=True,
    )


def card(slide, x, y, w, h, title, body, *, accent=None, fill=None, title_size=14, body_size=11.5):
    accent = accent or b.TEAL
    rounded(slide, x, y, w, h, fill or b.WHITE, line=b.GRID, radius=0.08, shadow=False)
    d.rect(slide, x, y, Inches(0.08), h, accent)
    hin = h / 914400
    win = w / 914400
    if hin < 1.05 or (win < 2.05 and hin < 1.3):
        d.text(slide, title, x + Inches(0.16), y + Inches(0.1), w - Inches(0.28), h - Inches(0.16),
               size=min(title_size, 11.0), color=b.NAVY, bold=True, font=b.FONT_H,
               anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        return
    title_size = min(title_size, 13.5)
    body_size = min(body_size, 11.3)
    title_h = Inches(0.42)
    body_y = y + Inches(0.6)
    d.text(slide, title, x + Inches(0.18), y + Inches(0.1), w - Inches(0.3), title_h,
           size=title_size, color=b.NAVY, bold=True, font=b.FONT_H, shrink=True)
    d.text(slide, body, x + Inches(0.18), body_y, w - Inches(0.32), y + h - body_y - Inches(0.1),
           size=body_size, color=b.MUTED, font=b.FONT, shrink=True)


def badge(slide, n, x, y, color=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.34), Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color or b.CORAL
    shp.line.fill.background()
    shp.text = str(n)
    set_text_fit(shp, 11)
    for p in shp.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = b.WHITE
            r.font.bold = True


def arrow(slide, x, y, w, h, *, color=None, rotate=0):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color or b.TEAL
    shp.line.fill.background()
    shp.rotation = rotate
    return shp


def header(slide, title, subtitle=None):
    d.rect(slide, 0, 0, d.W, Inches(0.14), b.TEAL)
    d.text(slide, title, d.M, Inches(0.34), d.CW, Inches(0.58), size=29, color=b.NAVY,
           bold=True, font=b.FONT_H, shrink=True)
    d.rect(slide, d.M, Inches(1.05), Inches(1.3), Inches(0.05), b.TEAL)
    if subtitle:
        d.text(slide, subtitle, d.M, Inches(1.17), d.CW, Inches(0.36), size=14.4, color=b.MUTED,
               font=b.FONT, shrink=True)


def footer(slide, page, total):
    d.footer(slide, page, total)


def content_slide(page, total, title, subtitle, body, bullets, diagram, takeaway):
    slide = d.slide(fill=b.WHITE)
    header(slide, title, subtitle)

    # Left narrative panel.
    rounded(slide, Inches(0.62), Inches(1.58), Inches(4.95), Inches(4.72), b.SOFT, line=b.GRID, radius=0.08)
    d.text(slide, "Narrative", Inches(0.9), Inches(1.78), Inches(4.15), Inches(0.28), size=13.5,
           color=b.ACCENT, bold=True, font=b.FONT_H)
    d.text(slide, body, Inches(0.9), Inches(2.16), Inches(4.32), Inches(1.82), size=14.8,
           color=b.INK, font=b.FONT, ls=1.02)
    if bullets:
        bullet_runs = [
            {"text": line, "bullet": True, "size": 13.4, "space_before": 6, "color": b.INK, "font": b.FONT}
            for line in bullets
        ]
        d.text(slide, bullet_runs, Inches(0.92), Inches(4.08), Inches(4.22), Inches(1.92),
               size=13.4, color=b.INK, font=b.FONT, ls=1.0)

    # Right diagram panel.
    rounded(slide, Inches(5.85), Inches(1.58), Inches(6.85), Inches(4.72), b.WHITE, line=b.GRID, radius=0.08)
    draw_diagram(diagram, slide, Inches(6.02), Inches(1.84), Inches(6.45), Inches(4.22))

    # Bottom takeaway bar.
    d.rect(slide, Inches(0.62), Inches(6.42), Inches(12.08), Inches(0.54), b.NAVY)
    d.text(slide, takeaway, Inches(0.86), Inches(6.51), Inches(11.62), Inches(0.32), size=12.8,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    footer(slide, page, total)


def section_slide(page, total, section_no, title, subtitle, bullets):
    slide = d.slide(fill=b.NAVY)
    d.rect(slide, 0, 0, Inches(0.18), d.H, b.TEAL)
    d.text(slide, f"SECTION {section_no:02d}", Inches(0.72), Inches(1.06), Inches(3.0), Inches(0.35),
           size=15, color=b.TEAL, bold=True, font=b.FONT_H)
    d.text(slide, title, Inches(0.72), Inches(1.5), Inches(9.8), Inches(1.34), size=38,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    d.rect(slide, Inches(0.72), Inches(2.86), Inches(1.8), Inches(0.06), b.GOLD)
    d.text(slide, subtitle, Inches(0.72), Inches(3.18), Inches(8.1), Inches(0.64), size=19,
           color=b.LIGHT_TEAL, font=b.FONT, shrink=True)
    x = Inches(0.82)
    for i, line in enumerate(bullets, 1):
        card(slide, x + Inches((i - 1) * 3.75), Inches(4.35), Inches(3.35), Inches(1.22),
             f"{i}. {line[0]}", line[1], accent=b.TEAL if i != 2 else b.GOLD,
             fill=b.NAVY_2, title_size=15, body_size=11.8)
        # Recolor card text placed on dark cards.
        for shape in list(slide.shapes)[-2:]:
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.color.rgb == b.NAVY:
                            r.font.color.rgb = b.WHITE
                        elif r.font.color.rgb == b.MUTED:
                            r.font.color.rgb = b.LIGHT_TEAL
    footer(slide, page, total)


def cover(page, total):
    slide = d.slide(fill=b.NAVY)
    d.rect(slide, Inches(8.65), 0, Inches(4.68), d.H, b.NAVY_2)
    d.rect(slide, Inches(8.65), 0, Inches(0.08), d.H, b.TEAL)
    d.text(slide, "AGENT-NATIVE", Inches(9.1), Inches(0.52), Inches(3.1), Inches(0.3),
           size=18, color=b.TEAL, bold=True, font=b.FONT_H, shrink=True)
    d.text(slide, "Context  Tools  Surface  Memory", Inches(9.1), Inches(0.9), Inches(3.2), Inches(0.24),
           size=10.8, color=b.LIGHT_TEAL, font=b.FONT, shrink=True)
    d.text(slide, "LOCAL REVIEWED REBUILD", Inches(0.74), Inches(0.82), Inches(4.9), Inches(0.34),
           size=15, color=b.TEAL, bold=True, font=b.FONT_H)
    d.text(slide, "The Best New AI Opportunity Nobody Is Talking About",
           Inches(0.74), Inches(1.38), Inches(7.35), Inches(1.75), size=35,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    d.text(slide, "Agent-native apps for Codex, Cursor, and bring-your-own-agent workflows",
           Inches(0.78), Inches(3.2), Inches(6.8), Inches(0.76), size=18.5,
           color=b.LIGHT_TEAL, font=b.FONT, shrink=True)
    d.rect(slide, Inches(0.78), Inches(4.08), Inches(2.1), Inches(0.07), b.GOLD)
    for i, chip in enumerate(["context", "tools", "shared surface", "memory loop"]):
        d.rect(slide, Inches(0.78 + i * 1.58), Inches(4.45), Inches(1.35), Inches(0.38),
               b.TEAL if i != 2 else b.GOLD)
        d.text(slide, chip, Inches(0.84 + i * 1.58), Inches(4.53), Inches(1.24), Inches(0.15),
               size=9.2, color=b.NAVY if i == 2 else b.WHITE, bold=True, font=b.FONT_H,
               align=PP_ALIGN.CENTER, shrink=True)
    draw_diagram("cover-hub", slide, Inches(8.98), Inches(0.82), Inches(3.92), Inches(5.8))
    d.text(slide, "Thesis: build apps that a user's own agent can enter, operate, and improve.",
           Inches(0.78), Inches(6.18), Inches(7.0), Inches(0.62), size=15.4,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    footer(slide, page, total)


def executive_summary(page, total):
    slide = d.slide(fill=b.WHITE)
    header(slide, "The opportunity is an inversion", "BLUF: the product surface moves inside the user's agent workflow.")
    d.rect(slide, Inches(0.68), Inches(1.58), Inches(11.95), Inches(0.98), b.NAVY)
    d.text(slide,
           "Agent-native apps make the product a shared surface for the user's agent.",
           Inches(0.95), Inches(1.78), Inches(11.35), Inches(0.44), size=16.4,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    cards = [
        ("1. The old model is bounded",
         "An app-owned assistant can only use the context and actions exposed inside that product."),
        ("2. The new model is portable",
         "The user's agent brings files, project memory, browser context, tools, and judgment into many apps."),
        ("3. The wedge is the surface",
         "The product opportunity is a reviewable workspace: canvas, document, board, sheet, or operational console."),
    ]
    for idx, (title, body) in enumerate(cards):
        card(slide, Inches(0.72 + idx * 4.03), Inches(2.82), Inches(3.58), Inches(2.02),
             title, body, accent=[b.CORAL, b.TEAL, b.GOLD][idx], title_size=14.2, body_size=12.2)
    d.rect(slide, Inches(0.72), Inches(5.28), Inches(11.88), Inches(0.94), b.SOFT)
    d.text(slide,
           "Strategic reading: do not build another isolated chatbot. Build a surface that can be operated through MCP, CLI, API, browser automation, and explicit human approval.",
           Inches(0.96), Inches(5.45), Inches(11.35), Inches(0.48), size=14.3,
           color=b.INK, bold=True, font=b.FONT_H, shrink=True)
    d.rect(slide, Inches(0.72), Inches(6.46), Inches(11.88), Inches(0.5), b.NAVY)
    d.text(slide, "The moat is not the chat box. The moat is context, tool access, and visible collaboration.",
           Inches(0.96), Inches(6.54), Inches(11.35), Inches(0.28), size=12.6,
           color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    footer(slide, page, total)


def draw_diagram(kind, slide, x, y, w, h):
    if kind == "cover-hub":
        card(slide, x, y, w, Inches(0.82), "User's agent", "The outer operating layer", accent=b.TEAL,
             fill=b.NAVY, title_size=13, body_size=10)
        card(slide, x, y + Inches(1.12), w, Inches(0.82), "App surface", "Browser, canvas, document, board", accent=b.GOLD,
             fill=b.NAVY, title_size=13, body_size=10)
        card(slide, x, y + Inches(2.24), w, Inches(0.82), "Tool contract", "MCP, CLI, API, browser actions", accent=b.TEAL,
             fill=b.NAVY, title_size=13, body_size=10)
        card(slide, x, y + Inches(3.36), w, Inches(0.82), "Memory loop", "Approved output becomes future context", accent=b.GOLD,
             fill=b.NAVY, title_size=13, body_size=10)
        arrow(slide, x + Inches(1.58), y + Inches(0.88), Inches(0.36), Inches(0.24), color=b.GOLD, rotate=90)
        arrow(slide, x + Inches(1.58), y + Inches(2.0), Inches(0.36), Inches(0.24), color=b.GOLD, rotate=90)
        arrow(slide, x + Inches(1.58), y + Inches(3.12), Inches(0.36), Inches(0.24), color=b.GOLD, rotate=90)
        label(slide, "BYO agent replaces app-owned assistant", x, y + Inches(4.78), w, Inches(0.44),
              size=12.5, color=b.LIGHT_TEAL, bold=True)
        return

    if kind == "old-model":
        rounded(slide, x + Inches(1.65), y + Inches(0.25), Inches(3.65), Inches(2.75), b.SOFT, line=b.GRID)
        label(slide, "One app boundary", x + Inches(1.85), y + Inches(0.42), Inches(3.2), Inches(0.25),
              size=13.5, color=b.NAVY, bold=True)
        card(slide, x + Inches(2.05), y + Inches(0.92), Inches(2.65), Inches(0.72), "Product UI", "User actions live here", accent=b.CORAL, body_size=10.8)
        card(slide, x + Inches(2.05), y + Inches(1.9), Inches(2.65), Inches(0.72), "App-owned AI", "Context is app-scoped", accent=b.CORAL, body_size=10.8)
        card(slide, x + Inches(0.1), y + Inches(0.95), Inches(1.25), Inches(0.82), "Local files", "outside", accent=b.GRID, body_size=10)
        card(slide, x + Inches(5.55), y + Inches(0.95), Inches(1.25), Inches(0.82), "Other apps", "outside", accent=b.GRID, body_size=10)
        card(slide, x + Inches(5.55), y + Inches(2.0), Inches(1.25), Inches(0.82), "Memory", "outside", accent=b.GRID, body_size=10)
        d.rect(slide, x + Inches(1.44), y + Inches(0.9), Inches(0.06), Inches(2.15), b.CORAL)
        d.rect(slide, x + Inches(5.38), y + Inches(0.9), Inches(0.06), Inches(2.15), b.CORAL)
        label(slide, "Context stops at the app edge", x + Inches(1.6), y + Inches(3.45), Inches(3.8), Inches(0.35),
              size=15, color=b.CORAL, bold=True)
        return

    if kind == "new-model":
        rounded(slide, x + Inches(0.25), y + Inches(0.22), Inches(6.3), Inches(3.7), b.SOFT, line=b.TEAL)
        label(slide, "User-owned agent workspace", x + Inches(0.55), y + Inches(0.42), Inches(5.7), Inches(0.3),
              size=14.5, color=b.NAVY, bold=True)
        card(slide, x + Inches(2.3), y + Inches(1.0), Inches(2.15), Inches(0.8), "App surface", "runs in browser", accent=b.GOLD, body_size=10.5)
        card(slide, x + Inches(0.55), y + Inches(2.28), Inches(1.45), Inches(0.72), "Files", "local truth", accent=b.TEAL, body_size=10.2)
        card(slide, x + Inches(2.62), y + Inches(2.28), Inches(1.45), Inches(0.72), "Tools", "MCP / CLI", accent=b.TEAL, body_size=10.2)
        card(slide, x + Inches(4.72), y + Inches(2.28), Inches(1.45), Inches(0.72), "Memory", "write-back", accent=b.TEAL, body_size=10.2)
        arrow(slide, x + Inches(1.95), y + Inches(1.23), Inches(0.45), Inches(0.24), color=b.GOLD)
        arrow(slide, x + Inches(4.38), y + Inches(1.23), Inches(0.45), Inches(0.24), color=b.GOLD, rotate=180)
        label(slide, "App inside agent flow", x + Inches(0.72), y + Inches(3.38), Inches(5.25), Inches(0.42),
              size=14, color=b.ACCENT, bold=True)
        return

    if kind == "why-now":
        items = [
            ("Local agents", "Codex/Cursor inspect files, run commands, and coordinate work."),
            ("Tool standard", "MCP exposes external systems as model-invoked tools."),
            ("Web surfaces", "Browser-hosted apps can be visible to human and agent."),
            ("Permission loops", "Sensitive actions can require explicit approval."),
        ]
        for i, (title, text) in enumerate(items):
            card(slide, x + Inches((i % 2) * 3.45), y + Inches((i // 2) * 1.75), Inches(3.1), Inches(1.25),
                 title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=13.8, body_size=11.3)
            badge(slide, i + 1, x + Inches((i % 2) * 3.45 + 2.62), y + Inches((i // 2) * 1.75 + 0.12), b.NAVY)
        label(slide, "Convergence creates the wedge", x + Inches(1.4), y + Inches(3.75), Inches(4.1), Inches(0.34),
              size=15, color=b.NAVY, bold=True)
        return

    if kind == "definition":
        rounded(slide, x + Inches(2.08), y + Inches(1.28), Inches(2.55), Inches(1.08), b.NAVY, line=b.TEAL)
        label(slide, "Agent-native app", x + Inches(2.25), y + Inches(1.52), Inches(2.2), Inches(0.38),
              size=16, color=b.WHITE, bold=True)
        items = [
            (0.1, 0.25, "Context", "reads project truth", b.TEAL),
            (4.95, 0.25, "Tools", "takes real actions", b.GOLD),
            (0.1, 3.0, "Surface", "work stays visible", b.ACCENT),
            (4.95, 3.0, "Memory", "quality compounds", b.CORAL),
        ]
        for ix, iy, title, text, color in items:
            card(slide, x + Inches(ix), y + Inches(iy), Inches(1.75), Inches(0.86), title, text,
                 accent=color, title_size=12.8, body_size=10.2)
        arrow(slide, x + Inches(1.82), y + Inches(0.56), Inches(0.6), Inches(0.2), color=b.TEAL)
        arrow(slide, x + Inches(4.42), y + Inches(0.56), Inches(0.6), Inches(0.2), color=b.GOLD, rotate=180)
        arrow(slide, x + Inches(1.82), y + Inches(3.28), Inches(0.6), Inches(0.2), color=b.ACCENT)
        arrow(slide, x + Inches(4.42), y + Inches(3.28), Inches(0.6), Inches(0.2), color=b.CORAL, rotate=180)
        return

    if kind == "job":
        card(slide, x + Inches(0.3), y + Inches(1.0), Inches(1.65), Inches(1.0), "Human", "judgment, taste, approvals", accent=b.GOLD, title_size=13.5)
        card(slide, x + Inches(2.55), y + Inches(0.45), Inches(1.75), Inches(1.0), "Agent", "context, planning, tool use", accent=b.TEAL, title_size=13.5)
        card(slide, x + Inches(4.9), y + Inches(1.0), Inches(1.65), Inches(1.0), "App", "surface, state, actions", accent=b.ACCENT, title_size=13.5)
        arrow(slide, x + Inches(1.95), y + Inches(1.28), Inches(0.55), Inches(0.24), color=b.GOLD)
        arrow(slide, x + Inches(4.32), y + Inches(1.28), Inches(0.55), Inches(0.24), color=b.TEAL)
        rounded(slide, x + Inches(1.65), y + Inches(2.85), Inches(3.62), Inches(0.82), b.NAVY, line=b.NAVY)
        label(slide, "Outcome: completed work, not generated chat", x + Inches(1.86), y + Inches(3.04),
              Inches(3.2), Inches(0.32), size=14, color=b.WHITE, bold=True)
        return

    if kind == "four-blocks":
        items = [
            ("Context", "What the agent needs to know before acting.", b.TEAL),
            ("Connectivity", "How the agent reads, writes, and invokes.", b.GOLD),
            ("Collaboration", "Where human and agent see the same object.", b.ACCENT),
            ("Compounding", "How approved work becomes future memory.", b.CORAL),
        ]
        for i, (title, text, color) in enumerate(items):
            card(slide, x + Inches((i % 2) * 3.45), y + Inches((i // 2) * 1.74), Inches(3.12), Inches(1.22),
                 title, text, accent=color, title_size=14.2, body_size=11.8)
        d.rect(slide, x + Inches(3.32), y + Inches(0.2), Inches(0.06), Inches(3.3), b.GRID)
        d.rect(slide, x + Inches(0.15), y + Inches(1.57), Inches(6.5), Inches(0.06), b.GRID)
        return

    if kind == "reference":
        rows = [
            ("1. Context sources", "files, docs, assets, strategy", b.TEAL),
            ("2. Agent host", "Codex, Cursor, Claude Code, custom shell", b.GOLD),
            ("3. Tool layer", "MCP, CLI, API, browser actions", b.ACCENT),
            ("4. Shared surface", "canvas, doc, sheet, board, console", b.CORAL),
            ("5. Write-back", "approved output updates memory", b.TEAL),
        ]
        for i, (title, text, color) in enumerate(rows):
            yy = y + Inches(0.15 + i * 0.78)
            card(slide, x + Inches(0.42), yy, Inches(5.9), Inches(0.55), title, text,
                 accent=color, title_size=12.2, body_size=9.8)
            if i < len(rows) - 1:
                arrow(slide, x + Inches(3.05), yy + Inches(0.58), Inches(0.32), Inches(0.2), color=b.GRID, rotate=90)
        return

    if kind == "context":
        labels = [
            ("Project files", "README, code, design.md"),
            ("Strategy docs", "brand, ICP, pricing"),
            ("Reference assets", "screenshots, exports, data"),
            ("Local memory", "decisions and corrections"),
        ]
        for i, (title, text) in enumerate(labels):
            card(slide, x + Inches(0.65 + i * 0.28), y + Inches(0.35 + i * 0.72),
                 Inches(5.65), Inches(0.72), title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i],
                 title_size=13, body_size=10.5)
        label(slide, "Context is the product's operating memory", x + Inches(0.95), y + Inches(3.55),
              Inches(5.1), Inches(0.34), size=14.6, color=b.NAVY, bold=True)
        return

    if kind == "tools":
        card(slide, x + Inches(2.36), y + Inches(1.55), Inches(2.1), Inches(0.86), "Agent", "chooses the next action", accent=b.TEAL, title_size=14)
        endpoints = [
            (0.2, 0.35, "MCP", "discover + invoke"),
            (4.95, 0.35, "CLI", "terminal control"),
            (0.2, 3.0, "API", "structured calls"),
            (4.95, 3.0, "Browser", "visible operations"),
        ]
        for i, (ix, iy, title, text) in enumerate(endpoints):
            card(slide, x + Inches(ix), y + Inches(iy), Inches(1.72), Inches(0.82), title, text,
                 accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=13, body_size=10.2)
        arrow(slide, x + Inches(1.92), y + Inches(0.62), Inches(0.7), Inches(0.22), color=b.GRID)
        arrow(slide, x + Inches(4.45), y + Inches(0.62), Inches(0.7), Inches(0.22), color=b.GRID, rotate=180)
        arrow(slide, x + Inches(1.92), y + Inches(3.27), Inches(0.7), Inches(0.22), color=b.GRID)
        arrow(slide, x + Inches(4.45), y + Inches(3.27), Inches(0.7), Inches(0.22), color=b.GRID, rotate=180)
        return

    if kind == "surface":
        rounded(slide, x + Inches(1.55), y + Inches(0.55), Inches(3.8), Inches(2.55), b.SOFT, line=b.TEAL)
        label(slide, "Shared work surface", x + Inches(2.0), y + Inches(0.78), Inches(2.9), Inches(0.3),
              size=14.2, color=b.NAVY, bold=True)
        card(slide, x + Inches(2.02), y + Inches(1.28), Inches(2.8), Inches(0.56), "Canvas / doc / sheet", "same state for both", accent=b.GOLD, body_size=9.6)
        card(slide, x + Inches(0.25), y + Inches(1.35), Inches(1.05), Inches(0.82), "Human", "reviews", accent=b.CORAL, title_size=12.5, body_size=9.4)
        card(slide, x + Inches(5.62), y + Inches(1.35), Inches(1.05), Inches(0.82), "Agent", "edits", accent=b.TEAL, title_size=12.5, body_size=9.4)
        arrow(slide, x + Inches(1.28), y + Inches(1.6), Inches(0.48), Inches(0.22), color=b.CORAL)
        arrow(slide, x + Inches(5.08), y + Inches(1.6), Inches(0.48), Inches(0.22), color=b.TEAL, rotate=180)
        label(slide, "Visible state makes correction cheap", x + Inches(1.0), y + Inches(3.48),
              Inches(5.0), Inches(0.44), size=13.6, color=b.ACCENT, bold=True)
        return

    if kind == "memory":
        steps = [
            ("Prompt", "desired outcome"),
            ("Work", "agent acts"),
            ("Review", "human approves"),
            ("Write-back", "context updates"),
        ]
        for i, (title, text) in enumerate(steps):
            card(slide, x + Inches(0.1 + i * 1.65), y + Inches(1.45), Inches(1.28), Inches(0.9),
                 title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=12.5, body_size=9.2)
            if i < 3:
                arrow(slide, x + Inches(1.39 + i * 1.65), y + Inches(1.76), Inches(0.38), Inches(0.18), color=b.GRID)
        arrow(slide, x + Inches(5.58), y + Inches(2.42), Inches(0.45), Inches(0.2), color=b.CORAL, rotate=90)
        arrow(slide, x + Inches(0.78), y + Inches(2.94), Inches(4.88), Inches(0.2), color=b.GRID, rotate=180)
        label(slide, "Approved work becomes reusable context", x + Inches(0.75), y + Inches(3.32),
              Inches(5.4), Inches(0.44), size=13.6, color=b.NAVY, bold=True)
        return

    if kind == "trust":
        steps = [
            ("Expose", "show tools"),
            ("Invoke", "show action"),
            ("Confirm", "ask when risky"),
            ("Apply", "write changes"),
        ]
        for i, (title, text) in enumerate(steps):
            card(slide, x + Inches(0.2 + i * 1.62), y + Inches(1.15), Inches(1.24), Inches(0.98),
                 title, text, accent=[b.TEAL, b.GOLD, b.CORAL, b.ACCENT][i], title_size=12.5, body_size=9.4)
            badge(slide, i + 1, x + Inches(0.64 + i * 1.62), y + Inches(0.62), [b.TEAL, b.GOLD, b.CORAL, b.ACCENT][i])
            if i < 3:
                arrow(slide, x + Inches(1.45 + i * 1.62), y + Inches(1.5), Inches(0.36), Inches(0.18), color=b.GRID)
        label(slide, "Trust is designed into the interface", x + Inches(0.65), y + Inches(3.18),
              Inches(5.8), Inches(0.44), size=13.7, color=b.NAVY, bold=True)
        return

    if kind == "magicpath":
        flow = [
            ("design.md", "brand + page rules", b.TEAL),
            ("Agent host", "reads local context", b.GOLD),
            ("MagicPath", "design canvas", b.ACCENT),
            ("Codebase", "pull design back", b.CORAL),
        ]
        for i, (title, text, color) in enumerate(flow):
            card(slide, x + Inches(0.12 + i * 1.64), y + Inches(1.35), Inches(1.3), Inches(0.92),
                 title, text, accent=color, title_size=11.9, body_size=9.1)
            if i < 3:
                arrow(slide, x + Inches(1.43 + i * 1.64), y + Inches(1.66), Inches(0.34), Inches(0.18), color=b.GRID)
        label(slide, "The agent brings the context", x + Inches(0.52),
              y + Inches(3.12), Inches(5.9), Inches(0.44), size=13.6, color=b.ACCENT, bold=True)
        return

    if kind == "notion":
        card(slide, x + Inches(0.35), y + Inches(0.75), Inches(1.7), Inches(0.9), "Project folder", "local plan, tasks, docs", accent=b.TEAL)
        card(slide, x + Inches(2.6), y + Inches(0.75), Inches(1.7), Inches(0.9), "Agent", "synthesizes and edits", accent=b.GOLD)
        card(slide, x + Inches(4.85), y + Inches(0.75), Inches(1.7), Inches(0.9), "Notion MCP", "read/write workspace", accent=b.ACCENT)
        arrow(slide, x + Inches(2.06), y + Inches(1.06), Inches(0.5), Inches(0.2), color=b.GRID)
        arrow(slide, x + Inches(4.32), y + Inches(1.06), Inches(0.5), Inches(0.2), color=b.GRID)
        rounded(slide, x + Inches(1.1), y + Inches(2.55), Inches(4.8), Inches(0.78), b.SOFT, line=b.GRID)
        label(slide, "Docs become agent-operable state",
              x + Inches(1.28), y + Inches(2.66), Inches(4.45), Inches(0.42), size=13.0, color=b.NAVY, bold=True)
        return

    if kind == "sheets":
        rows = ["Lead queue", "Invoice exceptions", "Content pipeline", "Ops forecast"]
        for i, row in enumerate(rows):
            d.rect(slide, x + Inches(0.6), y + Inches(0.52 + i * 0.52), Inches(3.0), Inches(0.42),
                   b.LIGHT_TEAL if i % 2 == 0 else b.SOFT)
            d.text(slide, row, x + Inches(0.78), y + Inches(0.58 + i * 0.52), Inches(2.2), Inches(0.22),
                   size=9.4, color=b.INK, font=b.FONT, shrink=True)
        card(slide, x + Inches(4.2), y + Inches(0.9), Inches(1.9), Inches(0.9), "Agent actions", "filter, update, explain, trigger", accent=b.TEAL)
        arrow(slide, x + Inches(3.72), y + Inches(1.3), Inches(0.45), Inches(0.2), color=b.GOLD)
        label(slide, "Rows become action surfaces", x + Inches(0.78), y + Inches(3.18),
              Inches(5.4), Inches(0.44), size=13.6, color=b.ACCENT, bold=True)
        return

    if kind == "legacy":
        names = ["Design suites", "Whiteboards", "Spreadsheets", "CRMs", "Finance ops", "Research tools"]
        for i, name in enumerate(names):
            card(slide, x + Inches(0.2 + (i % 3) * 2.12), y + Inches(0.45 + (i // 3) * 1.42),
                 Inches(1.76), Inches(0.92), name, "high-context surface", accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL, b.TEAL, b.GOLD][i],
                 title_size=12.1, body_size=9.6)
        label(slide, "Manual UI steps are the wedge", x + Inches(0.55),
              y + Inches(3.48), Inches(5.8), Inches(0.44), size=13.6, color=b.NAVY, bold=True)
        return

    if kind == "matrix":
        d.rect(slide, x + Inches(0.7), y + Inches(0.55), Inches(5.1), Inches(3.1), b.SOFT)
        d.rect(slide, x + Inches(3.25), y + Inches(0.55), Inches(2.55), Inches(1.55), b.LIGHT_TEAL)
        d.rect(slide, x + Inches(3.25), y + Inches(2.1), Inches(2.55), Inches(1.55), b.WHITE)
        d.rect(slide, x + Inches(0.7), y + Inches(2.1), Inches(2.55), Inches(1.55), b.WHITE)
        d.rect(slide, x + Inches(0.7), y + Inches(0.55), Inches(5.1), Inches(0.04), b.GRID)
        d.rect(slide, x + Inches(3.25), y + Inches(0.55), Inches(0.04), Inches(3.1), b.GRID)
        label(slide, "Best wedge", x + Inches(3.55), y + Inches(0.88), Inches(1.75), Inches(0.28), size=13.8, color=b.ACCENT, bold=True)
        label(slide, "Repeatable + visible", x + Inches(3.42), y + Inches(1.25), Inches(2.05), Inches(0.26), size=11.2, color=b.INK)
        label(slide, "Hard to review", x + Inches(1.08), y + Inches(0.94), Inches(1.7), Inches(0.24), size=11.3, color=b.MUTED)
        label(slide, "One-off work", x + Inches(1.08), y + Inches(2.72), Inches(1.7), Inches(0.24), size=11.3, color=b.MUTED)
        label(slide, "Visible state ->", x + Inches(2.05), y + Inches(3.82), Inches(2.2), Inches(0.22), size=10.8, color=b.MUTED)
        return

    if kind == "rubric":
        columns = ["Repeat", "Visible", "Context", "Action", "Risk"]
        rows = ["Design review", "Ops sheet", "Legal drafting"]
        for i, col in enumerate(columns):
            d.rect(slide, x + Inches(1.28 + i * 0.92), y + Inches(0.48), Inches(0.82), Inches(0.34), b.NAVY)
            label(slide, col, x + Inches(1.31 + i * 0.92), y + Inches(0.54), Inches(0.76), Inches(0.12),
                  size=7.8, color=b.WHITE, bold=True)
        for r, row in enumerate(rows):
            label(slide, row, x + Inches(0.1), y + Inches(1.0 + r * 0.72), Inches(1.05), Inches(0.22),
                  size=10.5, color=b.NAVY, bold=True, align=PP_ALIGN.LEFT)
            for c in range(5):
                color = [b.TEAL, b.GOLD, b.ACCENT, b.TEAL, b.CORAL][(r + c) % 5]
                d.rect(slide, x + Inches(1.36 + c * 0.92), y + Inches(0.93 + r * 0.72), Inches(0.55), Inches(0.38), color)
                label(slide, "H" if color != b.CORAL else "M", x + Inches(1.43 + c * 0.92), y + Inches(1.02 + r * 0.72),
                      Inches(0.4), Inches(0.12), size=8.8, color=b.WHITE if color != b.GOLD else b.NAVY, bold=True)
        label(slide, "Score, then test with tasks", x + Inches(1.0), y + Inches(3.48),
              Inches(4.9), Inches(0.44), size=13.1, color=b.ACCENT, bold=True)
        return

    if kind == "build-job":
        for i, (title, text, color) in enumerate([
            ("Job", "What outcome should human + agent complete? ", b.TEAL),
            ("Surface", "What object do they edit together?", b.GOLD),
            ("Agent role", "What judgment and actions belong to the agent?", b.ACCENT),
        ]):
            card(slide, x + Inches(0.34 + i * 2.12), y + Inches(1.24), Inches(1.78), Inches(1.08),
                 title, text.strip(), accent=color, title_size=13.2, body_size=10.6)
            if i < 2:
                arrow(slide, x + Inches(2.14 + i * 2.12), y + Inches(1.62), Inches(0.35), Inches(0.18), color=b.GRID)
        label(slide, "Vague job means vague tools", x + Inches(0.95), y + Inches(3.18),
              Inches(5.05), Inches(0.44), size=13.5, color=b.NAVY, bold=True)
        return

    if kind == "tool-contract":
        rows = [
            ("Entity", "Project, page, asset, task"),
            ("CRUD", "create, read, update, delete"),
            ("Permission", "safe, confirm, blocked"),
            ("Output", "state change, file, event"),
        ]
        for i, (title, text) in enumerate(rows):
            card(slide, x + Inches(0.55), y + Inches(0.45 + i * 0.8), Inches(5.7), Inches(0.55),
                 title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=11.8, body_size=9.8)
        label(slide, "Tool contract = agent API", x + Inches(1.0), y + Inches(3.68),
              Inches(5.0), Inches(0.44), size=13.5, color=b.ACCENT, bold=True)
        return

    if kind == "harness":
        layers = [
            ("Agent host", "Codex / Cursor / custom orchestrator"),
            ("Browser surface", "app visible beside the chat"),
            ("Tool bridge", "MCP / CLI / API / browser automation"),
            ("Audit trail", "logs, diffs, approvals, write-back"),
        ]
        for i, (title, text) in enumerate(layers):
            card(slide, x + Inches(0.72), y + Inches(0.55 + i * 0.84), Inches(5.4), Inches(0.58),
                 title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=12.2, body_size=9.8)
            if i < 3:
                arrow(slide, x + Inches(3.18), y + Inches(1.15 + i * 0.84), Inches(0.32), Inches(0.16), color=b.GRID, rotate=90)
        return

    if kind == "anti":
        items = [
            ("Router agent", "Agent only chooses a hard-coded workflow."),
            ("Silent actions", "State changes without visible review."),
            ("No parity", "User can do things the agent cannot."),
        ]
        for i, (title, text) in enumerate(items):
            card(slide, x + Inches(0.35 + i * 2.05), y + Inches(1.15), Inches(1.72), Inches(1.18),
                 title, text, accent=b.CORAL, title_size=12.5, body_size=10.2)
        label(slide, "Fake agentic = brittle automation", x + Inches(0.62),
              y + Inches(3.23), Inches(5.7), Inches(0.44), size=13.4, color=b.CORAL, bold=True)
        return

    if kind == "checklist":
        items = ["UI parity", "Full CRUD", "Visible tool use", "Approval gates", "Shared state", "Memory write-back"]
        for i, item in enumerate(items):
            d.rect(slide, x + Inches(0.4 + (i % 2) * 3.05), y + Inches(0.55 + (i // 2) * 0.9),
                   Inches(2.55), Inches(0.56), b.LIGHT_TEAL if i % 2 == 0 else b.SOFT)
            label(slide, item, x + Inches(0.65 + (i % 2) * 3.05), y + Inches(0.72 + (i // 2) * 0.9),
                  Inches(2.05), Inches(0.14), size=12.2, color=b.NAVY, bold=True)
        label(slide, "Launch when the agent completes a real job", x + Inches(0.75),
              y + Inches(3.43), Inches(5.55), Inches(0.44), size=13.4, color=b.ACCENT, bold=True)
        return

    if kind == "recommend":
        phases = [
            ("Week 1", "Pick wedge + surface"),
            ("Week 2", "Tool contract + harness"),
            ("Week 3", "Pilot tasks + QA"),
            ("Week 4", "Launch with examples"),
        ]
        for i, (title, text) in enumerate(phases):
            card(slide, x + Inches(0.18 + i * 1.62), y + Inches(1.18), Inches(1.26), Inches(0.95),
                 title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i], title_size=12.4, body_size=9.6)
            if i < 3:
                arrow(slide, x + Inches(1.46 + i * 1.62), y + Inches(1.5), Inches(0.35), Inches(0.16), color=b.GRID)
        label(slide, "One job, one surface, one loop", x + Inches(0.92), y + Inches(3.18),
              Inches(5.0), Inches(0.44), size=13.7, color=b.NAVY, bold=True)
        return

    if kind == "sources":
        sources = [
            ("OpenAI Codex", "Local agent can read, change, and run code; docs include MCP access."),
            ("MCP docs/spec", "Open standard for tools, data, workflows; human-in-loop guidance."),
            ("Notion MCP", "AI tools can read/write workspace based on permissions."),
            ("Cursor docs", "Agent, Rules, MCP, Skills, CLI are first-class docs areas."),
        ]
        for i, (title, text) in enumerate(sources):
            card(slide, x + Inches(0.28 + (i % 2) * 3.22), y + Inches(0.7 + (i // 2) * 1.45),
                 Inches(2.9), Inches(1.02), title, text, accent=[b.TEAL, b.GOLD, b.ACCENT, b.CORAL][i],
                 title_size=12.8, body_size=10.2)
        return

    # fallback
    label(slide, "Diagram intentionally uses native PowerPoint shapes", x + Inches(0.6), y + Inches(1.8),
          Inches(5.6), Inches(0.4), size=14, color=b.NAVY, bold=True)


SLIDES = [
    ("cover",),
    ("summary",),
    ("content", "Why this is happening now", "The enabling pieces have converged into a practical build pattern.",
     "Local coding agents, browser-hosted work surfaces, and MCP-style tool access now let software operate as part of the user's agent workflow instead of as an isolated destination.",
     ["Codex runs locally and can inspect or change project files.", "MCP standardizes connections to data, tools, and workflows.", "Web apps can provide a visible shared surface."],
     "why-now", "The opportunity exists because context, tools, and surfaces can now meet in one workflow."),
    ("section", 1, "The software model is changing", "From app-owned assistants to user-owned agents.", [
        ("Old boundary", "Each app owns its small assistant."),
        ("New workspace", "The user's agent carries context across tools."),
        ("Product shift", "The surface becomes the leverage point."),
    ]),
    ("content", "The old model is app with agent", "AI is bolted onto one product boundary.",
     "Most early AI apps added an assistant inside the existing product. That helps with app-local tasks, but it cannot naturally use the user's broader project context, local files, or other work systems.",
     ["The assistant is scoped to one application's data model.", "Connectors help, but the app still controls the agent boundary.", "The user must repeat context across many isolated assistants."],
     "old-model", "App-owned AI improves one silo; it does not become the user's operating layer."),
    ("content", "The new model is agent with app", "The user's own agent becomes the outer workflow.",
     "Agent-native apps invert the control model. The app is a surface inside the user's agent workspace, while the agent brings local context, tool access, memory, and approvals around it.",
     ["The user's agent can reference files and project history.", "The app exposes operations through MCP, CLI, API, or browser actions.", "The human sees the work and corrects the loop early."],
     "new-model", "The product is no longer only a destination; it becomes a surface agents can operate."),
    ("content", "Agent-native has a precise definition", "It is not the same thing as having a chatbot.",
     "An agent-native app is designed so an external user-owned agent can understand the state, take useful actions, and leave behind improved context. The app does not need to own the whole intelligence layer.",
     ["Context is readable by the agent.", "Actions are invocable through explicit tool contracts.", "State changes are visible and reviewable."],
     "definition", "Agent-native means agent-operable, human-visible, and memory-aware."),
    ("content", "The user buys completed work", "The job is shared between human, agent, and app.",
     "The user is not buying another chat interface. They are buying an outcome: a design explored, a document updated, a sheet reconciled, or a workflow completed with enough visibility to trust the result.",
     ["The human supplies judgment and approval.", "The agent supplies context handling and execution.", "The app supplies the durable work surface."],
     "job", "Frame the product around the completed job, not the presence of AI."),
    ("section", 2, "Architecture makes it real", "A strong agent-native app has four building blocks.", [
        ("Context", "What the agent knows."),
        ("Tools", "What the agent can do."),
        ("Surface", "Where work stays visible."),
    ]),
    ("content", "Four building blocks matter most", "The architecture is simple, but the details matter.",
     "Agent-native architecture combines context, connectivity, collaboration, and compounding memory. Weak products usually miss one of these blocks and collapse back into a chatbot or brittle automation.",
     ["Context prevents repeated prompting.", "Connectivity turns suggestions into actions.", "Compounding memory improves future runs."],
     "four-blocks", "If one block is missing, the agent cannot reliably finish real work."),
    ("content", "Reference architecture is a loop", "The product should be designed as an operating cycle.",
     "The agent reads context, plans work, invokes tools, edits the shared surface, then writes approved outcomes back to memory. That loop is the core system, regardless of the specific domain.",
     ["The agent host is not the whole product.", "The tool layer should be explicit and testable.", "The shared surface makes progress inspectable."],
     "reference", "Design the loop first; feature lists come after."),
    ("content", "Context is the first product layer", "Files and project truth are not optional.",
     "The biggest advantage of bring-your-own-agent software is that the agent can start with the user's actual project materials: local markdown, code, screenshots, brand rules, customer notes, and prior decisions.",
     ["Use predictable folders and filenames.", "Keep durable project context in markdown or structured data.", "Write decisions back after approval."],
     "context", "Context turns the app from a blank surface into a project-aware workspace."),
    ("content", "Tools convert intent into work", "MCP, CLI, API, and browser actions are the execution layer.",
     "Tool access is where agent-native products become useful. The agent needs small, composable primitives that map to real app operations, with safe permissions and clear results.",
     ["Expose read and write capabilities intentionally.", "Prefer complete CRUD for important entities.", "Use confirmation for risky operations."],
     "tools", "Do not hide the app behind one workflow-shaped tool; expose useful primitives."),
    ("content", "The surface must be shared", "Human and agent need to see the same object.",
     "A shared surface can be a design canvas, Notion page, spreadsheet, dashboard, board, or operational console. What matters is that both human and agent work on the same visible state.",
     ["Visual state makes mistakes easy to catch.", "The user can intervene without resetting the workflow.", "The agent can use the surface as feedback."],
     "surface", "Visible collaboration is the trust mechanism."),
    ("content", "Memory makes quality compound", "Approved work should improve the next run.",
     "Agent-native software gets stronger when completed work becomes reusable context: brand choices, page structures, prompt corrections, customer rules, and workflow decisions should feed the next task.",
     ["Capture decisions, not every token.", "Separate durable memory from raw logs.", "Make write-back explicit and reviewable."],
     "memory", "The product gets better when the project remembers what worked."),
    ("content", "Trust needs visible approval gates", "The MCP spec's human-in-loop guidance is product design advice.",
     "External tool use should be legible: users need to know which tools are exposed, when tools are invoked, and when a sensitive action requires confirmation. That is not bureaucracy; it is the interface for trust.",
     ["Show available tools and permissions.", "Make tool invocation visible in the UI.", "Require confirmation for irreversible actions."],
     "trust", "Trust is built through observable actions and reversible loops."),
    ("section", 3, "Where the wedges are", "The strongest opportunities look like familiar tools rebuilt for agents.", [
        ("Design", "Canvas work is visual and reviewable."),
        ("Knowledge", "Docs already hold context."),
        ("Operations", "Sheets and boards already drive action."),
    ]),
    ("content", "Design tools show the pattern", "MagicPath-style workflows make the agent useful without owning all context.",
     "In the transcript's example, the agent reads local design and brand context, operates a design canvas in the browser, then can bring the resulting direction back to the codebase. The design tool is the surface, not the only brain.",
     ["The app provides visual canvas state.", "The agent brings local files and project rules.", "The human reviews the design while it forms."],
     "magicpath", "Design is a strong wedge because output is visible, iterative, and context-heavy."),
    ("content", "Knowledge tools become live systems", "Notion-style workspaces are more than repositories.",
     "When an agent can read and write a knowledge workspace with permission, documents become operational surfaces. The app can hold plans, customer notes, specs, and task boards while the user's agent coordinates updates.",
     ["Notion MCP provides a read/write bridge based on permissions.", "Pages and databases become agent-operable entities.", "The human still owns review and workspace structure."],
     "notion", "Knowledge apps are attractive because the context and work surface already coexist."),
    ("content", "Sheets can become control planes", "Rows are visible, structured, and easy to verify.",
     "Spreadsheets are often the operating layer for small businesses. Agent-native sheet tools can let an agent reconcile, explain, classify, forecast, or trigger follow-up while every row remains inspectable.",
     ["Structured rows make state explicit.", "Humans can review before bulk changes.", "Actions can be scoped to small reversible batches."],
     "sheets", "Operational surfaces are good wedges when every action leaves a visible trace."),
    ("content", "Legacy surfaces are exposed", "The opportunity is not limited to new categories.",
     "Many incumbent tools still assume the human will manually move every object through the interface. Agent-native alternatives can start by targeting a narrow repeated workflow inside those surfaces.",
     ["Prioritize high-context, repeated work.", "Avoid tasks where correctness is invisible.", "Start where existing tools force context switching."],
     "legacy", "The best startup wedge may be a familiar surface redesigned for user-owned agents."),
    ("content", "Pick repeatable visible work", "The wedge should be easy to review before it is trusted.",
     "The strongest opportunities sit where the work happens often and the resulting state is visible. That lets the agent create leverage while the human can still understand and correct the output.",
     ["High repetition creates habit and ROI.", "Visible state reduces trust friction.", "Domain context creates defensibility."],
     "matrix", "If users cannot see whether the agent is right, the wedge is too risky for v1."),
    ("content", "Use a simple wedge rubric", "Score candidates before writing code.",
     "A good first product should score well on repetition, visible state, context richness, actionability, and manageable risk. These are judgment scores, not invented market metrics.",
     ["Favor workflows with repeated review cycles.", "Prefer surfaces where partial progress is useful.", "Deprioritize irreversible high-stakes automation."],
     "rubric", "A narrow but reviewable wedge beats a broad agent that users cannot trust."),
    ("section", 4, "How to build it", "Start with the job, then the surface, then the tool contract.", [
        ("Job", "Define the outcome."),
        ("Surface", "Make work visible."),
        ("Contract", "Expose safe actions."),
    ]),
    ("content", "Step one is the job", "Define what human and agent achieve together.",
     "Start with a concrete job-to-be-done: create a page design, reconcile a list, update a campaign plan, draft a board, or prepare a report. Then specify what the human and agent each contribute.",
     ["Name the completed artifact.", "Define what good looks like.", "List the decisions that require human judgment."],
     "build-job", "The app architecture should follow the job, not the AI demo."),
    ("content", "Step two is the tool contract", "The agent needs complete, safe primitives.",
     "Plan the MCP, CLI, API, and browser operations from day one. For each important entity, define how the agent can read, create, update, delete, validate, and explain changes.",
     ["Map tools to user vocabulary.", "Separate safe actions from approval-required actions.", "Return structured results the agent can reason about."],
     "tool-contract", "The tool contract is the hidden product surface for agents."),
    ("content", "Step three is the harness", "Test inside the actual agent workspace.",
     "Do not test only in your own app. Test inside Codex, Cursor, or the agent host your customer uses. The core question is whether the user's agent can complete the job with your surface and tools.",
     ["Open the app where the agent can see it.", "Run task prompts against real project context.", "Capture logs, diffs, approvals, and write-back."],
     "harness", "If it only works in your demo harness, it is not yet agent-native."),
    ("content", "Avoid fake agent-native patterns", "Some AI products look modern but keep the old architecture.",
     "The most common failure is using an agent as a router for hard-coded workflows. The product may talk like an agent, but it cannot handle new outcomes, share state, or compose tools in unexpected ways.",
     ["Avoid one giant workflow tool.", "Avoid silent actions that users cannot inspect.", "Avoid UI actions with no matching agent capability."],
     "anti", "Agent-native products expose capability; they do not trap intelligence behind brittle flows."),
    ("content", "Launch only after parity checks", "The acceptance test is a real agent task.",
     "A launchable agent-native app should pass a simple test: can the user's agent accomplish a valuable outcome inside the app, using real context, while the human can see and approve the important changes?",
     ["Every UI action needs an agent path.", "Every key entity needs complete CRUD.", "Every risky operation needs an approval gate."],
     "checklist", "Treat parity, visibility, and approval as launch requirements."),
    ("content", "Recommendation: build the narrow loop", "One wedge, one surface, one agent workflow.",
     "The right next move is a narrow v1: choose one repeated job, design the shared surface, expose the smallest complete tool contract, and test it inside the agent host your buyer already uses.",
     ["Do not start with a general-purpose agent.", "Do not over-invest in decorative AI UX.", "Do build the loop users can trust and repeat."],
     "recommend", "Win the first workflow, then let context and tool coverage expand."),
    ("content", "Sources and validation notes", "Claims are grounded in the transcript plus official product/protocol docs.",
     "This deck uses the supplied video transcript for the core thesis and current official documentation for product/protocol claims. No market-size numbers or adoption statistics were invented.",
     ["OpenAI Codex docs and GitHub repo.", "Model Context Protocol docs and tools specification.", "Notion MCP and Cursor documentation."],
     "sources", "Source-backed narrative matters more than unsupported hype."),
]


def build():
    total = len(SLIDES)
    for page, spec in enumerate(SLIDES, 1):
        kind = spec[0]
        if kind == "cover":
            cover(page, total)
        elif kind == "summary":
            executive_summary(page, total)
        elif kind == "section":
            _, sec, title, subtitle, bullets = spec
            section_slide(page, total, sec, title, subtitle, bullets)
        elif kind == "content":
            _, title, subtitle, body, bullets, diagram, takeaway = spec
            content_slide(page, total, title, subtitle, body, bullets, diagram, takeaway)
        else:
            raise ValueError(f"Unknown slide kind: {kind}")

    d.save(OUT)
    problems = validate_pptx(OUT)
    if problems:
        raise RuntimeError("\n".join(problems))
    print(OUT)


if __name__ == "__main__":
    build()
