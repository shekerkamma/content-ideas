#!/usr/bin/env python3
"""Build a Canva-style agent-native apps deck with explicit, self-explanatory content."""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

RUN_DIR = Path(__file__).resolve().parent
OUT = RUN_DIR / "agent-native-canva-reviewed.pptx"
ASSET = RUN_DIR / "diagram-assets-redrawn"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PHOTO_W = Inches(4.5)
PANEL_W = SLIDE_W - PHOTO_W

TEAL = RGBColor(0x00, 0xC9, 0xA7)
DARK_NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_2 = RGBColor(0x12, 0x24, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xF1)
ACCENT_TEAL = RGBColor(0x00, 0x9B, 0x82)
DARK_TEAL = RGBColor(0x00, 0x8F, 0x75)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
GOLD = RGBColor(0xFF, 0xB8, 0x00)
CORAL = RGBColor(0xE0, 0x5A, 0x6B)

PHOTOS = {
    "cover": ASSET / "hero-architecture.jpg",
    "old_new": ASSET / "inversion-diagram.jpg",
    "architecture": ASSET / "reference-architecture.jpg",
    "tools": ASSET / "tool-hub.jpg",
    "context": ASSET / "context-stack.jpg",
    "opportunity": ASSET / "opportunity-matrix.jpg",
    "workflow": ASSET / "workflow-loop.jpg",
}

STORY_MAP = {
    "The thesis": "The opportunity is software built for the user's own agent, not just for the user.",
    "Model inversion": "The old model adds AI inside one app. The new model puts the app inside the agent's workflow.",
    "Why now": "Agent hosts and tool connectors already exist. The missing step is product design around them.",
    "What makes something agent-native": "Agent-native software combines context, tools, a shared surface, and memory write-back.",
    "Reference architecture": "The system is a loop: context, agent, tools, surface, and write-back.",
    "Context layer": "Files and project notes stop the agent from relearning the same truth every session.",
    "Tool layer": "MCP, CLI, API, and browser access let the agent do real work.",
    "Shared surface": "The user and the agent need one visible object to edit together.",
    "Human in the loop": "Visible progress and cheap corrections make automation trustworthy.",
    "Where the opportunity is": "The best wedges already have repeat work and visible state.",
    "Design tools": "Design works because the work is visual, reviewable, and context-heavy.",
    "Docs and knowledge work": "Documents become live systems when the agent can write to them.",
    "Sheets and operations": "Sheets become control planes when row-level actions are explicit.",
    "How to pick a wedge": "Pick repeated, visible, and easy-to-verify work.",
    "Build steps 1-3": "Start with the job, define the surface, and support bring-your-own-agent.",
    "Build steps 4-5": "Choose a domain you know and map tools and permissions on day one.",
    "Anti-patterns": "Avoid routers, black boxes, and silent state changes.",
    "Launch checklist": "If the user, agent, and tools cannot see the same state, the wedge is too weak.",
    "Bottom line": "Build for the user's agent. Context compounds, tools stay open, and the work remains visible.",
}


def add_rect(slide, left, top, width, height, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, *, size=8, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.05, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * line_spacing)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, lines, *, size=8.4, color=INK, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = Pt(size * 0.92)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_picture_fit(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.28), SLIDE_H - Inches(0.4), Inches(4.4), Inches(0.24),
             "Agent-Native Apps | Canva-style Brief | June 2026",
             size=7, color=RGBColor(0x88, 0x88, 0x99))
    add_text(slide, SLIDE_W - Inches(1.05), SLIDE_H - Inches(0.4), Inches(0.8), Inches(0.24),
             f"{page_num} / {total}", size=7, color=RGBColor(0x88, 0x88, 0x99), align=PP_ALIGN.RIGHT)


def add_chip(slide, left, top, width, text, fill=ACCENT_TEAL):
    add_rect(slide, left, top, width, Inches(0.34), fill)
    add_text(slide, left + Inches(0.03), top + Inches(0.04), width - Inches(0.06), Inches(0.18),
             text, size=6.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_section_card(slide, left, top, width, height, title, lines, *, fill=SOFT):
    add_rect(slide, left, top, width, height, fill)
    add_rect(slide, left, top, width, Inches(0.34), DARK_NAVY)
    add_text(slide, left + Inches(0.08), top + Inches(0.05), width - Inches(0.16), Inches(0.16),
             title, size=9.4, color=WHITE, bold=True)
    add_bullets(slide, left + Inches(0.08), top + Inches(0.42), width - Inches(0.16), height - Inches(0.48), lines)


def add_callout(slide, left, top, width, text):
    add_rect(slide, left, top, width, Inches(0.56), DARK_NAVY)
    add_text(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.2),
             text, size=11.4, color=WHITE, bold=True, align=PP_ALIGN.LEFT)


def add_line(slide, x1, y1, x2, y2, color=ACCENT_TEAL, width=2):
    # Use a thin rotated rectangle rather than connector XML for better PPT compatibility.
    dx = x2 - x1
    dy = y2 - y1
    length = (dx**2 + dy**2) ** 0.5
    angle = 0 if length == 0 else math.degrees(math.atan2(dy, dx))
    left = x1 + dx / 2 - length / 2
    top = y1 + dy / 2 - Pt(width) / 2
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(length), int(Pt(width)))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.rotation = angle
    return line


def add_strip_bg(slide):
    add_rect(slide, PANEL_W, 0, PHOTO_W, SLIDE_H, SOFT)
    add_rect(slide, PANEL_W + Inches(0.12), Inches(0.12), PHOTO_W - Inches(0.24), SLIDE_H - Inches(0.24), WHITE)


def label_box(slide, left, top, width, height, title, body, fill=SOFT, accent=TEAL, title_size=10, body_size=7.2):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = Pt(1.2)
    add_rect(slide, left, top, Inches(0.08), height, accent)
    add_text(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.16),
             title, size=title_size, color=INK, bold=True)
    add_text(slide, left + Inches(0.12), top + Inches(0.28), width - Inches(0.24), height - Inches(0.32),
             body, size=body_size, color=MUTED)
    return box


def diagram_cover(slide):
    add_strip_bg(slide)
    center_x = PANEL_W + PHOTO_W / 2
    center_y = Inches(3.9)
    add_rect(slide, PANEL_W + Inches(0.32), Inches(0.55), PHOTO_W - Inches(0.64), Inches(0.4), TEAL)
    add_text(slide, PANEL_W + Inches(0.42), Inches(0.62), PHOTO_W - Inches(0.84), Inches(0.18),
             "AGENT-NATIVE", size=9, color=WHITE, bold=True)
    # hub
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.95), center_y - Inches(0.95), Inches(1.9), Inches(1.9))
    outer.fill.solid(); outer.fill.fore_color.rgb = LIGHT_TEAL; outer.line.color.rgb = GOLD; outer.line.width = Pt(2)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.68), center_y - Inches(0.68), Inches(1.36), Inches(1.36))
    inner.fill.solid(); inner.fill.fore_color.rgb = TEAL; inner.line.color.rgb = WHITE; inner.line.width = Pt(1)
    add_text(slide, center_x - Inches(0.44), center_y - Inches(0.12), Inches(0.88), Inches(0.2),
             "Agent", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # four blocks
    blocks = [
        (PANEL_W + Inches(0.35), Inches(1.2), Inches(1.2), Inches(0.62), "Files", "context"),
        (PANEL_W + Inches(2.1), Inches(1.2), Inches(1.2), Inches(0.62), "Memory", "write-back"),
        (PANEL_W + Inches(0.35), Inches(2.5), Inches(1.2), Inches(0.62), "Tools", "MCP / CLI"),
        (PANEL_W + Inches(2.1), Inches(2.5), Inches(1.2), Inches(0.62), "Surface", "canvas / docs"),
        (PANEL_W + Inches(0.35), Inches(4.55), Inches(1.2), Inches(0.62), "Browser", "shared view"),
        (PANEL_W + Inches(2.1), Inches(4.55), Inches(1.2), Inches(0.62), "Approval", "human loop"),
    ]
    for left, top, width, height, title, body in blocks:
        label_box(slide, left, top, width, height, title, body, fill=WHITE, accent=ACCENT_TEAL, title_size=9, body_size=5.8)
        add_line(slide, left + width, top + height/2, center_x - Inches(0.95), center_y, color=ACCENT_TEAL, width=2)
    add_rect(slide, PANEL_W + Inches(0.32), Inches(6.28), PHOTO_W - Inches(0.64), Inches(0.32), DARK_NAVY)
    add_text(slide, PANEL_W + Inches(0.4), Inches(6.34), PHOTO_W - Inches(0.8), Inches(0.15),
             "Build for the user's agent, not just the user.", size=8.4, color=WHITE, bold=True)


def diagram_inversion(slide):
    add_strip_bg(slide)
    add_text(slide, PANEL_W + Inches(0.35), Inches(0.28), Inches(2.9), Inches(0.18),
             "Model inversion", size=11, color=INK, bold=True)
    # top and bottom stacks
    label_box(slide, PANEL_W + Inches(0.35), Inches(0.75), Inches(3.45), Inches(1.2), "App with agent", "AI sits inside one product.\nContext stops at the boundary.", fill=WHITE, accent=CORAL, title_size=10, body_size=6.5)
    label_box(slide, PANEL_W + Inches(0.35), Inches(2.2), Inches(3.45), Inches(1.2), "Agent with app", "The user's agent enters the surface.\nContext and tools stay portable.", fill=WHITE, accent=TEAL, title_size=10, body_size=6.5)
    add_line(slide, PANEL_W + Inches(1.7), Inches(1.98), PANEL_W + Inches(1.7), Inches(2.18), color=GOLD, width=3)
    add_text(slide, PANEL_W + Inches(1.95), Inches(2.0), Inches(0.9), Inches(0.14), "shift", size=8, color=GOLD, bold=True)
    label_box(slide, PANEL_W + Inches(0.35), Inches(4.1), Inches(1.45), Inches(0.8), "Old", "Chatbot", fill=SOFT, accent=CORAL, title_size=9, body_size=6)
    label_box(slide, PANEL_W + Inches(2.0), Inches(4.1), Inches(1.45), Inches(0.8), "New", "Workspace", fill=SOFT, accent=TEAL, title_size=9, body_size=6)
    label_box(slide, PANEL_W + Inches(0.35), Inches(5.2), Inches(3.1), Inches(0.8), "Implication", "The moat moves to context, tool access, and visible collaboration.", fill=WHITE, accent=GOLD, title_size=9, body_size=6.2)


def diagram_layers(slide):
    add_strip_bg(slide)
    ys = [0.9, 1.7, 2.5, 3.3]
    titles = [("Context", "files / notes"), ("Agent loop", "Cursor / Codex"), ("Tools", "MCP / API"), ("Surface", "canvas / docs")]
    for y, (t, b), c in zip(ys, titles, [TEAL, GOLD, CORAL, ACCENT_TEAL]):
        label_box(slide, PANEL_W + Inches(0.33), Inches(y), Inches(3.35), Inches(0.56), t, b, fill=WHITE, accent=c, title_size=9, body_size=6)
    add_line(slide, PANEL_W + Inches(1.98), Inches(1.46), PANEL_W + Inches(1.98), Inches(1.66), color=TEAL, width=3)
    add_line(slide, PANEL_W + Inches(1.98), Inches(2.26), PANEL_W + Inches(1.98), Inches(2.46), color=GOLD, width=3)
    add_line(slide, PANEL_W + Inches(1.98), Inches(3.06), PANEL_W + Inches(1.98), Inches(3.26), color=CORAL, width=3)
    label_box(slide, PANEL_W + Inches(0.33), Inches(4.35), Inches(3.35), Inches(1.05), "Memory write-back", "Approved output returns as context for the next turn.\nThat is what makes quality compound over time.", fill=SOFT, accent=GOLD, title_size=9, body_size=6.2)
    label_box(slide, PANEL_W + Inches(0.33), Inches(5.65), Inches(3.35), Inches(0.75), "Outcome", "A loop, not a silo.", fill=WHITE, accent=TEAL, title_size=9, body_size=6.8)


def diagram_hub(slide):
    add_strip_bg(slide)
    center_x = PANEL_W + Inches(1.95)
    center_y = Inches(3.0)
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.65), center_y - Inches(0.65), Inches(1.3), Inches(1.3))
    hub.fill.solid(); hub.fill.fore_color.rgb = TEAL; hub.line.color.rgb = GOLD; hub.line.width = Pt(2)
    add_text(slide, center_x - Inches(0.38), center_y - Inches(0.1), Inches(0.76), Inches(0.18), "Agent", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    nodes = [
        (PANEL_W + Inches(0.32), Inches(0.9), "Files", "notes / assets", TEAL),
        (PANEL_W + Inches(2.32), Inches(0.9), "Memory", "history / rules", GOLD),
        (PANEL_W + Inches(0.32), Inches(2.1), "Browser", "shared view", ACCENT_TEAL),
        (PANEL_W + Inches(2.32), Inches(2.1), "Tools", "MCP / CLI / API", CORAL),
        (PANEL_W + Inches(0.32), Inches(4.35), "Surface", "canvas / docs", TEAL),
        (PANEL_W + Inches(2.32), Inches(4.35), "Approval", "human loop", GOLD),
    ]
    for left, top, t, b, c in nodes:
        label_box(slide, left, top, Inches(1.2), Inches(0.6), t, b, fill=WHITE, accent=c, title_size=9, body_size=5.8)
        add_line(slide, left + Inches(1.2), top + Inches(0.3), center_x - Inches(0.65), center_y, color=c, width=2)
    label_box(slide, PANEL_W + Inches(0.33), Inches(5.65), Inches(3.35), Inches(0.75), "Definition", "The product is a hub for action, not a chatbot shell.", fill=SOFT, accent=TEAL, title_size=9, body_size=6.6)


def diagram_loop(slide):
    add_strip_bg(slide)
    label_box(slide, PANEL_W + Inches(1.0), Inches(0.75), Inches(1.5), Inches(0.56), "Watch", "see state", fill=WHITE, accent=TEAL, title_size=10, body_size=6)
    label_box(slide, PANEL_W + Inches(2.15), Inches(2.4), Inches(1.55), Inches(0.56), "Correct", "small fix", fill=WHITE, accent=GOLD, title_size=10, body_size=6)
    label_box(slide, PANEL_W + Inches(0.92), Inches(4.0), Inches(1.55), Inches(0.56), "Compound", "learns", fill=WHITE, accent=CORAL, title_size=10, body_size=6)
    add_line(slide, PANEL_W + Inches(1.75), Inches(1.3), PANEL_W + Inches(2.1), Inches(2.35), color=TEAL, width=3)
    add_line(slide, PANEL_W + Inches(2.1), Inches(2.98), PANEL_W + Inches(1.7), Inches(4.0), color=GOLD, width=3)
    add_line(slide, PANEL_W + Inches(1.6), Inches(4.56), PANEL_W + Inches(1.3), Inches(1.28), color=CORAL, width=3)
    label_box(slide, PANEL_W + Inches(0.48), Inches(5.2), Inches(3.0), Inches(0.95), "Principle", "The loop must be visible, reversible, and reusable.", fill=SOFT, accent=TEAL, title_size=9, body_size=6.2)


def diagram_matrix(slide):
    add_strip_bg(slide)
    label_box(slide, PANEL_W + Inches(0.35), Inches(0.72), Inches(1.35), Inches(1.0), "Low value", "hidden\nrare", fill=WHITE, accent=CORAL, title_size=9, body_size=5.8)
    label_box(slide, PANEL_W + Inches(1.95), Inches(0.72), Inches(1.35), Inches(1.0), "Strong wedge", "visible\nrepeatable", fill=WHITE, accent=TEAL, title_size=9, body_size=5.8)
    label_box(slide, PANEL_W + Inches(0.35), Inches(2.0), Inches(1.35), Inches(1.0), "Weak wedge", "repeatable\nbut opaque", fill=WHITE, accent=GOLD, title_size=9, body_size=5.8)
    label_box(slide, PANEL_W + Inches(1.95), Inches(2.0), Inches(1.35), Inches(1.0), "Strong wedge", "shared surface\n + tool access", fill=WHITE, accent=ACCENT_TEAL, title_size=9, body_size=5.8)
    label_box(slide, PANEL_W + Inches(0.35), Inches(3.55), Inches(3.0), Inches(1.0), "What to choose", "Repeated work, visible state, and easy verification.", fill=SOFT, accent=GOLD, title_size=9, body_size=6.4)
    label_box(slide, PANEL_W + Inches(0.35), Inches(4.95), Inches(3.0), Inches(0.9), "Avoid", "One-off or invisible workflows.", fill=WHITE, accent=CORAL, title_size=9, body_size=6.4)


def diagram_default(slide, title):
    add_strip_bg(slide)
    label_box(slide, PANEL_W + Inches(0.35), Inches(0.8), Inches(3.0), Inches(0.85), title, "vector diagram", fill=WHITE, accent=TEAL, title_size=10, body_size=6.2)
    label_box(slide, PANEL_W + Inches(0.35), Inches(2.0), Inches(1.3), Inches(0.62), "Input", "context", fill=WHITE, accent=TEAL, title_size=9, body_size=6)
    label_box(slide, PANEL_W + Inches(1.95), Inches(2.0), Inches(1.3), Inches(0.62), "Process", "agent loop", fill=WHITE, accent=GOLD, title_size=9, body_size=6)
    label_box(slide, PANEL_W + Inches(0.35), Inches(3.0), Inches(1.3), Inches(0.62), "Tools", "actions", fill=WHITE, accent=CORAL, title_size=9, body_size=6)
    label_box(slide, PANEL_W + Inches(1.95), Inches(3.0), Inches(1.3), Inches(0.62), "Output", "visible state", fill=WHITE, accent=ACCENT_TEAL, title_size=9, body_size=6)
    add_line(slide, PANEL_W + Inches(1.65), Inches(2.31), PANEL_W + Inches(1.95), Inches(2.31), color=TEAL, width=2)
    add_line(slide, PANEL_W + Inches(1.0), Inches(2.62), PANEL_W + Inches(1.0), Inches(2.98), color=GOLD, width=2)
    add_line(slide, PANEL_W + Inches(2.6), Inches(2.62), PANEL_W + Inches(2.6), Inches(2.98), color=CORAL, width=2)
    label_box(slide, PANEL_W + Inches(0.35), Inches(4.4), Inches(3.0), Inches(1.2), "Takeaway", "The diagram should explain the system in one glance.", fill=SOFT, accent=TEAL, title_size=9, body_size=6.4)


def title_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)
    add_rect(slide, 0, Inches(2.8), SLIDE_W, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.95), Inches(7.1), Inches(0.3),
             "ARCHITECTURE BRIEF", size=12, color=LIGHT_TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(1.45), Inches(7.3), Inches(0.95),
             "Agent-Native Apps", size=42, color=WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(2.28), Inches(7.3), Inches(0.6),
             "Software designed for the user's own agent.", size=18, color=LIGHT_TEAL)
    add_rect(slide, Inches(0.8), Inches(2.62), Inches(1.55), Inches(0.05), TEAL)
    add_text(slide, Inches(0.8), Inches(3.0), Inches(7.2), Inches(0.35),
             "The opportunity is not another chatbot. It is software that becomes useful to an external agent.",
             size=14, color=WHITE, bold=True)
    chips = ["Cursor", "Codex", "MCP", "Browser", "Memory"]
    x = Inches(0.8)
    for chip in chips:
        add_chip(slide, x, Inches(5.15), Inches(0.92 if chip != "Browser" else 1.05), chip, fill=ACCENT_TEAL)
        x += Inches(1.08 if chip != "Browser" else 1.2)
    add_text(slide, Inches(0.8), Inches(6.15), Inches(4), Inches(0.25),
             "Build for the user's agent, not just the user.", size=11, color=RGBColor(0x88, 0x88, 0x99))
    diagram_cover(slide)
    add_footer(slide, 1, total)


def content_slide(prs, page_num, total, title, subtitle, photo_key, left_sections, right_sections, callout, chips=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, PANEL_W, SLIDE_H, TEAL)
    add_text(slide, Inches(0.4), Inches(0.18), PANEL_W - Inches(0.6), Inches(0.34),
             title, size=21, color=WHITE, bold=True)
    add_rect(slide, Inches(0.4), Inches(0.72), Inches(1.8), Inches(0.03), DARK_NAVY)
    add_text(slide, Inches(0.4), Inches(0.9), PANEL_W - Inches(0.7), Inches(0.22),
             subtitle, size=10.5, color=LIGHT_TEAL)
    story = STORY_MAP.get(title, "")
    if story:
        add_rect(slide, Inches(0.4), Inches(1.06), PANEL_W - Inches(0.7), Inches(0.46), SOFT)
        add_text(slide, Inches(0.48), Inches(1.12), PANEL_W - Inches(0.86), Inches(0.32),
                 story, size=7.6, color=INK)

    add_section_card(slide, Inches(0.4), Inches(1.6), Inches(3.95), Inches(2.38),
                     left_sections[0][0], left_sections[0][1], fill=SOFT)
    add_section_card(slide, Inches(4.42), Inches(1.6), Inches(3.95), Inches(2.38),
                     left_sections[1][0], left_sections[1][1], fill=WHITE)

    add_section_card(slide, Inches(0.4), Inches(4.15), Inches(3.95), Inches(1.36),
                     right_sections[0][0], right_sections[0][1], fill=SOFT)
    add_section_card(slide, Inches(4.42), Inches(4.15), Inches(3.95), Inches(1.36),
                     right_sections[1][0], right_sections[1][1], fill=WHITE)

    if chips:
        x = Inches(0.4)
        for chip in chips:
            width = Inches(max(0.92, min(1.9, len(chip) * 0.09 + 0.45)))
            add_chip(slide, x, Inches(5.7), width, chip, fill=DARK_TEAL)
            x += width + Inches(0.08)

    add_callout(slide, Inches(0.3), Inches(6.05), PANEL_W - Inches(0.5), callout)
    if title == "The thesis":
        diagram_cover(slide)
    elif title == "Model inversion":
        diagram_inversion(slide)
    elif title in {"Why now", "What makes something agent-native", "Reference architecture", "Build steps 1-3", "Build steps 4-5"}:
        diagram_layers(slide)
    elif title in {"Context layer", "Docs and knowledge work"}:
        diagram_hub(slide)
    elif title in {"Tool layer", "Sheets and operations"}:
        diagram_hub(slide)
    elif title in {"Shared surface", "Human in the loop", "Anti-patterns", "Launch checklist"}:
        diagram_loop(slide)
    elif title in {"Where the opportunity is", "How to pick a wedge"}:
        diagram_matrix(slide)
    elif title in {"Design tools"}:
        diagram_hub(slide)
    elif title == "Bottom line":
        diagram_cover(slide)
    else:
        diagram_default(slide, title)
    add_footer(slide, page_num, total)


def closing_slide(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)
    add_rect(slide, 0, Inches(2.85), SLIDE_W, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.95), Inches(7.3), Inches(0.35),
             "BOTTOM LINE", size=12, color=LIGHT_TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(1.4), Inches(7.1), Inches(1.0),
             "Build for the user's agent, not just the user.", size=34, color=WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(2.45), Inches(7.2), Inches(0.6),
             "Design a shared surface, expose the tools, and let the agent improve with context.",
             size=16, color=LIGHT_TEAL, bold=True)
    add_section_card(slide, Inches(0.8), Inches(3.2), Inches(3.9), Inches(1.4), "Next step",
                     ["Pick one repeated job that already lives in a tool-rich workflow.",
                      "Map the shared surface the user and agent can both see.",
                      "Expose the tool path before you polish the UI."],
                     fill=SOFT)
    add_section_card(slide, Inches(4.9), Inches(3.2), Inches(3.9), Inches(1.4), "What to avoid",
                     ["Do not build a router agent.",
                      "Do not hide state changes.",
                      "Do not trap the agent inside a closed app."],
                     fill=WHITE)
    add_section_card(slide, Inches(9.0), Inches(3.2), Inches(3.6), Inches(1.4), "Rule of thumb",
                     ["If the agent cannot see, act, and learn, the product is only AI-colored software."],
                     fill=SOFT)
    add_callout(slide, Inches(0.8), Inches(5.05), Inches(11.7), "The opportunity is software that becomes useful to an external agent.")
    add_text(slide, Inches(0.8), Inches(6.05), Inches(11.6), Inches(0.2),
             "Agent-native apps compound context, not just clicks.", size=12, color=WHITE, bold=True)
    add_picture_fit(slide, PHOTOS["workflow"], SLIDE_W - PHOTO_W, 0, PHOTO_W, SLIDE_H)
    add_footer(slide, page_num, total)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 20

    title_slide(prs, total)

    slides = [
        (2, "The thesis", "The workflow changes when the agent owns the loop", "cover",
         [("What changes", [
             "The user brings the agent. The app becomes a surface, not a silo.",
             "The same workspace must support human review and agent action.",
             "The best product is the one the user's agent can actually finish in."
         ]),
          ("What to build", [
             "Local context files, memory, and assets the agent can re-read.",
             "Tool paths that let the agent act through the app, not around it.",
             "A visible shared surface where the user can watch and correct."
         ])],
         [("Why this matters", [
             "If the agent cannot reach the work, the product is only a chatbot.",
             "If the agent cannot see the same state, collaboration breaks down."
         ]),
          ("Design rule", [
             "Start with the shared surface and the tool surface together.",
             "Treat context as a product asset, not an implementation detail."
         ])],
         "The architecture is a loop, not a router.", ["Context", "Surface", "Tools"]),
        (3, "Model inversion", "App with an agent vs. agent with an app", "old_new",
         [("Old model", [
             "The app embeds a chatbot inside a closed workspace.",
             "Context stops at the app boundary, so the model can assist but not finish.",
             "The user still does the hard part outside the product."
         ]),
          ("New model", [
             "The user's agent owns the context and brings it into the app.",
             "The app becomes one surface among many in a larger workflow.",
             "The user and the agent can act on the same object at the same time."
         ])],
         [("What changes", [
             "Moat shifts to context, tools, workflow design, and memory.",
             "The app must be readable by agents from outside the app."
         ]),
          ("Result", [
             "The app becomes a workspace, not a silo.",
             "The model is useful because it can operate on live state."
         ])],
         "Old model: the app contains the agent. New model: the agent contains the app surface.", ["App AI", "BYO agent"]),
        (4, "Why now", "Three shifts make the category viable today", "tools",
         [("User agents are real", [
             "People are already working inside Cursor, Codex, and similar hosts.",
             "OpenAI's Codex CLI can read, change, and run code locally.",
             "Codex's in-app browser gives a shared view of rendered pages."
         ]),
          ("Apps expose tools", [
             "MCP, CLI, API, and browser access make external action practical.",
             "Notion MCP can read and write workspace content with permissions.",
             "Codex can connect to MCP servers from the CLI or IDE extension."
         ])],
         [("Context is portable", [
             "Files, docs, assets, and memory can live outside one product.",
             "The agent can carry that context across multiple surfaces."
         ]),
          ("Evidence", [
             "OpenAI Codex CLI, Codex app browser, and Notion MCP show the workflow is already here.",
             "The category is not theoretical; the tooling already exists."
         ])],
         "The opportunity opens when agents can bring context into an app and act through tools.", ["Codex CLI", "Notion MCP", "Browser"]),
        (5, "What makes something agent-native", "The product needs four ingredients before it is truly BYO-agent friendly", "architecture",
         [("Required ingredients", [
             "Context: local files, notes, assets, and reusable memory.",
             "Tools: MCP, CLI, API, or browser access that the agent can call.",
             "Surface: a shared workspace both human and agent can edit."
         ]),
          ("What is missing without them", [
             "Without context, the agent repeats instructions.",
             "Without tools, the agent can only suggest work.",
             "Without a shared surface, the work stays invisible."
         ])],
         [("Design outcome", [
             "The agent can act, the user can observe, and both can correct.",
             "The system becomes better because the workspace remembers."
         ]),
          ("Definition", [
             "Agent-native software is software designed for the user's agent.",
             "The app is the surface, not the intelligence container."
         ])],
         "Agent-native means shared context + tool access + visible collaboration + compounding memory.", ["Context", "Tools", "Surface", "Memory"]),
        (6, "Reference architecture", "Context → agent loop → shared surface → outcome → memory", "architecture",
         [("Context layer", [
             "Local files, brand rules, task notes, and project assets.",
             "This is the durable knowledge the agent should re-use.",
             "It makes the product consistent across runs."
         ]),
          ("Execution layer", [
             "Cursor, Codex, or another host runs the agent loop.",
             "The agent calls tools instead of hand-waving the work.",
             "The loop can check, correct, and continue."
         ])],
         [("Shared surface", [
             "The canvas, document, sheet, or browser where work is visible.",
             "This is where humans can review and intervene."
         ]),
          ("Memory / write-back", [
             "The output becomes context for the next turn.",
             "That write-back is what lets quality compound."
         ])],
         "Read left to right: context feeds the agent, tools move it into the app, and write-back turns output into future context.", ["Context", "Agent", "Surface", "Memory"]),
        (7, "Context layer", "The system improves because the agent can re-read what it already learned", "context",
         [("What belongs here", [
             "Markdown notes, decision logs, brand rules, and design specs.",
             "Assets, examples, and the history of what was approved.",
             "Any project truth the agent should not re-discover."
         ]),
          ("Why it matters", [
             "Context removes repetitive prompting and keeps output consistent.",
             "It compounds value over time instead of resetting every session."
         ])],
         [("Operational rule", [
             "Write reusable knowledge into files the agent can read.",
             "Treat those files as product surface area."
         ]),
          ("Result", [
             "The agent gets smarter without shipping code.",
             "The user sees the benefit in fewer corrections."
         ])],
         "Context is the moat because it compounds across tasks and sessions.", ["Files", "Notes", "Assets", "Memory"]),
        (8, "Tool layer", "MCP, CLI, API, and browser access are the unlock", "tools",
         [("Tool paths", [
             "MCP for structured integrations across clients.",
             "CLI for local actions, scripts, and developer workflows.",
             "API or browser access when the app surface itself is the workspace."
         ]),
          ("What parity means", [
             "If a user can do it, the agent needs a path to do it too.",
             "The action should succeed or fail visibly and immediately."
         ])],
         [("Design rule", [
             "Expose atomic primitives, not one giant workflow box.",
             "Keep domain logic in the agent loop, not hidden in the tool."
         ]),
          ("Why it matters", [
             "Tool access is a product decision, not just plumbing.",
             "Parity is what makes the app actually usable by an external agent."
         ])],
         "The best tool surfaces make actions explicit, composable, and auditable.", ["MCP", "CLI", "API", "Browser"]),
        (9, "Shared surface", "The canvas, doc, sheet, or browser is where collaboration becomes legible", "architecture",
         [("Good surfaces", [
             "Canvas for visual work.",
             "Docs for narrative and planning work.",
             "Sheets for structured and tabular work."
         ]),
          ("Design rules", [
             "The user and agent must edit the same object.",
             "The state must be visible, reversible, and reviewable."
         ])],
         [("What to expose", [
             "History, approval, comments, and change tracking.",
             "The current state and the last meaningful action."
         ]),
          ("Why it matters", [
             "A shared surface lets the user watch the agent work.",
             "It turns AI from a hidden process into a collaboration layer."
         ])],
         "The surface is the product; the agent only makes it useful.", ["Canvas", "Docs", "Sheets", "Browser"]),
        (10, "Human in the loop", "The collaboration loop should be part of the product, not a side effect", "workflow",
         [("Watch", [
             "The user can see each change live.",
             "Progress should be visible, not hidden.",
             "The agent should show enough state to be trusted."
         ]),
          ("Correct", [
             "Feedback should be cheap and immediate.",
             "Corrections become new context for the next run."
         ])],
         [("Compounding", [
             "The loop teaches the system what good looks like.",
             "A better feedback loop makes the product better over time."
         ]),
          ("Why it matters", [
             "Automation without visibility is fragile.",
             "Collaboration needs review and correction built in."
         ])],
         "Visible feedback is what turns automation into collaboration.", ["Watch", "Correct", "Compound"]),
        (11, "Where the opportunity is", "The strongest surfaces already carry context and already host repeated work", "opportunity",
         [("Best-fit surfaces", [
             "Canvas tools like Miro because the work is already collaborative.",
             "Docs and knowledge bases because context already lives there.",
             "Sheets and suites because structured data is already the workflow."
         ]),
          ("Why they win", [
             "They already have a shared object to work on.",
             "They already have users who need reviewable state."
         ])],
         [("Opportunity logic", [
             "The best wedge is where the app can hold context and expose actions.",
             "Large incumbents still have room because they already own the workflow."
         ]),
          ("Priority signal", [
             "Look for repeated, visible, auditable work.",
             "That is where agent-native software has room to matter."
         ])],
         "The surface is the product; the agent just amplifies what is already there.", ["Miro", "Notion", "Sheets", "Adobe"]),
        (12, "Design tools", "Cursor + Magic Path shows the pattern in practice", "cover",
         [("How it works", [
             "Open the design app inside the agent host.",
             "Reference local brand files and the design spec.",
             "Let the agent create or revise the page while you watch."
         ]),
          ("Why it matters", [
             "The design stays consistent because the agent sees the local context.",
             "The result can be pulled straight back into the codebase."
         ])],
         [("What this proves", [
             "The app does not need its own embedded agent to be useful.",
             "The app only needs to be an excellent surface for a user's agent."
         ]),
          ("Practical benefit", [
             "Variation is faster because the agent can operate visually and locally.",
             "Feedback is immediate because both sides see the canvas."
         ])],
         "The design tool becomes a shared workspace instead of a separate destination.", ["Cursor", "Magic Path", "design.md"]),
        (13, "Docs and knowledge work", "Notion is already close to agent-native because the workspace is the data model", "context",
         [("Why Notion fits", [
             "It already stores working context, decisions, and task state.",
             "Notion MCP lets agents read and write workspace content.",
             "The browser gives a shared view for humans and agents."
         ]),
          ("How to design it", [
             "Keep source docs, status, and decisions in one place.",
             "Use page history as a memory layer the agent can reuse."
         ])],
         [("Good pattern", [
             "Ask the agent to summarize, update, and link instead of only answer.",
             "Make the output visible in the same workspace."
         ]),
          ("Why it works", [
             "Docs turn into live work surfaces instead of static pages.",
             "The tool becomes more useful because the content stays connected."
         ])],
         "Documents become agent-native when they are writable, searchable, and reusable.", ["Notion MCP", "Workspace", "History"]),
        (14, "Sheets and operations", "Structured tables become a control plane when the agent can safely edit them", "tools",
         [("Why Sheets fits", [
             "The data is already structured and auditable.",
             "Many teams already run operations inside spreadsheets.",
             "Formula workflows and row edits are repeatable."
         ]),
          ("Design rule", [
             "Expose row-level actions, not just whole-sheet edits.",
             "Keep lineage visible so users can trust the changes."
         ])],
         [("What the agent can do", [
             "Reconcile rows, update status, and flag anomalies.",
             "Summarize patterns, then write the result back into the sheet."
         ]),
          ("Why it matters", [
             "The agent becomes useful in a control plane the team already uses.",
             "The workflow gains speed without losing auditability."
         ])],
         "Spreadsheets become agent-native when row actions are explicit and reversible.", ["Rows", "Audit", "Reconcile"]),
        (15, "How to pick a wedge", "The best opportunity already has a repeated job, a visible surface, and a clear output", "opportunity",
         [("Good wedge characteristics", [
             "Repeated: the same job happens again and again.",
             "Visible: the work can be reviewed in a shared surface.",
             "Easy to verify: success is obvious to the user."
         ]),
          ("Bad wedge characteristics", [
             "One-off: no repeatable pattern to compound.",
             "Hidden: the user cannot see what the agent changed."
         ])],
         [("Selection rule", [
             "Start where the workflow already exists.",
             "Do not ask users to invent a new habit before the product works."
         ]),
          ("Why this matters", [
             "The best wedge is already expensive enough to fix.",
             "It is also narrow enough to ship without boiling the ocean."
         ])],
         "If the job is not repeated, visible, and auditable, it is probably not the wedge.", ["Repeated", "Visible", "Auditable"]),
        (16, "Build steps 1-3", "Pick the job, define the surface, and build for BYO agent", "old_new",
         [("1-2: job and surface", [
             "Start with one outcome the user and agent should achieve together.",
             "Choose the shared workspace both sides can edit.",
             "Define the review flow before you define the UI chrome."
         ]),
          ("3: BYO agent", [
             "Expose tools so any agent can use the app.",
             "Do not lock the product to one model or one vendor."
         ])],
         [("Practical implication", [
             "The product design starts with the shared surface and the tool surface.",
             "The UI should reveal work, not hide it."
         ]),
          ("Failure mode", [
             "If the app is just an agent front-end, it is replaceable.",
             "If the app is a workspace, it becomes sticky."
         ])],
         "Start with the job and the surface before you start writing the agent loop.", ["Job", "Surface", "BYO agent"]),
        (17, "Build steps 4-5", "Build from expertise and plan the tool map on day one", "tools",
         [("4-5: expertise and tools", [
             "Choose a domain you already know deeply enough to judge quality.",
             "Scope MCP, CLI, API, and browser access before shipping.",
             "Define permissions, failure states, and write-back behavior early."
         ]),
          ("What to specify now", [
             "Which surfaces are read/write.",
             "What gets logged or audited.",
             "How the user corrects the agent when it is wrong."
         ])],
         [("Why expertise matters", [
             "You need domain truth to know whether the agent is doing good work.",
             "The feedback loop is stronger when you can judge the output."
         ]),
          ("Why tools matter", [
             "Tool access is what makes the product actually usable.",
             "The app becomes agent-native only when the tools are explicit."
         ])],
         "Architecture decisions belong at the start, not after the UI is finished.", ["Expertise", "Tools", "Permissions"]),
        (18, "Anti-patterns", "The common ways teams build something that is not actually agent-native", "workflow",
         [("Do not build", [
             "A router agent that only selects fixed functions.",
             "A workflow-shaped black box that hides judgment.",
             "A product that writes state somewhere the user cannot see."
         ]),
          ("Why they fail", [
             "They reduce the agent to a caller.",
             "They break trust because the work is not legible.",
             "They prevent the user from correcting the outcome."
         ])],
         [("The fix", [
             "Use atomic tools, shared surfaces, and visible state changes.",
             "Keep the collaboration loop in the product."
         ]),
          ("Rule of thumb", [
             "If the code already solved the interesting part, the agent is just a caller.",
             "That is not an agent-native product."
         ])],
         "Legibility beats cleverness. The user should always know what changed.", ["Router", "Black box", "Silent state"]),
        (19, "Launch checklist", "Use these questions before you commit to the wedge", "context",
         [("Go / no-go", [
             "Does the user already work here today?",
             "Can the agent reach the tools?",
             "Can both see the same state?"
         ]),
          ("Confidence check", [
             "Can the user review and correct the output quickly?",
             "Will the context compound after each run?"
         ])],
         [("If the answer is no", [
             "You are asking for behavior change before product value.",
             "The wedge probably needs more architecture work."
         ]),
          ("If the answer is yes", [
             "The product has a real shot at becoming sticky.",
             "The agent can improve without turning into a separate silo."
         ])],
         "If any answer is no, the idea needs more architecture before it needs more code.", ["User works here", "Tools reachable", "Shared state", "Context compounds"]),
        (20, "Bottom line", "The category is software redesigned for the user's agent", "workflow",
         [("What to remember", [
             "The best new AI opportunity is not another chatbot.",
             "It is software that becomes useful to an external agent."
         ]),
          ("What to do next", [
             "Choose one repeated job.",
             "Map the shared surface.",
             "Expose the tools early."
         ])],
         [("Why this wins", [
             "Context compounds instead of resetting.",
             "The user can watch and correct the work.",
             "The app becomes part of the workflow, not a separate destination."
         ]),
          ("Final rule", [
             "Build for the user's agent, not just the user."
         ])],
         "Agent-native apps compound context, not just clicks.", ["User's agent", "Shared surface", "Tool access"]),
    ]

    for page_num, title, subtitle, photo_key, left, right, callout, chips in slides:
        content_slide(prs, page_num, total, title, subtitle, photo_key, left, right, callout, chips=chips)

    prs.save(OUT)
    print(f"Saved {OUT} with {total} slides")


if __name__ == "__main__":
    main()
