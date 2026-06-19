"""Build Founder's Build Stack PPTX from slide data."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "founders-build-stack.pptx")

BG      = RGBColor(0x08, 0x0b, 0x11)
BG_ELEV = RGBColor(0x0f, 0x14, 0x1d)
INK     = RGBColor(0xee, 0xf2, 0xf7)
SOFT    = RGBColor(0xae, 0xb8, 0xc7)
MUTED   = RGBColor(0x69, 0x74, 0x8a)
TEAL    = RGBColor(0x2d, 0xd4, 0xbf)
SKY     = RGBColor(0x38, 0xbd, 0xf8)
MAGENTA = RGBColor(0xe8, 0x79, 0xf9)
AMBER   = RGBColor(0xf6, 0xb9, 0x4b)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # bg
    sh = s.shapes.add_shape(1, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG; sh.line.fill.background()
    # accent bar
    bar = s.shapes.add_shape(1, 0, H - Inches(0.06), W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    return s

def txt(s, text, left, top, width, height, size=Pt(18), color=INK, bold=False,
        italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = size; run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return tb

def kicker(s, t, top=Inches(0.75)):
    txt(s, t.upper(), Inches(0.8), top, Inches(11), Inches(0.35),
        size=Pt(9), color=TEAL, bold=True)

def heading(s, t, top=Inches(1.25), size=Pt(34)):
    txt(s, t, Inches(0.8), top, Inches(11.5), Inches(1.3),
        size=size, color=INK, bold=True, wrap=True)

def bullets(s, items, top=Inches(2.8), size=Pt(17), color=INK):
    tb = s.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(3.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = f"→  {item}"
        run.font.size = size; run.font.color.rgb = color
        p.space_after = Pt(5)

def col_box(s, head, title, items, tag, tag_color, left, top=Inches(2.65), w=Inches(5.5), h=Inches(3.8)):
    box = s.shapes.add_shape(1, left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = BG_ELEV
    box.line.color.rgb = RGBColor(0x25, 0x2d, 0x3d)
    txt(s, head.upper(), left+Inches(0.2), top+Inches(0.15), w-Inches(0.4), Inches(0.3),
        size=Pt(8), color=TEAL, bold=True)
    txt(s, title, left+Inches(0.2), top+Inches(0.5), w-Inches(0.4), Inches(0.45),
        size=Pt(17), color=INK, bold=True)
    tb = s.shapes.add_textbox(left+Inches(0.2), top+Inches(1.05), w-Inches(0.4), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = f"→  {item}"
        run.font.size = Pt(13); run.font.color.rgb = SOFT
        p.space_after = Pt(3)
    txt(s, tag.upper(), left+Inches(0.2), top+h-Inches(0.5), w-Inches(0.4), Inches(0.3),
        size=Pt(8), color=tag_color, bold=True)

def two_col(s, lh, lt, li, ltag, ltc, rh, rt, ri, rtag, rtc, top=Inches(2.65)):
    col_box(s, lh, lt, li, ltag, ltc, Inches(0.8), top)
    col_box(s, rh, rt, ri, rtag, rtc, Inches(6.75), top)

def tbl(s, headers, rows, top=Inches(2.4), col_widths=None):
    nc = len(headers)
    cw = col_widths or [Inches(11.5/nc)]*nc
    lx = Inches(0.8); rh = Inches(0.5)
    x = lx
    for h, w in zip(headers, cw):
        b = s.shapes.add_shape(1, x, top, w, Inches(0.38))
        b.fill.solid(); b.fill.fore_color.rgb = BG_ELEV; b.line.fill.background()
        txt(s, h.upper(), x+Inches(0.1), top+Inches(0.06), w-Inches(0.2), Inches(0.28),
            size=Pt(8), color=TEAL, bold=True)
        x += w
    for ri, row in enumerate(rows):
        rt = top + Inches(0.38) + ri*rh
        rb = BG_ELEV if ri%2==0 else BG
        x = lx
        for cell, w in zip(row, cw):
            b = s.shapes.add_shape(1, x, rt, w, rh-Inches(0.04))
            b.fill.solid(); b.fill.fore_color.rgb = rb; b.line.fill.background()
            color = TEAL if (str(cell).startswith("$") and cell!="$0") else INK
            if str(cell) in ("$0","free","free tier","$0"):
                color = TEAL
            txt(s, str(cell), x+Inches(0.1), rt+Inches(0.1), w-Inches(0.2), rh-Inches(0.15),
                size=Pt(13), color=color)
            x += w

def pn(s, n):
    txt(s, str(n), Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
        size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)

# ── 1: Cover ──────────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Founder's Build Stack · Product Brief · June 2026", top=Inches(1.1))
txt(s, "Freelance Designer\nClient Portal", Inches(0.8), Inches(1.65), Inches(11.5), Inches(2.4),
    size=Pt(52), color=TEAL, bold=True, wrap=True)
txt(s, "Brief → Feedback → Invoice → Paid. Built in 30 days on Next.js + Supabase + Stripe.",
    Inches(0.8), Inches(4.1), Inches(10), Inches(1.0), size=Pt(20), color=SOFT)
pn(s, 1)

# ── 2: The Problem ────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 01 · The Problem")
heading(s, "Freelancers run client work on duct tape")
bullets(s, [
    "Client briefs arrive in email threads with no structure",
    "Deliverable feedback is scattered across Loom links, Slack DMs, PDF markups",
    "Invoices go out as manual Stripe payment links copy-pasted into Gmail",
    "The 'system' is a folder of forwarded emails and a Notion page nobody reads",
], top=Inches(2.7))
bq = s.shapes.add_textbox(Inches(1.0), Inches(5.1), Inches(11.2), Inches(1.1))
bq.text_frame.word_wrap = True
p = bq.text_frame.paragraphs[0]
run = p.add_run()
run.text = '"I spend 4 hours a week just tracking down approvals and chasing payments. That\'s a full billable day gone."'
run.font.size = Pt(17); run.font.color.rgb = SOFT; run.font.italic = True
bar = s.shapes.add_shape(1, Inches(0.7), Inches(5.1), Inches(0.05), Inches(1.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
pn(s, 2)

# ── 3: Market ─────────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 01 · The Market")
heading(s, "The tools that exist are built for agencies, not solos")
two_col(s,
    "Existing Options", "HoneyBook · Dubsado · 17hats",
    ["$200–$400/mo", "Built for 5-person agencies", "Onboarding takes a week", "80% of features unused by solos"],
    "OVERKILL + OVERPRICED", MAGENTA,
    "The Gap", "$29/mo · 3 features · done",
    ["Brief intake", "Deliverable feedback", "Invoice + payment", "Nothing else"],
    "THE PRODUCT WE'RE BUILDING", TEAL,
)
pn(s, 3)

# ── 4: Validation ─────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 02 · Validation")
heading(s, "Problem Validator · 25 / 30 · VALIDATED")
scores = [("Problem realness", 9, 10), ("Solution fit", 8, 10), ("Buying signal", 8, 10)]
for i, (label, score, max_) in enumerate(scores):
    top = Inches(2.7) + i * Inches(0.75)
    txt(s, label, Inches(0.8), top, Inches(3.5), Inches(0.4), size=Pt(16), color=SOFT)
    bar_w = Inches(6.0 * score / max_)
    bar = s.shapes.add_shape(1, Inches(4.5), top+Inches(0.13), bar_w, Inches(0.14))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    txt(s, f"{score}/{max_}", Inches(11.0), top, Inches(1.2), Inches(0.4),
        size=Pt(14), color=TEAL, bold=True)
bullets(s, [
    "3 named designers confirmed they would pay right now",
    "Existing alternatives solve the wrong problem at the wrong price",
    "Founder is a freelance designer — dog-fooding from day 1",
], top=Inches(5.0), size=Pt(16))
pn(s, 4)

# ── 5: ICP ────────────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 02 · ICP")
heading(s, "Ideal Customer · 93 / 100")
tbl(s,
    ["Criterion", "Score", "Profile"],
    [
        ["Industry / vertical",   "20/20", "Solo freelance designer"],
        ["Company size + revenue","18/20", "1–2 person studio, $5K–$12K/mo"],
        ["Decision-maker",        "15/15", "They sign their own checks"],
        ["Trigger event",         "13/15", "Just hit 8+ active clients OR lost client to messy comms"],
        ["Budget / WTP",          "13/15", "Pays $29–49/mo if it saves 3+ hrs/week"],
        ["Cultural fit",          "14/15", "Uses Figma, Notion, Linear — early adopter"],
    ],
    top=Inches(2.3),
    col_widths=[Inches(3.2), Inches(1.4), Inches(6.7)],
)
pn(s, 5)

# ── 6: Scope ──────────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 03 · Scope")
heading(s, "What ships in v1 · What doesn't")
two_col(s,
    "TIER 1 · Ships Day 30", "6 features",
    ["Brief intake form (client fills)", "Deliverable file upload",
     "Client feedback on deliverables", "Stripe invoicing + payment",
     "Email notifications (Resend)", "Designer dashboard + client magic-link"],
    "FULL LOOP IN 30 DAYS", TEAL,
    "TIER 2 · Post-launch", "Deferred",
    ["Revision count tracking", "Project timeline / Gantt view",
     "── TIER 3 · Cut ──", "In-app chat (email handles it)"],
    "SCOPE DISCIPLINE = ON-TIME SHIP", MAGENTA,
)
pn(s, 6)

# ── 7: Build vs Buy ───────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 03 · Build vs Buy")
heading(s, "Only pay for what you can't build")
tbl(s,
    ["Component", "Decision", "Tool", "Monthly Cost"],
    [
        ["File storage", "BUILD", "Supabase Storage", "$0"],
        ["Auth",         "BUILD", "Supabase Auth",    "$0"],
        ["Email",        "BUY ($0)", "Resend free tier", "$0"],
        ["Payments",     "BUY (no alt)", "Stripe Connect Express", "2.9% + 30¢"],
        ["Forms",        "BUILD", "Custom React (4hr)", "$0"],
        ["Database",     "BUILD", "Supabase Postgres",  "$0"],
    ],
    top=Inches(2.4),
    col_widths=[Inches(2.3), Inches(2.2), Inches(3.8), Inches(3.0)],
)
txt(s, "Total infra cost for first 3 months: $0 + Stripe transaction fees",
    Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
    size=Pt(15), color=TEAL, bold=True)
pn(s, 7)

# ── 8: Architecture ───────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 04 · Architecture")
heading(s, "Stack · Key Decisions")
stack = ["Next.js 14 App Router", "Supabase Postgres + Auth + Storage", "Vercel", "Stripe Connect Express", "Resend"]
sx = Inches(0.8)
for chip in stack:
    cw = Inches(len(chip) * 0.115 + 0.4)
    b = s.shapes.add_shape(1, sx, Inches(2.55), cw, Inches(0.42))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x08, 0x28, 0x25)
    b.line.color.rgb = RGBColor(0x1a, 0x6b, 0x5f)
    txt(s, chip, sx+Inches(0.15), Inches(2.6), cw-Inches(0.3), Inches(0.32),
        size=Pt(12), color=TEAL)
    sx += cw + Inches(0.15)
bullets(s, [
    "No client accounts — magic token per project (link = access, zero signup friction)",
    "Stripe Connect Express — money flows directly to designer, no money transmitter risk",
    "Signed URLs — deliverable files private; client gets 1-hour URL via API only",
    "Service role key server-only — client components never touch privileged data",
    "Append-only versioning — upload v2, client sees v1 + v2 in sequence",
], top=Inches(3.2), size=Pt(16))
pn(s, 8)

# ── 9: Data Model ─────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 04 · Data Model")
heading(s, "Schema in 5 tables")
tbl(s,
    ["Table", "Key Columns", "Notes"],
    [
        ["designers",    "id, slug, stripe_account_id",        "Linked to Supabase Auth user"],
        ["projects",     "status, magic_token, client_email",  "Status machine: brief_pending → active → paid"],
        ["briefs",       "answers (jsonb)",                    "Client fills; triggers designer notification"],
        ["deliverables", "file_path, version",                 "Supabase Storage; append-only"],
        ["feedback",     "comment, deliverable_id",            "Timestamped; linked to deliverable version"],
        ["invoices",     "amount_cents, status",               "Idempotent Stripe webhook on paid"],
    ],
    top=Inches(2.3),
    col_widths=[Inches(2.3), Inches(4.0), Inches(5.0)],
)
pn(s, 9)

# ── 10: 30-Day Roadmap ────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 05 · Roadmap")
heading(s, "30 Days to Launch")
weeks = [
    ("Week 1 · Days 1–5", "Foundation", ["Supabase schema + RLS", "Auth + magic-link", "Brief intake form"]),
    ("Week 2 · Days 6–12", "Core Loop",  ["File upload (Storage)", "Feedback UI", "Designer dashboard"]),
    ("Week 3 · Days 13–19","Revenue",    ["Stripe Connect Express", "Invoice creation", "Email notifications"]),
    ("Week 4 · Days 20–30","Ship",       ["Mobile QA (375px)", "Prod deploy", "3 beta users onboarded"]),
]
bw = Inches(2.7); bh = Inches(3.5); top = Inches(2.7)
for i, (week, title, items) in enumerate(weeks):
    lx = Inches(0.7) + i * (bw + Inches(0.22))
    b = s.shapes.add_shape(1, lx, top, bw, bh)
    b.fill.solid(); b.fill.fore_color.rgb = BG_ELEV
    b.line.color.rgb = RGBColor(0x25, 0x2d, 0x3d)
    txt(s, week.upper(), lx+Inches(0.18), top+Inches(0.18), bw-Inches(0.36), Inches(0.3),
        size=Pt(8), color=TEAL, bold=True)
    txt(s, title, lx+Inches(0.18), top+Inches(0.55), bw-Inches(0.36), Inches(0.45),
        size=Pt(17), color=INK, bold=True)
    tb = s.shapes.add_textbox(lx+Inches(0.18), top+Inches(1.1), bw-Inches(0.36), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = f"→  {item}"
        run.font.size = Pt(13); run.font.color.rgb = SOFT
        p.space_after = Pt(4)
pn(s, 10)

# ── 11: Checkpoints ───────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 05 · Checkpoints")
heading(s, "5 Gates Between Idea and Launch")
tbl(s,
    ["Day", "Gate", "Pass Criteria"],
    [
        ["5",  "Client access", "Client receives magic link → views project page"],
        ["12", "Full loop",     "Brief → upload → feedback works end-to-end"],
        ["19", "Revenue",       "First real Stripe payment processed"],
        ["26", "Polish",        "All TIER1 features shipped, zero known blockers"],
        ["30", "Launch",        "3 beta users onboarded, first invoice paid"],
    ],
    top=Inches(2.4),
    col_widths=[Inches(1.0), Inches(2.5), Inches(7.8)],
)
pn(s, 11)

# ── 12: GTM ───────────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 06 · Go-to-Market")
heading(s, "How the first 10 customers arrive")
bullets(s, [
    "Day 1: Post in Hexagon Slack + Dribbble Discord — 'built this for myself, want beta?'",
    "Week 1: Founder uses it on their own client work — screenshot the experience",
    "Week 2: 3 beta users onboarded free — get raw feedback before charging",
    "Week 4: First paid plan activated ($29/mo) — social proof for wider launch",
    "Month 2: Product Hunt launch + 'how I built this in 30 days' Twitter thread",
], top=Inches(2.7), size=Pt(17))
bq = s.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.2), Inches(0.85))
bq.text_frame.word_wrap = True
p = bq.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Distribution advantage: you're selling to your own community. Designers trust designers, not SaaS companies."
run.font.size = Pt(16); run.font.color.rgb = SOFT; run.font.italic = True
bar = s.shapes.add_shape(1, Inches(0.7), Inches(5.5), Inches(0.05), Inches(0.85))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
pn(s, 12)

# ── 13: Revenue Model ─────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 06 · Revenue Model")
heading(s, "Simple. Predictable. Scalable.")
two_col(s,
    "Pricing", "$29/mo per designer",
    ["Unlimited projects", "Unlimited clients", "All features included", "No per-client seats"],
    "$348 ARR PER CUSTOMER", TEAL,
    "Path to $10K MRR", "Milestones",
    ["Month 1: 3 beta users (free)", "Month 2: 10 paid → $290/mo",
     "Month 6: 50 paid → $1,450/mo", "Month 12: 350 paid → $10,150/mo"],
    "350 SOLO DESIGNERS = $10K MRR", SKY,
)
pn(s, 13)

# ── 14: AI Roadmap ────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Chapter 07 · AI Roadmap")
heading(s, "Two AI features queued for post-launch")
two_col(s,
    "Auto-Invoice Generator", "Brief → Amount Suggestion",
    ["Reads scope, deadline, budget from brief answers",
     "Suggests invoice amount with confidence range",
     "Model: Claude Haiku — <$0.001/call",
     "Rate limit: 10 suggestions/user/day"],
    "LOW COST · HIGH VALUE", TEAL,
    "Feedback Sentiment Summary", "After Client Comments",
    ["Surfaces: 'Client is happy / has concerns about X'",
     "Saves designer from reading 12 comments before a call",
     "Model: Claude Haiku — <$0.001/call",
     "All 4 AI validation tests: PASS"],
    "POST-LAUNCH · MONTH 2", SKY,
)
pn(s, 14)

# ── 15: Closing ───────────────────────────────────────────────────────────────
s = slide()
kicker(s, "Next Step", top=Inches(1.1))
txt(s, "Ship it in 30 days.", Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.6),
    size=Pt(56), color=TEAL, bold=True)
txt(s, "Schema is designed. Stack is chosen. First week's code is drafted.\nDay 1 starts with  npx create-next-app  and a Supabase project. The rest follows the plan.",
    Inches(0.8), Inches(3.6), Inches(10.5), Inches(1.4), size=Pt(20), color=SOFT, wrap=True)
pn(s, 15)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: 15")
