"""Build SaaS Replacement Audit PPTX from audit data using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT = os.path.join(os.path.dirname(__file__), "saas-replacement-audit.pptx")

# ── Palette (Aurora Glass dark) ──────────────────────────────────────────────
BG      = RGBColor(0x08, 0x0b, 0x11)
BG_ELEV = RGBColor(0x0f, 0x14, 0x1d)
INK     = RGBColor(0xee, 0xf2, 0xf7)
SOFT    = RGBColor(0xae, 0xb8, 0xc7)
MUTED   = RGBColor(0x69, 0x74, 0x8a)
TEAL    = RGBColor(0x2d, 0xd4, 0xbf)
SKY     = RGBColor(0x38, 0xbd, 0xf8)
MAGENTA = RGBColor(0xe8, 0x79, 0xf9)
AMBER   = RGBColor(0xf6, 0xb9, 0x4b)
WHITE   = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    shape.zorder = 0

def accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.06), W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

def kicker(slide, text, top=Inches(1.0)):
    tb = slide.shapes.add_textbox(Inches(0.8), top, Inches(11), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text.upper()
    run.font.size = Pt(9); run.font.color.rgb = TEAL; run.font.bold = True

def heading(slide, text, top=Inches(1.5), size=Pt(36), color=INK):
    tb = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.size = size; run.font.color.rgb = color; run.font.bold = True

def body_text(slide, lines, top=Inches(2.9), color=INK, size=Pt(18)):
    tb = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = f"→  {line}" if not line.startswith("→") else line
        run.font.size = size; run.font.color.rgb = color
        p.space_after = Pt(6)

def stat_box(slide, number, label, left=Inches(4), top=Inches(2.8)):
    tb = slide.shapes.add_textbox(left, top, Inches(5), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = number
    run.font.size = Pt(80); run.font.color.rgb = TEAL; run.font.bold = True
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run(); run2.text = label.upper()
    run2.font.size = Pt(10); run2.font.color.rgb = MUTED; run2.font.bold = True

def two_col(slide, left_head, left_title, left_items, left_tag, left_tag_color,
                    right_head, right_title, right_items, right_tag, right_tag_color,
                    top=Inches(2.7)):
    col_w = Inches(5.5); gap = Inches(0.4)
    for idx, (head, title, items, tag, tag_color, left_x) in enumerate([
        (left_head, left_title, left_items, left_tag, left_tag_color, Inches(0.8)),
        (right_head, right_title, right_items, right_tag, right_tag_color, Inches(0.8)+col_w+gap),
    ]):
        box = slide.shapes.add_shape(1, left_x, top, col_w, Inches(3.6))
        box.fill.solid(); box.fill.fore_color.rgb = BG_ELEV
        box.line.color.rgb = RGBColor(0x25, 0x2d, 0x3d)

        tb = slide.shapes.add_textbox(left_x+Inches(0.2), top+Inches(0.15), col_w-Inches(0.4), Inches(0.3))
        tf = tb.text_frame; p = tf.paragraphs[0]
        run = p.add_run(); run.text = head.upper()
        run.font.size = Pt(8); run.font.color.rgb = TEAL; run.font.bold = True

        tb2 = slide.shapes.add_textbox(left_x+Inches(0.2), top+Inches(0.5), col_w-Inches(0.4), Inches(0.45))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        run2 = p2.add_run(); run2.text = title
        run2.font.size = Pt(18); run2.font.color.rgb = INK; run2.font.bold = True

        tb3 = slide.shapes.add_textbox(left_x+Inches(0.2), top+Inches(1.05), col_w-Inches(0.4), Inches(2.0))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        first = True
        for item in items:
            p3 = tf3.paragraphs[0] if first else tf3.add_paragraph()
            first = False
            run3 = p3.add_run(); run3.text = f"→  {item}"
            run3.font.size = Pt(14); run3.font.color.rgb = SOFT
            p3.space_after = Pt(4)

        tb4 = slide.shapes.add_textbox(left_x+Inches(0.2), top+Inches(3.1), col_w-Inches(0.4), Inches(0.3))
        tf4 = tb4.text_frame; p4 = tf4.paragraphs[0]
        run4 = p4.add_run(); run4.text = tag.upper()
        run4.font.size = Pt(8); run4.font.color.rgb = tag_color; run4.font.bold = True

def table_slide(slide, headers, rows, top=Inches(2.5), col_widths=None):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(11.5/n_cols)] * n_cols
    left = Inches(0.8)
    row_h = Inches(0.52)

    # Header row
    x = left
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        box = slide.shapes.add_shape(1, x, top, cw, Inches(0.4))
        box.fill.solid(); box.fill.fore_color.rgb = BG_ELEV
        box.line.fill.background()
        tb = slide.shapes.add_textbox(x+Inches(0.1), top+Inches(0.05), cw-Inches(0.2), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run(); run.text = h.upper()
        run.font.size = Pt(9); run.font.color.rgb = TEAL; run.font.bold = True
        x += cw

    # Data rows
    for r_idx, row in enumerate(rows):
        row_top = top + Inches(0.4) + r_idx * row_h
        row_bg = BG_ELEV if r_idx % 2 == 0 else BG
        x = left
        for i, (cell, cw) in enumerate(zip(row, col_widths)):
            box = slide.shapes.add_shape(1, x, row_top, cw, row_h - Inches(0.04))
            box.fill.solid(); box.fill.fore_color.rgb = row_bg
            box.line.fill.background()

            color = INK
            if cell.startswith("$") and cell != "$0":
                color = TEAL
            elif cell in ("REPLACE",):
                color = TEAL
            elif cell in ("KEEP",):
                color = MUTED
            elif cell in ("NEGOTIATE",):
                color = AMBER
            elif cell in ("CONSOLIDATE", "AUDIT USAGE"):
                color = MAGENTA

            tb = slide.shapes.add_textbox(x+Inches(0.1), row_top+Inches(0.1), cw-Inches(0.2), row_h-Inches(0.15))
            p = tb.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(cell)
            run.font.size = Pt(13); run.font.color.rgb = color
            if cell in ("REPLACE","KEEP","NEGOTIATE","CONSOLIDATE","AUDIT USAGE"):
                run.font.bold = True
            x += cw

def page_num(slide, n):
    tb = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = str(n)
    run.font.size = Pt(9); run.font.color.rgb = MUTED

# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.0))
tf = tb.text_frame; p = tf.paragraphs[0]
run = p.add_run(); run.text = "SaaS Replacement Audit"
run.font.size = Pt(52); run.font.bold = True; run.font.color.rgb = TEAL

kicker(s, "Founder's Build Stack · Cost Intelligence", top=Inches(1.1))

tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(10), Inches(1.0))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
run2 = p2.add_run(); run2.text = "15 tools · $33,288/yr current spend · $29,500 in 3-year savings identified"
run2.font.size = Pt(20); run2.font.color.rgb = SOFT

page_num(s, 1)

# ── SLIDE 2: The Problem ──────────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 01 · The Problem", top=Inches(0.7))
heading(s, "You're renting infrastructure you should own", top=Inches(1.2), size=Pt(34))
body_text(s, [
    "Average Series A startup: $2,500–$5,000/mo in SaaS tools",
    "Of that, 30–45% is replaceable with owned code in under 30 days",
    "The remainder: negotiable, consolidatable, or simply unaudited",
], top=Inches(2.7))
# Pull quote
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2))
tf = tb.text_frame; p = tf.paragraphs[0]
run = p.add_run(); run.text = '"You don\'t have a SaaS problem. You have a build-vs-buy decision you never made."'
run.font.size = Pt(18); run.font.color.rgb = SOFT; run.font.italic = True
bar = s.shapes.add_shape(1, Inches(0.6), Inches(5.0), Inches(0.05), Inches(1.2))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
page_num(s, 2)

# ── SLIDE 3: Stack Overview ───────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 01 · The Stack")
heading(s, "15 tools audited · $2,774/mo", size=Pt(32))
two_col(s,
    "Tools Reviewed", "15 Tools",
    ["Retool, Airtable, Zapier, Segment", "Intercom, Datadog, Mixpanel, Postman",
     "Loom, Notion, Linear, Calendly", "Algolia, SendGrid, Typeform"],
    "$33,288 / YR", TEAL,
    "Verdict Distribution", "5 Buckets",
    ["REPLACE: 5 tools", "KEEP: 4 tools", "NEGOTIATE: 3 tools",
     "CONSOLIDATE: 1 tool", "AUDIT USAGE: 1 tool"],
    "Action on 10 of 15", MAGENTA,
)
page_num(s, 3)

# ── SLIDE 4: Classification Table ────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 02 · Classification")
heading(s, "5-Bucket Framework", size=Pt(32))
table_slide(s,
    ["Bucket", "Criteria", "Tools"],
    [
        ["REPLACE",       "Build cost < 3-yr SaaS AND complexity ≤ medium", "Retool, Zapier, Airtable, Segment, Typeform"],
        ["KEEP",          "Strategic, integrated, or cheap to ignore",        "Algolia, Linear, Notion, Loom, Calendly"],
        ["NEGOTIATE",     "Overpriced; cheaper alternatives exist",           "Datadog, Intercom, SendGrid"],
        ["CONSOLIDATE",   "Overlapping functionality",                        "Mixpanel → drop, use Segment data"],
        ["AUDIT USAGE",   "Unclear if actively used; no clear owner",         "Postman (free tier likely enough)"],
    ],
    top=Inches(2.4),
    col_widths=[Inches(1.8), Inches(4.8), Inches(4.7)],
)
page_num(s, 4)

# ── SLIDE 5: Top 3 by Savings ─────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 02 · Replace Candidates")
heading(s, "Top 3 by 3-Year Savings", size=Pt(32))
table_slide(s,
    ["Rank", "Tool", "Annual SaaS", "Build Cost", "Break-even", "3-Yr Savings"],
    [
        ["#1", "Retool",  "$7,200", "$8,700", "16.6 mo", "$15,132"],
        ["#2", "Zapier",  "$2,388", "$1,380",  "7.1 mo",  "$6,524"],
        ["#3", "Airtable","$2,880", "$4,350", "21.4 mo",  "$5,183"],
    ],
    top=Inches(2.5),
    col_widths=[Inches(0.8), Inches(1.8), Inches(1.9), Inches(1.9), Inches(2.0), Inches(2.9)],
)
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4))
p = tb.text_frame.paragraphs[0]
run = p.add_run(); run.text = "Intercom reclassified → NEGOTIATE (build cost exceeds 3-yr SaaS at current complexity)"
run.font.size = Pt(13); run.font.color.rgb = MUTED; run.font.italic = True
page_num(s, 5)

# ── SLIDE 6: Plan 1 — Retool ──────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 03 · Replacement Plan 1")
heading(s, "Retool → Custom Next.js Admin Panel", size=Pt(30))
two_col(s,
    "The Build", "Next.js + Supabase + shadcn/ui",
    ["User list + impersonation", "Subscription management",
     "NextAuth admin role guard", "Deploy on existing Vercel project"],
    "2 WEEKS · 1 ENGINEER", TEAL,
    "The Math", "$15,132 saved over 3 years",
    ["3-yr SaaS cost: $23,832", "Build cost: $8,700 one-time",
     "Break-even: month 17", "Hosting: $0 (existing Vercel)"],
    "HIGHEST VALUE BUILD", MAGENTA,
)
page_num(s, 6)

# ── SLIDE 7: Plan 2 — Zapier ──────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 03 · Replacement Plan 2")
heading(s, "Zapier → n8n Self-Hosted", size=Pt(30))
two_col(s,
    "The Migration", "n8n on $5/mo Hetzner VPS",
    ["Docker install: 1 hour", "10 zaps → n8n workflows: 3 days",
     "1-week parallel shadow run", "Cancel Zapier on day 8"],
    "3 DAYS · 0 NEW CODE", TEAL,
    "The Math", "$6,524 saved over 3 years",
    ["3-yr SaaS cost: $7,904", "Total migration cost: $1,380",
     "Break-even: month 8", "Fastest payback in the stack"],
    "FASTEST WIN · 7.1 MO PAYBACK", SKY,
)
page_num(s, 7)

# ── SLIDE 8: Plan 3 — Airtable ────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 03 · Replacement Plan 3")
heading(s, "Airtable → Supabase CRM", size=Pt(30))
two_col(s,
    "The Build", "Supabase Schema + Admin Views",
    ["contacts / companies / deals / activities tables",
     "Pipeline view in the admin panel",
     "CSV import from Airtable export",
     "RLS: sales team sees only their deals"],
    "1 WEEK · SHARED WITH PLAN #1", TEAL,
    "The Math", "$5,183 saved over 3 years",
    ["3-yr SaaS cost: $9,533",
     "Build cost: $4,350 (→ $1,500 with Plan #1 scaffold)",
     "Break-even: month 22 (→ month 7 shared)",
     ""],
    "SYNERGISTIC WITH PLAN #1", MAGENTA,
)
page_num(s, 8)

# ── SLIDE 9: Quick Wins ───────────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 04 · Quick Wins")
heading(s, "Negotiate + Consolidate First (Months 1–3)", size=Pt(28))
body_text(s, [
    "Postman → downgrade to free tier · save $588/yr",
    "Mixpanel → cancel, use Segment data in Supabase queries · save $1,068/yr",
    "Datadog → push for SMB tier or migrate to Vercel observability · save $1,476/yr",
    "Intercom → benchmark Crisp, push for 20% reduction · save $718/yr",
    "SendGrid → migrate to Resend (5× cheaper) or negotiate · save est. $600/yr",
], top=Inches(2.6), size=Pt(17))
stat_box(s, "$4,450", "Quick Win Annual Savings", left=Inches(4.0), top=Inches(5.2))
page_num(s, 9)

# ── SLIDE 10: 12-Month Roadmap ────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 05 · Roadmap")
heading(s, "12-Month Action Plan", size=Pt(32))
table_slide(s,
    ["Month", "Action", "Tool", "Annual Savings"],
    [
        ["1",    "Audit + downgrade to free",     "Postman",            "$588"],
        ["1",    "Cancel (consolidate)",           "Mixpanel",           "$1,068"],
        ["2",    "Negotiate -20% / -30%",          "Intercom + Datadog", "$2,194"],
        ["3–4",  "Migrate to n8n self-hosted",     "Zapier",             "$2,388"],
        ["4–7",  "Build custom admin panel",       "Retool",             "$7,200"],
        ["5–8",  "Build Supabase CRM (with #1)",   "Airtable",           "$2,880"],
        ["9–12", "Build event pipeline",           "Segment",            "$1,440"],
    ],
    top=Inches(2.3),
    col_widths=[Inches(1.2), Inches(4.5), Inches(3.3), Inches(2.3)],
)
page_num(s, 10)

# ── SLIDE 11: Summary ─────────────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Chapter 06 · Summary")
heading(s, "The Numbers", size=Pt(36))
two_col(s,
    "Before", "$33,288 / year",
    ["15 tools, $2,774/mo", "5 redundant / replaceable",
     "3 overpriced", "1 possibly unused"],
    "CURRENT STATE", MAGENTA,
    "After (Year 3)", "$16,970 / year",
    ["3 tools replaced with owned code", "3 negotiated down",
     "1 consolidated, 1 cancelled", "Stack is leaner + you own the moat"],
    "$29,500 SAVED OVER 3 YEARS", TEAL,
)
page_num(s, 11)

# ── SLIDE 12: Closing ─────────────────────────────────────────────────────────
s = blank_slide(prs); bg(s); accent_bar(s)
kicker(s, "Next Action", top=Inches(1.1))
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.6))
tf = tb.text_frame; p = tf.paragraphs[0]
run = p.add_run(); run.text = "Start with Zapier."
run.font.size = Pt(56); run.font.bold = True; run.font.color.rgb = TEAL

tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(10), Inches(1.2))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "7.1-month payback. 3 days of work. n8n self-hosted on a $5/mo VPS.\nShip it this week — then tackle Retool in month 4."
run2.font.size = Pt(20); run2.font.color.rgb = SOFT
page_num(s, 12)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: 12")
