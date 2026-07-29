#!/usr/bin/env python3
"""Build the 34-slide semiconductor startup landscape as a native branded PPTX."""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

from PIL import Image, ImageOps
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


RUN = Path(__file__).resolve().parent
ROOT = RUN.parents[1]
sys.path.insert(0, str(ROOT / "skills" / "branded-pptx-deck" / "scripts"))

from pptxkit import Deck  # noqa: E402


CONTENT_PATH = RUN / os.environ.get("DECK_CONTENT_FILE", "outputs/deck-content.json")
CONTENT = json.loads(CONTENT_PATH.read_text(encoding="utf-8"))
SLIDES = CONTENT["slides"]
SOURCE_BY_ID = {source["id"]: source for source in CONTENT["sources"]}
TOTAL = len(SLIDES)
OUT = RUN / os.environ.get("DECK_OUTPUT_FILE", "semiconductor-startups-2026-draft.pptx")
VISUALS = json.loads((RUN / "research" / "visual-assets.json").read_text(encoding="utf-8"))
ASSET_BY_COMPANY = {item["company"].upper(): item for item in VISUALS["assets"]}
DERIVED_ASSETS = RUN / "assets" / "derived"


def set_notes(slide, data: dict) -> None:
    urls = [SOURCE_BY_ID[item]["path"] for item in data["evidence"]]
    asset = ASSET_BY_COMPANY.get(data.get("company", "").upper())
    if asset:
        urls.append(asset["source_page"])
    narrative = data.get("claim") or data.get("subtitle", "")
    slide.notes_slide.notes_text_frame.text = (
        f"Speaker implication: {narrative}\n\n"
        f"Visible citation: {data.get('citation', '')}\n\n"
        "Evidence URLs:\n- " + "\n- ".join(urls)
    )


def add_site_visual(d: Deck, slide, company_name: str, *, x: float, y: float, w: float, h: float) -> None:
    """Add a normalized, editable picture shape sourced from the company website."""
    asset = ASSET_BY_COMPANY[company_name.upper()]
    source = RUN / asset["local_path"]
    DERIVED_ASSETS.mkdir(parents=True, exist_ok=True)
    destination = DERIVED_ASSETS / f"{source.stem}-profile.jpg"
    target = (1200, 1050)
    contain = company_name.upper() in {"TENSTORRENT", "FRACTILE", "ETCHED", "LIGHTMATTER", "MATX", "XSIGHT LABS"}
    with Image.open(source).convert("RGB") as image:
        if contain:
            normalized = ImageOps.contain(image, target, Image.Resampling.LANCZOS)
            canvas = Image.new("RGB", target, "white")
            canvas.paste(normalized, ((target[0] - normalized.width) // 2, (target[1] - normalized.height) // 2))
        else:
            canvas = ImageOps.fit(image, target, Image.Resampling.LANCZOS, centering=(0.5, 0.5))
        canvas.save(destination, quality=92, optimize=True)
    d.rect(
        slide,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        d.b.WHITE,
        line=d.b.GRID,
        radius=0.06,
        shadow=True,
    )
    slide.shapes.add_picture(
        str(destination),
        Inches(x + 0.08),
        Inches(y + 0.08),
        Inches(w - 0.16),
        Inches(h - 0.56),
    )
    d.rect(slide, Inches(x + 0.08), Inches(y + h - 0.48), Inches(w - 0.16), Inches(0.40), d.b.NAVY)
    d.text(
        slide,
        f"OFFICIAL SITE VISUAL  ·  {company_name.upper()}",
        Inches(x + 0.20),
        Inches(y + h - 0.36),
        Inches(w - 0.40),
        Inches(0.16),
        size=9,
        color=d.b.TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
        shrink=True,
    )


def add_citation(d: Deck, slide, data: dict, page: int, *, dark: bool = False) -> None:
    color = RGBColor(0xA9, 0xB4, 0xC0) if dark else d.b.MUTED
    d.text(
        slide,
        data.get("citation", ""),
        d.M,
        Inches(6.73),
        Inches(11.15),
        Inches(0.24),
        size=9,
        color=color,
        italic=True,
        shrink=True,
    )
    d.footer(slide, page, TOTAL, dark=dark)
    set_notes(slide, data)


def add_kicker(d: Deck, slide, text: str, *, dark: bool = False) -> None:
    d.text(
        slide,
        text,
        d.M,
        Inches(0.30),
        Inches(6.6),
        Inches(0.28),
        size=12,
        color=d.b.TEAL if dark else d.b.ACCENT,
        bold=True,
    )


def add_header(d: Deck, slide, data: dict, *, title_w=None) -> None:
    add_kicker(d, slide, data.get("kicker", data.get("company", "")))
    d.text(
        slide,
        data["title"],
        d.M,
        Inches(0.67),
        title_w or d.CW,
        Inches(0.62),
        size=27,
        color=d.b.NAVY,
        bold=True,
        shrink=True,
    )
    d.rect(slide, d.M, Inches(1.35), Inches(1.25), Inches(0.045), d.b.TEAL)
    d.text(
        slide,
        data.get("subtitle", ""),
        d.M,
        Inches(1.48),
        d.CW,
        Inches(0.36),
        size=13.2,
        color=d.b.MUTED,
        shrink=True,
    )


def card(
    d: Deck,
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    text: str,
    *,
    accent=None,
    fill=None,
    value: str | None = None,
    small: bool = False,
) -> None:
    b = d.b
    accent = accent or b.TEAL
    fill = fill or b.WHITE
    d.rect(
        slide,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        fill,
        line=b.GRID,
        radius=0.08,
        shadow=True,
    )
    d.rect(slide, Inches(x), Inches(y), Inches(0.07), Inches(h), accent)
    d.text(
        slide,
        label,
        Inches(x + 0.22),
        Inches(y + 0.17),
        Inches(w - 0.42),
        Inches(0.24),
        size=12,
        color=accent,
        bold=True,
        shrink=True,
    )
    text_top = y + 0.54
    if value is not None:
        d.text(
            slide,
            value,
            Inches(x + 0.22),
            Inches(y + 0.48),
            Inches(w - 0.42),
            Inches(0.42),
            size=24 if not small else 20,
            color=b.NAVY,
            bold=True,
            shrink=True,
        )
        text_top = y + 1.02
    d.text(
        slide,
        text,
        Inches(x + 0.22),
        Inches(text_top),
        Inches(w - 0.42),
        Inches(h - (text_top - y) - 0.16),
        size=12.2 if small else 13.2,
        color=b.INK,
        shrink=True,
        ls=1.02,
    )


def add_table(
    d: Deck,
    slide,
    headers: list[str],
    rows: list[list[str]],
    *,
    x: float = 0.62,
    y: float = 1.98,
    w: float = 12.1,
    h: float = 4.55,
    font_size: float = 12.0,
) -> None:
    shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = shape.table
    widths = [0.19, 0.26, 0.28, 0.27] if len(headers) == 4 else None
    if widths:
        for idx, share in enumerate(widths):
            table.columns[idx].width = int(Inches(w) * share)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = d.b.NAVY
        cell.text = header
    for row_idx, row in enumerate(rows, 1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = d.b.WHITE if row_idx % 2 else d.b.SOFT
            cell.text = str(value)
    for row_idx in range(len(rows) + 1):
        for col_idx in range(len(headers)):
            cell = table.cell(row_idx, col_idx)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = d.b.FONT
                    run.font.size = Pt(font_size)
                    run.font.bold = row_idx == 0 or col_idx == 0
                    run.font.color.rgb = d.b.WHITE if row_idx == 0 else d.b.INK


def cover(d: Deck, data: dict) -> None:
    b = d.b
    s = d.slide(fill=b.NAVY)
    d.rect(s, Inches(10.55), 0, Inches(2.78), d.H, b.NAVY_2)
    d.rect(s, Inches(10.45), 0, Inches(0.10), d.H, b.TEAL)
    for idx, color in enumerate([b.TEAL, b.GOLD, b.CORAL]):
        d.rect(
            s,
            Inches(11.10 + idx * 0.54),
            Inches(1.05 + idx * 0.44),
            Inches(0.34),
            Inches(4.4 - idx * 0.55),
            color,
            radius=0.08,
        )
    d.text(s, data["kicker"], d.M, Inches(0.85), Inches(8.8), Inches(0.34), size=12, color=b.TEAL, bold=True)
    d.text(s, data["title"], d.M, Inches(1.46), Inches(8.9), Inches(1.55), size=44, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, d.M, Inches(3.32), Inches(1.55), Inches(0.06), b.GOLD)
    d.text(s, data["subtitle"], d.M, Inches(3.62), Inches(8.8), Inches(1.08), size=19, color=b.LIGHT_TEAL, shrink=True)
    d.text(
        s,
        data.get("strapline", f"10 companies  ·  4 bottlenecks  ·  {TOTAL}-slide strategy brief"),
        d.M,
        Inches(5.15),
        Inches(8.8),
        Inches(0.42),
        size=13,
        color=b.GOLD,
        bold=True,
    )
    add_citation(d, s, data, 1, dark=True)


def executive(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data)
    xs = [0.62, 4.57, 8.52]
    colors = [d.b.TEAL, d.b.GOLD, d.b.CORAL]
    for x, col, item in zip(xs, colors, data["columns"]):
        card(d, s, x, 2.08, 3.58, 2.84, item["label"], item["text"], accent=col)
    d.rect(s, Inches(0.62), Inches(5.22), Inches(11.88), Inches(0.90), d.b.NAVY, radius=0.05)
    d.text(s, "BOTTOM LINE", Inches(0.90), Inches(5.40), Inches(1.35), Inches(0.26), size=12, color=d.b.TEAL, bold=True)
    d.text(s, data["takeaway"], Inches(2.25), Inches(5.33), Inches(9.83), Inches(0.43), size=15, color=d.b.WHITE, bold=True, shrink=True)
    add_citation(d, s, data, data["id"])


def market(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    d.rect(s, Inches(0.62), Inches(2.02), Inches(4.20), Inches(3.95), d.b.NAVY, radius=0.08)
    d.text(s, data["metric"], Inches(1.0), Inches(2.60), Inches(3.45), Inches(0.92), size=50, color=d.b.GOLD, bold=True, align=PP_ALIGN.CENTER)
    d.text(s, data["metric_label"], Inches(1.08), Inches(3.70), Inches(3.30), Inches(0.82), size=16, color=d.b.WHITE, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    d.text(s, "GARTNER FORECAST", Inches(1.30), Inches(5.10), Inches(2.86), Inches(0.28), size=12, color=d.b.TEAL, bold=True, align=PP_ALIGN.CENTER)
    for idx, point in enumerate(data["points"], 1):
        y = 2.12 + (idx - 1) * 1.28
        d.rect(s, Inches(5.22), Inches(y), Inches(0.56), Inches(0.56), d.b.TEAL, radius=0.4)
        d.text(s, str(idx), Inches(5.22), Inches(y + 0.10), Inches(0.56), Inches(0.24), size=12, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, point, Inches(6.02), Inches(y - 0.02), Inches(6.10), Inches(0.68), size=15, color=d.b.INK, bold=True, shrink=True)
    add_citation(d, s, data, data["id"])


def landscape(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data)
    colors = [d.b.TEAL, d.b.ACCENT, d.b.GOLD, d.b.CORAL]
    for idx, (lane, color) in enumerate(zip(data["lanes"], colors)):
        y = 2.03 + idx * 1.08
        d.rect(s, Inches(0.62), Inches(y), Inches(2.50), Inches(0.82), color, radius=0.06)
        d.text(s, lane["label"], Inches(0.84), Inches(y + 0.21), Inches(2.05), Inches(0.28), size=12, color=d.b.NAVY if color != d.b.CORAL else d.b.WHITE, bold=True, shrink=True)
        d.rect(s, Inches(3.30), Inches(y), Inches(4.10), Inches(0.82), d.b.WHITE, line=d.b.GRID, radius=0.04)
        d.text(s, lane["companies"], Inches(3.56), Inches(y + 0.20), Inches(3.60), Inches(0.28), size=14, color=d.b.NAVY, bold=True, shrink=True)
        d.text(s, lane["note"], Inches(7.70), Inches(y + 0.16), Inches(4.45), Inches(0.42), size=13, color=d.b.MUTED, shrink=True)
    add_citation(d, s, data, data["id"])


def four_cards(d: Deck, data: dict, key: str, *, label_key: str = "label", text_key: str = "text") -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data)
    items = data[key]
    colors = [d.b.TEAL, d.b.GOLD, d.b.CORAL, d.b.ACCENT]
    for idx, item in enumerate(items):
        x = 0.62 + (idx % 2) * 6.02
        y = 2.02 + (idx // 2) * 2.16
        label = item[label_key]
        if "n" in item:
            label = f"{item['n']}  {label}"
        body = item.get(text_key, "")
        if "companies" in item:
            body = f"{item['companies']}\n\n{body}"
        card(d, s, x, y, 5.70, 1.86, label, body, accent=colors[idx], small=False)
    add_citation(d, s, data, data["id"])


def capital(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    colors = [d.b.TEAL, d.b.GOLD, d.b.ACCENT, d.b.CORAL] * 2
    for idx, item in enumerate(data["signals"]):
        x = 0.62 + (idx % 4) * 3.06
        y = 2.02 + (idx // 4) * 2.03
        card(d, s, x, y, 2.82, 1.72, item["company"], item["label"], accent=colors[idx], value=item["value"], small=True)
    d.text(s, "Signals of investor conviction—not a valuation or quality ranking.", Inches(0.70), Inches(6.20), Inches(11.8), Inches(0.28), size=12, color=d.b.MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_citation(d, s, data, data["id"])


def thesis(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    colors = [d.b.TEAL, d.b.GOLD, d.b.CORAL]
    for idx, (item, color) in enumerate(zip(data["arrows"], colors)):
        x = 0.72 + idx * 4.10
        d.rect(s, Inches(x), Inches(2.28), Inches(3.55), Inches(2.28), d.b.SOFT, line=d.b.GRID, radius=0.08, shadow=True)
        d.text(s, item["label"], Inches(x + 0.24), Inches(2.55), Inches(3.02), Inches(0.30), size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, item["text"], Inches(x + 0.30), Inches(3.10), Inches(2.90), Inches(0.82), size=16, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER, shrink=True)
        if idx < 2:
            d.text(s, "→", Inches(x + 3.62), Inches(2.98), Inches(0.38), Inches(0.62), size=26, color=d.b.ACCENT, bold=True, align=PP_ALIGN.CENTER)
    d.rect(s, Inches(0.72), Inches(5.02), Inches(11.90), Inches(0.92), d.b.NAVY, radius=0.05)
    d.text(s, data["takeaway"], Inches(1.02), Inches(5.25), Inches(11.30), Inches(0.36), size=15, color=d.b.WHITE, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    add_citation(d, s, data, data["id"])


def section(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.NAVY)
    d.rect(s, 0, 0, Inches(0.18), d.H, d.b.TEAL)
    add_kicker(d, s, data["kicker"], dark=True)
    d.text(s, data["title"], Inches(0.88), Inches(1.72), Inches(10.8), Inches(1.28), size=40, color=d.b.WHITE, bold=True, shrink=True)
    d.rect(s, Inches(0.90), Inches(3.26), Inches(1.75), Inches(0.06), d.b.GOLD)
    d.text(s, data["subtitle"], Inches(0.90), Inches(3.70), Inches(10.30), Inches(0.80), size=20, color=d.b.LIGHT_TEAL, shrink=True)
    add_citation(d, s, data, data["id"], dark=True)


def company(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data, title_w=Inches(9.55))
    d.rect(s, Inches(10.52), Inches(0.28), Inches(2.18), Inches(0.46), d.b.NAVY, radius=0.14)
    d.text(s, data["tag"], Inches(10.62), Inches(0.37), Inches(1.98), Inches(0.22), size=12, color=d.b.TEAL, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    colors = [d.b.TEAL, d.b.GOLD, d.b.ACCENT, d.b.CORAL]
    labels = ["WEDGE", "PROOF", "ICP + GTM", "RISK"]
    source_keys = ["WEDGE", "PROOF", "INSERTION", "RISK"]
    add_site_visual(d, s, data["company"], x=0.62, y=2.02, w=4.02, h=4.02)
    for idx, (label, source_key, color) in enumerate(zip(labels, source_keys, colors)):
        x = 4.88 + (idx % 2) * 3.90
        y = 2.02 + (idx // 2) * 2.08
        card(d, s, x, y, 3.66, 1.82, label, data["profile"][source_key], accent=color, small=True)
    add_citation(d, s, data, data["id"])


def comparison(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    rows = data["rows"]
    font = 12.0 if len(rows) > 5 else 12.5
    add_table(d, s, data["headers"], rows, y=1.98, h=4.55, font_size=font)
    add_citation(d, s, data, data["id"])


def matrix(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    x0, y0, w, h = 1.20, 2.05, 10.70, 4.05
    d.rect(s, Inches(x0), Inches(y0), Inches(w), Inches(h), d.b.SOFT, line=d.b.GRID)
    d.rect(s, Inches(x0 + w / 2), Inches(y0), Inches(0.025), Inches(h), d.b.GRID)
    d.rect(s, Inches(x0), Inches(y0 + h / 2), Inches(w), Inches(0.025), d.b.GRID)
    d.rect(s, Inches(x0 + w / 2), Inches(y0), Inches(w / 2), Inches(h / 2), d.b.LIGHT_TEAL)
    d.text(s, "MORE COMMERCIAL EVIDENCE  →", Inches(4.80), Inches(6.15), Inches(3.85), Inches(0.28), size=12, color=d.b.MUTED, bold=True, align=PP_ALIGN.CENTER)
    d.text(s, "MORE SPECIALIZED  →", Inches(0.14), Inches(3.36), Inches(0.92), Inches(0.68), size=12, color=d.b.MUTED, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    label_offsets = {
        "Fractile": (0.15, 0.18),
        "MatX": (-1.25, -0.27),
        "Tenstorrent": (-0.75, 0.18),
        "Cornelis": (0.15, -0.27),
    }
    for idx, (label, x, y) in enumerate(data["points"]):
        px = x0 + x * w
        py = y0 + (1 - y) * h
        color = d.b.TEAL if x >= 0.55 else d.b.GOLD
        d.rect(s, Inches(px - 0.12), Inches(py - 0.12), Inches(0.24), Inches(0.24), color, radius=0.5)
        dx, dy = label_offsets.get(label, (-0.58, -0.27 if idx % 2 else 0.16))
        d.text(s, label, Inches(px + dx), Inches(py + dy), Inches(1.16), Inches(0.28), size=12, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    add_citation(d, s, data, data["id"])


def economics(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data)
    for idx, item in enumerate(data["models"]):
        y = 2.02 + idx * 0.82
        d.rect(s, Inches(0.62), Inches(y), Inches(2.35), Inches(0.62), d.b.NAVY, radius=0.04)
        d.text(s, item["label"], Inches(0.82), Inches(y + 0.16), Inches(1.95), Inches(0.27), size=12, color=d.b.TEAL, bold=True, align=PP_ALIGN.CENTER, shrink=True)
        d.text(s, item["value"], Inches(3.24), Inches(y + 0.13), Inches(3.30), Inches(0.28), size=13.5, color=d.b.NAVY, bold=True, shrink=True)
        d.text(s, item["risk"], Inches(6.65), Inches(y + 0.13), Inches(5.45), Inches(0.28), size=13, color=d.b.MUTED, shrink=True)
    add_citation(d, s, data, data["id"])


def engagement(d: Deck, data: dict) -> None:
    rows = [[item["icp"], item["proof"], item["offer"], item["partner"]] for item in data["engagements"]]
    data = {**data, "headers": ["ICP", "REQUIRED PROOF", "COMMERCIAL OFFER", "PARTNER ROLE"], "rows": rows}
    comparison(d, data)


def gtm_scorecard(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    for idx, (num, label, text) in enumerate(data["gates"]):
        col = idx % 2
        row = idx // 2
        x = 0.62 + col * 6.02
        y = 2.00 + row * 1.40
        d.rect(s, Inches(x), Inches(y), Inches(0.70), Inches(0.70), d.b.TEAL if col == 0 else d.b.GOLD, radius=0.4)
        d.text(s, num, Inches(x), Inches(y + 0.18), Inches(0.70), Inches(0.26), size=12, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, label, Inches(x + 0.92), Inches(y + 0.04), Inches(2.05), Inches(0.27), size=12, color=d.b.ACCENT, bold=True)
        d.text(s, text, Inches(x + 0.92), Inches(y + 0.38), Inches(4.62), Inches(0.45), size=12.5, color=d.b.INK, shrink=True)
    add_citation(d, s, data, data["id"])


def watchlist(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.SOFT)
    add_header(d, s, data)
    colors = [d.b.GOLD, d.b.CORAL, d.b.TEAL, d.b.ACCENT]
    for idx, item in enumerate(data["watch"]):
        y = 2.02 + idx * 1.02
        d.rect(s, Inches(0.62), Inches(y), Inches(2.55), Inches(0.78), colors[idx], radius=0.05)
        d.text(s, item["stage"], Inches(0.82), Inches(y + 0.21), Inches(2.15), Inches(0.28), size=12, color=d.b.NAVY if idx != 1 else d.b.WHITE, bold=True, align=PP_ALIGN.CENTER, shrink=True)
        d.text(s, item["companies"], Inches(3.48), Inches(y + 0.12), Inches(3.40), Inches(0.30), size=13.5, color=d.b.NAVY, bold=True, shrink=True)
        d.text(s, item["event"], Inches(7.02), Inches(y + 0.12), Inches(5.12), Inches(0.38), size=12.8, color=d.b.MUTED, shrink=True)
    add_citation(d, s, data, data["id"])


def implications(d: Deck, data: dict) -> None:
    four_cards(d, data, "plays")


def risks(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    for idx, (label, text) in enumerate(data["risks"], 1):
        y = 2.01 + (idx - 1) * 0.84
        d.text(s, f"0{idx}", Inches(0.70), Inches(y + 0.10), Inches(0.50), Inches(0.26), size=13, color=d.b.GOLD, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, label, Inches(1.42), Inches(y + 0.06), Inches(1.45), Inches(0.28), size=12, color=d.b.ACCENT, bold=True)
        d.text(s, text, Inches(3.10), Inches(y + 0.04), Inches(9.00), Inches(0.40), size=13.5, color=d.b.INK, shrink=True)
        d.rect(s, Inches(1.38), Inches(y + 0.57), Inches(10.72), Inches(0.018), d.b.GRID)
    add_citation(d, s, data, data["id"])


def conclusion(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.NAVY)
    add_kicker(d, s, data["kicker"], dark=True)
    d.text(s, data["title"], d.M, Inches(0.90), Inches(11.80), Inches(0.94), size=34, color=d.b.WHITE, bold=True, shrink=True)
    d.text(s, data["subtitle"], d.M, Inches(1.90), Inches(11.20), Inches(0.45), size=15, color=d.b.LIGHT_TEAL, shrink=True)
    colors = [d.b.TEAL, d.b.GOLD, d.b.CORAL]
    for idx, (item, color) in enumerate(zip(data["moves"], colors)):
        x = 0.72 + idx * 4.10
        d.rect(s, Inches(x), Inches(2.75), Inches(3.62), Inches(1.48), d.b.NAVY_2, line=color, radius=0.06)
        d.text(s, item["label"], Inches(x + 0.24), Inches(2.98), Inches(3.12), Inches(0.28), size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, item["text"], Inches(x + 0.26), Inches(3.43), Inches(3.08), Inches(0.52), size=13.5, color=d.b.WHITE, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    d.rect(s, Inches(0.72), Inches(4.78), Inches(11.90), Inches(0.94), d.b.TEAL, radius=0.05)
    d.text(s, data["takeaway"], Inches(1.05), Inches(5.04), Inches(11.25), Inches(0.40), size=16, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    add_citation(d, s, data, data["id"], dark=True)


def methodology(d: Deck, data: dict) -> None:
    s = d.slide(fill=d.b.WHITE)
    add_header(d, s, data)
    for idx, line in enumerate(data["method"], 1):
        y = 2.00 + (idx - 1) * 0.78
        d.rect(s, Inches(0.70), Inches(y), Inches(0.52), Inches(0.52), d.b.TEAL, radius=0.4)
        d.text(s, str(idx), Inches(0.70), Inches(y + 0.12), Inches(0.52), Inches(0.25), size=12, color=d.b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, line, Inches(1.48), Inches(y + 0.05), Inches(10.75), Inches(0.35), size=13.5, color=d.b.INK, shrink=True)
    d.rect(s, Inches(0.70), Inches(5.30), Inches(11.60), Inches(0.70), d.b.SOFT, line=d.b.GRID, radius=0.04)
    d.text(
        s,
        "Primary domains: crn.com · gartner.com · axelera.ai · tenstorrent.com · nextsilicon.com · "
        "fractile.ai · etched.com · d-matrix.ai · cornelis.com · lightmatter.co · matx.com · xsightlabs.com",
        Inches(0.94),
        Inches(5.51),
        Inches(11.10),
        Inches(0.28),
        size=12,
        color=d.b.MUTED,
        align=PP_ALIGN.CENTER,
        shrink=True,
    )
    add_citation(d, s, data, data["id"])


def build() -> None:
    d = Deck(footer="Semiconductor Startup Landscape · 2026")
    handlers = {
        "cover": cover,
        "executive": executive,
        "market": market,
        "landscape": landscape,
        "lens": lambda deck, data: four_cards(deck, data, "steps"),
        "capital": capital,
        "thesis": thesis,
        "business-models": lambda deck, data: four_cards(deck, data, "models"),
        "icp-map": lambda deck, data: four_cards(deck, data, "icps", text_key="need"),
        "gtm-motions": lambda deck, data: four_cards(deck, data, "motions"),
        "section": section,
        "company": company,
        "comparison": comparison,
        "matrix": matrix,
        "buyer-priority": comparison,
        "economics": economics,
        "engagement": engagement,
        "gtm-scorecard": gtm_scorecard,
        "watchlist": watchlist,
        "implications": implications,
        "risks": risks,
        "conclusion": conclusion,
        "methodology": methodology,
    }
    for data in SLIDES:
        handlers[data["layout"]](d, data)
    if d.n != TOTAL or TOTAL < 30:
        raise RuntimeError(f"expected at least 30 slides, built {d.n} of {TOTAL}")
    d.save(OUT)


if __name__ == "__main__":
    build()
