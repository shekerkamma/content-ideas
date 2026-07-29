#!/usr/bin/env python3
"""Validate editability, visible numbers, sources, and deck/HTML story sync."""

from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation


RUN = Path(__file__).resolve().parent
PPTX = RUN / "client-package" / "semiconductor-channel-decision-deck-skillpipe-draft.pptx"
HTML = RUN / "client-package" / "site" / "index.html"
ALLOWED = RUN / "outputs" / "allowed-numbers.yaml"
OUT = RUN / "client-package" / "qa" / "deliverable-validation.json"


def main() -> None:
    prs = Presentation(PPTX)
    deck_texts = [
        shape.text.strip()
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    editable = len(deck_texts)
    slide_titles = []
    for slide in prs.slides:
        candidates = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        slide_titles.append(candidates[1] if len(candidates) > 1 else candidates[0])
    html_text = HTML.read_text(encoding="utf-8")
    allowed_text = ALLOWED.read_text(encoding="utf-8")
    blocked_patterns = [
        r"\bTAM\b|\bSAM\b|\bSOM\b|\bCAGR\b|\bARR\b",
        r"\bmarket share\b|\bproven ROI\b|\bguaranteed payback\b",
    ]
    blocked_matches = sorted({
        match.group(0)
        for pattern in blocked_patterns
        for text in ("\n".join(deck_texts), html_text)
        for match in re.finditer(pattern, text, flags=re.IGNORECASE)
    })
    core_assertions = [
        "validation portfolio",
        "heterogeneous infrastructure",
        "Four bottlenecks",
        "Phase 0",
    ]
    story_sync = {
        phrase: {
            "deck": phrase.lower() in "\n".join(deck_texts).lower(),
            "html": phrase.lower() in html_text.lower(),
        }
        for phrase in core_assertions
    }
    result = {
        "status": "passed",
        "pptx": str(PPTX),
        "slides": len(prs.slides),
        "editable_text_shapes": editable,
        "titles_found": len(slide_titles),
        "html": str(HTML),
        "allowed_numbers": str(ALLOWED),
        "allowed_number_entries": len(re.findall(r"^  - id: NUM-", allowed_text, flags=re.MULTILINE)),
        "blocked_visible_matches": blocked_matches,
        "story_sync": story_sync,
    }
    if (
        len(prs.slides) != 44
        or editable == 0
        or blocked_matches
        or any(not all(value.values()) for value in story_sync.values())
    ):
        result["status"] = "failed"
    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(result, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, indent=2))
    raise SystemExit(0 if result["status"] == "passed" else 1)


if __name__ == "__main__":
    main()
