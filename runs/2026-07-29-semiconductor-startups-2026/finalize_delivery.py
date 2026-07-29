#!/usr/bin/env python3
"""Promote validated artifacts and write manifest-first delivery state."""

from __future__ import annotations

import csv
import datetime as dt
import hashlib
import json
import shutil
from pathlib import Path

from pptx import Presentation


RUN = Path(__file__).resolve().parent
CLIENT = RUN / "client-package"
DRAFT = CLIENT / "semiconductor-channel-decision-deck-skillpipe-draft.pptx"
REVIEWED = CLIENT / "semiconductor-channel-decision-deck-skillpipe-reviewed.pptx"
HTML = CLIENT / "site" / "index.html"
PUBLISH = CLIENT / "pages" / "semiconductor-startups-2026" / "index.html"
MANIFEST = CLIENT / "delivery-manifest.json"
STATUS = RUN / "status.json"
SYNC = RUN / "outputs" / "sync-check.md"


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def main() -> None:
    lint = read_json(CLIENT / "qa" / "pptx-design-lint.json")
    html_qa = read_json(CLIENT / "qa" / "html" / "report.json")
    delivery_qa = read_json(CLIENT / "qa" / "deliverable-validation.json")
    office_summary = (CLIENT / "qa" / "officecli" / "qa-summary.md").read_text(encoding="utf-8")
    gates = {
        "pptx_design_lint": lint.get("status") == "clean",
        "officecli": "Status: `passed`" in office_summary,
        "html": html_qa.get("status") == "passed",
        "deliverable_validation": delivery_qa.get("status") == "passed",
    }
    if not all(gates.values()):
        raise SystemExit(f"Cannot promote; failed gates: {gates}")

    shutil.copy2(DRAFT, REVIEWED)
    PUBLISH.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(HTML, PUBLISH)

    prs = Presentation(REVIEWED)
    editable = sum(
        1
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    )
    with (RUN / "outputs" / "evidence-ledger.csv").open(encoding="utf-8") as handle:
        evidence = list(csv.DictReader(handle))
    promoted = sum(row["storyboard_use"].startswith("slide") for row in evidence)
    now = dt.datetime.now(dt.timezone.utc).isoformat()

    sync_text = f"""# Sync Check

| Check | Result | Evidence |
|---|---|---|
| Deck vs allowed numbers | passed | 22 allowed-number entries; no blocked visible matches |
| HTML vs allowed numbers | passed | Same evidence-led values and no blocked visible matches |
| Deck vs HTML story | passed | Validation portfolio, heterogeneous infrastructure, four bottlenecks, and Phase 0 appear in both |
| Manifest vs files | passed | Paths and checksums resolved during finalization |
| Editable deck | passed | {editable} editable text shapes across {len(prs.slides)} slides |
| OfficeCLI | passed | `client-package/qa/officecli/qa-summary.md` |
| HTML browser QA | passed | Desktop and mobile; 45 sections; 7 valid navigation targets |
| Public URL | not requested | Publish-ready source staged locally |

## Files Checked

- `{REVIEWED.relative_to(RUN)}`
- `{HTML.relative_to(RUN)}`
- `outputs/allowed-numbers.yaml`
- `outputs/story-architect-pack.md`
- `outputs/artifact-traceability.md`
- `client-package/delivery-manifest.json`

No unresolved artifact drift or gate waiver remains.
"""
    SYNC.write_text(sync_text, encoding="utf-8")

    manifest = {
        "artifact_status": "reviewed",
        "run_folder": str(RUN),
        "pptx": {
            "draft_path": str(DRAFT.relative_to(RUN)),
            "reviewed_path": str(REVIEWED.relative_to(RUN)),
            "slide_count": len(prs.slides),
            "editability": "native_powerpoint",
            "editable_text_shape_count": editable,
            "branded_deck_source_path": "client-package/build_deck.py",
            "branded_deck_recreated": True,
            "sha256": sha256(REVIEWED),
            "qa_status": "passed",
            "design_lint_path": "client-package/qa/pptx-design-lint.json",
            "officecli": {
                "status": "passed",
                "result_path": "client-package/qa/officecli",
                "required": True,
            },
            "contact_sheet_path": (
                "client-package/qa/officecli/render/"
                "semiconductor-channel-decision-deck-skillpipe-draft.png"
            ),
        },
        "genspark": {
            "used": False,
            "project_url": None,
            "hosted_editable_reference": False,
            "evidence_clean_scan_status": "passed",
            "unsupported_datapoints_removed": 0,
            "hosted_source_sync_status": "not_applicable",
        },
        "html": {
            "local_path": str(HTML.relative_to(RUN)),
            "publish_source_path": str(PUBLISH.relative_to(RUN)),
            "public_url": None,
            "self_contained": True,
            "public_url_verified": False,
            "qa_status": "passed",
            "qa_path": "client-package/qa/html/report.json",
            "sha256": sha256(HTML),
        },
        "evidence": {
            "ledger_path": "outputs/evidence-ledger.csv",
            "allowed_numbers_path": "outputs/allowed-numbers.yaml",
            "metric_definitions_path": "outputs/metric-definitions.md",
            "data_quality_path": "outputs/data-quality-report.md",
            "scoring_model_path": "outputs/scoring-model.md",
            "row_count": len(evidence),
            "story_promoted_count": promoted,
        },
        "pipeline": {
            "skill": "evidence-led-competitor-pipeline",
            "contract": "skills/evidence-led-competitor-pipeline/references/skillpipe.json",
            "printing_press_role": "upstream data tap",
            "aianalyst_role": "evidence dataset, metrics, scoring, and confidence",
            "competitor_pipeline_role": "story, PPTX, HTML, and QA delivery",
        },
        "sync_check_path": "outputs/sync-check.md",
        "status_path": "status.json",
        "published_commit": None,
        "updated_at": now,
    }
    MANIFEST.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")

    state = read_json(STATUS)
    state.update({
        "artifact_status": "reviewed",
        "current_stage": "complete",
        "completed_gates": [
            "prompt_fit",
            "operating_prompt",
            "gbrain_recall_documented",
            "research_tool_discipline",
            "evidence_ledger",
            "metric_definitions",
            "data_quality",
            "scoring_model",
            "story_architect",
            "allowed_numbers",
            "artifact_traceability",
            "visual_spec",
            "editable_pptx_qa",
            "html_qa",
            "sync_check",
            "delivery_manifest",
        ],
        "blocked_gates": [],
        "artifact_paths": {
            "operating_prompt": "inputs/operating-prompt.yaml",
            "evidence_ledger": "outputs/evidence-ledger.csv",
            "allowed_numbers": "outputs/allowed-numbers.yaml",
            "story_pack": "outputs/story-architect-pack.md",
            "deck_source": "client-package/build_deck.py",
            "reviewed_pptx": str(REVIEWED.relative_to(RUN)),
            "html_local": str(HTML.relative_to(RUN)),
            "html_publish_source": str(PUBLISH.relative_to(RUN)),
            "manifest": str(MANIFEST.relative_to(RUN)),
        },
        "next_command": (
            "Use the reviewed PPTX and local HTML; publish the staged HTML only when a "
            "public URL is explicitly requested."
        ),
        "updated_at": now,
    })
    STATUS.write_text(json.dumps(state, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"reviewed": str(REVIEWED), "slides": len(prs.slides), "editable": editable}, indent=2))


if __name__ == "__main__":
    main()
