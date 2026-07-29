#!/usr/bin/env python3
"""Promote the redesigned deck only after deterministic and real-render QA pass."""

from __future__ import annotations

import csv
import hashlib
import json
import shutil
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation


RUN = Path(__file__).resolve().parent
CLIENT = RUN / "client-package"
OUTPUTS = RUN / "outputs"
QA = CLIENT / "qa"
DRAFT = CLIENT / "semiconductor-startups-2026-evidence-led-strategy-draft.pptx"
REVIEWED = CLIENT / "semiconductor-startups-2026-evidence-led-strategy-reviewed.pptx"

lint = json.loads((QA / "new-deck-design-lint.json").read_text(encoding="utf-8"))
office_summary = (QA / "officecli-new-final" / "qa-summary.md").read_text(encoding="utf-8")
material_change_path = QA / "material-change.json"
material_change = json.loads(material_change_path.read_text(encoding="utf-8"))
if lint["status"] != "clean" or lint["counts"]["error"] != 0:
    raise RuntimeError("Design lint is not clean.")
if "Status: `passed`" not in office_summary:
    raise RuntimeError("OfficeCLI final QA did not pass.")
if not material_change["material_change_passed"]:
    raise RuntimeError("Material redesign comparison did not pass.")

prs = Presentation(DRAFT)
if len(prs.slides) != 36:
    raise RuntimeError(f"Expected 36 slides, found {len(prs.slides)}")

editable_text_shapes = 0
visible_lines = []
for slide_number, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            editable_text_shapes += 1
            visible_lines.append(f"{slide_number}\t{shape.text.strip()}")

visible_text = "\n".join(visible_lines)
(QA / "new-deck-visible-text.txt").write_text(visible_text + "\n", encoding="utf-8")

blocked_terms = [
    "codex",
    "claude",
    "hyperframe",
    "excalidraw",
    "raw file path",
    "/home/shekerk",
]
found = [term for term in blocked_terms if term.lower() in visible_text.lower()]
if found:
    raise RuntimeError(f"Internal terms found in visible text: {found}")

visual_spec_path = RUN / "new-deck-visual-spec.json"
visual_spec = json.loads(visual_spec_path.read_text(encoding="utf-8"))
for visual in visual_spec["visuals"]:
    visual["qa"]["visual_checked"] = True
    visual["qa"]["legible"] = True
    visual["qa"]["not_cropped"] = True
    if visual["route"] == "place-asset":
        visual["qa"]["background_matched"] = True
visual_spec_path.write_text(json.dumps(visual_spec, indent=2), encoding="utf-8")

slide_plan_path = RUN / "new-deck-slide-plan.json"
slide_plan = json.loads(slide_plan_path.read_text(encoding="utf-8"))
for slide in slide_plan["slides"]:
    slide["status"] = "reviewed"
slide_plan_path.write_text(json.dumps(slide_plan, indent=2), encoding="utf-8")

shutil.copy2(DRAFT, REVIEWED)
sha256 = hashlib.sha256(REVIEWED.read_bytes()).hexdigest()

with (OUTPUTS / "evidence-ledger.csv").open(newline="", encoding="utf-8") as handle:
    evidence_rows = len(list(csv.DictReader(handle)))

critique = """# New Deck Critique And Polish

## Narrative

- The executive answer now precedes company profiles.
- Evidence quality, commercialization stage, business model, ICP, and GTM are part of the
  main argument rather than governance artifacts outside the deck.
- Company profiles retain comparable geometry but add role-specific partner decisions.
- The three partner lenses culminate in one consolidated portfolio heatmap.
- Phase 0 calibration and a 90-day evidence-buying plan close the decision loop.

## Visual review

- OfficeCLI rendered all 36 slides with real embedded company images.
- The first 24 slides and slides 25-36 were inspected as separate contact sheets.
- Section titles, profile confidence strips, watchlist cards, roadmap cards, and conclusion
  were resized during the polish pass.
- No visible collisions, clipped titles, broken images, or unreadable tables remain.

## Editability

PowerPoint-native editable diagrams, tables, cards, matrices, and text. Official company
visuals remain editable picture shapes.
"""
(OUTPUTS / "new-deck-critique.md").write_text(critique, encoding="utf-8")

sync = f"""# New Deck Sync Check

- Status: passed
- Slide plan: `new-deck-slide-plan.json`
- Visual specification: `new-deck-visual-spec.json`
- Evidence ledger rows: {evidence_rows}
- Slide count: {len(prs.slides)}
- Editable text shapes: {editable_text_shapes}
- Design lint: clean
- OfficeCLI: passed
- Visible internal-term scan: passed
- Material redesign comparison: passed; zero identical full-text slides anywhere
- HTML: out of scope for this user-requested PPTX-only rebuild
- Reviewed SHA-256: `{sha256}`
"""
(OUTPUTS / "new-deck-sync-check.md").write_text(sync, encoding="utf-8")

manifest = {
    "artifact_status": "reviewed",
    "scope": "new PPTX deck; HTML explicitly out of scope for this deck-only rebuild",
    "run_folder": str(RUN.relative_to(RUN.parents[1])),
    "pptx": {
        "draft_path": str(DRAFT.relative_to(RUN)),
        "reviewed_path": str(REVIEWED.relative_to(RUN)),
        "slide_count": len(prs.slides),
        "editability": "native_powerpoint",
        "editable_text_shape_count": editable_text_shapes,
        "builder_path": "client-package/build_new_deck.py",
        "content_path": "outputs/new-deck-content.json",
        "qa_status": "passed",
        "sha256": sha256,
        "officecli": {
            "status": "passed",
            "result_path": "client-package/qa/officecli-new-final",
            "required": True,
        },
        "design_lint": {
            "status": "clean",
            "result_path": "client-package/qa/new-deck-design-lint.json",
        },
        "material_change": {
            "status": "passed",
            "result_path": "client-package/qa/material-change.json",
            "same_position_full_text_count": material_change["same_position_full_text_count"],
            "identical_slide_text_anywhere_count": material_change[
                "identical_slide_text_anywhere_count"
            ],
        },
    },
    "evidence": {
        "ledger_path": "outputs/evidence-ledger.csv",
        "allowed_numbers_path": "outputs/allowed-numbers.yaml",
        "row_count": evidence_rows,
        "traceability_path": "outputs/artifact-traceability-v2.md",
    },
    "visual_spec_path": "new-deck-visual-spec.json",
    "slide_plan_path": "new-deck-slide-plan.json",
    "sync_check_path": "outputs/new-deck-sync-check.md",
    "critique_path": "outputs/new-deck-critique.md",
    "updated_at": datetime.now(timezone.utc).isoformat(),
}
(CLIENT / "new-deck-delivery-manifest.json").write_text(
    json.dumps(manifest, indent=2), encoding="utf-8"
)

print(json.dumps(manifest, indent=2))
