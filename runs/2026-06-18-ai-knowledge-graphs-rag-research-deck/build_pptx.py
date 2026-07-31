"""Build Aurora-Glass-themed PPTX for 'RAG Has an Amnesia Problem'."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x08, 0x0b, 0x11)
BG_ELEV = RGBColor(0x0f, 0x14, 0x1d)
INK     = RGBColor(0xee, 0xf2, 0xf7)
SOFT    = RGBColor(0xae, 0xb8, 0xc7)
MUTED   = RGBColor(0x69, 0x74, 0x8a)
TEAL    = RGBColor(0x2d, 0xd4, 0xbf)
SKY     = RGBColor(0x38, 0xbd, 0xf8)
MAGENTA = RGBColor(0xe8, 0x79, 0xf9)
AMBER   = RGBColor(0xf6, 0xb9, 0x4b)
LINE    = RGBColor(0x1e, 0x26, 0x35)

# ── Slide size: 16:9 widescreen ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(blank_layout)
    # Dark background rectangle
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.name = "__bg__"
    # Bottom gradient bar (teal → sky → magenta) — approximate with teal
    bar = s.shapes.add_shape(1, 0, H - Pt(4), W, Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    return s

def txbox(slide, text, x, y, w, h, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, italic=False, mono=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if mono:
        run.font.name = "Courier New"
    return tb

def label(slide, text, x, y, w, color=MUTED):
    txbox(slide, text, x, y, w, Inches(0.3), size=8, color=color, mono=True)

def glass_rect(slide, x, y, w, h, radius=False):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = BG_ELEV
    r.line.color.rgb = LINE
    r.line.width = Pt(0.5)
    return r

def folio(slide, n, total=10):
    label(slide, f"{n} / {total}", Inches(0.5), H - Inches(0.4), Inches(2), color=MUTED)
    label(slide, "AI Retrieval Paradigms", W - Inches(2.8), H - Inches(0.4), Inches(2.6), color=MUTED)

def top_labels(slide, left_text, right_text="AI Retrieval Paradigms"):
    label(slide, left_text,  Inches(0.5),          Inches(0.25), Inches(3), color=MUTED)
    label(slide, right_text, W - Inches(3.2), Inches(0.25), Inches(3), color=MUTED)

def kicker(slide, text, x=Inches(0.9), y=Inches(1.0)):
    txbox(slide, text, x, y, Inches(10), Inches(0.4), size=9, color=TEAL, mono=True)

def gradient_title(slide, text, x, y, w, h, size=48):
    """Approximate gradient title — use teal for gradient feel."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    for i, word in enumerate(text.split('\n')):
        if i > 0:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = word
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = TEAL  # single-color stand-in for gradient
    return tb

def bullet_list(slide, items, x, y, w, h, neg_indices=(), pos_indices=()):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        prefix = run = None
        pre = p.add_run()
        pre.font.size = Pt(11)
        pre.font.name = "Courier New"
        if i in neg_indices:
            pre.text = "✕  "
            pre.font.color.rgb = MAGENTA
        elif i in pos_indices:
            pre.text = "✓  "
            pre.font.color.rgb = TEAL
        else:
            pre.text = "→  "
            pre.font.color.rgb = TEAL
        run = p.add_run()
        run.text = item
        run.font.size = Pt(11)
        run.font.color.rgb = INK
    return tb

def col_box(slide, head, title, items, tag, tag_color,
            x, y, w, h, neg=(), pos=()):
    glass_rect(slide, x, y, w, h)
    cy = y + Inches(0.2)
    txbox(slide, head, x + Inches(0.2), cy, w, Inches(0.25), size=8, color=TEAL, mono=True)
    cy += Inches(0.3)
    txbox(slide, title, x + Inches(0.2), cy, w - Inches(0.4), Inches(0.45), size=16, bold=True, color=INK)
    cy += Inches(0.5)
    bullet_list(slide, items, x + Inches(0.2), cy, w - Inches(0.4),
                h - Inches(1.6), neg_indices=neg, pos_indices=pos)
    # tag at bottom
    txbox(slide, tag, x + Inches(0.2), y + h - Inches(0.45),
          w - Inches(0.4), Inches(0.35), size=8, color=tag_color, mono=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "OPENKB")
kicker(s, "KNOWLEDGE RETRIEVAL  ·  ENTERPRISE AI  ·  2024",
       x=Inches(0.9), y=Inches(1.6))
gradient_title(s, "RAG Has an\nAmnesia Problem",
               Inches(0.9), Inches(2.1), Inches(10), Inches(2.8), size=54)
txbox(s,
      "Why compilation-based retrieval outperforms vector similarity search for enterprise knowledge management\n"
      "— and when knowledge graphs win instead.",
      Inches(0.9), Inches(5.0), Inches(9), Inches(1.6), size=14, color=SOFT)
folio(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Thesis: Three paradigms
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 01")
txbox(s, "Three ways to answer a question",
      Inches(0.9), Inches(0.7), Inches(11), Inches(1.2),
      size=34, bold=True, color=INK)
txbox(s,
      "Enterprise AI retrieval has converged on three paradigms. Each bets on a different moment to pay the cost "
      "of understanding — at query time, at schema time, or at ingest time. The moment you choose determines "
      "everything about faithfulness, speed, and auditability.",
      Inches(0.9), Inches(1.85), Inches(10.5), Inches(1.4), size=14, color=INK)

py = Inches(3.4)
for pill_text, pill_color, pill_bg in [
    ("VECTOR RAG — FAST, SHALLOW",        MAGENTA, RGBColor(0x1a, 0x0d, 0x1a)),
    ("KNOWLEDGE GRAPHS — DEEP, COSTLY",   SKY,     RGBColor(0x0a, 0x14, 0x1f)),
    ("COMPILATION — HOLISTIC, AUDITABLE", TEAL,    RGBColor(0x07, 0x18, 0x16)),
]:
    pill = s.shapes.add_shape(1, Inches(0.9), py, Inches(3.6), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = pill_bg
    pill.line.color.rgb = pill_color
    pill.line.width = Pt(0.75)
    txbox(s, pill_text, Inches(1.05), py + Pt(5), Inches(3.4), Inches(0.32),
          size=8, color=pill_color, mono=True)
    py += Inches(0.58)

folio(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Compare: Vector RAG vs OpenKB
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 01")
txbox(s, "Vector RAG vs Compilation",
      Inches(0.9), Inches(0.65), Inches(10), Inches(0.7), size=26, bold=True, color=INK)

COL_W = Inches(5.9)
COL_H = Inches(5.3)
COL_Y = Inches(1.45)

col_box(s,
        head="TRADITIONAL", title="Vector RAG",
        items=["Retrieves isolated K-nearest chunks at query time",
               "No cross-document memory between queries",
               "Contradictions in source material undetected",
               "Synthesis is whatever the chunks happen to contain"],
        tag="AMNESIA PROBLEM", tag_color=MAGENTA,
        x=Inches(0.5), y=COL_Y, w=COL_W, h=COL_H,
        neg=(0,1,2,3))

# divider line
div = s.shapes.add_shape(1, Inches(6.45), COL_Y + Inches(0.2), Pt(2), COL_H - Inches(0.4))
div.fill.solid(); div.fill.fore_color.rgb = TEAL
div.line.fill.background()

col_box(s,
        head="COMPILATION-BASED", title="OpenKB Model",
        items=["LLM synthesises full corpus at ingest — once",
               "Concepts, entities, cross-references baked in",
               "Contradictions flagged at compile time",
               "Human-readable, diffable, auditable wiki artifact"],
        tag="BUILD-STEP RETRIEVAL", tag_color=TEAL,
        x=Inches(6.55), y=COL_Y, w=COL_W, h=COL_H,
        pos=(0,1,2,3))

folio(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Data: "0"
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 01")
txbox(s, "CROSS-REFERENCES BUILT BY TRADITIONAL RAG",
      Inches(3.5), Inches(1.2), Inches(6.3), Inches(0.35), size=8,
      color=MUTED, mono=True, align=PP_ALIGN.CENTER)
txbox(s, "0",
      Inches(3.0), Inches(1.6), Inches(7.3), Inches(3.2),
      size=160, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txbox(s, "EVERY QUERY STARTS COLD",
      Inches(3.5), Inches(4.75), Inches(6.3), Inches(0.35), size=9,
      color=TEAL, mono=True, align=PP_ALIGN.CENTER)
txbox(s,
      "K-nearest chunk retrieval returns passages, not understanding. No memory of how concepts connect "
      "across documents, no contradiction detection, no synthesis that spans the full corpus. Each query is an island.",
      Inches(2.2), Inches(5.2), Inches(8.9), Inches(1.2), size=13, color=INK, align=PP_ALIGN.CENTER)
folio(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Chapter 02
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "OPENKB")
txbox(s, "CHAPTER", Inches(0.9), Inches(1.4), Inches(3), Inches(0.3),
      size=8, color=MUTED, mono=True)
txbox(s, "02", Inches(0.9), Inches(1.7), Inches(5), Inches(2.5),
      size=130, bold=True, color=TEAL)
txbox(s, "Knowledge Graphs\nvs Vector Search",
      Inches(0.9), Inches(4.25), Inches(11), Inches(2.0),
      size=38, bold=True, color=INK)
folio(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Compare: KG vs Vector Search
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 02")
txbox(s, "Knowledge Graphs vs Vector Search",
      Inches(0.9), Inches(0.65), Inches(10), Inches(0.7), size=26, bold=True, color=INK)

col_box(s,
        head="NEO4J · NEPTUNE · GOOGLE KG", title="Knowledge Graphs",
        items=["Multi-hop reasoning via graph traversal",
               "Explicit entity + relationship modelling",
               "High faithfulness for structured domains",
               "Slow ingest — schema, ontology, entity resolution",
               "Brittle on dynamic, fast-moving corpora"],
        tag="STRUCTURED / STATIC DOMAINS", tag_color=SKY,
        x=Inches(0.5), y=COL_Y, w=COL_W, h=COL_H,
        pos=(0,1,2), neg=(3,4))

div = s.shapes.add_shape(1, Inches(6.45), COL_Y + Inches(0.2), Pt(2), COL_H - Inches(0.4))
div.fill.solid(); div.fill.fore_color.rgb = TEAL
div.line.fill.background()

col_box(s,
        head="PINECONE · WEAVIATE · PGVECTOR", title="Vector Search",
        items=["Fast ingest — embed and index, no schema",
               "Low setup complexity, any document type",
               "No entity relationships — flat similarity only",
               "Poor cross-document synthesis",
               "Zero contradiction detection"],
        tag="BROAD / SHALLOW", tag_color=MAGENTA,
        x=Inches(6.55), y=COL_Y, w=COL_W, h=COL_H,
        pos=(0,1), neg=(2,3,4))

folio(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Thesis: PageIndex
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 02")
txbox(s, "The page  is  the concept",
      Inches(0.9), Inches(0.7), Inches(11), Inches(1.2),
      size=38, bold=True, color=INK)
# Highlight "is"
txbox(s, "is",
      Inches(3.55), Inches(0.7), Inches(1.0), Inches(1.2),
      size=38, bold=True, color=TEAL)
txbox(s,
      "PageIndex (OpenKB) indexes compiled concept and entity pages — not raw chunks. "
      "Each page is already a synthesised, cross-referenced understanding of the concept. "
      "The retrieved unit knows everything the corpus says about the topic. "
      "No embedding pipeline, no vector infrastructure required.",
      Inches(0.9), Inches(1.95), Inches(10.5), Inches(1.6), size=14, color=INK)

py = Inches(3.85)
for pill_text, pill_color in [
    ("VECTORLESS RETRIEVAL",        TEAL),
    ("SEMANTIC STRUCTURE AT INGEST", SKY),
    ("NO EMBEDDING COST",            SOFT),
]:
    pill = s.shapes.add_shape(1, Inches(0.9), py, Inches(3.3), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = BG_ELEV
    pill.line.color.rgb = pill_color; pill.line.width = Pt(0.75)
    txbox(s, pill_text, Inches(1.05), py + Pt(5), Inches(3.1), Inches(0.32),
          size=8, color=pill_color, mono=True)
    py += Inches(0.58)

folio(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Compare: GraphRAG vs Compilation
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 03")
txbox(s, "GraphRAG vs Compilation Model",
      Inches(0.9), Inches(0.65), Inches(10), Inches(0.7), size=26, bold=True, color=INK)

col_box(s,
        head="MICROSOFT RESEARCH · 2024", title="GraphRAG",
        items=["Local + global community graph structure",
               "Graph-aware retrieval at query time",
               "Hybrid: knowledge graph + vector search",
               "Strong for domains with pre-existing ontologies",
               "Requires schema design and vector infrastructure"],
        tag="STRUCTURED DOMAINS", tag_color=SKY,
        x=Inches(0.5), y=COL_Y, w=COL_W, h=COL_H,
        pos=(3,), neg=(4,))

div = s.shapes.add_shape(1, Inches(6.45), COL_Y + Inches(0.2), Pt(2), COL_H - Inches(0.4))
div.fill.solid(); div.fill.fore_color.rgb = TEAL
div.line.fill.background()

col_box(s,
        head="VECTIFYAI · OPENKB", title="Compilation Model",
        items=["Flat wiki — concepts + entities + cross-refs",
               "Vectorless retrieval via PageIndex",
               "No schema, no ontology, no embedding pipeline",
               "Human-readable, diffable, fully auditable",
               "Built for dynamic, growing corpora"],
        tag="DYNAMIC CORPORA", tag_color=TEAL,
        x=Inches(6.55), y=COL_Y, w=COL_W, h=COL_H,
        pos=(3,4))

folio(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Quote
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "CHAPTER 03")
# Teal rule
rule = s.shapes.add_shape(1, Inches(1.4), Inches(1.6), Pt(4), Inches(3.8))
rule.fill.solid(); rule.fill.fore_color.rgb = TEAL
rule.line.fill.background()

txbox(s,
      '"Synthesis reflects everything the system has consumed\n— not just random K-nearest chunks."',
      Inches(1.9), Inches(1.6), Inches(10), Inches(2.6),
      size=24, italic=True, color=SOFT)
txbox(s,
      "OPENKB  ·  INSPIRED BY KARPATHY — COMPILATION-BASED RETRIEVAL AS A FIRST-CLASS BUILD STEP",
      Inches(1.9), Inches(4.4), Inches(10), Inches(0.4),
      size=8, color=MUTED, mono=True)
folio(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Closing
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
top_labels(s, "OPENKB")
kicker(s, "THE KNOWLEDGE-FIRST FUTURE", x=Inches(0.9), y=Inches(1.5))
gradient_title(s, "Build. Compile.\nKnow.",
               Inches(0.9), Inches(2.0), Inches(11), Inches(2.4), size=54)
txbox(s,
      "Stop retrieving isolated chunks. Treat document ingestion as a compilation step — and let your "
      "knowledge base compound across every document you add. RAG for breadth. Graphs for structure. "
      "Compilation for synthesis.",
      Inches(0.9), Inches(4.5), Inches(10), Inches(1.4), size=14, color=SOFT)

chip_x = Inches(0.9)
chip_y = Inches(6.1)
for chip_text in ["github.com/VectifyAI/OpenKB", "PageIndex — vectorless retrieval", "shekerkamma/OpenKB"]:
    chip_w = Inches(2.9)
    chip = s.shapes.add_shape(1, chip_x, chip_y, chip_w, Inches(0.38))
    chip.fill.solid(); chip.fill.fore_color.rgb = BG_ELEV
    chip.line.color.rgb = TEAL; chip.line.width = Pt(0.75)
    txbox(s, chip_text, chip_x + Inches(0.1), chip_y + Pt(3),
          chip_w - Inches(0.2), Inches(0.3), size=8, color=TEAL, mono=True)
    chip_x += chip_w + Inches(0.2)

folio(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
out = "/home/sheke/content-ideas/runs/2026-06-18-ai-knowledge-graphs-rag-research-deck/ai-knowledge-graphs-rag-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
