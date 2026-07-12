#!/usr/bin/env python3
"""Build a polished agent-native architecture deck with diagram-led slides."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx.dml.color import RGBColor  # noqa: E402

PPTXKIT_DIR = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
if not PPTXKIT_DIR.exists():
    PPTXKIT_DIR = Path("/home/shekerk/snap/codex/64/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT_DIR))

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptxkit import Deck  # noqa: E402

RUN_DIR = Path(__file__).resolve().parent
OUT = RUN_DIR / "agent-native-architecture-reviewed.pptx"
ASSET_DIR = RUN_DIR / "diagram-assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)

d = Deck(footer="Agent-Native Apps | Architecture Brief | June 2026")
b = d.b
TOTAL = 21

C = {
    "WHITE": RGBColor.from_string("FFFFFF"),
    "NAVY": RGBColor.from_string("1F3B6E"),
    "NAVY_CARD": RGBColor.from_string("0D3B5E"),
    "RED": RGBColor.from_string("E31837"),
    "LIGHT_BG": RGBColor.from_string("F5F7FA"),
    "BORDER": RGBColor.from_string("D0D5DD"),
    "MUTED": RGBColor.from_string("6B7280"),
    "CHARCOAL": RGBColor.from_string("1A1A1A"),
    "CODE_BG": RGBColor.from_string("F0F2F5"),
}


def mcolor(color):
    if isinstance(color, str):
        return color
    return "#" + str(color)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def pcolor(color):
    if isinstance(color, RGBColor):
        return color
    if isinstance(color, str):
        return RGBColor.from_string(color.replace("#", ""))
    raise TypeError(f"Unsupported color type: {type(color)!r}")


def card(s, left, top, w, h, header, body, *, fill=C["WHITE"], highlight=False, body_size=11.4):
    line = C["RED"] if highlight else C["BORDER"]
    d.rect(s, left, top, w, h, pcolor(fill), line=pcolor(line), radius=0.04)
    d.rect(s, left, top, w, Inches(0.38), pcolor(C["NAVY_CARD"]), line=pcolor(C["NAVY_CARD"]), radius=0.04)
    d.text(s, header, left + Inches(0.12), top + Inches(0.06), w - Inches(0.24), Inches(0.24),
           size=12, color=b.WHITE, bold=True, shrink=True)
    d.text(s, body, left + Inches(0.12), top + Inches(0.5), w - Inches(0.24), h - Inches(0.58),
           size=body_size, color=b.INK, shrink=True)


def banner(s, text):
    d.rect(s, 0, Inches(6.82), d.W, Inches(0.68), pcolor(C["NAVY_CARD"]))
    d.text(s, text, Inches(0.4), Inches(6.9), d.CW, Inches(0.28), size=15, color=b.WHITE,
           bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, shrink=True)


def pic(s, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    s.shapes.add_picture(str(path), left, top, **kwargs)


def _mpl():
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Circle

    return plt, FancyArrowPatch, FancyBboxPatch, Circle


def save_fig(fig, path):
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.03)
    import matplotlib.pyplot as plt

    plt.close(fig)
    return path


def rounded_box(ax, xy, w, h, text, *, fill="#FFFFFF", edge=C["BORDER"], color=C["CHARCOAL"], size=12, weight="bold"):
    from matplotlib.patches import FancyBboxPatch

    x, y = xy
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.08",
                                facecolor=mcolor(fill), edgecolor=mcolor(edge), linewidth=1.8))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=size, color=mcolor(color), fontweight=weight)


def arrow(ax, a, b_, color=C["BORDER"], label=None, label_offset=(0, 0)):
    from matplotlib.patches import FancyArrowPatch

    ax.add_patch(FancyArrowPatch(a, b_, arrowstyle="->", mutation_scale=15, linewidth=2.0, color=mcolor(color)))
    if label:
        ax.text((a[0] + b_[0]) / 2 + label_offset[0], (a[1] + b_[1]) / 2 + label_offset[1], label,
                ha="center", va="center", fontsize=9.5, color=mcolor(C["MUTED"]), fontweight="bold")


def diagram_hero():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.1, 4.4), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    rounded_box(ax, (0.4, 2.25), 2.55, 1.05, "Context files", fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"])
    rounded_box(ax, (0.7, 4.05), 2.0, 0.75, "Memory", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (0.7, 0.6), 2.0, 0.75, "Assets", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    ax.add_patch(Circle((5.05, 3.0), 1.05, facecolor=mcolor(C["NAVY_CARD"]), edgecolor=mcolor(C["RED"]), linewidth=2.2))
    ax.text(5.05, 3.22, "AGENT", ha="center", va="center", fontsize=16, color=mcolor(C["WHITE"]), fontweight="bold")
    ax.text(5.05, 2.82, "HOST", ha="center", va="center", fontsize=11.5, color=mcolor(C["RED"]), fontweight="bold")
    rounded_box(ax, (7.3, 4.05), 3.0, 0.82, "MCP  ·  CLI  ·  API", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"])
    rounded_box(ax, (7.3, 2.65), 3.0, 0.82, "App surface", fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"])
    rounded_box(ax, (7.3, 1.25), 3.0, 0.82, "Shared workspace", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"])
    arrow(ax, (2.95, 3.05), (3.98, 3.05), color=C["RED"], label="context")
    arrow(ax, (6.05, 3.05), (7.25, 4.4), color=C["BORDER"], label="tools", label_offset=(0.1, 0.12))
    arrow(ax, (6.05, 2.95), (7.25, 3.05), color=C["BORDER"], label="surface")
    arrow(ax, (6.05, 2.4), (7.25, 1.65), color=C["BORDER"], label="write-back", label_offset=(0.25, -0.08))
    ax.text(6.0, 5.45, "The user brings the agent; the app becomes the surface it uses.",
            ha="center", va="center", fontsize=12.4, color=mcolor(C["NAVY"]), fontweight="bold")
    ax.text(6.0, 0.28, "Shared context + tool access + visible collaboration = agent-native app.",
            ha="center", va="center", fontsize=10.9, color=mcolor(C["MUTED"]), style="italic")
    return save_fig(fig, ASSET_DIR / "hero-architecture.png")


def diagram_inversion():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(10.2, 4.3), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.5)
    ax.axis("off")
    # left
    ax.add_patch(FancyBboxPatch((0.4, 0.5), 5.2, 4.2, boxstyle="round,pad=0.02,rounding_size=0.12",
                                facecolor=mcolor(C["WHITE"]), edgecolor=mcolor(C["BORDER"]), linewidth=1.8))
    ax.add_patch(FancyBboxPatch((0.6, 4.4), 4.8, 0.38, boxstyle="round,pad=0.02,rounding_size=0.08",
                                facecolor=mcolor(C["RED"]), edgecolor=mcolor(C["RED"]), linewidth=1.0))
    ax.text(3.0, 4.59, "App with an agent", ha="center", va="center", fontsize=14.5, color=mcolor(C["WHITE"]), fontweight="bold")
    rounded_box(ax, (1.0, 2.0), 2.9, 1.0, "Product UI", fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"], size=13)
    ax.add_patch(Circle((4.25, 2.5), 0.5, facecolor=mcolor(C["RED"]), edgecolor=mcolor(C["RED"])))
    ax.text(4.25, 2.5, "AI", ha="center", va="center", fontsize=13, color=mcolor(C["WHITE"]), fontweight="bold")
    ax.text(1.0, 1.05, "Context stays inside the app.", fontsize=11, color=mcolor(C["MUTED"]))
    # right
    ax.add_patch(FancyBboxPatch((6.5, 0.5), 5.2, 4.2, boxstyle="round,pad=0.02,rounding_size=0.12",
                                facecolor=mcolor(C["WHITE"]), edgecolor=mcolor(C["BORDER"]), linewidth=1.8))
    ax.add_patch(FancyBboxPatch((6.7, 4.4), 4.8, 0.38, boxstyle="round,pad=0.02,rounding_size=0.08",
                                facecolor=mcolor(C["NAVY"]), edgecolor=mcolor(C["NAVY"]), linewidth=1.0))
    ax.text(9.1, 4.59, "Agent with an app", ha="center", va="center", fontsize=14.5, color=mcolor(C["WHITE"]), fontweight="bold")
    ax.add_patch(Circle((7.8, 2.5), 0.5, facecolor=mcolor(C["NAVY_CARD"]), edgecolor=mcolor(C["RED"]), linewidth=1.5))
    ax.text(7.8, 2.7, "HOST", ha="center", va="center", fontsize=10, color=mcolor(C["RED"]), fontweight="bold")
    ax.text(7.8, 2.43, "AGENT", ha="center", va="center", fontsize=11, color=mcolor(C["WHITE"]), fontweight="bold")
    rounded_box(ax, (9.2, 2.0), 2.0, 1.0, "App surface", fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"], size=13)
    arrow(ax, (4.75, 2.5), (6.25, 2.5), color=C["RED"])
    arrow(ax, (8.3, 2.5), (9.1, 2.5), color=C["NAVY"], label="surface", label_offset=(0.0, 0.22))
    rounded_box(ax, (8.6, 0.9), 2.7, 0.72, "Context + tools + memory", fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"], size=10.8)
    ax.text(6.9, 1.05, "The app becomes one surface among many.", fontsize=11, color=mcolor(C["MUTED"]))
    ax.text(6.0, 5.05, "The inversion: the agent owns context; the app is the workspace.", ha="center",
            va="center", fontsize=12.3, color=mcolor(C["NAVY"]), fontweight="bold")
    return save_fig(fig, ASSET_DIR / "inversion-diagram.png")


def diagram_reference():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(10.2, 4.6), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.2)
    ax.axis("off")
    zones = [
        (0.4, 3.6, 11.2, 0.95, "Context layer"),
        (0.4, 2.35, 11.2, 0.95, "Execution layer"),
        (0.4, 1.1, 11.2, 0.95, "Shared surface"),
    ]
    for x, y, w, h, label in zones:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.08",
                                    facecolor=mcolor(C["LIGHT_BG"]), edgecolor=mcolor(C["BORDER"]), linewidth=1.2))
        ax.text(x + 0.18, y + h - 0.14, label, fontsize=11.5, color=mcolor(C["NAVY"]), fontweight="bold", va="top")
    rounded_box(ax, (0.65, 3.88), 2.1, 0.5, "context.md", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (3.0, 3.88), 2.5, 0.5, "project files", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (5.85, 3.88), 2.7, 0.5, "assets / memory", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (0.65, 2.62), 2.4, 0.5, "Cursor / Codex", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (3.3, 2.62), 2.4, 0.5, "MCP / CLI / API", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (5.95, 2.62), 2.6, 0.5, "agent loop", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (0.65, 1.37), 2.4, 0.5, "canvas / doc", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (3.3, 1.37), 2.4, 0.5, "sheet / dashboard", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    rounded_box(ax, (5.95, 1.37), 2.6, 0.5, "outcomes", fill=C["WHITE"], edge=C["BORDER"], color=C["NAVY"], size=11)
    arrow(ax, (1.7, 3.88), (1.7, 3.16), color=C["RED"], label="context", label_offset=(0.0, 0.05))
    arrow(ax, (4.25, 3.88), (4.25, 3.16), color=C["BORDER"], label="task", label_offset=(0.0, 0.05))
    arrow(ax, (7.2, 3.88), (7.2, 3.16), color=C["BORDER"], label="memory", label_offset=(0.0, 0.05))
    arrow(ax, (3.05, 2.87), (3.25, 2.87), color=C["BORDER"])
    arrow(ax, (5.7, 2.87), (5.9, 2.87), color=C["BORDER"])
    arrow(ax, (3.05, 1.62), (3.25, 1.62), color=C["BORDER"])
    arrow(ax, (5.7, 1.62), (5.9, 1.62), color=C["RED"], label="write-back", label_offset=(0.0, -0.12))
    ax.text(6.0, 0.32, "Reference architecture: context → agent loop → shared surface → outcome → memory.",
            ha="center", va="center", fontsize=11, color=mcolor(C["MUTED"]), style="italic")
    return save_fig(fig, ASSET_DIR / "reference-architecture.png")


def diagram_tools():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.5, 5.2), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")
    center = Circle((5, 3), 0.95, facecolor=mcolor(C["NAVY_CARD"]), edgecolor=mcolor(C["RED"]), linewidth=2.2)
    ax.add_patch(center)
    ax.text(5, 3.2, "AGENT", ha="center", va="center", fontsize=15, color=mcolor(C["WHITE"]), fontweight="bold")
    ax.text(5, 2.9, "HOST", ha="center", va="center", fontsize=11, color=mcolor(C["RED"]), fontweight="bold")
    nodes = [
        (1.0, 4.4, "MCP", C["LIGHT_BG"]),
        (7.95, 4.4, "CLI", C["WHITE"]),
        (1.0, 1.1, "API", C["LIGHT_BG"]),
        (7.95, 1.1, "Browser", C["WHITE"]),
    ]
    for x, y, txt, fill in nodes:
        rounded_box(ax, (x, y), 1.9, 0.82, txt, fill=fill, edge=C["BORDER"], color=C["NAVY"], size=12)
    arrow(ax, (4.25, 3.7), (2.85, 4.75), color=C["RED"])
    arrow(ax, (5.75, 3.7), (7.85, 4.75), color=C["BORDER"])
    arrow(ax, (4.25, 2.35), (2.85, 1.5), color=C["BORDER"])
    arrow(ax, (5.75, 2.35), (7.85, 1.5), color=C["RED"])
    ax.text(5, 5.55, "Tool access is the unlock: atomic primitives, not workflow-shaped black boxes.",
            ha="center", va="center", fontsize=11.7, color=mcolor(C["NAVY"]), fontweight="bold")
    ax.text(5, 0.35, "Parity comes from making every useful action reachable through tools.",
            ha="center", va="center", fontsize=10.8, color=mcolor(C["MUTED"]), style="italic")
    return save_fig(fig, ASSET_DIR / "tool-hub.png")


def diagram_context_stack():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.8, 5.0), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")
    layers = [
        ("Local files", "markdown, assets, docs"),
        ("Project memory", "notes, history, context"),
        ("Connected tools", "browser, CLI, API"),
        ("Shared surface", "canvas, doc, sheet"),
        ("User oversight", "review, correction"),
    ]
    y = 4.9
    for i, (title, body) in enumerate(layers):
        fill = C["WHITE"] if i % 2 else C["LIGHT_BG"]
        rounded_box(ax, (1.0, y), 6.6, 0.68, f"{title}  ·  {body}", fill=fill, edge=C["BORDER"], color=C["NAVY"], size=11.6)
        y -= 0.86
    arrow(ax, (8.3, 5.2), (8.3, 1.8), color=C["RED"], label="quality grows", label_offset=(0.2, 0))
    ax.text(8.65, 3.5, "context\ncompounds", ha="center", va="center", fontsize=12.2, color=mcolor(C["NAVY"]), fontweight="bold")
    ax.text(4.3, 0.35, "The system improves over time because the agent can reuse accumulated context.",
            ha="center", va="center", fontsize=11, color=mcolor(C["MUTED"]), style="italic")
    return save_fig(fig, ASSET_DIR / "context-stack.png")


def diagram_opportunity_matrix():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(8.8, 5.1), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_facecolor("white")
    ax.axhline(0.5, color=mcolor(C["BORDER"]), lw=1.2)
    ax.axvline(0.5, color=mcolor(C["BORDER"]), lw=1.2)
    points = [
        ("Miro", 0.82, 0.83, C["NAVY"]),
        ("Notion", 0.74, 0.74, C["RED"]),
        ("MagicPath", 0.79, 0.66, C["NAVY"]),
        ("Sheets", 0.63, 0.56, C["RED"]),
        ("Adobe", 0.88, 0.43, C["NAVY"]),
    ]
    for label, x, y, col in points:
        ax.scatter(x, y, s=230, color=mcolor(col), alpha=0.9, edgecolors="white", linewidths=1.5, zorder=3)
        ax.annotate(label, (x, y), xytext=(0, 11), textcoords="offset points", ha="center",
                    fontsize=10.5, color=mcolor(C["NAVY"]), fontweight="bold")
    ax.text(0.02, 0.97, "Agent-native fit", ha="left", va="top", fontsize=11.2, color=mcolor(C["MUTED"]))
    ax.text(0.98, 0.02, "Shared surface richness", ha="right", va="bottom", fontsize=11.2, color=mcolor(C["MUTED"]))
    ax.text(0.97, 0.96, "Surface + tools = opportunity", ha="right", va="top", fontsize=11.5, color=mcolor(C["RED"]), fontweight="bold")
    for sp in ax.spines.values():
        sp.set_color(mcolor(C["BORDER"]))
    plt.tight_layout(pad=0.4)
    return save_fig(fig, ASSET_DIR / "opportunity-matrix.png")


def diagram_loop():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.8, 4.6), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5.5)
    ax.axis("off")
    steps = [
        (1.0, 3.7, "Intent"),
        (4.0, 4.0, "Agent acts"),
        (7.0, 3.7, "User watches"),
        (4.0, 1.1, "User corrects"),
    ]
    for x, y, txt in steps:
        rounded_box(ax, (x, y), 2.8, 0.92, txt, fill=C["LIGHT_BG"], edge=C["BORDER"], color=C["NAVY"], size=12.5)
    arrow(ax, (3.8, 4.15), (4.0, 4.45), color=C["BORDER"])
    arrow(ax, (6.8, 4.45), (7.0, 4.15), color=C["BORDER"])
    arrow(ax, (7.0, 3.25), (5.4, 2.0), color=C["RED"])
    arrow(ax, (4.0, 2.0), (2.1, 3.25), color=C["BORDER"])
    ax.text(5, 2.75, "Supervised collaboration", ha="center", va="center", fontsize=12.6,
            color=mcolor(C["NAVY"]), fontweight="bold")
    ax.text(5, 0.35, "The human stays in the loop and can correct the work live.",
            ha="center", va="center", fontsize=10.9, color=mcolor(C["MUTED"]), style="italic")
    return save_fig(fig, ASSET_DIR / "workflow-loop.png")


HERO = diagram_hero()
INVERSION = diagram_inversion()
REFERENCE = diagram_reference()
TOOLS = diagram_tools()
CONTEXT = diagram_context_stack()
MATRIX = diagram_opportunity_matrix()
LOOP = diagram_loop()


def title_slide():
    s = d.slide(fill=b.WHITE)
    d.rect(s, 0, 0, d.W, Inches(0.14), C["NAVY"])
    d.text(s, "ARCHITECTURE BRIEF", d.M, Inches(0.5), Inches(3.5), Inches(0.24), size=12,
           color=b.MUTED, bold=True)
    d.text(s, "Agent-Native Apps", d.M, Inches(1.0), Inches(6.8), Inches(0.7), size=38,
           color=b.NAVY, bold=True, font=b.FONT_H, shrink=True)
    d.text(s, "Software designed for the user's own agent.",
           d.M, Inches(1.8), Inches(6.8), Inches(0.28), size=15, color=b.MUTED, shrink=True)
    d.rect(s, d.M, Inches(2.4), Inches(1.6), Inches(0.06), pcolor(C["RED"]))
    d.text(s, "Cursor + Codex + shared context + tool access", d.M, Inches(2.55), Inches(6.4), Inches(0.3),
           size=14, color=b.NAVY, bold=True)
    pic(s, HERO, Inches(7.8), Inches(1.0), width=Inches(5.0))
    banner(s, "Build for the user's agent, not just the user.")
    d.footer(s, 1, TOTAL)
    notes(s, "This title slide frames the opportunity as an architecture problem: the user, context, tools, and surface must all line up.")


def thesis_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "The thesis", "The workflow changes when the agent owns the loop")
    card(s, d.M, Inches(1.72), Inches(4.0), Inches(2.4), "Old model",
         "App embeds a chatbot. The model sees only local app context, so it can answer questions but cannot finish the job.\n\nResult: the user still does the hard part outside the app.",
         fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(4.25), Inches(1.72), Inches(4.0), Inches(2.4), "New model",
         "The user's agent owns the context and uses the app as one surface among many.\n\nResult: the app becomes a place where work actually gets done, with the user's broader memory attached.",
         fill=C["WHITE"], highlight=True, body_size=11.4)
    card(s, d.M + Inches(8.5), Inches(1.72), Inches(4.0), Inches(2.4), "What changes",
         "Moat shifts to shared context, tool access, workflow design, and persistent memory.\n\nIf those four pieces are weak, the model layer is just decoration.",
         fill=C["LIGHT_BG"], body_size=11.4)
    d.rect(s, d.M, Inches(4.55), d.CW, Inches(0.7), C["NAVY_CARD"], radius=0.03)
    d.text(s, "If the app is not usable by the user's agent, it is only partially agent-native.",
           d.M + Inches(0.18), Inches(4.72), Inches(12.0), Inches(0.24), size=13.2, color=b.WHITE,
           bold=True, shrink=True, anchor=MSO_ANCHOR.MIDDLE)
    banner(s, "The architecture is a loop, not a router.")
    d.footer(s, 2, TOTAL)
    notes(s, "Use this slide to explain the inversion and the strategic implication.")


def inversion_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Model inversion", "App with an agent vs. agent with an app")
    pic(s, INVERSION, d.M, Inches(1.5), width=d.CW)
    d.rect(s, d.M, Inches(5.92), d.CW, Inches(0.52), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.03)
    d.text(s, "Old model: the app contains the agent. New model: the agent contains the app surface and brings its own context.",
           d.M + Inches(0.12), Inches(6.04), Inches(12.05), Inches(0.18), size=10.8, color=b.NAVY, bold=True, shrink=True)
    d.footer(s, 3, TOTAL)
    notes(s, "The left side is the old pattern. The right side is the agent-native pattern.")


def principles_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Five principles", "The checklist that separates real agent-native systems from AI-colored software")
    titles = [
        ("Parity", "Anything the user can do in UI, the agent can do through tools.\n\nIf parity is missing, the agent is only a suggestion engine."),
        ("Granularity", "Tools are atomic primitives. Features come from prompt-defined outcomes.\n\nKeep the actions small so agents can combine them."),
        ("Composability", "New capabilities come from new prompt + tool combinations.\n\nThe product gets smarter when the same pieces work in new sequences."),
        ("Emergent capability", "The system should handle requests you did not pre-script.\n\nThat is the difference between automation and a real agent surface."),
        ("Improvement over time", "Context and memory should make the system better with use.\n\nIf it does not learn from the workspace, it will not compound value."),
    ]
    xs = [d.M, d.M + Inches(2.5), d.M + Inches(5.0), d.M + Inches(7.5), d.M + Inches(10.0)]
    widths = [Inches(2.2)] * 5
    for i, ((head, body), x, w) in enumerate(zip(titles, xs, widths)):
        d.rect(s, x, Inches(1.95), w, Inches(3.2), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.04)
        d.rect(s, x, Inches(1.95), w, Inches(0.42), pcolor(C["NAVY_CARD"]))
        d.text(s, head, x + Inches(0.08), Inches(2.02), w - Inches(0.16), Inches(0.22), size=11.5,
               color=b.WHITE, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.1), Inches(2.5), w - Inches(0.2), Inches(2.3), size=11.0,
               color=b.INK, shrink=True)
    banner(s, "If the product cannot satisfy these five checks, it is still too constrained.")
    d.footer(s, 4, TOTAL)


def architecture_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Reference architecture", "Context → agent loop → shared surface → outcome → memory")
    pic(s, REFERENCE, d.M, Inches(1.58), width=d.CW)
    d.rect(s, d.M, Inches(5.92), d.CW, Inches(0.52), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.03)
    d.text(s, "Read left to right: context feeds the agent, tools move it into the app, and write-back turns output into future context.",
           d.M + Inches(0.12), Inches(6.04), Inches(12.05), Inches(0.18), size=10.8, color=b.NAVY, bold=True, shrink=True)
    d.footer(s, 5, TOTAL)


def tool_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Tool layer", "MCP, CLI, API, and browser access are the unlock")
    pic(s, TOOLS, d.M, Inches(1.65), width=Inches(6.85))
    card(s, d.M + Inches(7.1), Inches(1.85), Inches(5.2), Inches(1.2), "Parity",
         "If a user can do it, the agent needs a path to do it too.\n\nWithout that path, the feature is just a demo.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.1), Inches(3.2), Inches(5.2), Inches(1.2), "Granularity",
         "Expose primitives; keep domain logic in the agent loop.\n\nThat is what lets any agent chain actions together.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.1), Inches(4.55), Inches(5.2), Inches(1.2), "CRUD completeness",
         "Every entity needs create, read, update, and delete.\n\nIf delete is missing, the agent cannot safely clean up its own work.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "Tool access is a product decision, not an implementation detail.")
    d.footer(s, 6, TOTAL)


def context_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Context layer", "Local files, project memory, and workspace state compound over time")
    pic(s, CONTEXT, d.M, Inches(1.7), width=Inches(6.9))
    card(s, d.M + Inches(7.2), Inches(1.75), Inches(5.15), Inches(1.25), "Shared workspace",
         "The agent and user work in the same files and see the same state.\n\nThat shared state is what makes collaboration legible.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(3.15), Inches(5.15), Inches(1.25), "context.md pattern",
         "Write accumulated knowledge into files the agent can re-read.\n\nExamples: brand rules, architecture notes, prior decisions, and approved defaults.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(4.55), Inches(5.15), Inches(1.25), "Improvement over time",
         "Better memory makes the system smarter without shipping code.\n\nThe product improves because the workspace remembers what good looks like.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "Context is the moat. It compounds.")
    d.footer(s, 7, TOTAL)


def anti_patterns_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Anti-patterns", "The common ways teams build something that is not actually agent-native")
    cards = [
        ("Agent as router", "The agent only chooses which fixed function to call.\n\nThat means it can route work, but it cannot reason through it."),
        ("Workflow-shaped tools", "A single tool hides judgment and prevents composability.\n\nYou want primitives, not a black box with a nicer prompt."),
        ("Sandbox split-brain", "Agent output goes to a separate workspace the user doesn't see.\n\nIf the user cannot verify state, trust breaks down fast."),
        ("Silent actions", "State changes do not show up in the UI immediately.\n\nVisible feedback is how humans know the agent did the right thing."),
        ("Heuristic completion", "The system guesses the task is done instead of signaling it.\n\nCompletion needs explicit status, not vibes."),
    ]
    ys = [1.9, 2.9, 3.9, 4.9, 5.9]
    for (title, body), y in zip(cards, ys):
        d.rect(s, d.M, Inches(y), d.CW, Inches(0.72), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.04)
        d.rect(s, d.M, Inches(y), Inches(0.08), Inches(0.72), pcolor(C["RED"]))
        d.text(s, title, d.M + Inches(0.18), Inches(y + 0.09), Inches(3.0), Inches(0.2), size=12,
               color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, d.M + Inches(3.35), Inches(y + 0.09), Inches(8.8), Inches(0.2), size=11.0,
               color=b.INK, shrink=True)
    banner(s, "If the code already solved the interesting part, the agent is just a caller.")
    d.footer(s, 8, TOTAL)


def opportunity_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Where the opportunity is", "Surfaces that already carry context are best positioned")
    pic(s, MATRIX, d.M, Inches(1.82), width=Inches(7.0))
    card(s, d.M + Inches(7.2), Inches(1.9), Inches(5.0), Inches(1.25), "Miro / canvas tools",
         "Shared boards can hold markdown, diagrams, and prototypes.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(3.35), Inches(5.0), Inches(1.25), "Notion / docs / Sheets",
         "Knowledge bases and structured tables are easy to connect.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(4.8), Inches(5.0), Inches(1.25), "Adobe-class suites",
         "Large incumbents still have room for agent-native workflows.\n\nThey already own expensive, repeatable work where better context matters.", fill=C["LIGHT_BG"], body_size=11.4)
    d.rect(s, d.M, Inches(5.28), d.CW, Inches(0.42), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.03)
    d.text(s, "Best first wedges: canvas tools, docs, sheets, CRM/ops work, and design workflows.",
           d.M + Inches(0.12), Inches(5.38), Inches(12.05), Inches(0.16), size=10.5, color=b.NAVY, bold=True, shrink=True)
    banner(s, "The surface is the product. The agent only makes it useful.")
    d.footer(s, 9, TOTAL)


def why_now_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Why now", "Three shifts make agent-native software viable today")
    cards = [
        ("User agents are real", "People are already working inside Cursor, Codex, and similar hosts.\n\nOpenAI's Codex CLI can read, change, and run code locally."),
        ("Apps expose tools", "MCP, CLI, API, and browser access make external action practical.\n\nNotion MCP can read and write workspace content with permissions."),
        ("Context is portable", "Files, docs, assets, and memory can live outside a single product.\n\nCodex's in-app browser gives a shared view of rendered pages."),
    ]
    x = d.M
    for i, (title, body) in enumerate(cards):
        card(s, x, Inches(2.0), Inches(3.9), Inches(1.75), title, body,
             fill=C["LIGHT_BG"] if i != 1 else C["WHITE"], highlight=(i == 1), body_size=11.4)
        x += Inches(4.1)
    d.rect(s, d.M, Inches(4.35), d.CW, Inches(0.92), pcolor(C["NAVY_CARD"]))
    d.text(s, "The category opens when the agent can bring its own context into an app, not just consume the app's AI.",
           d.M + Inches(0.18), Inches(4.55), Inches(12.0), Inches(0.26), size=13.0, color=b.WHITE,
           bold=True, shrink=True, anchor=MSO_ANCHOR.MIDDLE)
    d.text(s, "Evidence: OpenAI Codex CLI, OpenAI Codex app browser, and Notion MCP show the workflow is already here.",
           d.M, Inches(5.55), d.CW, Inches(0.22), size=10.8, color=b.MUTED, italic=True, shrink=True)
    banner(s, "The opportunity is a workspace designed for the user's agent.")
    d.footer(s, 10, TOTAL)


def build_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "How to build", "Five steps to find and ship an agent-native app")
    steps = [
        ("1. Pick the job", "Define the outcome the user and agent achieve together."),
        ("2. Define the surface", "Choose the canvas or workspace both can work on."),
        ("3. Build for BYO agent", "Make the app useful inside Cursor/Codex-style hosts."),
        ("4. Start from expertise", "Choose a domain you already know deeply."),
        ("5. Plan tools early", "Scope MCP / CLI / API access on day one."),
    ]
    y = Inches(1.9)
    for i, (title, body) in enumerate(steps):
        fill = C["LIGHT_BG"] if i % 2 == 0 else C["WHITE"]
        d.rect(s, d.M, y, d.CW, Inches(0.72), pcolor(fill), line=pcolor(C["BORDER"]), radius=0.04)
        d.rect(s, d.M, y, Inches(0.6), Inches(0.72), pcolor(C["NAVY_CARD"]))
        d.text(s, str(i + 1), d.M, y + Inches(0.1), Inches(0.6), Inches(0.25), size=16, color=b.WHITE,
               bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, title, d.M + Inches(0.8), y + Inches(0.08), Inches(3.8), Inches(0.24), size=12.5,
               color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, d.M + Inches(4.95), y + Inches(0.08), Inches(7.0), Inches(0.24), size=11.4,
               color=b.INK, shrink=True)
        y += Inches(0.82)
    banner(s, "The product changes when you start with the shared surface and the tool surface together.")
    d.footer(s, 11, TOTAL)


def job_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Step 1. Pick the job", "Start with one outcome the user and agent should achieve together")
    card(s, d.M, Inches(1.82), Inches(5.2), Inches(2.8), "Job to be done",
         "A good wedge is narrow enough to evaluate and broad enough to repeat.\n\nExamples:\n• design a page\n• update a CRM record\n• reconcile a sheet",
         fill=C["LIGHT_BG"], highlight=True, body_size=11.4)
    card(s, d.M + Inches(5.45), Inches(1.82), Inches(3.1), Inches(1.25), "Success criteria",
         "The user can see the result, approve it, and correct it quickly.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(5.45), Inches(3.2), Inches(3.1), Inches(1.25), "What the agent needs",
         "Inputs, permissions, and a visible surface that matches the task.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(5.45), Inches(4.58), Inches(3.1), Inches(1.25), "Failure mode",
         "If the agent cannot complete the work end-to-end, the job is too large.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(8.75), Inches(1.82), Inches(3.55), Inches(4.01), "Decision rule",
         "Choose a job that already lives in a tool-rich workflow.\n\nThe best jobs are repeated, visible, and easy to verify.",
         fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "Do not start with a model feature. Start with a repeated job.")
    d.footer(s, 12, TOTAL)


def surface_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Step 2. Define the surface", "Pick the shared workspace both sides can edit")
    surface_cards = [
        ("Canvas", "Use when the task is visual or spatial.", True),
        ("Doc", "Use when the task is narrative or collaborative.", False),
        ("Sheet", "Use when the task is structured and tabular.", False),
        ("Browser", "Use when the task needs live app interaction.", False),
    ]
    x = d.M
    for title, body, highlight in surface_cards:
        card(s, x, Inches(2.05), Inches(2.95), Inches(2.0), title, body,
             fill=C["LIGHT_BG"] if not highlight else C["WHITE"], highlight=highlight, body_size=11.4)
        x += Inches(3.05)
    d.rect(s, d.M, Inches(4.6), d.CW, Inches(0.92), pcolor(C["NAVY_CARD"]))
    d.text(s, "The shared surface is where the user watches the agent work, not where the model hides the work.",
           d.M + Inches(0.2), Inches(4.82), Inches(12.0), Inches(0.26), size=13.0, color=b.WHITE,
           bold=True, shrink=True, anchor=MSO_ANCHOR.MIDDLE)
    banner(s, "The surface should make the collaboration obvious.")
    d.footer(s, 13, TOTAL)


def byo_agent_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Step 3. Build for BYO agent", "Design for external agents first, not an embedded chatbot")
    pic(s, INVERSION, d.M, Inches(1.65), width=Inches(7.0))
    card(s, d.M + Inches(7.35), Inches(1.85), Inches(4.95), Inches(1.25), "No lock-in",
         "Do not depend on a proprietary agent living inside the app.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.35), Inches(3.25), Inches(4.95), Inches(1.25), "Tool first",
         "Expose actions through MCP, CLI, API, and browser flows.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.35), Inches(4.65), Inches(4.95), Inches(1.25), "Shared context",
         "Let the agent inherit local files, notes, and memory.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "The app should be the place the agent works, not the place it gets trapped.")
    d.footer(s, 14, TOTAL)


def expertise_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Step 4. Start from expertise", "Use a domain you understand deeply enough to judge output quality")
    card(s, d.M, Inches(1.85), Inches(3.15), Inches(3.15), "Your advantage",
         "You know the workflow, the failure cases, the language, and the acceptable tradeoffs.",
         fill=C["LIGHT_BG"], highlight=True, body_size=11.4)
    card(s, d.M + Inches(3.95), Inches(1.85), Inches(3.65), Inches(3.15), "Good wedges",
         "Product design\nOperations\nMarketing\nFinance\nSupport",
         fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.9), Inches(1.85), Inches(4.4), Inches(3.15), "What to avoid",
         "A generic productivity app with no domain truth.\n\nIf you cannot tell good from bad output, the feedback loop will be weak.",
         fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "Expertise is what lets you define the right shape for the agent.")
    d.footer(s, 15, TOTAL)


def tools_early_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Step 5. Plan tools early", "Scope MCP / CLI / API before you build the UI around them")
    pic(s, TOOLS, d.M, Inches(1.6), width=Inches(6.8))
    card(s, d.M + Inches(7.1), Inches(1.82), Inches(5.2), Inches(1.22), "MCP",
         "Best when the product needs structured tool access across clients.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.1), Inches(3.18), Inches(5.2), Inches(1.22), "CLI",
         "Best when local agents need a simple command-line control path.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.1), Inches(4.54), Inches(5.2), Inches(1.22), "API / browser",
         "Best when the agent must operate inside the same surface as the human.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "If you cannot draw the tool map on day one, the product is still under-scoped.")
    d.footer(s, 16, TOTAL)


def workflow_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Human-in-the-loop", "The collaboration loop is part of the product")
    pic(s, LOOP, d.M, Inches(1.7), width=Inches(6.9))
    card(s, d.M + Inches(7.2), Inches(1.85), Inches(5.0), Inches(1.25), "Watch",
         "The user sees what the agent is doing in real time.", fill=C["LIGHT_BG"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(3.25), Inches(5.0), Inches(1.25), "Correct",
         "The user can steer the outcome before the work is locked in.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(7.2), Inches(4.65), Inches(5.0), Inches(1.25), "Compound",
         "Corrections become context the agent can reuse on the next run.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "Visible feedback is what turns automation into collaboration.")
    d.footer(s, 17, TOTAL)


def examples_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Examples worth studying", "The strongest candidates already have shared context and shared surfaces")
    card(s, d.M, Inches(1.9), Inches(3.75), Inches(3.0), "Miro",
         "Canvas-first workspaces can hold diagrams, markdown, and prototypes.", fill=C["LIGHT_BG"], highlight=True, body_size=11.4)
    card(s, d.M + Inches(4.05), Inches(1.9), Inches(3.75), Inches(3.0), "Notion",
         "Docs plus MCP make it easy to move between knowledge and action.", fill=C["WHITE"], body_size=11.4)
    card(s, d.M + Inches(8.0), Inches(1.9), Inches(3.75), Inches(3.0), "Sheets / Adobe",
         "Structured tables and design suites already have rich user workflows.", fill=C["LIGHT_BG"], body_size=11.4)
    banner(s, "The best opportunities sit where the surface already has meaning.")
    d.footer(s, 18, TOTAL)


def launch_checklist_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Launch checklist", "Use this before you commit to the wedge")
    qs = [
        ("Can the user work here today?", "If not, you are asking for behavior change first."),
        ("Can the agent reach the tools?", "There must be a real tool path, not a prompt illusion."),
        ("Can both see the same state?", "Shared surface means shared visibility."),
        ("Can the system improve with use?", "If context cannot compound, the product will stall."),
    ]
    y = Inches(1.9)
    for i, (q, a) in enumerate(qs):
        card(s, d.M, y, d.CW, Inches(0.88), q, a, fill=C["LIGHT_BG"] if i % 2 == 0 else C["WHITE"],
             highlight=(i == 0), body_size=11.4)
        y += Inches(0.98)
    banner(s, "If any answer is no, the idea needs more architecture before it needs more code.")
    d.footer(s, 19, TOTAL)


def decision_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Go / no-go filter", "The founder checklist for this category")
    qs = [
        ("User already works here?", "If no one works in this surface today, you may be creating behavior from scratch."),
        ("Context exists outside the app?", "Agent-native wins when the real project context is broader than one database."),
        ("Can the agent reach tools?", "MCP, CLI, API, or browser access should be explicit."),
        ("Can humans watch/correct?", "The visible feedback loop must be part of the product."),
    ]
    y = Inches(1.9)
    for i, (q, a) in enumerate(qs):
        card(s, d.M, y, d.CW, Inches(0.92), q, a, fill=C["LIGHT_BG"] if i % 2 == 0 else C["WHITE"],
             highlight=(i == 0), body_size=11.4)
        y += Inches(1.02)
    banner(s, "If you cannot answer yes to all four, the wedge is probably not ready.")
    d.footer(s, 20, TOTAL)


def close_slide():
    s = d.slide(fill=b.WHITE)
    d.header(s, "Bottom line", "The best new AI opportunity is software redesigned for the user's agent")
    d.rect(s, d.M, Inches(1.75), d.CW, Inches(1.0), pcolor(C["NAVY_CARD"]))
    d.text(s, "Build for the user's agent, not just the user.", d.M + Inches(0.2), Inches(1.95), d.CW - Inches(0.4),
           Inches(0.4), size=20, color=b.WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    d.text(s, "Design a shared surface, expose the tools, and let the agent improve with context.",
           d.M, Inches(3.0), d.CW, Inches(0.35), size=15, color=b.NAVY, bold=True)
    d.rect(s, d.M, Inches(3.55), Inches(12.0), Inches(1.0), pcolor(C["LIGHT_BG"]), line=pcolor(C["BORDER"]), radius=0.04)
    d.text(s, "The opportunity is not another chatbot.\nIt is software that becomes useful to an external agent.",
           d.M + Inches(0.18), Inches(3.78), Inches(11.6), Inches(0.52), size=15, color=b.INK, bold=True,
           shrink=True, anchor=MSO_ANCHOR.MIDDLE)
    for i, txt in enumerate(["Choose a domain you know", "Map the shared surface", "Expose the tools early"]):
        x = d.M + Inches(i * 4.08)
        d.rect(s, x, Inches(5.0), Inches(3.75), Inches(0.56), pcolor(C["WHITE"]), line=pcolor(C["BORDER"]), radius=0.04)
        d.text(s, txt, x + Inches(0.1), Inches(5.12), Inches(3.55), Inches(0.2), size=12.0, color=b.NAVY,
               bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
    banner(s, "Agent-native apps are the next wedge because they compound context, not just clicks.")
    d.footer(s, 21, TOTAL, dark=False)


title_slide()
thesis_slide()
inversion_slide()
principles_slide()
architecture_slide()
tool_slide()
context_slide()
anti_patterns_slide()
opportunity_slide()
why_now_slide()
build_slide()
job_slide()
surface_slide()
byo_agent_slide()
expertise_slide()
tools_early_slide()
workflow_slide()
examples_slide()
launch_checklist_slide()
decision_slide()
close_slide()

OUT.parent.mkdir(parents=True, exist_ok=True)
d.save(OUT)
