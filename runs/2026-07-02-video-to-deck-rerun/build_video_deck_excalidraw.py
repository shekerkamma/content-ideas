#!/usr/bin/env python3
import sys
from pathlib import Path

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt

ROOT = Path("/home/shekerk/content-ideas")
RUN = ROOT / "runs/2026-07-02-video-to-deck-rerun"

sys.path.insert(0, str(ROOT / "skills/branded-pptx-deck/scripts"))
from pptxkit import Deck, PP_ALIGN, RGBColor


OUT = RUN / "ai-agents-new-saas-video-deck-excalidraw-draft.pptx"

DIAGRAMS = {
    "unit": RUN / "ai-agents-new-saas-unit-economics.png",
    "scorecard": RUN / "ai-agents-new-saas-workflow-scorecard.png",
    "shadow": RUN / "ai-agents-new-saas-shadow-human.png",
    "workflows": RUN / "ai-agents-new-saas-agentic-workflows.png",
    "wrapper": RUN / "ai-agents-new-saas-wrapper-saas.png",
    "pilot": RUN / "ai-agents-new-saas-pilot-productize.png",
}


def bullets(items, size=12.5):
    return [{"text": item, "bullet": True, "space_before": 6, "size": size} for item in items]


def note(d, s, text, dark=False):
    color = RGBColor(0xC6, 0xD0, 0xDA) if dark else d.b.MUTED
    d.text(s, text, d.M, Inches(6.73), Inches(10.4), Inches(0.20),
           size=8.2, color=color, shrink=True)


def pill(d, s, x, y, text, fill=None, color=None):
    fill = fill or d.b.LIGHT_TEAL
    color = color or d.b.DARK_TEAL
    d.rect(s, x, y, Inches(1.78), Inches(0.34), fill, radius=0.12)
    d.text(s, text, x + Inches(0.10), y + Inches(0.085), Inches(1.58), Inches(0.12),
           size=8.4, color=color, bold=True, align=PP_ALIGN.CENTER, shrink=True)


def action_header(d, s, title, kicker=None, dark=False):
    b = d.b
    title_color = b.WHITE if dark else b.NAVY
    kicker_color = b.TEAL if dark else b.DARK_TEAL
    if kicker:
        d.text(s, kicker.upper(), d.M, Inches(0.38), Inches(5.2), Inches(0.18),
               size=9.5, color=kicker_color, bold=True, shrink=True)
    d.text(s, title, d.M, Inches(0.66), d.CW, Inches(0.48),
           size=24, color=title_color, bold=True, shrink=True)
    d.rect(s, d.M, Inches(1.24), Inches(1.25), Inches(0.045), b.TEAL)


def diagram_slide(d, total, page, title, frame, image_key, takeaways):
    b = d.b
    s = d.slide(fill=b.WHITE)
    action_header(d, s, title, f"Recreated from {frame} in Excalidraw")
    d.picture_centered(s, DIAGRAMS[image_key], top=Inches(1.42), width=Inches(10.95),
                       max_bottom=Inches(6.20))
    d.rect(s, Inches(10.48), Inches(1.52), Inches(2.22), Inches(3.75),
           b.SOFT, line=b.GRID, radius=0.08)
    d.text(s, "Why this matters", Inches(10.70), Inches(1.76), Inches(1.78), Inches(0.18),
           size=10.5, color=b.NAVY, bold=True, shrink=True)
    d.text(s, bullets(takeaways, size=9.4), Inches(10.70), Inches(2.15),
           Inches(1.72), Inches(2.78), color=b.INK, shrink=True)
    note(d, s, f"Deck visual is rendered from ai-agents-new-saas-{image_key_name(image_key)}.excalidraw, not a source hyperframe screenshot.")
    d.footer(s, page, total)


def image_key_name(key):
    return {
        "unit": "unit-economics",
        "scorecard": "workflow-scorecard",
        "shadow": "shadow-human",
        "workflows": "agentic-workflows",
        "wrapper": "wrapper-saas",
        "pilot": "pilot-productize",
    }[key]


def card(d, s, x, y, w, h, title, body, fill=None):
    b = d.b
    d.rect(s, x, y, w, h, fill or b.SOFT, line=b.GRID, radius=0.08)
    d.text(s, title, x + Inches(0.16), y + Inches(0.16), w - Inches(0.32),
           Inches(0.22), size=11.2, color=b.NAVY, bold=True, shrink=True)
    d.text(s, body, x + Inches(0.16), y + Inches(0.52), w - Inches(0.32),
           h - Inches(0.64), size=9.8, color=b.INK, shrink=True)


def connector(d, s, x1, y1, x2, y2):
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = d.b.TEAL
    line.line.width = Pt(2)
    return line


def build():
    missing = [str(p) for p in DIAGRAMS.values() if not p.exists()]
    if missing:
        raise SystemExit("Missing Excalidraw exports:\n" + "\n".join(missing))

    d = Deck(footer="Video-to-Deck | Excalidraw recreation package | Draft")
    b = d.b
    total = 10

    s = d.slide(fill=b.NAVY)
    action_header(d, s, "AI Agents Are The New SaaS", "Video-to-deck rebuilt from hyperframes", dark=True)
    d.text(s,
           "A narrow workflow with a paycheck attached becomes an agent business when the operator workflow is observed, encoded, controlled, piloted, and then productized.",
           d.M, Inches(1.55), Inches(5.9), Inches(0.95), size=17, color=b.WHITE, shrink=True)
    d.text(s, "What changed in this rebuild", d.M, Inches(3.05), Inches(3.2), Inches(0.28),
           size=13, color=b.TEAL, bold=True)
    d.text(s, bullets([
        "No original hyperframe screenshots inserted.",
        "Every meaningful visual frame has an Excalidraw source file.",
        "Presenter-only frames are excluded.",
        "Markdown remains audit/source material, not the deliverable.",
    ], 11.2), d.M, Inches(3.42), Inches(4.8), Inches(1.85), size=11.2,
           color=b.SOFT, shrink=True)
    s.shapes.add_picture(str(DIAGRAMS["unit"]), Inches(6.65), Inches(1.55), width=Inches(5.85))
    s.shapes.add_picture(str(DIAGRAMS["scorecard"]), Inches(6.65), Inches(4.10), width=Inches(5.85))
    note(d, s, "Status: draft. Visual layer is Excalidraw-generated; editable .excalidraw sources are delivered alongside.", dark=True)
    d.footer(s, 1, total, dark=True)

    s = d.slide(fill=b.WHITE)
    action_header(d, s, "The Storyline Comes From The Video, Not From A Text Dump",
                  "Storyboard synthesis")
    d.text(s, "BLUF", d.M, Inches(1.62), Inches(1.2), Inches(0.24),
           size=11.5, color=b.TEAL, bold=True)
    d.text(s,
           "The opportunity is not “build an agent.” The opportunity is to find a narrow, repeated, expensive human workflow and productize it after the pattern repeats.",
           d.M, Inches(1.95), Inches(6.4), Inches(0.72), size=15.5,
           color=b.NAVY, bold=True, shrink=True)
    stages = [
        ("1", "Economic wedge", "Unit economics compress"),
        ("2", "Where to search", "Paid workflow scorecard"),
        ("3", "How to build", "Shadow the human"),
        ("4", "How to trust", "Agentic control room"),
        ("5", "How to sell", "Pilot first, productize second"),
    ]
    x0 = Inches(0.78)
    y = Inches(3.15)
    for i, (num, title, body) in enumerate(stages):
        x = x0 + Inches(i * 2.48)
        d.rect(s, x, y, Inches(2.05), Inches(1.05), b.SOFT if i % 2 == 0 else b.LIGHT_TEAL,
               line=b.GRID, radius=0.08)
        d.text(s, num, x + Inches(0.14), y + Inches(0.12), Inches(0.3), Inches(0.20),
               size=11, color=b.TEAL, bold=True)
        d.text(s, title, x + Inches(0.48), y + Inches(0.13), Inches(1.30), Inches(0.18),
               size=9.5, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.16), y + Inches(0.52), Inches(1.68), Inches(0.28),
               size=8.7, color=b.INK, shrink=True)
        if i < len(stages) - 1:
            connector(d, s, x + Inches(2.05), y + Inches(0.52), x + Inches(2.38), y + Inches(0.52))
    note(d, s, "Storyboard source: ai-agents-new-saas-storyboard.md; ai-analyst unavailable, strategy-consulting fallback used.")
    d.footer(s, 2, total)

    diagram_slide(d, total, 3, "Agents Win When Unit Economics Compress", "frame_0001",
                  "unit", [
                      "AI is compelling where labor cost is frequent and measurable.",
                      "Speed and cost both improve, creating a budget wedge.",
                      "The buyer case starts with operating economics, not novelty.",
                  ])
    diagram_slide(d, total, 4, "Pick A Workflow With A Paycheck Attached", "frame_0003",
                  "scorecard", [
                      "Search where frequency, pain, tools, and budget overlap.",
                      "The first niche should have repeatable jobs with a clear finish line.",
                      "A scorecard beats brainstorming generic agent ideas.",
                  ])
    diagram_slide(d, total, 5, "Shadow The Human Before You Build", "frame_0004",
                  "shadow", [
                      "The real spec lives in observed work, exceptions, and approvals.",
                      "Ten to twenty workflow runs reveal hidden rules.",
                      "The agent must know when it is stuck.",
                  ])
    diagram_slide(d, total, 6, "Agentic Workflows Need A Control Pattern", "frame_0006",
                  "workflows", [
                      "Agents require predefined operating patterns.",
                      "Sequential flow handles repeatable paths.",
                      "Hierarchical flow handles review and escalation.",
                  ])
    diagram_slide(d, total, 7, "The Wrapper Turns Work Into SaaS", "frame_0007",
                  "wrapper", [
                      "The agent does work; the wrapper creates trust.",
                      "Logs, approvals, settings, handoffs, analytics, and policy matter.",
                      "Control surfaces make the product buyable.",
                  ])
    diagram_slide(d, total, 8, "Sell The Pilot Like Labor, Then Productize", "frames_0008/0009",
                  "pilot", [
                      "Start as a service promise against a specific job.",
                      "Use simple pricing buyers already understand.",
                      "Productize only after the same workflow repeats.",
                  ])

    s = d.slide(fill=b.WHITE)
    action_header(d, s, "The 30-Day Path Turns Proof Into Product",
                  "Transcript synthesis")
    phases = [
        ("Days 1-3", "Pick one niche\nList 20 annoying jobs\nScore for pain + budget"),
        ("Days 4-7", "Shadow real operators\nExtract trigger/rules/escalation\nDefine success metric"),
        ("Week 2", "Build narrow pilot\nUse real tickets/calls/leads\nCreate control room"),
        ("Week 3", "Sell 2-3 pilots\nSame niche\nSame workflow\nSame promise"),
        ("Week 4", "Package the pattern\nLock repeatable rules\nPrepare SaaS wrapper"),
    ]
    x = Inches(0.72)
    for i, (title, body) in enumerate(phases):
        fill = b.LIGHT_TEAL if i in (1, 3) else b.SOFT
        card(d, s, x + Inches(i * 2.45), Inches(1.82), Inches(2.05), Inches(2.0),
             title, body, fill)
        if i < 4:
            connector(d, s, x + Inches(i * 2.45 + 2.05), Inches(2.82),
                      x + Inches(i * 2.45 + 2.36), Inches(2.82))
    d.rect(s, Inches(1.05), Inches(4.55), Inches(10.9), Inches(0.74), b.NAVY, radius=0.08)
    d.text(s, "Decision gate: productize only after the same painful workflow repeats across buyers.",
           Inches(1.35), Inches(4.82), Inches(10.25), Inches(0.18),
           size=12.2, color=b.WHITE, bold=True, align=PP_ALIGN.CENTER, shrink=True)
    note(d, s, "This slide is synthesized from transcript notes; it is not a missing hyperframe substitute.")
    d.footer(s, 9, total)

    s = d.slide(fill=b.WHITE)
    action_header(d, s, "Every Meaningful Hyperframe Has An Artifact",
                  "Delivery package map")
    rows = [
        ("0001", "unit economics", "ai-agents-new-saas-unit-economics.excalidraw", "slide 3"),
        ("0003", "workflow scorecard", "ai-agents-new-saas-workflow-scorecard.excalidraw", "slide 4"),
        ("0004", "shadow human", "ai-agents-new-saas-shadow-human.excalidraw", "slide 5"),
        ("0006", "agentic workflows", "ai-agents-new-saas-agentic-workflows.excalidraw", "slide 6"),
        ("0007", "wrapper SaaS", "ai-agents-new-saas-wrapper-saas.excalidraw", "slide 7"),
        ("0008/0009", "pilot to product", "ai-agents-new-saas-pilot-productize.excalidraw", "slide 8"),
    ]
    widths = [Inches(0.95), Inches(2.25), Inches(5.2), Inches(1.05)]
    x0, y0 = Inches(0.75), Inches(1.72)
    d.rect(s, x0, y0, sum(widths), Inches(0.34), b.NAVY, radius=0.05)
    headers = ["Frame", "Concept", "Editable source", "Deck"]
    x = x0
    for h, w in zip(headers, widths):
        d.text(s, h, x + Inches(0.08), y0 + Inches(0.085), w - Inches(0.14), Inches(0.12),
               size=8.8, color=b.WHITE, bold=True, shrink=True)
        x += w
    for i, row in enumerate(rows):
        y = y0 + Inches(0.42 + i * 0.47)
        fill = b.SOFT if i % 2 == 0 else b.WHITE
        d.rect(s, x0, y, sum(widths), Inches(0.38), fill, line=b.GRID)
        x = x0
        for val, w in zip(row, widths):
            d.text(s, val, x + Inches(0.08), y + Inches(0.09), w - Inches(0.14),
                   Inches(0.12), size=8.2, color=b.INK, bold=w == widths[0], shrink=True)
            x += w
    pill(d, s, Inches(9.95), Inches(1.82), "Screenshots excluded", RGBColor(0xF8, 0xE7, 0xEA), b.CORAL)
    pill(d, s, Inches(9.95), Inches(2.30), "Excalidraw source", b.LIGHT_TEAL, b.DARK_TEAL)
    pill(d, s, Inches(9.95), Inches(2.78), "PPTX packaged", RGBColor(0xF7, 0xEF, 0xD6), RGBColor(0x9A, 0x6A, 0x00))
    d.text(s,
           "The PPTX embeds rendered Excalidraw previews for presentation quality. Edit the actual drawings in the paired .excalidraw files, then re-export/rebuild.",
           Inches(0.78), Inches(5.30), Inches(9.6), Inches(0.52), size=12.3,
           color=b.NAVY, bold=True, shrink=True)
    note(d, s, "Manifest: ai-agents-new-saas-hyperframes.md. Visual spec: ai-agents-new-saas-visual-spec.json.")
    d.footer(s, 10, total)

    d.save(OUT)


if __name__ == "__main__":
    build()
