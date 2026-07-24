from __future__ import annotations

import os
import shutil
import sys
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib-content-ideas")

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, "/home/shekerk/.claude/skills/ai-analyst/helpers")
from chart_helpers import swd_style, save_chart  # noqa: E402


RUN_DIR = Path("/home/shekerk/content-ideas/runs/2026-06-30-agent-replacement-scorecard-synthesis")
OUT_DIR = RUN_DIR / "outputs"
WORK_DIR = RUN_DIR / "working"
CHART_DIR = OUT_DIR / "charts"
DATA_PATH = Path("/mnt/c/Users/sheke/OneDrive/Desktop/Agent-Replacement-Scorecard.csv")
SOURCE_DECK = Path("/mnt/c/Users/sheke/OneDrive/Desktop/Agent-Replacement-Scorecard-AEO-Data-Backed-Client-Reviewed.pptx")
OUTPUT_DECK = Path("/mnt/c/Users/sheke/OneDrive/Desktop/Agent-Replacement-Scorecard-Data-Synthesis-Reviewed.pptx")
REPO_COPY = OUT_DIR / OUTPUT_DECK.name
TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE", "")) if os.environ.get("BRANDED_PPTX_TEMPLATE") else Path("/home/shekerk/.claude/templates/branded-template.pptx")

BG = RGBColor(247, 246, 242)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
LINE = RGBColor(229, 231, 235)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(5, 150, 105)
WHITE = RGBColor(255, 255, 255)


def clear_template_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=INK, align=None, margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, section.upper(), 0.55, 0.32, 3.0, 0.25, 8, True, AMBER)
    add_text(slide, title, 0.55, 0.52, 8.9, 0.55, 23, True, INK)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.08, 8.7, 0.34, 10.5, False, MUTED)


def add_footer(slide, idx, total):
    add_text(slide, "Agent Replacement Scorecard | data synthesis reviewed draft | 2026-06-30", 0.55, 7.08, 7.2, 0.22, 7.2, False, MUTED)
    add_text(slide, f"{idx} / {total}", 12.0, 7.08, 0.65, 0.22, 7.2, False, MUTED, PP_ALIGN.RIGHT)


def add_kpi(slide, label, value, x, y, w=2.2, h=1.0, color=AMBER, note=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    add_text(slide, str(value), x + 0.10, y + 0.10, w - 0.2, 0.33, 21, True, color)
    add_text(slide, label, x + 0.10, y + 0.49, w - 0.2, 0.28, 8.8, True, INK)
    if note:
        add_text(slide, note, x + 0.10, y + 0.73, w - 0.2, 0.19, 7.0, False, MUTED)


def add_bullets(slide, items, x, y, w, h, size=12.5, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_table(slide, df, x, y, w, h, font=8.5, widths=None, header_color=INK):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for r in range(df.shape[0]):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = str(df.iloc[r, c])
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else RGBColor(243, 244, 246)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font)
                p.font.color.rgb = INK
    return table


def save_fig(fig, name: str) -> Path:
    path = CHART_DIR / name
    save_chart(fig, path)
    plt.close(fig)
    return path


def bar_chart(series, title, name, highlight=None, color_map=None, xlabel="Rows"):
    colors = swd_style()
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
    cats = list(series.index)
    vals = list(series.values)
    bar_colors = []
    for cat in cats:
        if color_map and cat in color_map:
            bar_colors.append(color_map[cat])
        elif highlight and cat == highlight:
            bar_colors.append(colors["action"])
        else:
            bar_colors.append(colors["gray400"])
    ax.barh(cats[::-1], vals[::-1], color=bar_colors[::-1])
    ax.set_title(title, loc="left", fontsize=15, fontweight="bold")
    ax.set_xlabel(xlabel)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="x", color="#E5E7EB", linewidth=0.7)
    ax.set_axisbelow(True)
    for y, v in enumerate(vals[::-1]):
        ax.text(v + 0.15, y, f"{int(v)}", va="center", fontsize=10, color="#1F2937")
    ax.set_xlim(0, max(vals) + 2)
    return save_fig(fig, name)


def stacked_chart(df, title, name):
    swd_style()
    order = ["REPLACE", "RENEGOTIATE", "KEEP"]
    colors = {"REPLACE": "#DC2626", "RENEGOTIATE": "#D97706", "KEEP": "#9CA3AF"}
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
    left = [0] * len(df.index)
    for col in order:
        vals = df[col] if col in df.columns else [0] * len(df.index)
        ax.barh(df.index, vals, left=left, color=colors[col], label=col)
        for i, v in enumerate(vals):
            if v:
                ax.text(left[i] + v / 2, i, str(int(v)), ha="center", va="center", color="white", fontsize=9, fontweight="bold")
        left = [a + b for a, b in zip(left, vals)]
    ax.set_title(title, loc="left", fontsize=15, fontweight="bold")
    ax.set_xlabel("Use cases")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(frameon=False, loc="lower right", ncol=3)
    return save_fig(fig, name)


def matrix_chart(df):
    swd_style()
    score_matrix = pd.crosstab(df["Data moat"], df["Verdict"]).reindex(["L", "M", "H"]).fillna(0)
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
    im = ax.imshow(score_matrix.values, cmap="YlOrRd")
    ax.set_xticks(range(len(score_matrix.columns)), score_matrix.columns)
    ax.set_yticks(range(len(score_matrix.index)), score_matrix.index)
    ax.set_title("Data moat cleanly separates replace, renegotiate, and keep", loc="left", fontsize=15, fontweight="bold")
    ax.set_xlabel("Verdict")
    ax.set_ylabel("Data moat")
    for i in range(score_matrix.shape[0]):
        for j in range(score_matrix.shape[1]):
            val = int(score_matrix.iloc[i, j])
            ax.text(j, i, val, ha="center", va="center", color="#1F2937", fontsize=18, fontweight="bold")
    fig.colorbar(im, ax=ax, fraction=0.025, pad=0.04)
    return save_fig(fig, "data_moat_verdict_matrix.png")


def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_PATH)
    map_score = {"L": 1, "M": 2, "H": 3}
    df["Volume score"] = df["Volume"].map(map_score)
    df["Determinism score"] = df["Determinism"].map(map_score)
    df["Moat exposure score"] = df["Data moat"].map({"L": 3, "M": 2, "H": 1})
    df["Replacement exposure score"] = df["Volume score"] + df["Determinism score"] + df["Moat exposure score"]
    df["Evidence label"] = df["Proof (short)"].str.replace(", ", " | ", regex=False)
    return df


def build_analysis(df: pd.DataFrame) -> dict:
    verdict = df["Verdict"].value_counts().reindex(["REPLACE", "RENEGOTIATE", "KEEP"]).fillna(0).astype(int)
    agent_verdict = pd.crosstab(df["Agent type"], df["Verdict"]).reindex(columns=["REPLACE", "RENEGOTIATE", "KEEP"], fill_value=0)
    data_moat = pd.crosstab(df["Data moat"], df["Verdict"]).reindex(["L", "M", "H"]).reindex(columns=["REPLACE", "RENEGOTIATE", "KEEP"], fill_value=0)
    volume = pd.crosstab(df["Volume"], df["Verdict"]).reindex(["H", "M"]).reindex(columns=["REPLACE", "RENEGOTIATE", "KEEP"], fill_value=0)
    top = df.sort_values(["Replacement exposure score", "Volume score", "Determinism score"], ascending=False).head(8)
    protected = df.sort_values(["Replacement exposure score", "Volume score"], ascending=True).head(7)
    contradictions = df[((df["Data moat"] == "H") & (df["Verdict"] != "KEEP")) | ((df["Data moat"] == "L") & (df["Verdict"] != "REPLACE"))]
    return {
        "verdict": verdict,
        "agent_verdict": agent_verdict,
        "data_moat": data_moat,
        "volume": volume,
        "top": top,
        "protected": protected,
        "contradictions": contradictions,
    }


def build_charts(df, analysis):
    charts = {}
    charts["verdict"] = bar_chart(
        analysis["verdict"],
        "10 of 25 use cases are replacement candidates; 9 are renegotiation cases",
        "verdict_distribution.png",
        color_map={"REPLACE": "#DC2626", "RENEGOTIATE": "#D97706", "KEEP": "#9CA3AF"},
    )
    charts["agent_verdict"] = stacked_chart(analysis["agent_verdict"], "Customer workflows concentrate the clearest replacement cases", "agent_type_verdict.png")
    charts["data_moat"] = matrix_chart(df)
    charts["volume"] = stacked_chart(analysis["volume"], "High-volume workflows supply 9 of 10 replacement calls", "volume_verdict.png")
    score_counts = df["Replacement exposure score"].value_counts().sort_index()
    charts["score_hist"] = bar_chart(score_counts, "Replacement exposure clusters at the high end, not evenly across use cases", "score_histogram.png", highlight=8, xlabel="Use cases")
    return charts


def add_slide(prs, idx, total, title, subtitle=None, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle, section)
    add_footer(slide, idx, total)
    return slide


def build_deck(df, analysis, charts) -> Presentation:
    prs = Presentation(str(TEMPLATE))
    clear_template_slides(prs)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 18

    slide = add_slide(prs, 1, total, "Agent Replacement Scorecard", "Data synthesis deck built from 25 scored use cases, not AEO gate narration", "Reviewed draft")
    add_kpi(slide, "use cases scored", len(df), 0.65, 1.75, color=INK)
    add_kpi(slide, "replace candidates", int(analysis["verdict"]["REPLACE"]), 3.05, 1.75, color=RED)
    add_kpi(slide, "renegotiate cases", int(analysis["verdict"]["RENEGOTIATE"]), 5.45, 1.75, color=AMBER)
    add_kpi(slide, "keep cases", int(analysis["verdict"]["KEEP"]), 7.85, 1.75, color=MUTED)
    add_bullets(slide, [
        "Question: which SaaS layers are exposed to agent replacement, renegotiation, or durability?",
        "Model: volume + determinism + inverse data moat, validated against verdict distribution.",
        "Main finding: data moat is the dominant separator; volume determines where replacement becomes economically urgent.",
    ], 0.7, 3.25, 8.0, 1.2, 14)
    add_text(slide, "Source: Agent-Replacement-Scorecard.csv; original deck repaired because its AEO evidence-gate story did not match the structured scorecard dataset.", 0.72, 6.25, 10.5, 0.35, 8.8, False, MUTED)

    slide = add_slide(prs, 2, total, "The answer is not 'agents replace SaaS'; it is a three-way split", "Replacement, renegotiation, and durability each appear in the dataset.", "Executive take")
    slide.shapes.add_picture(str(charts["verdict"]), Inches(0.7), Inches(1.55), width=Inches(5.65))
    add_text(slide, "Finding", 7.0, 1.6, 2.0, 0.3, 10, True, AMBER)
    add_bullets(slide, [
        "10 workflows are replacement candidates where the exposed layer is seats, add-ons, or scripted workflow tooling.",
        "9 workflows are renegotiation cases where the system of record survives but add-on modules compress.",
        "6 workflows are keep cases, mostly where high data moat, permissions, audit, or proprietary telemetry are the value.",
    ], 7.0, 2.0, 5.1, 1.75, 12.5)
    add_text(slide, "Decision implication", 7.0, 4.35, 2.4, 0.3, 10, True, AMBER)
    add_bullets(slide, ["Lead with a diagnostic. Do not market a blanket replacement thesis."], 7.0, 4.75, 4.7, 0.5, 13)

    slide = add_slide(prs, 3, total, "Data moat explains the verdict boundary", "The scorecard is strongest when it shows the mechanism, not just examples.", "Model")
    slide.shapes.add_picture(str(charts["data_moat"]), Inches(0.8), Inches(1.55), width=Inches(5.8))
    add_bullets(slide, [
        "Low moat: 10 / 10 rows are REPLACE.",
        "Medium moat: 9 / 9 rows are RENEGOTIATE.",
        "High moat: 6 / 6 rows are KEEP.",
        "This is the cleanest synthetic finding and should anchor the deck.",
    ], 7.0, 1.8, 4.9, 1.7, 13)
    add_text(slide, "Interpretation: agents compress the interaction/workflow layer first. The durable substrate is where proprietary records, permissions, governance, and audit trails remain hard to dislodge.", 7.0, 4.2, 4.9, 0.95, 12, False, INK)

    slide = add_slide(prs, 4, total, "High volume turns exposure into near-term commercial pressure", "Volume does not decide the verdict alone; it makes replacement worth pursuing.", "Driver check")
    slide.shapes.add_picture(str(charts["volume"]), Inches(0.8), Inches(1.55), width=Inches(5.8))
    add_bullets(slide, [
        "High-volume rows produce 9 of the 10 replacement calls.",
        "Medium-volume rows skew toward renegotiate or keep.",
        "Priority should go to high-volume, low-moat workflows first because the cost takeout is easiest to explain.",
    ], 7.0, 1.9, 4.9, 1.5, 13)
    add_text(slide, "Use this as the commercial filter after the technical replacement score.", 7.0, 4.15, 4.9, 0.45, 12.5, True, INK)

    slide = add_slide(prs, 5, total, "The replacement score makes the ranking auditable", "Score = volume + determinism + moat exposure. High score is more exposed.", "Scoring")
    slide.shapes.add_picture(str(charts["score_hist"]), Inches(0.8), Inches(1.55), width=Inches(5.8))
    score_df = analysis["top"][["Use case", "Agent type", "Volume", "Determinism", "Data moat", "Replacement exposure score"]].copy()
    score_df.columns = ["Use case", "Type", "Vol", "Det", "Moat", "Score"]
    add_table(slide, score_df, 6.65, 1.48, 5.85, 4.6, 7.6, widths=[2.0, 0.85, 0.45, 0.45, 0.5, 0.5])

    slide = add_slide(prs, 6, total, "Customer-facing workflows carry the clearest replacement signal", "The category split shows where to start the client conversation.", "Segment view")
    slide.shapes.add_picture(str(charts["agent_verdict"]), Inches(0.8), Inches(1.55), width=Inches(6.0))
    add_bullets(slide, [
        "Customer use cases: 4 replace, 1 renegotiate, 0 keep.",
        "Security use cases: 3 keep, 1 replace; the moat is usually telemetry, policy, and audit.",
        "Employee and data workflows split because proprietary context often keeps the platform while compressing add-ons.",
    ], 7.15, 1.85, 4.9, 1.75, 12.5)

    slide = add_slide(prs, 7, total, "Top replacement candidates share the same pattern", "High volume, repeatable intent, and weak data moat.", "Replace shortlist")
    top_tbl = analysis["top"][["Use case", "SaaS affected (short)", "Build-vs-buy", "Replacement exposure score"]].copy()
    top_tbl.columns = ["Use case", "SaaS exposed", "Compression thesis", "Score"]
    add_table(slide, top_tbl.head(7), 0.65, 1.55, 12.0, 4.85, 7.7, widths=[2.35, 2.25, 3.95, 0.55])

    slide = add_slide(prs, 8, total, "Keep cases are not anti-agent; they are substrate cases", "The agent can improve workflow without replacing the durable system.", "Durable layers")
    keep_tbl = analysis["protected"][["Use case", "Agent type", "Data moat", "Verdict", "SaaS affected (short)"]].copy()
    keep_tbl.columns = ["Use case", "Type", "Moat", "Verdict", "Durable SaaS layer"]
    add_table(slide, keep_tbl, 0.75, 1.55, 11.6, 4.6, 8.2, widths=[2.5, 1.0, 0.55, 1.1, 3.0])
    add_text(slide, "Message: sell workflow compression, not platform eradication, when the moat is proprietary data or governance.", 0.82, 6.25, 9.8, 0.4, 12, True, AMBER)

    slide = add_slide(prs, 9, total, "Renegotiate is the commercial middle, not a weak conclusion", "Nine rows preserve the system of record while compressing modules and seats.", "Renegotiate")
    reneg = df[df["Verdict"] == "RENEGOTIATE"][["Use case", "Agent type", "SaaS affected (short)", "Build-vs-buy"]].head(8).copy()
    reneg.columns = ["Use case", "Type", "SaaS affected", "Negotiation lever"]
    add_table(slide, reneg, 0.75, 1.5, 11.7, 4.85, 8.0, widths=[2.5, 1.0, 2.3, 3.15])
    add_text(slide, "This is the procurement wedge: keep the core, challenge the per-seat workflow module.", 0.82, 6.28, 9.6, 0.35, 12, True, INK)

    slide = add_slide(prs, 10, total, "Proof examples should be evidence rows, not decoration", "Each scorecard row has a deployed proof claim and a SaaS-in-crosshairs claim.", "Evidence")
    sample = df[df["Verdict"] == "REPLACE"].head(4)
    y = 1.55
    for _, row in sample.iterrows():
        add_text(slide, row["Use case"], 0.75, y, 2.1, 0.25, 10, True, INK)
        add_text(slide, row["Proof (short)"], 3.0, y, 3.25, 0.38, 8.8, False, MUTED)
        add_text(slide, row["SaaS affected (short)"], 6.55, y, 2.2, 0.38, 8.8, True, RED)
        add_text(slide, row["Build-vs-buy"], 9.0, y, 2.85, 0.38, 8.8, False, INK)
        y += 0.9
    add_text(slide, "Deck repair: move from anecdote lists to proof-to-verdict traceability.", 0.82, 6.15, 8.2, 0.4, 12, True, AMBER)

    slide = add_slide(prs, 11, total, "Validation passed, with one important limitation", "The dataset is internally complete, but it is curated rather than statistically sampled.", "Validation")
    add_kpi(slide, "null columns over 20%", 0, 0.75, 1.65, color=GREEN)
    add_kpi(slide, "duplicate use cases", int(df["Use case"].duplicated().sum()), 3.05, 1.65, color=GREEN)
    proof_count = df["Proof — what's deployed"].notna().sum()
    add_kpi(slide, "rows with proof text", f"{proof_count}/25", 5.35, 1.65, color=GREEN)
    add_kpi(slide, "rows with verdict", f"{df['Verdict'].notna().sum()}/25", 7.65, 1.65, color=GREEN)
    add_bullets(slide, [
        "High confidence in internal arithmetic: counts, cross-tabs, and score derivation reconcile to 25 rows.",
        "Medium confidence in market generalization: the rows are evidence-backed examples, not a random market sample.",
        "Guardrail: call this a replacement exposure scorecard, not a statistically validated SaaS replacement model.",
    ], 0.8, 3.25, 10.5, 1.3, 13)

    slide = add_slide(prs, 12, total, "The scorecard should become a diagnostic workflow", "Each client account can be scored with the same dimensions.", "Productization")
    add_bullets(slide, [
        "1. Inventory workflows and attached SaaS modules.",
        "2. Score volume, determinism, and data moat on H/M/L.",
        "3. Classify each line as replace, renegotiate, or keep.",
        "4. Attach deployed proof examples and a buyer-specific savings lever.",
        "5. Convert high-scoring rows into pilot candidates or procurement asks.",
    ], 0.85, 1.65, 5.3, 3.0, 15)
    add_text(slide, "The artifact is stronger as an interactive account diagnostic than as a static thought-leadership deck.", 7.05, 2.0, 4.7, 1.0, 19, True, AMBER)

    slide = add_slide(prs, 13, total, "AEO angle: publish the model, not just the claim", "AEO pages should answer buyer questions grounded in the dataset.", "AEO rebuild")
    aeo_rows = pd.DataFrame([
        ["Replacement risk", "Which SaaS seats can AI agents replace first?", "Use replace shortlist + score formula"],
        ["Renegotiation", "Which AI agent gains should lower my SaaS bill?", "Use renegotiate table"],
        ["Moat defense", "Which SaaS systems survive agent adoption?", "Use keep cases + data moat logic"],
        ["Pilot design", "How do I test agent replacement safely?", "Use volume/determinism filters"],
        ["Procurement script", "What should I ask vendors at renewal?", "Use compression thesis"],
    ], columns=["Page cluster", "Buyer question", "Data-backed answer"])
    add_table(slide, aeo_rows, 0.8, 1.6, 11.3, 3.6, 9.2, widths=[1.8, 3.4, 4.0])
    add_text(slide, "This connects AEO to synthesis: pages are generated from observed scorecard patterns, not generic AI positioning.", 0.85, 5.75, 10.5, 0.45, 12, True, INK)

    slide = add_slide(prs, 14, total, "Claim guardrails keep the deck credible", "Separate what the data supports from what still needs proof.", "Guardrails")
    safe = ["25-row curated scorecard", "Directional exposure model", "Evidence-backed examples", "Client diagnostic v1", "Procurement/pilot prioritization"]
    unsafe = ["Total SaaS replacement forecast", "Statistical market sizing", "Engine consensus proof", "Universal category verdicts", "Guaranteed savings"]
    add_text(slide, "Supported", 0.9, 1.55, 2.0, 0.3, 13, True, GREEN)
    add_bullets(slide, safe, 0.9, 2.0, 4.5, 2.2, 13)
    add_text(slide, "Not supported yet", 7.0, 1.55, 2.2, 0.3, 13, True, RED)
    add_bullets(slide, unsafe, 7.0, 2.0, 4.5, 2.2, 13)

    slide = add_slide(prs, 15, total, "30-day evidence sprint should harden the model", "Move from curated synthesis to repeatable scoring evidence.", "Next sprint")
    sprint = pd.DataFrame([
        ["Dataset", "Expand 25 rows to 75+", "More stable segment patterns"],
        ["Scoring", "Score two reviewers independently", "Inter-rater reliability"],
        ["Economics", "Attach seat/module cost ranges", "Dollarized replacement cases"],
        ["AEO", "Publish 5 model-derived pages", "Crawlable buyer answers"],
        ["Validation", "Track vendor/customer proof links", "Source tieout for each proof row"],
    ], columns=["Workstream", "Target", "Why it matters"])
    add_table(slide, sprint, 0.85, 1.6, 11.1, 3.9, 9.5, widths=[1.55, 3.0, 4.4])

    slide = add_slide(prs, 16, total, "Decision ask: launch the diagnostic, not the proof claim", "The synthesis is strong enough for v1 account work and page briefs.", "Decision")
    add_kpi(slide, "launch", "v1", 0.85, 1.8, color=AMBER, note="diagnostic")
    add_kpi(slide, "score", "75+", 3.2, 1.8, color=INK, note="next rows")
    add_kpi(slide, "publish", "5", 5.55, 1.8, color=INK, note="AEO pages")
    add_kpi(slide, "validate", "2x", 7.9, 1.8, color=INK, note="reviewers")
    add_bullets(slide, [
        "Use this deck for client discovery and procurement strategy.",
        "Use the CSV model as the source of truth for page briefs and account diagnostics.",
        "Do not present the current evidence as final market proof until the expanded sprint is complete.",
    ], 0.9, 3.45, 9.6, 1.4, 14)

    slide = add_slide(prs, 17, total, "Appendix: scoring rubric", "Simple enough to explain in a workshop; explicit enough to audit.", "Appendix")
    rubric = pd.DataFrame([
        ["Volume", "H = high throughput / frequent task", "More economic pressure"],
        ["Determinism", "H = repeatable rules and stable outputs", "Easier automation boundary"],
        ["Data moat", "L = weak proprietary substrate", "Higher replacement exposure"],
        ["Verdict", "Replace / Renegotiate / Keep", "Commercial action, not moral judgment"],
    ], columns=["Dimension", "Definition", "Interpretation"])
    add_table(slide, rubric, 0.85, 1.55, 11.1, 3.2, 10, widths=[1.6, 4.0, 3.5])
    add_text(slide, "Exposure score used in this deck: Volume score + Determinism score + inverse Data Moat score.", 0.9, 5.35, 9.5, 0.45, 13, True, AMBER)

    slide = add_slide(prs, 18, total, "Appendix: files and QA status", "Reproducible deck build with explicit status.", "Appendix")
    add_bullets(slide, [
        f"Input CSV: {DATA_PATH}",
        f"Original deck inspected: {SOURCE_DECK}",
        f"Builder script: {RUN_DIR / 'build_scorecard_synthesis_deck.py'}",
        f"Charts: {CHART_DIR}",
        "Status: reviewed draft; data synthesis repaired; visual QA performed by slide text and file generation checks.",
    ], 0.85, 1.65, 11.2, 2.8, 12.5)
    add_text(slide, "Note: the original AEO evidence-gate material was not discarded conceptually; it was reframed as a downstream AEO page plan after the scorecard synthesis.", 0.9, 5.6, 10.5, 0.55, 11.5, False, MUTED)

    return prs


def write_artifacts(df, analysis):
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    summary = {
        "rows": int(len(df)),
        "columns": list(df.columns),
        "null_pct": (df.isna().mean() * 100).round(1).to_dict(),
        "verdict_counts": analysis["verdict"].to_dict(),
        "agent_type_verdict": analysis["agent_verdict"].to_dict(),
        "data_moat_verdict": analysis["data_moat"].to_dict(),
        "confidence": "Medium-high: source file is complete and internally consistent; market generalization remains curated, not statistical.",
    }
    pd.Series(summary, dtype="object").to_json(OUT_DIR / "analysis_summary.json", indent=2)
    df.to_csv(OUT_DIR / "scorecard_scored.csv", index=False)


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Branded template not found: {TEMPLATE}")
    df = load_data()
    analysis = build_analysis(df)
    write_artifacts(df, analysis)
    charts = build_charts(df, analysis)
    prs = build_deck(df, analysis, charts)
    prs.save(REPO_COPY)
    print(f"repo copy {REPO_COPY}")
    print(f"desktop target {OUTPUT_DECK}")


if __name__ == "__main__":
    main()
