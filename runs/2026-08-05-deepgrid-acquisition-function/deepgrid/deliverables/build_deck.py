#!/usr/bin/env python3
"""Build the DeepGrid client-acquisition strategy deck.

Charts are NATIVE python-pptx charts (editable in PowerPoint, Excel-backed) —
matplotlib is not installed on this host, and native charts are the better
artifact anyway.

Kept in the run folder so QA fixes stay reproducible (CLAUDE.md PPTX gate #8).
"""
import sys
from pathlib import Path

KIT = Path("/home/sheke/content-ideas/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(KIT))

from pptx.util import Inches, Pt, Emu                      # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR            # noqa: E402
from pptx.chart.data import CategoryChartData              # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION  # noqa: E402
from pptxkit import Deck, Brand, hx                        # noqa: E402

OUT = Path(__file__).resolve().parent
d = Deck(footer="DeepGrid Semi — client acquisition operating system · Level 3 build · 5 Aug 2026")
b = d.b
YRS = ["FY27", "FY28", "FY29", "FY30", "FY31", "FY32"]
TOTAL = 15
_page = {"n": 0}


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def page(s, dark=False):
    _page["n"] += 1
    d.footer(s, _page["n"], TOTAL, dark=dark)


def style_chart(gf, *, colors, legend=True, num_fmt='0.0"%"', font=9,
                gap=60, smooth=False, markers=False):
    ch = gf.chart
    ch.font.size = Pt(font)
    ch.font.name = b.FONT
    ch.has_title = False
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(font)
    else:
        ch.has_legend = False
    for i, ser in enumerate(ch.series):
        col = colors[i % len(colors)]
        if ch.chart_type in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
            ser.format.line.color.rgb = col
            ser.format.line.width = Pt(2.5)
            ser.smooth = smooth
        else:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = col
    try:
        ch.value_axis.has_major_gridlines = True
        ch.value_axis.major_gridlines.format.line.color.rgb = b.GRID
        ch.value_axis.major_gridlines.format.line.width = Pt(0.5)
        ch.value_axis.tick_labels.number_format = num_fmt
        ch.value_axis.tick_labels.number_format_is_linked = False
        ch.value_axis.tick_labels.font.size = Pt(font)
        ch.value_axis.format.line.fill.background()
        ch.category_axis.tick_labels.font.size = Pt(font)
        ch.category_axis.format.line.color.rgb = b.GRID
    except Exception:
        pass
    try:
        ch.plots[0].gap_width = gap
    except Exception:
        pass
    return ch


def add_chart(s, kind, cats, series, left, top, w, h, **kw):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(kind, left, top, w, h, cd)
    return style_chart(gf, **kw)


def table(s, headers, rows, left, top, width, *, col_w, row_h=Inches(0.34),
          head_fill=None, highlight=None, size=11.5, left_cols=(0,)):
    """Lightweight grid from rects + text. highlight: set of row indices.
    left_cols: column indices rendered left-aligned — use for prose cells,
    since centring a long sentence in a table cell is hard to scan."""
    head_fill = head_fill or b.NAVY
    x = left
    d.rect(s, left, top, width, row_h, head_fill)
    for i, htxt in enumerate(headers):
        d.text(s, htxt, x + Inches(0.10), top + Inches(0.05), col_w[i] - Inches(0.16),
               row_h - Inches(0.08), size=size, color=b.WHITE, bold=True, shrink=True,
               align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += col_w[i]
    y = top + row_h
    for r, row in enumerate(rows):
        hl = highlight and r in highlight
        d.rect(s, left, y, width, row_h,
               hx("FFF4D6") if hl else (b.SOFT if r % 2 == 0 else b.WHITE))
        if hl:
            d.rect(s, left, y, Inches(0.055), row_h, b.GOLD)
        x = left
        for i, cell in enumerate(row):
            d.text(s, str(cell), x + Inches(0.10), y + Inches(0.05), col_w[i] - Inches(0.16),
                   row_h - Inches(0.08), size=size, color=b.INK,
                   bold=bool(hl and i == 0), shrink=True,
                   align=PP_ALIGN.LEFT if i in left_cols else PP_ALIGN.CENTER)
            x += col_w[i]
        y += row_h
    return y


def cards(s, items, top, *, cols=3, height=Inches(1.6), gap=Inches(0.22), size=11.5):
    cw = (d.CW - gap * (cols - 1)) / cols
    for i, (title, body) in enumerate(items):
        r, c = divmod(i, cols)
        x = d.M + c * (cw + gap)
        y = top + r * (height + gap)
        d.rect(s, x, y, cw, height, b.WHITE, line=b.GRID, radius=0.06, shadow=True)
        d.rect(s, x, y, cw, Inches(0.05), b.TEAL)
        d.text(s, title, x + Inches(0.18), y + Inches(0.20), cw - Inches(0.36),
               Inches(0.46), size=13, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.18), y + Inches(0.70), cw - Inches(0.36),
               height - Inches(0.86), size=size, color=b.MUTED, shrink=True)


def kpi(s, items, top, *, height=Inches(1.5)):
    n = len(items)
    gap = Inches(0.22)
    cw = (d.CW - gap * (n - 1)) / n
    for i, (num, label, note, col) in enumerate(items):
        x = d.M + i * (cw + gap)
        d.rect(s, x, top, cw, height, b.NAVY, radius=0.05)
        d.text(s, num, x + Inches(0.16), top + Inches(0.16), cw - Inches(0.32), Inches(0.62),
               size=32, color=col, bold=True, font=b.FONT_H, shrink=True)
        d.text(s, label, x + Inches(0.16), top + Inches(0.80), cw - Inches(0.32), Inches(0.30),
               size=11.5, color=b.WHITE, bold=True, shrink=True)
        d.text(s, note, x + Inches(0.16), top + Inches(1.10), cw - Inches(0.32), Inches(0.32),
               size=9.5, color=b.LIGHT_TEAL, shrink=True)


def kpi_at(s, items, left, top, width, *, height=Inches(1.4), gap=Inches(0.18)):
    """KPI row confined to an explicit box — use whenever a chart shares the slide,
    so the cards cannot be laid over the plot area."""
    n = len(items)
    cw = (width - gap * (n - 1)) / n
    for i, (num, label, note, col) in enumerate(items):
        x = left + i * (cw + gap)
        d.rect(s, x, top, cw, height, b.NAVY, radius=0.05)
        d.text(s, num, x + Inches(0.15), top + Inches(0.14), cw - Inches(0.30), Inches(0.56),
               size=27, color=col, bold=True, font=b.FONT_H, shrink=True)
        d.text(s, label, x + Inches(0.15), top + Inches(0.73), cw - Inches(0.30), Inches(0.28),
               size=11, color=b.WHITE, bold=True, shrink=True)
        d.text(s, note, x + Inches(0.15), top + Inches(1.01), cw - Inches(0.30), Inches(0.28),
               size=9.5, color=b.LIGHT_TEAL, shrink=True)


def divider(s, num, title, subtitle):
    d.rect(s, 0, 0, d.W, d.H, b.NAVY)
    d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
    d.text(s, f"SECTION {num}", Inches(0.85), Inches(2.6), Inches(6), Inches(0.4),
           size=14, color=b.TEAL, bold=True)
    d.text(s, title, Inches(0.85), Inches(3.05), Inches(11), Inches(1.0),
           size=42, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    d.rect(s, Inches(0.85), Inches(4.15), Inches(1.6), Inches(0.05), b.TEAL)
    d.text(s, subtitle, Inches(0.85), Inches(4.42), Inches(10.5), Inches(0.6),
           size=15, color=b.LIGHT_TEAL, shrink=True)


def badge(s, txt, left, top, fill, *, w=Inches(1.5), fg=None):
    d.rect(s, left, top, w, Inches(0.30), fill, radius=0.3)
    d.text(s, txt, left, top + Inches(0.045), w, Inches(0.24), size=9.5,
           color=fg or b.WHITE, bold=True, align=PP_ALIGN.CENTER)




ASSET = Path(__file__).resolve().parent / "assets"


def photo_band(s, src, left, top, width, height):
    """Place an image cropped to fill a band — used for divider side panels."""
    from PIL import Image
    tgt = ASSET / "fitted" / f"band_{src.stem}_{int(width)}x{int(height)}.png"
    tgt.parent.mkdir(parents=True, exist_ok=True)
    if not tgt.exists():
        im = Image.open(src).convert("RGB")
        want = width / height
        have = im.width / im.height
        if have > want:
            nw = int(im.height * want)
            im = im.crop(((im.width - nw) // 2, 0, (im.width + nw) // 2, im.height))
        else:
            nh = int(im.width / want)
            im = im.crop((0, (im.height - nh) // 2, im.width, (im.height + nh) // 2))
        im.save(tgt)
    s.shapes.add_picture(str(tgt), left, top, width, height)


def photo(s, name, left, top, width, height):
    """Place a harvested product photo, cropped to fill the box."""
    from PIL import Image
    src = ASSET / "product" / f"{name}.png"
    if not src.exists():
        d.rect(s, left, top, width, height, b.SOFT, line=b.GRID)
        return
    tgt = ASSET / "fitted" / f"{name}_{int(width)}x{int(height)}.png"
    tgt.parent.mkdir(parents=True, exist_ok=True)
    if not tgt.exists():
        im = Image.open(src).convert("RGB")
        want = width / height
        have = im.width / im.height
        if have > want:
            nw = int(im.height * want)
            im = im.crop(((im.width - nw) // 2, 0, (im.width + nw) // 2, im.height))
        else:
            nh = int(im.width / want)
            im = im.crop((0, (im.height - nh) // 2, im.width, (im.height + nh) // 2))
        im.save(tgt)
    s.shapes.add_picture(str(tgt), left, top, width, height)
# ── COVER ───────────────────────────────────────────────────────────────────
s = d.slide(fill=b.NAVY)
bg = ASSET / "harvested" / "deck_image1.png"
if bg.exists():
    s.shapes.add_picture(str(bg), 0, 0, d.W, d.H)
d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
d.text(s, "CLIENT ACQUISITION OPERATING SYSTEM", Inches(0.85), Inches(1.80),
       Inches(9.0), Inches(0.4), size=13, color=b.TEAL, bold=True)
d.text(s, "How DeepGrid wins clients", Inches(0.85), Inches(2.30), Inches(8.6),
       Inches(1.2), size=44, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
d.rect(s, Inches(0.85), Inches(3.75), Inches(1.6), Inches(0.05), b.TEAL)
d.text(s, "Who we sell to, what we sell, what we may claim, and how a lead travels — "
          "built as files the team and the tooling both read",
       Inches(0.85), Inches(4.05), Inches(9.2), Inches(0.75), size=15,
       color=b.LIGHT_TEAL, shrink=True)
for i, chip in enumerate(["2 segments now", "price by function",
                          "9 claims retired", "1 lead state machine"]):
    x = Inches(0.85) + i * Inches(2.6)
    d.rect(s, x, Inches(5.30), Inches(2.45), Inches(0.44), b.NAVY_2, radius=0.25, line=b.ACCENT)
    d.text(s, chip, x, Inches(5.42), Inches(2.45), Inches(0.26), size=10,
           color=b.LIGHT_TEAL, align=PP_ALIGN.CENTER)
page(s, dark=True)
notes(s, """
This deck documents an operating system, not a strategy opinion. Everything in it is written
into six context files that sit in the run folder — the deck is the readable version, the
files are what actually removes work.
The test that matters is not whether this deck is persuasive. It is whether a cold session,
given only those six files, produces outbound a founder would send.
""")

# ── THE TENSION ─────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "The knowledge exists. It is just not usable.",
         "Four documents, fourteen contradictions, nine retired claims — and no operating layer")
rows = [
    ["FY32 gross margin", "84% (IM) · 88% blended (July) · 72% kit (BP-1A)", "72% kit only; 88% needs 120% on non-kit"],
    ["Defence revenue", "“₹1 Cr” (IM, July deck)", "₹23.01L delivered, different entity, MCEME IP"],
    ["Mandate date", "“Live since April 2026”", "GSR 184(E) was a draft; AEBS binds 1 Jan 2027"],
    ["Compute claim", "“39.3 TOPS measured on FPGA”", "Artix-7 ceilings ~1.8 TOPS. It is a derivation"],
    ["Price advantage", "“3–7×” (IM) · “12.9×” (July)", "~2.8–6.5× vs imports; ~2.3× loaded on cost"],
]
table(s, ["Where they disagree", "What is in circulation", "What survives verification"], rows,
      d.M, Inches(1.82), d.CW, col_w=[Inches(2.9), Inches(4.6), Inches(4.63)],
      row_h=Inches(0.54), left_cols=(0, 1, 2))
d.rect(s, d.M, Inches(5.20), Inches(6.4), Inches(1.5), hx("FDECEE"), radius=0.04)
d.rect(s, d.M, Inches(5.20), Inches(0.06), Inches(1.5), b.CORAL)
d.text(s, "What this costs, per message", d.M + Inches(0.22), Inches(5.38), Inches(5.9),
       Inches(0.32), size=14, color=b.NAVY, bold=True)
d.text(s, "Every outbound email, proposal and call prep needs a human to reconcile four "
          "documents and recall which claims were retired. Pay that tax wrongly once and the "
          "buyer's evaluator finds it — costing the credibility of the true claims too.",
       d.M + Inches(0.22), Inches(5.74), Inches(5.9), Inches(0.85), size=12,
       color=b.INK, shrink=True)
d.rect(s, Inches(7.25), Inches(5.20), Inches(5.48), Inches(1.5), b.NAVY, radius=0.04)
d.text(s, "Why this is not a document problem", Inches(7.48), Inches(5.38), Inches(5.0),
       Inches(0.32), size=14, color=b.TEAL, bold=True)
d.text(s, "The thinking is good — the July deck has a real ICP scorecard, two-track "
          "positioning and a banned-claims list. It has just never been written somewhere a "
          "person or a tool can execute against. “Read BP-1A §7” is a research task, not a "
          "context file.",
       Inches(7.48), Inches(5.74), Inches(5.0), Inches(0.85), size=12,
       color=b.LIGHT_TEAL, shrink=True)
page(s)
notes(s, """
Open on the problem, not the answer. Without this slide the six files look like bureaucracy;
with it they look like the obvious fix.
The table is deliberately not a criticism of the team. Every row represents a document that
was right when written and was then superseded without the earlier one being withdrawn. That
is normal, and it is exactly what a context layer prevents.
The right-hand panel matters for the room: the thinking is genuinely good. What is missing is
somewhere for it to live that a person or a tool can act on.
""")

# ── THE SYSTEM ──────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "The system is a ladder: two prompts, one chain, six files",
         "Levels 1-3 built · gate run four times · the fourth tested the segment we can actually sell")
rows = [
    ["L1 prompts/", "Qualify-lead + first-contact, four-part structure", "BUILT — unvalidated"],
    ["L2 chains/", "5-step outreach round, checkpoint after qualify", "BUILT — never run"],
    ["01-offer.md", "Five-line ladder AD0–AD4, each with its evidence state", "FIXED — recovered from BP-1A"],
    ["02-icp.md", "Who buys, who does not, and the prioritisation test", "Ready"],
    ["03-voice.md", "5 verbatim samples from BP-1B + banned-words list", "FIXED — was never blocked"],
        ["04-proof.md", "Every claim we may make, with evidence state", "Founder approval on claimable set"],
    ["05-objections.md", "10 road objections + a 7-item AGV gap register", "Partial — AGV set needs real calls"],
    ["06-process.md", "Lead states, stages, owners, approval gates", "Ready"],
]
table(s, ["File", "What it carries", "Status"], rows, d.M, Inches(1.82), d.CW,
      col_w=[Inches(2.6), Inches(6.4), Inches(3.13)], highlight={2},
      row_h=Inches(0.38), left_cols=(0, 1, 2))
d.rect(s, d.M, Inches(5.44), Inches(6.4), Inches(1.5), b.NAVY, radius=0.04)
d.text(s, "Why files and not a document", d.M + Inches(0.22), Inches(5.62), Inches(5.9),
       Inches(0.32), size=14, color=b.TEAL, bold=True)
d.text(s, "Today every outbound message requires a human to reconcile four documents that "
          "disagree, and to remember which claims were retired. Files move that from a "
          "per-message task to a one-time write — and any tooling reads the same source the "
          "team does.",
       d.M + Inches(0.22), Inches(6.00), Inches(5.9), Inches(0.85), size=12,
       color=b.LIGHT_TEAL, shrink=True)
d.rect(s, Inches(7.25), Inches(5.44), Inches(5.48), Inches(1.5), hx("FFF4D6"), radius=0.04)
d.rect(s, Inches(7.25), Inches(5.44), Inches(0.06), Inches(1.5), b.GOLD)
d.text(s, "Round 4 — testing the sellable segment", Inches(7.48), Inches(5.62), Inches(5.0), Inches(0.32),
       size=14, color=b.NAVY, bold=True)
d.text(s, "Rounds 1-3 all tested outbound to a port — the segment we had just demoted to design-partner. Round 4 tested government/PSU + AD0, the one thing sellable today. It recovered the AD0 Tower spec, a second GeM record, and a price defence that exists. It also caught a scoring model that summed to 12 out of 10.",
       Inches(7.48), Inches(6.00), Inches(5.0), Inches(0.85), size=12,
       color=b.INK, shrink=True)
page(s)
notes(s, """
The framework this comes from is explicit that a deliverable which removes no work from
someone's calendar is a toy. So the deliverable is the files; this slide is the map.
Note the blocked row. Voice cannot be invented — it needs three real emails and three real
posts a founder actually wrote. Inventing one would produce the model's default voice, which
is the exact failure Level 3 exists to fix.
""")

# ── DIV WHO ─────────────────────────────────────────────────────────────────
s = d.slide(fill=b.NAVY)
bgp = ASSET / "product" / "truck.png"
if bgp.exists():
    photo_band(s, bgp, Inches(8.0), 0, Inches(5.333), d.H)
d.rect(s, 0, 0, Inches(8.0), d.H, b.NAVY)
d.rect(s, Inches(7.94), 0, Inches(0.06), d.H, b.TEAL)
d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
d.text(s, "PART 01", Inches(0.85), Inches(2.95), Inches(6), Inches(0.34),
       size=14, color=b.TEAL, bold=True)
d.text(s, "Who we sell to", Inches(0.85), Inches(3.38), Inches(6.6), Inches(0.9),
       size=40, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
d.rect(s, Inches(0.85), Inches(4.32), Inches(1.6), Inches(0.05), b.TEAL)
d.text(s, "Two segments now, and the disqualifiers that save weeks",
       Inches(0.85), Inches(4.60), Inches(6.6), Inches(0.5), size=15,
       color=b.LIGHT_TEAL, shrink=True)
page(s, dark=True)
notes(s, "Part one — targeting.")


# ── WHO WE SELL TO ──────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "One segment sells today. The rest need something we do not hold",
         "From 02-icp.md · the prioritisation test is what the buyer requires us to already hold")
rows = [
    ["Government / PSU / defence", "Foreign silicon excluded by rule; GeM access proven",
     "Eligibility, not certification", "NOW"],
    ["Plant, yard and port operators", "Private land — no homologation gate at all",
     "A working demo and a site reference — we have neither", "DESIGN PARTNER"],
    ["Large organised fleets (50+)", "Real safety P&L, single decision-maker",
     "Certified AEBS — not held until 2027", "DESIGN PARTNER"],
    ["OEM / Tier-1 line-fit", "One relationship covers thousands of vehicles",
     "Homologation + a relationship we lack", "START NOW, CLOSES LATE"],
]
table(s, ["Segment", "Why it is addressable", "What it requires us to hold", "Priority"],
      rows, d.M, Inches(1.82), d.CW,
      col_w=[Inches(3.0), Inches(3.8), Inches(3.5), Inches(1.83)], highlight={0},
      row_h=Inches(0.56), left_cols=(0, 1, 2))
d.rect(s, d.M, Inches(4.76), d.CW, Inches(1.9), hx("FDECEE"), radius=0.04)
d.rect(s, d.M, Inches(4.76), Inches(0.06), Inches(1.9), b.CORAL)
d.text(s, "Disqualifiers — these matter more than the qualifiers",
       d.M + Inches(0.25), Inches(4.94), Inches(8), Inches(0.32), size=14,
       color=b.NAVY, bold=True)
d.text(s, [
    {"text": "Sub-five-truck operators (~3.5M). No purchasing process, price anchored at "
             "₹4,500–11,000. Not addressable directly at any sane cost.", "size": 11.5,
     "color": b.INK, "bullet": True},
    {"text": "Anyone buying purely to pass inspection — they will buy a ₹15,000–40,000 box "
             "and they are right to.", "size": 11.5, "color": b.INK, "bullet": True,
     "space_before": 6},
    {"text": "Anyone needing certification today, or a quantified insurance outcome. Walk "
             "rather than invent one.", "size": 11.5, "color": b.INK, "bullet": True,
     "space_before": 6},
], d.M + Inches(0.25), Inches(5.33), d.CW - Inches(0.5), Inches(1.2), shrink=True)
page(s)
notes(s, """
The prioritisation rule is deliberately narrow and it is the whole slide: what does this
buyer require us to already hold?
Government and yard operators require eligibility and a demo. Both exist today. Everything
else waits on certification or a relationship we do not have.
The red panel is the half most ICP work omits. Knowing who wastes your time saves more weeks
than knowing who to target, and the third bullet is a discipline point — walking away from a
buyer who needs a number we cannot evidence is cheaper than inventing it.
""")

# ── PRICE LADDER ────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "Two contests, and only one of them is winnable on price",
         "The ladder the buyer actually sees — from 02-icp.md and 05-objections.md")
add_chart(s, XL_CHART_TYPE.BAR_CLUSTERED,
          ["AIS-140 GPS device", "“AIS-compliant” box", "DeepGrid AD2 kit",
           "Mobileye EyeQ", "Continental ARS", "Nvidia Orin"],
          [("Price (₹ lakh)", [0.08, 0.28, 2.30, 6.50, 10.0, 15.0])],
          d.M, Inches(1.82), Inches(7.3), Inches(4.2),
          colors=[b.TEAL], legend=False, num_fmt='0.0"L"', gap=45)
d.rect(s, Inches(8.15), Inches(1.82), Inches(4.58), Inches(4.2), b.NAVY, radius=0.04)
d.text(s, "How to use this in a room", Inches(8.38), Inches(2.04), Inches(4.1),
       Inches(0.34), size=15, color=b.TEAL, bold=True)
d.text(s, [
    {"text": "Against imports we win 2.8–6.5×. That contest is settled and does not need "
             "re-arguing.", "size": 11.5, "color": b.WHITE},
    {"text": "Against the compliance box we are ~20× the anchor. Price is not an argument "
             "we can win here.", "size": 11.5, "color": b.GOLD, "space_before": 10},
    {"text": "So never open a fleet conversation on price. Open on certified AEBS — what "
             "the ₹30,000 box physically cannot do.", "size": 11, "color": b.LIGHT_TEAL,
     "space_before": 10},
    {"text": "Keep the import comparison for OEM and Tier-1 rooms, where the alternative "
             "bid really is ₹6.5L and up.", "size": 11, "color": b.LIGHT_TEAL,
     "space_before": 10},
], Inches(8.38), Inches(2.52), Inches(4.1), Inches(3.3), shrink=True)
badge(s, "VERIFIED", d.M, Inches(6.22), b.ACCENT, w=Inches(1.3))
d.text(s, "AIS-140 street pricing confirmed independently Aug 2026; import prices per "
          "company deck, not benchmarked", Inches(2.05), Inches(6.26), Inches(10),
       Inches(0.3), size=10, color=b.MUTED, italic=True, shrink=True)
page(s)
notes(s, """
One axis, deliberately, so the asymmetry is impossible to miss: DeepGrid sits far closer to
the compliance box than to the imports, but the buyer's mental anchor is the box.
The instruction is the second bullet. Opening on price invites the comparison we lose.
Opening on certified AEBS moves the conversation to ground where a ₹30,000 device has no
answer at all.
""")


# ── THE PRODUCT LADDER ──────────────────────────────────────────────────────
s = d.slide()
d.header(s, "Six products, not one price range",
         "From 01-offer.md · AD0-AD2 is driver assistance · AD3 is it rented · AD4 is driverless")
rows = [
    ["AD0 basic", "Smart mirror / 360\u00b0 surround. Cameras replace mirrors; driver drives.",
     "\u20b90.50L \u26a0", "Proven on FPGA", "SELL"],
    ["AD0 Tower", "Pod: 3\u00d7 8MP cameras + 4D 77GHz radar + thermal + DMS + compute.",
     "\u20b92.3L \u26a0", "Perception measured", "SELL"],
    ["AD2", "Full truck kit \u2014 compute + software + sensors. Warns and can brake.",
     "\u20b92.0\u20132.5L", "Claim outruns silicon", "SELL"],
    ["AD3", "The AD2 kit rented, not sold. Same hardware, different contract.",
     "\u20b90.66 Cr/yr", "Plan only, FY28+", "PARTNER"],
    ["AD4 Heavy", "Driverless heavy transport AGV >5T. No driver in the cab.",
     "\u20b93.50 Cr", "No site, no reference", "PARTNER"],
    ["AD4 Seaport", "Driverless container-yard AGV \u2014 stack to quay. No driver.",
     "\u20b90.45 Cr", "No site, no reference", "PARTNER"],
]
table(s, ["Line", "What it actually is", "ASP", "Evidence today", "Motion"], rows,
      d.M, Inches(2.0), d.CW,
      col_w=[Inches(1.5), Inches(5.5), Inches(1.5), Inches(2.6), Inches(1.03)],
      highlight={0, 1}, row_h=Inches(0.50), left_cols=(0, 1))
d.rect(s, d.M, Inches(5.66), d.CW, Inches(1.06), hx("FFF4D6"), radius=0.04)
d.rect(s, d.M, Inches(5.66), Inches(0.06), Inches(1.06), b.GOLD)
d.text(s, "\u26a0 Two products are both called AD0: \u20b90.50L in the ADAS table, \u20b92.3L for the Smart "
          "Mirror Tower \u2014 4.6\u00d7 apart, and only the Tower carries radar, thermal and DMS. Confirm "
          "which is being quoted before any AD0 outbound. The wider \u20b90.50L-to-\u20b93.5 Cr spread is a "
          "ladder, not a discount band: AD0 and AD2 are what you sell, AD3 and AD4 are design partners only.",
       d.M + Inches(0.25), Inches(5.82), d.CW - Inches(0.5), Inches(0.78),
       size=12, color=b.INK, bold=True, shrink=True)
page(s)
notes(s, """
This slide exists because it was missing, and its absence broke the cold-chat gate test.
The files described AD4 as "off-highway autonomy, \u20b90.45-3.5 Cr" and never said what it does.
A cold assistant asked to write to a container port could not describe the product, because
the description was not there. It is recovered here from the BP-1A ADAS product-line table.
The important line is the one between AD2 and AD4. AD0 and AD2 are driver assistance: a human
drives and the system helps. AD4 removes the driver. Different buyer, different sales cycle,
and a completely different proof burden \u2014 a driverless multi-tonne vehicle operating near
people needs a safety case that does not exist yet.
Note the evidence column. Everything below AD2 is a plan. 96 percent of the AD4 revenue line
is an unbuilt product, so treat AD4 conversations as discovery, never as pipeline.
""")

# ── WHAT WE SELL ────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "Price the function, give away the warnings",
         "From 01-offer.md · the pricing change that removes the buyer's best objection")
photo(s, "module-m2", d.M, Inches(1.86), Inches(2.9), Inches(1.75))
photo(s, "compute-box", Inches(3.68), Inches(1.86), Inches(2.9), Inches(1.75))
photo(s, "detection", Inches(6.86), Inches(1.86), Inches(2.9), Inches(1.75))
photo(s, "truck", Inches(10.04), Inches(1.86), Inches(2.69), Inches(1.75))
rows = [
    ["Certified AEBS (AIS-162)", "Radar, brake actuation, proving-ground homologation",
     "A ₹30k camera box cannot do this at any price", "PRICE IT"],
    ["Driver monitoring (DMS)", "In-cabin perception, AIS-184 addressable",
     "Hard to approximate; fleet-relevant", "PRICE IT"],
    ["LDW / BSD / warnings", "Warning-only functions",
     "Exactly what the cheap box does adequately", "RIDE FREE"],
]
table(s, ["Function", "What it is", "Why it prices — or does not", "Action"], rows,
      d.M, Inches(3.85), d.CW,
      col_w=[Inches(2.9), Inches(4.0), Inches(4.0), Inches(1.23)], highlight={0, 1},
      row_h=Inches(0.54), left_cols=(0, 1, 2))
d.rect(s, d.M, Inches(6.14), d.CW, Inches(0.62), hx("FFF4D6"), radius=0.04)
d.rect(s, d.M, Inches(6.14), Inches(0.06), Inches(0.62), b.GOLD)
d.text(s, "Sell the bundle and the buyer compares ₹2.3L to ₹30,000 on the box's terms. Sell "
          "the function and the box has no answer. Never open a fleet conversation on price.",
       d.M + Inches(0.25), Inches(6.28), d.CW - Inches(0.5), Inches(0.4),
       size=12, color=b.INK, bold=True, shrink=True)
page(s)
notes(s, """
This is the single most useful change in the whole system and it costs nothing to adopt.
Neither of DeepGrid's plans prices by function — everything is quoted as a kit, which forces
a bundle-versus-box comparison that the cheap box wins.
Split it. AEBS needs radar, brake actuation and homologation; a camera box cannot deliver it
regardless of price. Driver monitoring is similarly hard to fake. Those two carry the price.
The warning functions are what the cheap box does fine. Giving them away removes the buyer's
strongest comparison and costs no margin that was ever winnable.
""")

# ── TRACKS ──────────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "Two tracks, two proof burdens, never the same room",
         "From 02-icp.md and 05-objections.md · positioning discipline")
photo(s, "robots", d.M, Inches(1.86), Inches(6.0), Inches(1.7))
photo(s, "radar-viz", Inches(6.98), Inches(1.86), Inches(5.75), Inches(1.7))
d.rect(s, d.M, Inches(3.68), Inches(6.0), Inches(2.55), b.NAVY, radius=0.04)
d.text(s, "TRACK A · DEFENCE, PSU, YARDS — SELL NOW", d.M + Inches(0.22), Inches(3.86),
       Inches(5.5), Inches(0.3), size=11.5, color=b.TEAL, bold=True)
d.text(s, [
    {"text": "Say: sovereign data, on-die HSM, Indian-designed, shipping on FPGA, GeM "
             "record, ₹23.01L delivered.", "size": 11.5, "color": b.WHITE},
    {"text": "Never say: the retrofit price line, fleet ROI, or anything about insurance.",
     "size": 11.5, "color": b.CORAL, "space_before": 8},
    {"text": "Gate: a named programme and a named SI partner by Month 3, or this is a "
             "services business rather than a silicon one.", "size": 11,
     "color": b.LIGHT_TEAL, "space_before": 8},
], d.M + Inches(0.22), Inches(4.22), Inches(5.5), Inches(1.85), shrink=True)
d.rect(s, Inches(6.98), Inches(3.68), Inches(5.75), Inches(2.55), b.NAVY_2, radius=0.04)
d.text(s, "TRACK B · FLEETS — RECRUIT DESIGN PARTNERS", Inches(7.20), Inches(3.86),
       Inches(5.3), Inches(0.3), size=11.5, color=b.GOLD, bold=True)
d.text(s, [
    {"text": "Say: the 1 Oct 2027 AEBS requirement, one kit covering AEBS and DMS, "
             "influence over what gets certified for your fleet.", "size": 11.5,
     "color": b.WHITE},
    {"text": "Never say: sovereign data, defence customers, or a 2026 deadline.",
     "size": 11.5, "color": b.CORAL, "space_before": 8},
    {"text": "Gate: two design partners committing fleet access and a measurement baseline "
             "by Month 4.", "size": 11, "color": b.LIGHT_TEAL, "space_before": 8},
], Inches(7.20), Inches(4.22), Inches(5.3), Inches(1.85), shrink=True)
d.rect(s, d.M, Inches(6.35), d.CW, Inches(0.5), hx("FFF4D6"), radius=0.04)
d.rect(s, d.M, Inches(6.35), Inches(0.06), Inches(0.5), b.GOLD)
d.text(s, "A blurred statement reads as a company that has not chosen. The two stories are "
          "individually strong and jointly weak.",
       d.M + Inches(0.25), Inches(6.46), d.CW - Inches(0.5), Inches(0.3),
       size=11.5, color=b.INK, bold=True, shrink=True)
page(s)
notes(s, """
The discipline rule is the easiest to break under pressure and usually with good intentions —
a fleet meeting is going badly and the defence traction is right there as proof somebody buys
this.
Resist it. Defence proof does not transfer; it signals the fleet product is a side project.
In reverse it is worse: a procurement officer who hears a per-unit retrofit pitch
re-categorises the whole conversation as commodity buying.
Both gates are revised from the original plan because verification showed the originals were
either unreachable or measuring the wrong thing.
""")

# ── WHAT WE MAY CLAIM ───────────────────────────────────────────────────────
s = d.slide()
d.header(s, "What we may claim — and the nine claims now retired",
         "From 04-proof.md · this file is the ceiling on every claim the system may make")
d.text(s, "May be claimed", d.M, Inches(1.80), Inches(6.2), Inches(0.3),
       size=13, color=b.NAVY, bold=True)
table(s, ["Claim", "State"], [
    ["GeM procurement record held", "CONTRACTED"],
    ["₹23.01L delivered — Robot Training", "CONTRACTED*"],
    ["YOLOv11n at 40 fps on FPGA", "MEASURED"],
    ["Attention head at 24.25 ms", "MEASURED"],
    ["DGS001 running on FPGA", "SILICON"],
    ["15 provisional patents filed", "FILED"],
], d.M, Inches(2.14), Inches(6.2), col_w=[Inches(4.5), Inches(1.7)],
    highlight={0, 2, 3}, row_h=Inches(0.42), left_cols=(0,))
d.text(s, "* sits in Deepgrid Datacentre, not Semi. IP is MCEME-owned. Not ADAS revenue — "
          "this caveat must always travel with it.",
       d.M, Inches(5.05), Inches(6.2), Inches(0.5), size=10, color=b.MUTED,
       italic=True, shrink=True)
d.text(s, "Retired — never say", Inches(7.05), Inches(1.80), Inches(5.68), Inches(0.3),
       size=13, color=b.NAVY, bold=True)
table(s, ["Retired claim", "Why"], [
    ["“39.3 TOPS measured”", "Artix-7 ceilings ~1.8 TOPS"],
    ["“84% / 88% gross margin”", "Retracted; blend needs 120%"],
    ["“12.9× cheaper”", "Their ASP vs our die cost"],
    ["“₹1 Cr defence revenue”", "₹23.01L; rest is L1"],
    ["“Transformer VLA on-chip”", "Bandwidth-bound"],
    ["“Mandate live since Apr 2026”", "That rule was a draft"],
    ["“Mandate-ready” · “ASIL-D”", "Neither is held"],
    ["Any insurance / accident %", "No baseline, no actuary"],
], Inches(7.05), Inches(2.14), Inches(5.68), col_w=[Inches(3.1), Inches(2.58)],
    row_h=Inches(0.42), left_cols=(0, 1))
page(s)
notes(s, """
This is the page that protects the business, and it belongs in front of anyone client-facing
before their first meeting.
The left column is short on purpose. Two contracted items, two measured numbers, one live
demo, one filing. That is a real and defensible position for a company at this stage — and
far stronger than a longer list that collapses on inspection.
The right column exists because every one of those claims is disprovable by a competent
evaluator in a single meeting. Being caught once costs the credibility of everything else in
the room, including the true things on the left.
""")

# ── DIV HOW ─────────────────────────────────────────────────────────────────
s = d.slide(fill=b.NAVY)
bgp = ASSET / "product" / "detection.png"
if bgp.exists():
    photo_band(s, bgp, Inches(8.0), 0, Inches(5.333), d.H)
d.rect(s, 0, 0, Inches(8.0), d.H, b.NAVY)
d.rect(s, Inches(7.94), 0, Inches(0.06), d.H, b.TEAL)
d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
d.text(s, "PART 02", Inches(0.85), Inches(2.95), Inches(6), Inches(0.34),
       size=14, color=b.TEAL, bold=True)
d.text(s, "How a lead travels", Inches(0.85), Inches(3.38), Inches(6.6), Inches(0.9),
       size=40, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
d.rect(s, Inches(0.85), Inches(4.32), Inches(1.6), Inches(0.05), b.TEAL)
d.text(s, "One intake path, one state machine, written-down approvals",
       Inches(0.85), Inches(4.60), Inches(6.6), Inches(0.5), size=15,
       color=b.LIGHT_TEAL, shrink=True)
page(s, dark=True)
notes(s, "Part two — the pipeline.")


# ── HOW A LEAD TRAVELS ──────────────────────────────────────────────────────
s = d.slide()
d.header(s, "One intake path, one state machine, one record per lead",
         "From 06-process.md · this is the foundation everything automated sits on")
states = ["new", "researched", "contacted", "replied", "booked", "proposed", "won / lost / dead"]
w = (d.CW - Inches(0.36)) / len(states)
for i, st in enumerate(states):
    x = d.M + i * (w + Inches(0.06))
    fill = b.ACCENT if i < 6 else b.NAVY
    d.rect(s, x, Inches(1.90), w, Inches(0.62), fill, radius=0.06)
    d.text(s, st, x, Inches(2.06), w, Inches(0.3), size=10.5, color=b.WHITE,
           bold=True, align=PP_ALIGN.CENTER, shrink=True)
d.text(s, "Every lead enters the same way regardless of source — referral, inbound, event, "
          "cold, GeM tender alert. Same record shape, same first step. Multiple intake paths "
          "is the most common reason these systems collapse.",
       d.M, Inches(2.68), d.CW, Inches(0.5), size=11.5, color=b.MUTED, shrink=True)
rows = [
    ["1–2", "Source and qualify", "Commercial lead", "Passes ICP; disqualifiers cleared"],
    ["3", "Research", "Commercial lead", "Buying trigger found, or “none found” recorded"],
    ["4–5", "Contact and discovery", "Commercial lead", "Requirement, timeline, decision unit"],
    ["6", "Technical review", "CTO", "Their evaluator has seen demo and proof kit"],
    ["7–8", "Proposal and close", "Commercial + founder", "Priced by function; founder approved"],
    ["9", "Deploy", "Engineering", "Installed; baseline measurement started"],
]
table(s, ["Stage", "What happens", "Owner", "Exit criterion"], rows, d.M, Inches(3.32),
      d.CW, col_w=[Inches(0.9), Inches(3.4), Inches(2.6), Inches(5.23)],
      row_h=Inches(0.42), left_cols=(0, 1, 2, 3))
d.rect(s, d.M, Inches(6.28), d.CW, Inches(0.55), hx("FDECEE"), radius=0.04)
d.rect(s, d.M, Inches(6.28), Inches(0.06), Inches(0.55), b.CORAL)
d.text(s, "Human approval, not negotiable: anything sent to a client · any price · any "
          "contract · any claim not in 04-proof.md · any lead scoring above 9.",
       d.M + Inches(0.25), Inches(6.41), d.CW - Inches(0.5), Inches(0.32),
       size=11.5, color=b.INK, bold=True, shrink=True)
page(s)
notes(s, """
The state machine looks trivial and is the thing most acquisition systems get wrong.
One record, one state, one store. If two people — or later, two agents — disagree about
where a lead is, the store is broken and no amount of process fixes it.
One intake path matters just as much. The temptation is to handle referrals differently from
cold leads because referrals feel special. That is exactly how the record shape diverges and
the pipeline stops being countable.
The approval line at the bottom is written down in advance deliberately, so it is never a
judgement call in the moment.
""")

# ── LADDER CHART ────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "Ten levels, built in order — we are at three",
         "Each level must remove a task from somebody's calendar, or it is deleted")
add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED,
          ["L1\nPrompt", "L2\nChain", "L3\nBrain", "L4\nSkill", "L5\nDeliverable",
           "L6\nConnect", "L7\nAgent", "L8\nTeam", "L9\nAlways on", "L10\nFunction"],
          [("Build effort (hours)", [0.3, 1, 2, 2, 3, 1, 4, 12, 8, 80])],
          d.M, Inches(1.82), Inches(7.5), Inches(3.9),
          colors=[b.TEAL], legend=False, num_fmt='0"h"', gap=40)
d.text(s, "Effort is not the constraint — sequence is. Level 7 built on a broken Level 3 is "
          "just a faster mess.", d.M, Inches(5.82), Inches(7.5), Inches(0.4),
       size=10.5, color=b.MUTED, shrink=True)
kpi_at(s, [("L3", "Building now", "the Brain", b.GOLD),
           ("L6", "Blocked", "no CRM exists", b.CORAL)],
       Inches(8.25), Inches(1.86), Inches(4.48))
d.rect(s, Inches(8.25), Inches(3.48), Inches(4.48), Inches(2.75), b.NAVY, radius=0.04)
d.text(s, "Why Level 3 first", Inches(8.48), Inches(3.70), Inches(4.0), Inches(0.34),
       size=15, color=b.TEAL, bold=True)
d.text(s, [
    {"text": "Levels 1 and 2 are trivial once the Brain exists — and writing them first "
             "would encode the contradictions between the four source documents.",
     "size": 11.5, "color": b.WHITE},
    {"text": "Level 6 is genuinely blocked: there is no CRM or shared store, so a connected "
             "system has nowhere to read and write.", "size": 11, "color": b.LIGHT_TEAL,
     "space_before": 10},
    {"text": "That is a purchase decision, not something to engineer around.",
     "size": 11, "color": b.GOLD, "space_before": 10},
], Inches(8.48), Inches(4.14), Inches(4.0), Inches(1.95), shrink=True)
page(s)
notes(s, """
The chart is build effort, and it deliberately shows that effort is not what governs the
order — sequence is.
Level 10 dwarfs everything else at roughly two to four weeks on top of all the rest, which is
why jumping to it is the classic failure.
The blocked marker on Level 6 is the honest constraint. No CRM means no connection layer, and
no connection layer means agents have nothing real to act on. Fix that before anything above
it is worth attempting.
""")


# ── DIV NEXT ────────────────────────────────────────────────────────────────
s = d.slide(fill=b.NAVY)
bgp = ASSET / "product" / "robots.png"
if bgp.exists():
    photo_band(s, bgp, Inches(8.0), 0, Inches(5.333), d.H)
d.rect(s, 0, 0, Inches(8.0), d.H, b.NAVY)
d.rect(s, Inches(7.94), 0, Inches(0.06), d.H, b.TEAL)
d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
d.text(s, "PART 03", Inches(0.85), Inches(2.95), Inches(6), Inches(0.34),
       size=14, color=b.TEAL, bold=True)
d.text(s, "What happens next", Inches(0.85), Inches(3.38), Inches(6.6), Inches(0.9),
       size=40, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
d.rect(s, Inches(0.85), Inches(4.32), Inches(1.6), Inches(0.05), b.TEAL)
d.text(s, "Three founder decisions gate everything else",
       Inches(0.85), Inches(4.60), Inches(6.6), Inches(0.5), size=15,
       color=b.LIGHT_TEAL, shrink=True)
page(s, dark=True)
notes(s, "Part three — execution.")


# ── NEXT STEPS ──────────────────────────────────────────────────────────────
s = d.slide()
d.header(s, "What happens next", "Item 1 is one sentence from the founder and it blocks the whole port motion")
rows = [
    ["1", "Answer: is AD4 Seaport a vehicle or a retrofit kit?", "Founder", "Today", "Decides replace-your-fleet vs convert-it; blocks all port outbound"],
    ["2", "Supply one real COLD email to a stranger", "Founder", "This week", "5 samples recovered from BP-1B; all are replies, none are cold"],
    ["3", "Founder approves the claimable set in 04-proof.md", "Founder", "This week", "Nothing goes out until this is signed"],
    ["4", "Founder approves pricing in 01-offer.md", "Founder", "This week", "Price-by-function is a commitment, not a description"],
    ["5", "Re-run the cold-chat gate test after items 1–4", "Commercial lead", "Week 2", "Run twice: guardrails pass, outbound still fails"],
    ["6", "Withdraw the June IM from circulation", "Founder", "Immediately", "It carries six retired claims"],
    ["7", "Resolve the AD3 rent-vs-sale price conflict", "Founder", "Week 2", "₹0.66 Cr/yr is ~26× the AD2 sale price"],
    ["8", "Name a defence programme and an SI partner", "Commercial lead", "Weeks 1–12", "Track A has a proposition and no pipeline"],
    ["9", "Choose a CRM or shared store", "Founder", "Month 2", "Unblocks Level 6 and everything above"],
    ["10", "Build the AD4 safety case, or stop pitching AD4", "Founder", "Before any port pilot", "No funded functional-safety work exists for a driverless vehicle"],
]
table(s, ["#", "Action", "Owner", "When", "Why"], rows, d.M, Inches(1.82), d.CW,
      col_w=[Inches(0.55), Inches(4.6), Inches(2.1), Inches(1.6), Inches(3.28)],
      highlight={0, 1, 2, 3}, row_h=Inches(0.36), left_cols=(0, 1, 4))
d.rect(s, d.M, Inches(6.14), d.CW, Inches(0.62), hx("FFF4D6"), radius=0.04)
d.rect(s, d.M, Inches(6.14), Inches(0.06), Inches(0.62), b.GOLD)
d.text(s, "Items 1–4 are founder decisions and they gate everything else. Item 1 is a single sentence. Item 10 is new: the driverless segment has no objection set and no safety case — until it does, AD4 is not pitchable as a product.",
       d.M + Inches(0.25), Inches(6.28), d.CW - Inches(0.5), Inches(0.4),
       size=12, color=b.INK, bold=True, shrink=True)
page(s)
notes(s, """
Three of the eight need a founder specifically, and they are the three that gate the rest.
Voice samples cannot be delegated or invented. The claimable set is a commitment about what
the company will assert publicly. Pricing by function is a real pricing decision, not a
presentation choice.
Item five costs nothing and should happen today — the June IM is the most polished document
in the set and carries six claims that have been retired.
Item eight is the honest weakness in what has been built: the objection responses are drafted
from verified evidence rather than taken from real conversations, which is the opposite of
what the framework asks for. They are a starting point to be replaced, not a finished file.
""")

out = OUT / "DeepGrid-client-acquisition-system-draft.pptx"
d.save(out, validate=True)
print(f"slides: {d.n}")
