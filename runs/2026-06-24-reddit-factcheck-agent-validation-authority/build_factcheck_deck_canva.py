#!/usr/bin/env python3
"""Build the client-ready fact-check dossier from the branded Canva template."""

from __future__ import annotations

import json
import os
import re
import sys
from collections import Counter, defaultdict
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPTXKIT = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT))
from pptxkit import validate_pptx  # noqa: E402

RUN_DIR = Path(__file__).resolve().parent
CLAIM_PACK = RUN_DIR / "claim-pack.json"
EVIDENCE_PACK = RUN_DIR / "reddit-evidence-pack.json"
TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE") or "/home/shekerk/.claude/templates/branded-template.pptx")
OUT = RUN_DIR / "reddit-factcheck-dossier-branded-canva-evidence-draft.pptx"

NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_2 = RGBColor(0x12, 0x24, 0x3A)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
DARK_TEAL = RGBColor(0x00, 0x8F, 0x75)
GOLD = RGBColor(0xFF, 0xB8, 0x00)
AMBER = RGBColor(0xF2, 0xA8, 0x3B)
CORAL = RGBColor(0xE0, 0x5A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)


def load_claims() -> list[dict]:
    return json.loads(CLAIM_PACK.read_text(encoding="utf-8"))["claims"]


def load_evidence() -> dict:
    return json.loads(EVIDENCE_PACK.read_text(encoding="utf-8")) if EVIDENCE_PACK.exists() else {"claims": [], "source_files": []}


def evidence_by_claim(evidence: dict) -> dict[str, dict]:
    return {item["claim_id"]: item for item in evidence.get("claims", [])}


def verdict_counts(evidence: dict) -> Counter:
    return Counter(item.get("suggested_reddit_verdict", "pending") for item in evidence.get("claims", []))


def source_stats(evidence: dict) -> tuple[int, int]:
    source_files = [Path(path) for path in evidence.get("source_files", [])]
    comments = 0
    for path in source_files:
        if not path.exists():
            continue
        data = json.loads(path.read_text(encoding="utf-8"))
        comments += len(data.get("comments", []))
    return len(source_files), comments


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def trim_to_first_n(prs: Presentation, n: int) -> None:
    for index in range(len(prs.slides._sldIdLst) - 1, n - 1, -1):
        rid = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[index]


def keep_template_base(shape) -> bool:
    is_picture_rail = shape.shape_type == 13 and shape.left >= Inches(8.4)
    is_left_background = shape.left <= 1 and shape.top <= 1 and shape.width >= Inches(8.7)
    is_cover_rule = shape.left <= 1 and shape.width >= Inches(8.7) and shape.height <= Inches(0.08)
    return bool(is_picture_rail or is_left_background or is_cover_rule)


def clear_to_template_base(slide) -> None:
    for shape in list(slide.shapes):
        if keep_template_base(shape):
            continue
        remove_shape(shape)


def rect(slide, x, y, w, h, color, radius: bool = False, line=None):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, *, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    lines = value if isinstance(value, list) else [value]
    for index, item in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        if isinstance(item, dict):
            content = item.get("text", "")
            paragraph.space_before = Pt(item.get("space_before", 0))
            item_size = item.get("size", size)
            item_color = item.get("color", color)
            item_bold = item.get("bold", bold)
            if item.get("bullet"):
                content = "- " + content
        else:
            content = str(item)
            item_size = size
            item_color = color
            item_bold = bold
        run = paragraph.add_run()
        run.text = content
        run.font.name = "Calibri"
        run.font.size = Pt(item_size)
        run.font.bold = item_bold
        run.font.color.rgb = item_color
    return shape


def short(value: str, limit: int = 115) -> str:
    value = re.sub(r"\s+", " ", value.replace("*", "")).strip()
    return value if len(value) <= limit else value[: limit - 1].rstrip() + "..."


def claim_id(value: str) -> str:
    match = re.search(r"(\d+)$", value)
    return f"C-{int(match.group(1)):03d}" if match else value.upper()


def slide_num(locator: str) -> int:
    match = re.search(r"\b(?:Slide|Section)\s+(\d+)\b", locator, re.I)
    return int(match.group(1)) if match else 0


def evidence_label(claim: dict) -> str:
    if claim["evidence_need"] == "primary_required":
        return "Primary source"
    if claim["evidence_need"] == "mixed":
        return "Mixed"
    return "Reddit"


def priority_score(claim: dict) -> int:
    body = claim["claim"].lower()
    score = 0
    if claim["primary_source_required"]:
        score += 4
    if claim["evidence_need"] == "mixed":
        score += 2
    if re.search(r"\d|%|\$|sr 26-2|gartner|occ|fdic|federal|nist|iso", body):
        score += 3
    if any(term in body for term in ["klarna", "wendy", "mercedes", "home depot", "big 4"]):
        score += 2
    return score


def side_panel(slide, title: str, lines: list[str], status: str = "DRAFT") -> None:
    rect(slide, Inches(8.84), 0, Inches(4.49), Inches(7.5), NAVY_2)
    rect(slide, Inches(8.84), 0, Inches(0.08), Inches(7.5), TEAL)
    text(slide, title, Inches(9.22), Inches(0.55), Inches(3.3), Inches(0.72), size=21, color=WHITE, bold=True)
    y = 1.55
    for line in lines:
        text(slide, line, Inches(9.22), Inches(y), Inches(3.38), Inches(0.42), size=11.5, color=WHITE)
        y += 0.58
    rect(slide, Inches(9.22), Inches(6.15), Inches(2.6), Inches(0.34), GOLD)
    text(slide, status, Inches(9.31), Inches(6.23), Inches(2.2), Inches(0.16), size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def footer(slide, page: int, total: int, dark: bool = False) -> None:
    color = WHITE if dark else INK
    text(slide, "Reddit fact-check dossier | confidential", Inches(0.28), Inches(7.07), Inches(4.8), Inches(0.22), size=7.5, color=color)
    text(slide, f"{page} / {total}", Inches(7.3), Inches(7.07), Inches(0.7), Inches(0.22), size=7.5, color=color, align=PP_ALIGN.RIGHT)


def cover(slide, claims: list[dict], evidence: dict) -> None:
    counts = verdict_counts(evidence)
    source_count, comment_count = source_stats(evidence)
    text(slide, "REDDIT FACT-CHECK DOSSIER", Inches(0.8), Inches(0.88), Inches(5.8), Inches(0.32), size=13, color=TEAL, bold=True)
    text(slide, "Agent Validation & Assurance Authority", Inches(0.8), Inches(1.42), Inches(6.9), Inches(1.0), size=31, color=WHITE, bold=True)
    rect(slide, Inches(0.8), Inches(2.73), Inches(1.45), Inches(0.05), TEAL)
    text(
        slide,
        "Client-ready draft built from the branded Canva template. This version includes Reddit HTML evidence, claim-level scoring, and source traceability.",
        Inches(0.8),
        Inches(3.12),
        Inches(6.7),
        Inches(0.78),
        size=13.5,
        color=WHITE,
    )
    y = 4.62
    for label in [
        f"{len(claims)} extracted claims",
        f"{counts['supported_by_reddit_evidence']} Reddit-supported",
        f"{counts['no_reddit_evidence_found']} no Reddit evidence",
    ]:
        rect(slide, Inches(0.8), Inches(y), Inches(2.05), Inches(0.36), DARK_TEAL, radius=True)
        text(slide, label, Inches(0.91), Inches(y + 0.08), Inches(1.82), Inches(0.18), size=6.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        y += 0.46
    text(slide, f"Built {date.today().isoformat()} | Status: DRAFT", Inches(0.8), Inches(6.22), Inches(4.3), Inches(0.25), size=9, color=GOLD, bold=True)
    side_panel(slide, "Evidence attached", [f"{source_count} Reddit threads", f"{comment_count} comments parsed", "Primary facts still need official sources."])
    footer(slide, 1, 8, dark=True)


def executive_summary(slide, claims: list[dict], evidence: dict) -> None:
    counts = verdict_counts(evidence)
    text(slide, "Executive Verdict", Inches(0.4), Inches(0.28), Inches(7.6), Inches(0.56), size=20, color=NAVY, bold=True)
    text(slide, "Reddit validation exists; hard facts still need primary-source review.", Inches(0.4), Inches(0.84), Inches(7.0), Inches(0.28), size=9.4, color=INK)
    metrics = [
        (str(counts["supported_by_reddit_evidence"]), "Reddit-supported", GOLD),
        (str(counts["weak_reddit_support"]), "weak Reddit support", NAVY),
        (str(counts["no_reddit_evidence_found"]), "no Reddit evidence", CORAL),
        (str(counts["requires_primary_source_corroboration"]), "primary required", AMBER),
    ]
    x = 0.42
    for value, label, color in metrics:
        rect(slide, Inches(x), Inches(1.35), Inches(1.72), Inches(0.94), NAVY)
        text(slide, value, Inches(x + 0.12), Inches(1.48), Inches(0.62), Inches(0.34), size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, label, Inches(x + 0.78), Inches(1.45), Inches(0.82), Inches(0.36), size=6.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += 1.93
    bullets = [
        "Old Reddit HTML was parsed into evidence JSON.",
        "Claims are scored against post and comment text.",
        "Reddit supports language, pain, and objections; hard facts need official sources.",
        "Status stays draft until human evidence review is complete.",
    ]
    text(slide, [{"text": item, "bullet": True, "space_before": 7, "size": 10.8} for item in bullets], Inches(0.45), Inches(2.52), Inches(7.1), Inches(2.1), size=10.8, color=INK)
    source_count, comment_count = source_stats(evidence)
    side_panel(slide, "Evidence base", [f"{source_count} source threads", f"{comment_count} comments", "Old Reddit HTML extraction"])
    footer(slide, 2, 8)


def evidence_limits(slide) -> None:
    text(slide, "Reddit Evidence Has Clear Limits", Inches(0.4), Inches(0.28), Inches(7.6), Inches(0.42), size=22, color=NAVY, bold=True)
    columns = [
        ("Reddit can support", ["Operator pain", "Practitioner language", "Vendor complaints", "Objections and demand"], NAVY),
        ("Reddit cannot prove alone", ["Market-size numbers", "Regulatory meaning", "Funding or headcount", "Named-company facts"], CORAL),
        ("Output standard", ["Claim IDs", "Source metadata", "Short rationale", "No opinion-only labels"], GOLD),
    ]
    x = 0.42
    for title, items, color in columns:
        rect(slide, Inches(x), Inches(1.12), Inches(2.38), Inches(3.48), SOFT)
        rect(slide, Inches(x), Inches(1.12), Inches(2.38), Inches(0.12), color)
        text(slide, title, Inches(x + 0.16), Inches(1.45), Inches(2.05), Inches(0.38), size=11.5, color=NAVY, bold=True)
        text(slide, [{"text": item, "bullet": True, "space_before": 9, "size": 10.2} for item in items], Inches(x + 0.2), Inches(2.05), Inches(1.95), Inches(1.7), color=INK)
        x += 2.6
    side_panel(slide, "Rule", ["Use Reddit as community evidence.", "Use official sources for hard facts.", "Label gaps explicitly."])
    footer(slide, 3, 8)


def triage(slide, claims: list[dict]) -> None:
    text(slide, "Claim Triage Prioritizes Reddit Review", Inches(0.4), Inches(0.28), Inches(7.6), Inches(0.42), size=22, color=NAVY, bold=True)
    text(slide, "The biggest review load sits in market-language and operator-pain claims.", Inches(0.4), Inches(0.78), Inches(7), Inches(0.22), size=10.2, color=INK)
    by_slide: dict[int, list[dict]] = defaultdict(list)
    for claim in claims:
        by_slide[slide_num(claim["locator"])].append(claim)
    rows = sorted(by_slide.items(), key=lambda item: item[0])[:14]
    max_count = max(len(items) for _, items in rows)
    y = 1.22
    for slide_number, items in rows:
        primary = sum(1 for item in items if item["evidence_need"] == "primary_required")
        mixed = sum(1 for item in items if item["evidence_need"] == "mixed")
        reddit = len(items) - primary - mixed
        text(slide, f"S{slide_number}", Inches(0.46), Inches(y - 0.02), Inches(0.35), Inches(0.15), size=7.5, color=INK, bold=True)
        width = 5.75 * len(items) / max_count
        rect(slide, Inches(0.86), Inches(y), Inches(width), Inches(0.13), NAVY)
        if mixed:
            rect(slide, Inches(0.86 + width), Inches(y), Inches(0.18 * mixed), Inches(0.13), AMBER)
        if primary:
            rect(slide, Inches(0.86 + width + (0.18 * mixed)), Inches(y), Inches(0.18 * primary), Inches(0.13), CORAL)
        text(slide, f"{len(items)} | R{reddit} M{mixed} P{primary}", Inches(6.8), Inches(y - 0.02), Inches(0.82), Inches(0.14), size=5.4, color=INK)
        y += 0.33
    side_panel(slide, "Signal", ["Teal/navy bars show Reddit review load.", "Coral marks primary-source risk.", "Amber means mixed evidence."])
    footer(slide, 4, 8)


def priority_claims(slide, claims: list[dict], evidence: dict) -> None:
    by_id = evidence_by_claim(evidence)
    text(slide, "Claims That Still Need Work", Inches(0.4), Inches(0.28), Inches(7.6), Inches(0.42), size=21, color=NAVY, bold=True)
    priority = sorted(
        claims,
        key=lambda claim: (
            by_id.get(claim["id"], {}).get("suggested_reddit_verdict") in {"no_reddit_evidence_found", "requires_primary_source_corroboration"},
            priority_score(claim),
        ),
        reverse=True,
    )[:5]
    y = 1.05
    for claim in priority:
        verdict = by_id.get(claim["id"], {}).get("suggested_reddit_verdict", "pending")
        color = CORAL if verdict in {"no_reddit_evidence_found", "requires_primary_source_corroboration"} else AMBER if "weak" in verdict else NAVY
        rect(slide, Inches(0.44), Inches(y), Inches(7.25), Inches(0.68), SOFT)
        rect(slide, Inches(0.44), Inches(y), Inches(0.08), Inches(0.68), color)
        text(slide, claim_id(claim["id"]), Inches(0.6), Inches(y + 0.18), Inches(0.5), Inches(0.16), size=6.6, color=color, bold=True)
        text(slide, short(claim["claim"], 60), Inches(1.14), Inches(y + 0.11), Inches(4.85), Inches(0.34), size=6.4, color=INK, bold=True)
        text(slide, verdict.replace("_", " ")[:24], Inches(6.1), Inches(y + 0.17), Inches(1.3), Inches(0.22), size=5.2, color=MUTED)
        y += 0.8
    side_panel(slide, "Gate", ["Fix unsupported claims.", "Attach official sources for numbers and rules.", "Then mark reviewed."])
    footer(slide, 5, 8)


def discovery(slide, evidence: dict) -> None:
    text(slide, "Reddit Evidence Sources", Inches(0.72), Inches(0.82), Inches(6.8), Inches(0.72), size=22, color=WHITE, bold=True)
    text(slide, "Curated old.reddit threads were parsed into JSON and scored against the claim pack.", Inches(0.72), Inches(1.52), Inches(6.7), Inches(0.42), size=10.5, color=WHITE)
    source_files = [Path(path) for path in evidence.get("source_files", [])][:7]
    y = 2.2
    for index, path in enumerate(source_files, 1):
        data = json.loads(path.read_text(encoding="utf-8"))
        post = data.get("post", {})
        rect(slide, Inches(0.72), Inches(y), Inches(6.7), Inches(0.4), NAVY_2)
        text(slide, f"R{index}", Inches(0.86), Inches(y + 0.12), Inches(0.48), Inches(0.12), size=6.0, color=TEAL, bold=True)
        text(slide, short(f"r/{post.get('subreddit','?')} - {post.get('title','')}", 76), Inches(1.38), Inches(y + 0.1), Inches(5.55), Inches(0.15), size=6.4, color=WHITE)
        y += 0.52
    side_panel(slide, "Traceability", ["Source URLs are in the report.", "Every claim has evidence status.", "No opinion-only labels."])
    footer(slide, 6, 8, dark=True)


def artifacts(slide) -> None:
    text(slide, "Delivery Artifacts", Inches(0.4), Inches(0.28), Inches(7.6), Inches(0.42), size=23, color=NAVY, bold=True)
    cards = [
        ("Annotated source deck", "Original review copy with claim rails and Fact Check Notes.", "PPTX", TEAL),
        ("Evidence report", "Markdown audit trail with claims, queries, and limitations.", "MD", GOLD),
        ("Machine artifacts", "JSON source for repeatable scoring and regeneration.", "JSON", AMBER),
    ]
    y = 1.18
    for title, body, file_name, color in cards:
        rect(slide, Inches(0.45), Inches(y), Inches(7.15), Inches(0.82), SOFT)
        rect(slide, Inches(0.45), Inches(y), Inches(0.09), Inches(0.82), color)
        text(slide, title, Inches(0.66), Inches(y + 0.12), Inches(1.92), Inches(0.28), size=8.2, color=NAVY, bold=True)
        text(slide, body, Inches(2.74), Inches(y + 0.14), Inches(2.5), Inches(0.32), size=6.6, color=INK)
        text(slide, file_name, Inches(5.65), Inches(y + 0.18), Inches(0.75), Inches(0.16), size=8.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        y += 1.0
    side_panel(slide, "Use case", ["Source deck: editing", "Dossier: executive readout", "JSON: reproducibility"])
    footer(slide, 7, 8)


def decision(slide) -> None:
    text(slide, "DECISION", Inches(0.4), Inches(0.5), Inches(1.5), Inches(0.2), size=10, color=TEAL, bold=True)
    text(slide, "Move From Claim Extraction To Evidence Review", Inches(0.4), Inches(1.0), Inches(7.1), Inches(0.88), size=26, color=NAVY, bold=True)
    bullets = [
        "Edit unsupported claims in the annotated source deck.",
        "Keep Reddit-backed language where it reflects operator pain.",
        "Attach official sources for SR 26-2, Gartner, named-company facts, and numbers.",
    ]
    text(slide, [{"text": item, "bullet": True, "space_before": 8, "size": 10.8} for item in bullets], Inches(0.55), Inches(2.18), Inches(6.9), Inches(1.75), color=INK)
    rect(slide, Inches(0.55), Inches(5.18), Inches(7.0), Inches(0.42), GOLD)
    text(slide, "Current status: EVIDENCE DRAFT - Reddit evidence attached, primary-source review pending.", Inches(0.72), Inches(5.29), Inches(6.4), Inches(0.22), size=7.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    side_panel(slide, "Reviewed when", ["Primary facts sourced", "Noisy matches removed", "Deck QA passed"])
    footer(slide, 8, 8)


def main() -> int:
    claims = load_claims()
    evidence = load_evidence()
    prs = Presentation(str(TEMPLATE))
    trim_to_first_n(prs, 8)
    for slide in prs.slides:
        clear_to_template_base(slide)

    cover(prs.slides[0], claims, evidence)
    executive_summary(prs.slides[1], claims, evidence)
    evidence_limits(prs.slides[2])
    triage(prs.slides[3], claims)
    priority_claims(prs.slides[4], claims, evidence)
    discovery(prs.slides[5], evidence)
    artifacts(prs.slides[6])
    decision(prs.slides[7])

    prs.save(str(OUT))
    problems = validate_pptx(OUT)
    if problems:
        raise RuntimeError("\n".join(problems))
    print(f"Saved {OUT} using template {TEMPLATE}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
