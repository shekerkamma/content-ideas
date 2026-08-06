#!/usr/bin/env python3
"""Build the DeepGrid client-acquisition strategy deck.

Charts are NATIVE python-pptx charts (editable in PowerPoint, Excel-backed) —
matplotlib is not installed on this host, and native charts are the better
artifact anyway.

Kept in the run folder so QA fixes stay reproducible (CLAUDE.md PPTX gate #8).
"""
import sys
from pathlib import Path

KIT = Path("/home/sheke/content-ideas/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(KIT))

from pptx.util import Inches, Pt, Emu                      # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR            # noqa: E402
from pptx.chart.data import CategoryChartData              # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION  # noqa: E402
from pptxkit import Deck, Brand, hx                        # noqa: E402

OUT = Path(__file__).resolve().parent
d = Deck(footer="DeepGrid Semi — client acquisition strategy · July-corrected basis · 5 Aug 2026")
b = d.b
YRS = ["FY27", "FY28", "FY29", "FY30", "FY31", "FY32"]
TOTAL = 42
_page = {"n": 0}


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def page(s, dark=False):
    _page["n"] += 1
    d.footer(s, _page["n"], TOTAL, dark=dark)


def style_chart(gf, *, colors, legend=True, num_fmt='0.0"%"', font=9,
                gap=60, smooth=False, markers=False):
    ch = gf.chart
    ch.font.size = Pt(font)
    ch.font.name = b.FONT
    ch.has_title = False
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(font)
    else:
        ch.has_legend = False
    for i, ser in enumerate(ch.series):
        col = colors[i % len(colors)]
        if ch.chart_type in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
            ser.format.line.color.rgb = col
            ser.format.line.width = Pt(2.5)
            ser.smooth = smooth
        else:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = col
    try:
        ch.value_axis.has_major_gridlines = True
        ch.value_axis.major_gridlines.format.line.color.rgb = b.GRID
        ch.value_axis.major_gridlines.format.line.width = Pt(0.5)
        ch.value_axis.tick_labels.number_format = num_fmt
        ch.value_axis.tick_labels.number_format_is_linked = False
        ch.value_axis.tick_labels.font.size = Pt(font)
        ch.value_axis.format.line.fill.background()
        ch.category_axis.tick_labels.font.size = Pt(font)
        ch.category_axis.format.line.color.rgb = b.GRID
    except Exception:
        pass
    try:
        ch.plots[0].gap_width = gap
    except Exception:
        pass
    return ch


def add_chart(s, kind, cats, series, left, top, w, h, **kw):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(kind, left, top, w, h, cd)
    return style_chart(gf, **kw)


def table(s, headers, rows, left, top, width, *, col_w, row_h=Inches(0.34),
          head_fill=None, highlight=None, size=11.5, left_cols=(0,)):
    """Lightweight grid from rects + text. highlight: set of row indices.
    left_cols: column indices rendered left-aligned — use for prose cells,
    since centring a long sentence in a table cell is hard to scan."""
    head_fill = head_fill or b.NAVY
    x = left
    d.rect(s, left, top, width, row_h, head_fill)
    for i, htxt in enumerate(headers):
        d.text(s, htxt, x + Inches(0.10), top + Inches(0.05), col_w[i] - Inches(0.16),
               row_h - Inches(0.08), size=size, color=b.WHITE, bold=True, shrink=True,
               align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += col_w[i]
    y = top + row_h
    for r, row in enumerate(rows):
        hl = highlight and r in highlight
        d.rect(s, left, y, width, row_h,
               hx("FFF4D6") if hl else (b.SOFT if r % 2 == 0 else b.WHITE))
        if hl:
            d.rect(s, left, y, Inches(0.055), row_h, b.GOLD)
        x = left
        for i, cell in enumerate(row):
            d.text(s, str(cell), x + Inches(0.10), y + Inches(0.05), col_w[i] - Inches(0.16),
                   row_h - Inches(0.08), size=size, color=b.INK,
                   bold=bool(hl and i == 0), shrink=True,
                   align=PP_ALIGN.LEFT if i in left_cols else PP_ALIGN.CENTER)
            x += col_w[i]
        y += row_h
    return y


def cards(s, items, top, *, cols=3, height=Inches(1.6), gap=Inches(0.22), size=11.5):
    cw = (d.CW - gap * (cols - 1)) / cols
    for i, (title, body) in enumerate(items):
        r, c = divmod(i, cols)
        x = d.M + c * (cw + gap)
        y = top + r * (height + gap)
        d.rect(s, x, y, cw, height, b.WHITE, line=b.GRID, radius=0.06, shadow=True)
        d.rect(s, x, y, cw, Inches(0.05), b.TEAL)
        d.text(s, title, x + Inches(0.18), y + Inches(0.20), cw - Inches(0.36),
               Inches(0.46), size=13, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.18), y + Inches(0.70), cw - Inches(0.36),
               height - Inches(0.86), size=size, color=b.MUTED, shrink=True)


def kpi(s, items, top, *, height=Inches(1.5)):
    n = len(items)
    gap = Inches(0.22)
    cw = (d.CW - gap * (n - 1)) / n
    for i, (num, label, note, col) in enumerate(items):
        x = d.M + i * (cw + gap)
        d.rect(s, x, top, cw, height, b.NAVY, radius=0.05)
        d.text(s, num, x + Inches(0.16), top + Inches(0.16), cw - Inches(0.32), Inches(0.62),
               size=32, color=col, bold=True, font=b.FONT_H, shrink=True)
        d.text(s, label, x + Inches(0.16), top + Inches(0.80), cw - Inches(0.32), Inches(0.30),
               size=11.5, color=b.WHITE, bold=True, shrink=True)
        d.text(s, note, x + Inches(0.16), top + Inches(1.10), cw - Inches(0.32), Inches(0.32),
               size=9.5, color=b.LIGHT_TEAL, shrink=True)


def kpi_at(s, items, left, top, width, *, height=Inches(1.4), gap=Inches(0.18)):
    """KPI row confined to an explicit box — use whenever a chart shares the slide,
    so the cards cannot be laid over the plot area."""
    n = len(items)
    cw = (width - gap * (n - 1)) / n
    for i, (num, label, note, col) in enumerate(items):
        x = left + i * (cw + gap)
        d.rect(s, x, top, cw, height, b.NAVY, radius=0.05)
        d.text(s, num, x + Inches(0.15), top + Inches(0.14), cw - Inches(0.30), Inches(0.56),
               size=27, color=col, bold=True, font=b.FONT_H, shrink=True)
        d.text(s, label, x + Inches(0.15), top + Inches(0.73), cw - Inches(0.30), Inches(0.28),
               size=11, color=b.WHITE, bold=True, shrink=True)
        d.text(s, note, x + Inches(0.15), top + Inches(1.01), cw - Inches(0.30), Inches(0.28),
               size=9.5, color=b.LIGHT_TEAL, shrink=True)


def divider(s, num, title, subtitle):
    d.rect(s, 0, 0, d.W, d.H, b.NAVY)
    d.rect(s, 0, 0, Inches(0.18), d.H, b.TEAL)
    d.text(s, f"SECTION {num}", Inches(0.85), Inches(2.6), Inches(6), Inches(0.4),
           size=14, color=b.TEAL, bold=True)
    d.text(s, title, Inches(0.85), Inches(3.05), Inches(11), Inches(1.0),
           size=42, color=b.WHITE, bold=True, font=b.FONT_H, shrink=True)
    d.rect(s, Inches(0.85), Inches(4.15), Inches(1.6), Inches(0.05), b.TEAL)
    d.text(s, subtitle, Inches(0.85), Inches(4.42), Inches(10.5), Inches(0.6),
           size=15, color=b.LIGHT_TEAL, shrink=True)


def badge(s, txt, left, top, fill, *, w=Inches(1.5), fg=None):
    d.rect(s, left, top, w, Inches(0.30), fill, radius=0.3)
    d.text(s, txt, left, top + Inches(0.045), w, Inches(0.24), size=9.5,
           color=fg or b.WHITE, bold=True, align=PP_ALIGN.CENTER)




ASSET = Path(__file__).resolve().parent / "assets"


def photo(s, name, left, top, width, height):
    """Place a harvested product photo, cropped to fill the box."""
    from PIL import Image
    src = ASSET / "product" / f"{name}.png"
    if not src.exists():
        d.rect(s, left, top, width, height, b.SOFT, line=b.GRID)
        return
    tgt = ASSET / "fitted" / f"{name}_{int(width)}x{int(height)}.png"
    tgt.parent.mkdir(parents=True, exist_ok=True)
    if not tgt.exists():
        im = Image.open(src).convert("RGB")
        want = width / height
        have = im.width / im.height
        if have > want:
            nw = int(im.height * want)
            im = im.crop(((im.width - nw) // 2, 0, (im.width + nw) // 2, im.height))
        else:
            nh = int(im.width / want)
            im = im.crop((0, (im.height - nh) // 2, im.width, (im.height + nh) // 2))
        im.save(tgt)
    s.shapes.add_picture(str(tgt), left, top, width, height)
