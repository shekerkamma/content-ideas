#!/usr/bin/env python3
"""Branded deck for the video: The BEST New AI Opportunity Nobody Is Talking About.

Source input: YouTube video transcript supplied in the task prompt. The watch
preflight was attempted, but local video tooling is unavailable in this
environment, so the transcript is the usable evidence source for the deck.
"""

import sys
from pathlib import Path

PPTXKIT_DIR = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
if not PPTXKIT_DIR.exists():
    PPTXKIT_DIR = Path("/home/shekerk/snap/codex/64/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT_DIR))
from pptx.util import Inches  # noqa: E402
from pptxkit import Deck, PP_ALIGN, MSO_ANCHOR  # noqa: E402


RUN_DIR = Path(__file__).resolve().parent
OUT = RUN_DIR / "agent-native-opportunity-reviewed.pptx"
TOTAL = 22

d = Deck(footer="Agent-Native Opportunity Brief | June 2026")
b = d.b


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullets(items, size=12.5):
    return [{"text": t, "bullet": True, "size": size, "space_before": 5} for t in items]


def card(
    s,
    left,
    top,
    w,
    h,
    title,
    body,
    *,
    fill=None,
    accent=None,
    title_color=None,
    body_color=None,
    body_size=12,
):
    fill = fill or b.SOFT
    accent = accent or b.TEAL
    d.rect(s, left, top, w, h, fill, radius=0.06)
    d.rect(s, left, top, Inches(0.07), h, accent)
    d.text(
        s,
        title,
        left + Inches(0.18),
        top + Inches(0.14),
        w - Inches(0.36),
        Inches(0.28),
        size=14,
        bold=True,
        color=title_color or b.NAVY,
        shrink=True,
    )
    d.text(
        s,
        body,
        left + Inches(0.18),
        top + Inches(0.48),
        w - Inches(0.36),
        h - Inches(0.58),
        size=body_size,
        color=body_color or b.INK,
        shrink=True,
    )


def label(s, text, left, top, w, fill, color):
    d.rect(s, left, top, w, Inches(0.42), fill, radius=0.08)
    d.text(
        s,
        text,
        left + Inches(0.08),
        top + Inches(0.04),
        w - Inches(0.16),
        Inches(0.32),
        size=10.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        shrink=True,
    )


def footer(s, page, dark=False):
    d.footer(s, page, TOTAL, dark=dark)


DIAGRAM_DIR = RUN_DIR / "diagram-assets"
DIAGRAM_DIR.mkdir(parents=True, exist_ok=True)


def pic(s, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    s.shapes.add_picture(str(path), left, top, **kwargs)


def _save_fig(fig, path):
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.02)
    fig.clf()
    return path


def _mpl():
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Circle

    return plt, FancyArrowPatch, FancyBboxPatch, Circle


NAVY = "#0A1628"
NAVY_2 = "#12243A"
TEAL = "#00C9A7"
LIGHT_TEAL = "#E0F7F1"
GOLD = "#FFB800"
CORAL = "#E05A6B"
SOFT = "#F4F7F8"
INK = "#1B2B3C"
MUTED = "#5B6B7C"
GRID = "#D9DFE5"
WHITE = "#FFFFFF"


def make_cover_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(7.0, 4.2), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")
    nodes = [
        (1.1, 2.2, 2.2, 1.1, "Docs"),
        (1.7, 4.2, 2.6, 1.1, "Memory"),
        (6.0, 4.1, 2.6, 1.1, "Tools"),
        (6.4, 2.1, 2.3, 1.1, "Apps"),
    ]
    for x, y, w, h, txt in nodes:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.15",
                                     facecolor=LIGHT_TEAL, edgecolor=TEAL, linewidth=1.8))
        ax.text(x + w / 2, y + h / 2, txt, ha="center", va="center", fontsize=13, color=INK, fontweight="bold")
    ax.add_patch(Circle((4.85, 3.15), 1.15, facecolor=NAVY_2, edgecolor=TEAL, linewidth=2.2))
    ax.text(4.85, 3.33, "USER'S", ha="center", va="center", fontsize=13, color=TEAL, fontweight="bold")
    ax.text(4.85, 2.95, "AGENT", ha="center", va="center", fontsize=18, color=WHITE, fontweight="bold")
    arrows = [
        ((3.3, 2.75), (3.8, 3.05)),
        ((3.85, 4.45), (4.0, 3.8)),
        ((6.0, 4.55), (5.7, 3.9)),
        ((6.2, 2.65), (5.95, 3.0)),
    ]
    for a, b_ in arrows:
        ax.add_patch(FancyArrowPatch(a, b_, arrowstyle="->", mutation_scale=14,
                                      linewidth=2.2, color=GOLD))
    ax.text(4.85, 0.55, "Shared surface: work happens inside the tools the agent can reach", ha="center",
            va="center", fontsize=11.5, color=MUTED, fontstyle="italic")
    return _save_fig(fig, DIAGRAM_DIR / "cover-agent-hub.png")


def make_inversion_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(10.2, 4.2), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.5)
    ax.axis("off")
    # Left panel: app with embedded agent
    ax.add_patch(FancyBboxPatch((0.4, 0.5), 5.0, 4.3, boxstyle="round,pad=0.03,rounding_size=0.18",
                                facecolor=SOFT, edgecolor=CORAL, linewidth=2.0))
    ax.text(2.9, 4.45, "App with a chatbot", ha="center", va="center", fontsize=15, color=NAVY, fontweight="bold")
    ax.add_patch(FancyBboxPatch((1.0, 2.05), 2.8, 1.15, boxstyle="round,pad=0.02,rounding_size=0.12",
                                facecolor=WHITE, edgecolor=CORAL, linewidth=1.8))
    ax.text(2.4, 2.62, "Product UI", ha="center", va="center", fontsize=13, color=INK, fontweight="bold")
    ax.add_patch(Circle((4.05, 2.65), 0.48, facecolor=CORAL, edgecolor=CORAL))
    ax.text(4.05, 2.65, "AI", ha="center", va="center", fontsize=13, color=WHITE, fontweight="bold")
    ax.text(1.02, 1.02, "Only narrow app context\nClosed by design", ha="left", va="bottom",
            fontsize=11.2, color=MUTED)
    # Right panel: agent with app access
    ax.add_patch(FancyBboxPatch((6.6, 0.5), 5.0, 4.3, boxstyle="round,pad=0.03,rounding_size=0.18",
                                facecolor=LIGHT_TEAL, edgecolor=TEAL, linewidth=2.0))
    ax.text(9.1, 4.45, "Agent with app access", ha="center", va="center", fontsize=15, color=NAVY, fontweight="bold")
    ax.add_patch(Circle((7.6, 2.55), 0.55, facecolor=NAVY_2, edgecolor=TEAL, linewidth=2.0))
    ax.text(7.6, 2.75, "USER", ha="center", va="center", fontsize=10, color=TEAL, fontweight="bold")
    ax.text(7.6, 2.45, "AGENT", ha="center", va="center", fontsize=13, color=WHITE, fontweight="bold")
    app = FancyBboxPatch((9.1, 1.95), 2.6, 1.15, boxstyle="round,pad=0.02,rounding_size=0.12",
                         facecolor=WHITE, edgecolor=TEAL, linewidth=1.8)
    ax.add_patch(app)
    ax.text(10.4, 2.53, "App surface", ha="center", va="center", fontsize=13, color=INK, fontweight="bold")
    ax.add_patch(FancyArrowPatch((8.15, 2.55), (9.0, 2.55), arrowstyle="->", mutation_scale=14,
                                 linewidth=2.5, color=GOLD))
    ax.add_patch(FancyBboxPatch((8.8, 1.0), 2.9, 0.72, boxstyle="round,pad=0.02,rounding_size=0.1",
                                facecolor=SOFT, edgecolor=TEAL, linewidth=1.4))
    ax.text(10.25, 1.36, "Context + tools + memory", ha="center", va="center", fontsize=10.8, color=MUTED)
    ax.text(9.1, 1.02, "Can orchestrate from Cursor / Codex / browser", ha="center", va="bottom",
            fontsize=10.8, color=MUTED)
    ax.add_patch(FancyArrowPatch((8.2, 3.25), (10.1, 3.25), arrowstyle="->", mutation_scale=14,
                                 linewidth=2.0, color=TEAL))
    ax.text(6.0, 5.0, "The inversion: the agent owns the context, not the app.", ha="center",
            va="center", fontsize=12.2, color=NAVY, fontweight="bold")
    return _save_fig(fig, DIAGRAM_DIR / "inversion-diagram.png")


def make_thesis_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(10.2, 4.4), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5)
    ax.axis("off")
    boxes = [
        (0.4, 1.5, 2.9, 1.7, "App with\na chatbot", SOFT, CORAL),
        (4.6, 1.5, 2.9, 1.7, "Agent with\napp access", LIGHT_TEAL, TEAL),
        (8.8, 1.5, 2.9, 1.7, "Agent-native\nproduct", SOFT, NAVY),
    ]
    for x, y, w, h, title, fill, edge in boxes:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.03,rounding_size=0.18",
                                    facecolor=fill, edgecolor=edge, linewidth=2.0))
        ax.text(x + w / 2, y + 1.15, title, ha="center", va="center", fontsize=15, color=NAVY, fontweight="bold")
    ax.add_patch(FancyArrowPatch((3.35, 2.35), (4.55, 2.35), arrowstyle="->", mutation_scale=16,
                                 linewidth=2.4, color=GOLD))
    ax.add_patch(FancyArrowPatch((7.55, 2.35), (8.75, 2.35), arrowstyle="->", mutation_scale=16,
                                 linewidth=2.4, color=GOLD))
    ax.text(6.0, 4.15, "Moat shifts from embedded AI to context + integrations + workflow design",
            ha="center", va="center", fontsize=12.2, color=MUTED, fontweight="bold")
    ax.text(6.0, 0.55, "The user brings the agent; the app becomes the surface it uses.",
            ha="center", va="center", fontsize=11.4, color=MUTED, fontstyle="italic")
    return _save_fig(fig, DIAGRAM_DIR / "thesis-diagram.png")


def make_context_stack_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.5, 5.1), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6.2)
    ax.axis("off")
    layers = [
        ("Local files", "Markdown, assets, project docs", LIGHT_TEAL),
        ("Project memory", "Artifacts, notes, remembered decisions", SOFT),
        ("Connected tools", "Browser, CLI, API, service connectors", LIGHT_TEAL),
        ("Shared surface", "Canvas, doc, sheet, dashboard", SOFT),
        ("User oversight", "Review, correction, approval", LIGHT_TEAL),
    ]
    y = 4.9
    for title, body, fill in layers:
        ax.add_patch(FancyBboxPatch((1.0, y), 6.9, 0.72, boxstyle="round,pad=0.02,rounding_size=0.1",
                                    facecolor=fill, edgecolor=TEAL, linewidth=1.5))
        ax.text(1.3, y + 0.37, title, ha="left", va="center", fontsize=12.8, color=NAVY, fontweight="bold")
        ax.text(3.9, y + 0.37, body, ha="left", va="center", fontsize=10.6, color=INK)
        y -= 0.9
    ax.add_patch(FancyArrowPatch((8.4, 5.3), (8.4, 1.7), arrowstyle="->", mutation_scale=18,
                                 linewidth=2.4, color=GOLD))
    ax.text(8.75, 3.5, "quality\nimproves", ha="center", va="center", fontsize=12.5, color=NAVY, fontweight="bold")
    ax.text(8.45, 5.55, "stack", ha="center", va="bottom", fontsize=10.8, color=MUTED)
    ax.text(4.45, 0.45, "The app is the middle layer of a larger workflow stack.",
            ha="center", va="center", fontsize=11.3, color=MUTED, fontstyle="italic")
    return _save_fig(fig, DIAGRAM_DIR / "context-stack.png")


def make_tool_hub_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.6, 5.0), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")
    center = Circle((5.0, 3.0), 0.85, facecolor=NAVY_2, edgecolor=TEAL, linewidth=2.2)
    ax.add_patch(center)
    ax.text(5.0, 3.18, "USER", ha="center", va="center", fontsize=11, color=TEAL, fontweight="bold")
    ax.text(5.0, 2.88, "AGENT", ha="center", va="center", fontsize=15, color=WHITE, fontweight="bold")
    nodes = [
        (1.2, 4.5, "MCP", LIGHT_TEAL, TEAL),
        (8.1, 4.5, "CLI", SOFT, CORAL),
        (1.2, 1.0, "API", LIGHT_TEAL, GOLD),
        (8.1, 1.0, "Browser", SOFT, NAVY),
    ]
    for x, y, txt, fill, edge in nodes:
        ax.add_patch(FancyBboxPatch((x, y), 1.9, 0.85, boxstyle="round,pad=0.02,rounding_size=0.1",
                                    facecolor=fill, edgecolor=edge, linewidth=1.8))
        ax.text(x + 0.95, y + 0.43, txt, ha="center", va="center", fontsize=13, color=NAVY, fontweight="bold")
    arrows = [((4.35, 3.55), (2.95, 4.9)), ((5.65, 3.55), (8.05, 4.9)), ((4.35, 2.5), (2.95, 1.45)), ((5.65, 2.5), (8.05, 1.45))]
    for a, b_ in arrows:
        ax.add_patch(FancyArrowPatch(a, b_, arrowstyle="->", mutation_scale=15, linewidth=2.1, color=GOLD))
    ax.text(5.0, 5.45, "Tool access is what makes the shared surface operational.", ha="center",
            va="center", fontsize=11.9, color=MUTED, fontweight="bold")
    ax.text(5.0, 0.4, "MCP, CLI, API, and browser automation are not details — they are the unlock.", ha="center",
            va="center", fontsize=10.9, color=MUTED, fontstyle="italic")
    return _save_fig(fig, DIAGRAM_DIR / "tool-hub.png")


def make_workflow_loop_diagram():
    plt, FancyArrowPatch, FancyBboxPatch, Circle = _mpl()
    fig, ax = plt.subplots(figsize=(8.8, 4.6), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5.5)
    ax.axis("off")
    steps = [
        (1.0, 3.7, "1. Intent", LIGHT_TEAL, TEAL),
        (4.1, 4.0, "2. Agent acts", SOFT, CORAL),
        (7.0, 3.7, "3. User watches", LIGHT_TEAL, TEAL),
        (4.1, 1.0, "4. User corrects", SOFT, GOLD),
    ]
    for x, y, label, fill, edge in steps:
        ax.add_patch(FancyBboxPatch((x, y), 2.8, 0.95, boxstyle="round,pad=0.02,rounding_size=0.1",
                                    facecolor=fill, edgecolor=edge, linewidth=1.8))
        ax.text(x + 1.4, y + 0.48, label, ha="center", va="center", fontsize=12.6, color=NAVY, fontweight="bold")
    arrows = [
        ((3.85, 4.15), (4.1, 4.45)),
        ((6.9, 4.45), (7.0, 4.15)),
        ((7.0, 3.2), (5.5, 1.95)),
        ((4.1, 1.95), (2.0, 3.2)),
    ]
    for a, b_ in arrows:
        ax.add_patch(FancyArrowPatch(a, b_, arrowstyle="->", mutation_scale=15, linewidth=2.2, color=GOLD))
    ax.text(5.0, 2.75, "Supervised collaboration", ha="center", va="center", fontsize=13.0,
            color=NAVY, fontweight="bold")
    ax.text(5.0, 0.45, "The human stays in the loop and can correct the work live.", ha="center",
            va="center", fontsize=11.0, color=MUTED, fontstyle="italic")
    return _save_fig(fig, DIAGRAM_DIR / "workflow-loop.png")


COVER_IMG = make_cover_diagram()
INVERSION_IMG = make_inversion_diagram()
THESIS_IMG = make_thesis_diagram()
STACK_IMG = make_context_stack_diagram()
TOOL_HUB_IMG = make_tool_hub_diagram()
LOOP_IMG = make_workflow_loop_diagram()
MATRIX_IMG = d.chart_matrix(
    [
        ("Miro", 0.84, 0.86, TEAL),
        ("Notion", 0.72, 0.73, NAVY),
        ("MagicPath", 0.79, 0.67, CORAL),
        ("Sheets", 0.62, 0.58, GOLD),
        ("Adobe", 0.86, 0.44, NAVY_2),
    ],
    DIAGRAM_DIR / "opportunity-matrix.png",
    xlabel="Shared surface richness",
    ylabel="Agent-native fit",
    note="Surface + tools = opportunity",
)


# 1. Cover
s = d.slide(fill=b.NAVY)
d.rect(s, 0, 0, d.W, Inches(0.16), b.TEAL)
d.rect(s, Inches(10.55), 0, Inches(2.78), d.H, b.NAVY_2)
d.rect(s, Inches(10.55), 0, Inches(0.08), d.H, b.TEAL)
pic(s, COVER_IMG, Inches(10.62), Inches(0.55), width=Inches(2.55))
d.text(s, "VIDEO-TO-STRATEGY DECK | 22-SLIDE BRIEF", d.M, Inches(0.92), Inches(8), Inches(0.26),
       size=13, color=b.TEAL, bold=True)
d.text(
    s,
    "The BEST New AI Opportunity\nNobody Is Talking About",
    d.M,
    Inches(1.38),
    Inches(9.0),
    Inches(1.95),
    size=38,
    color=b.WHITE,
    bold=True,
    font=b.FONT_H,
    shrink=True,
)
d.text(
    s,
    "Why agent-native apps win when users bring their own agent into tools like Cursor and Codex",
    d.M,
    Inches(3.62),
    Inches(8.65),
    Inches(0.62),
    size=16,
    color=b.MUTED,
    shrink=True,
)
d.rect(s, d.M, Inches(4.35), Inches(1.42), Inches(0.06), b.TEAL)
d.text(
    s,
    "Source: YouTube transcript, Chris / June 2026",
    d.M,
    Inches(4.58),
    Inches(7.8),
    Inches(0.28),
    size=11.5,
    color=b.MUTED,
)
for i, txt in enumerate(["Context-rich apps", "BYO-agent model", "Super-app surfaces", "MCP / CLI / API"]):
    label(s, txt, d.M + Inches(i * 2.02), Inches(5.55), Inches(1.82), b.NAVY_2, b.WHITE)
footer(s, 1, dark=True)
notes(
    s,
    "This deck translates the transcript into a practical product thesis. The core claim "
    "is that the next opportunity is not another app with a chatbot inside it; it is an "
    "app built so the user's own agent can use it as a shared surface. The deck frames the "
    "model shift, the operating implications, the five-step build process, and the kinds of "
    "products that fit the pattern. The goal is to turn a video idea into a concise strategy "
    "artifact that can be reused in internal discussion or founder planning.",
)


# 2. Model inversion
s = d.slide(fill=b.WHITE)
d.header(s, "The model inverted", "App with an agent vs. agent with an app")
pic(s, INVERSION_IMG, d.M, Inches(1.6), width=d.CW, height=Inches(4.2))
d.rect(s, d.M, Inches(6.0), d.CW, Inches(0.55), b.GOLD, radius=0.04)
d.text(
    s,
    "Implication: the moat moves from embedded AI features to external context, tool access, and shared workflow design.",
    d.M + Inches(0.15),
    Inches(6.09),
    Inches(12.4),
    Inches(0.25),
    size=13.5,
    color=b.NAVY,
    bold=True,
    anchor=MSO_ANCHOR.MIDDLE,
    shrink=True,
)
footer(s, 2)
notes(
    s,
    "The transcript contrasts two software patterns. In the old pattern, the app embeds an "
    "agent. That gives the product a narrow kind of intelligence, but the context remains "
    "trapped inside the application. In the newer pattern, the user's own agent sits outside "
    "the app and uses the app as a surface. That changes the buying center, the product "
    "architecture, and the value proposition: the product now has to work well with the "
    "user's broader context rather than only its own database.",
)


# 3. Why now
s = d.slide(fill=b.WHITE)
d.header(s, "Why this matters now", "Three forces make agent-native apps viable")
rows = [
    ("Super-app workbenches", "Cursor and Codex are becoming the places where knowledge work happens."),
    ("Context wins", "Local files, memories, and project context shape better outcomes than app-only memory."),
    ("Tool connectivity", "MCP, CLI, and API access let an external agent actually do the work in the app."),
    ("Visual feedback loop", "Users can watch the agent act inside the surface and correct it immediately."),
]
y = Inches(1.9)
for i, (title, body) in enumerate(rows):
    fill = b.NAVY if i % 2 == 0 else b.SOFT
    tcol = b.WHITE if i % 2 == 0 else b.NAVY
    d.rect(s, d.M, y, d.CW, Inches(0.95), fill, radius=0.05)
    d.text(s, title, d.M + Inches(0.2), y + Inches(0.12), Inches(3.1), Inches(0.28),
           size=13.8, bold=True, color=tcol, shrink=True)
    d.text(s, body, d.M + Inches(3.2), y + Inches(0.12), Inches(8.9), Inches(0.45),
           size=12.2, color=tcol if i % 2 == 0 else b.INK, shrink=True, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.08)
footer(s, 3)
notes(
    s,
    "The timing argument in the transcript is simple. First, the work surface is shifting "
    "into super-apps like Cursor and Codex. Second, the winning context no longer lives only "
    "inside a product; it lives in local files, project memory, and connected tools. Third, "
    "agent-native apps are practical because there are now standard ways for an agent to "
    "connect to an app through MCP, CLI, or API. The result is a new product category in "
    "which the app is designed to be used by the user's agent, not just by the user directly.",
)


# 4. Use-case realization
s = d.slide(fill=b.WHITE)
d.header(s, "Agent-native app anatomy", "Detailed use-case realization")
card(
    s,
    d.M,
    Inches(1.72),
    Inches(3.65),
    Inches(4.05),
    "Challenge",
    "Users want an app that can work across their full project context, not just the content stored inside one product.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=12.4,
)
card(
    s,
    d.M + Inches(4.2),
    Inches(1.72),
    Inches(4.05),
    Inches(3.65),
    "Solution",
    "Design the app for the user's own agent. Let the app expose a shared surface that works inside Cursor, Codex, or another orchestrator.",
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=12.4,
)
card(
    s,
    d.M + Inches(8.4),
    Inches(1.72),
    Inches(4.05),
    Inches(3.65),
    "How it works",
    "The agent reads local markdown, assets, and memory, then uses the app through browser, CLI, or API while the user watches the result.",
    fill=b.SOFT,
    accent=b.NAVY_2,
    body_size=12.4,
)
d.rect(s, d.M, Inches(5.55), d.CW, Inches(0.9), b.NAVY, radius=0.05)
d.text(
    s,
    "Stack: local context + shared surface + tool connectivity + persistent memory",
    d.M + Inches(0.18),
    Inches(5.72),
    Inches(11.8),
    Inches(0.3),
    size=13.5,
    color=b.WHITE,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
label(s, "Users: founders, designers, operators, analysts", d.M, Inches(6.55), Inches(3.75), b.TEAL, b.NAVY)
label(s, "Organizations: teams that need real context", d.M + Inches(3.95), Inches(6.55), Inches(3.9), b.NAVY_2, b.WHITE)
label(s, "Surface: browser, canvas, sheet, doc, or dashboard", d.M + Inches(8.0), Inches(6.55), Inches(4.25), b.GOLD, b.NAVY)
footer(s, 4)
notes(
    s,
    "This slide captures the product anatomy implied by the transcript. The challenge is "
    "that users need the product to reflect project-specific context that goes beyond the "
    "app. The solution is to make the app a shared surface for the user's agent, not a "
    "closed system. The implementation model is straightforward: the agent reads local "
    "documentation, assets, and memory, then connects to the app through browser control, "
    "CLI, or API. The user stays in the loop and can inspect the work as it happens.",
)


# 5. Build plan
s = d.slide(fill=b.WHITE)
d.header(s, "Five steps to start building", "The transcript's build playbook")
steps = [
    ("1. Pick the job to be done", "Define the outcome the user and agent achieve together."),
    ("2. Define the shared surface", "Choose the canvas both sides collaborate on at the same time."),
    ("3. Build for BYO agent", "Make the app usable from any agent host, not just your own."),
    ("4. Start from your expertise", "Choose a domain you already know deeply."),
    ("5. Plan MCP / CLI / tools early", "Scope integrations on day one so the collaboration is real."),
]
y = Inches(1.95)
for idx, (title, body) in enumerate(steps):
    d.rect(s, d.M, y, Inches(0.8), Inches(0.72), b.TEAL if idx % 2 == 0 else b.NAVY, radius=0.05)
    d.text(s, str(idx + 1), d.M, y + Inches(0.1), Inches(0.8), Inches(0.35),
           size=18, color=b.WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    d.rect(s, d.M + Inches(0.95), y, Inches(11.2), Inches(0.72), b.SOFT if idx % 2 == 0 else b.WHITE, radius=0.05)
    d.text(s, title, d.M + Inches(1.15), y + Inches(0.08), Inches(3.95), Inches(0.24),
           size=13.3, color=b.NAVY, bold=True, shrink=True)
    d.text(s, body, d.M + Inches(5.2), y + Inches(0.08), Inches(6.45), Inches(0.24),
           size=12.1, color=b.INK, shrink=True)
    y += Inches(0.88)
footer(s, 5)
notes(
    s,
    "The five-step framework is the practical takeaway from the video. Step one is to "
    "define the job to be done in terms of a joint outcome for user and agent. Step two is "
    "to define the shared surface where the collaboration happens. Step three is to build "
    "for the bring-your-own-agent model rather than embedding your own agent in a closed "
    "box. Step four is to start where your own expertise gives you a strong wedge. Step "
    "five is to decide the MCP, CLI, and tool list at the beginning so the app can actually "
    "work with agents from day one.",
)


# 6. Opportunity map
s = d.slide(fill=b.WHITE)
d.header(s, "Where the opportunity is", "Examples of agent-native surfaces already visible")
card(
    s,
    d.M,
    Inches(1.85),
    Inches(3.82),
    Inches(1.65),
    "Miro / canvas tools",
    "Shared whiteboards with agent-readable artifacts.",
    fill=b.SOFT,
    accent=b.TEAL,
    body_size=11.3,
)
card(
    s,
    d.M + Inches(3.98),
    Inches(1.85),
    Inches(3.82),
    Inches(1.65),
    "Notion",
    "Knowledge stores external agents can read and write.",
    fill=b.LIGHT_TEAL,
    accent=b.NAVY,
    body_size=11.3,
)
card(
    s,
    d.M + Inches(7.36),
    Inches(1.85),
    Inches(3.82),
    Inches(1.65),
    "MagicPath / Paper",
    "Visual design canvases for agent-led composition.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.3,
)
card(
    s,
    d.M,
    Inches(3.75),
    Inches(5.78),
    Inches(1.55),
    "Google Sheets",
    "Structured data surfaces with connectors.",
    fill=b.LIGHT_TEAL,
    accent=b.GOLD,
    body_size=11.3,
)
card(
    s,
    d.M + Inches(6.35),
    Inches(3.75),
    Inches(5.82),
    Inches(1.55),
    "Adobe-class tools",
    "Incumbents still have room for an agent-native workflow.",
    fill=b.SOFT,
    accent=b.NAVY_2,
    body_size=11.3,
)
label(s, "Opportunity: every surface that already holds context can be agent-native", d.M, Inches(5.7), Inches(12.1), b.NAVY, b.WHITE)
footer(s, 6)
notes(
    s,
    "The transcript names examples that already show the pattern. Miro is shifting toward "
    "shared canvases with markdown, Mermaid diagrams, and HTML prototypes. Notion already "
    "supports MCP-style connectivity and browser access. MagicPath and similar tools let an "
    "agent work directly on a design canvas. Google Sheets is an obvious structured-data "
    "surface with connectors. The broader opportunity is larger than these examples: any "
    "incumbent tool with a surface and an integration path can be recast as agent-native.",
)


# 7. What to build
s = d.slide(fill=b.WHITE)
d.header(s, "What to build", "A concise decision filter")
card(
    s,
    d.M,
    Inches(1.85),
    Inches(2.2),
    Inches(3.8),
    "Build if",
    bullets(
        [
            "The user already works in an existing surface",
            "Your domain knowledge is better than generic tooling",
            "The app can expose MCP, CLI, or API access",
        ]
    ),
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=11.6,
)
card(
    s,
    d.M + Inches(4.1),
    Inches(1.85),
    Inches(3.8),
    Inches(2.2),
    "Avoid if",
    bullets(
        [
            "You need a closed system to feel valuable",
            "The agent has no real external tools",
            "The workflow cannot benefit from project context",
        ]
    ),
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.6,
)
card(
    s,
    d.M + Inches(8.2),
    Inches(1.85),
    Inches(4.0),
    Inches(2.2),
    "Product rule",
    "The app should be useful to the user's agent even when the product's own agent layer is minimal or absent.",
    fill=b.NAVY,
    accent=b.GOLD,
    title_color=b.WHITE,
    body_color=b.LIGHT_TEAL,
    body_size=12.3,
)
d.rect(s, d.M, Inches(4.65), d.CW, Inches(0.9), b.GOLD, radius=0.04)
d.text(
    s,
    "Bottom line: the opportunity is software redesigned for the user's own agent.",
    d.M + Inches(0.18),
    Inches(4.82),
    Inches(12.2),
    Inches(0.3),
    size=13.2,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 7)
notes(
    s,
    "This slide converts the talk into a decision rule. Build where the user already has a "
    "surface, where your domain expertise is real, where the agent can connect through MCP, "
    "CLI, or API, and where local context and memory matter. Avoid products that only work "
    "when closed around their own mini-agent. The product should be useful to the user's "
    "agent even if your embedded AI features are deliberately small. That is the strategic "
    "difference the video is arguing for.",
)


# 8. Strategic takeaway
s = d.slide(fill=b.NAVY)
d.rect(s, 0, 0, d.W, Inches(0.16), b.TEAL)
d.text(s, "Strategic takeaway", d.M, Inches(0.92), Inches(3.2), Inches(0.26),
       size=13, color=b.TEAL, bold=True)
d.text(
    s,
    "Agent-native apps are the new wedge",
    d.M,
    Inches(1.36),
    Inches(9.3),
    Inches(0.42),
    size=30,
    color=b.WHITE,
    bold=True,
    font=b.FONT_H,
    shrink=True,
)
d.text(
    s,
    "Design for the user's agent, not just the user",
    d.M,
    Inches(2.02),
    Inches(7.8),
    Inches(0.28),
    size=19,
    color=b.TEAL,
    bold=True,
)
card(
    s,
    d.M,
    Inches(3.0),
    Inches(12.0),
    Inches(1.62),
    "Summary",
    "The opportunity is to make products that live inside Cursor, Codex, and similar super-apps, with context and tool access designed in from the start.",
    fill=b.NAVY_2,
    accent=b.TEAL,
    title_color=b.WHITE,
    body_color=b.LIGHT_TEAL,
    body_size=12.2,
)
footer(s, 8, dark=True)
notes(
    s,
    "This is a summary checkpoint rather than the end of the deck. The opportunity is not another AI feature. It is a "
    "product designed so the user's own agent can bring local context, memory, and tools to "
    "the app. That shifts the product from isolated automation to collaborative workflow. "
    "For founders and product teams, the practical next step is to pick a domain you know "
    "well, identify the shared surface, and map the MCP, CLI, or API surface before you "
    "build the first version.",
)


# 9. Thesis diagram
s = d.slide(fill=b.WHITE)
d.header(s, "The thesis in one diagram", "From embedded AI to agent-native surfaces")
pic(s, THESIS_IMG, d.M, Inches(1.55), width=d.CW, height=Inches(4.25))
d.rect(s, d.M, Inches(6.0), d.CW, Inches(0.55), b.GOLD, radius=0.04)
d.text(
    s,
    "Moat moves outward: context, integrations, workflow design, and shared surfaces matter more than embedded AI alone.",
    d.M + Inches(0.18),
    Inches(6.09),
    Inches(12.0),
    Inches(0.22),
    size=13.2,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 9)
notes(
    s,
    "This slide captures the core inversion in the transcript. The old model puts a chatbot inside "
    "software. The new model assumes the user's own agent will use software as a surface. That changes "
    "where value accumulates: not in a thin AI layer, but in the context the app can share, the tools it "
    "can expose, and the way the workflow is designed around the agent's actions.",
)


# 10. Hallmarks
s = d.slide(fill=b.WHITE)
d.header(s, "What agent-native means", "Six hallmarks that show up in the transcript")
hallmarks = [
    ("Shared surface", "A canvas, doc, sheet, or dashboard both user and agent can see."),
    ("BYO agent", "It works inside Cursor, Codex, or another orchestrator."),
    ("Project context", "Local files and artifacts matter as much as app data."),
    ("Tool access", "MCP, CLI, API, or browser access lets the agent act."),
    ("Live correction", "The human can watch the work and steer it immediately."),
    ("Persistent memory", "The workflow does not reset every time the session ends."),
]
positions = [
    (d.M, Inches(1.9)),
    (d.M + Inches(4.0), Inches(1.9)),
    (d.M + Inches(8.0), Inches(1.9)),
    (d.M, Inches(3.95)),
    (d.M + Inches(4.0), Inches(3.95)),
    (d.M + Inches(8.0), Inches(3.95)),
]
for (title, body), (x, y) in zip(hallmarks, positions):
    card(s, x, y, Inches(3.75), Inches(1.68), title, body, fill=b.LIGHT_TEAL if x == d.M + Inches(4.0) else b.SOFT, accent=b.TEAL if x != d.M + Inches(4.0) else b.CORAL, body_size=10.8)
d.rect(s, d.M, Inches(6.05), d.CW, Inches(0.55), b.NAVY, radius=0.04)
d.text(
    s,
    "If these six are true, the app is genuinely agent-native.",
    d.M + Inches(0.18),
    Inches(6.14),
    Inches(12.0),
    Inches(0.22),
    size=12.7,
    color=b.WHITE,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 10)
notes(
    s,
    "The transcript keeps returning to a small set of conditions. The product has to expose a shared "
    "surface, allow the user's own agent to enter, carry project context, accept tools, support live "
    "correction, and preserve memory across sessions. Those six conditions are a useful shorthand for "
    "deciding whether a product is actually agent-native or just AI-colored.",
)


# 11. Super-apps
s = d.slide(fill=b.WHITE)
d.header(s, "Why super-app workbenches matter", "Cursor and Codex become the new front door to work")
card(
    s,
    d.M,
    Inches(1.9),
    Inches(5.85),
    Inches(5.8),
    "What changes",
    bullets(
        [
            "People begin work in the agent host, not in your product.",
            "The host carries the broader context and chooses which app to invoke.",
            "The app becomes one surface inside a larger workflow, not the center of gravity.",
        ],
        size=11.7,
    ),
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=11.7,
)
card(
    s,
    d.M + Inches(6.0),
    Inches(1.9),
    Inches(5.8),
    Inches(5.85),
    "What that means for teams",
    bullets(
        [
            "Optimize handoff: docs, permissions, and context sharing.",
            "Design for orchestrators like Cursor and Codex, not only for the end user.",
            "Assume the app will be used by an agent even when the UI is not the starting point.",
        ],
        size=11.7,
    ),
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.7,
)
label(s, "Super-apps are the new interface layer for knowledge work", d.M, Inches(6.95), Inches(12.0), b.GOLD, b.NAVY)
footer(s, 11)
notes(
    s,
    "The transcript names Cursor and Codex as examples of super-apps where knowledge work is increasingly "
    "happening. That matters because the product no longer owns the full working session. It is now a tool "
    "that has to earn a place inside a broader agent-driven workbench. For founders, this means the product's "
    "first job is to make itself useful to the orchestrator.",
)


# 12. Context stack
s = d.slide(fill=b.WHITE)
d.header(s, "The context stack", "The agent's quality comes from more than one database")
pic(s, STACK_IMG, d.M, Inches(1.78), width=Inches(6.4), height=Inches(4.45))
card(
    s,
    d.M + Inches(6.75),
    Inches(1.95),
    Inches(5.05),
    Inches(4.35),
    "Why it matters",
    "The transcript argues that local project context and connected tools outperform app-only memory. The app becomes more useful when it can combine what lives on disk, what the team remembers, and what the agent can reach.",
    fill=b.NAVY,
    accent=b.GOLD,
    title_color=b.WHITE,
    body_color=b.LIGHT_TEAL,
    body_size=11.9,
)
d.rect(s, d.M, Inches(6.33), d.CW, Inches(0.55), b.NAVY_2, radius=0.04)
d.text(
    s,
    "Think of the product as the middle layer of a larger workflow stack, not the whole stack.",
    d.M + Inches(0.18),
    Inches(6.42),
    Inches(12.0),
    Inches(0.22),
    size=12.5,
    color=b.WHITE,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 12)
notes(
    s,
    "A major point in the transcript is that the winning context is not only inside the app. It lives in "
    "local files, memories, and connected tools. The app should therefore be designed as part of a context "
    "stack: the agent can bring information from outside, act on the app's surface, and leave useful artifacts "
    "behind for future sessions.",
)


# 13. Tool access
s = d.slide(fill=b.WHITE)
d.header(s, "Tool access is the unlock", "MCP, CLI, API, and browser access make the app usable")
pic(s, TOOL_HUB_IMG, d.M, Inches(1.72), width=Inches(6.9), height=Inches(4.35))
card(
    s,
    d.M + Inches(7.15),
    Inches(1.95),
    Inches(5.2),
    Inches(4.0),
    "No tool access means no agent-native workflow.",
    "The transcript treats integrations as foundational, not optional.",
    fill=b.NAVY,
    accent=b.TEAL,
    title_color=b.WHITE,
    body_color=b.LIGHT_TEAL,
    body_size=11.8,
)
d.rect(s, d.M + Inches(7.15), Inches(6.05), Inches(5.2), Inches(0.55), b.GOLD, radius=0.04)
d.text(
    s,
    "Give the agent at least one reliable way to act.",
    d.M + Inches(7.33),
    Inches(6.16),
    Inches(4.84),
    Inches(0.2),
    size=12.4,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 13)
notes(
    s,
    "The agent-native idea only works if the agent can actually do something in the product. That is why the "
    "transcript keeps mentioning MCP, CLI, API, and browser access. These are not implementation details; they "
    "are what make the shared surface operational rather than symbolic.",
)


# 14. Human in the loop
s = d.slide(fill=b.WHITE)
d.header(s, "The human stays in the loop", "The visual feedback loop is part of the product")
pic(s, LOOP_IMG, d.M, Inches(1.72), width=d.CW, height=Inches(4.2))
d.rect(s, d.M, Inches(5.98), d.CW, Inches(0.8), b.GOLD, radius=0.04)
d.text(
    s,
    "This is supervised collaboration, not invisible automation.",
    d.M + Inches(0.18),
    Inches(6.14),
    Inches(12.0),
    Inches(0.24),
    size=13.0,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 14)
notes(
    s,
    "The transcript is clear that users should be able to watch the agent work and correct it live. That makes "
    "the feedback loop part of the product. The app is not trying to hide the machine; it is trying to make the "
    "machine legible and steerable inside a shared surface.",
)


# 15. Opportunity taxonomy
s = d.slide(fill=b.WHITE)
d.header(s, "Where the opportunity shows up", "The transcript points to surfaces that already hold context")
pic(s, MATRIX_IMG, d.M, Inches(1.82), width=Inches(7.1), height=Inches(4.45))
card(
    s,
    d.M + Inches(7.35),
    Inches(1.95),
    Inches(5.35),
    Inches(1.45),
    "Collaborative surfaces",
    "Canvas and knowledge tools sit in the high-fit quadrant because they already hold project context.",
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=10.9,
)
card(
    s,
    d.M + Inches(7.35),
    Inches(3.58),
    Inches(5.35),
    Inches(1.45),
    "Structured surfaces",
    "Tables and creative suites are strong because the agent can transform, organize, and refine work.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=10.9,
)
d.rect(s, d.M, Inches(6.12), d.CW, Inches(0.55), b.NAVY, radius=0.04)
d.text(
    s,
    "The common pattern: a surface already holds context, and an external agent can now use it directly.",
    d.M + Inches(0.18),
    Inches(6.22),
    Inches(12.0),
    Inches(0.22),
    size=12.8,
    color=b.WHITE,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 15)
notes(
    s,
    "The opportunity map is not just a list of random products. It is a set of surface types that already carry "
    "project context and can therefore be recast as agent-native. The transcript keeps returning to these kinds "
    "of environments because they are already where work happens.",
)


# 16. Miro
s = d.slide(fill=b.WHITE)
d.header(s, "Miro / canvas tools", "Shared whiteboards with agent-readable artifacts")
card(
    s,
    d.M,
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "What exists",
    "A shared canvas that can hold markdown, Mermaid diagrams, and HTML prototypes.",
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(4.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Why it fits",
    "The board is already a collaboration surface, so an external agent can contribute to it visibly.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(8.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Agent behavior",
    "Generate, edit, and refine visual artifacts on the canvas while the user watches and corrects.",
    fill=b.SOFT,
    accent=b.NAVY,
    body_size=11.2,
)
d.rect(s, d.M, Inches(5.0), d.CW, Inches(0.7), b.GOLD, radius=0.04)
d.text(
    s,
    "Canvas tools are one of the clearest examples of software becoming a shared agent surface.",
    d.M + Inches(0.18),
    Inches(5.13),
    Inches(12.0),
    Inches(0.22),
    size=12.8,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 16)
notes(
    s,
    "The transcript explicitly calls out Miro-style tools. The key detail is that these systems can hold "
    "markdown, diagrams, and prototypes, which makes them readable and editable by an agent. The result is a "
    "shared visual workspace rather than a locked-in whiteboard.",
)


# 17. Notion
s = d.slide(fill=b.WHITE)
d.header(s, "Notion", "Knowledge stores external agents can read and write")
card(
    s,
    d.M,
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "What exists",
    "A knowledge base that already stores the team's docs, tasks, and project memory.",
    fill=b.SOFT,
    accent=b.NAVY,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(4.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Why it fits",
    "The transcript notes MCP-connected access and browser-based use, which makes it easy for agents to work with it.",
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(8.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Agent behavior",
    "Read pages, update task context, and keep the team's project memory in sync across sessions.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.2,
)
footer(s, 17)
notes(
    s,
    "Notion is a strong example because it already sits on top of project memory. The transcript specifically "
    "references MCP-connected access and browser use, which means an external agent can read and write in the "
    "knowledge layer rather than being trapped inside a separate chatbot.",
)


# 18. MagicPath / Paper
s = d.slide(fill=b.WHITE)
d.header(s, "MagicPath / Paper", "Visual design canvases for agent-led composition")
card(
    s,
    d.M,
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "What exists",
    "Design-oriented surfaces where visual composition is the primary interaction pattern.",
    fill=b.LIGHT_TEAL,
    accent=b.CORAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(4.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Why it fits",
    "The transcript treats these as an obvious place for an external agent to shape the canvas visually.",
    fill=b.SOFT,
    accent=b.TEAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(8.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Agent behavior",
    "Generate layouts, move elements, and refine compositions while the designer supervises the result.",
    fill=b.SOFT,
    accent=b.NAVY,
    body_size=11.2,
)
footer(s, 18)
notes(
    s,
    "The transcript mentions MagicPath and Paper as examples of design surfaces where an external agent can shape "
    "the canvas visually. This is a different category from a generic text app: the value is in composition and "
    "visual judgment, which makes the feedback loop central.",
)


# 19. Google Sheets
s = d.slide(fill=b.WHITE)
d.header(s, "Google Sheets", "Structured data surfaces with existing connectors")
card(
    s,
    d.M,
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "What exists",
    "A structured table with familiar formulas, rows, and connector-friendly workflows.",
    fill=b.SOFT,
    accent=b.GOLD,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(4.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Why it fits",
    "The transcript treats Sheets as the obvious case where agents can manipulate structured project data.",
    fill=b.LIGHT_TEAL,
    accent=b.TEAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(8.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Agent behavior",
    "Populate, clean, transform, and validate tables while leaving the human free to inspect changes.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.2,
)
footer(s, 19)
notes(
    s,
    "Google Sheets appears in the transcript as a simple but powerful example: it is already a structured data "
    "surface with obvious connections to the rest of the workflow. That makes it an easy place for an external "
    "agent to do useful work without requiring a new interaction model.",
)


# 20. Adobe-class tools
s = d.slide(fill=b.WHITE)
d.header(s, "Adobe-class tools", "Large incumbents still have room for an agent-native workflow")
card(
    s,
    d.M,
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "What exists",
    "Deep creative suites with rich assets, repeated workflows, and many approval checkpoints.",
    fill=b.LIGHT_TEAL,
    accent=b.NAVY_2,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(4.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Why it fits",
    "The transcript argues that even incumbents still have room for an agent-native version of the workflow.",
    fill=b.SOFT,
    accent=b.CORAL,
    body_size=11.2,
)
card(
    s,
    d.M + Inches(8.0),
    Inches(1.95),
    Inches(3.9),
    Inches(3.9),
    "Agent behavior",
    "Assist with asset handling, repetitive edits, layout variations, and review cycles inside a familiar suite.",
    fill=b.SOFT,
    accent=b.GOLD,
    body_size=11.2,
)
footer(s, 20)
notes(
    s,
    "The transcript does not claim that incumbents disappear. Instead, it says there is still room for an "
    "agent-native version of the workflow inside large creative suites. That is a valuable signal for founders "
    "because it shows how big the surface area still is, even inside mature categories.",
)


# 21. Founder filter
s = d.slide(fill=b.WHITE)
d.header(s, "A practical founder filter", "Ask these questions before you build")
questions = [
    ("1. Where does the user already work?", "If the answer is nowhere, the product may be trying to create a behavior that doesn't exist."),
    ("2. What context lives outside the app?", "Agent-native products win when the real project context is broader than one database."),
    ("3. Which tools must the agent reach?", "MCP, CLI, API, or browser access should be explicit from the start."),
    ("4. How will the human correct it?", "If the user cannot see and steer the work, the product loses the feedback loop."),
]
qy = Inches(1.9)
for i, (title, body) in enumerate(questions):
    card(
        s,
        d.M,
        qy,
        d.CW,
        Inches(0.88),
        title,
        body,
        fill=b.LIGHT_TEAL if i % 2 == 0 else b.SOFT,
        accent=b.TEAL if i % 2 == 0 else b.CORAL,
        body_size=10.8,
    )
    qy += Inches(1.0)
d.rect(s, d.M, Inches(6.15), d.CW, Inches(0.58), b.GOLD, radius=0.04)
d.text(
    s,
    "If you cannot answer all four, the wedge is probably not ready yet.",
    d.M + Inches(0.18),
    Inches(6.25),
    Inches(12.0),
    Inches(0.22),
    size=12.8,
    color=b.NAVY,
    bold=True,
    shrink=True,
    anchor=MSO_ANCHOR.MIDDLE,
)
footer(s, 21)
notes(
    s,
    "This slide turns the transcript into a real decision test. The key questions are whether a user already "
    "works in a surface, whether the important context sits outside the app, whether the agent can reach the "
    "necessary tools, and whether a human can inspect and correct the work.",
)


# 22. Closing
s = d.slide(fill=b.NAVY)
d.rect(s, 0, 0, d.W, Inches(0.16), b.TEAL)
d.text(s, "The next move", d.M, Inches(0.92), Inches(2.7), Inches(0.26),
       size=13, color=b.TEAL, bold=True)
d.text(
    s,
    "Build for the user's agent",
    d.M,
    Inches(1.34),
    Inches(9.2),
    Inches(0.45),
    size=30,
    color=b.WHITE,
    bold=True,
    font=b.FONT_H,
    shrink=True,
)
d.text(
    s,
    "Design the app as a shared surface inside Cursor, Codex, and similar workbenches.",
    d.M,
    Inches(2.0),
    Inches(8.4),
    Inches(0.3),
    size=18,
    color=b.TEAL,
    bold=True,
    shrink=True,
)
card(
    s,
    d.M,
    Inches(3.05),
    Inches(12.0),
    Inches(1.55),
    "Summary",
    "The best new AI opportunity is software that becomes useful to the user's own agent: context-aware, tool-connected, and designed around a visible workflow.",
    fill=b.NAVY_2,
    accent=b.TEAL,
    title_color=b.WHITE,
    body_color=b.LIGHT_TEAL,
    body_size=12.0,
)
for i, txt in enumerate(["Choose a domain you know", "Map the shared surface", "Expose tools early"]):
    label(s, txt, d.M + Inches(i * 4.1), Inches(5.0), Inches(3.55), b.SOFT if i % 2 == 0 else b.NAVY_2, b.NAVY if i % 2 == 0 else b.WHITE)
footer(s, 22, dark=True)
notes(
    s,
    "The closing message is simple: the opportunity is software redesigned for the user's own agent. For founders, "
    "the practical path is to choose a domain you understand, define the shared surface, and map the tools before "
    "shipping the first version.",
)


OUT.parent.mkdir(parents=True, exist_ok=True)
d.save(OUT)
