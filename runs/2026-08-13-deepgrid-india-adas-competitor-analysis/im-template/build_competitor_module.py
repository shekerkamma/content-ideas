#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT=Path(__file__).resolve().parent
RUN=ROOT.parent
WEB=RUN/'assets/web'; AS=ROOT/'assets'
OUT=ROOT/'DeepGrid-Semi-Competitor-Dossiers-IM-Design-draft.pptx'
NAVY='081525'; INK='101827'; TEAL='04B3C7'; WHITE='FFFFFF'; PALE='F1F7F9'; GRAY='65707E'; LINE='D8E2E7'; RED='C62828'; AMBER='C96A0A'; GREEN='16835C'

def rgb(h): return RGBColor.from_string(h)
def box(s,x,y,w,h,fill=WHITE,line=None,r=0):
 sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill); sh.line.color.rgb=rgb(line or fill); return sh
def text(s,t,x,y,w,h,size=12,color=INK,bold=False,serif=False,align='left',italic=False):
 tb=s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0); tf.vertical_anchor=MSO_ANCHOR.TOP
 p=tf.paragraphs[0]; p.text=t; p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
 for r in p.runs: r.font.name='Georgia' if serif else 'Arial'; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(color)
 return tb
def image_crop(s,path,x,y,w,h):
 path=Path(path)
 if not path.exists(): return box(s,x,y,w,h,PALE,LINE)
 with Image.open(path) as im: iw,ih=im.size
 target=w/h; ratio=iw/ih
 pic=s.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
 if ratio>target:
  crop=(iw-ih*target)/2/iw; pic.crop_left=pic.crop_right=crop
 else:
  crop=(ih-iw/target)/2/ih; pic.crop_top=pic.crop_bottom=crop
 return pic
def chrome(s,section,title,sub,page,dark=False):
 bg=NAVY if dark else WHITE; fg=WHITE if dark else INK
 s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(bg)
 text(s,section.upper(),.28,.18,6,.18,7.5,TEAL,True); text(s,title,.28,.40,12.7,.66,19,fg,True,True); text(s,sub,.3,1.09,12.5,.22,9,('A8B8C6' if dark else GRAY)); box(s,.28,1.38,12.78,.025,TEAL if dark else INK)
 box(s,.28,7.14,12.78,.01,('2A3A4B' if dark else LINE)); text(s,'DEEPGRID SEMI · INDIA ADAS COMPETITOR DOSSIER · CONFIDENTIAL',.28,7.23,9,.12,5.8,('8EA0B2' if dark else GRAY),True); text(s,f'{page:02d}',12.55,7.23,.5,.12,5.8,('8EA0B2' if dark else GRAY),True,'', 'right')
def status_color(s): return {'HIGH':RED,'MED–HIGH':AMBER,'MEDIUM':AMBER,'LOW':GREEN,'STRUCTURAL':RED}.get(s,GREEN)

C=[
dict(name='Starkenn',threat='HIGH',headline='Starkenn owns the strongest India field narrative—but not the perception-silicon control point',offer='Radar-led BrakeSafe AEBS plus mining safety and fleet analytics',proof='Company reports 2,500+ vehicles and deployment across mining, fleet and commercial use',arena='Government · mining · retrofit',boundary='AIS-162 type approval and independent homologation remain unverified',response='Partner for radar and field access; counter only where vision VRU proof and lifecycle economics are visible',asset=WEB/'starkenn-brakesafe-img-1.png',source='starkenn.com/brakesafe · starkenn.com/mining · company claims, accessed Aug 2026'),
dict(name='Gahan AI',threat='HIGH',headline='Gahan could become the domestic radar benchmark—but its qualification narrative needs primary proof',offer='DRWIG indigenous 77 GHz radar platform on TI mmWave silicon',proof='Company claims TRL-8, AEC-Q100 completion, TI partnership and OEM pilots',arena='OEM qualification · defence · infrastructure',boundary='Certificate scope, qualified component and named programme status are not independently evidenced',response='Verify first; benchmark radar weather performance and OEM qualification gates',asset=WEB/'gahan-home-fullpage.png',source='gahanai.com/about · claims are company-attributed, accessed Aug 2026'),
dict(name='drivebuddyAI',threat='MEDIUM',headline='drivebuddyAI converts Indian-road data into a credible integrated compliance platform',offer='Camera–radar ADAS, DMS and fleet intelligence in one operating stack',proof='ARAI-certified AIS-184 DMS; company reports 1bn km data and live HCV demonstrations at ADAS Test City',arena='Fleet · compliance · HCV demonstration',boundary='DMS certification does not establish AEBS type approval or braking accountability',response='Diligence OEM integration and fusion performance; avoid conflating DMS proof with AEBS readiness',asset=WEB/'drivebuddy-demo-img-9.jpg',source='Roadzen investor release · ARAI update · drivebuddyAI demonstration, 2024–2026'),
dict(name='ZF CVCS India',threat='HIGH',headline='ZF’s advantage is not perception alone—it already owns the braking channel and OEM safety case',offer='Integrated CV ADAS tied to braking, actuation, validation and homologation',proof='Two disclosed Indian CV OEM ADAS nominations; India validation and production-readiness evidence',arena='N2/N3 OEM programmes',boundary='Programme customer, volumes and unit economics remain undisclosed',response='Attach DGS001 as a bounded camera-perception subsystem; do not compete for the full safety case',asset=AS/'zf-nomination-1.png',source='ZF CVCS India exchange filings, 3 Dec 2025 and 31 Mar 2026'),
dict(name='Aptiv',threat='HIGH',headline='Aptiv has converted a scalable global platform into an India commercial-vehicle design win',offer='Gen 6 radar, smart camera, STRADVISION vision and modular ADAS software',proof='Official selection by a leading commercial-vehicle OEM for future trucks and buses in India',arena='OEM platform · 14 models / 30+ variants',boundary='Named OEM, production economics and DeepGrid-comparable perception metrics are not public',response='Treat as benchmark and potential attach route; win on local data, bounded cost and integration speed',asset=WEB/'aptiv-gen6-img-1.png',source='Aptiv official release, 6 Jan 2026'),
dict(name='STRADVISION',threat='MED–HIGH',headline='STRADVISION attacks DeepGrid’s core layer: efficient vision software on constrained automotive compute',offer='SVNet perception software designed for scalable, resource-constrained hardware',proof='Official India commercial-vehicle supply programme with a global OEM',arena='OEM perception software',boundary='India programme volumes, hardware bill and customer identity are undisclosed',response='Benchmark accuracy-per-watt and Indian ODD; defend with owned compute economics only after measured proof',asset=AS/'stradvision-page.png',source='STRADVISION newsroom, 16 Apr 2026'),
dict(name='Sterling Tools × MINIEYE',threat='MEDIUM',headline='Sterling–MINIEYE combines production ADAS IP with a local Tier-1 route into Indian OEMs',offer='Full-stack ADAS/DMS from MINIEYE; localization and OEM engagement from Sterling',proof='Formal MOU disclosed to NSE/BSE; covers passenger and commercial vehicles',arena='OEM channel · localization',boundary='MOU is market entry—not a disclosed India design win or production nomination',response='Track conversion from MOU to programme; counter through faster local proof and braking-channel partners',asset=AS/'sterling-page.png',source='Sterling Tools NSE/BSE filing and MINIEYE HKEX announcement, Jan 2026'),
dict(name='bitsensing',threat='LOW',headline='bitsensing validates retrofit demand—but its warning-led kit stops short of an AEBS safety case',offer='Radar–camera commercial-vehicle ADAS kit with driver warnings and easy retrofit',proof='Official 500+ bus expansion objective after Korea pilot; UNECE R151/R159 support',arena='Aftermarket fleet warning',boundary='No India deployment or automatic-braking responsibility evidenced',response='Use as demand and installation benchmark; monitor India entry rather than divert core OEM effort',asset=AS/'bitsensing-page.png',source='bitsensing official launch, 5 Feb 2026'),
dict(name='Netrasemi',threat='LOW',headline='Netrasemi is a structural silicon option, not yet a proven India commercial-vehicle ADAS rival',offer='A2000 indigenous edge-AI SoC for vision, DMS, surround view and fleet use cases',proof='Successful silicon bring-up; early customer trials reported in automotive and surveillance',arena='Edge compute · component layer',boundary='Automotive qualification, CV design win and 2027 production readiness remain open',response='Monitor qualification and production; compare silicon economics only when automotive proof emerges',asset=AS/'netrasemi-page.png',source='Netrasemi official solutions and news pages, May 2026'),
dict(name='Incumbent Tier-1 cohort',threat='STRUCTURAL',headline='Bosch, Continental and Valeo win by controlling the complete OEM integration boundary',offer='Sensors, compute, software, braking/chassis integration, validation and global supply',proof='Bosch and Continental offer CV ADAS stacks; Valeo has an India CV nomination and localized VSS360 plan',arena='OEM accountability · industrial scale',boundary='Cohort view masks company-specific India programme differences',response='Do not displace the stack; attach at perception, data or cost control points with explicit safety boundaries',asset=AS/'bosch-page.png',source='Bosch and Continental official CV pages · Valeo official India nomination, May 2026'),
]

DETAIL={
'Starkenn':dict(chain=['150 m radar','BrakeSafe AEBS','Vehicle integration','Mining safety stack','Fleet analytics'],strengths=['India field access','Harsh-environment radar','Retrofit installation'],limits=['Homologation not evidenced','Merchant-silicon economics','No disclosed OEM nomination'],pitch=['HSE: add vision-classified VRUs','Engineering: fuse camera + radar','Procurement: compare lifecycle compute']),
'Gahan AI':dict(chain=['TI mmWave SoC','DRWIG radar','Radar perception','OEM qualification claim','Defence / mobility routes'],strengths=['Indigenous radar narrative','Weather resilience','Claimed OEM pilots'],limits=['Qualification scope unclear','No public production nomination','Claims mostly self-published'],pitch=['OEM: verify certificate scope','Engineering: benchmark monsoon clutter','Board: monitor qualification conversion']),
'drivebuddyAI':dict(chain=['Camera + radar','ADAS / DMS fusion','Driver-risk analytics','ARAI AIS-184 proof','Fleet / OEM interface'],strengths=['Indian-road data','Certified DMS','Live HCV demonstration'],limits=['DMS ≠ AEBS approval','Braking boundary unclear','OEM nomination undisclosed'],pitch=['Fleet: lead with risk analytics','OEM: separate DMS and AEBS gates','Engineering: test false interventions']),
'ZF CVCS India':dict(chain=['Radar + camera','ADAS software','Brake actuation','Homologation','OEM series programme'],strengths=['Braking-system control','Disclosed nominations','India validation stack'],limits=['Closed supplier architecture','Economics undisclosed','Perception source not differentiating'],pitch=['OEM: attach, do not replace','Engineering: offer bounded perception','Purchasing: prove cost-down curve']),
'Aptiv':dict(chain=['Gen 6 radar','Smart camera','STRADVISION vision','Modular ADAS stack','Multi-model OEM rollout'],strengths=['Turn-key platform','Indian CV design win','Global industrial scale'],limits=['Named OEM undisclosed','Global-stack cost burden','Local-data advantage contestable'],pitch=['OEM: sell faster local tuning','Engineering: benchmark accuracy/watt','Purchasing: quantify subsystem savings']),
'STRADVISION':dict(chain=['Camera input','SVNet perception','Hardware abstraction','Safety-certified software','OEM production integration'],strengths=['Efficient vision software','Automotive process maturity','India supply programme'],limits=['No owned sensor/actuation','Depends on compute partner','India economics undisclosed'],pitch=['Engineering: benchmark on same silicon','OEM: contrast local ODD corpus','Board: defend only measured advantage']),
'Sterling Tools × MINIEYE':dict(chain=['MINIEYE ADAS/DMS','India perception tuning','Sterling application engineering','Local manufacturing route','OEM commercialisation'],strengths=['Existing OEM relationships','Production-grade foreign IP','Local industrial footprint'],limits=['MOU not design win','India tuning still ahead','Partner dependency'],pitch=['OEM: ask programme status','Purchasing: expose imported-IP economics','BD: race to braking-channel alliance']),
'bitsensing':dict(chain=['NXP radar','Camera context','Warning fusion','Driver display','Retrofit fleet deployment'],strengths=['Simple installation','Radar-camera package','500-bus expansion plan'],limits=['Warning-led, not AEBS','No India operating proof','Actuation absent'],pitch=['Fleet: compare warning adoption','OEM: distinguish retrofit from series','Sales: use as market validation']),
'Netrasemi':dict(chain=['A2000 edge SoC','Vision acceleration','Automotive use cases','Early customer trials','2027 production target'],strengths=['India silicon narrative','Broad edge-AI applicability','Successful bring-up'],limits=['Qualification not public','No CV design win','Production ramp ahead'],pitch=['Engineering: compare toolchain','OEM: ask qualification roadmap','Board: monitor, do not overreact']),
'Incumbent Tier-1 cohort':dict(chain=['Multi-modal sensors','Central compute','ADAS software','Brake / chassis control','Global OEM accountability'],strengths=['Complete safety case','Industrial capacity','Entrenched OEM access'],limits=['Higher organisational cost','Slower India-specific iteration','Closed-stack incentives'],pitch=['OEM: sell the missing local layer','Tier-1: make integration bounded','Board: fund attach proof, not displacement']),
}
EXEC={
'Starkenn':dict(verdict='The immediate India field competitor—and the most plausible radar-fusion partner',posture='PARTNER + COUNTER',confidence='MEDIUM',usecases=[('MINING AEBS','Dust, vibration and low visibility','Radar-led BrakeSafe + fleet controls'),('GOVT / PSU RETROFIT','Fast fitment and local service','Made-in-India radar AEBS stack'),('CV VRU PROTECTION','Pedestrian and cyclist risk','Radar base; vision remains attach opportunity')]),
'Gahan AI':dict(verdict='A high-potential domestic radar challenger whose qualification claims must clear diligence',posture='VERIFY + BENCHMARK',confidence='LOW–MED',usecases=[('OEM FRONT RADAR','All-weather ranging','DRWIG 77 GHz radar platform'),('TWO-WHEELER SAFETY','Dense mixed traffic','Micro-Doppler and proximity claims'),('DEFENCE / INFRA','Sovereign sensing','India-owned radar narrative')]),
'drivebuddyAI':dict(verdict='The strongest integrated DMS and fleet-data proposition—not yet proven as an AEBS system owner',posture='DILIGENCE',confidence='MED–HIGH',usecases=[('DMS / DDAWS','Fatigue and distraction','ARAI-certified AIS-184 DMS'),('FLEET RISK','Driver behaviour and coaching','Continuously learning analytics'),('HCV COMPLIANCE','Multiple warning functions','Camera–radar integrated platform')]),
'ZF CVCS India':dict(verdict='The reference system partner because it controls braking, validation and OEM accountability',posture='ATTACH',confidence='HIGH',usecases=[('SERIES AEBS','N2/N3 braking intervention','Integrated ADAS + brake actuation'),('BSIS / MOIS','VRU protection around trucks','Multi-sensor CV safety suite'),('OEM HOMOLOGATION','Programme-level release','Validation and safety-case ownership')]),
'Aptiv':dict(verdict='A turn-key platform benchmark already converting India regulation into multi-model OEM rollout',posture='ATTACH / BENCHMARK',confidence='HIGH',usecases=[('MULTI-MODEL CV ADAS','14 models / 30+ variants','Modular Gen 6 platform'),('REGULATORY COMPLIANCE','2027 safety functions','Radar + smart camera + software'),('INDIA ODD','Mixed roads and traffic','Localized platform tuning')]),
'STRADVISION':dict(verdict='The closest software-layer substitute for DeepGrid on constrained automotive compute',posture='BENCHMARK',confidence='HIGH',usecases=[('FRONT VISION','Objects, lanes and VRUs','SVNet perception software'),('LOW-COST COMPUTE','Resource-constrained ECUs','Optimized software architecture'),('OEM SCALE','Cross-platform deployment','Automotive process and integration stack')]),
'Sterling Tools × MINIEYE':dict(verdict='A credible local-channel vehicle for proven foreign ADAS IP—but still at MOU stage in India',posture='TRACK CONVERSION',confidence='HIGH',usecases=[('OEM ADAS / DMS','Passenger and CV compliance','MINIEYE full-stack solutions'),('INDIA LOCALIZATION','Mixed-traffic perception','Sterling tuning and engineering'),('LOCAL SUPPLY','OEM procurement access','Tier-1 manufacturing relationships')]),
'bitsensing':dict(verdict='A practical retrofit warning benchmark with limited direct overlap with certified AEBS',posture='MONITOR',confidence='HIGH',usecases=[('BUS RETROFIT','Blind spots and collision alerts','Radar–camera ADAS kit'),('FLEET MODERNIZATION','Minimal installation downtime','Controller, display and sensors'),('URBAN VRU WARNING','Pedestrians and motorcycles','BSIS, MOIS and surround warnings')]),
'Netrasemi':dict(verdict='A future indigenous compute alternative—not yet a commercial-vehicle ADAS competitor',posture='MONITOR QUALIFICATION',confidence='MEDIUM',usecases=[('DMS / OMS COMPUTE','On-device vision inference','A2000 edge-AI SoC'),('SURROUND VIEW','Multi-camera processing','Integrated vision acceleration'),('FLEET EDGE AI','Local analytics and telematics','Low-latency edge compute')]),
'Incumbent Tier-1 cohort':dict(verdict='The structural OEM gatekeepers; DeepGrid must enter through their open subsystem boundaries',posture='ATTACH / BOUNDARY',confidence='HIGH',usecases=[('FULL CV ADAS','End-to-end safety functions','Sensors, compute and actuation'),('OEM ACCOUNTABILITY','Global programme release','Validation and industrial supply'),('LOCALIZED VSS360','India CV safety suite','Valeo one-box camera and radar architecture')]),
}
ANATOMY={
'Starkenn':dict(identity='India-native vehicle-safety integrator built around retrofit collision mitigation for heavy vehicles, with mining and government fleets as its most visible beachhead.',focus='Commercial vehicles, mine haulage, buses and fleet operators that need a locally installed safety layer without waiting for a new-vehicle programme.',stack=['77 GHz / long-range radar sensing','BrakeSafe collision logic and AEB intervention','Vehicle CAN / brake interface and retrofit fitment','Attention DMS, Terralock and load / alcohol modules','Stark-I fleet telemetry and analytics'],model='Sell a fitted system plus adjacent fleet-safety modules; use field access, installation capability and local service to expand account value.',proof='Company reports 2,500+ fitted vehicles and deployments spanning mining, fleet and commercial use.',unknown='Independent AIS-162 type approval, production-grade safety architecture and named OEM nomination are not evidenced.'),
'Gahan AI':dict(identity='Bengaluru radar venture positioning indigenous automotive sensing as a sovereign alternative for mobility, two-wheelers, infrastructure and defence.',focus='OEM and strategic programmes that value all-weather ranging, India-owned radar IP and a domestic engineering relationship.',stack=['TI mmWave radar silicon','DRWIG 77 GHz radar hardware','Range / velocity / angle and micro-Doppler processing','Radar perception and object-tracking software','Vehicle, two-wheeler and infrastructure interfaces'],model='Develop radar modules and perception IP, then convert pilots and qualification claims into OEM or strategic programme design-ins.',proof='Company states TRL-8 maturity, OEM pilots, TI collaboration and automotive qualification progress.',unknown='Certificate holder, qualified component, AEC-Q100 scope and named production nomination require primary verification.'),
'drivebuddyAI':dict(identity='India-focused ADAS and driver-intelligence platform combining camera, radar, DMS and fleet analytics around a continuously learning road-data layer.',focus='Heavy commercial fleets and OEMs seeking driver monitoring, warning functions, compliance evidence and operational risk reduction in one platform.',stack=['Road-facing cameras and radar','Cabin-facing DMS sensing','ADAS / DMS perception and sensor fusion','Driver-risk scoring and event analytics','Cloud fleet portal, coaching and OEM interfaces'],model='Deploy hardware into fleets, capture Indian-road data, monetize analytics and use operating proof to move upstream toward OEM integration.',proof='ARAI-certified AIS-184 DMS, reported large-scale road-data corpus, HCV demonstrations and a reported multi-thousand-truck contract.',unknown='AEBS type approval, brake-actuation responsibility and exact OEM series-production boundary remain unproven.'),
'ZF CVCS India':dict(identity='Commercial-vehicle systems Tier-1 that joins perception to braking, actuation, validation and OEM release responsibility.',focus='N2/N3 OEM programmes that require multiple mandated safety functions, system-level homologation and a supplier able to own the braking interface.',stack=['Radar, camera and vehicle sensors','CV ADAS perception and decision software','Electronic braking / actuation control','Vehicle integration and India validation','Homologation, production release and lifecycle supply'],model='Win platform nominations by bundling safety functions with entrenched braking relationships and programme-level engineering accountability.',proof='Two disclosed Indian CV OEM ADAS nominations, extensive India validation and a stated multi-feature safety suite.',unknown='Customers, unit economics, programme volumes and the precise external perception-supplier boundary are undisclosed.'),
'Aptiv':dict(identity='Global ADAS platform supplier bringing a modular Gen 6 sensing, compute and software stack into a broad Indian commercial-vehicle rollout.',focus='Large OEM platforms that need one architecture reused across many truck and bus models while meeting future safety requirements.',stack=['Gen 6 radar and smart-camera sensors','STRADVISION SVNet vision software','Sensor fusion and modular ADAS applications','Vehicle network / compute integration','Global validation, manufacturing and programme support'],model='Reuse a standardized platform across high-variant OEM portfolios, localize it through India engineering and deliver at automotive scale.',proof='Official India CV OEM selection covering 14 models and more than 30 variants.',unknown='Named customer, certification status, production pricing and directly comparable India-ODD performance are not public.'),
'STRADVISION':dict(identity='Automotive perception-software specialist selling deployable vision intelligence rather than a complete sensor-to-brake system.',focus='OEMs and Tier-1s needing production-grade camera perception on low-cost or resource-constrained automotive compute.',stack=['Camera image ingestion','SVNet neural-network perception','Objects, lanes, freespace and VRU outputs','Hardware abstraction and SoC optimization','Safety process, tooling and OEM integration'],model='License and integrate perception software across hardware platforms; scale through Tier-1 and OEM production programmes.',proof='Official commercial-vehicle India supply programme and production use through global automotive partners.',unknown='India volumes, compute bill, local data depth and same-hardware performance versus DeepGrid are undisclosed.'),
'Sterling Tools × MINIEYE':dict(identity='Market-entry alliance pairing MINIEYE’s production ADAS / DMS portfolio with Sterling’s India Tier-1 relationships, localization and manufacturing route.',focus='Indian passenger- and commercial-vehicle OEMs seeking established ADAS IP with a domestic commercial and industrial interface.',stack=['MINIEYE cameras, radar and DMS hardware','Perception, fusion and driver-monitoring software','India calibration and application engineering','Sterling manufacturing / supply-chain interface','OEM integration and commercial account access'],model='Import proven technology and localize engineering, sourcing and customer access through an established Indian component supplier.',proof='Formal exchange-disclosed MOU covering ADAS and DMS for Indian passenger and commercial vehicles.',unknown='The MOU has not yet established a disclosed India design win, localized product baseline or production certification.'),
'bitsensing':dict(identity='Radar-specialist offering a practical camera–radar warning kit that modernizes existing commercial fleets without assuming the brake safety case.',focus='Bus and fleet operators needing fast retrofit coverage for blind spots, forward collision and vulnerable-road-user warnings.',stack=['NXP-based radar sensors','Camera context and object classification','Warning-level sensor fusion','Cab display / audible driver alerts','Retrofit controller, harness and installation kit'],model='Sell through fleet and channel partners with low installation downtime; expand vehicle coverage through repeatable retrofit packages.',proof='Official commercial-vehicle kit, Korean bus pilot and stated 500+ bus expansion objective; India channel relationship reported.',unknown='Automatic braking, India fleet deployments and OEM series-production responsibility are not evidenced.'),
'Netrasemi':dict(identity='Indian fabless edge-AI semiconductor company building general-purpose on-device vision compute that could enable automotive applications.',focus='OEMs and device makers seeking indigenous, low-latency AI compute for DMS, surround view, telematics and broader edge workloads.',stack=['A2000 edge-AI SoC','CPU / AI acceleration and vision processing','Multi-camera and on-device inference','SDK, model deployment and application tooling','Board / module integration into customer products'],model='Sell silicon and development platforms to solution builders; convert engineering trials into volume programmes across several edge sectors.',proof='Successful silicon bring-up, fundraising and company-reported early trials in automotive and surveillance use cases.',unknown='AEC-Q100, ISO 26262 readiness, commercial-vehicle design wins and production ramp timing remain open.'),
'Incumbent Tier-1 cohort':dict(identity='Bosch, Continental and Valeo represent the full-stack incumbent model: sensors, compute, software, actuation and global programme accountability.',focus='OEM platforms prioritizing one accountable supplier, validated safety functions, industrial capacity and worldwide lifecycle support.',stack=['Radar, cameras and surround sensing','Central / domain compute and ADAS software','Vehicle network, braking and chassis interfaces','Simulation, validation and homologation','Global manufacturing, quality and OEM programme teams'],model='Bundle system breadth with embedded OEM relationships, amortize engineering across platforms and defend the integration boundary.',proof='Established global CV safety portfolios, India engineering footprints and disclosed India-specific activity including Valeo’s CV nomination.',unknown='Cohort-level comparison masks vendor-specific economics, architecture openness and exact India programme status.'),
}
ANATOMY_TITLES={
'Starkenn':'Starkenn connects radar-led AEB to fleet operations—field integration is its advantage',
'Gahan AI':'Gahan links indigenous 77 GHz radar to OEM applications—but qualification is the conversion gate',
'drivebuddyAI':'drivebuddyAI unifies road sensing, DMS and fleet analytics around an India-data loop',
'ZF CVCS India':'ZF connects sensing to braking and homologation—owning the complete CV safety boundary',
'Aptiv':'Aptiv turns modular sensing and software into a multi-model OEM platform',
'STRADVISION':'STRADVISION supplies the efficient vision layer inside a partner-owned vehicle stack',
'Sterling Tools × MINIEYE':'Sterling–MINIEYE joins production ADAS IP to an Indian Tier-1 route to market',
'bitsensing':'bitsensing packages radar–camera warnings for rapid commercial-fleet retrofit',
'Netrasemi':'Netrasemi provides indigenous edge-AI compute—not yet an automotive safety system',
'Incumbent Tier-1 cohort':'Incumbent Tier-1s win by owning the complete sensor-to-actuation integration boundary',
}
VIS={
'Starkenn':[WEB/'starkenn-brakesafe-img-1.png',WEB/'starkenn-brakesafe-img-2.png',WEB/'starkenn-brakesafe-img-3.png'],
'Gahan AI':[AS/'gahan-page-1.png',AS/'gahan-page-2.png',AS/'gahan-page-3.png'],
'drivebuddyAI':[AS/'curated/drivebuddy-hcv.jpg',AS/'drivebuddy-page-2.png',AS/'drivebuddy-page-3.png'],
'ZF CVCS India':[AS/'curated/zf-truck.jpg',AS/'zf-nomination-1.png',AS/'curated/zf-truck.jpg'],
'Aptiv':[AS/'curated/aptiv-vehicle.jpg',WEB/'aptiv-gen6-img-1.png',AS/'curated/aptiv-vehicle.jpg'],
'STRADVISION':[AS/'stradvision-page.png']*3,
'Sterling Tools × MINIEYE':[AS/'curated/sterling-signing.jpg',AS/'sterling-page.png',AS/'curated/sterling-signing.jpg'],
'bitsensing':[AS/'curated/bitsensing-kit.png',AS/'curated/bitsensing-kit.png',AS/'curated/bitsensing-kit.png'],
'Netrasemi':[AS/'curated/netrasemi-chip.jpg',AS/'curated/netrasemi-chip.jpg',AS/'curated/netrasemi-chip.jpg'],
'Incumbent Tier-1 cohort':[AS/'curated/continental-truck.jpg',AS/'curated/valeo-vss360.png',AS/'curated/continental-truck.jpg'],
}

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
layout=prs.slide_layouts[6]
page=0
for c in C:
 d=DETAIL[c['name']]
 e=EXEC[c['name']]
 page+=1; s=prs.slides.add_slide(layout); chrome(s,'EXECUTIVE SUMMARY',e['verdict'],f"{c['name']} · management posture: {e['posture']}",page,True)
 image_crop(s,VIS[c['name']][0],.34,1.64,4.0,3.06)
 box(s,.34,4.70,4.0,.48,TEAL); text(s,e['posture'],.52,4.84,3.64,.16,10,NAVY,True,'','center')
 summary=[('WHY IT MATTERS',c['proof'],GREEN),('WHERE IT WINS',c['arena'],TEAL),('WHAT IS NOT PROVEN',c['boundary'],AMBER)]
 for i,(lab,val,col) in enumerate(summary):
  y=1.68+i*1.02; box(s,4.68,y,8.32,.84,'10263A','28445B',1); box(s,4.68,y,.07,.84,col); text(s,lab,4.93,y+.14,1.75,.18,8,col,True); text(s,val,6.70,y+.11,6.0,.48,11.5,WHITE,True)
 box(s,4.68,4.80,8.32,.78,WHITE,WHITE,1); text(s,'DEEPGRID POSITION',4.94,5.02,1.65,.16,8,TEAL,True); text(s,c['response'],6.62,4.92,6.05,.4,12.5,INK,True,True)
 box(s,.34,5.78,12.66,.78,'10263A','28445B',1); text(s,'EXECUTIVE DECISION',.60,6.01,1.65,.16,8,TEAL,True); text(s,f"{e['posture']} · Evidence confidence: {e['confidence']} · Release resources only against a named buyer, proof metric and dated gate.",2.30,5.92,10.35,.38,12,WHITE,True)
 s.notes_slide.notes_text_frame.text=f"Executive summary for {c['name']}. Verdict: {e['verdict']}. Evidence: {c['source']}. Boundary: {c['boundary']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'COMPETITOR PROFILE',c['headline'],f"{c['name']} · strongest arena: {c['arena']}",page)
 image_crop(s,VIS[c['name']][0],.28,1.62,5.05,3.25); box(s,.28,4.88,5.05,.55,NAVY); text(s,'SOURCE-DERIVED VISUAL',.48,5.02,1.7,.14,7,TEAL,True); text(s,c['source'],2.05,5.00,3.05,.32,7,WHITE)
 text(s,'WHAT THEY SELL',5.68,1.67,2,.15,8,TEAL,True); text(s,c['offer'],5.68,1.92,7.3,.55,16,INK,True,True)
 cards=[('OPERATING PROOF',c['proof'],GREEN),('BUYER ARENA',c['arena'],TEAL),('EVIDENCE BOUNDARY',c['boundary'],AMBER)]
 for i,(lab,val,col) in enumerate(cards):
  y=2.68+i*.82; box(s,5.68,y,7.35,.68,PALE,LINE,1); box(s,5.68,y,.06,.68,col); text(s,lab,5.9,y+.11,1.55,.14,7.5,col,True); text(s,val,7.35,y+.10,5.42,.43,9.5,INK)
 box(s,.28,5.67,12.75,1.05,NAVY); text(s,'DEEPGRID RESPONSE',.55,5.92,1.8,.2,8,TEAL,True); text(s,c['response'],2.35,5.82,9.25,.52,13,WHITE,True,True); text(s,'THREAT',11.65,6.30,.62,.22,7,'A8B8C6',True); text(s,c['threat'],11.82,6.51,1.05,.18,8.5,status_color(c['threat']),True,'','right')
 s.notes_slide.notes_text_frame.text=f"Evidence status and sources: {c['source']}\nBoundary: {c['boundary']}\nResponse: {c['response']}"

 a=ANATOMY[c['name']]
 page+=1; s=prs.slides.add_slide(layout); chrome(s,'COMPETITOR ANATOMY',ANATOMY_TITLES[c['name']],'What it does · customer focus · technology stack · operating model',page)
 box(s,.30,1.61,4.02,1.20,NAVY,NAVY,1); text(s,'WHAT IT IS',.54,1.85,1.05,.16,8,TEAL,True); text(s,a['identity'],.54,2.10,3.54,.54,10.8,WHITE,True)
 box(s,.30,2.98,4.02,1.28,PALE,LINE,1); text(s,'CUSTOMER & USE-CASE FOCUS',.54,3.20,2.05,.16,8,TEAL,True); text(s,a['focus'],.54,3.48,3.54,.58,10.4,INK)
 box(s,.30,4.43,4.02,1.26,PALE,LINE,1); text(s,'HOW IT MAKES MONEY / SCALES',.54,4.65,2.25,.16,8,AMBER,True); text(s,a['model'],.54,4.93,3.54,.58,10.4,INK)
 text(s,'TECHNOLOGY & DELIVERY STACK',4.68,1.66,3.0,.18,9,TEAL,True)
 for i,label in enumerate(a['stack']):
  y=2.02+i*.58; col=[TEAL,'1A91A3','176C80',AMBER,RED][i]
  box(s,4.68,y,5.04,.44,'F5F9FA',LINE,1); box(s,4.68,y,.08,.44,col); text(s,f'0{i+1}',4.88,y+.11,.34,.16,9,col,True,True); text(s,label,5.30,y+.09,4.16,.23,10.2,INK,True)
  if i<4: text(s,'↓',7.03,y+.42,.35,.18,10,GRAY,True,'','center')
 box(s,10.02,1.62,3.00,1.96,'10263A','28445B',1); text(s,'OPERATING PROOF',10.25,1.86,1.45,.16,8,GREEN,True); text(s,a['proof'],10.25,2.19,2.54,1.08,11.2,WHITE,True)
 box(s,10.02,3.78,3.00,1.91,'FFF7EC','E8C79F',1); text(s,'WHAT REMAINS UNKNOWN',10.25,4.02,1.70,.16,8,AMBER,True); text(s,a['unknown'],10.25,4.35,2.54,1.04,10.8,INK,True)
 box(s,.30,5.98,12.72,.66,NAVY,NAVY,1); text(s,'SO WHAT FOR DEEPGRID',.54,6.18,1.70,.16,8,TEAL,True); text(s,c['response'],2.30,6.10,10.36,.34,11.7,WHITE,True)
 s.notes_slide.notes_text_frame.text=f"Detailed competitor narrative for {c['name']}. Source basis: {c['source']}. Company-attributed claims remain bounded by: {a['unknown']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'THREAT MECHANISM',f"{c['name']} wins only if its control point becomes the buyer’s deciding criterion",'Threat mechanism → falsifier → management action',page,True)
 box(s,8.9,1.62,4.15,2.5,'10263A','28445B',1); text(s,f"BUYER DECISION\n{c['arena']}",9.14,1.91,3.64,.78,15,WHITE,True,True); text(s,'The threat becomes material only when this arena—not standalone perception—sets the supplier decision.',9.14,2.94,3.58,.62,10.5,'B9C8D8')
 stages=[('CONTROL POINT',f"{d['chain'][0]} → {d['chain'][-1]}",TEAL),('CONVERSION MECHANISM',c['arena'],AMBER),('DISPLACEMENT RISK','DeepGrid sits outside the chosen programme boundary',RED)]
 for i,(lab,val,col) in enumerate(stages):
  x=.35+i*2.75; box(s,x,1.78,2.48,2.18,'10263A','28445B',1); text(s,f'0{i+1}',x+.18,1.98,.4,.3,18,col,True,True); text(s,lab,x+.18,2.42,2.1,.2,8,col,True); text(s,val,x+.18,2.78,2.08,.62,11,WHITE,True)
  if i<2: text(s,'→',x+2.47,2.52,.3,.3,20,TEAL,True,'','center')
 box(s,.35,4.40,12.7,.86,'10263A','28445B',1); text(s,'WHY THIS MATTERS',.58,4.60,1.55,.15,8,TEAL,True); text(s,f"The buyer can prefer {c['name']} when {c['arena'].lower()} matters more than standalone perception ownership.",2.2,4.51,10.45,.42,11,WHITE)
 box(s,.35,5.42,12.7,.48,'10263A','28445B',1); text(s,'FALSIFIER',.58,5.57,1.1,.14,8,AMBER,True); text(s,c['boundary'],1.68,5.50,10.95,.26,10.5,WHITE)
 box(s,.35,6.02,12.7,.55,WHITE,WHITE,1); text(s,'MANAGEMENT ACTION',.58,6.19,1.65,.15,8,TEAL,True); text(s,c['response'],2.25,6.10,10.3,.3,12,INK,True)
 s.notes_slide.notes_text_frame.text=f"Threat mechanism for {c['name']}. Fact boundary: {c['boundary']}. Source: {c['source']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'PRODUCT & CONTROL POINTS',f"{c['name']} packages five control points—but only some create a defensible moat",'Follow the system from sensing to the buyer relationship',page)
 image_crop(s,VIS[c['name']][1],.28,1.62,3.05,4.36)
 for i,label in enumerate(d['chain']):
  x=3.58+i*1.85; col=[TEAL,'1A91A3','176C80',AMBER,RED][i]
  box(s,x,2.02,1.52,1.55,PALE,col,1); text(s,f'0{i+1}',x+.14,2.18,.3,.25,15,col,True,True); text(s,label,x+.14,2.67,1.22,.55,10,INK,True)
  if i<4: text(s,'→',x+1.53,2.61,.3,.3,19,TEAL,True,'','center')
 box(s,3.58,3.95,9.45,.78,NAVY); text(s,'WHERE VALUE ACCRUES',3.82,4.13,1.5,.14,7.5,TEAL,True); text(s,c['offer'],5.22,4.06,7.5,.42,12,WHITE,True,True)
 box(s,3.58,4.96,9.45,1.02,PALE,LINE,1); text(s,'DEEPGRID OPENING',3.82,5.18,1.5,.14,7.5,AMBER,True); text(s,c['response'],5.22,5.10,7.48,.5,11,INK,True)
 box(s,.28,6.22,12.75,.48,NAVY); text(s,'SALES MESSAGE',.52,6.37,1.25,.12,7,TEAL,True); text(s,f"Do not sell against the whole stack; isolate the control point that {c['name']} cannot defend economically or locally.",1.78,6.31,10.9,.24,11,WHITE,True)
 s.notes_slide.notes_text_frame.text=f"Product chain synthesis for {c['name']}. Source: {c['source']}. Interpretation is bounded by: {c['boundary']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'BUYER BATTLECARD',f"{c['name']} wins on {d['strengths'][0].lower()}—and breaks at three proof or control boundaries",'Sales battlecard: concede what is true, then move the buying criterion',page,True)
 for col,(lab,items,color) in enumerate([('CONCEDE',d['strengths'],GREEN),('CHALLENGE',d['limits'],AMBER)]):
  x=.34+col*6.35; box(s,x,1.7,6.05,3.18,'10263A','28445B',1); text(s,lab,x+.22,1.94,1.3,.18,9,color,True)
  for i,it in enumerate(items):
   text(s,f'0{i+1}',x+.22,2.38+i*.68,.35,.2,10,color,True,True); text(s,it,x+.68,2.34+i*.68,4.96,.36,12,WHITE,True)
 box(s,.34,4.98,12.4,.55,'10263A','28445B',1); text(s,'EVIDENCE BOUNDARY',.58,5.15,1.55,.14,8,AMBER,True); text(s,c['boundary'],2.18,5.10,10.2,.28,10.5,WHITE)
 box(s,.34,5.72,12.4,.72,WHITE,WHITE,1); text(s,'REFRAME',.58,5.96,1.05,.15,8,TEAL,True); text(s,c['response'],1.66,5.86,10.7,.4,13,INK,True,True)
 s.notes_slide.notes_text_frame.text=f"Sales battlecard for {c['name']}. Concessions are evidence-aligned. Challenges must not be stated as disqualifications without buyer validation. Source: {c['source']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'SALES PLAY',f"DeepGrid can change the {c['name']} conversation with three persona-specific moves",'Lead with the buyer’s objective; end with a measurable proof request',page)
 personas=['BUSINESS / PROCUREMENT','CHIEF ENGINEER','PROGRAMME / SAFETY']
 for i,(persona,msg) in enumerate(zip(personas,d['pitch'])):
  y=1.68+i*1.22; col=[TEAL,AMBER,GREEN][i]; box(s,.34,y,12.4,1.02,PALE,LINE,1); box(s,.34,y,.07,1.02,col); text(s,persona,.62,y+.15,2.05,.17,8.5,col,True); text(s,msg,2.72,y+.12,4.0,.36,13,INK,True,True); text(s,'PROOF ASK',7.05,y+.15,1.0,.14,7,GRAY,True); text(s,['Lifecycle cost + local support','Same-ODD benchmark + integration effort','Named responsibility matrix + release gate'][i],8.02,y+.12,4.35,.4,10.5,INK)
 box(s,.34,5.48,12.4,.68,NAVY); text(s,'STOP / ESCALATE',.62,5.70,1.35,.12,7,TEAL,True); text(s,f"Escalate only when the buyer values DeepGrid’s bounded perception advantage and accepts the evidence plan.",2.05,5.62,10.2,.28,11,WHITE,True)
 box(s,.28,6.35,12.75,.35,TEAL); text(s,'NEXT MEETING OUTPUT',.52,6.45,1.5,.12,7,NAVY,True); text(s,'A named owner, a vehicle or programme, the comparison metric, data rights and a dated go / no-go gate.',2.06,6.41,10.6,.18,10,NAVY,True)
 s.notes_slide.notes_text_frame.text=f"Persona sales play for {c['name']}. Response: {c['response']}. Evidence boundary: {c['boundary']}"

 page+=1; s=prs.slides.add_slide(layout); chrome(s,'USE-CASE CONCLUSION',f"{c['name']} is differentiated in three use cases—but each creates a specific DeepGrid opening",'Close the competitor discussion with a use-case choice, not a generic threat label',page)
 for i,(uc,pain,solution) in enumerate(e['usecases']):
  x=.34+i*4.22; col=[TEAL,AMBER,GREEN][i]
  box(s,x,1.72,3.92,3.12,PALE,LINE,1); box(s,x,1.72,3.92,.08,col)
  text(s,f'0{i+1}',x+.22,1.98,.52,.35,22,col,True,True); text(s,uc,x+.86,2.03,2.78,.3,12,INK,True)
  text(s,'BUYER PROBLEM',x+.22,2.60,1.3,.16,7.5,GRAY,True); text(s,pain,x+.22,2.86,3.45,.45,12,INK,True,True)
  text(s,'UNIQUE SOLUTION',x+.22,3.50,1.3,.16,7.5,col,True); text(s,solution,x+.22,3.77,3.45,.55,12.5,INK,True,True)
 box(s,.34,5.12,12.58,.68,NAVY); text(s,'WHERE DEEPGRID FITS',.60,5.34,1.65,.16,8,TEAL,True); text(s,c['response'],2.30,5.23,10.25,.38,12.5,WHITE,True,True)
 box(s,.34,5.98,12.58,.62,WHITE,TEAL,1); text(s,'COMMERCIAL TEST',.60,6.18,1.5,.15,8,TEAL,True); text(s,f"Select one named {c['arena'].lower()} use case; agree the buyer metric, responsibility boundary and 90-day stop / scale gate.",2.15,6.09,10.42,.34,11.5,INK,True)
 s.notes_slide.notes_text_frame.text=f"Use-case conclusion for {c['name']}. Use cases: {e['usecases']}. DeepGrid response: {c['response']}. Source: {c['source']}"
prs.save(OUT); print(OUT)
