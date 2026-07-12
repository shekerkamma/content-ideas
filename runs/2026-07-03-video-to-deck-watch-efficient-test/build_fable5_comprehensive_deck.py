#!/usr/bin/env python3
from pathlib import Path
import re
import sys

ROOT = Path("/home/shekerk/content-ideas")
RUN = ROOT / "runs/2026-07-03-video-to-deck-watch-efficient-test"
sys.path.insert(0, str(ROOT / "skills/branded-pptx-deck/scripts"))

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt
from pptxkit import Deck, Inches, PP_ALIGN, RGBColor


OUT = RUN / "fable-5-comprehensive-deck-draft.pptx"


def hx(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_footer(d, s, page, total):
    d.text(s, f"{page} / {total}", d.W - Inches(1.2), Inches(7.08), Inches(0.65), Inches(0.18),
           size=8.2, color=d.b.MUTED, bold=True, align=PP_ALIGN.RIGHT)


def header(d, s, title, subtitle=None, kicker=None):
    b = d.b
    d.rect(s, 0, 0, d.W, Inches(0.14), b.TEAL)
    if kicker:
        d.text(s, kicker.upper(), d.M, Inches(0.24), Inches(4.9), Inches(0.16),
               size=8.6, color=b.DARK_TEAL, bold=True, shrink=True)
    d.text(s, title, d.M, Inches(0.44), d.CW, Inches(0.72),
           size=20.8, color=b.NAVY, bold=True, shrink=True)
    d.rect(s, d.M, Inches(1.14), Inches(1.35), Inches(0.045), b.TEAL)
    if subtitle:
        d.text(s, subtitle, d.M, Inches(1.28), d.CW, Inches(0.34), size=10.6, color=b.MUTED, shrink=True)


def label(d, s, text, x, y, w, h, *, fill=None, color=None, size=10.5, bold=True, line=None):
    b = d.b
    fill = fill or b.SOFT
    color = color or b.NAVY
    d.rect(s, x, y, w, h, fill, line=line or b.GRID, radius=0.04)
    d.text(s, text, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), h - Inches(0.12),
           size=size, color=color, bold=bold, shrink=True, align=PP_ALIGN.CENTER)


def card(d, s, title, body, x, y, w, h, accent=None, dark=False):
    b = d.b
    fill = b.NAVY if dark else b.WHITE
    line = accent or b.GRID
    d.rect(s, x, y, w, h, fill, line=line, radius=0.05, shadow=True)
    d.rect(s, x, y, Inches(0.07), h, accent or b.TEAL)
    d.text(s, title, x + Inches(0.22), y + Inches(0.20), w - Inches(0.42), Inches(0.32),
           size=12.2, color=b.WHITE if dark else b.NAVY, bold=True, shrink=True)
    d.text(s, body, x + Inches(0.22), y + Inches(0.64), w - Inches(0.42), h - Inches(0.78),
           size=8.8, color=hx("DCE6EF") if dark else b.INK, shrink=True)


def arrow(d, s, x1, y1, x2, y2, color=None):
    color = color or d.b.TEAL
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.75)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def bullet_box(d, s, title, items, x, y, w, h, accent=None):
    b = d.b
    d.rect(s, x, y, w, h, b.SOFT, line=b.GRID, radius=0.05, shadow=True)
    d.text(s, title, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), Inches(0.30),
           size=10.9, color=accent or b.NAVY, bold=True, shrink=True)
    d.text(s, [{"text": item, "bullet": True, "space_before": 2, "size": 7.0} for item in items],
           x + Inches(0.24), y + Inches(0.54), w - Inches(0.48), h - Inches(0.62),
           color=b.INK, shrink=True)


def takeaway(d, s, text):
    b = d.b
    d.rect(s, Inches(0.72), Inches(6.38), Inches(11.88), Inches(0.45), b.NAVY, radius=0.04)
    d.text(s, "NEXT MOVE", Inches(0.92), Inches(6.47), Inches(0.96), Inches(0.18),
           size=7.8, color=b.GOLD, bold=True, shrink=True)
    d.text(s, text, Inches(1.86), Inches(6.45), Inches(10.30), Inches(0.22),
           size=8.8, color=b.WHITE, bold=True, shrink=True)


def cover(d, total):
    b = d.b
    s = d.slide(fill=b.NAVY)
    d.rect(s, Inches(9.2), 0, Inches(4.15), d.H, b.NAVY_2)
    d.rect(s, Inches(9.2), 0, Inches(0.08), d.H, b.TEAL)
    d.text(s, "EXECUTIVE BRIEFING", d.M, Inches(0.76), Inches(4.4), Inches(0.18),
           size=9.5, color=b.TEAL, bold=True, shrink=True)
    d.text(s, "Use Fable 5 Before It Disappears", d.M, Inches(1.20), Inches(7.8), Inches(1.20),
           size=30, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, d.M, Inches(2.38), Inches(1.45), Inches(0.05), b.TEAL)
    d.text(s, "A comprehensive operating deck for turning temporary model access into durable product and workflow assets.",
           d.M, Inches(2.74), Inches(7.5), Inches(0.42), size=12.4, color=hx("DCE6EF"), shrink=True)
    items = [("1", "Improve reusable skills"), ("2", "Stress-test existing features"), ("3", "Plan ambitious systems")]
    for i, (n, txt) in enumerate(items):
        y = Inches(3.62 + i * 0.78)
        d.rect(s, d.M, y, Inches(0.42), Inches(0.42), b.TEAL, radius=0.05)
        d.text(s, n, d.M, y + Inches(0.11), Inches(0.42), Inches(0.12), size=10, color=b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, txt, d.M + Inches(0.62), y + Inches(0.06), Inches(4.8), Inches(0.24), size=11, color=b.WHITE, bold=True, shrink=True)
    label(d, s, "Temporary access", Inches(9.85), Inches(1.10), Inches(2.55), Inches(0.58), fill=b.WHITE, color=b.NAVY, size=12)
    label(d, s, "Reusable assets", Inches(9.85), Inches(2.30), Inches(2.55), Inches(0.58), fill=b.WHITE, color=b.NAVY, size=12)
    label(d, s, "Implementation plans", Inches(9.85), Inches(3.50), Inches(2.55), Inches(0.58), fill=b.WHITE, color=b.NAVY, size=12)
    arrow(d, s, Inches(11.12), Inches(1.76), Inches(11.12), Inches(2.20), b.GOLD)
    arrow(d, s, Inches(11.12), Inches(2.96), Inches(11.12), Inches(3.40), b.GOLD)
    d.text(s, "Output: assets that remain useful after the access window closes.",
           Inches(9.75), Inches(5.00), Inches(2.78), Inches(0.72), size=9.8, color=b.GOLD, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    add_footer(d, s, 1, total)


def exec_summary(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Use the window on reusable leverage, not random prompts", "The video's central advice is to spend scarce access on work that compounds.", "Executive Summary")
    d.rect(s, d.M, Inches(1.62), d.CW, Inches(0.75), b.NAVY, radius=0.04)
    d.text(s, "BOTTOM LINE", d.M + Inches(0.24), Inches(1.82), Inches(1.05), Inches(0.20), size=7.6, color=b.GOLD, bold=True)
    d.text(s, "The best use of Fable 5 is to improve durable assets: skills you use daily, features that must work, and architecture decisions with broad context.",
           d.M + Inches(1.38), Inches(1.70), d.CW - Inches(1.6), Inches(0.46), size=9.0, color=b.WHITE, bold=True, shrink=True)
    card(d, s, "Where it shines", "Large context, ambiguous systems, and work that needs cross-file or cross-workflow reasoning.", Inches(0.75), Inches(2.78), Inches(3.75), Inches(1.48), b.TEAL)
    card(d, s, "Where to spend it", "Assets and decisions that persist after the window closes, not throwaway experimentation.", Inches(4.80), Inches(2.78), Inches(3.75), Inches(1.48), b.GOLD)
    card(d, s, "What to codify", "Model routing, improved skill instructions, feature punch lists, and implementation-grade plans.", Inches(8.85), Inches(2.78), Inches(3.75), Inches(1.48), b.CORAL)
    bullet_box(d, s, "Executive takeaway", [
        "Start with the highest-leverage reusable asset.",
        "Provide full context, not a thin prompt.",
        "Turn the output into a checklist, spec, or reusable instruction.",
    ], Inches(1.18), Inches(4.78), Inches(4.9), Inches(1.25), b.DARK_TEAL)
    bullet_box(d, s, "Delivery standard", [
        "Every recommendation must map to a specific workflow.",
        "Every workflow needs an owner and next action.",
        "Every output must be saved in the system of work.",
    ], Inches(7.18), Inches(4.78), Inches(4.9), Inches(1.25), b.DARK_TEAL)
    add_footer(d, s, page, total)


def thesis(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Fable is best when the task needs broad context and durable output", "The value is not generic cleverness; it is context absorption plus usable decisions.", "Operating Thesis")
    x0, y = Inches(0.86), Inches(2.04)
    steps = [
        ("Broad context", "Files, notes, product intent, prior decisions"),
        ("Reasoning pass", "Find the pattern, contradiction, or missing path"),
        ("Durable artifact", "Updated skill, implementation plan, or decision map"),
        ("Codified workflow", "Reuse it in daily work after access ends"),
    ]
    for i, (t, body) in enumerate(steps):
        x = x0 + Inches(i * 3.05)
        d.rect(s, x, y, Inches(2.34), Inches(1.34), b.SOFT, line=b.GRID, radius=0.05, shadow=True)
        d.text(s, t, x + Inches(0.15), y + Inches(0.18), Inches(2.04), Inches(0.18), size=11.8, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.15), y + Inches(0.58), Inches(2.04), Inches(0.42), size=9.0, color=b.INK, shrink=True)
        if i < len(steps) - 1:
            arrow(d, s, x + Inches(2.38), y + Inches(0.67), x + Inches(2.88), y + Inches(0.67), b.TEAL)
    d.rect(s, Inches(1.2), Inches(4.30), Inches(10.9), Inches(1.32), b.NAVY, radius=0.04)
    d.text(s, "Decision rule", Inches(1.48), Inches(4.58), Inches(1.4), Inches(0.14), size=8.8, color=b.GOLD, bold=True)
    d.text(s, "Use Fable when losing context would change the answer. Use smaller tools when the task is narrow, mechanical, or already decomposed.",
           Inches(2.88), Inches(4.48), Inches(8.82), Inches(0.26), size=12.5, color=b.WHITE, bold=True, shrink=True)
    add_footer(d, s, page, total)


def three_moves(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The video resolves into three high-value moves", "Each move turns temporary model access into reusable operating leverage.", "Storyboard")
    moves = [
        ("01", "Improve reusable skills", "Daily-driver instructions, skill chains, and workflow helpers."),
        ("02", "Stress-test existing features", "Find root causes and missed implementation paths before launch."),
        ("03", "Plan ambitious systems", "Research and design features that require many moving parts."),
    ]
    for i, (n, t, body) in enumerate(moves):
        x = Inches(0.86 + i * 4.08)
        d.rect(s, x, Inches(2.02), Inches(3.42), Inches(2.54), b.WHITE, line=b.GRID, radius=0.05, shadow=True)
        d.text(s, n, x + Inches(0.22), Inches(2.26), Inches(0.7), Inches(0.22), size=14, color=b.TEAL, bold=True)
        d.text(s, t, x + Inches(0.22), Inches(2.70), Inches(2.9), Inches(0.58), size=12.8, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.22), Inches(3.54), Inches(2.9), Inches(0.58), size=8.8, color=b.INK, shrink=True)
    takeaway(d, s, "Run the moves in order: improve the tools, harden the product, then plan the next complex build.")
    add_footer(d, s, page, total)


def move1(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Move 1: improve skills that shape daily execution", "The first demonstrated use case is a reusable skill-library improvement pass.", "Reusable Assets")
    d.rect(s, Inches(0.90), Inches(1.86), Inches(3.05), Inches(3.72), b.NAVY, radius=0.05, shadow=True)
    d.text(s, "Why this matters", Inches(1.18), Inches(2.16), Inches(2.45), Inches(0.22), size=15, color=b.WHITE, bold=True, shrink=True)
    d.text(s, [{"text": x, "bullet": True, "space_before": 7, "size": 10.2} for x in [
        "Skills are daily drivers, so improvements compound.",
        "Good skills capture intent, context, and handoffs.",
        "The output survives after temporary access closes.",
    ]], Inches(1.18), Inches(2.72), Inches(2.40), Inches(1.80), color=hx("DCE6EF"), shrink=True)
    label(d, s, "Current skill", Inches(4.85), Inches(2.05), Inches(2.10), Inches(0.70), fill=b.SOFT)
    label(d, s, "Fable review", Inches(7.60), Inches(2.05), Inches(2.10), Inches(0.70), fill=b.LIGHT_TEAL)
    label(d, s, "Improved instruction", Inches(10.35), Inches(2.05), Inches(2.10), Inches(0.70), fill=b.SOFT)
    arrow(d, s, Inches(6.98), Inches(2.40), Inches(7.52), Inches(2.40))
    arrow(d, s, Inches(9.73), Inches(2.40), Inches(10.27), Inches(2.40))
    bullet_box(d, s, "What Fable should inspect", [
        "Trigger description and activation clarity.",
        "Dependency handoffs between related skills.",
        "Edge cases where the skill can drift.",
        "Concrete acceptance criteria for outputs.",
    ], Inches(4.55), Inches(3.36), Inches(3.55), Inches(1.72), b.DARK_TEAL)
    bullet_box(d, s, "What gets codified", [
        "Shorter instructions.",
        "Sharper routing rules.",
        "Explicit deliverable gates.",
        "Reusable workflow sequence.",
    ], Inches(8.70), Inches(3.36), Inches(3.55), Inches(1.72), b.DARK_TEAL)
    takeaway(d, s, "Do not spend the window on isolated prompts; spend it on the instructions that run your work repeatedly.")
    add_footer(d, s, page, total)


def skill_workflow(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Skill improvement becomes a repeatable four-step workflow", "The demo flow can be standardized across any high-use skill or plugin.", "Workflow")
    steps = [
        ("Select model", "Route the session to Fable for the review pass."),
        ("Load context", "Provide the skill, related helpers, and desired behavior."),
        ("Compare findings", "Separate useful changes from generic advice."),
        ("Codify updates", "Patch instructions, rerun, and keep the improved version."),
    ]
    for i, (t, body) in enumerate(steps):
        x = Inches(0.86 + i * 3.08)
        d.rect(s, x, Inches(2.02), Inches(2.38), Inches(1.62), b.WHITE, line=b.GRID, radius=0.05, shadow=True)
        d.rect(s, x, Inches(2.02), Inches(2.38), Inches(0.13), b.TEAL)
        d.text(s, t, x + Inches(0.16), Inches(2.34), Inches(2.05), Inches(0.20), size=11.8, color=b.NAVY, bold=True, shrink=True)
        d.text(s, body, x + Inches(0.16), Inches(2.74), Inches(2.05), Inches(0.45), size=8.8, color=b.INK, shrink=True)
        if i < 3:
            arrow(d, s, x + Inches(2.42), Inches(2.83), x + Inches(2.92), Inches(2.83))
    d.rect(s, Inches(1.18), Inches(4.42), Inches(10.96), Inches(1.14), b.SOFT, line=b.GRID, radius=0.04)
    d.text(s, "Quality bar", Inches(1.45), Inches(4.64), Inches(1.15), Inches(0.22), size=8.3, color=b.DARK_TEAL, bold=True)
    d.text(s, "A good result changes behavior: clearer trigger, fewer ambiguous handoffs, stronger edge-case handling, and a cleaner next-run checklist.",
           Inches(2.58), Inches(4.54), Inches(9.18), Inches(0.38), size=10.0, color=b.NAVY, bold=True, shrink=True)
    add_footer(d, s, page, total)


def quality_criteria(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Good skill feedback improves intent, handoffs, and edge-case behavior", "The goal is not a longer instruction; the goal is a more reliable operator.", "Quality Criteria")
    rows = [
        ("Intent", "The trigger and user promise are clear enough to route correctly."),
        ("Context", "The skill knows what files, examples, and constraints matter."),
        ("Handoffs", "Upstream and downstream skills receive explicit artifacts."),
        ("Failure handling", "Blocked states, fallback paths, and status labels are defined."),
        ("Verification", "The skill requires evidence before marking work complete."),
    ]
    x, y = Inches(1.0), Inches(1.78)
    for i, (left, right) in enumerate(rows):
        yy = y + Inches(i * 0.78)
        d.rect(s, x, yy, Inches(2.25), Inches(0.56), b.NAVY if i == 0 else b.SOFT, line=b.GRID)
        d.rect(s, x + Inches(2.25), yy, Inches(8.90), Inches(0.56), b.WHITE, line=b.GRID)
        d.text(s, left, x + Inches(0.14), yy + Inches(0.13), Inches(1.88), Inches(0.22),
               size=8.8, color=b.WHITE if i == 0 else b.NAVY, bold=True, shrink=True)
        d.text(s, right, x + Inches(2.45), yy + Inches(0.13), Inches(8.25), Inches(0.24),
               size=8.5, color=b.INK, shrink=True)
    takeaway(d, s, "Treat skill updates as product work: fewer words, clearer behavior, better verification.")
    add_footer(d, s, page, total)


def move2(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Move 2: stress-test existing features before they reach users", "The second use case turns Fable into a project reviewer for systems already in motion.", "Product Readiness")
    card(d, s, "Input", "A project with a feature that almost works but has unresolved failures.", Inches(0.86), Inches(2.0), Inches(3.25), Inches(1.5), b.TEAL)
    card(d, s, "Review lens", "Ask for root cause, architecture fit, tool usage, and why prior fixes did not solve it.", Inches(5.04), Inches(2.0), Inches(3.25), Inches(1.5), b.GOLD)
    card(d, s, "Output", "A ranked fix plan with the real bottleneck, implementation sequence, and verification steps.", Inches(9.22), Inches(2.0), Inches(3.25), Inches(1.5), b.CORAL)
    arrow(d, s, Inches(4.16), Inches(2.75), Inches(4.96), Inches(2.75))
    arrow(d, s, Inches(8.34), Inches(2.75), Inches(9.14), Inches(2.75))
    d.rect(s, Inches(1.18), Inches(4.35), Inches(10.96), Inches(1.12), b.NAVY, radius=0.04)
    d.text(s, "Practical standard", Inches(1.44), Inches(4.58), Inches(1.38), Inches(0.22), size=7.9, color=b.GOLD, bold=True)
    d.text(s, "A useful review explains why the previous attempts failed and what specific change should be made first.",
           Inches(2.82), Inches(4.48), Inches(8.9), Inches(0.40), size=10.6, color=b.WHITE, bold=True, shrink=True)
    add_footer(d, s, page, total)


def feature_workflow(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Feature review should move from symptom to implementation sequence", "The recurring trap is fixing visible errors while missing the actual system behavior.", "Diagnostic Workflow")
    lanes = [
        ("Symptom", "What keeps failing?"),
        ("System path", "Where does the work actually run?"),
        ("Tool use", "Which tools are used or ignored?"),
        ("Root cause", "What explains repeated failure?"),
        ("Fix order", "What changes first, second, third?"),
    ]
    for i, (t, body) in enumerate(lanes):
        x = Inches(0.74 + i * 2.52)
        fill = b.NAVY if i == 3 else b.SOFT
        col = b.WHITE if i == 3 else b.NAVY
        label(d, s, t, x, Inches(2.00), Inches(1.88), Inches(0.55), fill=fill, color=col, size=10.4)
        d.text(s, body, x, Inches(2.74), Inches(1.88), Inches(0.44), size=7.9, color=b.INK, shrink=True, align=PP_ALIGN.CENTER)
        if i < len(lanes) - 1:
            arrow(d, s, x + Inches(1.94), Inches(2.28), x + Inches(2.38), Inches(2.28))
    bullet_box(d, s, "Review questions", [
        "What did previous attempts assume incorrectly?",
        "Which server or client path owns the behavior?",
        "Where should logging, events, or state changes happen?",
        "What proves the fix actually worked?",
    ], Inches(1.05), Inches(4.22), Inches(5.15), Inches(1.34), b.DARK_TEAL)
    bullet_box(d, s, "Expected output", [
        "A precise diagnosis.",
        "A sequenced fix plan.",
        "A verification checklist.",
        "A clear owner for follow-through.",
    ], Inches(7.10), Inches(4.22), Inches(5.15), Inches(1.34), b.DARK_TEAL)
    add_footer(d, s, page, total)


def move3(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Move 3: use Fable for ambitious systems that need many moving parts", "The third use case is planning and research where context fragmentation causes bad decisions.", "Complex Planning")
    d.rect(s, Inches(0.86), Inches(1.90), Inches(11.60), Inches(3.38), b.SOFT, line=b.GRID, radius=0.05)
    components = [
        ("Notes", Inches(1.20), Inches(2.35)),
        ("Research", Inches(3.25), Inches(2.35)),
        ("Existing code", Inches(5.30), Inches(2.35)),
        ("Architecture", Inches(7.35), Inches(2.35)),
        ("User experience", Inches(9.40), Inches(2.35)),
    ]
    for t, x, y in components:
        label(d, s, t, x, y, Inches(1.45), Inches(0.55), fill=b.WHITE, size=9.4)
    for i in range(4):
        arrow(d, s, Inches(2.68 + i * 2.05), Inches(2.62), Inches(3.14 + i * 2.05), Inches(2.62))
    d.rect(s, Inches(4.26), Inches(3.72), Inches(4.88), Inches(0.74), b.NAVY, radius=0.04)
    d.text(s, "Integrated implementation plan", Inches(4.56), Inches(3.93), Inches(4.28), Inches(0.26), size=11.2, color=b.WHITE, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    arrow(d, s, Inches(6.70), Inches(3.02), Inches(6.70), Inches(3.64), b.GOLD)
    takeaway(d, s, "Use the larger context window when the answer depends on combining notes, docs, code, and product intent.")
    add_footer(d, s, page, total)


def streaming_workflow(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The complex example is a server-side streaming workflow", "The model is asked to reason across backend events, token streaming, and the user-facing experience.", "Architecture Thinking")
    y = Inches(2.05)
    boxes = [
        ("Server functions", "Where work executes", Inches(0.92)),
        ("Event stream", "Progress and state changes", Inches(3.46)),
        ("Token stream", "Generated response updates", Inches(6.00)),
        ("Front end", "User-visible status and result", Inches(8.54)),
        ("Controls", "Pause, interrupt, resume", Inches(11.08)),
    ]
    for i, (t, body, x) in enumerate(boxes):
        d.rect(s, x, y, Inches(1.72), Inches(1.38), b.WHITE, line=b.GRID, radius=0.05, shadow=True)
        d.text(s, t, x + Inches(0.13), y + Inches(0.20), Inches(1.46), Inches(0.18), size=10.2, color=b.NAVY, bold=True, shrink=True, align=PP_ALIGN.CENTER)
        d.text(s, body, x + Inches(0.13), y + Inches(0.62), Inches(1.46), Inches(0.36), size=8.0, color=b.INK, shrink=True, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            arrow(d, s, x + Inches(1.76), y + Inches(0.69), x + Inches(2.42), y + Inches(0.69))
    bullet_box(d, s, "Planning questions", [
        "Which events should stream to the user?",
        "Which state changes need persistence?",
        "How do pause and interrupt cases behave?",
        "What feedback proves the system is still working?",
    ], Inches(1.22), Inches(4.38), Inches(4.90), Inches(1.28), b.DARK_TEAL)
    bullet_box(d, s, "Output needed", [
        "Component responsibilities.",
        "Event vocabulary.",
        "Implementation order.",
        "Verification cases.",
    ], Inches(7.22), Inches(4.38), Inches(4.90), Inches(1.28), b.DARK_TEAL)
    add_footer(d, s, page, total)


def comparison(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The advantage is fewer context breaks and better cross-cutting reasoning", "Smaller chunks still work, but they can miss system-wide interactions.", "Operating Comparison")
    d.rect(s, Inches(0.92), Inches(1.88), Inches(5.46), Inches(3.70), b.SOFT, line=b.GRID, radius=0.05, shadow=True)
    d.rect(s, Inches(6.96), Inches(1.88), Inches(5.46), Inches(3.70), b.LIGHT_TEAL, line=b.GRID, radius=0.05, shadow=True)
    d.text(s, "Small-chunk approach", Inches(1.24), Inches(2.20), Inches(4.78), Inches(0.22), size=15, color=b.NAVY, bold=True, shrink=True)
    d.text(s, "Integrated planning approach", Inches(7.28), Inches(2.20), Inches(4.78), Inches(0.22), size=15, color=b.NAVY, bold=True, shrink=True)
    left = ["Break problem into many parts", "Research each part separately", "Manually merge recommendations", "Risk conflicting assumptions"]
    right = ["Load broader context at once", "Reason across code, docs, and intent", "Return a sequenced plan", "Expose contradictions earlier"]
    d.text(s, [{"text": x, "bullet": True, "space_before": 8, "size": 10.5} for x in left],
           Inches(1.28), Inches(2.88), Inches(4.72), Inches(1.74), color=b.INK, shrink=True)
    d.text(s, [{"text": x, "bullet": True, "space_before": 8, "size": 10.5} for x in right],
           Inches(7.32), Inches(2.88), Inches(4.72), Inches(1.74), color=b.INK, shrink=True)
    takeaway(d, s, "Use Fable when the risk is not a missing answer, but a fragmented answer.")
    add_footer(d, s, page, total)


def caught(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The most useful findings are the ones that change the plan", "The video highlights several issue classes that matter before implementation.", "Findings Pattern")
    items = [
        ("Repeated loops", "Prior fixes kept circling without addressing the real behavior."),
        ("Pause and interrupt handling", "Control paths needed clearer state and user-facing behavior."),
        ("Incorrect assumptions", "Some implementation ideas were factually wrong for the stack."),
        ("Missed areas", "The review surfaced parts of the system not previously inspected."),
    ]
    for i, (t, body) in enumerate(items):
        x = Inches(0.92 + (i % 2) * 6.0)
        y = Inches(1.90 + (i // 2) * 1.58)
        card(d, s, t, body, x, y, Inches(5.20), Inches(1.10), [b.TEAL, b.GOLD, b.CORAL, b.DARK_TEAL][i])
    d.rect(s, Inches(1.22), Inches(5.42), Inches(10.86), Inches(0.58), b.NAVY, radius=0.04)
    d.text(s, "A good model pass should not just sound smart; it should change the implementation order.",
           Inches(1.52), Inches(5.54), Inches(10.26), Inches(0.30), size=9.4, color=b.WHITE, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    add_footer(d, s, page, total)


def action_plan(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Turn the remaining window into a five-day asset sprint", "The operating plan converts temporary access into saved, reusable improvements.", "Action Plan")
    days = [
        ("Day 1", "Pick assets", "Choose one high-use skill and one product feature."),
        ("Day 2", "Improve skills", "Run the review, patch instructions, and retest."),
        ("Day 3", "Review feature", "Diagnose root cause and build the fix sequence."),
        ("Day 4", "Plan complex work", "Feed notes, code, docs, and desired behavior."),
        ("Day 5", "Codify", "Save plans, update checklists, and lock the workflow."),
    ]
    for i, (day, title, body) in enumerate(days):
        x = Inches(0.70 + i * 2.50)
        d.rect(s, x, Inches(2.08), Inches(2.05), Inches(2.42), b.WHITE, line=b.GRID, radius=0.05, shadow=True)
        d.rect(s, x, Inches(2.08), Inches(2.05), Inches(0.42), b.NAVY)
        d.text(s, day, x + Inches(0.18), Inches(2.18), Inches(1.65), Inches(0.18), size=7.8, color=b.GOLD, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, title, x + Inches(0.18), Inches(2.76), Inches(1.65), Inches(0.30), size=10.8, color=b.NAVY, bold=True, shrink=True, align=PP_ALIGN.CENTER)
        d.text(s, body, x + Inches(0.18), Inches(3.26), Inches(1.65), Inches(0.68), size=7.2, color=b.INK, shrink=True, align=PP_ALIGN.CENTER)
    d.rect(s, Inches(1.24), Inches(5.30), Inches(10.84), Inches(0.60), b.SOFT, line=b.GRID, radius=0.04)
    d.text(s, "Success condition: by the end of the sprint, the work is embedded in skills, specs, checklists, or implementation plans.",
           Inches(1.52), Inches(5.44), Inches(10.28), Inches(0.26), size=9.4, color=b.NAVY, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    add_footer(d, s, page, total)


def conclusion(d, page, total):
    b = d.b
    s = d.slide(fill=b.NAVY)
    d.text(s, "Conclusion", d.M, Inches(0.72), Inches(3.0), Inches(0.22), size=10, color=b.TEAL, bold=True)
    d.text(s, "Use Fable for the work that keeps paying back", d.M, Inches(1.12), Inches(8.2), Inches(1.08),
           size=25.5, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, d.M, Inches(2.20), Inches(1.35), Inches(0.05), b.TEAL)
    d.text(s, "The video's practical message is simple: do not spend scarce access on isolated experiments. Spend it on the skills, product reviews, and complex plans that become part of the operating system.",
           d.M, Inches(2.58), Inches(7.35), Inches(0.96), size=11.2, color=hx("DCE6EF"), shrink=True)
    for i, txt in enumerate(["Improve daily-driver skills", "Stress-test existing features", "Plan hard systems with full context"]):
        y = Inches(4.12 + i * 0.64)
        d.rect(s, d.M, y, Inches(0.36), Inches(0.36), b.TEAL, radius=0.04)
        d.text(s, str(i + 1), d.M, y + Inches(0.09), Inches(0.36), Inches(0.10), size=8.5, color=b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, txt, d.M + Inches(0.56), y + Inches(0.03), Inches(5.5), Inches(0.30), size=8.9, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, Inches(8.7), Inches(1.45), Inches(3.35), Inches(3.95), b.NAVY_2, line=b.TEAL, radius=0.05, shadow=True)
    d.text(s, "Recommended next move", Inches(9.05), Inches(1.92), Inches(2.65), Inches(0.24), size=14, color=b.GOLD, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    d.text(s, "Pick the highest-leverage reusable asset today and turn Fable's output into a committed workflow change.",
           Inches(9.08), Inches(2.54), Inches(2.58), Inches(1.08), size=11.0, color=b.WHITE, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    d.text(s, f"{page} / {total}", d.W - Inches(1.2), Inches(7.08), Inches(0.65), Inches(0.18),
           size=9, color=hx("889098"), bold=True, align=PP_ALIGN.RIGHT)


def scan_visible_text(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    banned = [
        "transcript", "hyperframe", "excalidraw", "youtube", "source", "validation",
        "synthesis", "audit", "codex", "claude", "/home/", "runs/",
    ]
    hits = [term for term in banned if term in text.lower()]
    time_hits = re.findall(r"\b\d{1,2}:\d{2}\b", text)
    if hits or time_hits:
        raise RuntimeError(f"Visible text scan failed: terms={hits}, time_hits={time_hits[:5]}")


def main():
    d = Deck(footer="")
    total = 15
    cover(d, total)
    exec_summary(d, 2, total)
    thesis(d, 3, total)
    three_moves(d, 4, total)
    move1(d, 5, total)
    skill_workflow(d, 6, total)
    quality_criteria(d, 7, total)
    move2(d, 8, total)
    feature_workflow(d, 9, total)
    move3(d, 10, total)
    streaming_workflow(d, 11, total)
    comparison(d, 12, total)
    caught(d, 13, total)
    action_plan(d, 14, total)
    conclusion(d, 15, total)
    d.save(OUT)
    scan_visible_text(OUT)
    print(f"Visible text scan passed: {OUT}")


if __name__ == "__main__":
    main()
