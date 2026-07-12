#!/usr/bin/env python3
from pathlib import Path
import re
import sys

ROOT = Path("/home/shekerk/content-ideas")
RUN = ROOT / "runs/2026-07-03-video-to-deck-watch-efficient-test"
sys.path.insert(0, str(ROOT / "skills/branded-pptx-deck/scripts"))

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt
from pptxkit import Deck, Inches, PP_ALIGN, RGBColor

OUT = RUN / "fable-5-comprehensive-deck-v2-draft.pptx"


def hx(s):
    s = s.strip("#")
    return RGBColor(int(s[:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def header(d, s, section, title, subtitle=""):
    b = d.b
    d.rect(s, 0, 0, d.W, Inches(0.13), b.TEAL)
    d.text(s, section.upper(), d.M, Inches(0.20), Inches(5.0), Inches(0.26), size=6.6, color=b.DARK_TEAL, bold=True, shrink=True)
    d.text(s, title, d.M, Inches(0.42), d.CW, Inches(0.86), size=16.2, color=b.NAVY, bold=True, shrink=True)
    d.rect(s, d.M, Inches(1.19), Inches(1.25), Inches(0.045), b.TEAL)
    if subtitle:
        d.text(s, subtitle, d.M, Inches(1.34), d.CW, Inches(0.30), size=8.7, color=b.MUTED, shrink=True)


def footer(d, s, page, total, dark=False):
    col = hx("889098") if dark else d.b.MUTED
    d.text(s, f"{page} / {total}", d.W - Inches(1.18), Inches(7.08), Inches(0.65), Inches(0.18), size=8, color=col, bold=True, align=PP_ALIGN.RIGHT)


def card(d, s, title, body, x, y, w, h, accent=None, fill=None, dark=False):
    b = d.b
    fill = fill or (b.NAVY if dark else b.WHITE)
    accent = accent or b.TEAL
    d.rect(s, x, y, w, h, fill, line=accent if not dark else b.TEAL, radius=0.04, shadow=True)
    d.rect(s, x, y, Inches(0.06), h, accent)
    d.text(s, title, x + Inches(0.20), y + Inches(0.16), w - Inches(0.40), Inches(0.26), size=10.8, color=b.WHITE if dark else b.NAVY, bold=True, shrink=True)
    d.text(s, body, x + Inches(0.20), y + Inches(0.56), w - Inches(0.40), h - Inches(0.68), size=8.0, color=hx("DCE6EF") if dark else b.INK, shrink=True)


def bullets(d, s, items, x, y, w, h, size=8.5, color=None):
    d.text(s, [{"text": item, "bullet": True, "space_before": 4, "size": size} for item in items], x, y, w, h, color=color or d.b.INK, shrink=True)


def arrow(d, s, x1, y1, x2, y2, color=None):
    color = color or d.b.TEAL
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.6)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass


def note_bar(d, s, label, text, y=6.28):
    b = d.b
    d.rect(s, Inches(0.72), Inches(y), Inches(11.88), Inches(0.50), b.NAVY, radius=0.04)
    d.text(s, label.upper(), Inches(0.94), Inches(y + 0.16), Inches(1.10), Inches(0.14), size=7.2, color=b.GOLD, bold=True, shrink=True)
    d.text(s, text, Inches(2.02), Inches(y + 0.13), Inches(10.1), Inches(0.20), size=8.5, color=b.WHITE, bold=True, shrink=True)


def lane(d, s, title, body, x, y, w=2.0, h=0.74, fill=None):
    b = d.b
    d.rect(s, x, y, Inches(w), Inches(h), fill or b.SOFT, line=b.GRID, radius=0.035, shadow=True)
    d.text(s, title, x + Inches(0.12), y + Inches(0.08), Inches(w - 0.24), Inches(0.24), size=7.2, color=b.NAVY, bold=True, shrink=True, align=PP_ALIGN.CENTER)
    d.text(s, body, x + Inches(0.12), y + Inches(0.34), Inches(w - 0.24), Inches(0.26), size=5.6, color=b.INK, shrink=True, align=PP_ALIGN.CENTER)


def cover(d, total):
    b = d.b
    s = d.slide(fill=b.NAVY)
    d.rect(s, Inches(8.85), 0, Inches(4.48), d.H, b.NAVY_2)
    d.rect(s, Inches(8.85), 0, Inches(0.08), d.H, b.TEAL)
    d.text(s, "EXECUTIVE BRIEFING", d.M, Inches(0.72), Inches(4.5), Inches(0.16), size=9, color=b.TEAL, bold=True)
    d.text(s, "How to use Fable 5 while the window is still open", d.M, Inches(1.14), Inches(7.6), Inches(1.06), size=30, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, d.M, Inches(2.40), Inches(1.30), Inches(0.05), b.TEAL)
    d.text(s, "A rebuilt deck with articulated narrative, demo-specific findings, and editable workflow diagrams.", d.M, Inches(2.74), Inches(7.2), Inches(0.36), size=11.2, color=hx("DCE6EF"), shrink=True)
    for i, text in enumerate(["Improve daily-driver skills", "Diagnose existing product loops", "Plan ambitious systems with full context"]):
        y = Inches(3.65 + i * 0.70)
        d.rect(s, d.M, y, Inches(0.34), Inches(0.34), b.TEAL, radius=0.035)
        d.text(s, str(i + 1), d.M, y + Inches(0.08), Inches(0.34), Inches(0.10), size=7.5, color=b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, text, d.M + Inches(0.52), y + Inches(0.05), Inches(5.5), Inches(0.18), size=10.6, color=b.WHITE, bold=True, shrink=True)
    lane(d, s, "Scarce access", "Do not spend it on throwaway prompts", Inches(9.65), Inches(1.15), w=2.65, h=0.82, fill=b.WHITE)
    lane(d, s, "High-context pass", "Hold intent, details, and system state", Inches(9.65), Inches(2.45), w=2.65, h=0.82, fill=b.WHITE)
    lane(d, s, "Durable artifact", "Skill patch, fix plan, or spec", Inches(9.65), Inches(3.75), w=2.65, h=0.82, fill=b.WHITE)
    arrow(d, s, Inches(10.98), Inches(2.02), Inches(10.98), Inches(2.38), b.GOLD)
    arrow(d, s, Inches(10.98), Inches(3.32), Inches(10.98), Inches(3.68), b.GOLD)
    footer(d, s, 1, total, dark=True)


def exec_summary(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Executive Summary", "Spend the window on assets that compound after access closes", "The opening argument is clear: avoid random one-shots and invest Fable into reusable work.")
    d.rect(s, d.M, Inches(1.70), d.CW, Inches(0.86), b.NAVY, radius=0.035)
    d.text(s, "BLUF", d.M + Inches(0.22), Inches(1.98), Inches(0.62), Inches(0.18), size=7.0, color=b.GOLD, bold=True)
    d.text(s, "Fable 5 preserves intent across many details, then turns that understanding into a reusable operating artifact.", d.M + Inches(0.96), Inches(1.86), d.CW - Inches(1.15), Inches(0.34), size=8.7, color=b.WHITE, bold=True, shrink=True)
    data = [
        ("Skills", "Use Fable to improve the instructions that run daily work: GTM strategy, positioning, pricing, landing pages, and design handoffs."),
        ("Existing systems", "Use Fable to find root causes when repeated patches do not explain why an agent loop or product feature keeps failing."),
        ("New builds", "Use Fable to merge notes, code, research, and product intent into a plan for complex features such as streaming UX."),
    ]
    for i, (t, body) in enumerate(data):
        card(d, s, t, body, Inches(0.80 + i * 4.08), Inches(2.95), Inches(3.55), Inches(1.72), [b.TEAL, b.GOLD, b.CORAL][i])
    note_bar(d, s, "Decision", "The question is not whether Fable is impressive; it is which work deserves the expensive context window.", y=5.55)
    footer(d, s, page, total)


def allocation(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The opening frames an allocation choice", "The speaker contrasts throwaway experiments with three concrete uses that produce durable value.")
    lane(d, s, "Option A", "One-shot designs and throwaway projects", Inches(1.1), Inches(2.35), w=3.4, h=1.00)
    lane(d, s, "Option B", "Three reusable workflows that keep paying back", Inches(8.8), Inches(2.35), w=3.4, h=1.00, fill=b.LIGHT_TEAL)
    d.text(s, "vs.", Inches(6.16), Inches(2.63), Inches(0.78), Inches(0.18), size=15, color=b.NAVY, bold=True, align=PP_ALIGN.CENTER)
    card(d, s, "Why Option B wins", "The access window is short and the model becomes expensive afterward. That makes durable outputs more valuable than isolated demonstrations.", Inches(2.0), Inches(4.20), Inches(9.35), Inches(1.12), b.TEAL)
    note_bar(d, s, "Implication", "Prioritize work where Fable's answer can be saved, reused, implemented by another model, or encoded into a skill.", y=6.15)
    footer(d, s, page, total)


def thesis(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Fable's edge is seeing the goal through the details", "The argument emphasizes context, intent, edge cases, and carrying critical details into implementation.")
    steps = [
        ("Context", "Full skill, codebase, notes, and user goal"),
        ("Intent", "What the builder is actually trying to accomplish"),
        ("Detail", "Small seams where important behavior slips"),
        ("Artifact", "A plan detailed enough for downstream execution"),
    ]
    for i, (t, body) in enumerate(steps):
        x = Inches(0.85 + i * 3.08)
        lane(d, s, t, body, x, Inches(2.30), w=2.25, h=1.06, fill=b.LIGHT_TEAL if i == 2 else b.SOFT)
        if i < 3:
            arrow(d, s, x + Inches(2.28), Inches(2.83), x + Inches(2.90), Inches(2.83))
    card(d, s, "The practical standard", "The model pass is successful only if it changes the plan: clearer skill handoffs, a real root cause, or a build plan that avoids known failure modes.", Inches(1.30), Inches(4.35), Inches(10.65), Inches(1.18), b.GOLD)
    footer(d, s, page, total)


def skill_context(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Use case 1: improve skills because they shape daily execution", "The demo starts with a go-to-market skill library used for strategy, positioning, pricing, and landing pages.")
    card(d, s, "What the library does", "It helps users choose channels, position against alternatives, set pricing, structure landing pages, and write copy.", Inches(0.82), Inches(1.95), Inches(3.72), Inches(1.72), b.TEAL)
    card(d, s, "Why Fable is relevant", "The skill needs to translate real-world marketing rules into enforceable downstream behavior, not just generate plausible text.", Inches(4.85), Inches(1.95), Inches(3.72), Inches(1.72), b.GOLD)
    card(d, s, "What should improve", "The output should become easier to hand to design and code tools without losing brand voice, customer language, or page intent.", Inches(8.88), Inches(1.95), Inches(3.72), Inches(1.72), b.CORAL)
    note_bar(d, s, "Readout", "The speaker is not using Fable to write one landing page; he is using it to improve the system that creates many landing pages.", y=5.35)
    footer(d, s, page, total)


def skill_findings(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The skill review found two practical failure modes", "The process was strong, but important details were not surviving the journey into downstream outputs.")
    card(d, s, "Failure mode 1: detail loss", "Inputs such as brand voice, customer language, and design specifics were created upstream but were not consistently carried into later phases.", Inches(0.92), Inches(2.05), Inches(5.45), Inches(1.65), b.TEAL)
    card(d, s, "Failure mode 2: missing adapters", "Different design tools require different prompt formats. A general plan needs a handoff layer that converts it into the right shape for each tool.", Inches(6.95), Inches(2.05), Inches(5.45), Inches(1.65), b.GOLD)
    lane(d, s, "Raw plan", "Good strategic content", Inches(1.35), Inches(4.35), w=1.8)
    lane(d, s, "Adapter", "Tool-specific handoff", Inches(4.15), Inches(4.35), w=1.8, fill=b.LIGHT_TEAL)
    lane(d, s, "Design tool", "Receives the right format", Inches(6.95), Inches(4.35), w=1.8)
    lane(d, s, "Memorable output", "Higher wow factor", Inches(9.75), Inches(4.35), w=1.8)
    for x in [3.18, 5.98, 8.78]:
        arrow(d, s, Inches(x), Inches(4.72), Inches(x + 0.78), Inches(4.72))
    footer(d, s, page, total)


def skill_operating_model(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The skill workflow should end in a codified operating change", "Fable can orchestrate the review; implementation can be delegated once the intent and details are captured.")
    items = [
        ("Review high-use skill", "Choose the instructions that affect daily work."),
        ("Answer grounding questions", "Clarify what should improve and what good output means."),
        ("Extract failure modes", "Look for lost detail, weak handoffs, and edge-case drift."),
        ("Patch the system", "Update skill rules, adapters, and verification gates."),
    ]
    for i, (t, body) in enumerate(items):
        x = Inches(0.75 + i * 3.13)
        lane(d, s, t, body, x, Inches(2.25), w=2.35, h=1.10, fill=b.LIGHT_TEAL if i == 2 else b.SOFT)
        if i < 3:
            arrow(d, s, x + Inches(2.38), Inches(2.80), x + Inches(2.95), Inches(2.80))
    note_bar(d, s, "Management point", "Skills are product assets. Reviewing them is not housekeeping; it is improving the engine that produces future work.", y=5.35)
    footer(d, s, page, total)


def product_context(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Use case 2: diagnose existing features before launch", "The product demo is a mobile app whose agent loop became unreliable after more autonomous behavior was added.")
    card(d, s, "Situation", "The first app version made a few model calls. Bugs began after the builder added a custom agent loop and richer behavior.", Inches(0.82), Inches(2.0), Inches(3.7), Inches(1.55), b.TEAL)
    card(d, s, "Problem", "Repeated fixes created many commits but did not solve the core issue. The system kept patching symptoms.", Inches(4.85), Inches(2.0), Inches(3.7), Inches(1.55), b.GOLD)
    card(d, s, "Business need", "The app was nearing market, so the feature needed to become robust rather than merely patched.", Inches(8.88), Inches(2.0), Inches(3.7), Inches(1.55), b.CORAL)
    note_bar(d, s, "Implication", "Use Fable when the team is too close to the bug and needs a higher-level system diagnosis.", y=5.45)
    footer(d, s, page, total)


def product_diagnosis(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The root cause was state living in the wrong place", "Fable identified a mismatch: the server-side loop was stateless, while the conversation needed continuity.")
    lane(d, s, "Client side", "Rendering plus too much orchestration", Inches(1.00), Inches(2.25), w=2.7, h=1.0)
    lane(d, s, "Server loop", "Stateless execution", Inches(5.30), Inches(2.25), w=2.7, h=1.0, fill=b.LIGHT_TEAL)
    lane(d, s, "Conversation", "Needs memory across turns", Inches(9.60), Inches(2.25), w=2.7, h=1.0)
    arrow(d, s, Inches(3.75), Inches(2.75), Inches(5.20), Inches(2.75))
    arrow(d, s, Inches(8.05), Inches(2.75), Inches(9.50), Inches(2.75))
    card(d, s, "Needle in the haystack", "A server-side checkpointer should persist state so interrupts, clarifications, and user replies continue the same execution with full context.", Inches(1.35), Inches(4.15), Inches(10.62), Inches(1.15), b.CORAL)
    footer(d, s, page, total)


def product_vs_patching(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Fable changed the diagnosis from patching bugs to redesigning the loop", "The demo contrasts duct-tape fixes with identifying the underlying system pattern.")
    d.rect(s, Inches(0.95), Inches(1.95), Inches(5.3), Inches(3.55), b.SOFT, line=b.GRID, radius=0.04, shadow=True)
    d.rect(s, Inches(7.05), Inches(1.95), Inches(5.3), Inches(3.55), b.LIGHT_TEAL, line=b.GRID, radius=0.04, shadow=True)
    d.text(s, "Symptom patching", Inches(1.25), Inches(2.25), Inches(4.7), Inches(0.25), size=13, color=b.NAVY, bold=True, shrink=True)
    d.text(s, "System diagnosis", Inches(7.35), Inches(2.25), Inches(4.7), Inches(0.25), size=13, color=b.NAVY, bold=True, shrink=True)
    bullets(d, s, ["Fixes each visible failure", "Creates more commits", "Misses state ownership", "Bug returns in a new form"], Inches(1.25), Inches(2.90), Inches(4.55), Inches(1.35), size=8.6)
    bullets(d, s, ["Explains why failures repeat", "Moves state to the right layer", "Defines continuity and checkpointing", "Creates a spec for downstream implementation"], Inches(7.35), Inches(2.90), Inches(4.55), Inches(1.35), size=8.6)
    note_bar(d, s, "Decision", "When repeated patches do not converge, stop debugging locally and ask for a system-level explanation.", y=6.05)
    footer(d, s, page, total)


def spec_handoff(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The output should become a spec another model can implement", "The speaker explicitly uses Fable to create proposals, while a less expensive model can execute the plan.")
    steps = [
        ("Fable pass", "Diagnose root cause and preserve intent"),
        ("Spec proposal", "Capture fixes, rationale, and critical details"),
        ("Implementation", "Delegate execution to a cheaper model or subagent"),
        ("Verification", "Confirm continuity, tool use, and user correction flows"),
    ]
    for i, (t, body) in enumerate(steps):
        x = Inches(0.86 + i * 3.08)
        lane(d, s, t, body, x, Inches(2.28), w=2.35, h=1.08, fill=b.LIGHT_TEAL if i == 1 else b.SOFT)
        if i < 3:
            arrow(d, s, x + Inches(2.40), Inches(2.82), x + Inches(2.92), Inches(2.82))
    card(d, s, "Why this matters", "The high-cost model does not need to do every implementation step. Its job is to produce a plan that carries forward intent and all critical details.", Inches(1.35), Inches(4.35), Inches(10.6), Inches(1.12), b.GOLD)
    footer(d, s, page, total)


def planning_context(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Use case 3: plan new systems that require merged context", "The planning example combines background knowledge, goals, notes, feature ideas, code, and fresh research.")
    inputs = [("Background", "What the builder knows"), ("Goal", "What the feature must accomplish"), ("Notes", "Past thinking and constraints"), ("Code", "Current system reality"), ("Research", "Modern implementation options")]
    for i, (t, body) in enumerate(inputs):
        x = Inches(0.70 + i * 2.50)
        lane(d, s, t, body, x, Inches(2.18), w=1.92, h=0.95, fill=b.SOFT)
        if i < 4:
            arrow(d, s, x + Inches(1.94), Inches(2.66), x + Inches(2.32), Inches(2.66))
    card(d, s, "Fable's planning role", "Marry the builder's intent with external research and the current codebase, then return a plan that can survive implementation handoff.", Inches(1.35), Inches(4.20), Inches(10.62), Inches(1.20), b.TEAL)
    footer(d, s, page, total)


def streaming(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The concrete planning problem is streaming backend work into the UI", "The app needs both token streaming and event streaming so the interface reflects what the agent is doing.")
    lane(d, s, "Server work", "Tools, milestones, state changes", Inches(0.95), Inches(2.20), w=2.25, h=1.0)
    lane(d, s, "Event stream", "What is happening now", Inches(3.95), Inches(2.20), w=2.25, h=1.0, fill=b.LIGHT_TEAL)
    lane(d, s, "Token stream", "Response appears progressively", Inches(6.95), Inches(2.20), w=2.25, h=1.0, fill=b.LIGHT_TEAL)
    lane(d, s, "Front end", "Render progress and result", Inches(9.95), Inches(2.20), w=2.25, h=1.0)
    for x in [3.22, 6.22, 9.22]:
        arrow(d, s, Inches(x), Inches(2.70), Inches(x + 0.66), Inches(2.70))
    bullets(d, s, ["Tokens should type out rather than appear as one block.", "Events should reveal tool calls and important markers.", "The UI and backend capability must be designed together."], Inches(1.40), Inches(4.35), Inches(10.6), Inches(0.90), size=9.2)
    footer(d, s, page, total)


def comparison(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "The model comparison exposes the kinds of mistakes Fable prevented", "The speaker reran a similar planning prompt elsewhere, then used Fable to identify what the other plan got wrong.")
    risks = [
        ("Reintroduced bugs", "A plan would recreate failure modes already fixed."),
        ("Broken interruption flow", "Pause and user response behavior was backwards."),
        ("Wrong platform facts", "Some assumptions about Supabase capability were incorrect."),
        ("Missed code areas", "Relevant parts of the codebase were not explored."),
    ]
    for i, (t, body) in enumerate(risks):
        x = Inches(0.88 + (i % 2) * 6.02)
        y = Inches(2.02 + (i // 2) * 1.55)
        card(d, s, t, body, x, y, Inches(5.2), Inches(1.05), [b.TEAL, b.GOLD, b.CORAL, b.DARK_TEAL][i])
    note_bar(d, s, "Implication", "Use Fable on plans where a plausible answer can still be dangerously wrong.", y=5.65)
    footer(d, s, page, total)


def selection_matrix(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Use Fable when the task has high context, high consequence, or high ambiguity", "The examples create a practical filter for choosing which problems deserve the expensive pass.")
    rows = [
        ("Daily-driver skill", "High reuse", "High leverage", "Run now"),
        ("Existing product loop", "Launch risk", "Root cause unclear", "Run now"),
        ("Complex new feature", "Many moving parts", "Research plus code", "Run now"),
        ("Simple one-off task", "Low reuse", "Narrow scope", "Use cheaper model"),
    ]
    x, y = Inches(0.95), Inches(1.95)
    widths = [3.2, 2.55, 3.0, 2.3]
    headers = ["Problem type", "Why it matters", "Fable fit", "Decision"]
    for j, h in enumerate(headers):
        left = x + Inches(sum(widths[:j]))
        d.rect(s, left, y, Inches(widths[j]), Inches(0.42), b.NAVY, line=b.NAVY)
        d.text(s, h, left + Inches(0.10), y + Inches(0.11), Inches(widths[j] - 0.2), Inches(0.10), size=7.6, color=b.WHITE, bold=True, shrink=True)
    for i, row in enumerate(rows):
        yy = y + Inches(0.42 + i * 0.58)
        for j, val in enumerate(row):
            left = x + Inches(sum(widths[:j]))
            fill = b.LIGHT_TEAL if j == 3 and "Run" in val else b.WHITE
            d.rect(s, left, yy, Inches(widths[j]), Inches(0.58), fill, line=b.GRID)
            d.text(s, val, left + Inches(0.10), yy + Inches(0.16), Inches(widths[j] - 0.20), Inches(0.16), size=7.8, color=b.NAVY if j in [0, 3] else b.INK, bold=j in [0, 3], shrink=True)
    note_bar(d, s, "Rule", "The more the answer depends on preserving intent across many details, the stronger the Fable fit.", y=5.70)
    footer(d, s, page, total)


def action_plan(d, page, total):
    b = d.b
    s = d.slide(fill=b.WHITE)
    header(d, s, "Convert the remaining window into an asset sprint", "The closing recommendation is to choose one or all three high-value use cases before access becomes expensive.")
    days = [
        ("1", "Pick targets", "One skill, one existing feature, one ambitious system."),
        ("2", "Run skill review", "Patch detail flow and tool adapters."),
        ("3", "Run feature diagnosis", "Find root cause and spec the fix."),
        ("4", "Run complex plan", "Merge notes, code, research, and UX goal."),
        ("5", "Codify outputs", "Save specs, patches, checklists, and next actions."),
    ]
    for i, (num, title, body) in enumerate(days):
        x = Inches(0.70 + i * 2.50)
        d.rect(s, x, Inches(2.10), Inches(2.05), Inches(2.25), b.WHITE, line=b.GRID, radius=0.04, shadow=True)
        d.rect(s, x, Inches(2.10), Inches(2.05), Inches(0.36), b.NAVY)
        d.text(s, f"Day {num}", x + Inches(0.16), Inches(2.20), Inches(1.70), Inches(0.10), size=7.0, color=b.GOLD, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, title, x + Inches(0.14), Inches(2.72), Inches(1.74), Inches(0.25), size=9.4, color=b.NAVY, bold=True, shrink=True, align=PP_ALIGN.CENTER)
        d.text(s, body, x + Inches(0.14), Inches(3.23), Inches(1.74), Inches(0.48), size=6.6, color=b.INK, shrink=True, align=PP_ALIGN.CENTER)
    note_bar(d, s, "Success", "The week ends with changed systems, not just an interesting chat.", y=5.55)
    footer(d, s, page, total)


def conclusion(d, page, total):
    b = d.b
    s = d.slide(fill=b.NAVY)
    d.text(s, "Conclusion", d.M, Inches(0.72), Inches(3.0), Inches(0.16), size=9, color=b.TEAL, bold=True)
    d.text(s, "Use Fable where losing the thread would change the outcome", d.M, Inches(1.14), Inches(8.1), Inches(0.95), size=27, color=b.WHITE, bold=True, shrink=True)
    d.rect(s, d.M, Inches(2.32), Inches(1.3), Inches(0.05), b.TEAL)
    d.text(s, "The argument is strongest when read as an operating memo: Fable is the high-context planner for the few tasks where intent, state, handoffs, and edge cases must stay connected.", d.M, Inches(2.70), Inches(7.3), Inches(0.72), size=11.2, color=hx("DCE6EF"), shrink=True)
    for i, t in enumerate(["Improve the system that creates outputs", "Find the root cause behind repeated failures", "Plan the build that is too interconnected to chunk casually"]):
        y = Inches(4.15 + i * 0.58)
        d.rect(s, d.M, y, Inches(0.30), Inches(0.30), b.TEAL, radius=0.03)
        d.text(s, str(i + 1), d.M, y + Inches(0.07), Inches(0.30), Inches(0.09), size=6.8, color=b.NAVY, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, t, d.M + Inches(0.48), y + Inches(0.03), Inches(6.4), Inches(0.20), size=9.4, color=b.WHITE, bold=True, shrink=True)
    card(d, s, "Recommended next move", "Pick the highest-consequence asset and turn Fable's answer into a committed change today.", Inches(9.0), Inches(2.02), Inches(3.08), Inches(2.08), b.TEAL, fill=b.NAVY_2, dark=True)
    footer(d, s, page, total, dark=True)


def scan_visible_text(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    banned = ["transcript", "hyperframe", "excalidraw", "youtube", "source", "validation", "synthesis", "audit", "codex", "claude", "/home/", "runs/"]
    hits = [term for term in banned if term in text.lower()]
    time_hits = re.findall(r"\b\d{1,2}:\d{2}\b", text)
    if hits or time_hits:
        raise RuntimeError(f"Visible text scan failed: terms={hits}, time_hits={time_hits[:5]}")


def main():
    d = Deck()
    total = 16
    slides = [
        cover, exec_summary, allocation, thesis,
        skill_context, skill_findings, skill_operating_model,
        product_context, product_diagnosis, product_vs_patching, spec_handoff,
        planning_context, streaming, comparison, selection_matrix, action_plan, conclusion,
    ]
    total = len(slides)
    for idx, fn in enumerate(slides, 1):
        if fn is cover:
            fn(d, total)
        else:
            fn(d, idx, total)
    d.save(OUT)
    scan_visible_text(OUT)
    print(f"Visible text scan passed: {OUT}")


if __name__ == "__main__":
    main()
