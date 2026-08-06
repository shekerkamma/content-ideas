#!/usr/bin/env python3
from __future__ import annotations

import json, sys
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[3]
RUN = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT / "skills/branded-pptx-deck/scripts"))
from pptxkit import Brand, Deck

TEMPLATE = Path("/mnt/c/Users/sheke/OneDrive/Desktop/DeepGrid-Semi-Investor-Briefing-Client-Ready-v8-reviewed.pptx")
OUT = RUN / "Client-Acquisition-Function-Operating-Briefing-draft.pptx"

B = Brand(
    NAVY=RGBColor(0x0A,0x16,0x28), NAVY_2=RGBColor(0x23,0x36,0x4C),
    TEAL=RGBColor(0x00,0xB4,0xD8), ACCENT=RGBColor(0x0A,0x93,0x96),
    DARK_TEAL=RGBColor(0x12,0x85,0x5B), LIGHT_TEAL=RGBColor(0xE8,0xF8,0xFC),
    GOLD=RGBColor(0xF2,0xA8,0x3B), AMBER=RGBColor(0xB4,0x53,0x09),
    CORAL=RGBColor(0xD9,0x4A,0x5A), WHITE=RGBColor(0xFF,0xFF,0xFF),
    SOFT=RGBColor(0xF8,0xFA,0xFC), INK=RGBColor(0x1E,0x29,0x3B),
    MUTED=RGBColor(0x5B,0x6B,0x7C), GRID=RGBColor(0xE2,0xE8,0xF0),
    FONT="Aptos", FONT_H="Aptos Display")

# kind, section, action title, subtitle, supporting points
SLIDES = [
("cover","CLIENT ACQUISITION FUNCTION","Build a Revenue System That Learns Before It Automates","From inconsistent prompts to a governed, always-on operating function",["58-slide operator briefing","Native, editable PowerPoint","Human-controlled external action"]),
("summary","EXECUTIVE DECISION","The next acquisition advantage is operational reliability—not more outreach volume","Approve a 90-day staged build that earns automation one gate at a time",["Start with one recurring calendar task","Compound context and review evidence","Scale only after watched runs pass"]),
("story","THE ARGUMENT","Five moves turn scattered prompts into a dependable acquisition function","The whole deck in one decision sequence",["Diagnose the current level","Stabilize a single task","Build the shared business brain","Separate specialist roles","Govern, measure, and compound"]),
("divider","01 · WHY CHANGE","Most acquisition automation scales inconsistency","Volume multiplies weak targeting, unsupported personalization, and hidden judgment",[]),
("cards","THE FAILURE MODE","Prompt libraries create activity without creating operating leverage","A prompt is useful—but it is not a process, control system, or learning loop",["Context resets: each run rediscovers the offer and buyer","Quality drifts: ‘good’ is not defined or measured","Handoffs vanish: decisions remain trapped in chat","Automation risk: sending arrives before reliability"]),
("compare","OPERATING REALITY","The bottleneck is not copy generation; it is coordinated judgment","Every acquisition stage depends on decisions that must become explicit",["Targeting | Is this account worth attention?","Relevance | What changed now?","Message | Which pain and proof fit?","Progression | What is the safest next move?","Learning | Which outcomes change the rules?"]),
("cards","ECONOMICS","Bad automation creates negative operating leverage","Cheap generation can make expensive mistakes arrive faster",["False positives consume seller time","Generic outreach damages scarce reputation","Unsupported claims create commercial risk","Unlogged exceptions prevent learning"]),
("flow","SYSTEM THESIS","A reliable function compounds five assets across every cycle","Prompts are only the first asset",["Evidence","Context","Artifacts","Decisions","Learning"]),
("score","MATURITY DIAGNOSIS","‘Always on’ is Level 9; most teams begin with unverified Level-1 fragments","Name the current level from evidence—not aspiration",["0 | Manual / undocumented","1 | Saved prompt","2 | Chain","3 | Business brain","4–5 | Skill + deliverable","6–8 | Connection + agents","9–10 | Always-on function"]),
("divider","02 · THE LADDER","Ten levels prevent premature autonomy","Each level removes one calendar task and earns the right to build the next",[]),
("gate","LEVEL 1 · SAVED PROMPT","Stabilize one narrow task before building a workflow","Gate: at least 8 of 10 representative cases need no substantive edit",["Fixed inputs and output","Explicit acceptance rubric","Fresh-session testing","Failure cases retained"]),
("gate","LEVEL 2 · CHAIN","Make intermediate reasoning inspectable and restartable","Gate: ordered stages, named artifacts, and a human checkpoint",["Research brief","Fit decision","Message angle","Draft","Review decision"]),
("gate","LEVEL 3 · BUSINESS BRAIN","Consistency comes from shared context—not longer instructions","Gate: six current files produce consistent work in a fresh session",["Offer","Ideal customer","Voice","Proof","Objections","Process"]),
("gate","LEVEL 4 · SKILL","Package the repeatable method as an invocable capability","Gate: explicit and implicit invocation both satisfy the output contract",["Trigger conditions","Inputs and workflow","Output contract","Guardrails","Validation"]),
("gate","LEVEL 5 · DELIVERABLE MACHINE","Convert analysis into a finished artifact—not another chat response","Gate: structured output in under 10 minutes with format and content checks",["Account brief","Personalized sequence","Call preparation","Proposal draft","Weekly pipeline review"]),
("gate","LEVEL 6 · CONNECTION","Observe systems before writing to them","Gate: one read-only week, then one bounded draft-only write with rollback",["CRM read","Calendar read","Research sources","Draft queue","Action log"]),
("gate","LEVEL 7 · AGENT","Give one bounded role ownership of one outcome","Gate: junior-team-member quality with explicit stop conditions",["One mission","Fixed inputs","Reviewable output","No sending","Clear escalation"]),
("gate","LEVEL 8 · TEAM","Separate discovery, persuasion, progression, and control","Gate: fixed file handoffs and an enforceable Auditor veto",["Scout","Writer","Closer","Auditor","Human approver"]),
("gate","LEVEL 9 · ALWAYS ON","Scheduling is earned through watched reliability","Gate: five watched runs and two stable weeks",["Idempotency","Retry limits","Alerts","Volume bounds","Zero unapproved sends"]),
("gate","LEVEL 10 · FUNCTION","A function coordinates states, owners, economics, and learning","Gate: governed state machine plus a 20% traffic pilot for two weeks",["Single intake","Stage ownership","Approval policy","Exception handling","Weekly feedback loop"]),
("divider","03 · THE BUSINESS BRAIN","Six files make every acquisition decision context-aware","This is the durable commercial memory behind prompts, agents, and channels",[]),
("cards","OFFER","The offer file defines the commercial promise and its boundaries","Prevent attractive copy from outrunning delivery reality",["Buyer outcome and urgency","Mechanism and differentiation","Scope and exclusions","Pricing constraints","Proof required before claims"]),
("cards","IDEAL CUSTOMER","The ICP file identifies situations—not just firmographics","Buying triggers and disqualifiers sharpen attention",["Company and role fit","Trigger events","Cost of inaction","Existing substitutes","Hard disqualifiers"]),
("cards","VOICE","Voice is a decision system, not a list of adjectives","Examples and banned patterns make tone testable",["Vocabulary to use","Phrases to avoid","Sentence rhythm","Good and bad examples","Channel-specific adjustments"]),
("cards","PROOF","Every commercial claim needs an evidence object","Confidence rises when proof includes limitations",["Claim","Evidence","Source and date","Applicability","Known limitations"]),
("cards","OBJECTIONS","Objections are diagnostic signals—not rebuttal prompts","Separate the stated concern from the underlying decision risk",["Observed objection","Likely diagnosis","Response options","Supporting proof","Escalation condition"]),
("cards","PROCESS","The process file defines the operating system","States and exit criteria replace heroic memory",["Lifecycle states","Owners","Entry and exit criteria","Tools and artifacts","Approvals and metrics"]),
("flow","CONTEXT UPDATE LOOP","Commercial memory must change when reality changes","Every accepted edit becomes a candidate rule update",["Outcome observed","Reviewer explains edit","Pattern classified","Context rule proposed","Owner approves update"]),
("divider","04 · LIFECYCLE DESIGN","The system advances opportunities through explicit states","A state exists only when its entry evidence and exit decision are recorded",[]),
("flow","ACQUISITION STATE MACHINE","One intake coordinates research, drafting, approval, and progression","External action remains a human-controlled transition",["Candidate","Verified","Scored","Drafted","Approved","Sent","Engaged","Qualified","Advanced","Closed / Nurture"]),
("cards","INTAKE","A stable subject identity prevents duplicated and misdirected work","The intake record is the anchor for every downstream artifact",["Stable company and person IDs","Source URLs and retrieval date","Channel permissions","Existing relationship check","Assigned owner"]),
("score","FIT SCORING","Score the situation, not the logo","A good account without a buying trigger is not a priority",["Strategic fit | 25","Trigger strength | 25","Pain intensity | 20","Reachability | 15","Proof relevance | 15","Disqualifiers | hard stop"]),
("compare","TRIGGER LIBRARY","Events create timing; themes create relevance","Use observable change to justify attention",["Leadership change | new mandate","Funding / expansion | new capacity","Hiring pattern | capability gap","Regulation | non-discretionary need","Technology change | replacement window"]),
("cards","RESEARCH BRIEF","Research must answer a commercial decision","Collection without a decision question becomes trivia",["Why this account?","Why this buyer?","Why now?","Which proof transfers?","What would disconfirm fit?"]),
("cards","MESSAGE ARCHITECTURE","Personalization earns attention by compressing relevance","The first message should prove understanding, not explain the whole offer",["Observed change","Business implication","Relevant point of view","Credible proof","Low-friction next step"]),
("compare","SEQUENCE DESIGN","A sequence should increase information—not repeat the same ask","Each touch has a different job",["1 | Establish relevance","2 | Add a useful implication","3 | Offer proof or diagnostic","4 | Reframe the cost of delay","5 | Close the loop respectfully"]),
("cards","QUALIFICATION","Qualification protects delivery capacity and seller attention","Advance only when value, access, and timing are real",["Problem acknowledged","Economic consequence","Decision process visible","Access to required stakeholders","Mutual next step"]),
("flow","PROPOSAL PATH","A proposal is the output of accumulated decisions","Do not use a document to discover basic scope",["Qualified need","Solution shape","Commercial boundaries","Risk review","Human pricing approval","Client-ready proposal"]),
("divider","05 · THE FOUR-ROLE TEAM","Specialization raises quality and makes control visible","Each role owns one decision surface and communicates through files",[]),
("cards","SCOUT","Scout finds and verifies candidates—but never contacts them","Discovery quality is measured by accepted opportunities, not list size",["Search defined arenas","Verify identity and recency","Capture trigger evidence","Score fit and confidence","Reject weak candidates"]),
("cards","WRITER","Writer turns verified evidence into channel-ready drafts","Writer cannot invent proof, change the offer, or send",["Select message angle","Map evidence to pain","Draft channel variant","Label inference","Request missing context"]),
("cards","CLOSER","Closer prepares progression options without making commitments","Commercial judgment stays bounded by approved terms",["Summarize engagement","Diagnose next decision","Prepare meeting brief","Draft follow-up","Escalate pricing exceptions"]),
("cards","AUDITOR","Auditor protects evidence, policy, reputation, and learning","The veto is a feature—not an inconvenience",["Verify claims and sources","Check duplicate contact","Enforce voice and policy","Confirm approval status","Reject and explain defects"]),
("flow","FILE HANDOFFS","Agents cooperate through explicit artifacts—not hidden chat state","Every handoff is replayable and attributable",["Scored candidate","Evidence-backed draft","Audit decision","Human approval","Outcome record"]),
("compare","BOUNDARIES","Separation of duties prevents silent risk accumulation","No role can discover, persuade, approve, and send alone",["Scout cannot draft outreach","Writer cannot alter evidence","Closer cannot change commercial terms","Auditor cannot approve its own work","Automation cannot impersonate approval"]),
("divider","06 · CHANNEL PLAYBOOKS","One operating system can serve several channels","Channel mechanics vary; evidence, approval, and learning remain shared",[]),
("cards","LINKEDIN","LinkedIn works when relevance is visible before the pitch","Protect relationship capital with restrained volume",["Connection context first","One observed trigger","No fabricated familiarity","Draft-only automation","Human-controlled send"]),
("cards","EMAIL","Email earns attention through specificity and low cognitive load","Optimize for qualified replies, not opens",["Clear sender identity","Evidence-led subject","One idea per message","Plain next step","Suppression and opt-out rules"]),
("cards","PARTNERSHIPS","Partner acquisition needs mutual economics—not generic outreach","The unit of analysis is the joint value proposition",["Shared customer problem","Complementary capability","Economic exchange","Proof plan","Named joint motion"]),
("cards","REFERRALS","Referrals become systematic when the ask is concrete","Make it easy to identify and introduce the right person",["Trigger after delivered value","Describe the exact situation","Provide forwarding language","Protect the relationship","Close the feedback loop"]),
("cards","CONTENT-LED","Content should create diagnostic conversations","The acquisition system turns repeated buyer questions into useful assets",["Mine objections and calls","Publish decision frameworks","Link proof to buyer jobs","Capture intent signals","Route engagement to human follow-up"]),
("divider","07 · GOVERNANCE","The fastest safe system knows exactly when to stop","Control design is part of conversion quality—not a compliance appendix",[]),
("score","APPROVAL MATRIX","External consequence determines the approval threshold","Research and drafting can run; commitments cannot",["Internal research | automatic","Draft creation | automatic","CRM draft write | bounded","External send | human","Pricing / contract | authorized owner","Exception | named escalation"]),
("cards","STOP CONDITIONS","Low confidence should create a question—not a confident message","Stopping early protects brand and training data",["Ambiguous identity","Missing or stale evidence","Conflicting instructions","Policy or consent risk","Abnormal volume","Unsupported personalization"]),
("cards","FAILURE RECOVERY","Retries must not create duplicate commercial action","Every stage needs idempotency and rollback",["Stable artifact ID","Attempt counter","Duplicate suppression","Quarantine queue","Owner alert","Resume from failed stage"]),
("cards","OBSERVABILITY","A system is only ‘always on’ when its behavior is observable","Logs must explain what happened and why",["Input and source version","Decision and confidence","Artifact checksum","Reviewer change","External action ID","Outcome and latency"]),
("divider","08 · MEASUREMENT","Measure decision quality before message volume","The learning system improves when leading and lagging indicators connect",[]),
("score","MEASUREMENT TREE","Revenue is lagging; operating quality predicts whether it can scale","Track the complete path from attention to economics",["Coverage | qualified accounts","Quality | accepted drafts","Engagement | positive replies","Progression | qualified meetings","Economics | pipeline and revenue","Risk | rejects and exceptions"]),
("cards","QUALITY METRICS","Reviewer behavior is the fastest signal of system quality","Edits reveal missing context and weak rules",["No-substantive-edit rate","Evidence completeness","False-positive rate","Auditor rejection rate","Repeat-defect rate"]),
("cards","FUNNEL METRICS","Use conversion rates to locate the decision bottleneck","A low reply rate and a low qualification rate imply different fixes",["Verified → approved","Approved → sent","Sent → positive reply","Reply → qualified","Qualified → proposal","Proposal → win"]),
("cards","CAPACITY METRICS","Automation should release scarce judgment—not create review debt","Measure human effort per accepted outcome",["Research minutes per accepted account","Review minutes per approved draft","Touches per qualified meeting","Exception handling time","Seller hours per opportunity"]),
("cards","EXPERIMENT DESIGN","Change one variable at a time and retain the control","Learning requires comparability",["Declare hypothesis","Choose one segment","Freeze offer and proof","Define success threshold","Record disconfirming evidence"]),
("flow","WEEKLY LEARNING LOOP","Every week should leave the system smarter than it began","Accepted learning updates rules; anecdotes do not",["Review outcomes","Cluster edits and failures","Select one rule change","Run controlled test","Approve context update"]),
("divider","09 · 90-DAY BUILD","Move from prompt fragments to a governed pilot","The roadmap prioritizes reliability, context, and watched operation",[]),
("roadmap","DAYS 0–30","Prove one task and establish the commercial memory layer","Target: Levels 1–3",["Inventory the manual process","Build 10-case test set","Pass one prompt at 8/10","Create six context files","Run fresh-session consistency test"]),
("roadmap","DAYS 31–60","Package the method and produce finished deliverables","Target: Levels 4–5",["Define trigger and output contract","Build the acquisition skill","Generate account brief + message draft","Review every output","Measure time and edit rate"]),
("roadmap","DAYS 61–90","Connect safely and pilot bounded specialist roles","Target: Levels 6–8 only if gates pass",["Observe CRM read-only","Enable one draft-only write","Pilot Scout and Writer","Add Auditor veto","Keep all sending human-controlled"]),
("score","RESOURCE PLAN","A small accountable team beats an unowned automation program","Name owners before selecting tools",["Executive sponsor | priorities","Revenue owner | commercial rules","Operator | workflow and QA","Technical owner | connections and logs","Legal / security | policy gates"]),
("cards","RISK REGISTER","The main risks are commercial and operational—not model novelty","Design mitigations into the workflow",["Brand damage | approval + volume bounds","False claims | evidence objects + audit","Duplicate contact | stable IDs + suppression","Context drift | owners + review dates","Review overload | narrow scope + sampling"]),
("compare","BUILD VS BUY","Buy commodity execution; own commercial judgment and memory","Differentiation lives in context, decisions, and learning",["Buy | enrichment, CRM, scheduling, delivery","Configure | workflow, routing, permissions","Own | ICP, offer, proof, scoring, approval policy","Measure | edits, outcomes, exceptions"]),
("score","PROMOTION SCORECARD","A level advances only when proof is complete","Green status requires quality, control, and repeatability together",["Artifact complete","Acceptance test passed","Failure cases retained","Human owner named","Rollback verified","Learning captured"]),
("divider","10 · DECISION","Do not fund ‘autonomy’; fund the next proven capability","The roadmap can begin immediately without granting unsafe authority",[]),
("summary","THE 90-DAY ASK","Approve a staged pilot with one task, one owner, and one gate","Success is a reliable acquisition capability—not an impressive demo",["Choose the first calendar task","Name the commercial owner","Provide 10 historical cases","Freeze external sending","Review progress every week"]),
("cards","FIRST WORKSHOP","The first two hours should produce decisions and artifacts","Avoid starting with a tool tour",["Map the manual workflow","Select one narrow task","Define a good output","Identify approval points","Assemble the first 10 cases"]),
("close","CLOSE","Build the function in the order trust is earned","Research and draft automatically. Approve consequence deliberately. Compound every learning.",["NEXT MOVE","Select the first calendar task and schedule the Level-1 test","STATUS","Proposed operating model · Not yet an always-on function"]),
]

TOTAL=len(SLIDES)

def add_footer(d,s,n,dark=False): d.footer(s,n,TOTAL,dark=dark)

def card(d,s,x,y,w,h,head,body,accent=None):
    d.rect(s,x,y,w,h,B.WHITE,line=B.GRID,radius=.06)
    d.rect(s,x,y,Inches(.07),h,accent or B.TEAL)
    d.text(s,head,x+Inches(.2),y+Inches(.18),w-Inches(.35),Inches(.35),size=13,color=B.NAVY,bold=True,shrink=True)
    d.text(s,body,x+Inches(.2),y+Inches(.62),w-Inches(.35),h-Inches(.78),size=11.5,color=B.MUTED,shrink=True)

def build():
    d=Deck(B,footer="CLIENT ACQUISITION FUNCTION · OPERATING BRIEFING · AUGUST 2026")
    d.prs=Presentation(str(TEMPLATE)); d._wipe(); d.W=d.prs.slide_width; d.H=d.prs.slide_height; d.M=Inches(.6); d.CW=d.W-d.M*2
    for n,(kind,section,title,subtitle,items) in enumerate(SLIDES,1):
        dark=kind in {"cover","divider","close"}
        s=d.prs.slides.add_slide(d.prs.slide_layouts[0])
        d.rect(s,0,0,d.W,d.H,B.NAVY if dark else B.SOFT)
        if kind=="cover":
            d.rect(s,0,0,Inches(.18),d.H,B.TEAL); d.text(s,section,d.M,Inches(.7),Inches(6),Inches(.35),size=14,color=B.TEAL,bold=True)
            d.text(s,title,d.M,Inches(1.35),Inches(11.3),Inches(1.55),size=39,color=B.WHITE,bold=True,font=B.FONT_H,shrink=True)
            d.text(s,subtitle,d.M,Inches(3.15),Inches(10.7),Inches(.75),size=22,color=RGBColor(0x8D,0xB4,0xD4),shrink=True)
            for i,t in enumerate(items): card(d,s,Inches(.6+i*4.05),Inches(5.3),Inches(3.75),Inches(.85),str(i+1).zfill(2),t,B.TEAL)
        elif kind=="divider":
            d.text(s,section,d.M,Inches(1.0),Inches(6),Inches(.4),size=15,color=B.TEAL,bold=True)
            d.text(s,title,d.M,Inches(1.9),Inches(11.5),Inches(1.65),size=36,color=B.WHITE,bold=True,font=B.FONT_H,shrink=True)
            d.text(s,subtitle,d.M,Inches(3.75),Inches(10.8),Inches(1.0),size=20,color=RGBColor(0x8D,0xB4,0xD4),shrink=True)
            d.rect(s,d.M,Inches(5.35),Inches(3.0),Inches(.08),B.GOLD)
        elif kind=="close":
            d.text(s,section,d.M,Inches(.75),Inches(4),Inches(.35),size=14,color=B.TEAL,bold=True)
            d.text(s,title,d.M,Inches(1.5),Inches(11.8),Inches(1.15),size=34,color=B.WHITE,bold=True,shrink=True)
            d.text(s,subtitle,d.M,Inches(2.85),Inches(11.2),Inches(.8),size=19,color=RGBColor(0x8D,0xB4,0xD4),shrink=True)
            for i in range(0,len(items),2): card(d,s,Inches(.6+(i//2)*6.1),Inches(4.4),Inches(5.75),Inches(1.2),items[i],items[i+1],B.GOLD if i else B.TEAL)
        else:
            d.rect(s,0,0,d.W,Inches(.14),B.TEAL)
            d.text(s,title,d.M,Inches(.28),d.CW,Inches(.92),size=24,color=B.NAVY,bold=True,font=B.FONT_H,shrink=True)
            d.rect(s,d.M,Inches(1.22),Inches(1.45),Inches(.045),B.TEAL)
            d.text(s,subtitle,d.M,Inches(1.34),d.CW,Inches(.3),size=11,color=B.MUTED,shrink=True)
            if kind in {"flow","roadmap"}:
                count=len(items); cols=5 if count>6 else count; rows=(count+cols-1)//cols; gap=Inches(.15); w=(d.CW-gap*(cols-1))/cols; h=Inches(1.65) if rows>1 else Inches(2.35)
                for i,t in enumerate(items):
                    row=i//cols; col=i%cols; x=d.M+col*(w+gap); y=Inches(1.95)+row*(h+Inches(.25)); d.rect(s,x,y,w,h,B.NAVY_2 if i%2==0 else B.NAVY,radius=.06)
                    d.text(s,str(i+1).zfill(2),x+Inches(.14),y+Inches(.18),w-Inches(.28),Inches(.3),size=11.5,color=B.TEAL,bold=True)
                    d.text(s,t,x+Inches(.14),y+Inches(.65),w-Inches(.28),h-Inches(.8),size=12.5,color=B.WHITE,bold=True,shrink=True)
                    if col<cols-1 and i<count-1: d.text(s,"→",x+w-Inches(.05),y+Inches(.78),Inches(.25),Inches(.3),size=17,color=B.GOLD,bold=True)
            elif kind in {"score","compare"}:
                for i,t in enumerate(items):
                    y=Inches(1.9)+i*Inches(.68); d.rect(s,d.M,y,d.CW,Inches(.52),B.WHITE,line=B.GRID,radius=.03)
                    d.rect(s,d.M,y,Inches(.12),Inches(.52),B.TEAL if i%2==0 else B.ACCENT)
                    d.text(s,t,d.M+Inches(.28),y+Inches(.1),d.CW-Inches(.5),Inches(.3),size=13,color=B.INK,bold=(i==0),shrink=True)
            else:
                cols=2 if len(items)>3 else len(items); rows=(len(items)+cols-1)//cols; w=(d.CW-Inches(.3)*(cols-1))/cols; h=Inches(4.65)/rows-Inches(.18)
                for i,t in enumerate(items):
                    head,body=(t.split(" | ",1)+[""])[:2] if " | " in t else (f"{i+1:02d}",t)
                    x=d.M+(i%cols)*(w+Inches(.3)); y=Inches(1.85)+(i//cols)*(h+Inches(.28)); card(d,s,x,y,w,h,head,body,B.GOLD if kind=="gate" else B.TEAL)
        add_footer(d,s,n,dark)
    d.save(OUT)

    claims=[]; plans=[]; visuals=[]
    for i,(kind,section,title,subtitle,items) in enumerate(SLIDES,1):
        cid=f"claim-{i:02d}"; vid=f"visual-{i:02d}"
        claims.append({"id":cid,"text":title,"kind":"recommendation" if kind in {"summary","close","roadmap"} else "interpretation","confidence":"high","evidence_ids":[]})
        plans.append({"slide_id":i,"archetype":kind if kind in {"cover","comparison","roadmap","operating-model","executive-summary","section-divider","decision-scorecard","recommendation-ask"} else "three-card-argument","action_title":title,"audience_job":subtitle,"claim_ids":[cid],"evidence_ids":[],"visual_ids":[vid],"speaker_notes":{"talk_time_seconds":45,"narrative":subtitle,"demo_fallback":None},"accessibility":{"reading_order":["title","subtitle","content"],"visual_alt_text":f"Native editable {kind} layout supporting: {title}"},"status":"built"})
        visuals.append({"id":vid,"slide_ids":[i],"purpose":subtitle,"artifact_type":"diagram" if kind in {"flow","roadmap","score","compare"} else "component","evidence_status":"authored","route":"native","reason":"Structured business meaning is clearest and editable as native PowerPoint objects.","evidence_ids":[],"source":None,"asset":None,"editable_source":None,"prompt_file":None,"execution":None,"placement":None,"qa":{"visual_checked":False,"legible":False,"not_cropped":False,"background_matched":None,"provenance_recorded":True,"text_free":None,"not_evidence":None,"exact_fidelity":None}})
    (RUN/"slide-plan.json").write_text(json.dumps({"contract_version":"1.0","deck":{"name":"Building a Governed Client Acquisition Function","evidence_contract":"presentation-evidence.json","visual_contract":"visual-spec.json"},"claims":claims,"slides":plans},indent=2)+"\n")
    (RUN/"visual-spec.json").write_text(json.dumps({"contract_version":"1.0","deck":{"output_mode":"native","editability":"PPT-native editable diagrams","reference_deck":"slides.pptx"},"visuals":visuals},indent=2)+"\n")
    manifest={"skills":[{"name":"build-client-acquisition-function","application":"Content architecture, maturity gates, roles, controls, and compounding loop"},{"name":"present","application":"Governed presentation routing"},{"name":"presentation-source-bundle","application":"Reference-deck evidence contract"},{"name":"pptx-design-quality","application":"Design intent and QA gates"},{"name":"pptx-visual-spec","application":"Native visual routing"},{"name":"branded-pptx-deck","application":"Native editable PPTX construction"}]}
    (RUN/"skill-application-manifest.json").write_text(json.dumps(manifest,indent=2)+"\n")
    print(f"contracts: {TOTAL} slides")

if __name__=="__main__": build()
