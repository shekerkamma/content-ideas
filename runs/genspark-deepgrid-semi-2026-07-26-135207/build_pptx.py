#!/usr/bin/env python3
"""build_pptx.py — assemble rendered slide PNGs into a 16:9 image-per-slide PPTX.

Matches Genspark's export format (full-bleed 2560x1440 image per slide), but from
YOUR own rendered template. Run in WSL (python-pptx installed there).

  python3 build_pptx.py --png <png-dir> --out <deck.pptx>

The output is IMAGE-BASED (one picture per slide). Source-editability lives in the
deck HTML + theme.css, NOT in PowerPoint shapes. Declare this on delivery.
"""
import argparse, glob, os, sys
from pptx import Presentation
from pptx.util import Inches


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--png", default="build/png", help="dir of slide-*.png")
    ap.add_argument("--out", default="build/deck.pptx", help="output .pptx path")
    a = ap.parse_args()

    pngs = sorted(glob.glob(os.path.join(a.png, "slide-*.png")))
    if not pngs:
        sys.exit(f"No slide-*.png in {a.png} — run render.mjs first.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)

    os.makedirs(os.path.dirname(os.path.abspath(a.out)), exist_ok=True)
    prs.save(a.out)
    print(f"DONE: {len(pngs)} slides -> {a.out}")


if __name__ == "__main__":
    main()
