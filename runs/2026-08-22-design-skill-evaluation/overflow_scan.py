"""Replicate preview_pptx.py's overflow rule exactly and list every offender."""
import sys, textwrap
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

prs = Presentation(sys.argv[1])
bad = 0
for idx, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        total = 0.0
        for para in sh.text_frame.paragraphs:
            if not para.runs:
                continue
            text = "".join(r.text for r in para.runs)
            size = para.runs[0].font.size.pt if para.runs[0].font.size else 14
            max_chars = max(4, int(w / (size * 0.0085)))
            lines = textwrap.wrap(text, max_chars) or [text]
            sb = para.space_before.pt / 72.0 if para.space_before else 0
            total += sb + len(lines) * (size / 72.0) * 1.18
        if total > h + 0.03:
            bad += 1
            print(f"slide {idx:2d}  {sh.name:22s}  needs {total:.2f}in in {h:.2f}in  "
                  f":: {sh.text_frame.text.strip()[:52]!r}")
print(f"\n{bad} overflowing shape(s)")
