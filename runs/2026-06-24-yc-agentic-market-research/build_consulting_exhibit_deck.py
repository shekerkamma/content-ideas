#!/usr/bin/env python3
"""Build a McKinsey-inspired consulting exhibit deck from qualified Reddit evidence.

This intentionally uses consulting-style exhibit structure: action titles,
one main exhibit per slide, source notes, and short labels instead of prose
cards. It does not copy any proprietary template.
"""

from __future__ import annotations

import json
import os
import sys
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPTXKIT = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT))
from pptxkit import validate_pptx  # noqa: E402

RUN_DIR = Path(__file__).resolve().parent
TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE") or "/home/shekerk/.claude/templates/branded-template.pptx")
COMPANY_PACK = RUN_DIR / "yc-company-source-pack.json"
EVIDENCE_PACK = RUN_DIR / "reddit-evidence-pack.json"
OUT = RUN_DIR / "yc-agentic-market-research-consulting-exhibit-v3-draft.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x0A, 0x16, 0x28)
BLUE = RGBColor(0x1F, 0x5F, 0xB8)
TEAL = RGBColor(0x00, 0xA8, 0x8F)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
AMBER = RGBColor(0xD4, 0x8B, 0x00)
RED = RGBColor(0xC7, 0x3E, 0x3A)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5F, 0x6C, 0x7A)
GRID = RGBColor(0xD9, 0xDF, 0xE5)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)


def rect(slide, x, y, w, h, color, *, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, *, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    lines = value if isinstance(value, list) else [value]
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(item, dict):
            content = item["text"]
            size_i = item.get("size", size)
            color_i = item.get("color", color)
            bold_i = item.get("bold", bold)
            p.space_before = Pt(item.get("space_before", 0))
            if item.get("bullet"):
                content = "- " + content
        else:
            content = str(item)
            size_i, color_i, bold_i = size, color, bold
        r = p.add_run()
        r.text = content
        r.font.name = "Arial"
        r.font.size = Pt(size_i)
        r.font.bold = bold_i
        r.font.color.rgb = color_i
    return box


def slide_title(slide, title: str, kicker: str = "") -> None:
    text(slide, title, Inches(0.55), Inches(0.24), Inches(12.1), Inches(0.78), size=16.5, color=NAVY, bold=True)
    if kicker:
        text(slide, kicker, Inches(0.55), Inches(0.94), Inches(11.8), Inches(0.30), size=7.8, color=MUTED)
    rect(slide, Inches(0.55), Inches(1.22), Inches(1.15), Inches(0.04), TEAL)


def footer(slide, page: int, total: int, source: str = "Sources: YC local directory export; qualified Reddit evidence pack") -> None:
    text(slide, source, Inches(0.55), Inches(7.08), Inches(8.8), Inches(0.18), size=6.2, color=MUTED)
    text(slide, f"{page} / {total}", Inches(12.05), Inches(7.08), Inches(0.65), Inches(0.18), size=6.5, color=MUTED, align=PP_ALIGN.RIGHT)


def verdict_label(verdict: str) -> str:
    return {
        "qualified_reddit_support": "Qualified support",
        "weak_qualified_reddit_support": "Weak signal",
        "no_qualified_reddit_evidence": "Evidence gap",
    }.get(verdict, verdict.replace("_", " "))


def verdict_color(verdict: str):
    if verdict == "qualified_reddit_support":
        return GREEN
    if "weak" in verdict:
        return AMBER
    return RED


def truncate(value: str, n: int) -> str:
    value = " ".join(value.split())
    return value if len(value) <= n else value[: n - 1].rstrip() + "."


def cover(slide, rows: list[dict]) -> None:
    clear(slide)
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, Inches(0.0), Inches(0.0), Inches(0.18), H, TEAL)
    text(slide, "QUALIFIED REDDIT MARKET EVIDENCE", Inches(0.75), Inches(0.72), Inches(5.5), Inches(0.25), size=9.5, color=TEAL, bold=True)
    text(slide, "YC Agentic AI verticals: where Reddit supports the thesis", Inches(0.75), Inches(1.20), Inches(10.7), Inches(1.20), size=27, color=WHITE, bold=True)
    text(slide, "Consulting-exhibit draft | Not a raw scrape deck", Inches(0.75), Inches(2.72), Inches(6.0), Inches(0.30), size=12, color=WHITE)
    counts = Counter(r["suggested_reddit_verdict"] for r in rows)
    x = 0.75
    for value, label, color in [
        (counts["qualified_reddit_support"], "qualified support", GREEN),
        (counts["weak_qualified_reddit_support"], "weak signal", AMBER),
        (counts["no_qualified_reddit_evidence"], "evidence gap", RED),
    ]:
        rect(slide, Inches(x), Inches(4.65), Inches(2.25), Inches(0.78), color)
        text(slide, str(value), Inches(x + 0.15), Inches(4.80), Inches(0.55), Inches(0.28), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, label, Inches(x + 0.82), Inches(4.86), Inches(1.18), Inches(0.18), size=7.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += 2.55
    text(slide, f"Built {date.today().isoformat()} | Draft for review", Inches(0.75), Inches(6.55), Inches(4.5), Inches(0.22), size=8, color=TEAL, bold=True)
    footer(slide, 1, 9, source="Sources: YC local directory export; Reddit HTML extraction; qualification gate")


def executive_answer(slide, rows: list[dict]) -> None:
    clear(slide)
    slide_title(slide, "Reddit narrows the investable set to a small number of vertical workflow wedges", "Qualified Reddit evidence supports market pain in two wedges; the rest should be treated as hypotheses pending better sourcing")
    counts = Counter(r["suggested_reddit_verdict"] for r in rows)
    labels = [
        ("Qualified support", counts["qualified_reddit_support"], GREEN, "Use in narrative"),
        ("Weak signal", counts["weak_qualified_reddit_support"], AMBER, "Re-test"),
        ("Evidence gap", counts["no_qualified_reddit_evidence"], RED, "Do not claim"),
    ]
    x = 0.75
    for label, value, color, action in labels:
        rect(slide, Inches(x), Inches(1.62), Inches(3.55), Inches(1.08), WHITE, line=GRID)
        rect(slide, Inches(x), Inches(1.62), Inches(0.10), Inches(1.08), color)
        text(slide, str(value), Inches(x + 0.30), Inches(1.78), Inches(0.78), Inches(0.48), size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, label, Inches(x + 1.20), Inches(1.78), Inches(1.55), Inches(0.22), size=10.5, color=NAVY, bold=True)
        text(slide, action, Inches(x + 1.20), Inches(2.12), Inches(1.45), Inches(0.18), size=8.5, color=MUTED)
        x += 4.0
    text(slide, "Implication", Inches(0.75), Inches(3.25), Inches(1.2), Inches(0.24), size=12, color=NAVY, bold=True)
    bullets = [
        "Prioritize property management and construction review for deeper validation because Reddit contains qualified operator pain.",
        "Treat clinic communications and hotel operations as plausible but thin; collect more targeted operator threads before using them as proof.",
        "Avoid validating mortgage servicing, AR/accounting, insurance claims, and lending from this Reddit run; the current evidence pack does not support them.",
    ]
    text(slide, [{"text": b, "bullet": True, "space_before": 8, "size": 10.2} for b in bullets], Inches(0.95), Inches(3.64), Inches(10.8), Inches(1.70), size=10.2)
    footer(slide, 2, 9)


def evidence_funnel(slide, rows: list[dict]) -> None:
    clear(slide)
    total_rejected = sum(r.get("items_rejected_by_qualification_gate", 0) for r in rows)
    accepted = sum(r.get("evidence_count", 0) for r in rows)
    slide_title(slide, "The evidence gate prevents raw Reddit search results from becoming false validation", "Accepted items must match workflow, persona, and concrete pain signal")
    stages = [
        ("Raw Reddit items", f"{total_rejected + accepted:,}", "Posts/comments from discovery"),
        ("Rejected as noise", f"{total_rejected:,}", "Wrong subreddit, weak context, or generic keyword match"),
        ("Qualified evidence", f"{accepted}", "Workflow/persona/pain match retained"),
    ]
    x = 0.75
    for i, (label, value, note) in enumerate(stages):
        color = BLUE if i == 0 else RED if i == 1 else GREEN
        rect(slide, Inches(x), Inches(1.82), Inches(3.1), Inches(1.35), WHITE, line=color)
        text(slide, value, Inches(x + 0.25), Inches(2.02), Inches(2.55), Inches(0.35), size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, label, Inches(x + 0.25), Inches(2.48), Inches(2.55), Inches(0.22), size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(slide, note, Inches(x + 0.30), Inches(2.82), Inches(2.45), Inches(0.22), size=7.3, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            text(slide, ">", Inches(x + 3.25), Inches(2.30), Inches(0.35), Inches(0.30), size=20, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        x += 3.95
    checks = [
        ("Workflow", "same operating process"),
        ("Persona", "operator / practitioner context"),
        ("Pain", "manual work, delay, rework, workaround, objection"),
        ("Traceability", "subreddit, permalink, source JSON retained"),
    ]
    x = 0.95
    for label, note in checks:
        rect(slide, Inches(x), Inches(4.40), Inches(2.65), Inches(0.64), SOFT, line=GRID)
        text(slide, label, Inches(x + 0.15), Inches(4.53), Inches(0.85), Inches(0.18), size=8.2, color=NAVY, bold=True)
        text(slide, note, Inches(x + 1.08), Inches(4.49), Inches(1.34), Inches(0.26), size=5.8, color=MUTED)
        x += 2.90
    footer(slide, 3, 9)


def company_table(slide, companies: list[dict], rows_by_id: dict[str, dict], page: int, title: str) -> None:
    clear(slide)
    slide_title(slide, title, "YC facts are source fields; Reddit status is separate")
    headers = ["Company", "YC facts", "Workflow wedge", "Qualified Reddit readout", "Analyst implication"]
    xs = [0.62, 1.80, 3.15, 5.15, 7.38]
    ws = [1.05, 1.18, 1.78, 1.95, 4.85]
    for x, w, h in zip(xs, ws, headers):
        text(slide, h, Inches(x), Inches(1.34), Inches(w), Inches(0.20), size=6.7, color=MUTED, bold=True)
    y = 1.66
    for company in companies:
        row = rows_by_id[company["claim_id"]]
        color = verdict_color(row["suggested_reddit_verdict"])
        rect(slide, Inches(0.55), Inches(y - 0.06), Inches(12.05), Inches(0.72), WHITE, line=GRID)
        rect(slide, Inches(0.55), Inches(y - 0.06), Inches(0.07), Inches(0.72), color)
        text(slide, company["name"], Inches(xs[0]), Inches(y + 0.04), Inches(ws[0]), Inches(0.18), size=7.6, color=NAVY, bold=True)
        text(slide, f"{company['batch'].replace('Summer', 'S').replace('Winter', 'W').replace('Spring', 'Sp')} | team {company['team_size']}", Inches(xs[1]), Inches(y + 0.04), Inches(ws[1]), Inches(0.18), size=6.1, color=INK)
        text(slide, truncate(company["wedge"], 42), Inches(xs[2]), Inches(y + 0.04), Inches(ws[2]), Inches(0.20), size=6.4, color=INK, bold=True)
        text(slide, verdict_label(row["suggested_reddit_verdict"]), Inches(xs[3]), Inches(y + 0.04), Inches(ws[3]), Inches(0.20), size=6.4, color=color, bold=True)
        text(slide, truncate(company["analyst_positioning"], 118), Inches(xs[4]), Inches(y + 0.03), Inches(ws[4]), Inches(0.28), size=5.8, color=INK)
        y += 0.88
    footer(slide, page, 9)


def supported_wedges(slide, companies_by_id: dict[str, dict], rows: list[dict]) -> None:
    clear(slide)
    slide_title(slide, "Only two wedges have enough qualified Reddit evidence to use as proof points", "These should be framed as qualitative market evidence, not market-size or revenue validation")
    strong = [r for r in rows if r["suggested_reddit_verdict"] == "qualified_reddit_support"][:2]
    x = 0.75
    for row in strong:
        company = companies_by_id[row["claim_id"]]
        ev = row["evidence"][0]
        rect(slide, Inches(x), Inches(1.58), Inches(5.72), Inches(4.35), WHITE, line=GRID)
        rect(slide, Inches(x), Inches(1.58), Inches(5.72), Inches(0.10), GREEN)
        text(slide, company["name"], Inches(x + 0.25), Inches(1.90), Inches(1.7), Inches(0.25), size=14, color=NAVY, bold=True)
        text(slide, company["wedge"], Inches(x + 2.05), Inches(1.95), Inches(2.9), Inches(0.18), size=7.5, color=MUTED)
        labels = [
            ("YC wedge", company["workflow"]),
            ("Reddit observation", ev["excerpt"]),
            ("Implication", "Pain is observable in operator discussion; proceed to targeted interviews and primary-source sizing."),
        ]
        y = 2.50
        for label, body in labels:
            text(slide, label, Inches(x + 0.30), Inches(y), Inches(1.25), Inches(0.18), size=7.0, color=BLUE, bold=True)
            text(slide, truncate(body, 150), Inches(x + 1.55), Inches(y - 0.01), Inches(3.75), Inches(0.38), size=7.0, color=INK)
            y += 1.02
        text(slide, f"Source: r/{ev.get('subreddit', '')}; {row.get('items_rejected_by_qualification_gate', 0):,} raw items rejected", Inches(x + 0.30), Inches(5.50), Inches(4.9), Inches(0.18), size=6.2, color=MUTED)
        x += 6.05
    footer(slide, 6, 9)


def weak_and_gap(slide, companies_by_id: dict[str, dict], rows: list[dict]) -> None:
    clear(slide)
    slide_title(slide, "Several companies still require targeted proof before validation", "Weak Reddit signal is directionally useful; evidence gaps should stay out of the validated narrative")
    weak = [r for r in rows if "weak_qualified" in r["suggested_reddit_verdict"]][:5]
    gaps = [r for r in rows if r["suggested_reddit_verdict"] == "no_qualified_reddit_evidence"][:6]
    text(slide, "Weak signals: re-test with narrower subreddit queries", Inches(0.75), Inches(1.42), Inches(5.5), Inches(0.22), size=10.5, color=AMBER, bold=True)
    y = 1.86
    for row in weak:
        company = companies_by_id.get(row["claim_id"], {"name": row["claim_id"]})
        rect(slide, Inches(0.75), Inches(y - 0.04), Inches(5.65), Inches(0.38), WHITE, line=GRID)
        text(slide, company["name"], Inches(0.95), Inches(y + 0.05), Inches(1.10), Inches(0.14), size=6.7, color=NAVY, bold=True)
        text(slide, truncate(company.get("wedge", ""), 50), Inches(2.10), Inches(y + 0.05), Inches(3.65), Inches(0.14), size=6.1, color=INK)
        y += 0.48
    text(slide, "Evidence gaps: do not claim Reddit validation", Inches(7.05), Inches(1.42), Inches(4.7), Inches(0.22), size=10.5, color=RED, bold=True)
    y = 1.86
    for row in gaps:
        company = companies_by_id.get(row["claim_id"], {"name": row["claim_id"]})
        rect(slide, Inches(7.05), Inches(y - 0.04), Inches(5.05), Inches(0.38), WHITE, line=GRID)
        text(slide, company["name"], Inches(7.25), Inches(y + 0.05), Inches(1.10), Inches(0.14), size=6.7, color=NAVY, bold=True)
        text(slide, f"{row.get('items_rejected_by_qualification_gate', 0):,} raw items rejected", Inches(8.48), Inches(y + 0.05), Inches(2.2), Inches(0.14), size=6.1, color=RED, bold=True)
        y += 0.48
    rect(slide, Inches(0.75), Inches(5.72), Inches(11.35), Inches(0.56), SOFT, line=GRID)
    text(slide, "Next step: targeted operator evidence collection by vertical, then update the deck only when new items pass the same qualification gate.", Inches(1.00), Inches(5.90), Inches(10.6), Inches(0.16), size=8.5, color=NAVY, bold=True)
    footer(slide, 7, 9)


def research_plan(slide) -> None:
    clear(slide)
    slide_title(slide, "The next sprint should collect fewer, higher-quality operator sources", "The goal is not more Reddit volume; it is better-matched practitioner evidence")
    rows = [
        ("Property management", "Deepen", "Interview property managers; target r/PropertyManagement and r/Landlord around maintenance/tickets."),
        ("Construction review", "Deepen", "Collect contractor, architect, owner threads on drawing errors, rework, punch lists, and delays."),
        ("Clinic communications", "Re-test", "Target clinic managers/front desk workflows; separate dental, PT, and physician office evidence."),
        ("Hotel operations", "Re-test", "Target front-desk and revenue operations threads; filter out guest anecdote noise."),
        ("Mortgage / ins. / finance", "Rebuild", "Use trade forums, primary sources, and customer proof before using Reddit as evidence."),
    ]
    y = 1.48
    for vertical, action, note in rows:
        color = GREEN if action == "Deepen" else AMBER if action == "Re-test" else RED
        rect(slide, Inches(0.75), Inches(y), Inches(11.35), Inches(0.70), WHITE, line=GRID)
        rect(slide, Inches(0.75), Inches(y), Inches(0.08), Inches(0.70), color)
        text(slide, vertical, Inches(1.02), Inches(y + 0.17), Inches(1.8), Inches(0.18), size=8.2, color=NAVY, bold=True)
        text(slide, action, Inches(3.05), Inches(y + 0.17), Inches(0.85), Inches(0.18), size=7.2, color=color, bold=True)
        text(slide, note, Inches(4.10), Inches(y + 0.16), Inches(7.5), Inches(0.20), size=7.1, color=INK)
        y += 0.86
    footer(slide, 8, 9)


def appendix(slide, rows: list[dict]) -> None:
    clear(slide)
    slide_title(slide, "Appendix: qualified source traceability", "The full report and JSON evidence pack remain in the run folder")
    qualified = [r for r in rows if r.get("evidence")]
    y = 1.46
    for row in qualified[:8]:
        ev = row["evidence"][0]
        color = verdict_color(row["suggested_reddit_verdict"])
        rect(slide, Inches(0.75), Inches(y), Inches(11.35), Inches(0.48), WHITE, line=GRID)
        rect(slide, Inches(0.75), Inches(y), Inches(0.08), Inches(0.48), color)
        text(slide, row["claim_id"].replace("claim-", "C-"), Inches(0.98), Inches(y + 0.13), Inches(0.50), Inches(0.14), size=6.4, color=NAVY, bold=True)
        text(slide, f"r/{ev.get('subreddit', '')}", Inches(1.62), Inches(y + 0.13), Inches(1.70), Inches(0.14), size=6.3, color=BLUE, bold=True)
        text(slide, verdict_label(row["suggested_reddit_verdict"]), Inches(3.55), Inches(y + 0.13), Inches(1.45), Inches(0.14), size=6.3, color=color, bold=True)
        text(slide, truncate(ev.get("excerpt", ""), 130), Inches(5.25), Inches(y + 0.10), Inches(6.25), Inches(0.20), size=5.7, color=INK)
        y += 0.58
    footer(slide, 9, 9)


def build() -> Path:
    if not TEMPLATE.exists():
        raise SystemExit(f"Branded template not found: {TEMPLATE}")
    companies = load_json(COMPANY_PACK)["companies"]
    evidence = load_json(EVIDENCE_PACK)
    rows = evidence["claims"]
    rows_by_id = {r["claim_id"]: r for r in rows}
    companies_by_id = {c["claim_id"]: c for c in companies}

    prs = Presentation(str(TEMPLATE))
    while len(prs.slides) < 9:
        prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(len(prs.slides._sldIdLst) - 1, 8, -1):
        rid = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[i]
    prs.slide_width, prs.slide_height = W, H

    cover(prs.slides[0], rows)
    executive_answer(prs.slides[1], rows)
    evidence_funnel(prs.slides[2], rows)
    company_table(prs.slides[3], companies[:5], rows_by_id, 4, "Five operating wedges are plausible, but only one has qualified Reddit support")
    company_table(prs.slides[4], companies[5:], rows_by_id, 5, "Finance, claims, and lending wedges need stronger external proof before validation")
    supported_wedges(prs.slides[5], companies_by_id, rows)
    weak_and_gap(prs.slides[6], companies_by_id, rows)
    research_plan(prs.slides[7])
    appendix(prs.slides[8], rows)

    prs.save(OUT)
    problems = validate_pptx(OUT)
    if problems:
        raise SystemExit("PPTX validation failed:\n- " + "\n- ".join(problems))
    print(OUT)
    return OUT


if __name__ == "__main__":
    build()
