#!/usr/bin/env python3
"""Build a synthesis-first executive deck.

Slides are generated from semantic_synthesis_pack.json, not directly from
company records or Reddit excerpts. Source artifacts remain audit trail only.
"""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPTXKIT = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT))
from pptxkit import validate_pptx  # noqa: E402

RUN_DIR = Path(__file__).resolve().parent
TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE") or "/home/shekerk/.claude/templates/branded-template.pptx")
SYNTHESIS = RUN_DIR / "semantic_synthesis_pack.json"
OUT = RUN_DIR / "yc-agentic-market-research-semantic-synthesis-v4-draft.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x0A, 0x16, 0x28)
BLUE = RGBColor(0x1F, 0x5F, 0xB8)
TEAL = RGBColor(0x00, 0xA8, 0x8F)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
AMBER = RGBColor(0xD4, 0x8B, 0x00)
RED = RGBColor(0xC7, 0x3E, 0x3A)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5F, 0x6C, 0x7A)
GRID = RGBColor(0xD9, 0xDF, 0xE5)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)


def rect(slide, x, y, w, h, color, *, line=None, radius=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, *, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    for i, item in enumerate(value if isinstance(value, list) else [value]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(item, dict):
            content = item["text"]
            size_i = item.get("size", size)
            color_i = item.get("color", color)
            bold_i = item.get("bold", bold)
            p.space_before = Pt(item.get("space_before", 0))
            if item.get("bullet"):
                content = "- " + content
        else:
            content, size_i, color_i, bold_i = str(item), size, color, bold
        r = p.add_run()
        r.text = content
        r.font.name = "Arial"
        r.font.size = Pt(size_i)
        r.font.bold = bold_i
        r.font.color.rgb = color_i
    return box


def title(slide, msg: str, note: str = "") -> None:
    text(slide, msg, Inches(0.58), Inches(0.28), Inches(11.9), Inches(0.70), size=17, color=NAVY, bold=True)
    if note:
        text(slide, note, Inches(0.58), Inches(0.92), Inches(11.3), Inches(0.26), size=7.6, color=MUTED)
    rect(slide, Inches(0.58), Inches(1.23), Inches(1.05), Inches(0.04), TEAL)


def foot(slide, source: str = "Sources: semantic synthesis pack; Printing Press YC CLI output; qualified Reddit evidence pack") -> None:
    text(slide, source, Inches(0.58), Inches(7.10), Inches(10.8), Inches(0.16), size=5.8, color=MUTED)


def short(value: str, n: int) -> str:
    value = " ".join(value.split())
    return value if len(value) <= n else value[: n - 1].rstrip() + "."


def dot(slide, x, y, color, label, sub="") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text(slide, label, Inches(x + 0.22), Inches(y - 0.03), Inches(1.35), Inches(0.14), size=5.8, color=INK, bold=True)
    if sub:
        text(slide, sub, Inches(x + 0.22), Inches(y + 0.14), Inches(1.45), Inches(0.12), size=4.8, color=MUTED)


def cover(slide, s: dict) -> None:
    clear(slide)
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, 0, 0, Inches(0.18), H, TEAL)
    text(slide, "SEMANTIC SYNTHESIS", Inches(0.78), Inches(0.70), Inches(4.2), Inches(0.22), size=9, color=TEAL, bold=True)
    text(slide, s["bluf"], Inches(0.78), Inches(1.18), Inches(10.8), Inches(0.95), size=29, color=WHITE, bold=True)
    text(slide, s["core_thesis"], Inches(0.82), Inches(2.58), Inches(9.4), Inches(0.55), size=12.5, color=WHITE)
    for i, beat in enumerate(s["slide_story"][:3]):
        y = 4.28 + i * 0.55
        rect(slide, Inches(0.82), Inches(y), Inches(0.16), Inches(0.16), TEAL)
        text(slide, beat, Inches(1.12), Inches(y - 0.03), Inches(8.2), Inches(0.18), size=9.0, color=WHITE)
    foot(slide)


def thesis_tree(slide, s: dict) -> None:
    clear(slide)
    title(slide, "Agentic value concentrates where workflow failure is visible, costly, and repetitive", "This is the synthesis layer: the slide is organized by semantic pattern, not by source-record order")
    rect(slide, Inches(0.78), Inches(1.72), Inches(3.15), Inches(0.88), NAVY)
    text(slide, "Where agents can win", Inches(1.05), Inches(1.98), Inches(2.45), Inches(0.22), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    branches = [
        ("Coordination burden", "Handoffs across people, systems, and follow-up loops", GREEN),
        ("Rework prevention", "Errors become expensive only after work moves downstream", GREEN),
        ("Regulated operations", "High-value workflows, but proof must come from primary/operator sources", AMBER),
        ("Back-office fragmentation", "Clear workflow logic, but Reddit evidence is currently weak", AMBER),
    ]
    x = 0.70
    for label, desc, color in branches:
        rect(slide, Inches(x), Inches(3.16), Inches(2.88), Inches(1.32), WHITE, line=color)
        rect(slide, Inches(x), Inches(3.16), Inches(2.88), Inches(0.08), color)
        text(slide, label, Inches(x + 0.18), Inches(3.48), Inches(2.35), Inches(0.24), size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, Inches(x + 0.25), Inches(3.88), Inches(2.25), Inches(0.34), size=6.8, color=INK, align=PP_ALIGN.CENTER)
        x += 3.10
    foot(slide)


def opportunity_map(slide, s: dict) -> None:
    clear(slide)
    title(slide, "The best candidates sit where workflow specificity and Reddit evidence both score high", "The map uses synthesized scores from the source pack and evidence gate")
    x0, y0, w, h = 1.05, 1.62, 9.4, 4.65
    rect(slide, Inches(x0), Inches(y0), Inches(w), Inches(h), WHITE, line=GRID)
    text(slide, "Low Reddit evidence", Inches(x0), Inches(y0 + h + 0.15), Inches(1.6), Inches(0.15), size=6, color=MUTED)
    text(slide, "High Reddit evidence", Inches(x0 + w - 1.75), Inches(y0 + h + 0.15), Inches(1.6), Inches(0.15), size=6, color=MUTED, align=PP_ALIGN.RIGHT)
    text(slide, "High workflow specificity", Inches(x0 + w + 0.15), Inches(y0), Inches(1.4), Inches(0.2), size=6, color=MUTED)
    text(slide, "Low specificity", Inches(x0 + w + 0.15), Inches(y0 + h - 0.15), Inches(1.4), Inches(0.2), size=6, color=MUTED)
    rect(slide, Inches(x0 + w / 2), Inches(y0), Inches(0.01), Inches(h), GRID)
    rect(slide, Inches(x0), Inches(y0 + h / 2), Inches(w), Inches(0.01), GRID)
    for item in s["company_scores"]:
        x = x0 + (item["reddit_evidence"] - 0.8) / 2.4 * (w - 1.2)
        y = y0 + (3.25 - item["workflow_specificity"]) / 2.4 * (h - 0.7)
        color = GREEN if item["diligence_priority"] == "High" else AMBER if item["diligence_priority"] == "Medium" else RED
        dot(slide, x, y, color, item["company"], item["diligence_priority"])
    text(slide, "Priority zone", Inches(x0 + w - 2.1), Inches(y0 + 0.25), Inches(1.5), Inches(0.18), size=8, color=GREEN, bold=True)
    foot(slide)


def evidence_concentration(slide, s: dict) -> None:
    clear(slide)
    title(slide, "Qualified Reddit evidence clusters in two pain patterns, not across the full startup list", "This slide summarizes the evidence semantics instead of pasting Reddit excerpts")
    themes = [t for t in s["semantic_themes"] if t["theme"] in {"Coordination-heavy operations", "Rework and risk prevention"}]
    x = 0.82
    for theme in themes:
        rect(slide, Inches(x), Inches(1.65), Inches(5.65), Inches(4.40), WHITE, line=GREEN)
        rect(slide, Inches(x), Inches(1.65), Inches(5.65), Inches(0.10), GREEN)
        text(slide, theme["theme"], Inches(x + 0.25), Inches(1.98), Inches(4.8), Inches(0.24), size=13, color=NAVY, bold=True)
        text(slide, theme["synthesis"], Inches(x + 0.30), Inches(2.58), Inches(4.75), Inches(0.55), size=10.2, color=INK)
        text(slide, "Companies", Inches(x + 0.30), Inches(3.62), Inches(1.1), Inches(0.14), size=7, color=BLUE, bold=True)
        text(slide, ", ".join(theme["companies"]), Inches(x + 1.45), Inches(3.60), Inches(3.8), Inches(0.18), size=8.2, color=INK)
        text(slide, "Evidence read", Inches(x + 0.30), Inches(4.42), Inches(1.1), Inches(0.14), size=7, color=BLUE, bold=True)
        text(slide, theme["evidence_read"], Inches(x + 1.45), Inches(4.38), Inches(3.8), Inches(0.45), size=7.3, color=INK)
        x += 6.05
    foot(slide)


def heatmap(slide, s: dict) -> None:
    clear(slide)
    title(slide, "The diligence stack separates strong wedges from plausible but unvalidated ideas", "Scores are qualitative synthesis outputs: 3 = high, 2 = medium, 1 = low")
    cols = [("Workflow", "workflow_specificity"), ("Reddit", "reddit_evidence"), ("Integration", "integration_burden")]
    text(slide, "Company", Inches(0.80), Inches(1.42), Inches(1.25), Inches(0.18), size=7, color=MUTED, bold=True)
    for i, (label, _) in enumerate(cols):
        text(slide, label, Inches(2.62 + i * 1.10), Inches(1.42), Inches(0.85), Inches(0.18), size=7, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "Synthesis implication", Inches(6.15), Inches(1.42), Inches(4.6), Inches(0.18), size=7, color=MUTED, bold=True)
    y = 1.78
    for item in s["company_scores"]:
        rect(slide, Inches(0.72), Inches(y - 0.04), Inches(11.75), Inches(0.42), WHITE, line=GRID)
        text(slide, item["company"], Inches(0.88), Inches(y + 0.06), Inches(1.35), Inches(0.14), size=6.7, color=NAVY, bold=True)
        for i, (_, key) in enumerate(cols):
            val = item[key]
            color = GREEN if val == 3 else AMBER if val == 2 else RED
            rect(slide, Inches(2.74 + i * 1.10), Inches(y + 0.04), Inches(0.30), Inches(0.18), color)
        text(slide, short(item["synthesis"], 108), Inches(6.15), Inches(y + 0.04), Inches(5.85), Inches(0.16), size=5.8, color=INK)
        y += 0.51
    foot(slide)


def unvalidated(slide, s: dict) -> None:
    clear(slide)
    title(slide, "Most wedges are still hypotheses; the next research should target proof gaps", "This is the practical implication of semantic retrieval plus evidence qualification")
    groups = [t for t in s["semantic_themes"] if t["theme"] not in {"Coordination-heavy operations", "Rework and risk prevention"}]
    y = 1.62
    for theme in groups:
        rect(slide, Inches(0.82), Inches(y), Inches(11.4), Inches(1.05), WHITE, line=GRID)
        rect(slide, Inches(0.82), Inches(y), Inches(0.08), Inches(1.05), AMBER if "Back-office" in theme["theme"] or "Professional" in theme["theme"] else RED)
        text(slide, theme["theme"], Inches(1.08), Inches(y + 0.18), Inches(2.6), Inches(0.20), size=9.2, color=NAVY, bold=True)
        text(slide, ", ".join(theme["companies"]), Inches(3.90), Inches(y + 0.18), Inches(2.6), Inches(0.18), size=7.2, color=BLUE, bold=True)
        text(slide, theme["synthesis"], Inches(6.60), Inches(y + 0.14), Inches(5.05), Inches(0.24), size=6.8, color=INK)
        text(slide, theme["evidence_read"], Inches(6.60), Inches(y + 0.58), Inches(5.05), Inches(0.20), size=6.3, color=MUTED)
        y += 1.27
    foot(slide)


def decision(slide, s: dict) -> None:
    clear(slide)
    title(slide, "The decision is to narrow diligence, not to validate the whole list", "A useful evidence deck should change what we do next")
    rect(slide, Inches(0.95), Inches(1.65), Inches(11.25), Inches(1.05), NAVY)
    text(slide, s["decision"], Inches(1.25), Inches(1.92), Inches(10.45), Inches(0.40), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    actions = [
        ("Deepen now", "CentralComs, InspectMind AI", GREEN),
        ("Re-test narrowly", "Dodo, Lance, Klarify, Fazeshift, ProcIndex", AMBER),
        ("Hold out of validated story", "Kastle, Basepilot, Veritus", RED),
    ]
    x = 0.95
    for label, names, color in actions:
        rect(slide, Inches(x), Inches(3.45), Inches(3.45), Inches(1.35), WHITE, line=color)
        rect(slide, Inches(x), Inches(3.45), Inches(3.45), Inches(0.08), color)
        text(slide, label, Inches(x + 0.25), Inches(3.78), Inches(2.8), Inches(0.22), size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(slide, names, Inches(x + 0.30), Inches(4.25), Inches(2.75), Inches(0.28), size=7.2, color=INK, align=PP_ALIGN.CENTER)
        x += 3.85
    foot(slide)


def build() -> Path:
    s = load_json(SYNTHESIS)
    prs = Presentation(str(TEMPLATE))
    while len(prs.slides) < 7:
        prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(len(prs.slides._sldIdLst) - 1, 6, -1):
        rid = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[i]
    prs.slide_width, prs.slide_height = W, H
    cover(prs.slides[0], s)
    thesis_tree(prs.slides[1], s)
    opportunity_map(prs.slides[2], s)
    evidence_concentration(prs.slides[3], s)
    heatmap(prs.slides[4], s)
    unvalidated(prs.slides[5], s)
    decision(prs.slides[6], s)
    prs.save(OUT)
    problems = validate_pptx(OUT)
    if problems:
        raise SystemExit("PPTX validation failed:\n- " + "\n- ".join(problems))
    print(OUT)
    return OUT


if __name__ == "__main__":
    build()
