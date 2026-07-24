#!/usr/bin/env python3
"""Build a branded client-ready PPTX from qualified Reddit market evidence."""

from __future__ import annotations

import json
import os
import re
import sys
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPTXKIT = Path("/home/shekerk/.claude/skills/branded-pptx-deck/scripts")
sys.path.insert(0, str(PPTXKIT))
from pptxkit import validate_pptx  # noqa: E402

RUN_DIR = Path(__file__).resolve().parent
CLAIM_PACK = RUN_DIR / "claim-pack.json"
EVIDENCE_PACK = RUN_DIR / "reddit-evidence-pack.json"
COMPANY_PACK = RUN_DIR / "yc-company-source-pack.json"
TEMPLATE = Path(os.environ.get("BRANDED_PPTX_TEMPLATE") or "/home/shekerk/.claude/templates/branded-template.pptx")
OUT = RUN_DIR / "yc-agentic-market-research-qualified-reddit-branded-draft.pptx"
TOTAL_SLIDES = 9

NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_2 = RGBColor(0x12, 0x24, 0x3A)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
DARK_TEAL = RGBColor(0x00, 0x8F, 0x75)
GOLD = RGBColor(0xFF, 0xB8, 0x00)
AMBER = RGBColor(0xF2, 0xA8, 0x3B)
CORAL = RGBColor(0xE0, 0x5A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def load_companies() -> list[dict]:
    return load_json(COMPANY_PACK)["companies"]


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def trim_to_first_n(prs: Presentation, n: int) -> None:
    for index in range(len(prs.slides._sldIdLst) - 1, n - 1, -1):
        rid = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[index]


def keep_template_base(shape) -> bool:
    is_picture_rail = shape.shape_type == 13 and shape.left >= Inches(8.4)
    is_left_background = shape.left <= 1 and shape.top <= 1 and shape.width >= Inches(8.7)
    is_cover_rule = shape.left <= 1 and shape.width >= Inches(8.7) and shape.height <= Inches(0.08)
    return bool(is_picture_rail or is_left_background or is_cover_rule)


def clear_to_template_base(slide) -> None:
    for shape in list(slide.shapes):
        if keep_template_base(shape):
            continue
        remove_shape(shape)


def rect(slide, x, y, w, h, color, radius: bool = False, line=None):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, *, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    lines = value if isinstance(value, list) else [value]
    for index, item in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        if isinstance(item, dict):
            content = item.get("text", "")
            paragraph.space_before = Pt(item.get("space_before", 0))
            item_size = item.get("size", size)
            item_color = item.get("color", color)
            item_bold = item.get("bold", bold)
            if item.get("bullet"):
                content = "- " + content
        else:
            content = str(item)
            item_size = size
            item_color = color
            item_bold = bold
        run = paragraph.add_run()
        run.text = content
        run.font.name = "Calibri"
        run.font.size = Pt(item_size)
        run.font.bold = item_bold
        run.font.color.rgb = item_color
    return shape


def short(value: str, limit: int = 150) -> str:
    value = re.sub(r"\s+", " ", value.replace("*", "")).strip()
    return value if len(value) <= limit else value[: limit - 1].rstrip() + "..."


def claim_id(value: str) -> str:
    match = re.search(r"(\d+)$", value)
    return f"C-{int(match.group(1)):03d}" if match else value.upper()


def company_from_claim(claim: str) -> str:
    first = claim.split(" is ", 1)[0]
    if len(first) < 36 and not first.lower().startswith(("property managers", "hotel operators", "clinic front")):
        return first
    if claim.lower().startswith("property managers"):
        return "CentralComs wedge"
    if claim.lower().startswith("hotel operators"):
        return "Lance wedge"
    if claim.lower().startswith("clinic front"):
        return "Dodo wedge"
    return "Workflow claim"


def company_by_claim(companies: list[dict]) -> dict[str, dict]:
    return {company["claim_id"]: company for company in companies}


def short_label(value: str, limit: int = 18) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    return value if len(value) <= limit else value[: limit - 1].rstrip() + "."


def verdict_color(verdict: str):
    if verdict == "qualified_reddit_support":
        return TEAL
    if "weak_qualified" in verdict:
        return AMBER
    return CORAL


def verdict_label(verdict: str) -> str:
    return {
        "qualified_reddit_support": "Qualified Reddit support",
        "weak_qualified_reddit_support": "Weak qualified signal",
        "no_qualified_reddit_evidence": "No qualified Reddit evidence",
    }.get(verdict, verdict.replace("_", " "))


def footer(slide, page: int, total: int, dark: bool = False) -> None:
    color = WHITE if dark else MUTED
    text(slide, "Qualified Reddit market evidence | client-ready draft", Inches(0.32), Inches(7.08), Inches(4.9), Inches(0.2), size=7.5, color=color)
    text(slide, f"{page} / {total}", Inches(7.22), Inches(7.08), Inches(0.8), Inches(0.2), size=7.5, color=color, align=PP_ALIGN.RIGHT)


def side_panel(slide, title: str, lines: list[str], status: str = "DRAFT") -> None:
    rect(slide, Inches(8.84), 0, Inches(4.49), Inches(7.5), NAVY_2)
    rect(slide, Inches(8.84), 0, Inches(0.08), Inches(7.5), TEAL)
    text(slide, title, Inches(9.22), Inches(0.55), Inches(3.35), Inches(0.7), size=20, color=WHITE, bold=True)
    y = 1.48
    for line in lines:
        text(slide, line, Inches(9.22), Inches(y), Inches(3.38), Inches(0.42), size=11, color=WHITE)
        y += 0.58
    rect(slide, Inches(9.22), Inches(6.22), Inches(2.45), Inches(0.34), GOLD)
    text(slide, status, Inches(9.31), Inches(6.30), Inches(2.22), Inches(0.16), size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def content_bg(slide, color=SOFT) -> None:
    rect(slide, 0, 0, Inches(8.84), Inches(7.5), color)


def top_evidence(row: dict) -> dict | None:
    evidence = row.get("evidence") or []
    return evidence[0] if evidence else None


def cover(slide, claims: list[dict], rows: list[dict]) -> None:
    counts = Counter(row["suggested_reddit_verdict"] for row in rows)
    rejected = sum(row.get("items_rejected_by_qualification_gate", 0) for row in rows)
    text(slide, "QUALIFIED REDDIT MARKET EVIDENCE", Inches(0.8), Inches(0.82), Inches(6.2), Inches(0.32), size=13, color=TEAL, bold=True)
    text(slide, "YC Agentic AI Vertical Signals", Inches(0.8), Inches(1.32), Inches(7.25), Inches(1.15), size=28, color=WHITE, bold=True)
    rect(slide, Inches(0.8), Inches(2.72), Inches(1.45), Inches(0.05), TEAL)
    text(slide, "Reddit data is useful only after qualification. This deck separates real operator workflow evidence from raw search noise.", Inches(0.8), Inches(3.06), Inches(6.7), Inches(0.82), size=12.4, color=WHITE)
    y = 4.35
    for value, label, color in [
        (str(counts["qualified_reddit_support"]), "strong qualified claims", TEAL),
        (str(counts["weak_qualified_reddit_support"]), "weak signals", GOLD),
        (str(counts["no_qualified_reddit_evidence"]), "evidence gaps", CORAL),
    ]:
        rect(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.42), color, radius=True)
        text(slide, f"{value}  {label}", Inches(0.94), Inches(y + 0.1), Inches(2.2), Inches(0.18), size=7.4, color=NAVY if color == GOLD else WHITE, bold=True, align=PP_ALIGN.CENTER)
        y += 0.52
    text(slide, f"Built {date.today().isoformat()} | Status: DRAFT", Inches(0.8), Inches(6.15), Inches(4.3), Inches(0.25), size=9, color=GOLD, bold=True)
    side_panel(slide, "Evidence standard", [f"{len(claims)} claims scored", f"{rejected:,} raw items rejected", "Raw Reddit is not validation"])
    footer(slide, 1, TOTAL_SLIDES, dark=True)


def executive_verdict(slide, rows: list[dict]) -> None:
    content_bg(slide, SOFT)
    counts = Counter(row["suggested_reddit_verdict"] for row in rows)
    text(slide, "Executive Verdict", Inches(0.42), Inches(0.32), Inches(7.6), Inches(0.48), size=22, color=NAVY, bold=True)
    text(slide, "Use Reddit as a market-signal filter, not as blanket proof for every niche company.", Inches(0.42), Inches(0.86), Inches(7.3), Inches(0.28), size=10, color=INK)
    metrics = [
        (str(counts["qualified_reddit_support"]), "qualified support", TEAL),
        (str(counts["weak_qualified_reddit_support"]), "weak signal", AMBER),
        (str(counts["no_qualified_reddit_evidence"]), "unsupported by qualified Reddit", CORAL),
    ]
    x = 0.45
    for value, label, color in metrics:
        rect(slide, Inches(x), Inches(1.36), Inches(2.32), Inches(0.92), NAVY)
        text(slide, value, Inches(x + 0.16), Inches(1.50), Inches(0.7), Inches(0.34), size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, label, Inches(x + 0.90), Inches(1.48), Inches(1.16), Inches(0.34), size=6.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += 2.55
    bullets = [
        "Strongest qualified evidence: property management and construction review pain.",
        "Weak signals: clinic communications and hotel operations need better targeted threads before they are investment-grade evidence.",
        "No qualified Reddit support yet: therapy admin, mortgage servicing, AR/accounting, insurance claims, and consumer lending claims.",
        "Recommendation: ship no broad validation claim until a vertical has multiple qualified operator threads plus primary-source corroboration.",
    ]
    text(slide, [{"text": item, "bullet": True, "space_before": 8, "size": 10.5} for item in bullets], Inches(0.48), Inches(2.62), Inches(7.2), Inches(2.3), size=10.5, color=INK)
    side_panel(slide, "So what", ["Prioritize evidence-rich wedges", "Mark gaps explicitly", "Avoid Reddit-shaped filler"])
    footer(slide, 2, TOTAL_SLIDES)


def method_slide(slide) -> None:
    content_bg(slide, SOFT)
    text(slide, "How The Evidence Was Qualified", Inches(0.42), Inches(0.32), Inches(7.7), Inches(0.48), size=22, color=NAVY, bold=True)
    cards = [
        ("1", "Workflow match", "The post/comment must discuss the same operating workflow, not just a generic company or AI term."),
        ("2", "Persona match", "The source must plausibly come from or address the operator, practitioner, buyer, or affected user."),
        ("3", "Pain signal", "Accepted evidence must contain a concrete pain, workaround, delay, objection, manual process, or adoption signal."),
        ("4", "Noise rejection", "Broad or irrelevant Reddit threads are counted as rejected discovery, not as evidence."),
    ]
    y = 1.18
    for number, title, body in cards:
        rect(slide, Inches(0.55), Inches(y), Inches(7.15), Inches(0.90), WHITE, line=TEAL)
        rect(slide, Inches(0.72), Inches(y + 0.18), Inches(0.42), Inches(0.42), TEAL, radius=True)
        text(slide, number, Inches(0.83), Inches(y + 0.29), Inches(0.2), Inches(0.1), size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title, Inches(1.32), Inches(y + 0.12), Inches(1.9), Inches(0.24), size=11.5, color=NAVY, bold=True)
        text(slide, body, Inches(3.08), Inches(y + 0.12), Inches(4.35), Inches(0.50), size=8.0, color=INK)
        y += 1.08
    side_panel(slide, "Decision rule", ["No qualified item", "= evidence gap", "One item", "= weak signal", "Multiple items", "= qualified support"])
    footer(slide, 3, TOTAL_SLIDES)


def company_profiles(slide, companies: list[dict], rows_by_id: dict[str, dict], page: int, title: str) -> None:
    content_bg(slide, SOFT)
    text(slide, title, Inches(0.42), Inches(0.28), Inches(7.7), Inches(0.56), size=18, color=NAVY, bold=True)
    text(slide, "Each profile uses YC source fields for facts and a separate qualified Reddit status for market evidence.", Inches(0.42), Inches(0.78), Inches(7.3), Inches(0.34), size=8.2, color=MUTED)
    y = 1.22
    for company in companies:
        row = rows_by_id[company["claim_id"]]
        verdict = row["suggested_reddit_verdict"]
        color = verdict_color(verdict)
        rect(slide, Inches(0.52), Inches(y), Inches(7.18), Inches(0.86), WHITE, line=color)
        text(slide, company["name"], Inches(0.70), Inches(y + 0.12), Inches(1.10), Inches(0.22), size=8.4, color=NAVY, bold=True)
        text(slide, f"{company['batch']} | team {company['team_size']} | {company['status']}", Inches(0.70), Inches(y + 0.42), Inches(1.78), Inches(0.18), size=5.8, color=MUTED)
        text(slide, verdict_label(verdict), Inches(2.60), Inches(y + 0.12), Inches(1.25), Inches(0.24), size=6.2, color=color, bold=True)
        text(slide, short(company["wedge"], 42), Inches(3.92), Inches(y + 0.12), Inches(1.60), Inches(0.22), size=6.6, color=NAVY, bold=True)
        text(slide, short(company["analyst_positioning"], 128), Inches(3.92), Inches(y + 0.42), Inches(3.35), Inches(0.22), size=5.3, color=INK)
        y += 1.02
    side_panel(slide, "Readout", ["YC facts", "Analyst wedge", "Reddit status", "No raw scrape text"])
    footer(slide, page, TOTAL_SLIDES)


def strong_evidence(slide, claims_by_id: dict[str, dict], company_map: dict[str, dict], rows: list[dict]) -> None:
    content_bg(slide, SOFT)
    text(slide, "Where Reddit Actually Supports The Thesis", Inches(0.42), Inches(0.32), Inches(7.7), Inches(0.48), size=21, color=NAVY, bold=True)
    strong = [row for row in rows if row["suggested_reddit_verdict"] == "qualified_reddit_support"][:2]
    y = 1.14
    for row in strong:
        claim = claims_by_id[row["claim_id"]]
        company = company_map.get(row["claim_id"], {})
        ev = top_evidence(row) or {}
        rect(slide, Inches(0.55), Inches(y), Inches(7.2), Inches(1.65), NAVY)
        text(slide, f"{company.get('name', company_from_claim(claim['claim']))} | {company.get('wedge', '')}", Inches(0.78), Inches(y + 0.18), Inches(3.0), Inches(0.25), size=10.5, color=TEAL, bold=True)
        text(slide, short(company.get("analyst_positioning", claim["claim"]), 160), Inches(0.78), Inches(y + 0.55), Inches(3.0), Inches(0.58), size=7.8, color=WHITE)
        text(slide, "Qualified evidence", Inches(4.08), Inches(y + 0.20), Inches(1.6), Inches(0.18), size=7.5, color=GOLD, bold=True)
        text(slide, short(ev.get("excerpt", ""), 210), Inches(4.08), Inches(y + 0.50), Inches(3.35), Inches(0.62), size=7.2, color=WHITE)
        text(slide, f"r/{ev.get('subreddit', '')} | rejected raw items: {row.get('items_rejected_by_qualification_gate', 0):,}", Inches(4.08), Inches(y + 1.20), Inches(3.35), Inches(0.18), size=6.8, color=GOLD)
        y += 2.02
    side_panel(slide, "Use in narrative", ["Operator pain exists", "Cite as qualitative signal", "Still avoid TAM/funding claims"])
    footer(slide, 6, TOTAL_SLIDES)


def weak_and_gap_slide(slide, claims_by_id: dict[str, dict], company_map: dict[str, dict], rows: list[dict]) -> None:
    content_bg(slide, SOFT)
    text(slide, "What Should Not Be Claimed Yet", Inches(0.42), Inches(0.32), Inches(7.7), Inches(0.48), size=22, color=NAVY, bold=True)
    weak = [row for row in rows if "weak_qualified" in row["suggested_reddit_verdict"]]
    gaps = [row for row in rows if row["suggested_reddit_verdict"] == "no_qualified_reddit_evidence"]
    text(slide, "Weak signals", Inches(0.58), Inches(1.15), Inches(2.8), Inches(0.26), size=13, color=AMBER, bold=True)
    y = 1.55
    for row in weak[:5]:
        claim = claims_by_id[row["claim_id"]]
        company = company_map.get(row["claim_id"], {})
        ev = top_evidence(row) or {}
        label = company.get("name", company_from_claim(claim["claim"]))
        text(slide, f"{claim_id(row['claim_id'])} {label}: {short(ev.get('excerpt', 'single qualified item only'), 85)}", Inches(0.65), Inches(y), Inches(6.9), Inches(0.28), size=6.5, color=INK)
        y += 0.42
    text(slide, "Evidence gaps", Inches(0.58), Inches(3.92), Inches(2.8), Inches(0.26), size=13, color=CORAL, bold=True)
    y = 4.30
    for row in gaps[:6]:
        claim = claims_by_id[row["claim_id"]]
        company = company_map.get(row["claim_id"], {})
        label = company.get("name", company_from_claim(claim["claim"]))
        text(slide, f"{claim_id(row['claim_id'])} {label}: no qualified Reddit item after {row.get('items_rejected_by_qualification_gate', 0):,} rejects.", Inches(0.65), Inches(y), Inches(6.9), Inches(0.25), size=7.1, color=INK)
        y += 0.34
    side_panel(slide, "Action", ["Re-query targeted subreddits", "Collect primary sources", "Do not label as validated"])
    footer(slide, 7, TOTAL_SLIDES)


def recommendations(slide) -> None:
    content_bg(slide, SOFT)
    text(slide, "Recommended Next Research Sprint", Inches(0.42), Inches(0.32), Inches(7.7), Inches(0.48), size=22, color=NAVY, bold=True)
    rows = [
        ("Property management", "Keep", "Reddit shows real workflow pain around spreadsheets, maintenance, tenant/vendor coordination."),
        ("Construction review", "Keep", "Reddit shows delay, rework, punch-list, and contract friction in construction workflows."),
        ("Clinic communications", "Re-test", "Current source is plausible but thin; collect more clinic/front-desk operator threads."),
        ("Hotel operations", "Re-test", "Weak signal is not enough; collect call, booking, and guest workflow evidence specifically."),
        ("Mortgage / insurance / AR", "Gap", "Current Reddit run did not produce qualified evidence for these wedges."),
    ]
    y = 1.20
    for vertical, action, rationale in rows:
        color = TEAL if action == "Keep" else AMBER if action == "Re-test" else CORAL
        rect(slide, Inches(0.55), Inches(y), Inches(7.2), Inches(0.78), SOFT, line=color)
        text(slide, vertical, Inches(0.72), Inches(y + 0.12), Inches(1.65), Inches(0.38), size=6.6, color=NAVY, bold=True)
        text(slide, action, Inches(2.52), Inches(y + 0.12), Inches(0.85), Inches(0.28), size=6.8, color=color, bold=True)
        text(slide, rationale, Inches(3.42), Inches(y + 0.10), Inches(4.02), Inches(0.38), size=6.6, color=INK)
        y += 0.88
    side_panel(slide, "Deck rule", ["Only qualified evidence", "No raw Reddit filler", "Primary facts remain separate"])
    footer(slide, 8, TOTAL_SLIDES)


def appendix(slide, rows: list[dict]) -> None:
    content_bg(slide, SOFT)
    text(slide, "Source Appendix", Inches(0.42), Inches(0.32), Inches(7.7), Inches(0.48), size=22, color=NAVY, bold=True)
    text(slide, "Top qualified sources retained after the evidence gate. Full JSON and Markdown report remain in the run folder.", Inches(0.42), Inches(0.86), Inches(7.3), Inches(0.28), size=9.4, color=MUTED)
    qualified = [row for row in rows if row.get("evidence")]
    y = 1.35
    for row in qualified[:8]:
        ev = top_evidence(row) or {}
        rect(slide, Inches(0.55), Inches(y), Inches(7.2), Inches(0.48), SOFT, line=verdict_color(row["suggested_reddit_verdict"]))
        text(slide, claim_id(row["claim_id"]), Inches(0.72), Inches(y + 0.12), Inches(0.55), Inches(0.16), size=7.2, color=NAVY, bold=True)
        text(slide, short_label(f"r/{ev.get('subreddit', '')}", 17), Inches(1.38), Inches(y + 0.10), Inches(1.25), Inches(0.26), size=5.8, color=TEAL, bold=True)
        text(slide, f"qualified score {ev.get('score', '')} | raw file retained in audit trail", Inches(2.64), Inches(y + 0.12), Inches(4.7), Inches(0.24), size=5.8, color=INK)
        y += 0.58
    side_panel(slide, "Audit trail", ["reddit-evidence-pack.json", "reddit-factcheck-report.md", "builder script retained"])
    footer(slide, 9, TOTAL_SLIDES)


def build() -> Path:
    if not TEMPLATE.exists():
        raise SystemExit(f"Branded template not found: {TEMPLATE}")
    claims = load_json(CLAIM_PACK)["claims"]
    companies = load_companies()
    evidence = load_json(EVIDENCE_PACK)
    rows = evidence["claims"]
    rows_by_id = {row["claim_id"]: row for row in rows}
    claims_by_id = {claim["id"]: claim for claim in claims}
    company_map = company_by_claim(companies)

    prs = Presentation(str(TEMPLATE))
    trim_to_first_n(prs, TOTAL_SLIDES)
    while len(prs.slides) < TOTAL_SLIDES:
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide in prs.slides:
        clear_to_template_base(slide)

    cover(prs.slides[0], claims, rows)
    executive_verdict(prs.slides[1], rows)
    method_slide(prs.slides[2])
    company_profiles(prs.slides[3], companies[:5], rows_by_id, 4, "Company Profiles: Operating Wedges")
    company_profiles(prs.slides[4], companies[5:], rows_by_id, 5, "Company Profiles: Finance, Claims, Construction")
    strong_evidence(prs.slides[5], claims_by_id, company_map, rows)
    weak_and_gap_slide(prs.slides[6], claims_by_id, company_map, rows)
    recommendations(prs.slides[7])
    appendix(prs.slides[8], rows)

    prs.save(OUT)
    problems = validate_pptx(OUT)
    if problems:
        raise SystemExit("PPTX validation failed:\n- " + "\n- ".join(problems))
    print(OUT)
    return OUT


if __name__ == "__main__":
    build()
