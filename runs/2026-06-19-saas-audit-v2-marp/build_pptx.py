"""SaaS Audit v2 PPTX — voice layer applied."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "saas-audit-v2.pptx")

BG      = RGBColor(0x08, 0x0b, 0x11)
BG_ELEV = RGBColor(0x0f, 0x14, 0x1d)
BG_MATH = RGBColor(0x07, 0x1e, 0x1b)
INK     = RGBColor(0xee, 0xf2, 0xf7)
SOFT    = RGBColor(0xae, 0xb8, 0xc7)
MUTED   = RGBColor(0x69, 0x74, 0x8a)
TEAL    = RGBColor(0x2d, 0xd4, 0xbf)
SKY     = RGBColor(0x38, 0xbd, 0xf8)
MAGENTA = RGBColor(0xe8, 0x79, 0xf9)
AMBER   = RGBColor(0xf6, 0xb9, 0x4b)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bar = s.shapes.add_shape(1, 0, H - Inches(0.06), W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    return s

def t(s, text, l, top, w, h, size=Pt(18), color=INK, bold=False, italic=False,
      align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(l, top, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = size; r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    return tb

def kicker(s, text, top=Inches(0.72)):
    t(s, text.upper(), Inches(0.8), top, Inches(11.5), Inches(0.35),
      size=Pt(9), color=TEAL, bold=True)

def heading(s, text, top=Inches(1.2), size=Pt(33), color=INK):
    t(s, text, Inches(0.8), top, Inches(11.5), Inches(1.4),
      size=size, color=color, bold=True, wrap=True)

def bullets(s, items, top=Inches(2.75), size=Pt(17), color=INK):
    tb = s.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = f"→  {item}"
        r.font.size = size; r.font.color.rgb = color
        p.space_after = Pt(5)

def pn(s, n):
    t(s, str(n), Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
      size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)

def col_box(s, head, title, items, tag, tc, lx, top=Inches(2.6), w=Inches(5.5), h=Inches(3.9)):
    b = s.shapes.add_shape(1, lx, top, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = BG_ELEV
    b.line.color.rgb = RGBColor(0x22, 0x2b, 0x38)
    t(s, head.upper(), lx+Inches(0.2), top+Inches(0.18), w-Inches(0.4), Inches(0.28),
      size=Pt(8), color=TEAL, bold=True)
    t(s, title, lx+Inches(0.2), top+Inches(0.52), w-Inches(0.4), Inches(0.45),
      size=Pt(17), color=INK, bold=True)
    tb = s.shapes.add_textbox(lx+Inches(0.2), top+Inches(1.08), w-Inches(0.4), Inches(2.3))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = f"→  {item}"
        r.font.size = Pt(13); r.font.color.rgb = SOFT
        p.space_after = Pt(3)
    t(s, tag.upper(), lx+Inches(0.2), top+h-Inches(0.48), w-Inches(0.4), Inches(0.28),
      size=Pt(8), color=tc, bold=True)

def two_col(s, lh, lt, li, ltag, ltc, rh, rt, ri, rtag, rtc, top=Inches(2.6)):
    col_box(s, lh, lt, li, ltag, ltc, Inches(0.8), top)
    col_box(s, rh, rt, ri, rtag, rtc, Inches(6.83), top)

def tbl(s, headers, rows, top=Inches(2.35), cws=None):
    nc = len(headers)
    cws = cws or [Inches(11.5/nc)]*nc
    lx = Inches(0.8); rh = Inches(0.53)
    x = lx
    for h, w in zip(headers, cws):
        b = s.shapes.add_shape(1, x, top, w, Inches(0.36))
        b.fill.solid(); b.fill.fore_color.rgb = BG_ELEV; b.line.fill.background()
        t(s, h.upper(), x+Inches(0.1), top+Inches(0.05), w-Inches(0.2), Inches(0.28),
          size=Pt(8), color=TEAL, bold=True)
        x += w
    for ri2, row in enumerate(rows):
        rt = top + Inches(0.36) + ri2*rh
        rb = BG_ELEV if ri2 % 2 == 0 else BG
        x = lx
        for i, (cell, w) in enumerate(zip(row, cws)):
            b2 = s.shapes.add_shape(1, x, rt, w, rh-Inches(0.03))
            b2.fill.solid(); b2.fill.fore_color.rgb = rb; b2.line.fill.background()
            # color logic
            cs = str(cell)
            color = INK
            if cs in ("REPLACE",): color = TEAL
            elif cs in ("KEEP",): color = MUTED
            elif cs in ("NEGOTIATE",): color = AMBER
            elif cs in ("CONSOLIDATE",): color = SKY
            elif cs in ("AUDIT USAGE",): color = MAGENTA
            elif cs.startswith("$") and cs not in ("$0","$/mo"):
                color = TEAL
            t(s, cs, x+Inches(0.1), rt+Inches(0.09), w-Inches(0.2), rh-Inches(0.12),
              size=Pt(13), color=color, bold=(color != INK and i == 0))
            x += w

def math_box(s, lines, top=Inches(2.7), highlight_lines=None):
    highlight_lines = highlight_lines or []
    bh = Inches(0.32) * len(lines) + Inches(0.5)
    b = s.shapes.add_shape(1, Inches(0.8), top, Inches(11.5), bh)
    b.fill.solid(); b.fill.fore_color.rgb = BG_MATH
    b.line.color.rgb = RGBColor(0x1a, 0x6b, 0x5f)
    tb = s.shapes.add_textbox(Inches(1.1), top+Inches(0.22), Inches(11.0), bh-Inches(0.35))
    tf = tb.text_frame; tf.word_wrap = False
    first = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = TEAL if i in highlight_lines else SOFT
        r.font.bold = i in highlight_lines
        p.space_after = Pt(2)

def pullquote(s, text, top=Inches(5.2)):
    bar = s.shapes.add_shape(1, Inches(0.7), top, Inches(0.05), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    t(s, text, Inches(1.0), top, Inches(11.3), Inches(1.0),
      size=Pt(18), color=SOFT, italic=True, wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Cover
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "SaaS Replacement Audit · 15 tools · $33,288/yr", top=Inches(1.05))
t(s, "$30,500 is sitting\nin your SaaS stack.",
  Inches(0.8), Inches(1.65), Inches(11.5), Inches(2.4),
  size=Pt(50), color=TEAL, bold=True, wrap=True)
t(s, "$16,318 of your annual SaaS spend is replaceable or renegotiable.\nHere's the exit plan — month by month, with the math.",
  Inches(0.8), Inches(4.2), Inches(10.5), Inches(1.1), size=Pt(19), color=SOFT, wrap=True)
pn(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "The Problem")
heading(s, "You're paying $1,360/mo for things you could own")
bullets(s, [
    "Retool: $600/mo to display a table your engineers built in a day",
    "Zapier: $199/mo for 10 HTTP calls between services you own",
    "Airtable: $240/mo for a spreadsheet masquerading as a CRM",
    "Mixpanel: $89/mo for data you're already paying Segment to collect",
], top=Inches(2.65))
pullquote(s, '"You don\'t have a SaaS problem. You have a build-vs-buy decision you never made."')
pn(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: Classification
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "15 Tools · 5 Verdicts")
heading(s, "Only 5 tools actually earn their invoice", size=Pt(30))
tbl(s,
    ["Tool", "$/mo", "Verdict", "One-line reason"],
    [
        ["Retool",   "$600", "REPLACE",     "Paying $600/mo to host a table. Own it for $8,700 once."],
        ["Airtable", "$240", "REPLACE",     "Spreadsheet as CRM. A Supabase schema and an afternoon."],
        ["Zapier",   "$199", "REPLACE",     "10 HTTP calls. n8n, $5/mo VPS, 3 days."],
        ["Segment",  "$120", "REPLACE",     "200 lines of code + Supabase insert. Paying for a pipe."],
        ["Intercom", "$299", "NEGOTIATE",   "Crisp exists at $25/mo. One call saves $718/yr."],
        ["Datadog",  "$410", "NEGOTIATE",   "Vercel + Supabase covers 80%. Push for startup tier."],
        ["Mixpanel", "$89",  "CONSOLIDATE", "Paying twice for event tracking. Cancel this week."],
        ["Algolia",  "$500", "KEEP",        "Search is core. Supabase FTS breaks at 50K rows."],
    ],
    top=Inches(2.2),
    cws=[Inches(1.5), Inches(0.85), Inches(1.7), Inches(7.2)],
)
pn(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: Top 3 ranked
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "Replace · Ranked by Savings")
heading(s, "3 builds. $26,839 back over 3 years.", size=Pt(32))
tbl(s,
    ["#", "Replace", "Annual SaaS", "Build cost", "Break-even", "3-yr savings"],
    [
        ["1", "Retool → Next.js admin panel", "$7,200", "$8,700",        "Month 17",              "$15,132"],
        ["2", "Zapier → n8n self-hosted",     "$2,388", "$1,380",        "Month 8 ← fastest",    "$6,524"],
        ["3", "Airtable → Supabase CRM",      "$2,880", "$4,350→$1,500*","Month 22→Month 7*",     "$5,183"],
    ],
    top=Inches(2.5),
    cws=[Inches(0.5), Inches(3.5), Inches(1.7), Inches(1.9), Inches(2.0), Inches(1.7)],
)
t(s, "*Build #3 alongside #1 — shared scaffold cuts $2,850 off the cost and moves break-even from month 22 to month 7.",
  Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.5), size=Pt(13), color=MUTED, italic=True)
pn(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Math — Zapier (fastest win)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "The Math · Build #2 · Fastest Win")
heading(s, "Zapier → n8n: $6,524 back, break-even month 8")
math_box(s, [
    "3-Year SaaS:   $199 × 12 = $2,388  →  × 1.10 = $2,627  →  × 1.10 = $2,890  →  Total: $7,905",
    "",
    "Build cost:    n8n VPS ($5/mo × 36)         =  $180",
    "               Migration (8hr × $150/hr)     =  $1,200",
    "               Total: $1,380",
    "",
    "Break-even:    $1,380 ÷ ($199 − $5) = 7.1 months",
    "After month 8: Zapier never touches your bank account again.",
], top=Inches(2.65), highlight_lines=[0, 4, 6, 7])
pullquote(s, '"This is the fastest money in the stack. 3 days of work. Ship it this week."', top=Inches(5.5))
pn(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Math — Retool (highest value)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "The Math · Build #1 · Highest Value")
heading(s, "Retool → Next.js admin panel: $15,132 back")
math_box(s, [
    "3-Year SaaS:   $600 × 12 = $7,200  →  × 1.10 = $7,920  →  × 1.10 = $8,712  →  Total: $23,832",
    "",
    "Build cost:    40hr × $150/hr = $6,000 one-time",
    "               Maintenance (15%/yr × 3yr) = $2,700",
    "               Hosting: $0 — existing Vercel deployment",
    "               Total: $8,700",
    "",
    "Break-even:    $8,700 ÷ ($600 − $75/mo maint) = Month 17",
    "After month 17: you own this forever.",
], top=Inches(2.65), highlight_lines=[0, 5, 7, 8])
t(s, "Stack: Next.js 14 + Supabase RLS + shadcn/ui + NextAuth admin role · 2 weeks · 1 engineer",
  Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.5), size=Pt(13), color=MUTED)
pn(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Build plan — Retool
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "Build Plan #1 · Retool Replacement")
heading(s, "Ship the 3 things CS actually uses — nothing else")
two_col(s,
    "What ships in 30 days", "Not all 10 Retool apps. These 3.",
    ["User list — search, filter, impersonate",
     "Subscription management — plan, status, cancel",
     "Activity log — what they did, when"],
    "2 WEEKS · 1 ENGINEER · $0 HOSTING", TEAL,
    "Migration — 4 steps, no data loss", "Audit → Build → Shadow → Cancel",
    ["Audit: which Retool apps get opened weekly",
     "Build: replicate those 2–3, nothing else",
     "Shadow: run both for 2 weeks",
     "CS signs off → cancel Retool"],
    "DAY 1: AUDIT.  MONTH 2: CANCEL.", MAGENTA,
)
pn(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: Quick wins — zero code
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "Quick Wins · Zero Code")
heading(s, "$4,450/yr this month — no engineers required", size=Pt(30))
bullets(s, [
    "Cancel Mixpanel this week. Segment already has this data. $1,068/yr gone immediately.",
    "Downgrade Postman to free. Email support. If someone objects, buy them one seat ($49/yr). Save $539.",
    "Call Intercom. Say Crisp costs $25/mo. They will negotiate. Target: $239/mo. Save $718/yr.",
    "Call Datadog. Say you're evaluating Grafana Cloud. Target: 30% reduction. Save $1,476/yr.",
    "Migrate to Resend — or call SendGrid and threaten to. Save $600/yr either way.",
], top=Inches(2.6), size=Pt(16))
pn(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: 12-month action plan
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "The Full Picture")
heading(s, "From $33,288 to $16,370 — in 12 months", size=Pt(30))
tbl(s,
    ["Month", "The move", "Tool", "Saves"],
    [
        ["This week",    "Cancel. No migration.",                       "Mixpanel",         "$1,068/yr"],
        ["Week 2",       "One email to support.",                       "Postman",           "$588/yr"],
        ["Month 1",      "One call each. Name the competitor price.",   "Intercom + Datadog","$2,194/yr"],
        ["Months 2–3",   "Migrate to n8n. 3 days. Fastest money here.", "Zapier",           "$2,388/yr from mo 4"],
        ["Months 3–7",   "Build Next.js admin panel.",                  "Retool",           "$7,200/yr from mo 8"],
        ["Months 4–7",   "Build Supabase CRM alongside it.",            "Airtable",         "$2,880/yr from mo 8"],
        ["Months 8–12",  "Build event pipeline when admin panel ships.", "Segment",          "$1,440/yr from mo 12"],
    ],
    top=Inches(2.25),
    cws=[Inches(1.9), Inches(4.3), Inches(2.7), Inches(2.4)],
)
pn(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: Closing
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
kicker(s, "Next Action", top=Inches(1.05))
t(s, "Start with Zapier.\nThis week.",
  Inches(0.8), Inches(1.65), Inches(11.5), Inches(2.4),
  size=Pt(52), color=TEAL, bold=True, wrap=True)
t(s, "3 days. $1,380 total cost. $6,524 back over 3 years. Break-even: month 8.\n"
     "n8n on a $5/mo Hetzner VPS. One Docker command. The rest is just recreating HTTP calls you already own.",
  Inches(0.8), Inches(4.2), Inches(11.0), Inches(1.4), size=Pt(18), color=SOFT, wrap=True)
pn(s, 10)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: 10")
