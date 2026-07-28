import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
from copy import deepcopy
from pathlib import Path

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise unittest.SkipTest(f"optional PPTX QA dependencies unavailable: {exc}") from exc


ROOT = Path(__file__).resolve().parents[1]
SKILL = ROOT / "skills" / "pptx-design-quality"


def load_linter():
    path = SKILL / "scripts" / "lint_pptx.py"
    spec = importlib.util.spec_from_file_location("pptx_design_lint", path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def add_text(slide, text, left, top, width, height, size):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.text = text
    run = shape.text_frame.paragraphs[0].runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(20, 30, 40)
    return shape


def build_deck(path, *, defective=False):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for index in range(2):
        slide = presentation.slides.add_slide(blank)
        add_text(slide, f"Action title {index + 1}", 0.6, 0.4, 8, 0.7, 30)
        add_text(slide, "Concise supporting evidence.", 0.6, 1.5, 5, 1.0, 16)
    if defective:
        slide = presentation.slides[0]
        add_text(slide, "Tiny text", 12.8, 7.2, 1.2, 0.6, 7)
    presentation.save(path)


def fill_brief(path):
    path.write_text(
        """# Deck brief

## Audience
Executive sponsors familiar with the program.

## Decision
Approve the recommended pilot.

## Narrative promise
The evidence supports a narrow, measurable launch.

## Voice and tone
Direct, calm, and evidence-led.

## Anti-references
Avoid generic gradients and dense dashboard walls.

## Evidence and editability
Use cited primary evidence and editable native charts.

## Success criteria
The decision and next step are unmistakable.
"""
    )


def test_context_initializer_and_validator(tmp_path):
    run_dir = tmp_path / "run"
    init_script = SKILL / "scripts" / "init_deck_context.py"
    validate_script = SKILL / "scripts" / "validate_deck_context.py"

    created = subprocess.run(
        [
            sys.executable,
            str(init_script),
            "--run",
            str(run_dir),
            "--title",
            "Test deck",
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    assert created.returncode == 0
    assert (run_dir / "deck-brief.md").exists()
    design = json.loads((run_dir / "deck-design.json").read_text())
    assert design["deck"]["name"] == "Test deck"

    untouched = subprocess.run(
        [
            sys.executable,
            str(validate_script),
            str(run_dir / "deck-brief.md"),
            str(run_dir / "deck-design.json"),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    assert untouched.returncode == 1
    assert "empty section" in untouched.stderr

    fill_brief(run_dir / "deck-brief.md")
    validated = subprocess.run(
        [
            sys.executable,
            str(validate_script),
            str(run_dir / "deck-brief.md"),
            str(run_dir / "deck-design.json"),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    assert validated.returncode == 0, validated.stderr


def test_context_validator_requires_waiver_reason(tmp_path):
    run_dir = tmp_path / "run"
    subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "init_deck_context.py"),
            "--run",
            str(run_dir),
            "--title",
            "Test deck",
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    design_path = run_dir / "deck-design.json"
    design = json.loads(design_path.read_text())
    design["qa"]["ignore_rules"] = ["TEXT_TOO_SMALL"]
    design_path.write_text(json.dumps(design))

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "validate_deck_context.py"),
            str(run_dir / "deck-brief.md"),
            str(design_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    assert result.returncode == 1
    assert "missing reason for ignored rule TEXT_TOO_SMALL" in result.stderr


def test_context_validator_rejects_whitespace_waiver_reason(tmp_path):
    run_dir = tmp_path / "run"
    subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "init_deck_context.py"),
            "--run",
            str(run_dir),
            "--title",
            "Test deck",
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    fill_brief(run_dir / "deck-brief.md")
    design_path = run_dir / "deck-design.json"
    design = json.loads(design_path.read_text())
    design["qa"]["ignore_rules"] = ["TEXT_TOO_SMALL"]
    design["qa"]["waivers"] = {"TEXT_TOO_SMALL": "   "}
    design_path.write_text(json.dumps(design))

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "validate_deck_context.py"),
            str(run_dir / "deck-brief.md"),
            str(design_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "TEXT_TOO_SMALL" in result.stderr


def test_linter_accepts_a_clean_deck(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "clean.pptx"
    build_deck(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert report["status"] == "clean"
    assert report["findings"] == []


def test_linter_reports_small_and_out_of_bounds_text(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "defective.pptx"
    build_deck(deck_path, defective=True)

    report = linter.lint_presentation(deck_path, fast=True)
    rule_ids = {finding["rule_id"] for finding in report["findings"]}

    assert report["status"] == "findings"
    assert "TEXT_TOO_SMALL" in rule_ids
    assert "SHAPE_OUT_OF_BOUNDS" in rule_ids


def test_linter_rejects_empty_and_blank_presentations(tmp_path):
    linter = load_linter()
    empty_path = tmp_path / "empty.pptx"
    empty = Presentation()
    empty.slide_width = Inches(13.333)
    empty.slide_height = Inches(7.5)
    empty.save(empty_path)

    empty_report = linter.lint_presentation(
        empty_path, config={"deck": {"output_mode": "image-per-slide"}}, fast=True
    )
    assert any(
        finding["rule_id"] == "DECK_NO_SLIDES"
        for finding in empty_report["findings"]
    )

    blank_path = tmp_path / "blank.pptx"
    blank = Presentation()
    blank.slide_width = Inches(13.333)
    blank.slide_height = Inches(7.5)
    blank.slides.add_slide(blank.slide_layouts[6])
    blank.save(blank_path)

    blank_report = linter.lint_presentation(
        blank_path, config={"deck": {"output_mode": "image-per-slide"}}, fast=True
    )
    assert any(
        finding["rule_id"] == "SLIDE_EMPTY"
        for finding in blank_report["findings"]
    )


def test_linter_honors_explicit_rule_waivers(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "defective.pptx"
    build_deck(deck_path, defective=True)
    config = {
        "qa": {
            "ignore_rules": ["TEXT_TOO_SMALL"],
            "waivers": {"TEXT_TOO_SMALL": "Approved source footnote."},
        }
    }

    report = linter.lint_presentation(deck_path, config=config, fast=True)

    assert all(
        finding["rule_id"] != "TEXT_TOO_SMALL" for finding in report["findings"]
    )
    assert any(
        finding["rule_id"] == "TEXT_TOO_SMALL"
        for finding in report["ignored_findings"]
    )


def test_linter_does_not_waive_errors_or_undocumented_warnings(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "defective.pptx"
    build_deck(deck_path, defective=True)
    config = {
        "qa": {
            "ignore_rules": ["SHAPE_OUT_OF_BOUNDS", "TEXT_TOO_SMALL"],
            "waivers": {"SHAPE_OUT_OF_BOUNDS": "Cannot waive an error."},
        }
    }

    report = linter.lint_presentation(deck_path, config=config, fast=True)
    active_rule_ids = {finding["rule_id"] for finding in report["findings"]}

    assert "SHAPE_OUT_OF_BOUNDS" in active_rule_ids
    assert "TEXT_TOO_SMALL" in active_rule_ids
    assert report["status"] == "findings"


def test_linter_resolves_inherited_placeholder_font_size(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "inherited.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Inherited title"
    slide.placeholders[1].text = "Inherited body text"
    presentation.save(deck_path)

    report = linter.lint_presentation(
        deck_path,
        config={
            "layout": {"aspect_ratio": "4:3"},
            "typography": {"min_body_font_pt": 40},
        },
        fast=True,
    )

    assert any(
        finding["rule_id"] == "TEXT_TOO_SMALL"
        and finding["shape"] == "Content Placeholder 2"
        for finding in report["findings"]
    )


def test_linter_recurses_into_grouped_shapes(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "grouped.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    add_text(group, "Grouped title", 0.6, 0.4, 8, 0.7, 30)
    add_text(group, "Grouped tiny text", 0.6, 1.5, 4, 0.5, 7)
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)
    rule_ids = {finding["rule_id"] for finding in report["findings"]}

    assert "SLIDE_MISSING_TITLE" not in rule_ids
    assert "TEXT_TOO_SMALL" in rule_ids


def test_linter_applies_group_transform_to_child_bounds(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "scaled-group.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    add_text(group, "Grouped title", 0, 0, 2, 0.7, 30)
    add_text(group, "Evidence", 2.2, 0, 2, 0.7, 16)
    group.left = Inches(12.5)
    group.top = Inches(0.4)
    group.width = Inches(3)
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert any(
        finding["rule_id"] == "SHAPE_OUT_OF_BOUNDS"
        for finding in report["findings"]
    )


def test_linter_applies_group_flip_to_child_bounds(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "flipped-group.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    add_text(group, "Flipped child", 0, 0, 0.5, 0.7, 30)
    group.left = Inches(10.5)
    group.top = Inches(0.4)
    group.width = Inches(3)
    group._element.xfrm.flipH = True
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert any(
        finding["rule_id"] == "SHAPE_OUT_OF_BOUNDS"
        for finding in report["findings"]
    )


def test_overflow_estimator_preserves_paragraph_boundaries(tmp_path):
    linter = load_linter()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = add_text(slide, "one", 0.5, 0.5, 6, 0.5, 16)
    for _ in range(8):
        paragraph = shape.text_frame.add_paragraph()
        paragraph.text = "one"
        paragraph.runs[0].font.size = Pt(16)

    assert linter._overflow_risk(shape, slide)


def test_linter_inspects_native_table_cells(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "table.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_text(slide, "Table title", 0.6, 0.4, 8, 0.7, 30)
    table = slide.shapes.add_table(
        2, 2, Inches(0.8), Inches(1.5), Inches(8), Inches(2)
    ).table
    for row in table.rows:
        for cell in row.cells:
            cell.text = "Dense evidence"
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert any(
        finding["rule_id"] == "TEXT_TOO_SMALL" and "cell" in finding["shape"]
        for finding in report["findings"]
    )


def test_linter_uses_merged_table_cell_span_geometry(tmp_path):
    linter = load_linter()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        2, 3, Inches(1), Inches(1), Inches(9), Inches(2)
    )
    table_shape.table.cell(0, 0).merge(table_shape.table.cell(1, 2))
    view = linter.ShapeView(
        table_shape,
        table_shape.left,
        table_shape.top,
        table_shape.width,
        table_shape.height,
        0,
    )

    cells = list(linter._table_cells(view))
    merged = next(cell for cell in cells if cell.name.endswith("cell 1,1"))

    assert merged.width == table_shape.width
    assert merged.height == table_shape.height


def test_linter_uses_visible_pixels_for_cropped_image_dpi(tmp_path):
    linter = load_linter()
    image_path = tmp_path / "source.png"
    Image.new("RGB", (600, 600), "white").save(image_path)
    deck_path = tmp_path / "cropped.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_text(slide, "Image title", 0.6, 0.4, 8, 0.7, 30)
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.5), Inches(4), Inches(4)
    )
    picture.crop_left = 0.25
    picture.crop_right = 0.25
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "IMAGE_LOW_DPI" for finding in report["findings"]
    )


def test_linter_applies_caption_font_threshold_below_visual(tmp_path):
    linter = load_linter()
    image_path = tmp_path / "visual.png"
    Image.new("RGB", (1200, 600), "white").save(image_path)
    deck_path = tmp_path / "caption.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_text(slide, "Visual title", 0.6, 0.4, 8, 0.7, 30)
    slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.5), Inches(6), Inches(3)
    )
    caption = add_text(
        slide, "Figure 1. Supporting evidence.", 1, 4.55, 6, 0.3, 10
    )
    caption.name = "Figure caption"
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert all(
        not (
            finding["rule_id"] == "TEXT_TOO_SMALL"
            and finding["shape"] == "Figure caption"
        )
        for finding in report["findings"]
    )


def test_linter_resolves_explicit_theme_colors_for_contrast(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "theme-color.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = add_text(slide, "Theme title", 0.6, 0.4, 8, 0.7, 30)
    run = shape.text_frame.paragraphs[0].runs[0]
    run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_honors_slide_color_map_override(tmp_path):
    linter = load_linter()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    master_map = slide.slide_layout.slide_master._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}clrMap"
    )
    override = slide._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr"
    )
    for child in list(override):
        override.remove(child)
    explicit = OxmlElement("a:overrideClrMapping")
    for key, value in master_map.items():
        explicit.set(key, value)
    explicit.set("accent1", "accent2")
    override.append(explicit)

    assert linter._theme_color(slide, "accent1") == linter._theme_color(
        slide, "accent2"
    )


def test_linter_uses_underlying_shape_fill_for_contrast(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "underlying-fill.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(8), Inches(0.7)
    )
    panel.fill.solid()
    panel.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    shape = add_text(slide, "Transparent title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.theme_color = (
        MSO_THEME_COLOR.ACCENT_1
    )
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_does_not_treat_small_accent_as_text_background(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "small-accent.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape = add_text(slide, "Dark title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_resolves_inherited_slide_background_for_contrast(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "inherited-background.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = add_text(slide, "Invisible title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255
    )
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_resolves_explicit_master_background_for_contrast(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "master-background.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.slide_masters[0].background.fill.solid()
    presentation.slide_masters[0].background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = add_text(slide, "Readable title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255
    )
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert all(
        finding["rule_id"] != "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_uses_master_shape_as_text_background(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "master-shape-background.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    master_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height
    )
    master_panel.fill.solid()
    master_panel.fill.fore_color.rgb = RGBColor(0, 0, 0)
    presentation.slide_masters[0].shapes._spTree.append(
        deepcopy(master_panel._element)
    )
    slide.shapes._spTree.remove(master_panel._element)
    shape = add_text(slide, "Readable title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255
    )
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert all(
        finding["rule_id"] != "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_looks_through_transparent_solid_fills_for_contrast(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "transparent-fill.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape = add_text(slide, "Invisible title", 0.6, 0.4, 8, 0.7, 30)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    color_element = shape._element.spPr.solidFill[0]
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "0")
    color_element.append(alpha)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert any(
        finding["rule_id"] == "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_skips_contrast_when_text_overlays_a_picture(tmp_path):
    linter = load_linter()
    image_path = tmp_path / "background.png"
    Image.new("RGB", (1600, 900), "black").save(image_path)
    deck_path = tmp_path / "picture-background.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
    )
    shape = add_text(slide, "Readable white title", 0.6, 0.4, 8, 0.7, 30)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255
    )
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=False)

    assert all(
        finding["rule_id"] != "TEXT_LOW_CONTRAST"
        for finding in report["findings"]
    )


def test_linter_accounts_for_rotated_rendered_bounds(tmp_path):
    linter = load_linter()
    deck_path = tmp_path / "rotated.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = add_text(slide, "Rotated title", 0.1, 0.1, 3, 0.5, 30)
    shape.rotation = 45
    presentation.save(deck_path)

    report = linter.lint_presentation(deck_path, fast=True)

    assert any(
        finding["rule_id"] == "SHAPE_OUT_OF_BOUNDS"
        for finding in report["findings"]
    )


def test_image_per_slide_mode_skips_native_title_and_repetition_rules(tmp_path):
    linter = load_linter()
    image_path = tmp_path / "slide.png"
    Image.new("RGB", (1600, 900), "white").save(image_path)
    deck_path = tmp_path / "flattened.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for _ in range(3):
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
        )
    presentation.save(deck_path)

    report = linter.lint_presentation(
        deck_path, config={"deck": {"output_mode": "image-per-slide"}}, fast=True
    )

    rule_ids = {finding["rule_id"] for finding in report["findings"]}
    assert "SLIDE_MISSING_TITLE" not in rule_ids
    assert "LAYOUT_REPETITION" not in rule_ids


def test_linter_cli_reports_invalid_powerpoint_without_traceback(tmp_path):
    invalid = tmp_path / "invalid.pptx"
    invalid.write_text("not a PowerPoint package")

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "lint_pptx.py"),
            str(invalid),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "lint failed:" in result.stderr
    assert "Traceback" not in result.stderr


def test_linter_cli_reports_invalid_config_without_traceback(tmp_path):
    deck_path = tmp_path / "clean.pptx"
    build_deck(deck_path)
    config_path = tmp_path / "invalid-config.json"
    config_path.write_text('{"layout": {"max_words_per_slide": "many"}}')

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL / "scripts" / "lint_pptx.py"),
            str(deck_path),
            "--config",
            str(config_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "lint failed:" in result.stderr
    assert "Traceback" not in result.stderr


def test_html_design_gate_is_version_pinned():
    wrapper = (ROOT / "scripts" / "design-qa-detect.sh").read_text()

    assert "impeccable@3.4.0" in wrapper
    assert "impeccable@latest" not in wrapper


def test_integrated_skills_use_cross_host_quality_resolver():
    quality_skill = (SKILL / "SKILL.md").read_text()
    assert "PPTX_DESIGN_QUALITY_DIR" in quality_skill
    assert "CODEX_SKILLS_HOME" in quality_skill
    assert "CLAUDE_SKILLS_HOME" in quality_skill

    for relative_path in (
        "skills/branded-pptx-deck/SKILL.md",
        "skills/genspark-branded-deck/SKILL.md",
        "skills/pptx-visual-spec/portable-skills/marp/SKILL.md",
    ):
        integrated = (ROOT / relative_path).read_text()
        assert '$PPTX_DESIGN_QUALITY_DIR/scripts/lint_pptx.py' in integrated
        assert (
            '${CONTENT_IDEAS_DIR:-$HOME/content-ideas}/skills/'
            "pptx-design-quality/scripts/lint_pptx.py"
        ) not in integrated


def load_tests(loader, tests, pattern):
    """Expose the pytest-style regression functions to unittest-based CI."""
    suite = unittest.TestSuite()
    for name, function in sorted(globals().items()):
        if not name.startswith("test_") or not callable(function):
            continue

        def run_test(function=function):
            if function.__code__.co_argcount:
                with tempfile.TemporaryDirectory() as directory:
                    function(Path(directory))
            else:
                function()

        suite.addTest(unittest.FunctionTestCase(run_test, description=name))
    return suite
