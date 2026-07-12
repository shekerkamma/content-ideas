import importlib.util
import json
import sys
from pathlib import Path
from zipfile import ZipFile


REPO = Path(__file__).resolve().parents[1]
SCRIPT = REPO / "skills" / "reddit-new-factcheck" / "scripts" / "prepare_factcheck.py"
ANNOTATE_SCRIPT = REPO / "skills" / "reddit-new-factcheck" / "scripts" / "annotate_input_deck.py"
SKILL = REPO / "skills" / "reddit-new-factcheck" / "SKILL.md"


def load_module():
    spec = importlib.util.spec_from_file_location("prepare_factcheck", SCRIPT)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def load_annotate_module():
    spec = importlib.util.spec_from_file_location("annotate_input_deck", ANNOTATE_SCRIPT)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def test_markdown_claim_pack_and_report(tmp_path):
    mod = load_module()
    doc = tmp_path / "research.md"
    doc.write_text(
        "# Deck\n\n"
        "Gartner predicts that 40% of agentic AI projects will be cancelled by 2027.\n\n"
        "Practitioners complain that agent governance is unclear and risky.\n",
        encoding="utf-8",
    )
    out_dir = tmp_path / "out"

    assert mod.main(["--input", str(doc), "--topic", "agentic AI governance", "--out-dir", str(out_dir)]) == 0

    claims = json.loads((out_dir / "claim-pack.json").read_text(encoding="utf-8"))["claims"]
    assert len(claims) == 2
    assert claims[0]["primary_source_required"] is True
    assert claims[1]["evidence_need"] == "reddit_evidence"
    report = (out_dir / "reddit-factcheck-report.md").read_text(encoding="utf-8")
    assert "Reddit Fact-Check Report" in report
    assert "primary-source-required" in report


def test_reddit_evidence_pack_scores_thread_json(tmp_path):
    mod = load_module()
    doc = tmp_path / "research.md"
    doc.write_text(
        "Operators say agent governance and kill switches are blockers for production AI agents.\n",
        encoding="utf-8",
    )
    thread = tmp_path / "thread.json"
    thread.write_text(
        json.dumps(
            {
                "post": {"title": "AI agent governance", "content": "How do you handle kill switches?", "score": 42},
                "comments": [
                    {
                        "id": "t1_a",
                        "author": "user",
                        "text": "We cannot ship AI agents without governance, audit logs, and kill switches.",
                        "score": 18,
                        "depth": 0,
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    out_dir = tmp_path / "out"

    assert (
        mod.main(
            [
                "--input",
                str(doc),
                "--topic",
                "AI agent governance",
                "--out-dir",
                str(out_dir),
                "--reddit-json",
                str(thread),
            ]
        )
        == 0
    )

    pack = json.loads((out_dir / "reddit-evidence-pack.json").read_text(encoding="utf-8"))
    assert pack["claims"][0]["suggested_reddit_verdict"] == "qualified_reddit_support"
    assert pack["claims"][0]["evidence"][0]["id"] in {"post", "t1_a"}
    assert pack["claims"][0]["evidence"][0]["qualification"]["qualified"] is True


def test_reddit_evidence_gate_rejects_unrelated_subreddit_noise(tmp_path):
    mod = load_module()
    doc = tmp_path / "research.md"
    doc.write_text(
        "CentralComs is an AI agents product for property managers, focused on maintenance calls, ticket chasing, and owner, tenant, and vendor updates.\n",
        encoding="utf-8",
    )
    thread = tmp_path / "thread.json"
    thread.write_text(
        json.dumps(
            {
                "post": {
                    "title": "CentralComs movie trailer reaction",
                    "content": "The property in this scene had a tenant joke and a vendor cameo, but this is about a film.",
                    "score": 400,
                    "subreddit": "movies",
                },
                "comments": [],
            }
        ),
        encoding="utf-8",
    )
    out_dir = tmp_path / "out"

    assert (
        mod.main(
            [
                "--input",
                str(doc),
                "--topic",
                "property management",
                "--out-dir",
                str(out_dir),
                "--reddit-json",
                str(thread),
            ]
        )
        == 0
    )

    pack = json.loads((out_dir / "reddit-evidence-pack.json").read_text(encoding="utf-8"))
    claim = pack["claims"][0]
    assert claim["suggested_reddit_verdict"] == "no_qualified_reddit_evidence"
    assert claim["evidence"] == []
    assert claim["items_rejected_by_qualification_gate"] >= 1


def test_reddit_evidence_gate_accepts_operator_workflow_pain(tmp_path):
    mod = load_module()
    doc = tmp_path / "research.md"
    doc.write_text(
        "Property managers are overloaded by fragmented maintenance, ticket, and vendor coordination across AppFolio, Buildium, and Yardi.\n",
        encoding="utf-8",
    )
    thread = tmp_path / "thread.json"
    thread.write_text(
        json.dumps(
            {
                "post": {
                    "title": "What software do you use to manage rentals?",
                    "content": "I manage rental properties and I'm drowning in fragmented spreadsheets for AppFolio, Buildium, and Yardi exports, rent tracking, tenant info, maintenance tickets, vendor coordination, and expense logs.",
                    "score": 11,
                    "subreddit": "Landlord",
                },
                "comments": [
                    {
                        "id": "t1_prop",
                        "text": "My tenants put in maintenance tickets and I still track rent, vendors, and expenses in a spreadsheet.",
                        "score": 3,
                        "subreddit": "Landlord",
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    out_dir = tmp_path / "out"

    assert (
        mod.main(
            [
                "--input",
                str(doc),
                "--topic",
                "property management",
                "--out-dir",
                str(out_dir),
                "--reddit-json",
                str(thread),
            ]
        )
        == 0
    )

    pack = json.loads((out_dir / "reddit-evidence-pack.json").read_text(encoding="utf-8"))
    claim = pack["claims"][0]
    assert claim["suggested_reddit_verdict"] == "weak_qualified_reddit_support"
    assert claim["evidence"][0]["qualification"]["profile"] == "property_management"


def test_pptx_text_extraction(tmp_path):
    mod = load_module()
    pptx = tmp_path / "deck.pptx"
    slide_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:cSld><p:spTree><p:sp><p:txBody><a:p><a:r><a:t>AI agents need audit logs before deployment.</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld>
    </p:sld>"""
    with ZipFile(pptx, "w") as zf:
        zf.writestr("ppt/slides/slide1.xml", slide_xml)
    segments, method = mod.read_document(pptx)
    assert method == "pptx-text"
    assert segments[0].locator == "Slide 1"
    assert "audit logs" in segments[0].text


def test_annotate_input_deck_creates_review_copy(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    mod = load_annotate_module()
    pptx = tmp_path / "input.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(5), Inches(1)).text = (
        "Operators say agent governance blocks deployments."
    )
    prs.save(pptx)

    claim_pack = tmp_path / "claim-pack.json"
    claim_pack.write_text(
        json.dumps(
            {
                "claims": [
                    {
                        "id": "claim-001",
                        "claim": "Operators say agent governance blocks deployments.",
                        "locator": "Slide 1",
                        "reddit_query": "agent governance deployment blockers",
                        "primary_source_required": False,
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    evidence_pack = tmp_path / "reddit-evidence-pack.json"
    evidence_pack.write_text(
        json.dumps(
            {
                "claims": [
                    {
                        "claim_id": "claim-001",
                        "evidence_count": 1,
                        "suggested_reddit_verdict": "evidence_found",
                        "evidence": [
                            {
                                "source_json": "reddit-thread-1.json",
                                "excerpt": "Governance and approvals block our agent deployments.",
                            }
                        ],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    out = tmp_path / "review.pptx"

    summary = mod.annotate_deck(pptx, claim_pack, out, evidence_pack)

    assert summary["claims"] == 1
    assert summary["annotated_source_slides"] == 1
    assert out.exists()
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 4
    first_slide_text = "\n".join(
        run.text
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    second_slide_text = "\n".join(
        run.text
        for shape in prs.slides[1].shapes
        if getattr(shape, "has_text_frame", False)
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    assert "Fact Check Summary" in first_slide_text
    assert "What Changed In The Deck" in second_slide_text
    with ZipFile(out) as zf:
        slide_text = "\n".join(
            zf.read(name).decode("utf-8")
            for name in zf.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
    assert "FACT CHECK" in slide_text
    assert "Reddit-supported" in slide_text
    assert "Fact Check Notes" in slide_text


def test_skill_contract_requires_branded_pptx_delivery():
    text = SKILL.read_text(encoding="utf-8")
    compact_text = " ".join(text.split())
    assert "annotated copy of the original PPTX" in text
    assert "Prepend two executive summary slides" in text
    assert "BRANDED_PPTX_TEMPLATE" in text
    assert "fail closed" in compact_text
    assert "Do not create a separate summary deck by default" in text
    assert "template-native slide generation is required" in compact_text.lower()
    assert "Reddit evidence qualification gate" in text
    assert "Raw Reddit search output is not evidence" in text
    assert "Markdown and JSON" in text
    assert "audit trail" in text


def test_skill_contract_requires_annotated_input_deck():
    text = SKILL.read_text(encoding="utf-8")
    compact_text = " ".join(text.split())
    assert "annotated copy of the original PPTX" in text
    assert "<input-slug>-reddit-factchecked-reviewed.pptx" in text
    assert "annotate_input_deck.py" in text
    assert "Never overwrite the original file" in text
    assert "Fact checks must be evidence-backed" in text
    assert "The annotated original deck is the only default client-facing artifact" in compact_text
    assert "what to keep, what to rewrite, and what to replace" in compact_text
